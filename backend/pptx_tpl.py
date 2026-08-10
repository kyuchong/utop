"""
고객사 양식으로 결과서를 찍는다.

전에는 빈 문서에서 표와 글자를 파이썬으로 그렸다. 그러면 고객사가 쓰는
글꼴·표선 두께·머리글·로고가 하나도 안 맞아서, 받는 쪽에서 다시 옮겨
붙여야 했다. 시험 100건이면 그 옮겨 붙이기가 하루다.

여기서는 **고객사가 준 pptx 를 그대로 열고, 그 안의 장을 복제해서 값만
채운다.** 서식은 손대지 않으므로 저쪽 눈에는 자기네 양식 그대로다.

고객사마다 양식이 다르다. `templates/pptx/` 에 파일을 두고 `TEMPLATES` 에
한 줄 적으면 늘어난다 — 채우는 자리(표의 몇 행 몇 열)는 양식마다 다르므로
그 지도도 함께 적는다.
"""

from __future__ import annotations

import copy
import re
from pathlib import Path
from typing import Any, Iterable

TPL_DIR = Path(__file__).resolve().parent / "templates" / "pptx"


class Spot:
    """
    값을 채울 자리 하나.

    `표 이름` 이 아니라 **몇 번째 표의 몇 행 몇 열** 로 가리킨다. 고객사
    양식에는 도형 이름이 없거나 `Text Box 46` 처럼 뜻 없는 이름이라,
    이름으로 찾으면 양식이 조금만 바뀌어도 엉뚱한 데를 채운다.
    """

    def __init__(self, tbl: int, row: int, col: int, keep: bool = False):
        self.tbl = tbl
        self.row = row
        self.col = col
        # 원래 글자를 지우지 않고 뒤에 잇는다(제목 칸처럼 라벨이 붙은 자리)
        self.keep = keep


# 고객사별 양식.
#
#  · file    — templates/pptx 안의 파일 이름
#  · first   — 시험 하나의 **첫 장**이 될 본보기 장 번호(0-기준)
#  · more    — 결과가 넘칠 때 덧붙일 **이어지는 장**의 본보기 번호
#  · spots   — 채울 자리. 아래 이름은 `render` 가 넘기는 값의 이름이다
TEMPLATES: dict[str, dict[str, Any]] = {
    "lguplus": {
        "label": "LG유플러스",
        "file": "lguplus.pptx",
        "first": 0,
        "more": 1,
        "spots": {
            # 1행: TC_ID │ 번호 │ REQ ID │ (요구사항) │ 시험항목 │ (이름)
            "tc_id": Spot(0, 0, 1),
            "req_id": Spot(0, 0, 4),
            "tc_name": Spot(0, 0, 6),
            # 3행 왼쪽 — 시험 규격
            "spec": Spot(0, 2, 0),
            # 5행 왼쪽 — 시험 방법
            "method": Spot(0, 4, 0),
            # 5행 오른쪽 — 시험 결과 (첫 장은 「뒷면 참조」)
            "result_head": Spot(0, 4, 6),
            # 6행 — 비고
            "note": Spot(0, 5, 2),
        },
        # 이어지는 장 — 2행은 「시험 결과」 **제목띠**이고 본문은 3행이다.
        # 3행 하나가 높이 11.78cm 로 가로 8칸이 통째로 합쳐진 한 판이라,
        # 결과가 길어도 여기에 다 들어간다. 제목띠에 넣으면 한 줄짜리
        # 띠에 결과가 욱여넣어져 표가 무너진다 — 처음에 그렇게 잡았었다.
        "more_spots": {
            "tc_id": Spot(0, 0, 1),
            "req_id": Spot(0, 0, 4),
            "tc_name": Spot(0, 0, 6),
            "result": Spot(0, 2, 0),
            "note": Spot(0, 3, 2),
        },
        # 이어지는 장의 결과 판이 담는 줄 수. 넘으면 장을 하나 더 만든다.
        "more_lines": 34,
    },
}


def list_templates() -> list[dict[str, str]]:
    """화면이 고를 수 있는 양식들. 파일이 없는 것은 빼고 준다."""
    out = []
    for key, t in TEMPLATES.items():
        if (TPL_DIR / str(t["file"])).exists():
            out.append({"id": key, "label": str(t.get("label") or key)})
    return out


# ── 장 복제 ────────────────────────────────────────────────────────
#
# python-pptx 에는 「장 복제」 가 없다. 장의 XML 을 통째로 베끼고, 그 장이
# 걸고 있는 것들(그림·차트)의 연결(rels)을 새 장에 다시 걸어 준다.
# 이것을 안 하면 그림이 사라지거나 엉뚱한 그림이 붙는다.
def clone_slide(prs, src):
    layout = src.slide_layout
    new = prs.slides.add_slide(layout)

    # 레이아웃이 넣어 준 빈 개체 틀을 걷어낸다 — 베껴 올 것과 겹쳐서
    # 빈 상자가 위에 뜬다
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)

    # 걸린 것(그림·차트)을 새 장에도 건다.
    #
    # 새 장에서 받는 rId 는 본보기 장의 것과 다를 수 있다. 베껴 온 XML 은
    # 옛 rId 를 가리키므로, **옛 rId → 새 rId** 로 갈아 끼워야 한다.
    # 안 갈면 그림이 사라지거나 엉뚱한 그림이 붙는다.
    ridmap: dict[str, str] = {}
    for rid, rel in src.part.rels.items():
        if rel.reltype.endswith("/slideLayout"):
            continue
        if rel.is_external:
            ridmap[rid] = new.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            ridmap[rid] = new.part.rels.get_or_add(rel.reltype, rel._target)

    R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for key in ("id", "embed", "link", "pict", "dm", "lo", "qs", "cs"):
                k = R + key
                if k in node.attrib and node.attrib[k] in ridmap:
                    node.attrib[k] = ridmap[node.attrib[k]]
        new.shapes._spTree.append(el)
    return new


def move_slide(prs, slide, to: int) -> None:
    """장 차례 옮기기 — 만든 것은 늘 맨 뒤에 붙으므로 제자리로 보낸다."""
    ids = prs.slides._sldIdLst
    RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    rid = None
    for r, rel in prs.part.rels.items():
        if rel.reltype.endswith("/slide") and not rel.is_external and rel._target is slide.part:
            rid = r
            break
    el = next((x for x in list(ids) if x.get(RID) == rid), None)
    if el is None:
        return
    ids.remove(el)
    ids.insert(to, el)


def drop_slide(prs, idx: int) -> None:
    """
    본보기로만 쓰고 결과에는 남기지 않을 장을 뺀다.

    목록(`sldIdLst`)에서만 빼면 장은 안 보이지만 **파일 안에는 남는다** —
    4장짜리 결과서에 6장이 들어 있게 된다. 연결(rel)까지 끊어야 저장할 때
    같이 빠진다.
    """
    ids = prs.slides._sldIdLst
    lst = list(ids)
    if not (0 <= idx < len(lst)):
        return
    el = lst[idx]
    rid = el.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    ids.remove(el)
    if rid:
        try:
            prs.part.drop_rel(rid)
        except Exception:
            pass


# ── 글자 채우기 ────────────────────────────────────────────────────
def _first_run(para):
    """그 문단의 첫 글자 조각. 서식(글꼴·크기·색)이 거기 붙어 있다."""
    return para.runs[0] if para.runs else None


def set_text(cell, text: str, keep: bool = False) -> None:
    """
    칸에 글자를 넣되 **서식을 지키며** 넣는다.

    `cell.text = ...` 로 넣으면 python-pptx 가 문단을 새로 만들면서 글꼴과
    크기를 잃는다. 고객사 양식은 그 서식이 전부라, 첫 조각의 서식을 그대로
    두고 글자만 갈아 끼운다. 줄이 여럿이면 첫 문단을 본떠 늘린다.
    """
    tf = cell.text_frame
    lines = str(text or "").split("\n")
    base = tf.paragraphs[0]
    run = _first_run(base)
    if run is None:
        run = base.add_run()
    if keep:
        run.text = (run.text or "") + lines[0]
    else:
        run.text = lines[0]
    # 첫 문단에 조각이 여럿이면 나머지는 지운다 — 옛 글자가 뒤에 남는다
    for extra in list(base.runs)[1:]:
        extra._r.getparent().remove(extra._r)
    # 옛 문단도 지운다
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    for ln in lines[1:]:
        np = copy.deepcopy(base._p)
        base._p.getparent().append(np)
        from pptx.text.text import _Paragraph

        para = _Paragraph(np, tf)
        r = _first_run(para)
        if r is None:
            r = para.add_run()
        r.text = ln
        for extra in list(para.runs)[1:]:
            extra._r.getparent().remove(extra._r)


def tables_of(slide) -> list:
    """이 장의 표들 — 그룹 안에 든 것까지."""
    out: list = []

    def walk(shapes: Iterable):
        for sh in shapes:
            if getattr(sh, "has_table", False):
                out.append(sh.table)
            if sh.shape_type == 6 or hasattr(sh, "shapes"):  # 그룹
                try:
                    walk(sh.shapes)
                except Exception:
                    pass

    walk(slide.shapes)
    return out


def fill(slide, spots: dict[str, Spot], values: dict[str, str]) -> None:
    """자리 지도대로 값을 채운다. 값이 없는 자리는 손대지 않는다."""
    tbls = tables_of(slide)
    for name, sp in spots.items():
        if name not in values:
            continue
        if sp.tbl >= len(tbls):
            continue
        tbl = tbls[sp.tbl]
        try:
            cell = tbl.cell(sp.row, sp.col)
        except Exception:
            continue
        set_text(cell, values[name], sp.keep)


def safe(s: str) -> str:
    return re.sub(r'[\\/:*?"<>|]+', "_", str(s or "")).strip() or "report"
