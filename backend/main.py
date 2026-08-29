import os
import sys
import json
import re
import asyncio
import httpx
import platform
import socket
import subprocess
import threading
import contextvars
from uuid import uuid4 as _uuid4
from datetime import datetime
from pathlib import Path
from typing import Optional
from dotenv import load_dotenv
load_dotenv(Path(__file__).parent.parent / ".env")

import paramiko
from fastapi import FastAPI, WebSocket, WebSocketDisconnect, HTTPException, UploadFile, File, Form, Request, Response
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from fastapi.middleware.gzip import GZipMiddleware
from pydantic import BaseModel
import anthropic
import engine
import db  # PostgreSQL 접속 레이어 (커넥션 풀 + CRUD 헬퍼)
import id_migrate  # 옛 ID → 모델그룹 기준 ID 옮기기

# ───────────────────────────────────────────
# 경로 설정
# ───────────────────────────────────────────
BASE_DIR = Path(__file__).parent.parent  # backend/ 의 상위 = nettest/
DATA_DIR = BASE_DIR / "data"
FRONTEND_DIR = BASE_DIR / "frontend"
DEVICES_FILE = DATA_DIR / "devices" / "devices.json"
PROCEDURES_FILE = DATA_DIR / "procedures" / "procedures.json"
RESULTS_DIR = DATA_DIR / "results"
LLMS_FILE = DATA_DIR / "integrations" / "llms.json"
CUSTOM_FIELDS_FILE = DATA_DIR / "config" / "custom_fields.json"
DEVICE_CATALOG_FILE = DATA_DIR / "state" / "device_catalog.json"
PERMISSIONS_FILE = DATA_DIR / "config" / "permissions.json"
PROMPTS_FILE = DATA_DIR / "config" / "prompts.json"
CONFLUENCE_FILE = DATA_DIR / "integrations" / "confluence.json"
JIRA_FILE = DATA_DIR / "integrations" / "jira.json"
DEFECT_CLASS_FILE = DATA_DIR / "config" / "issue_defect_class.json"   # 이슈키 → defect 분류(현장장애/상용망검증)
HELP_FILE = DATA_DIR / "config" / "help.json"
CHAT_SESS_FILE = DATA_DIR / "state" / "chat_sessions.json"
GLOBAL_PARAMS_FILE = DATA_DIR / "config" / "global_params.json"
RSC_MANPOWER_FILE = DATA_DIR / "state" / "manpower.json"
RSC_PROJECTS_FILE = DATA_DIR / "state" / "projects.json"

app = FastAPI(title="NetTest Automation")
# 응답 gzip 압축 (>= 500 bytes 자동) — JSON 은 압축률 매우 높음. 브라우저는 자동으로 Accept-Encoding: gzip 보냄.
# 단, SSE 스트리밍 경로는 gzip 대상에서 제외 — gzip 은 청크를 버퍼링해서 한꺼번에 flush 하므로 스트리밍이 죽음.
# 스트리밍 경로는 요청 시 Accept-Encoding 헤더를 서버 진입 직전에 제거해 GZipMiddleware 가 skip 하도록 유도한다.
_SSE_PATH_PREFIXES = ("/api/chat/local/stream", "/api/dify/chat", "/api/chat/stream", "/api/jira/ask-stream", "/api/run-cli-stream", "/api/ping-stream")

@app.middleware("http")
async def _disable_gzip_for_sse(request, call_next):
    try:
        if any(request.url.path.startswith(p) for p in _SSE_PATH_PREFIXES):
            # 하위 GZipMiddleware 는 Accept-Encoding 에 'gzip' 이 있으면 압축. 헤더를 비우면 압축 skip.
            _hdrs = dict(request.scope.get("headers") or [])
            new_hdrs = [(k, v) for (k, v) in (request.scope.get("headers") or []) if k != b"accept-encoding"]
            request.scope["headers"] = new_hdrs
    except Exception:
        pass
    return await call_next(request)

app.add_middleware(GZipMiddleware, minimum_size=500, compresslevel=5)

# 정적 파일 커스텀 마운트 — 캐시 버스터(?v=xxx)를 이미 쓰므로 강력한 브라우저 캐시 허용.
# 매 새로고침마다 조건부 GET (200-500ms 왕복 대기) 을 안 함 → 페이지 재접속이 빨라짐.
class _CachedStatic(StaticFiles):
    async def get_response(self, path, scope):
        resp = await super().get_response(path, scope)
        try:
            qs = (scope.get("query_string") or b"").decode("ascii", errors="ignore")
            # ?v=... 캐시버스터 있으면 1년 immutable (내용 바뀌면 v 값이 바뀌므로 안전)
            if "v=" in qs:
                resp.headers["Cache-Control"] = "public, max-age=31536000, immutable"
            else:
                # 버스터 없으면 짧게 (10분) 만 — 조건부 GET 은 절약하되 갱신 확인 여지 남김
                resp.headers["Cache-Control"] = "public, max-age=600"
        except Exception:
            pass
        return resp

app.mount("/static", _CachedStatic(directory=str(FRONTEND_DIR / "static")), name="static")


# ══════════════════════════════════════════════════════════════════════
# API 응답 캐시 헤더
#
# 전에는 이 목록 전부에 max-age=10 + stale-while-revalidate=60 을 걸고
# 「실시간 정확도는 WebSocket 이 보장한다」 고 적어 두었다. 그게 틀렸다 —
# **WebSocket 이 시켜서 다시 읽는 것도 같은 캐시를 지나간다.** 남이 저장해서
# 다시 읽어도 브라우저가 묵은 응답을 그대로 내주니, 화면이 늘 한 판씩
# 늦었다. 남이 고친 것을 보려면 새로고침을 눌러야 했다.
#
# 그래서 둘로 가른다.
#
#  · 같이 고치는 자료 — 캐시하지 않는다. 여럿이 한 시험을 놓고 일하는
#    도구에서 10초 묵은 값은 10초짜리 오답이다
#  · 잘 안 바뀌는 설정 — 짧게 캐시한다. 로고·도움말 같은 것
#
# `public` 도 `private` 로 바꾼다. 로그인한 사람에 따라 달라지는 응답을
# 공용 캐시에 담게 두면 안 된다.
# ══════════════════════════════════════════════════════════════════════
_LIVE_PATHS = (
    "/api/tc",              # ?meta=1 목록 및 단건
    "/api/cycle",           # ?meta=1 목록 및 단건
    "/api/req",             # 목록 및 단건
    "/api/manuals",
    "/api/board",
    "/api/racks",
    "/api/devices",
    "/api/procedures",
    "/api/folders",
    "/api/cycle-folders",
    "/api/manual-folders",
    "/api/custom-fields",
    "/api/permissions",
    "/api/global-params",
    # 사람이 그 자리에서 고치는 설정 — 30초 캐시가 「저장했는데 옛 이름이
    # 다시 보인다」 를 만들었다(지적). 고치자마자 다시 읽는 자료다.
    "/api/llms",
    "/api/prompts",
    # 브랜딩 — 고치고 새로고침하면 옛 값이 돌아왔다(지적: 크기 변경이 안 된다,
    # 사진 제거가 안 된다). 30초 캐시가 방금 저장한 것을 덮고 있었다.
    # 로고·이름은 자주 읽히지만 그 몇 KB 를 아끼자고 「저장이 안 되는 화면」
    # 을 만들 수는 없다.
    "/api/branding",
    # 조직도 — 고치자마자 다시 읽는 자료다. Cache-Control 을 아예 안 붙여
    # 두었더니 브라우저가 제 나름대로 캐시해, 지운 마디가 도로 보이고 방금
    # 넣은 사람이 안 보였다(지적: 대표이사 자리가 안 채워진다). 위 형제들이
    # 겪은 그 덫이다.
    "/api/org",
)

_CACHEABLE_PATHS = (
    "/api/device-catalog",
    "/api/help",
    "/api/page-ai",
    "/api/dify/assistants",
    "/api/ui-options",
    "/api/org-options",
    "/api/jira/config",
)

_CACHE_EXCLUDE_SUFFIX = ("/run-history", "/snapshots", "/ui-options")
@app.middleware("http")
async def _api_cache_headers(request, call_next):
    resp = await call_next(request)
    try:
        method = request.method
        path = request.url.path
        # GET 만, 그리고 캐시 대상 경로 (하위 경로 포함 매치)
        if method == "GET":
            # 실시간 반영이 필요한 하위 리소스는 캐시 대상에서 제외 (run-history/snapshots 등)
            if any(seg in path for seg in _CACHE_EXCLUDE_SUFFIX):
                if "cache-control" not in {k.lower() for k in resp.headers.keys()}:
                    resp.headers["Cache-Control"] = "no-store"
                return resp
            _has = "cache-control" in {k.lower() for k in resp.headers.keys()}
            # 같이 고치는 자료 — 늘 서버에 물어본다
            for p in _LIVE_PATHS:
                if path == p or path.startswith(p + "/") or path.startswith(p + "?"):
                    if not _has:
                        resp.headers["Cache-Control"] = "no-store"
                    return resp
            for p in _CACHEABLE_PATHS:
                if path == p or path.startswith(p + "/") or path.startswith(p + "?"):
                    # 이미 다른 미들웨어·엔드포인트가 Cache-Control 지정했으면 존중
                    if not _has:
                        resp.headers["Cache-Control"] = "private, max-age=30"
                    break
    except Exception:
        pass
    return resp

# 데이터 파일 초기화
# 데이터 루트가 빈 볼륨/새 설치일 수 있으므로 부모 폴더를 먼저 만든다.
# (이 줄이 없으면 도커 첫 기동 때 FileNotFoundError 로 import 자체가 실패한다)
for f, default in [
    (LLMS_FILE, {"llms": []}),
]:
    if not f.exists():
        f.parent.mkdir(parents=True, exist_ok=True)
        with open(f, "w", encoding="utf-8") as fp:
            import json as _json
            _json.dump(default, fp, ensure_ascii=False, indent=2)

# 등록된 것이 모델을 안 들고 있을 때 쓰는 이름 — 한 곳에 둔다
CLAUDE_FALLBACK_MODEL = "claude-sonnet-4-5-20250929"

# Anthropic 클라이언트 (API 키 없으면 None)
_api_key = os.environ.get("ANTHROPIC_API_KEY", "")
claude_client = anthropic.Anthropic(api_key=_api_key) if _api_key else None

# WebSocket 연결 관리
active_connections: list[WebSocket] = []

# ───────────────────────────────────────────
# 유틸
# ───────────────────────────────────────────
def load_json(path: Path) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def save_json(path: Path, data: dict):
    # 상위 폴더가 없으면 만든다. 새로 클론한 곳에는 data/state 가 없어서
    # import 단계의 초기화(FOLDERS_FILE 등)가 FileNotFoundError 로 죽었다.
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def _run_async(coro):
    """더미 — 예전 shim 코드가 참조하는 곳이 있을 수 있어 남겨둠. 새 이벤트 루프에서 실행."""
    try:
        return asyncio.run(coro)
    except RuntimeError:
        # 이미 loop 있는 컨텍스트 → 스레드에서 별도 loop
        import concurrent.futures as _cf
        with _cf.ThreadPoolExecutor(max_workers=1) as _ex:
            return _ex.submit(asyncio.run, coro).result(timeout=15)

# 소식 → 수정 이력. 한 곳(broadcast)에서 받아 적으면 저장 지점 여덟
# 군데를 따로 고칠 일이 없고, 새 소식이 생겨도 여기 한 줄이다.
_AUDIT_MAP = {
    "tc_updated": ("tc", "tcid", "updated"),
    "tc_deleted": ("tc", "tcid", "deleted"),
    "req_updated": ("req", "req_id", "updated"),
    "req_deleted": ("req", "req_id", "deleted"),
    "cycle_updated": ("cycle", "cycle_id", "updated"),
    "defect_updated": ("defect", "id", "updated"),
    "tc_run_history_new": ("tc", "tcid", "run"),
}


async def broadcast(message: dict):
    # 수정 이력 — 접속자가 없어도 남긴다 (알림 종·감사가 나중에 읽는다)
    try:
        m = _AUDIT_MAP.get(str(message.get("type") or ""))
        if m:
            kind, key, action = m
            extra = ""
            if message.get("type") == "tc_run_history_new":
                extra = f" PASS {message.get('pass', 0)} FAIL {message.get('fail', 0)}"
            await db.audit_add(kind, str(message.get(key) or ""), action + extra,
                               str(message.get("user") or ""))
    except Exception:
        pass  # 이력이 소식을 막으면 안 된다

    # 모든 접속자에게 병렬 전송 — 순차 await 로 하면 접속자 수만큼 지연 누적 (10명이면 delete API 응답이 왕복 10회만큼 늦어짐)
    if not active_connections:
        return
    async def _one(ws):
        try: await ws.send_json(message)
        except Exception: pass
    await asyncio.gather(*(_one(ws) for ws in list(active_connections)), return_exceptions=True)

engine.broadcast = broadcast

# ── 동시 접속 presence + 편집 제어권 ──
ws_state = {}          # id(ws) -> {"ws":ws, "user":str|None, "page":str|None}
page_controller = {}   # page -> user (제어권 보유자)

def _presence_users(page):
    seen = []
    for s in ws_state.values():
        u = s.get("user")
        if s.get("page") == page and u and u not in seen:
            seen.append(u)
    return seen

async def _broadcast_presence(page):
    if not page:
        return
    users = _presence_users(page)
    ctrl = page_controller.get(page)
    if ctrl not in users:                       # 제어자가 떠났으면 첫 접속자에게 자동 양도
        ctrl = users[0] if users else None
        if ctrl:
            page_controller[page] = ctrl
        elif page in page_controller:
            del page_controller[page]
    await broadcast({"type": "presence", "page": page, "users": users, "controller": ctrl})

_tc_running: dict = {}   # tcid -> {"user": str, "at": ts} — 지금 자동 실행 중인 시험


async def _broadcast_tc_running(tcid: str, user: str, on: bool):
    """이 시험을 **누가 지금 돌리고 있나** — 보고 있는 모두에게. 실행한 사람만
    「진행중」이 보이던 것을 남들도 보게 한다(지시)."""
    await broadcast({"type": "tc_running", "tcid": tcid, "user": user, "on": bool(on)})


async def _broadcast_focus(page):
    """이 화면에서 누가 어느 항목을 보고 있나 — {항목번호: [사람…]}"""
    if not page:
        return
    at = {}
    for st in list(ws_state.values()):
        if st.get("page") != page:
            continue
        f = st.get("focus")
        u = st.get("user")
        if f is None or not u:
            continue
        at.setdefault(str(f), [])
        if u not in at[str(f)]:
            at[str(f)].append(u)
    await broadcast({"type": "focus", "page": page, "at": at})


@app.post("/api/broadcast-reload")
async def broadcast_reload(payload: dict = None):
    # 접속 중인 모든 클라이언트에 강제 새로고침 신호 브로드캐스트.
    # payload: {"delay_sec": int (기본 3), "message": str (기본 '관리자가 새로고침을 요청했습니다')}
    _p = payload or {}
    try: _delay = int(_p.get("delay_sec", 3) or 3)
    except Exception: _delay = 3
    _msg = str(_p.get("message", "") or "관리자가 새로고침을 요청했습니다")
    await broadcast({"type": "force_reload", "delay_sec": _delay, "message": _msg})
    return {"ok": True, "targets": len(active_connections), "delay_sec": _delay}

app.include_router(engine.router)

# ───────────────────────────────────────────
# 연결 상태 체크
# ───────────────────────────────────────────
def check_tcp(ip: str, port: int, timeout: float = 2.0) -> bool:
    try:
        s = socket.create_connection((ip, port), timeout=timeout)
        s.close()
        return True
    except Exception:
        return False

def check_ssh(ip: str, port: int, username: str, password: str, timeout: float = 5.0) -> bool:
    try:
        client = paramiko.SSHClient()
        client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        client.connect(ip, port=port, username=username, password=password, timeout=timeout)
        client.close()
        return True
    except Exception:
        return False

def check_telnet(ip: str, port: int, timeout: float = 3.0) -> bool:
    return check_tcp(ip, port, timeout)

# ───────────────────────────────────────────
# SSH / Telnet 명령 실행
# ───────────────────────────────────────────
def ssh_exec(ip: str, port: int, username: str, password: str, command: str) -> str:
    client = paramiko.SSHClient()
    client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
    client.connect(ip, port=port, username=username, password=password, timeout=10)
    stdin, stdout, stderr = client.exec_command(command)
    output = stdout.read().decode("utf-8", errors="replace")
    err = stderr.read().decode("utf-8", errors="replace")
    client.close()
    return output + err

def tcl_exec(script_path: str) -> str:
    try:
        result = subprocess.run(
            ["tclsh", script_path],
            capture_output=True, text=True, timeout=60
        )
        return result.stdout + result.stderr
    except FileNotFoundError:
        return "[오류] tclsh 가 설치되어 있지 않거나 PATH에 없습니다."
    except subprocess.TimeoutExpired:
        return "[오류] TCL 스크립트 실행 시간 초과 (60초)"

# ───────────────────────────────────────────
# 라우터 - 페이지
# ───────────────────────────────────────────
@app.get("/")
async def index():
    return FileResponse(
        str(FRONTEND_DIR / "index.html"),
        headers={"Cache-Control": "no-cache, no-store, must-revalidate", "Pragma": "no-cache"}
    )

# ───────────────────────────────────────────
# 라우터 - LLM 관리
# ───────────────────────────────────────────
def init_llms_file():
    if not LLMS_FILE.exists():
        save_json(LLMS_FILE, {"llms": []})

@app.get("/api/llms")
async def get_llms():
    init_llms_file()
    return load_json(LLMS_FILE)

# ── 자원 관리: 인원 투입(M/M) · 프로젝트 ──
@app.get("/api/resource/manpower")
async def get_manpower():
    d = _kv_load_sync("manpower", {})
    return d if isinstance(d, dict) else {}

def _rsc_backup_kv(key: str):
    """저장 전 백업 (DB 저장이라 파일 백업 대신 data/backups/ 폴더에 JSON 스냅샷 저장)."""
    try:
        _cur = _kv_load_sync(key, None)
        if not isinstance(_cur, dict): return
        nrows = sum(len((p or {}).get("rows", []) or []) for p in (_cur.get("pages", {}) or {}).values())
        nrows += len(_cur.get("rows", []) or [])
        if nrows == 0: return
        bdir = DATA_DIR / "backups"; bdir.mkdir(exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d-%H%M%S")
        save_json(bdir / f"{key}-{ts}.json", _cur)
        # 최근 30개 회전
        baks = sorted(bdir.glob(f"{key}-*.json"))
        for old in baks[:-30]:
            try: old.unlink()
            except Exception: pass
    except Exception:
        pass

def _rsc_backup(path: Path):
    """레거시 파일 기반 백업 — DB 이전 후엔 _rsc_backup_kv 를 씀. 호환용 shim."""
    _rsc_backup_kv(path.stem)

@app.post("/api/resource/manpower")
async def save_manpower(payload: dict):
    _rsc_backup_kv("manpower")             # 저장 전 자동 백업
    _kv_save_sync("manpower", payload)
    return {"ok": True}

@app.get("/api/resource/projects")
async def get_projects():
    if not RSC_PROJECTS_FILE.exists():
        save_json(RSC_PROJECTS_FILE, {})
    return load_json(RSC_PROJECTS_FILE)

@app.post("/api/resource/projects")
async def save_projects(payload: dict):
    save_json(RSC_PROJECTS_FILE, payload)
    return {"ok": True}

class LLMCreate(BaseModel):
    name: str
    type: str
    endpoint: str
    model: str
    apikey: str = ""
    max_tokens: int = 4096
    context_size: int = 262144
    temperature: float = 0.7
    # 파라미터(선택) — 프론트 편집 화면에서 저장
    completion_mode: str = "chat"
    top_p: Optional[float] = None
    top_k: Optional[int] = None
    presence_penalty: Optional[float] = None
    frequency_penalty: Optional[float] = None
    # 고급 옵션 — 저장 시 유실되던 필드들
    compat_mode: str = "openai"
    thinking_mode: str = "none"
    function_call_type: str = "not_support"
    stream_function_call: str = "not_support"
    vision_support: str = "not_support"
    structured_output: str = "not_support"
    stream_mode_auth: str = "not_use"
    stream_delimiter: str = "\n\n"
    system_prompt: str = ""
    greeting: str = ""
    placeholder: str = ""
    uses: list = []
    status: str = "active"
    field_prompts: dict = {}
    kb_group: str = ""   # 지식 검색 노출 그룹: '' | general | kb | jira | external

@app.post("/api/llms")
async def add_llm(llm: LLMCreate):
    init_llms_file()
    data = load_json(LLMS_FILE)
    new_id = f"llm{len(data['llms'])+1:03d}_{int(datetime.now().timestamp())}"
    new_llm = {"id": new_id, **llm.model_dump()}
    data["llms"].append(new_llm)
    save_json(LLMS_FILE, data)
    return {"success": True, "llm": new_llm}

# ★ /api/llms/{llm_id} 보다 반드시 앞에 위치해야 함
@app.post("/api/llms/import")
async def import_llms(payload: dict):
    """localStorage LLM 목록 일괄 import"""
    init_llms_file()
    data = load_json(LLMS_FILE)
    llms = payload.get("llms", [])
    added = 0
    for llm in llms:
        exists = any(
            l.get("name") == llm.get("name") and l.get("endpoint") == llm.get("endpoint")
            for l in data["llms"]
        )
        if not exists:
            new_id = llm.get("id") or f"llm{len(data['llms'])+1:03d}_{int(datetime.now().timestamp())}"
            data["llms"].append({"id": new_id, **{k:v for k,v in llm.items() if k!="id"}})
            added += 1
    save_json(LLMS_FILE, data)
    return {"success": True, "added": added, "total": len(data["llms"])}

@app.post("/api/llms/reorder")
async def reorder_llms(payload: dict):
    """LLM 목록 순서 재배치 (드래그)"""
    init_llms_file()
    data = load_json(LLMS_FILE)
    ids = payload.get("ids", [])
    order = {lid: i for i, lid in enumerate(ids)}
    data["llms"].sort(key=lambda l: order.get(l.get("id"), 9999))
    save_json(LLMS_FILE, data)
    return {"success": True, "count": len(data["llms"])}

@app.put("/api/llms/{llm_id}")
async def update_llm(llm_id: str, llm: LLMCreate):
    init_llms_file()
    data = load_json(LLMS_FILE)
    for i, l in enumerate(data["llms"]):
        if l["id"] == llm_id:
            data["llms"][i] = {"id": llm_id, **llm.model_dump()}
            save_json(LLMS_FILE, data)
            return {"success": True}
    raise HTTPException(404, "LLM을 찾을 수 없습니다")

@app.delete("/api/llms/{llm_id}")
async def delete_llm(llm_id: str):
    init_llms_file()
    data = load_json(LLMS_FILE)
    data["llms"] = [l for l in data["llms"] if l["id"] != llm_id]
    save_json(LLMS_FILE, data)
    return {"success": True}


@app.post("/api/anthropic/models")
async def anthropic_models(body: dict):
    """Claude 에서 쓸 수 있는 모델 목록 — **모델명을 손으로 치지 않게**(지시).

    저장 전에도 물어볼 수 있어야 해서 키를 몸통으로 받는다(주소에 실으면
    서버 로그에 키가 남는다). 키가 없거나 못 닿으면 **아는 모델**을 돌려준다 —
    목록이 비면 고를 수가 없다.
    """
    import httpx
    fallback = ["claude-opus-5", "claude-sonnet-5", "claude-haiku-4-5-20251001"]
    key = str((body or {}).get("apikey") or "").strip()
    base = str((body or {}).get("endpoint") or "https://api.anthropic.com").rstrip("/")
    if base.endswith("/v1"):
        base = base[:-3]
    if not key:
        return {"ok": True, "models": fallback, "source": "기본 목록 (키를 넣으면 실제 목록을 읽습니다)"}
    try:
        async with httpx.AsyncClient(timeout=15) as c:
            r = await c.get(
                base + "/v1/models",
                headers={"x-api-key": key, "anthropic-version": "2023-06-01"},
            )
        if r.status_code == 200:
            ids = [str(m.get("id") or "") for m in (r.json().get("data") or []) if m.get("id")]
            return {"ok": True, "models": ids or fallback, "source": "Anthropic"}
        if r.status_code in (401, 403):
            return {"ok": False, "models": fallback, "error": f"키를 받지 않았습니다 ({r.status_code})"}
        return {"ok": False, "models": fallback, "error": f"{r.status_code} · {r.text[:120]}"}
    except Exception as e:
        return {"ok": False, "models": fallback, "error": f"닿지 못했습니다 — {str(e)[:120]}"}


@app.post("/api/openai/models")
async def openai_models(body: dict):
    """OpenAI 호환 서버에서 고를 모델 목록 — Claude 와 같은 손놀림으로(지시).

    OpenAI 규격은 `GET {주소}/models` 한 번이면 된다. 랩의 vLLM 도 같은
    말을 하므로 사내 서버에도 그대로 통한다 — 주소만 다르다.

    키를 몸통으로 받는 까닭은 Anthropic 쪽과 같다: 주소에 실으면 서버
    로그에 키가 남고, 저장 전에도 물어볼 수 있어야 한다.
    """
    import httpx
    fallback = ["gpt-4o", "gpt-4o-mini", "gpt-4.1", "gpt-4.1-mini", "o3-mini"]
    key = str((body or {}).get("apikey") or "").strip()
    base = str((body or {}).get("endpoint") or "https://api.openai.com/v1").rstrip("/")
    if not base.endswith("/v1") and "api.openai.com" in base:
        base += "/v1"
    if not key and "api.openai.com" in base:
        # 공식 주소는 키 없이 목록을 안 준다 — 아는 것부터 보여 준다
        return {"ok": True, "models": fallback, "source": "기본 목록 (키를 넣으면 실제 목록을 읽습니다)"}
    try:
        headers = {"Authorization": "Bearer " + key} if key else {}
        async with httpx.AsyncClient(timeout=15) as c:
            r = await c.get(base + "/models", headers=headers)
        if r.status_code == 200:
            ids = [str(m.get("id") or "") for m in (r.json().get("data") or []) if m.get("id")]
            # 임베딩·음성·이미지까지 다 나온다 — 글 짓는 것만 남긴다
            chat = [m for m in ids if not any(
                w in m for w in ("embedding", "whisper", "tts", "dall-e", "moderation", "audio", "realtime")
            )]
            return {"ok": True, "models": sorted(chat) or ids or fallback,
                    "source": "OpenAI" if "api.openai.com" in base else base}
        if r.status_code in (401, 403):
            return {"ok": False, "models": fallback, "error": f"키를 받지 않았습니다 ({r.status_code})"}
        return {"ok": False, "models": fallback, "error": f"{r.status_code} · {r.text[:120]}"}
    except Exception as e:
        return {"ok": False, "models": fallback, "error": f"닿지 못했습니다 — {str(e)[:120]}"}


@app.post("/api/llms/{llm_id}/test")
async def test_llm(llm_id: str):
    """저장된 설정으로 실제로 한 번 불러 본다.

    주소·모델명·키가 맞는지는 눌러 보기 전에는 알 수 없다. 사내 서버는
    자주 내려가 있어서, 답이 안 나올 때 '설정이 틀렸나 서버가 죽었나' 를
    가려주지 않으면 매번 처음부터 의심하게 된다.

    /v1/models 로 먼저 확인한다 — 토큰을 쓰지 않고 인증과 주소만 본다.
    그것이 없는 서버(일부 vLLM 구성)를 위해 짧은 chat 요청으로 한 번 더 시도한다.
    """
    import httpx

    init_llms_file()
    data = load_json(LLMS_FILE)
    llm = next((l for l in data["llms"] if l.get("id") == llm_id), None)
    if llm is None:
        raise HTTPException(404, "LLM 을 찾을 수 없습니다")

    base = (llm.get("endpoint") or "").rstrip("/")
    ltype = str(llm.get("type") or "").lower()
    model = (llm.get("model") or "").strip()

    # ── Claude(Anthropic) 는 말이 다르다 ──────────────────────────
    #
    # OpenAI 계열은 `/v1/models` 에 Bearer 를 얹지만, Anthropic 은 `x-api-key`
    # 와 `anthropic-version` 을 쓴다. 그래서 여기까지 오면 늘 실패했고, 설정이
    # 맞는데도 「연결 테스트」 가 안 된다는 말을 들었다(지적).
    if ltype in ("claude", "anthropic") or "anthropic.com" in base:
        key = str(llm.get("apikey") or "").strip()
        if not key:
            return {"ok": False, "detail": "API Key 가 비어 있습니다 (sk-ant-… 키를 넣으세요)"}
        url = (base or "https://api.anthropic.com").rstrip("/")
        if url.endswith("/v1"):
            url = url[:-3]
        try:
            async with httpx.AsyncClient(timeout=15) as client:
                r = await client.get(
                    url + "/v1/models",
                    headers={"x-api-key": key, "anthropic-version": "2023-06-01"},
                )
            if r.status_code == 200:
                names = [str(m.get("id") or "") for m in (r.json().get("data") or [])]
                if model and names and model not in names:
                    return {
                        "ok": False,
                        "detail": f"키는 맞는데 '{model}' 모델이 없습니다",
                        "models": names[:20],
                    }
                return {"ok": True, "detail": f"연결됨 (모델 {len(names)}개)", "models": names[:20]}
            if r.status_code in (401, 403):
                return {"ok": False, "detail": f"키를 받지 않았습니다 ({r.status_code})"}
            return {"ok": False, "detail": f"{r.status_code} · {r.text[:160]}"}
        except Exception as e:
            return {"ok": False, "detail": f"닿지 못했습니다 — {str(e)[:160]}"}

    if not base:
        return {"ok": False, "detail": "엔드포인트가 비어 있습니다"}
    # 사람은 보통 .../v1 까지 적어 둔다. 안 적었으면 붙여 준다.
    root = base[:-3].rstrip("/") if base.endswith("/v1") else base
    headers = {}
    if llm.get("apikey"):
        headers["Authorization"] = f"Bearer {llm['apikey']}"

    async with httpx.AsyncClient(timeout=12) as client:
        try:
            r = await client.get(f"{root}/v1/models", headers=headers)
            if r.status_code == 200:
                names = [m.get("id") for m in (r.json().get("data") or []) if m.get("id")]
                if model and names and model not in names:
                    # 주소는 맞는데 모델 이름이 틀린 경우. 가장 흔한 실수다.
                    return {
                        "ok": False,
                        "detail": f"서버는 응답하는데 '{model}' 모델이 없습니다",
                        "models": names[:20],
                    }
                return {"ok": True, "detail": f"연결됨 (모델 {len(names)}개)", "models": names[:20]}
            if r.status_code in (401, 403):
                return {"ok": False, "detail": f"인증 실패 ({r.status_code}) — API Key 를 확인하세요"}
        except Exception as e:
            return {"ok": False, "detail": f"연결하지 못했습니다 — {e}"}

        # /v1/models 가 없는 서버. 가장 짧은 요청으로 한 번 더 본다.
        try:
            r = await client.post(
                f"{root}/v1/chat/completions",
                headers={**headers, "Content-Type": "application/json"},
                json={"model": model, "messages": [{"role": "user", "content": "ping"}],
                      "max_tokens": 1},
            )
            if r.status_code == 200:
                return {"ok": True, "detail": "연결됨 (chat 응답 확인)"}
            return {"ok": False, "detail": f"응답 코드 {r.status_code} — {r.text[:200]}"}
        except Exception as e:
            return {"ok": False, "detail": f"연결하지 못했습니다 — {e}"}

# ─────────────────── 사용자 관리 / 인증 ───────────────────
import hashlib as _hashlib
import secrets as _secrets

USERS_FILE = DATA_DIR / "state" / "users.json"
ROLES = ["관리자", "담당", "팀장", "팀원"]
SESSIONS = {}  # token -> {username, role, name, ts}  ← in-memory 캐시 (DB 는 sessions 테이블)
SESSIONS_FILE = DATA_DIR / "state" / "sessions.json"   # 레거시 — 시작 시 DB 로 1회 마이그레이션 후 미사용
_SESSION_TTL = 60 * 60 * 24 * 30  # 로그인 세션 유지 기간(30일)

def _save_one_session(sid: str):
    """세션 하나를 DB 로 저장. asyncpg 는 async 라 sync 컨텍스트에서 호출 시 fire-and-forget."""
    _s = SESSIONS.get(sid)
    if not _s: return
    try:
        _loop = _MAIN_LOOP or asyncio.get_event_loop()
        if _loop:
            asyncio.run_coroutine_threadsafe(
                db.session_upsert(sid, _s, None, _s.get("username", "")),
                _loop
            )
    except Exception:
        pass

def _delete_one_session(sid: str):
    """세션 하나를 DB 에서 삭제."""
    try:
        _loop = _MAIN_LOOP or asyncio.get_event_loop()
        if _loop:
            asyncio.run_coroutine_threadsafe(db.session_delete(sid), _loop)
    except Exception:
        pass

async def _load_sessions_from_db():
    """서버 시작 시 DB 에서 세션 로드 → in-memory SESSIONS 채움. TTL 초과분은 스킵 + DB 에서도 정리."""
    global SESSIONS
    try:
        _all = await db.sessions_all()
        _now = datetime.now().timestamp()
        _fresh = {}
        _expired = []
        for sid, s in (_all or {}).items():
            if not isinstance(s, dict): continue
            if (_now - float(s.get("ts", 0) or 0)) < _SESSION_TTL:
                _fresh[sid] = s
            else:
                _expired.append(sid)
        SESSIONS = _fresh
        # 만료 세션 정리 (백그라운드)
        for sid in _expired:
            try: await db.session_delete(sid)
            except Exception: pass
    except Exception as e:
        print(f"[_load_sessions_from_db] failed: {e}", flush=True)

async def _migrate_sessions_file_to_db():
    """레거시 sessions.json 이 있으면 DB 로 1회 이전 후 파일은 legacy 폴더로 이동."""
    try:
        if not SESSIONS_FILE.exists(): return
        _d = load_json(SESSIONS_FILE)
        if not isinstance(_d, dict) or not _d: return
        _now = datetime.now().timestamp()
        _cnt = 0
        for sid, s in _d.items():
            if not isinstance(s, dict): continue
            if (_now - float(s.get("ts", 0) or 0)) >= _SESSION_TTL: continue
            # DB 에 없으면 삽입 (이미 있으면 file 값이 오래된 것일 가능성 → skip)
            _existing = await db.session_get(sid)
            if _existing is None:
                try:
                    await db.session_upsert(sid, s, None, s.get("username", ""))
                    _cnt += 1
                except Exception: pass
        # 파일은 legacy 로 이동 (혹시 몰라 보존)
        _legacy = DATA_DIR / "legacy"
        _legacy.mkdir(exist_ok=True)
        _dst = _legacy / ("sessions-migrated-" + datetime.now().strftime("%Y%m%d-%H%M%S") + ".json")
        try: SESSIONS_FILE.rename(_dst)
        except Exception: pass
        if _cnt > 0:
            print(f"[migrate] sessions.json → DB: {_cnt}건 이전, 파일은 {_dst.name} 로 백업", flush=True)
    except Exception as e:
        print(f"[_migrate_sessions_file_to_db] failed: {e}", flush=True)

def _hash_pw(password: str, salt: str) -> str:
    return _hashlib.pbkdf2_hmac("sha256", str(password).encode("utf-8"), str(salt).encode("utf-8"), 100000).hex()

# ── users 데이터 저장소: DB(app_kv 'users') 로 이전. 파일(users.json) 은 레거시 백업용만.
#    기존 코드 21군데의 _users_load_sync()/save_json(USERS_FILE, ...) 를 wrapper 로 흡수.
_USERS_CACHE = {"users": []}   # in-memory 캐시 (읽기 부하 최소화)
_USERS_LOADED = False

def _users_load_sync():
    """DB 에서 users 로드 → 캐시. 이미 캐시 있으면 캐시 반환 (매 요청 DB 왕복 방지).
    ★ 요청 처리 스레드가 곧 이벤트 루프 스레드라 run_coroutine_threadsafe(same_loop) 는 deadlock.
       그래서 async 컨텍스트에서 안전한 진입은 오직 캐시 반환 뿐이고, DB 로드는 startup 훅의
       _init_users_async 가 미리 캐시를 채워둔 뒤에만 유효하다. 캐시가 비었다면 파일 fallback 만."""
    global _USERS_CACHE, _USERS_LOADED
    if _USERS_LOADED:
        return _USERS_CACHE
    # 캐시 없음 = startup 이전 (드문 경로) → 파일에서만 시도, DB 접근 금지 (deadlock 회피)
    _data = None
    if USERS_FILE.exists():
        try: _data = load_json(USERS_FILE)
        except Exception: _data = None
    if not isinstance(_data, dict) or "users" not in _data:
        _data = {"users": []}
    _USERS_CACHE = _data
    _USERS_LOADED = True
    return _data

def _users_reload_sync():
    """캐시 무효화 후 다시 로드 (다른 프로세스가 DB 를 갱신했을 때)."""
    global _USERS_LOADED
    _USERS_LOADED = False
    return _users_load_sync()

def _users_save_sync(data: dict):
    """users 데이터 저장 → DB + 캐시 갱신."""
    global _USERS_CACHE
    if not isinstance(data, dict): return
    _USERS_CACHE = data
    try:
        _loop = _MAIN_LOOP
        if _loop and _loop.is_running():
            asyncio.run_coroutine_threadsafe(db.kv_set("users", data), _loop)
        else:
            asyncio.run(db.kv_set("users", data))
    except Exception as e:
        print(f"[_users_save_sync] DB write failed: {e}", flush=True)

async def _migrate_users_file_to_db():
    """레거시 users.json → DB 로 1회 이전.
    안전장치: 파일 users 수 > DB users 수 이면 파일이 정본 → 병합 (파일 계정 우선, DB 추가 계정 유지)."""
    try:
        if not USERS_FILE.exists(): return
        try: _fdata = load_json(USERS_FILE)   # 파일 직접 로드 (캐시/wrapper 우회)
        except Exception as _fe:
            print(f"[migrate] users.json 파일 읽기 실패: {_fe}", flush=True)
            return
        if not isinstance(_fdata, dict) or not isinstance(_fdata.get("users"), list) or not _fdata["users"]:
            return
        _file_users = _fdata["users"]
        _existing = await db.kv_get("users", None)
        if isinstance(_existing, dict) and isinstance(_existing.get("users"), list) and _existing["users"]:
            _db_users = _existing["users"]
            # 파일이 DB 보다 많은 계정을 가지면 → 파일을 정본으로. DB 에만 있는 계정은 뒤에 추가.
            _file_names = {u.get("username") for u in _file_users if u.get("username")}
            _extra = [u for u in _db_users if u.get("username") and u.get("username") not in _file_names]
            if len(_file_users) > len(_db_users) or _extra:
                _merged = list(_file_users) + _extra
                await db.kv_set("users", {"users": _merged})
                print(f"[migrate] users.json({len(_file_users)}) + DB({len(_db_users)}) → 병합 {len(_merged)}명 (파일 계정 우선)", flush=True)
            else:
                print(f"[migrate] users: DB({len(_db_users)}) 유지 (파일={len(_file_users)}, 이미 최신)", flush=True)
        else:
            # DB 비어있음 → 파일 통째로 이전
            await db.kv_set("users", _fdata)
            print(f"[migrate] users.json → DB ({len(_file_users)}명)", flush=True)
        # 성공적으로 처리 완료 → 파일은 legacy 로 이동 (백업 보존)
        _legacy = DATA_DIR / "legacy"; _legacy.mkdir(exist_ok=True)
        _dst = _legacy / ("users-migrated-" + datetime.now().strftime("%Y%m%d-%H%M%S") + ".json")
        try: USERS_FILE.rename(_dst)
        except Exception: pass
    except Exception as e:
        print(f"[_migrate_users_file_to_db] failed: {e}", flush=True)

async def _init_users_async():
    """startup 훅용: DB 에 users 없으면 기본 admin 1개 생성 후 캐시 로드. async 컨텍스트라 wrapper 안 씀."""
    global _USERS_CACHE, _USERS_LOADED
    _data = await db.kv_get("users", None)
    if isinstance(_data, dict) and isinstance(_data.get("users"), list) and _data["users"]:
        _USERS_CACHE = _data; _USERS_LOADED = True
        return
    # DB 완전히 비었을 때만 admin 생성
    salt = _secrets.token_hex(8)
    admin = {
        "id": "admin", "username": "admin", "name": "관리자", "role": "관리자",
        "salt": salt, "password": _hash_pw("admin", salt), "active": True,
        "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }
    _new = {"users": [admin]}
    await db.kv_set("users", _new)
    _USERS_CACHE = _new; _USERS_LOADED = True
    print("[init_users_file] DB 비어있어 기본 admin 계정 생성", flush=True)

def init_users_file():
    """요청 컨텍스트용 sync 진입점. 캐시 있으면 skip. 캐시 없으면 sync wrapper 로 로드.
    (startup 훅은 대신 _init_users_async 를 씀 — 여긴 절대 부르면 안 됨: loop 재진입 deadlock)"""
    if _USERS_CACHE.get("users"): return
    _users_load_sync()
    # ★ admin 자동 생성은 async 진입점(_init_users_async) 에서만 — sync 에서는 하지 않음
    #   (sync 시점에 DB 못 읽으면 빈 캐시로 오판 → 원본 데이터를 admin 하나로 덮어쓰는 사고 발생)


# ══════════════════════════════════════════════════════════════════════
# 범용 KV wrapper — 파일(JSON) 을 app_kv(DB) 로 이전할 때 최소 리팩터링으로 쓰는 진입점.
# users 와 같은 캐시·안전장치 패턴을 KV key 별로 재사용.
# ══════════════════════════════════════════════════════════════════════
_KV_CACHE = {}          # key -> data (dict/list)
_KV_LOADED = {}         # key -> bool
_KV_FALLBACK_FILE = {}  # key -> Path (startup 이전에 캐시 없으면 이 파일에서 읽음)

def _kv_register_fallback(key: str, file_path):
    """이 key 의 파일 fallback 경로 등록. startup 훅에서 이전 미완료 시 캐시가 파일 로드."""
    _KV_FALLBACK_FILE[key] = file_path

def _kv_load_sync(key: str, default=None):
    """DB(app_kv) 에서 key 값 로드 → 캐시. 캐시 있으면 캐시 반환. 없으면 fallback 파일 로드.
    ★ sync 컨텍스트에서 loop 재진입 deadlock 방지 위해 DB 는 startup 훅 async 로만 채운다."""
    if _KV_LOADED.get(key):
        return _KV_CACHE.get(key, default)
    _data = None
    _fp = _KV_FALLBACK_FILE.get(key)
    if _fp is not None:
        try:
            if _fp.exists(): _data = load_json(_fp)
        except Exception: _data = None
    if _data is None: _data = default
    _KV_CACHE[key] = _data
    _KV_LOADED[key] = True
    return _data

def _kv_save_sync(key: str, data):
    """key 값 저장 → DB(app_kv) fire-and-forget + 캐시 갱신."""
    if data is None: return
    _KV_CACHE[key] = data
    try:
        _loop = _MAIN_LOOP
        if _loop and _loop.is_running():
            asyncio.run_coroutine_threadsafe(db.kv_set(key, data), _loop)
        else:
            asyncio.run(db.kv_set(key, data))
    except Exception as e:
        print(f"[_kv_save_sync '{key}'] DB write failed: {e}", flush=True)

async def _kv_init_async(key: str, file_path, sizeguard: bool = True):
    """startup 훅용: 파일 → DB 이전 (안전장치: 파일이 DB 보다 크면 파일이 정본).
    이전 완료 시 파일은 legacy 폴더로 이동. 캐시는 최종 값으로 채움."""
    global _KV_CACHE, _KV_LOADED
    try:
        _file_data = None
        _file_size = 0
        if file_path.exists():
            try:
                _file_data = load_json(file_path)
                _file_size = file_path.stat().st_size
            except Exception as _fe:
                print(f"[_kv_init '{key}'] 파일 읽기 실패: {_fe}", flush=True)
        _db_data = await db.kv_get(key, None)
        _db_size = len(json.dumps(_db_data, ensure_ascii=False, default=str)) if _db_data is not None else 0
        # 정책:
        # - 파일 없음 → DB 캐시로만 (이전 이미 완료 상태)
        # - 파일 있고 DB 없음 → 파일을 DB 로
        # - 둘 다 있고 파일이 크게 더 큼(sizeguard=True) → 파일이 정본, DB 덮어쓰기
        # - 둘 다 있고 DB 가 같거나 큼 → DB 유지, 파일은 legacy 로만 이동
        _chosen = None
        _reason = ""
        if _file_data is None:
            _chosen = _db_data if _db_data is not None else None
            _reason = "DB 만 존재 (이전 완료 상태)"
        elif _db_data is None:
            _chosen = _file_data
            await db.kv_set(key, _file_data)
            _reason = f"파일→DB 이전 (파일 {_file_size} bytes)"
        else:
            if sizeguard and _file_size > _db_size * 1.2:
                # 파일이 DB 보다 20%+ 크면 파일이 정본
                _chosen = _file_data
                await db.kv_set(key, _file_data)
                _reason = f"파일이 크므로 정본으로 채택 (file={_file_size}, db={_db_size})"
            else:
                _chosen = _db_data
                _reason = f"DB 유지 (file={_file_size}, db={_db_size})"
        # 파일 → legacy 로 이동
        if file_path.exists():
            _legacy = DATA_DIR / "legacy"; _legacy.mkdir(exist_ok=True)
            _dst = _legacy / (file_path.stem + "-migrated-" + datetime.now().strftime("%Y%m%d-%H%M%S") + file_path.suffix)
            try: file_path.rename(_dst)
            except Exception: pass
        _KV_CACHE[key] = _chosen
        _KV_LOADED[key] = True
        print(f"[migrate '{key}'] {_reason}", flush=True)
    except Exception as e:
        print(f"[_kv_init '{key}'] failed: {e}", flush=True)

def _public_user(u: dict) -> dict:
    return {k: v for k, v in u.items() if k not in ("password", "salt")}

def _find_user(username: str):
    init_users_file()
    for u in _users_load_sync()["users"]:
        if u.get("username") == username:
            return u
    return None

# 지금 요청의 세션. 미들웨어가 채운다.
#
# 옛 엔드포인트들은 토큰을 쿼리(?token=)로만 받는다(21곳쯤). 새 화면은
# Authorization 헤더로 보내므로 그대로 두면 전부 401 이 난다.
# 라우트를 하나씩 고치는 대신, 토큰 인자가 비었을 때 이 값으로 넘어가게
# 한다 — 미들웨어가 이미 확인한 세션이라 안전하고, 옛 화면의 ?token= 도
# 그대로 동작한다.
_CUR_SESSION: contextvars.ContextVar = contextvars.ContextVar("utop_session", default=None)


def _user_from_token(token: str):
    s = SESSIONS.get(token or "") or _CUR_SESSION.get()
    return _find_user(s.get("username")) if s else None

def _require_admin(token: str):
    u = _user_from_token(token or _REQ_TOKEN.get(""))
    if not u:
        raise HTTPException(401, "로그인이 필요합니다")
    if u.get("role") != "관리자":
        raise HTTPException(403, "관리자 권한이 필요합니다")
    return u

# ══════════════════════════════════════════════════════════════════════
# 인증 강제
#
# 지금까지 라우트 234개 중 어느 것도 로그인을 확인하지 않았다. 사내망이라
# 넘어갔지만 50명이 함께 쓰면 누가 무엇을 고쳤는지 알 수 없고, 장비 락도
# 편집 중 표시도 '누구' 를 알아야 만들 수 있다.
#
# 라우트마다 의존성을 붙이면 234곳을 고쳐야 하고 새 라우트에서 빠뜨리기
# 쉽다. 미들웨어에서 한 번에 막고, 열어둘 곳만 목록으로 둔다 —
# 기본이 '막힘' 이어야 새로 만든 라우트가 자동으로 보호된다.
# ══════════════════════════════════════════════════════════════════════

@app.get("/api/health")
async def api_health():
    """도커 헬스체크가 부르는 곳. 로그인 없이 열려 있다.

    프로세스가 살아 있는지만 보면 의미가 없다 — uvicorn 은 떠 있는데 DB 가
    안 붙어 모든 화면이 500 인 상태가 '정상' 으로 보고된다. 그래서 DB 까지
    한 번 찔러 보고, 안 되면 503 으로 답한다.
    """
    try:
        async with db.pool().acquire() as c:
            await c.fetchval("SELECT 1")
    except Exception as e:
        from fastapi.responses import JSONResponse
        return JSONResponse({"ok": False, "db": str(e)[:200]}, status_code=503)
    return {"ok": True, "db": True}


# 로그인 없이 열어두는 경로. 접두사로 비교한다.
_AUTH_PUBLIC = (
    "/api/login",
    # 서버끼리 부르는 자리라 사람의 세션이 없다. 대신 N2X_RELAY_KEY 로
    # 자기가 확인하고, 열쇠가 안 정해져 있으면 아예 거절한다.
    "/api/n2x/send",
    # 실행기(runner 컨테이너)가 부르는 자리. 사람의 세션이 없어서 RUNNER_KEY
    # 로 자기가 확인한다. 열쇠가 안 정해져 있으면 아예 거절한다.
    "/api/runner/",
    "/api/logout",
    "/api/health",
    "/api/req-images/",     # 마크다운 안 <img> 는 헤더를 못 붙인다
    # 로그인 **화면**이 이것으로 그려진다 — 로고·회사 사진·이름. 로그인
    # 전이라 세션이 없다(지적: 로그인 시 사진이 안 나온다). 비밀이 아니라
    # 회사 간판이라 열어 둔다.
    "/api/branding",
    "/openapi.json",
    "/docs",
    "/redoc",
)


# 지금 요청의 토큰 — 미들웨어가 담아 두고 관리자 검사(_require_admin)가
# 꺼내 쓴다.
#
# 화면은 토큰을 **헤더**로 보내는데, 관리자 자리들은 함수 인자(token: str = "")
# 로 받아 **주소의 ?token=** 만 봤다. 그래서 저장이 조용히 401 로 떨어지고
# 화면은 옛 값을 다시 보여 줬다(지적: 15 를 16 으로 고쳐도 되돌아간다).
# 자리마다 손대면 또 빠뜨린다 — 한 곳에서 받는다.
_REQ_TOKEN: "contextvars.ContextVar[str]" = contextvars.ContextVar("_REQ_TOKEN", default="")


def _token_from(request) -> str:
    """Authorization: Bearer 우선, 없으면 기존 방식(쿼리 ?token=)."""
    h = request.headers.get("authorization") or ""
    if h.lower().startswith("bearer "):
        return h[7:].strip()
    return request.query_params.get("token") or ""


async def _session_of(token: str):
    """메모리 캐시 → 없으면 PG. 워커가 여러 개면 로그인한 워커에만
    메모리 세션이 있으므로 DB 를 반드시 확인해야 한다."""
    if not token:
        return None
    s = SESSIONS.get(token)
    if s:
        return s
    try:
        row = await db.session_get(token)
    except Exception:
        return None
    if row:
        SESSIONS[token] = row      # 이 워커에도 캐시
        return row
    return None


@app.middleware("http")
async def _require_login(request, call_next):
    path = request.url.path
    # 중계 전용 서버는 N2X 창구와 상태 확인만 연다. DB 가 없으니 다른
    # 창구는 어차피 터지고, 열어 두면 이 서버가 시험 서버인 줄 알고
    # 붙었다가 알 수 없는 오류만 본다.
    if N2X_RELAY_ONLY and path.startswith("/api/") and not (
        path.startswith("/api/n2x/") or path.startswith("/api/health")
    ):
        from fastapi.responses import JSONResponse
        return JSONResponse(
            {"detail": "이 서버는 N2X 중계 전용입니다 — 시험 서버로 접속하세요"},
            status_code=503,
        )
    if request.method == "OPTIONS" or not path.startswith("/api/"):
        return await call_next(request)      # 화면·정적 파일은 통과
    if any(path.startswith(p) for p in _AUTH_PUBLIC):
        _REQ_TOKEN.set(_token_from(request))
        return await call_next(request)

    _tok = _token_from(request)
    _REQ_TOKEN.set(_tok)
    s = await _session_of(_tok)
    if not s:
        from fastapi.responses import JSONResponse
        return JSONResponse({"detail": "로그인이 필요합니다"}, status_code=401)

    # 아래 코드가 '누가 했는지' 를 알 수 있게 실어 보낸다
    request.state.user = s
    _CUR_SESSION.set(s)
    return await call_next(request)


class LoginReq(BaseModel):
    username: str
    password: str


# ══════════════ Jira 계정으로 로그인 ══════════════
#
# 사원이 모두 Jira 계정을 갖고 있어 **Jira 를 정본**으로 삼는다(합의).
# 우리는 비밀번호를 저장하지 않는다 — Jira 에서 바꾸면 그대로 따라간다.
# 로컬 계정은 안전망이다: Jira 가 죽었거나 admin 같은 비상 계정용.
#
# Jira 에는 아무것도 안 쓴다 — 로그인 한 번에 읽기 호출(myself) 하나뿐이다.

#: 연속 실패 잠그기 — 우리가 먼저 막아 Jira 까지 실패가 안 쌓이게 한다.
#: (Jira Server 는 실패가 쌓이면 CAPTCHA 를 걸어 그 사람이 웹에서 풀어야 한다)
_LOGIN_FAILS: dict = {}
#: 마지막으로 Jira 가 거절한 것 — 관리자 화면(계정 관리)에서만 본다.
#: 비밀번호는 담지 않는다. 아이디와 까닭·시각뿐이다.
_JIRA_LAST_FAIL: dict = {}
_LOGIN_FAIL_MAX = 3
_LOGIN_LOCK_SEC = 60


def _jira_verify_flag(cfg: dict) -> bool:
    """TLS 인증서를 검증할까 — **안 정했으면 검증한다.**

    설정 파일에 `verify: null` 이 들어 있는 경우가 있다(옛 화면이 남긴 값).
    `cfg.get("verify", True)` 는 그때 None 을 돌려주어 라이브러리마다 다르게
    해석된다 — 켜고 끈 기억이 없는데 동작이 달라진다.
    """
    v = cfg.get("verify")
    return True if v is None else bool(v)


def _jira_login_base(cfg: dict = None) -> str:
    """**로그인을 물어볼 Jira 주소.**

    이슈를 등록·조회하는 Jira 와 사람을 확인하는 Jira 가 다를 수 있다(지시:
    사내에 둘이다). 안 적었으면 이슈 쪽 주소를 그대로 쓴다 — 대개는 같다.
    """
    cfg = cfg if cfg is not None else _jira_cfg()
    return (str(cfg.get("login_url") or "").strip() or str(cfg.get("url") or "").strip()).rstrip("/")


def _jira_login_on() -> bool:
    cfg = _jira_cfg()
    return bool(cfg.get("login_enabled")) and bool(_jira_login_base(cfg))


async def _jira_verify_login(username: str, password: str) -> tuple:
    """그 사람의 ID/PW 로 Jira 에 물어본다.

    돌려주는 것: (Jira 가 아는 사람 정보 | None, 안 된 까닭).
    까닭이 'captcha' 면 Jira 가 사람 확인을 걸어 둔 것이라 우리가 풀 수 없다 —
    그 사람이 Jira 웹에 한 번 들어가 풀어야 한다.
    """
    import base64 as _b64
    import httpx
    cfg = _jira_cfg()
    base = _jira_login_base(cfg)          # ★ 로그인은 로그인용 주소로
    if not base:
        return None, "no-url"
    basic = _b64.b64encode(f"{username}:{password}".encode("utf-8")).decode("ascii")
    try:
        async with httpx.AsyncClient(timeout=12, verify=_jira_verify_flag(cfg)) as c:
            r = await c.get(
                base + "/rest/api/2/myself",
                headers={"Accept": "application/json", "Authorization": "Basic " + basic},
            )
    except Exception as exc:
        # Jira 가 안 뜨거나 망이 막혔다 — 로컬 계정으로 넘어간다.
        # 인증서 문제는 따로 말한다. 「닿지 못했습니다」 로 뭉개면 망을 뒤지게
        # 되는데, 실제로는 체크박스 하나(TLS 검증)나 인증서 갱신이면 끝난다.
        msg = str(exc)
        print(f"[jira-login] 붙지 못했습니다: {msg[:200]}", flush=True)
        if "CERTIFICATE" in msg.upper() or "SSL" in msg.upper():
            return None, "cert"
        return None, "unreachable"
    if r.status_code == 200:
        try:
            return r.json(), ""
        except Exception:
            return None, "bad-response"
    reason = str(r.headers.get("X-Seraph-LoginReason") or "")
    if "CAPTCHA" in reason.upper():
        return None, "captcha"
    return None, "denied"


def _jira_auto_create() -> bool:
    """모르는 사람이 Jira 로 들어오면 그 자리에서 계정을 만들까.

    사원이 모두 Jira 계정을 갖고 있으니 **회원가입을 따로 두지 않는다**(지시).
    그래서 기본은 **켜짐**이다. 명단에 있는 사람만 받고 싶은 곳(운영 정책)은
    계정 관리에서 끈다 — 그러면 관리자가 먼저 등록해야 들어온다.
    """
    cfg = _jira_cfg()
    if cfg.get("login_auto_create") is not None:
        return bool(cfg.get("login_auto_create"))
    v = str(os.environ.get("JIRA_AUTO_CREATE", "")).strip().lower()
    if v:
        return v in ("1", "true", "yes", "on")
    return True


def _upsert_jira_user(username: str, ju: dict) -> dict:
    """Jira 로 들어온 사람을 UTOP 사용자로. **비밀번호는 담지 않는다.**

    이미 있으면 이름·메일만 Jira 쪽으로 맞춘다. 관리자가 꺼 둔 계정(active
    False)을 여기서 되살리지는 않는다 — 끄는 것은 UTOP 의 결정이다.
    명단에 없고 자동 등록이 꺼져 있으면 **None** 이다(로그인도 막힌다).
    """
    data = _users_load_sync()
    name = str(ju.get("displayName") or "").strip()
    mail = str(ju.get("emailAddress") or "").strip()
    key = str(ju.get("key") or ju.get("accountId") or "").strip()
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    for x in data["users"]:
        if x.get("username") == username:
            if name:
                x["name"] = name
            if mail:
                x["email"] = mail
            if key:
                x["jira_key"] = key
            x["source"] = "jira"
            x["last_login"] = now
            _users_save_sync(data)
            return x
    if not _jira_auto_create():
        # 명단에 없는 사람은 여기서 끝난다 — 관리자가 계정 관리에서 먼저 등록한다
        print(f"[jira-login] 명단에 없어 막았습니다: {username}", flush=True)
        return None
    nu = {
        "id": username, "username": username, "name": name or username,
        "role": "팀원", "email": mail, "active": True, "source": "jira",
        "jira_key": key,
        "created_at": now, "last_login": now,
    }
    data["users"].append(nu)
    _users_save_sync(data)
    print(f"[jira-login] 새 사용자 등록: {username} ({name})", flush=True)
    return nu


def _touch_login(username: str) -> None:
    """마지막으로 들어온 때 — 계정 관리에서 「쓰는 사람·안 쓰는 사람」 을 가른다."""
    try:
        data = _users_load_sync()
        for x in data["users"]:
            if x.get("username") == username:
                x["last_login"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                _users_save_sync(data)
                return
    except Exception:
        pass


def _issue_session(u: dict) -> dict:
    """세션 하나 발급. 동시 로그인은 그대로 허용한다 —
    새로 들어왔다고 남의(내 다른 자리의) 세션을 끊지 않는다."""
    token = _secrets.token_hex(16)
    SESSIONS[token] = {
        "username": u["username"], "role": u.get("role"),
        "name": u.get("name"), "ts": datetime.now().timestamp(),
    }
    _save_one_session(token)   # 새 세션 하나만 저장 (전체 save 는 부하 큼)
    return {"token": token, "user": _public_user(u)}


@app.post("/api/login")
async def api_login(req: LoginReq):
    uname = req.username.strip()

    # ① 연속 실패로 잠긴 동안은 Jira 까지 가지 않는다
    lock = _LOGIN_FAILS.get(uname)
    if lock and lock.get("until", 0) > _t.time():
        left = int(lock["until"] - _t.time()) + 1
        raise HTTPException(429, f"로그인 시도가 많습니다 — {left}초 뒤에 다시 하세요")

    def _fail(msg: str):
        n = (lock.get("n", 0) if lock else 0) + 1
        _LOGIN_FAILS[uname] = {
            "n": n,
            "until": _t.time() + _LOGIN_LOCK_SEC if n >= _LOGIN_FAIL_MAX else 0,
        }
        raise HTTPException(401, msg)

    u = _find_user(uname)

    # ② **먼저 UTOP 비밀번호**를 본다.
    #
    #   Jira 를 먼저 부르면 세 가지가 한꺼번에 무너진다(지적: Jira 연동을 켜면
    #   기존 계정으로 로그인이 안 된다):
    #     · Jira 가 죽으면 admin 도 못 들어온다 — 되돌릴 손이 없어진다
    #     · 같은 아이디가 Jira 에도 있으면 실패가 쌓여 CAPTCHA 가 걸린다
    #     · 로그인마다 바깥 서버를 기다린다
    #   로컬 비밀번호는 우리 손 안에 있고 즉시 판가름 난다. 그것부터 본다.
    if (
        u
        and u.get("active", True)
        and u.get("password")
        and _hash_pw(req.password, u.get("salt", "")) == u.get("password")
    ):
        _LOGIN_FAILS.pop(uname, None)
        _touch_login(uname)
        return _issue_session(u)

    # ③ 안 맞으면 그때 Jira 에 물어본다 — 회원가입 없이 들어오는 길
    if _jira_login_on():
        ju, why = await _jira_verify_login(uname, req.password)
        if ju:
            if u and not u.get("active", True):
                raise HTTPException(401, "관리자가 꺼 둔 계정입니다 — 시스템 담당자에게 문의하세요")
            # Jira 에서 잠긴 계정은 UTOP 에도 못 들어온다 — Jira 가 정본이다
            if ju.get("active") is False:
                raise HTTPException(401, "Jira 에서 잠긴 계정입니다 — Jira 담당자에게 문의하세요")
            u2 = _upsert_jira_user(uname, ju)
            if not u2:
                raise HTTPException(
                    401,
                    "등록되지 않은 계정입니다 — 관리자에게 계정 등록을 요청하세요",
                )
            _LOGIN_FAILS.pop(uname, None)
            return _issue_session(u2)
        # 여기까지 왔으면 로컬 비밀번호도 Jira 도 아니다. 까닭은 남긴다 —
        # 「왜 안 들어가지나」 를 로그 없이 고칠 수는 없다. 비밀번호는 안 남긴다.
        print(f"[jira-login] 거절: {uname} — {why}", flush=True)
        _JIRA_LAST_FAIL.clear()
        _JIRA_LAST_FAIL.update(
            {"user": uname, "why": why, "at": datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
        )
        if why == "captcha":
            raise HTTPException(
                401,
                "Jira 가 사람 확인(CAPTCHA)을 걸었습니다 — Jira 웹에 한 번 로그인해 풀고 다시 시도하세요",
            )
        if why == "cert":
            raise HTTPException(
                401,
                "Jira 인증서에 문제가 있어 물어보지 못했습니다(만료 등) — "
                "SETUP → Jira 연동에서 「TLS 인증서 검증」 을 끄거나 인증서를 갱신하세요",
            )
        if why == "unreachable":
            raise HTTPException(
                401, "Jira 서버에 닿지 못했습니다 — UTOP 비밀번호가 있는 계정으로 들어오세요"
            )

    _fail("아이디 또는 비밀번호가 올바르지 않거나 비활성 계정입니다")

@app.post("/api/logout")
async def api_logout(payload: dict):
    _tok = payload.get("token", "")
    SESSIONS.pop(_tok, None)
    if _tok:
        _delete_one_session(_tok)   # 해당 세션 하나만 삭제
    return {"ok": True}

@app.get("/api/me")
async def api_me(request: Request, token: str = ""):
    # 미들웨어가 이미 세션을 확인해 request.state.user 에 넣어 뒀다.
    # 옛 화면은 아직 ?token= 으로 부르므로 그 경로도 남긴다.
    s = getattr(request.state, "user", None)
    u = _find_user(s.get("username")) if s else _user_from_token(token)
    if not u:
        raise HTTPException(401, "세션이 없습니다")
    # 상단바가 이름 뒤에 팀·소속담당을 적는다 — 그 값은 **조직도가 정본**이다
    w = _org_where(u.get("name") or u.get("username") or "")
    return {"user": {**_public_user(u), **{k: v for k, v in w.items() if v}}}

@app.post("/api/me/avatar")
async def api_me_avatar(payload: dict, token: str = ""):
    u = _user_from_token(token)
    if not u:
        raise HTTPException(401, "세션이 없습니다")
    av = str(payload.get("avatar") or "")
    if len(av) > 400000:
        raise HTTPException(400, "이미지가 너무 큽니다 — 더 작게 줄여 주세요")
    data = _users_load_sync()
    for x in data["users"]:
        if x.get("username") == u.get("username"):
            if av:
                x["avatar"] = av
            else:
                x.pop("avatar", None)
            break
    _users_save_sync(data)
    return {"ok": True, "user": _public_user(_find_user(u.get("username")))}

@app.post("/api/me/change-password")
async def api_me_change_password(payload: dict, token: str = ""):
    """본인 비밀번호 변경 — 현재 비밀번호 검증 후 새 비밀번호로 교체.
    성공 시 이 사용자의 모든 세션(다른 브라우저 포함)을 종료 → 프론트가 로그아웃 처리한다."""
    u = _user_from_token(token)
    if not u:
        raise HTTPException(401, "세션이 없습니다")
    cur = str(payload.get("current_password") or "")
    new = str(payload.get("new_password") or "")
    if not cur or not new:
        raise HTTPException(400, "현재 비밀번호와 새 비밀번호를 모두 입력하세요")
    if len(new) < 4:
        raise HTTPException(400, "새 비밀번호는 4자 이상이어야 합니다")
    if new == cur:
        raise HTTPException(400, "새 비밀번호가 현재 비밀번호와 같습니다")
    # 최신 상태 재로드 (다른 관리자가 방금 바꿨을 수 있으므로 파일에서 다시 읽음)
    data = _users_load_sync()
    target = None
    for x in data["users"]:
        if x.get("username") == u.get("username"):
            target = x
            break
    if not target:
        raise HTTPException(404, "사용자를 찾을 수 없습니다")
    # 현재 비밀번호 검증
    if _hash_pw(cur, target.get("salt", "")) != target.get("password"):
        raise HTTPException(401, "현재 비밀번호가 올바르지 않습니다")
    # 새 salt + 새 hash 로 교체 (salt 재생성으로 이전 hash 사용 불가)
    salt = _secrets.token_hex(8)
    target["salt"] = salt
    target["password"] = _hash_pw(new, salt)
    _users_save_sync(data)
    # 이 사용자의 모든 세션 종료 → 즉시 재로그인 필요
    uname = u.get("username")
    to_drop = [k for k, v in SESSIONS.items() if v.get("username") == uname]
    for k in to_drop:
        SESSIONS.pop(k, None)
        _delete_one_session(k)
    return {"ok": True}

@app.get("/api/users")
async def api_users(token: str = ""):
    _require_admin(token)
    return {
        "users": [
            {**_public_user(u), "retired": _is_retired(u), "org": _org_of(u)}
            for u in _users_load_sync()["users"]
        ],
        "roles": ROLES,
    }

@app.post("/api/users")
async def api_user_create(payload: dict, token: str = ""):
    _require_admin(token)
    data = _users_load_sync()
    uname = str(payload.get("username", "")).strip()
    if not uname:
        raise HTTPException(400, "아이디를 입력하세요")
    if any(u.get("username") == uname for u in data["users"]):
        raise HTTPException(400, "이미 존재하는 아이디입니다")
    role = payload.get("role") if payload.get("role") in ROLES else "팀원"
    email = str(payload.get("email", "")).strip()
    if not email:
        raise HTTPException(400, "이메일을 입력하세요 (필수)")
    if not _valid_email(email):
        raise HTTPException(400, "이메일 형식이 올바르지 않습니다")
    if not _allowed_email_domain(email):
        raise HTTPException(400, "@" + ALLOWED_EMAIL_DOMAIN + " 이메일만 등록할 수 있습니다")
    salt = _secrets.token_hex(8)
    pw = payload.get("password") or "1234"
    nu = {"id": uname, "username": uname, "name": payload.get("name") or uname, "role": role,
          "email": email,
          "company": str(payload.get("company", "")).strip(), "position": str(payload.get("position", "")).strip(),
          "duty": str(payload.get("duty", "")).strip(),
          "dept": str(payload.get("dept", "")).strip(), "team": str(payload.get("team", "")).strip(),
          "salt": salt, "password": _hash_pw(pw, salt), "active": bool(payload.get("active", True)),
          "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
    data["users"].append(nu)
    _users_save_sync(data)
    return {"ok": True, "user": _public_user(nu)}

@app.put("/api/users/{username}")
async def api_user_update(username: str, payload: dict, token: str = "", request: Request = None):
    _require_admin(token)
    data = _users_load_sync()
    for u in data["users"]:
        if u.get("username") == username:
            if payload.get("name") is not None:
                u["name"] = payload.get("name")
            if payload.get("company") is not None:
                u["company"] = str(payload.get("company") or "").strip()   # 회사
            if payload.get("dept") is not None:
                u["dept"] = str(payload.get("dept") or "").strip()   # 소속담당
            if payload.get("team") is not None:
                u["team"] = str(payload.get("team") or "").strip()   # 소속팀
            if payload.get("position") is not None:
                u["position"] = str(payload.get("position") or "").strip()   # 직책
            if payload.get("duty") is not None:
                u["duty"] = str(payload.get("duty") or "").strip()   # 보직
            if payload.get("email") is not None:
                _em = str(payload.get("email")).strip()
                if _em and not _valid_email(_em):
                    raise HTTPException(400, "이메일 형식이 올바르지 않습니다")
                if _em and not _allowed_email_domain(_em):
                    raise HTTPException(400, "@" + ALLOWED_EMAIL_DOMAIN + " 이메일만 사용할 수 있습니다")
                u["email"] = _em
            if payload.get("role") in ROLES:
                u["role"] = payload.get("role")
                # 관리자가 손으로 정한 역할은 **못**이 된다 — Jira 동기화(role_by=jira)가
                # 다음에 팀장으로 되돌리지 않게, jira 표식을 뗀다.
                u.pop("role_by", None)
            if payload.get("active") is not None:
                _was_pending = bool(u.get("pending"))
                u["active"] = bool(payload.get("active"))
                if u["active"]:
                    u["pending"] = False
                    _mc = _load_mail_cfg()
                    if _was_pending and u.get("email") and _mc.get("enabled"):
                        try:
                            _subj = _mc.get("approval_subject") or _DEFAULT_APPROVAL_SUBJECT
                            _tpl = _mc.get("approval_html") or _DEFAULT_APPROVAL_TPL
                            # app_url: 설정값 우선, 없으면 관리자가 접속한 주소(request.base_url)로 자동 유도 → 로그인 버튼 링크 생성
                            _app_url = (_mc.get("app_url") or "").strip()
                            if not _app_url and request is not None:
                                try:
                                    _app_url = str(request.base_url).strip().rstrip("/")
                                except Exception:
                                    _app_url = ""
                            _send_mail(u["email"], _subj,
                                       _render_mail_tpl(_tpl, u.get("name"), u.get("username"), u.get("email"),
                                                        _app_url, u.get("dept"), u.get("team"),
                                                        u.get("position"), u.get("duty")),
                                       html=True)
                        except Exception:
                            pass
            if payload.get("password"):
                salt = _secrets.token_hex(8)
                u["salt"] = salt
                u["password"] = _hash_pw(payload.get("password"), salt)
            _users_save_sync(data)
            return {"ok": True, "user": _public_user(u)}
    raise HTTPException(404, "사용자를 찾을 수 없습니다")

@app.delete("/api/users/{username}")
async def api_user_delete(username: str, token: str = ""):
    _require_admin(token)
    if username == "admin":
        raise HTTPException(400, "기본 관리자(admin)는 삭제할 수 없습니다")
    data = _users_load_sync()
    before = len(data["users"])
    data["users"] = [u for u in data["users"] if u.get("username") != username]
    if len(data["users"]) == before:
        raise HTTPException(404, "사용자를 찾을 수 없습니다")
    _users_save_sync(data)
    return {"ok": True}

def _org_of(u: dict) -> str:
    """이 사람의 조직.

    관리자가 정한 소속(dept)이 있으면 그것이 정본이다. 없으면 **이름 괄호**에서
    뽑는다 — 「강경묵(생산)」·「김대원(SW3)」 처럼 Jira 표시이름이 조직을 달고
    온다. 동기화를 다시 돌리지 않아도 조직별로 묶을 수 있게 하려는 것이다.
    """
    d = str(u.get("dept") or "").strip()
    if d:
        return d
    m = re.search(r"\(([^)]+)\)", str(u.get("name") or ""))
    if not m:
        return ""
    v = m.group(1).strip()
    # 「퇴사자」·「퇴사-비활성화불가」 는 조직이 아니다
    return "" if "퇴사" in v else v


def _org_where(name: str) -> dict:
    """조직도에서 이 사람의 **팀·담당**을 찾는다.

    소속담당은 조직도가 정본이다. dept 는 Jira 표시이름 꼬리에서 뽑은 값이라
    「전규종(검증)」 → 「검증」 이 되는데, 그 사람은 실제로 **품질보증담당**의
    장이다(지적). 두 값이 다르면 조직도를 따른다.

    길에서 이름이 「…팀」 으로 끝나는 마디를 팀으로, 「…담당」 으로 끝나는
    마디를 담당으로 본다. 겸임이면 처음 찾은 자리를 쓴다 — 어느 쪽인지는
    사람이 조직도에서 정한다.
    """
    key = re.split(r"[(\[_]", str(name or ""))[0].replace(" ", "")
    if not key:
        return {}
    org = _kv_load_sync("org_tree", None)
    if not isinstance(org, dict):
        return {}
    found: dict = {}

    def walk(n: dict, path: list):
        nonlocal found
        if found:
            return
        p = [*path, str(n.get("name") or "")]
        names = []
        t = str(n.get("lead") or "").strip()
        if t:
            i = t.rfind(" ")
            names.append(t[:i] if i > 0 else t)
        names += [str(m.get("name") or "") for m in (n.get("members") or [])]
        if any(re.split(r"[(\[_]", x)[0].replace(" ", "") == key for x in names):
            found = {
                "team": next((x for x in reversed(p) if x.endswith("팀")), ""),
                "dept": next((x for x in reversed(p) if x.endswith("담당")), ""),
                "org_path": " › ".join(p[1:]),
            }
            return
        for c in (n.get("children") or []):
            walk(c, p)

    walk(org, [])
    return found


def _is_retired(u: dict) -> bool:
    """나간 사람인가.

    **Jira 비활성만 보면 안 된다**(지적: 퇴사 계정이 안 지워진다). 실제 자료를
    보면 나간 사람이 계정은 살아 있고 **이름에 표시**만 달려 있다 —
    「김진보(퇴사자)」·「김대환(Bilab) (퇴사-비활성화불가)」 처럼. 그래서 이름의
    「퇴사」 표기도 함께 본다. 둘 중 하나면 나간 사람이다.
    """
    if u.get("jira_active") is False:
        return True
    return "퇴사" in str(u.get("name") or "")


@app.get("/api/org")
async def api_org_get():
    """조직도 — 회사 → 그룹 → 담당 → 팀 → 사람.

    사람은 `{name, rank}` 다. **직급(rank)은 Jira 에 없다**(확인함) — 사람이
    준 조직도가 정본이고, 여기 담아 둔다. 계정과는 **이름으로** 잇는다:
    계정 이름이 「강경묵(생산)」 처럼 꼬리를 달고 있어 괄호·밑줄 앞까지만 본다.
    """
    return {"org": _kv_load_sync("org_tree", None)}


@app.post("/api/org")
async def api_org_save(payload: dict, token: str = ""):
    """조직도 통째로 저장. 관리자만."""
    _require_admin(token)
    org = (payload or {}).get("org")
    if not isinstance(org, dict) or not org.get("name"):
        raise HTTPException(400, "조직도 모양이 아닙니다")
    _kv_save_sync("org_tree", org)
    return {"ok": True, **_apply_org_roles(org)}


@app.post("/api/org/member-role")
async def api_org_member_role(payload: dict, token: str = ""):
    """조직도의 **계정 없는 사람**에게 역할을 준다. 관리자만.

    조직도 204명 중 40명은 계정이 아예 없다(확인함 — 이름 표기가 달라 못
    이어진 것이 아니라, 그 이름이 계정 목록 어디에도 없다). Jira 계정이
    없으면 UTOP 을 안 쓰는 사람이라 계정을 만들어 줄 수도 없다.

    그런데도 역할은 적어 두어야 한다(지시) — 조직도는 「누가 무엇을 맡나」
    를 보는 표이지 「누가 UTOP 을 쓰나」 만 보는 표가 아니다. 그래서 역할을
    **조직도 그 사람 칸에** 담는다. 나중에 그 사람의 계정이 생기면 계정 쪽
    역할이 이 값을 덮는다 — 계정이 있으면 계정이 정본이다.
    """
    _require_admin(token)
    nm = str((payload or {}).get("name") or "").strip()
    role = str((payload or {}).get("role") or "").strip()
    if not nm:
        raise HTTPException(400, "이름이 없습니다")
    org = _kv_load_sync("org_tree", None)
    if not isinstance(org, dict):
        raise HTTPException(404, "조직도가 없습니다")

    hit = 0

    def walk(n: dict):
        nonlocal hit
        for m in (n.get("members") or []):
            if str(m.get("name") or "").strip() == nm:
                if role:
                    m["role"] = role
                else:
                    m.pop("role", None)
                hit += 1
        for c in (n.get("children") or []):
            walk(c)

    walk(org)
    if not hit:
        raise HTTPException(404, f"조직도에 「{nm}」 이(가) 없습니다")
    _kv_save_sync("org_tree", org)
    return {"ok": True, "hit": hit}


def _org_walk(n: dict):
    """마디를 하나씩 내준다 — 뿌리부터 깊이 우선."""
    yield n
    for c in (n.get("children") or []):
        yield from _org_walk(c)


def _org_at(org: dict, path: list) -> dict | None:
    """이름 길로 마디를 찾는다. 같은 이름이 여러 곳에 있어도(사업1담당 밑
    네트워크사업1팀 처럼) 길로 찾으면 헷갈리지 않는다."""
    cur = org
    for step in (path or [])[1:]:
        nxt = None
        for c in (cur.get("children") or []):
            if str(c.get("name")) == str(step):
                nxt = c
                break
        if nxt is None:
            return None
        cur = nxt
    return cur if not path or str(org.get("name")) == str(path[0]) else None


@app.post("/api/org/seed")
async def api_org_seed(payload: dict = None, token: str = ""):
    """이미지에 실린 조직도를 **손으로 심는다**. 관리자만.

    시작할 때 자동으로 심지만(비어 있을 때만), 그게 안 먹은 서버에서는
    확인할 길이 없었다 — 서버에 들어갈 수 없으면 「왜 안 됐나」 를 물을 데가
    없다(253 지적). 눌러서 심고 **결과를 눈으로 보게** 한다.

    이미 조직도가 있으면 안 덮는다. 덮으려면 force 를 줘야 한다 — 화면이
    먼저 물어본 뒤에 보낸다.
    """
    _require_admin(token)
    force = bool((payload or {}).get("force"))
    cur = _kv_load_sync("org_tree", None)
    seed = Path(__file__).parent / "seed" / "org_tree.json"
    if not seed.exists():
        raise HTTPException(404, "씨앗 파일이 이미지에 없습니다 — 코드를 다시 받으세요")
    if cur and not force:
        return {"ok": False, "had": True, "name": cur.get("name"),
                "detail": "이미 조직도가 있습니다"}
    org = json.loads(seed.read_text(encoding="utf-8"))
    _kv_save_sync("org_tree", org)

    people: set = set()

    def walk(n: dict):
        t = str(n.get("lead") or "").strip()
        if t:
            i = t.rfind(" ")
            people.add(t[:i] if i > 0 else t)
        for m in (n.get("members") or []):
            people.add(str(m.get("name") or ""))
        for c in (n.get("children") or []):
            walk(c)

    walk(org)
    people.discard("")
    return {"ok": True, "had": bool(cur), "nodes": len(list(_org_walk(org))),
            "people": len(people), **_apply_org_roles(org)}


@app.post("/api/org/node")
async def api_org_node_add(payload: dict, token: str = ""):
    """조직을 하나 만든다 — 고른 조직 **아래**에. 관리자만."""
    _require_admin(token)
    path = (payload or {}).get("path") or []
    name = str((payload or {}).get("name") or "").strip()
    if not name:
        raise HTTPException(400, "조직 이름이 없습니다")
    org = _kv_load_sync("org_tree", None)
    if not isinstance(org, dict):
        raise HTTPException(404, "조직도가 없습니다")
    at = _org_at(org, path)
    if at is None:
        raise HTTPException(404, "그 조직을 못 찾았습니다")
    kids = at.setdefault("children", [])
    if any(str(c.get("name")) == name for c in kids):
        raise HTTPException(400, f"「{name}」 은(는) 이미 있습니다")
    kids.append({"name": name})
    _kv_save_sync("org_tree", org)
    return {"ok": True}


@app.post("/api/org/rename")
async def api_org_rename(payload: dict, token: str = ""):
    """조직 이름을 바꾼다. 관리자만."""
    _require_admin(token)
    path = (payload or {}).get("path") or []
    name = str((payload or {}).get("name") or "").strip()
    if not name:
        raise HTTPException(400, "새 이름이 없습니다")
    org = _kv_load_sync("org_tree", None)
    if not isinstance(org, dict):
        raise HTTPException(404, "조직도가 없습니다")
    at = _org_at(org, path)
    if at is None:
        raise HTTPException(404, "그 조직을 못 찾았습니다")
    at["name"] = name
    _kv_save_sync("org_tree", org)
    return {"ok": True}


@app.post("/api/org/delete-node")
async def api_org_node_del(payload: dict, token: str = ""):
    """**빈 조직만** 지운다. 관리자만.

    사람이나 하위 조직이 든 마디를 지우면 그 사람들이 조직도에서 통째로
    사라진다 — 되돌릴 방법이 없다. 비었을 때만 지우게 해, 잘못 만든 것을
    치우는 데만 쓰이게 한다.
    """
    _require_admin(token)
    path = (payload or {}).get("path") or []
    if len(path) < 2:
        raise HTTPException(400, "맨 위 조직은 못 지웁니다")
    org = _kv_load_sync("org_tree", None)
    if not isinstance(org, dict):
        raise HTTPException(404, "조직도가 없습니다")
    parent = _org_at(org, path[:-1])
    if parent is None:
        raise HTTPException(404, "그 조직을 못 찾았습니다")
    gone = None
    for c in (parent.get("children") or []):
        if str(c.get("name")) == str(path[-1]):
            gone = c
            break
    if gone is None:
        raise HTTPException(404, "그 조직을 못 찾았습니다")
    if (gone.get("members") or []) or (gone.get("children") or []) or gone.get("lead"):
        raise HTTPException(400, "빈 조직만 지울 수 있습니다 — 먼저 사람을 옮기세요")
    parent["children"] = [c for c in parent["children"] if c is not gone]
    _kv_save_sync("org_tree", org)
    return {"ok": True}


@app.post("/api/org/move-member")
async def api_org_move_member(payload: dict, token: str = ""):
    """사람을 다른 조직으로 옮긴다. 관리자만.

    직급·역할은 사람에게 붙은 것이라 그대로 들고 간다 — 옮겼다고 직급이
    지워지면 옮기기가 두려운 기능이 된다.
    """
    _require_admin(token)
    nm = str((payload or {}).get("name") or "").strip()
    to = (payload or {}).get("to") or []
    if not nm:
        raise HTTPException(400, "이름이 없습니다")
    org = _kv_load_sync("org_tree", None)
    if not isinstance(org, dict):
        raise HTTPException(404, "조직도가 없습니다")
    # 옮길 곳이 비었으면 **조직도에서 뺀다**(「(조직도에 없음)」 을 고른 것).
    # 넣기만 되고 빼기가 없으면, 시험 삼아 넣어 본 사람을 되돌릴 길이 없다
    # (지적: 조직도에 없는 계정으로 다시 못 바꾼다).
    dest = _org_at(org, to) if to else None
    if to and dest is None:
        raise HTTPException(404, "옮길 조직을 못 찾았습니다")

    picked = None

    def strip(n: dict):
        nonlocal picked
        keep = []
        for m in (n.get("members") or []):
            if str(m.get("name") or "").strip() == nm and picked is None:
                picked = m
            else:
                keep.append(m)
        if n.get("members") is not None:
            n["members"] = keep
        for c in (n.get("children") or []):
            strip(c)

    strip(org)
    if dest is None:
        # 빼기 — 조직도에 없던 사람이면 이미 목적을 이룬 것이라 조용히 넘긴다
        if picked is not None:
            _kv_save_sync("org_tree", org)
        return {"ok": True, "to": None, "removed": picked is not None}
    if picked is None:
        # 조직도 어디에도 없던 사람 — **새로 넣는다.** 계정은 있는데 조직도에
        # 이름이 없는 사람이 45명이다(admin·qag 를 포함해). 없다고 물리면
        # 그 45명을 조직에 넣을 방법이 아예 없다.
        for n2 in _org_walk(org):
            L = str(n2.get("lead") or "").strip()
            i2 = L.rfind(" ")
            if (L[:i2] if i2 > 0 else L) == nm:
                raise HTTPException(400, f"「{nm}」 은(는) 그 조직의 장입니다 — 조직의 장을 바꾸세요")
        picked = {"name": nm}
    dest.setdefault("members", []).append(picked)
    _kv_save_sync("org_tree", org)
    return {"ok": True, "to": dest.get("name")}


def _apply_org_roles(org: dict) -> dict:
    """조직도의 **장**을 계정 역할 「담당」 으로 맞춘다.

    표에만 「담당」 이라 적고 실제 역할은 팀원이면, 눌러서 열어 본 사람이
    두 값을 보고 어느 쪽이 맞는지 알 수 없다(지적). 실제 역할을 바꾼다.

    **관리자는 절대 안 내린다** — 조직의 장이라고 관리자 권한을 뺏으면
    그 사람이 화면을 못 쓴다(실제로 전규종이 관리자다). 팀원·팀장만 올린다.

    표식(`role_by='org'`)을 남겨, 장에서 내려오면 **우리가 올린 것만** 되돌린다.
    관리자가 손으로 정한 역할은 건드리지 않는다 — Jira 동기화의 규칙과 같다.
    """
    leads: set = set()

    def walk(n: dict):
        t = str(n.get("lead") or "").strip()
        if t:
            i = t.rfind(" ")
            leads.add((t[:i] if i > 0 else t))
        for c in (n.get("children") or []):
            walk(c)

    walk(org)

    def key(v) -> str:
        return re.split(r"[(\[_]", str(v or ""))[0].replace(" ", "")

    lead_keys = {key(x) for x in leads} | {key(re.sub(r"\d+$", "", x)) for x in leads}
    data = _users_load_sync()
    up = down = 0
    for u in data["users"]:
        if u.get("role") == "관리자":
            continue  # 관리자는 안 내린다
        k = key(u.get("name") or u.get("username"))
        if k and k in lead_keys:
            if u.get("role") != "담당":
                u["role"] = "담당"
                u["role_by"] = "org"
                up += 1
        elif u.get("role") == "담당" and u.get("role_by") == "org":
            u["role"] = "팀원"
            u.pop("role_by", None)
            down += 1
    if up or down:
        _users_save_sync(data)
    return {"role_up": up, "role_down": down}


@app.post("/api/users/delete-retired")
async def api_users_delete_retired(token: str = ""):
    """Jira 에서 나간 사람(jira_active=False)을 **명단에서 지운다**(지시: 퇴사자
    필요 없음).

    기록은 안 깨진다 — 사이클의 실행자·담당자는 이름 **문자열**로 담겨 있어
    (FK 가 아니다) 계정을 지워도 그 기록의 이름은 남는다. admin 은 못 지운다.

    다음 동기화에서 되살아나지 않게, 동기화 쪽에서 비활성 신규는 안 만든다.
    """
    _require_admin(token)
    data = _users_load_sync()
    before = len(data["users"])
    kept, gone = [], []
    for u in data["users"]:
        if u.get("username") != "admin" and _is_retired(u):
            gone.append(u.get("username"))
        else:
            kept.append(u)
    data["users"] = kept
    _users_save_sync(data)
    return {"ok": True, "deleted": before - len(kept), "names": gone[:50]}


# init_users_file() 은 startup 훅에서 호출 (모듈 로드 시점엔 DB 풀 없음)

# ───────────────────────────────────────────
# 메일(SMTP) 설정 + 발송
# ───────────────────────────────────────────
MAIL_FILE = DATA_DIR / "integrations" / "mail.json"
_MAIL_DEFAULT = {"host": "", "port": 587, "username": "", "password": "",
                 "from_addr": "", "from_name": "ubiQuoss-TOP", "security": "starttls", "enabled": False}

def _load_mail_cfg() -> dict:
    cfg = dict(_MAIL_DEFAULT)
    try:
        if MAIL_FILE.exists():
            d = load_json(MAIL_FILE)
            if isinstance(d, dict):
                cfg.update(d)   # 저장된 모든 키 보존(approval_*, app_url, share_* 등 포함) + 기본값으로 누락 보완
    except Exception:
        pass
    return cfg

def _save_mail_cfg(cfg: dict):
    MAIL_FILE.parent.mkdir(parents=True, exist_ok=True)
    save_json(MAIL_FILE, cfg)

def _valid_email(addr: str) -> bool:
    import re as _re_mail
    return bool(_re_mail.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", str(addr or "").strip()))

# 가입 허용 이메일 도메인 (회사 메일만 허용)
ALLOWED_EMAIL_DOMAIN = "ubiquoss.com"
def _allowed_email_domain(addr: str) -> bool:
    return str(addr or "").strip().lower().endswith("@" + ALLOWED_EMAIL_DOMAIN)

def _send_mail(to_addrs, subject: str, body: str, html: bool = False):
    """SMTP로 메일 발송. to_addrs: str(콤마/세미콜론 구분) 또는 list. 실패 시 예외 발생."""
    import smtplib, ssl as _ssl
    from email.message import EmailMessage
    cfg = _load_mail_cfg()
    if not cfg.get("host"):
        raise RuntimeError("SMTP 서버가 설정되지 않았습니다 (시스템 → 메일 설정)")
    if isinstance(to_addrs, str):
        to_list = [a.strip() for a in to_addrs.replace(";", ",").split(",") if a.strip()]
    else:
        to_list = [a for a in (to_addrs or []) if a]
    if not to_list:
        raise RuntimeError("받는 사람이 없습니다")
    msg = EmailMessage()
    from_addr = cfg.get("from_addr") or cfg.get("username")
    msg["From"] = f'{cfg.get("from_name") or "ubiQuoss-TOP"} <{from_addr}>'
    msg["To"] = ", ".join(to_list)
    msg["Subject"] = subject
    if html:
        msg.set_content("이 메일은 HTML 형식입니다. HTML을 지원하는 클라이언트에서 열어주세요.")
        msg.add_alternative(body, subtype="html")
    else:
        msg.set_content(body)
    host = cfg["host"]; port = int(cfg.get("port") or 587); sec = str(cfg.get("security") or "starttls").lower()
    if sec == "ssl":
        ctx = _ssl.create_default_context()
        with smtplib.SMTP_SSL(host, port, timeout=20, context=ctx) as s:
            if cfg.get("username"):
                s.login(cfg["username"], cfg.get("password") or "")
            s.send_message(msg)
    else:
        with smtplib.SMTP(host, port, timeout=20) as s:
            s.ehlo()
            if sec == "starttls":
                s.starttls(context=_ssl.create_default_context()); s.ehlo()
            if cfg.get("username"):
                s.login(cfg["username"], cfg.get("password") or "")
            s.send_message(msg)
    return to_list

_DEFAULT_APPROVAL_SUBJECT = "[ubiQuoss-TOP] \U0001F389 가입이 승인되었습니다"

# 사이클 배정 알림 메일 기본 폼 (메일 설정 → 사이클 배정 폼에서 편집 가능)
# 플레이스홀더: {{assignee}} {{model}} {{vgroup}} {{version}} {{period}} {{count}} {{items}} {{app_url}} {{login_button}}
_DEFAULT_CYCLE_SUBJECT = "[ubiQuoss-TOP] 시험 사이클 배정 — {{model}} {{version}}"
_DEFAULT_CYCLE_TPL = """<!DOCTYPE html><html><body style="margin:0;padding:0;background:#eef1f6;">
<div style="font-family:'Malgun Gothic','맑은 고딕',Arial,sans-serif;max-width:960px;margin:0 auto;color:#1f2937;">
  <div style="background:linear-gradient(135deg,#2563eb,#4f8ae8);color:#fff;padding:18px 22px;border-radius:11px 11px 0 0;">
    <div style="font-size:18px;font-weight:800;">📋 시험 사이클이 배정되었습니다</div>
    <div style="font-size:12.5px;opacity:.92;margin-top:3px;">{{assignee}} 님, 아래 항목을 시험해 주세요.</div></div>
  <div style="border:1px solid #e3e8ef;border-top:none;border-radius:0 0 11px 11px;padding:20px 22px;">
    <table style="font-size:13px;line-height:1.7;margin-bottom:14px;">
      <tr><td style="color:#6b7280;padding-right:16px;">모델</td><td style="font-weight:700;">{{model}}</td></tr>
      <tr><td style="color:#6b7280;padding-right:16px;">버전 그룹</td><td style="font-weight:700;">{{vgroup}}</td></tr>
      <tr><td style="color:#6b7280;padding-right:16px;">버전</td><td style="font-weight:700;">{{version}}</td></tr>
      <tr><td style="color:#6b7280;padding-right:16px;">시험 기간</td><td>{{period}}</td></tr>
      <tr><td style="color:#6b7280;padding-right:16px;">시험 항목</td><td style="font-weight:700;color:#00875a;">{{count}} 건</td></tr>
    </table>
    <div style="font-size:13px;font-weight:800;color:#374151;margin-bottom:6px;">📝 시험 항목 목록</div>
    {{items}}
    {{login_button}}
    <div style="margin-top:16px;font-size:11px;color:#9ca3af;border-top:1px solid #eef0f4;padding-top:10px;">ubiQuoss-TOP 시험 자동화 플랫폼에서 자동 발송된 메일입니다.</div>
  </div>
</div></body></html>"""

_DEFAULT_APPROVAL_TPL = """<!DOCTYPE html><html><body style="margin:0;padding:0;background:#eef1f6;">
<table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="background:#eef1f6;padding:30px 12px;font-family:'Apple SD Gothic Neo','Malgun Gothic',Arial,sans-serif;">
<tr><td align="center">
<table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="max-width:520px;background:#ffffff;border-radius:16px;overflow:hidden;box-shadow:0 8px 30px rgba(20,40,80,0.10);">
<tr><td bgcolor="#2d6fd4" style="background-color:#2d6fd4;background:linear-gradient(135deg,#2d6fd4,#1b59bd);padding:32px 30px;text-align:center;">
<div style="font-size:23px;font-weight:800;color:#ffffff;letter-spacing:-0.4px;">ubi<span style="color:#ff90a6;">Q</span>uoss-TOP</div>
<div style="font-size:12px;color:#cfe0ff;margin-top:5px;">Ubiquoss Test Orchestration Platform</div>
</td></tr>
<tr><td style="text-align:center;padding:36px 30px 4px;">
<div style="font-size:48px;line-height:1;">\U0001F389</div>
<div style="font-size:22px;font-weight:800;color:#1c2942;margin-top:16px;">가입을 진심으로 환영합니다!</div>
<div style="font-size:14px;color:#5a6b85;margin-top:10px;line-height:1.75;"><b style="color:#2d6fd4;">{{name}}</b>님, 가입 신청이 <b>승인</b>되었습니다.<br>이제 ubiQuoss-TOP의 모든 기능을 사용하실 수 있습니다.</div>
</td></tr>
<tr><td style="padding:24px 30px 8px;">
<table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="background:#f5f8fd;border:1px solid #e1eaf7;border-radius:12px;">
<tr><td style="padding:6px 24px;">
<table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="font-size:14px;border-collapse:collapse;">
<tr><td style="padding:12px 0;border-bottom:1px solid #e6edf8;color:#8a99b5;width:92px;vertical-align:top;">아이디</td><td style="padding:12px 0;border-bottom:1px solid #e6edf8;font-weight:700;color:#1c2942;font-family:ui-monospace,monospace;">{{username}}</td></tr>
<tr><td style="padding:12px 0;border-bottom:1px solid #e6edf8;color:#8a99b5;vertical-align:top;">이름</td><td style="padding:12px 0;border-bottom:1px solid #e6edf8;font-weight:700;color:#1c2942;">{{name}}</td></tr>
<tr><td style="padding:12px 0;border-bottom:1px solid #e6edf8;color:#8a99b5;vertical-align:top;">소속담당</td><td style="padding:12px 0;border-bottom:1px solid #e6edf8;font-weight:700;color:#1c2942;">{{dept}}</td></tr>
<tr><td style="padding:12px 0;border-bottom:1px solid #e6edf8;color:#8a99b5;vertical-align:top;">소속팀</td><td style="padding:12px 0;border-bottom:1px solid #e6edf8;font-weight:700;color:#1c2942;">{{team}}</td></tr>
<tr><td style="padding:12px 0;border-bottom:1px solid #e6edf8;color:#8a99b5;vertical-align:top;">직책</td><td style="padding:12px 0;border-bottom:1px solid #e6edf8;font-weight:700;color:#1c2942;">{{position}}</td></tr>
<tr><td style="padding:12px 0;border-bottom:1px solid #e6edf8;color:#8a99b5;vertical-align:top;">보직</td><td style="padding:12px 0;border-bottom:1px solid #e6edf8;font-weight:700;color:#1c2942;">{{duty}}</td></tr>
<tr><td style="padding:12px 0;border-bottom:1px solid #e6edf8;color:#8a99b5;vertical-align:top;">메일</td><td style="padding:12px 0;border-bottom:1px solid #e6edf8;font-weight:700;color:#1c2942;">{{email}}</td></tr>
<tr><td style="padding:12px 0;color:#8a99b5;vertical-align:middle;">상태</td><td style="padding:12px 0;"><span style="display:inline-block;background:#e3f6ec;color:#00875a;font-size:12px;font-weight:700;padding:4px 13px;border-radius:20px;">&#10003; 승인 완료</span></td></tr>
</table>
</td></tr></table></td></tr>
<tr><td align="center" style="padding:22px 30px 36px;">{{login_button}}</td></tr>
<tr><td style="background:#f7f9fc;border-top:1px solid #eef1f6;padding:18px 30px;text-align:center;font-size:11.5px;color:#9aa7bd;line-height:1.7;">본 메일은 ubiQuoss-TOP 가입 승인에 따라 자동 발송되었습니다.<br>문의는 시스템 관리자에게 연락해 주세요.</td></tr>
</table></td></tr></table></body></html>"""

def _login_button_html(app_url: str = "") -> str:
    if app_url:
        url = app_url.strip().rstrip("/")
        return ('<a href="' + url + '" target="_blank" '
                'style="display:inline-block;background:#2d6fd4;color:#ffffff;text-decoration:none;'
                'font-size:15px;font-weight:700;padding:14px 40px;border-radius:10px;'
                'box-shadow:0 6px 16px rgba(45,111,212,0.35);">로그인 하러 가기 &rarr;</a>')
    return '<span style="font-size:13px;color:#8090ab;">로그인 페이지에서 로그인해 주세요.</span>'

def _repair_placeholders(html: str) -> str:
    """WYSIWYG 편집으로 {{ }} 안쪽에 끼어든 HTML 태그/엔티티 제거 → 깨진 플레이스홀더 복원."""
    if not html or "{" not in html:
        return html
    import re as _rp
    keys = "name|username|email|dept|team|position|duty|app_url|login_button"
    junk = r"(?:<[^>]*>|&nbsp;|\s)*"
    pat = r"\{\{" + junk + r"(" + keys + r")" + junk + r"\}\}"
    return _rp.sub(pat, lambda m: "{{" + m.group(1).lower() + "}}", html, flags=_rp.I)

def _render_mail_tpl(tpl: str, name: str = "", username: str = "", email: str = "", app_url: str = "",
                     dept: str = "", team: str = "", position: str = "", duty: str = "") -> str:
    """플레이스홀더 치환: {{name}} {{username}} {{email}} {{dept}} {{team}} {{position}} {{duty}} {{app_url}} {{login_button}}"""
    def _esc(s):
        return str(s or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    out = _repair_placeholders(tpl or "")
    out = out.replace("{{name}}", _esc(name or username or "회원"))
    out = out.replace("{{username}}", _esc(username or ""))
    out = out.replace("{{email}}", _esc(email or ""))
    out = out.replace("{{dept}}", _esc(dept or ""))
    out = out.replace("{{team}}", _esc(team or ""))
    out = out.replace("{{position}}", _esc(position or ""))
    out = out.replace("{{duty}}", _esc(duty or ""))
    out = out.replace("{{app_url}}", (app_url or "").strip().rstrip("/"))
    out = out.replace("{{login_button}}", _login_button_html(app_url))
    if "<html" not in out.lower():   # 위지윅이 본문 조각만 보낸 경우 메일 문서로 감쌈
        out = ('<!DOCTYPE html><html><body style="margin:0;padding:0;background:#eef1f6;">'
               + out + "</body></html>")
    return out


@app.get("/api/mail/config")
async def api_mail_config_get(token: str = ""):
    _require_admin(token)
    return {"config": _load_mail_cfg(),
            "default_approval_subject": _DEFAULT_APPROVAL_SUBJECT,
            "default_approval_html": _DEFAULT_APPROVAL_TPL,
            "default_cycle_subject": _DEFAULT_CYCLE_SUBJECT,
            "default_cycle_html": _DEFAULT_CYCLE_TPL}

@app.post("/api/mail/preview-approval")
async def api_mail_preview_approval(payload: dict, token: str = "", request: Request = None):
    """가입 승인 메일 미리보기 — 입력 HTML을 샘플 데이터로 렌더."""
    _require_admin(token)
    cfg = _load_mail_cfg()
    tpl = payload.get("html") or _DEFAULT_APPROVAL_TPL
    _app_url = (cfg.get("app_url") or "").strip()
    if not _app_url and request is not None:
        try:
            _app_url = str(request.base_url).strip().rstrip("/")
        except Exception:
            _app_url = ""
    html = _render_mail_tpl(tpl,
                            payload.get("name") or "홍길동",
                            payload.get("username") or "hong",
                            payload.get("email") or ("hong@" + ALLOWED_EMAIL_DOMAIN),
                            _app_url,
                            payload.get("dept") or "품질보증담당",
                            payload.get("team") or "QA팀",
                            payload.get("position") or "책임",
                            payload.get("duty") or "팀원")
    return {"ok": True, "html": html}

# ── 브랜딩(로고) ── 관리자가 직접 로고 이미지를 등록 (data URI base64)
BRANDING_FILE = DATA_DIR / "config" / "branding.json"

def _load_branding() -> dict:
    try:
        if BRANDING_FILE.exists():
            d = load_json(BRANDING_FILE)
            if isinstance(d, dict):
                return d
    except Exception:
        pass
    return {}

@app.get("/api/branding")
async def api_branding_get():
    b = _load_branding()
    return {"logo": b.get("logo") or "", "name_text": b.get("name_text") or "",
            "name_size": b.get("name_size") or "", "name_color": b.get("name_color") or "",
            "name_font": b.get("name_font") or "", "name_accent_color": b.get("name_accent_color") or "",
            "fab_greeting": b.get("fab_greeting") or "", "fab_quick": (b.get("fab_quick") if isinstance(b.get("fab_quick"), list) else []),
            "fab_prompt": b.get("fab_prompt") or "", "fab_rules": b.get("fab_rules") or "",
            "link_url": b.get("link_url") or "",
            # 로그인 화면 왼쪽 판 — 회사 건물 사진과 그 위에 얹는 글(지시)
            "login_image": b.get("login_image") or "",
            "login_title": b.get("login_title") or "",
            "login_sub": b.get("login_sub") or "",
            # 로그인 화면 로고·글자 — **메뉴 것과 따로다**(지시: 구분해).
            # 한 값을 둘이 나눠 쓰니 한쪽을 고치면 다른 쪽이 따라 바뀌었다.
            "login_logo": b.get("login_logo") or "",
            "login_size": b.get("login_size") or "",
            "login_color": b.get("login_color") or "",
            "login_accent_color": b.get("login_accent_color") or "",
            "login_font": b.get("login_font") or "",
            # 오른쪽 판(들어가는 자리) — 코드에 박혀 있던 문구를 뺀다(지시)
            "login_form_title": b.get("login_form_title") or "",
            "login_id_ph": b.get("login_id_ph") or "",
            "login_note": b.get("login_note") or "",
            "login_foot": b.get("login_foot") or "",
            "login_md": b.get("login_md") or "",
            "login_body_size": b.get("login_body_size") or "",
            "login_body_color": b.get("login_body_color") or "",
            "login_keep": b.get("login_keep") or ""}

@app.post("/api/branding")
async def api_branding_save(payload: dict, request: Request, token: str = ""):
    # 화면은 Authorization 헤더로 토큰을 보낸다. 여기서 쿼리(?token=)만 보아
    # 401 이 났고, 저장은 조용히 실패해 새로고침하면 옛 값이 돌아왔다
    # (지적: 15 를 16 으로 고쳐도 15 로 되돌아간다).
    _require_admin(token or _token_from(request))
    b = _load_branding()
    for k in ("name_text", "name_size", "name_color", "name_font", "name_accent_color", "link_url",
              "login_title", "login_sub", "login_size", "login_color", "login_accent_color",
              "login_font", "login_form_title", "login_id_ph", "login_note", "login_foot",
              "login_md", "login_body_size", "login_body_color", "login_keep"):
        if k in payload:
            b[k] = str(payload.get(k) or "")[:200]
    if "fab_greeting" in payload:
        b["fab_greeting"] = str(payload.get("fab_greeting") or "")[:1500]
    if "fab_quick" in payload:
        q = payload.get("fab_quick")
        b["fab_quick"] = [str(x)[:200] for x in q if str(x).strip()][:50] if isinstance(q, list) else []
    if "fab_prompt" in payload:
        b["fab_prompt"] = str(payload.get("fab_prompt") or "")[:4000]
    if "fab_rules" in payload:
        b["fab_rules"] = str(payload.get("fab_rules") or "")[:8000]
    save_json(BRANDING_FILE, b)
    return {"ok": True, "name_text": b.get("name_text") or "", "name_size": b.get("name_size") or "",
            "name_color": b.get("name_color") or "", "name_font": b.get("name_font") or "",
            "name_accent_color": b.get("name_accent_color") or ""}

@app.post("/api/branding/login-image")
async def api_branding_login_image(payload: dict, request: Request, token: str = ""):
    """로그인 화면 왼쪽 판에 깔 사진 — 회사 건물처럼 우리 것을 올린다(지시).

    남의 사진을 갖다 쓰지 않는다. 올리는 사람이 권리를 아는 사진이라야
    한다 — 그래서 자동으로 받아 오지 않고 **올리는 자리**만 둔다.
    """
    _require_admin(token or _token_from(request))
    img = str(payload.get("image") or "")
    if img and not img.startswith("data:image/"):
        raise HTTPException(400, "이미지 파일만 등록할 수 있습니다")
    if len(img) > 8_000_000:   # base64 약 6MB — 사진이라 로고보다 넉넉히
        raise HTTPException(400, "이미지가 너무 큽니다 (6MB 이하로 올려주세요)")
    b = _load_branding()
    b["login_image"] = img
    save_json(BRANDING_FILE, b)
    return {"ok": True}


@app.post("/api/branding/login-logo")
async def api_branding_login_logo(payload: dict, request: Request, token: str = ""):
    """로그인 화면 로고 — 메뉴 로고와 **따로** 둔다(지시)."""
    _require_admin(token or _token_from(request))
    logo = str(payload.get("logo") or "")
    if logo and not logo.startswith("data:image/"):
        raise HTTPException(400, "이미지 파일만 등록할 수 있습니다")
    if len(logo) > 4_000_000:
        raise HTTPException(400, "이미지가 너무 큽니다 (3MB 이하로 올려주세요)")
    b = _load_branding()
    b["login_logo"] = logo
    save_json(BRANDING_FILE, b)
    return {"ok": True}


@app.post("/api/branding/logo")
async def api_branding_logo(payload: dict, request: Request, token: str = ""):
    _require_admin(token or _token_from(request))
    logo = str(payload.get("logo") or "")
    if logo and not logo.startswith("data:image/"):
        raise HTTPException(400, "이미지 파일만 등록할 수 있습니다")
    if len(logo) > 4_000_000:   # base64 약 3MB 상한
        raise HTTPException(400, "이미지가 너무 큽니다 (3MB 이하로 올려주세요)")
    b = _load_branding()
    b["logo"] = logo
    save_json(BRANDING_FILE, b)
    return {"ok": True, "logo": logo}

# ── TC/REQ 공유 메일 ── 어떤 섹션을 포함할지 체크로 선택(관리자), 발송은 로그인 사용자
_DEFAULT_SHARE_SUBJECT = "[ubiQuoss-TOP] {id} {title}"
_DEFAULT_REQ_SECTIONS = {"info": True, "desc": True, "impl": True, "scenario": False, "tc": True}
_DEFAULT_TC_SECTIONS = {"info": True, "purpose": True, "topo": True, "traffic": True, "steps": True, "issue": True, "history": True, "cycle": True}

@app.get("/api/share-config")
async def api_share_config_get(token: str = ""):
    u = _user_from_token(token)
    if not u:
        raise HTTPException(401, "로그인이 필요합니다")
    cfg = _load_mail_cfg()
    rs = cfg.get("share_sections_req") or cfg.get("share_sections")   # 구버전 호환
    if not isinstance(rs, dict):
        rs = dict(_DEFAULT_REQ_SECTIONS)
    ts = cfg.get("share_sections_tc")
    if not isinstance(ts, dict):
        ts = dict(_DEFAULT_TC_SECTIONS)
    _osub = cfg.get("share_subject") or _DEFAULT_SHARE_SUBJECT   # 구버전 공통값 → 폴백
    _oin = cfg.get("share_intro") or ""
    _oout = cfg.get("share_outro") or ""
    return {
        "req": {
            "subject": cfg.get("share_req_subject") or _osub,
            "sections": rs,
            "intro": cfg.get("share_req_intro") if cfg.get("share_req_intro") is not None else _oin,
            "outro": cfg.get("share_req_outro") if cfg.get("share_req_outro") is not None else _oout,
        },
        "tc": {
            "subject": cfg.get("share_tc_subject") or _osub,
            "sections": ts,
            "intro": cfg.get("share_tc_intro") if cfg.get("share_tc_intro") is not None else _oin,
            "outro": cfg.get("share_tc_outro") if cfg.get("share_tc_outro") is not None else _oout,
        },
        "app_url": cfg.get("app_url") or "",
        "mail_enabled": bool(cfg.get("enabled")),
        "default_subject": _DEFAULT_SHARE_SUBJECT,
        "default_req_sections": _DEFAULT_REQ_SECTIONS,
        "default_tc_sections": _DEFAULT_TC_SECTIONS,
    }

@app.post("/api/share-config")
async def api_share_config_save(payload: dict, token: str = ""):
    _require_admin(token)
    cfg = _load_mail_cfg()
    req = payload.get("req")
    if isinstance(req, dict):
        if req.get("subject") is not None: cfg["share_req_subject"] = str(req.get("subject"))
        if req.get("intro") is not None: cfg["share_req_intro"] = str(req.get("intro"))
        if req.get("outro") is not None: cfg["share_req_outro"] = str(req.get("outro"))
        if isinstance(req.get("sections"), dict):
            cfg["share_sections_req"] = {str(k): bool(v) for k, v in req["sections"].items()}
    tc = payload.get("tc")
    if isinstance(tc, dict):
        if tc.get("subject") is not None: cfg["share_tc_subject"] = str(tc.get("subject"))
        if tc.get("intro") is not None: cfg["share_tc_intro"] = str(tc.get("intro"))
        if tc.get("outro") is not None: cfg["share_tc_outro"] = str(tc.get("outro"))
        if isinstance(tc.get("sections"), dict):
            cfg["share_sections_tc"] = {str(k): bool(v) for k, v in tc["sections"].items()}
    _save_mail_cfg(cfg)
    return {"ok": True}

@app.post("/api/share-mail")
async def api_share_mail(payload: dict, token: str = ""):
    u = _user_from_token(token)
    if not u:
        raise HTTPException(401, "로그인이 필요합니다")
    to = payload.get("to") or []
    if isinstance(to, str):
        to = [to]
    to = [str(x).strip() for x in to if str(x).strip()]
    if not to:
        raise HTTPException(400, "받는 사람을 입력하세요")
    subject = str(payload.get("subject") or "[ubiQuoss-TOP] 공유")
    html = str(payload.get("html") or "")
    if not html:
        raise HTTPException(400, "공유할 내용이 비어 있습니다")
    if not _load_mail_cfg().get("enabled"):
        raise HTTPException(400, "메일 발송이 꺼져 있습니다 (시스템 → 메일 설정)")
    try:
        sent = _send_mail(to, subject, html, html=True)
    except Exception as e:
        raise HTTPException(400, f"발송 실패: {e}")
    return {"ok": True, "sent": sent}

@app.post("/api/notify/cycle")
async def api_notify_cycle(payload: dict, token: str = ""):
    """사이클 생성 시 담당자에게 배정 알림 메일 발송 (메일 발송 토글 ON일 때 프론트가 호출)."""
    u = _user_from_token(token)
    if not u:
        raise HTTPException(401, "로그인이 필요합니다")
    assignee = str(payload.get("assignee") or "").strip()
    if not assignee:
        raise HTTPException(400, "담당자가 없습니다")
    vg = str(payload.get("version_group") or "")
    ver = str(payload.get("version") or "")
    # 담당자 이메일 조회 (이름 또는 아이디 매칭)
    email = ""
    try:
        for x in _users_load_sync().get("users", []):
            if assignee in (x.get("name"), x.get("username")) and x.get("email"):
                email = x["email"]; break
    except Exception:
        pass
    if not email:
        raise HTTPException(400, f"'{assignee}' 담당자의 이메일을 찾을 수 없습니다")
    _mc = _load_mail_cfg()
    if not _mc.get("enabled"):
        raise HTTPException(400, "메일 발송이 꺼져 있습니다 (시스템 → 메일 설정)")
    model = str(payload.get("model") or "")
    start = str(payload.get("start") or "")
    end = str(payload.get("end") or "")
    items = payload.get("items") or []
    app_url = str(_mc.get("app_url") or "").rstrip("/")
    import html as _h
    esc = lambda s: _h.escape(str(s or ""))
    # 시험 항목 표
    rows = ""
    for i, it in enumerate(items, 1):
        rows += (f'<tr>'
                 f'<td style="padding:6px 10px;border:1px solid #e3e8ef;text-align:center;color:#6b7280;white-space:nowrap;">{i}</td>'
                 f'<td style="padding:6px 10px;border:1px solid #e3e8ef;font-family:monospace;color:#2563eb;white-space:nowrap;">{esc(it.get("id"))}</td>'
                 f'<td style="padding:6px 10px;border:1px solid #e3e8ef;">{esc(it.get("name"))}</td>'
                 f'</tr>')
    if not rows:
        rows = '<tr><td colspan="3" style="padding:10px;border:1px solid #e3e8ef;color:#9ca3af;text-align:center;">배정된 시험 항목 없음</td></tr>'
    period = (f"{esc(start)} ~ {esc(end)}" if (start or end) else "-")
    # 시험 항목 목록 표 ({{items}} 치환용)
    items_html = ('<table style="border-collapse:collapse;width:100%;font-size:12.5px;">'
        '<thead><tr style="background:#f3f6fb;">'
        '<th style="padding:6px 10px;border:1px solid #e3e8ef;width:36px;">#</th>'
        '<th style="padding:6px 10px;border:1px solid #e3e8ef;text-align:left;width:220px;white-space:nowrap;">TC ID</th>'
        '<th style="padding:6px 10px;border:1px solid #e3e8ef;text-align:left;">시험명</th>'
        f'</tr></thead><tbody>{rows}</tbody></table>')
    # 메일 폼: 설정의 사이클 배정 폼(있으면) → 없으면 기본 폼. 플레이스홀더 치환.
    tpl = _mc.get("cycle_html") or _DEFAULT_CYCLE_TPL
    subj_tpl = _mc.get("cycle_subject") or _DEFAULT_CYCLE_SUBJECT
    def _fill(s):
        s = _repair_placeholders(s or "")
        s = s.replace("{{assignee}}", esc(assignee)).replace("{{model}}", esc(model or "-"))
        s = s.replace("{{vgroup}}", esc(vg or "-")).replace("{{version}}", esc(ver or "-"))
        s = s.replace("{{period}}", period).replace("{{count}}", str(len(items)))
        s = s.replace("{{items}}", items_html)
        s = s.replace("{{app_url}}", (app_url or "").strip().rstrip("/"))
        s = s.replace("{{login_button}}", (f'<a href="{esc(app_url)}" style="display:inline-block;margin-top:14px;padding:9px 20px;background:#2563eb;color:#fff;text-decoration:none;border-radius:7px;font-weight:700;">Test Workflow 열기 →</a>' if app_url else ""))
        return s
    subject = _fill(subj_tpl).strip()
    html = _fill(tpl)
    if "<html" not in html.lower():
        html = '<!DOCTYPE html><html><body style="margin:0;padding:0;background:#eef1f6;">' + html + '</body></html>'
    if payload.get("preview"):
        return {"ok": True, "preview": True, "to": email, "subject": subject, "html": html}
    try:
        sent = _send_mail([email], subject, html, html=True)
    except Exception as e:
        raise HTTPException(400, f"발송 실패: {e}")
    return {"ok": True, "sent": sent, "to": email}

@app.post("/api/mail/config")
async def api_mail_config_save(payload: dict, token: str = ""):
    _require_admin(token)
    cfg = _load_mail_cfg()
    for k in ("host", "username", "password", "from_addr", "from_name", "security", "app_url",
              "approval_subject", "approval_html", "cycle_subject", "cycle_html"):
        if payload.get(k) is not None:
            cfg[k] = str(payload.get(k))
    if payload.get("port") is not None:
        try:
            cfg["port"] = int(payload.get("port"))
        except Exception:
            pass
    if payload.get("enabled") is not None:
        cfg["enabled"] = bool(payload.get("enabled"))
    if str(cfg.get("security") or "").lower() not in ("starttls", "ssl", "none"):
        cfg["security"] = "starttls"
    _save_mail_cfg(cfg)
    return {"ok": True, "config": cfg}

@app.post("/api/mail/test")
async def api_mail_test(payload: dict, token: str = ""):
    _require_admin(token)
    to = str(payload.get("to", "")).strip()
    if not to:
        raise HTTPException(400, "받는 사람(테스트 수신 주소)을 입력하세요")
    try:
        sent = _send_mail(to, "[ubiQuoss-TOP] 메일 설정 테스트",
                          "ubiQuoss-TOP 메일(SMTP) 설정 테스트입니다.\n이 메일이 보이면 SMTP 발송이 정상 동작하는 것입니다.")
    except Exception as e:
        raise HTTPException(400, f"발송 실패: {e}")
    return {"ok": True, "sent": sent}

# ───────────────────────────────────────────
# 회원가입(관리자 승인) + @멘션 알림
# ───────────────────────────────────────────
@app.post("/api/signup")
async def api_signup(payload: dict):
    data = _users_load_sync()
    uname = str(payload.get("username", "")).strip()
    name = str(payload.get("name", "")).strip() or uname
    email = str(payload.get("email", "")).strip()
    company = str(payload.get("company", "")).strip()   # 회사
    dept = str(payload.get("dept", "")).strip()   # 소속담당 (예: 품질보증담당)
    team = str(payload.get("team", "")).strip()   # 소속팀 (예: QA팀)
    position = str(payload.get("position", "")).strip()   # 직책
    duty = str(payload.get("duty", "")).strip()   # 보직
    pw = payload.get("password") or ""
    if not uname:
        raise HTTPException(400, "아이디를 입력하세요")
    if not pw:
        raise HTTPException(400, "비밀번호를 입력하세요")
    if not company or not dept or not team or not position:
        raise HTTPException(400, "회사·소속담당·소속팀·직책을 모두 선택하세요")
    if not email or not _valid_email(email):
        raise HTTPException(400, "올바른 이메일을 입력하세요")
    if not _allowed_email_domain(email):
        raise HTTPException(400, "@" + ALLOWED_EMAIL_DOMAIN + " 이메일만 가입할 수 있습니다")
    if any(u.get("username") == uname for u in data["users"]):
        raise HTTPException(400, "이미 존재하는 아이디입니다")
    salt = _secrets.token_hex(8)
    nu = {"id": uname, "username": uname, "name": name, "role": "팀원", "email": email,
          "company": company, "position": position, "duty": duty, "dept": dept, "team": team,
          "salt": salt, "password": _hash_pw(pw, salt), "active": False, "pending": True,
          "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
    data["users"].append(nu)
    _users_save_sync(data)
    try:
        admins = [u.get("email") for u in data["users"] if u.get("role") == "관리자" and u.get("email")]
        if admins and _load_mail_cfg().get("enabled"):
            _send_mail(admins, "[ubiQuoss-TOP] 신규 가입 승인 요청",
                       f"신규 가입 신청이 있습니다.\n\n아이디: {uname}\n회사: {company}\n소속담당: {dept}\n소속팀: {team}\n직책: {position}\n이름: {name}\n이메일: {email}\n\n[시스템 → 사용자 관리]에서 승인해주세요.")
    except Exception:
        pass
    return {"ok": True}

# ── 조직 설정: 회사 ▸ 소속담당 ▸ 소속팀 (계층) + 직책·보직(평면) ──
ORG_FILE = DATA_DIR / "config" / "org_options.json"
_ORG_DEFAULT = {
    "companies": [
        {"name": "유비쿼스", "depts": [
            {"name": "품질보증담당", "teams": ["QA팀", "검증팀"]},
            {"name": "개발담당", "teams": ["SW개발팀", "HW개발팀", "시스템개발팀"]},
            {"name": "기술지원담당", "teams": ["기술지원팀"]},
        ]},
        {"name": "유비쿼스솔루션", "depts": [
            {"name": "영업담당", "teams": ["영업1팀", "영업2팀"]},
            {"name": "경영지원담당", "teams": ["경영지원팀"]},
        ]},
    ],
    "position": ["사원", "주임", "대리", "과장", "차장", "부장", "수석", "책임", "선임", "이사"],
    "duty": ["팀원", "파트장", "팀장", "그룹장", "본부장", "PM", "PL", "해당없음"],
}

def _org_clean_list(v):
    out = []
    if isinstance(v, list):
        for x in v:
            s = str(x).strip()
            if s and s not in out:
                out.append(s)
    return out

def _org_clean_companies(v):
    out = []; seen = set()
    if isinstance(v, list):
        for c in v:
            if not isinstance(c, dict):
                continue
            nm = str(c.get("name", "")).strip()
            if not nm or nm in seen:
                continue
            seen.add(nm)
            depts = []; dseen = set()
            for d in (c.get("depts") or []):
                if not isinstance(d, dict):
                    continue
                dn = str(d.get("name", "")).strip()
                if not dn or dn in dseen:
                    continue
                dseen.add(dn)
                depts.append({"name": dn, "teams": _org_clean_list(d.get("teams"))})
            out.append({"name": nm, "depts": depts})
    return out

def _load_org() -> dict:
    import copy as _copy
    if ORG_FILE.exists():
        try:
            d = json.loads(ORG_FILE.read_text(encoding="utf-8"))
            if isinstance(d, dict) and isinstance(d.get("companies"), list):
                return {"companies": _org_clean_companies(d.get("companies")),
                        "position": _org_clean_list(d.get("position")) or list(_ORG_DEFAULT["position"]),
                        "duty": _org_clean_list(d.get("duty")) or list(_ORG_DEFAULT["duty"])}
            # 구(舊) 평면 구조 → 계층 마이그레이션(각 회사에 모든 담당, 각 담당에 모든 팀)
            if isinstance(d, dict) and (d.get("company") or d.get("dept") or d.get("team")):
                comps = _org_clean_list(d.get("company")) or [_ORG_DEFAULT["companies"][0]["name"]]
                depts = _org_clean_list(d.get("dept")); teams = _org_clean_list(d.get("team"))
                tree = [{"name": c, "depts": [{"name": dn, "teams": list(teams)} for dn in depts]} for c in comps]
                return {"companies": tree,
                        "position": _org_clean_list(d.get("position")) or list(_ORG_DEFAULT["position"]),
                        "duty": _org_clean_list(d.get("duty")) or list(_ORG_DEFAULT["duty"])}
        except Exception:
            pass
    return _copy.deepcopy(_ORG_DEFAULT)

@app.get("/api/org-options")
async def org_options_get():
    org = _load_org()
    return {"ok": True, "companies": org["companies"], "position": org["position"],
            "duty": org["duty"], "company": [c["name"] for c in org["companies"]]}

@app.post("/api/org-options")
async def org_options_save(payload: dict, token: str = ""):
    _require_admin(token)
    cur = _load_org()
    if isinstance(payload.get("companies"), list):
        cur["companies"] = _org_clean_companies(payload["companies"])
    if isinstance(payload.get("position"), list):
        cur["position"] = _org_clean_list(payload["position"])
    if isinstance(payload.get("duty"), list):
        cur["duty"] = _org_clean_list(payload["duty"])
    ORG_FILE.write_text(json.dumps(cur, ensure_ascii=False), encoding="utf-8")
    return {"ok": True, "companies": cur["companies"], "position": cur["position"],
            "duty": cur["duty"], "company": [c["name"] for c in cur["companies"]]}

# ── 시험 절차 학습: 검증된 절차 스냅샷 저장소 (LLM few-shot + 회귀 baseline) ──
LEARNED_FILE = DATA_DIR / "state" / "learned_procedures.json"
def _load_learned() -> dict:
    d = _kv_load_sync("learned_procedures", {"items": []})
    if isinstance(d, dict) and isinstance(d.get("items"), list):
        return d
    return {"items": []}
def _save_learned(d):
    _kv_save_sync("learned_procedures", d)

@app.post("/api/learn/procedure")
async def learn_procedure_save(payload: dict, token: str = ""):
    import time as _tt
    store = _load_learned()
    item = {
        "id": "lp-" + str(int(_tt.time() * 1000)),
        "tcid": str(payload.get("tcid") or ""),
        "title": str(payload.get("title") or ""),
        "models": payload.get("models") if isinstance(payload.get("models"), list) else [],
        "role": str(payload.get("role") or ""),
        "vendor": str(payload.get("vendor") or ""),
        "steps": payload.get("steps") if isinstance(payload.get("steps"), list) else [],
        "outputs": payload.get("outputs") if isinstance(payload.get("outputs"), list) else [],
        "by": str(payload.get("by") or ""),
        "at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }
    if item["tcid"]:   # 같은 TC+모델그룹 재학습 → 덮어쓰기(모델그룹별 항목은 각각 보존)
        _msig = ",".join(sorted(item.get("models") or []))
        store["items"] = [it for it in store["items"]
                          if not (it.get("tcid") == item["tcid"] and ",".join(sorted(it.get("models") or [])) == _msig)]
    store["items"].insert(0, item)
    _save_learned(store)
    return {"ok": True, "id": item["id"], "count": len(store["items"])}

@app.get("/api/learn/procedures")
async def learn_procedures_list(model: str = "", role: str = "", q: str = "", limit: int = 50):
    items = _load_learned().get("items", [])
    def _match(it):
        if model and model not in (it.get("models") or []):
            return False
        if role and it.get("role") != role:
            return False
        if q:
            hay = (str(it.get("title", "")) + " " + " ".join(it.get("models") or []) + " "
                   + json.dumps(it.get("steps") or [], ensure_ascii=False)).lower()
            if q.lower() not in hay:
                return False
        return True
    out = [it for it in items if _match(it)][:max(1, min(limit, 500))]
    return {"ok": True, "items": out, "total": len(items)}

@app.delete("/api/learn/procedure/{lp_id}")
async def learn_procedure_delete(lp_id: str, token: str = ""):
    store = _load_learned()
    n0 = len(store["items"])
    store["items"] = [it for it in store["items"] if it.get("id") != lp_id]
    _save_learned(store)
    return {"ok": True, "removed": n0 - len(store["items"])}

#
# 용도별 프롬프트.
#
# 프롬프트를 코드에 박아 두면 한 글자 고치는 데도 배포를 해야 한다. 랩에
# 쓰는 말은 현장에서 자꾸 바뀌고(장비 계열이 늘거나 부르는 이름이 다르거나),
# 그때마다 사람이 기다려야 한다. 설정에 두고 화면에서 고친다.
#
# 여기 적은 것은 **기본값**이다. 설정에 없으면 이것을 쓴다 — 처음 쓰는
# 사람이 빈 화면을 보지 않도록.
#
LLM_PURPOSES: dict[str, dict] = {
    # ── 요구사항 ────────────────────────────────────────────────
    "req_intent": {
        "label": "REQ-Intent",
        "hint": "요구사항 › 구현내용(Intent) — 한 줄 요청을 구현의도로 다듬습니다.",
        "system": (
            "당신은 네트워크 장비 시험 조직의 요구사항 작성자다. 사람이 적은 "
            "짧은 요청을 **구현의도**로 다듬는다.\n"
            "규칙:\n"
            "1) 제목은 한 줄, 측정 가능한 동작으로 적는다.\n"
            "2) 설명에는 대상 장비·조건·기대 동작을 적는다.\n"
            "3) 검증 기준은 시험으로 확인할 수 있는 문장으로 적는다 — "
            "'빨라야 한다' 가 아니라 '몇 초 안에' 로.\n"
            "4) 지어내지 않는다. 사람이 말하지 않은 수치는 (확인 필요) 로 남긴다."
        ),
    },
    "req_coverage": {
        "label": "REQ-Coverage",
        "hint": "요구사항 › 구현내용을 근거로 **시험 항목 목록**을 뽑습니다.",
        "system": (
            "당신은 네트워크 장비 시험 설계자다. 주어진 요구사항의 구현의도를 읽고 "
            "그것을 덮는 **시험 항목 목록**을 만든다.\n"
            "규칙:\n"
            "1) 항목 하나는 한 가지만 확인한다. 여러 개를 한 항목에 묶지 마라.\n"
            "2) 이름은 '무엇을 어떻게 확인하는가' 로 적는다 — 명령 이름을 그대로 쓰지 마라.\n"
            "3) 정상 동작뿐 아니라 경계·실패 조건도 빠뜨리지 않는다.\n"
            "4) 구현의도에 없는 기능은 만들지 않는다.\n"
            "5) 스텝은 여기서 만들지 않는다 — 항목 이름과 목적까지다."
        ),
    },
    # ── 시험 항목 ───────────────────────────────────────────────
    "coverage_object": {
        "label": "Coverage-Object",
        "hint": "시험항목 › 시험 목적(Object)과 사전 준비 조건을 씁니다.",
        "system": (
            "당신은 네트워크 장비 시험 항목의 목적과 사전 조건을 쓴다.\n"
            "규칙:\n"
            "1) 목적은 '무엇을 확인하는 시험인가' 한두 문장으로.\n"
            "2) 사전 조건은 장비 상태·배선·설정을 줄로 나눠 적는다.\n"
            "3) 주어진 시험 이름·스텝·배선에 있는 사실만 쓴다. 지어내지 않는다.\n"
            "4) 결과서에 그대로 실리는 글이다 — 존댓말 없이 개조식으로."
        ),
    },
    # ── 토폴로지 — 목적을 적었으면 그다음이 배선이다(지시: 3번 뒤) ──
    "wiring": {
        "label": "Topology-Wiring",
        "hint": "토폴로지 › 말로 적은 랩 배선을 「장비 포트 ↔ 계측기 포트」 줄로 옮깁니다.",
        "system": (
            "당신은 네트워크 시험 랩의 배선을 정리한다. 사람이 말한 연결을 "
            "'장비 포트 ↔ 계측기 포트' 줄로 옮긴다.\n"
            "규칙:\n"
            "1) dev·meter 는 반드시 주어진 목록의 id 를 그대로 쓴다.\n"
            "2) port·meterPort 는 반드시 그 장비/계측기의 ports 목록에 있는 값을 그대로 쓴다.\n"
            "3) 목록에 없으면 그 줄은 만들지 않는다. 비슷한 이름을 지어내지 마라.\n"
            "4) 한 포트는 한 번만 쓴다.\n"
            "5) JSON 만 출력한다."
        ),
    },
    "coverage_manual": {
        "label": "Coverage-Manual",
        "hint": "시험항목 › Manual 절차 — 사람이 손으로 하는 순서를 씁니다.",
        "system": (
            "당신은 네트워크 장비 시험의 **수동 절차**를 쓴다. 사람이 손으로 하는 "
            "순서라 장비 명령이 아니라 **행동**으로 적는다.\n"
            "규칙:\n"
            "1) 한 줄에 한 가지 행동. '무엇을 한다 → 무엇을 본다' 차례로.\n"
            "2) 눈으로 확인할 수 있는 것만 적는다(LED·화면·소리·측정기 값).\n"
            "3) **판정 기준은 비워 둔다** — 돌려 본 뒤 사람이 정한다.\n"
            "4) 되돌리는 절차가 필요하면 마지막에 적는다."
        ),
    },
    "coverage_automation": {
        "label": "Coverage-Automation",
        "hint": "시험항목 › Automation 스텝 · AI 화면의 고급 갈래(절차 짓기)가 함께 씁니다.",
        "system": "",  # 비우면 코드가 든 긴 규칙(nl_test.py)을 그대로 쓴다
    },
    # ── 사이클 ──────────────────────────────────────────────────
    "cycle_summary": {
        "label": "Cycle-Test Summary",
        "hint": "사이클 실행 › 시험 진행 요약의 AI 요약을 씁니다.",
        "system": (
            "당신은 네트워크 장비 시험(QA) 결과 분석 전문가다. 주어진 회차 결과를 "
            "근거로 한국어 Markdown 보고서를 쓴다.\n"
            "규칙:\n"
            "1) 총평 한 문단 — 이 버전을 내보내도 되는가에 답한다.\n"
            "2) 전체·수동·자동 현황을 표로.\n"
            "3) 깨진 항목은 무엇이 왜 깨졌는지 묶어서 적는다. 스텝 번호를 밝힌다.\n"
            "4) 결과에 없는 것은 쓰지 않는다 — 미실행은 미실행이라고 적는다."
        ),
    },
    # ── 화면에 안 세우는 것 ─────────────────────────────────────
    # 「일곱 자리」 는 사람이 손보는 자리다(지시). 이것은 고를 것이 없는
    # 붙박이라 목록에서 감춘다 — 지우면 AI 「일반」 갈래가 시험을 못 고른다.
    "similar": {
        "hidden": True,
        "label": "닮은 시험 찾기",
        "hint": "AI 일반 갈래 — 말과 가장 가까운 기존 시험을 고릅니다.",
        "system": (
            "당신은 네트워크 시험 담당자다. 사람이 하려는 시험과 가장 가까운 것을 "
            "주어진 목록에서 고른다.\n"
            "규칙:\n"
            "1) 목록에 있는 tcid 만 쓴다. 새로 만들지 마라.\n"
            "2) 가까운 것부터 최대 3개.\n"
            "3) 가까운 것이 없으면 빈 배열을 준다. 억지로 채우지 마라.\n"
            "4) JSON 만 출력한다."
        ),
    },
}

# 옛 이름 → 새 이름. 저장해 둔 프롬프트가 이름이 바뀌었다고 사라지면 안 된다.
LLM_PURPOSE_ALIAS = {
    "requirement": "req_intent",
    "objective": "coverage_object",
    "steps": "coverage_automation",
}


def _prompt_of(purpose: str) -> dict:
    """이 용도에 쓸 프롬프트와 LLM. 설정에 있으면 그것, 없으면 기본값."""
    base = LLM_PURPOSES.get(purpose) or {}
    saved = {}
    if PROMPTS_FILE.exists():
        try:
            _all = load_json(PROMPTS_FILE).get("purposes") or {}
            saved = _all.get(purpose) or {}
            if not saved:
                # 이름을 바꾸기 전에 저장해 둔 것 — 그대로 이어 쓴다
                for _old, _new in LLM_PURPOSE_ALIAS.items():
                    if _new == purpose and _all.get(_old):
                        saved = _all[_old]
                        break
        except Exception:
            saved = {}
    return {
        "label": saved.get("label") or base.get("label") or purpose,
        "hint": base.get("hint") or "",
        "system": (saved.get("system") or "").strip() or base.get("system") or "",
        "llm": saved.get("llm") or "",
        # 채팅 화면 쪽 — 여는 말·입력칸 안내·추천 질문(지시)
        "greeting": saved.get("greeting") or base.get("greeting") or "",
        "placeholder": saved.get("placeholder") or base.get("placeholder") or "",
        "asks": list(saved.get("asks") or base.get("asks") or []),
    }


def _llm_pick(purpose: str = "", llm_id: str = ""):
    """
    쓸 수 있는 로컬 LLM 하나.

    고르는 규칙이 `/api/llm/generate` 안에 박혀 있어서, 다른 곳에서 LLM 을
    쓰려면 그 여든 줄을 통째로 베껴야 했다. 한 곳으로 뺀다.

    용도에 붙여 둔 LLM 이 있으면 그것을 먼저 쓴다 — 배선처럼 짧은 일에는
    작은 모델을, 절차 만들기에는 큰 모델을 붙일 수 있어야 한다.
    """
    init_llms_file()
    llms = load_json(LLMS_FILE).get("llms") or []

    def _ok(l):
        if not (l.get("status", "active") == "active" and l.get("endpoint")):
            return False
        t = str(l.get("type") or "").lower()
        return t in ("local", "vllm", "openai", "openai-compatible", "")

    # 화면에서 고른 것이 먼저다 — 설정의 기본값은 안 고른 사람을 위한 것
    want = llm_id or (_prompt_of(purpose).get("llm") if purpose else "")
    if want:
        for l in llms:
            if str(l.get("id") or "") == want and _ok(l):
                return l
    for l in llms:
        if _ok(l):
            return l
    return None


@app.get("/api/llm/purposes")
async def llm_purposes():
    """용도 목록 — 화면이 이것으로 설정 칸을 그린다. 기본 프롬프트도 함께 준다."""
    out = []
    for k, v in LLM_PURPOSES.items():
        if v.get("hidden"):
            continue          # 사람이 손볼 자리만 세운다(지시: 일곱 자리)
        cur = _prompt_of(k)
        out.append({
            "id": k,
            "label": cur["label"],
            "hint": v.get("hint") or "",
            "system": cur["system"],
            "llm": cur["llm"],
            "default": v.get("system") or "",
            # 채팅 화면에서 보이는 것 — 여는 말·입력칸 안내·추천 질문(지시)
            "greeting": cur.get("greeting") or "",
            "placeholder": cur.get("placeholder") or "",
            "asks": cur.get("asks") or [],
        })
    return {"purposes": out}


@app.post("/api/llm/purposes")
async def llm_purposes_save(data: dict):
    """용도별 프롬프트·LLM 저장. 다른 설정(prompts.json)은 건드리지 않는다."""
    cur = load_json(PROMPTS_FILE) if PROMPTS_FILE.exists() else {}
    ps = dict(cur.get("purposes") or {})
    for k, v in (data.get("purposes") or {}).items():
        if k not in LLM_PURPOSES:
            continue
        ps[k] = {
            "system": str((v or {}).get("system") or ""),
            "llm": str((v or {}).get("llm") or ""),
            # 채팅 화면 쪽 값 — 없으면 빈 것으로 둔다(옛 저장본 호환)
            "greeting": str((v or {}).get("greeting") or ""),
            "placeholder": str((v or {}).get("placeholder") or ""),
            "asks": [str(x) for x in ((v or {}).get("asks") or []) if str(x).strip()],
        }
    cur["purposes"] = ps
    PROMPTS_FILE.parent.mkdir(parents=True, exist_ok=True)
    save_json(PROMPTS_FILE, cur)
    return {"ok": True}


async def _llm_json(llm, sys_p, user_p, schema, timeout=120):
    """LLM 에게 JSON 하나를 받는다. `guided_json` 이 없는 판이면 한 번 더 물러선다."""
    import httpx
    body = {
        "model": llm.get("model") or "",
        "messages": [{"role": "system", "content": sys_p}, {"role": "user", "content": user_p}],
        "temperature": 0.1,
        "max_tokens": 1536,
        "guided_json": schema,
    }
    headers = {"Content-Type": "application/json"}
    if llm.get("apikey"):
        headers["Authorization"] = f"Bearer {llm['apikey']}"
    url = str(llm["endpoint"]).rstrip("/") + "/chat/completions"
    async with httpx.AsyncClient(timeout=timeout) as client:
        r = await client.post(url, headers=headers, json=body)
        if r.status_code != 200:
            body.pop("guided_json", None)
            body["response_format"] = {"type": "json_object"}
            r = await client.post(url, headers=headers, json=body)
        if r.status_code != 200:
            raise RuntimeError(f"LLM {r.status_code}: {r.text[:300]}")
        data = r.json()
        txt = (((data.get("choices") or [{}])[0].get("message") or {}).get("content") or "").strip()
    m = re.search(r"\{.*\}", txt, re.S)
    return json.loads(m.group(0) if m else txt)


@app.post("/api/llm/similar")
async def llm_similar(payload: dict):
    """
    목적 한 줄과 **닮은 시험**을 찾는다.

    이 시스템에는 이미 검증된 시험이 쌓여 있다. 어떤 것은 여덟 번씩
    돌았고 판정 기준도 그만큼 다듬어졌다. 그런데 새 시험을 만들 때 그것을
    쓰지 않고 매번 빈 화면에서 시작한다 — 가장 좋은 자산을 놀리고 있다.

    그래서 **짓지 않고 찾는다.** AI 는 「성능」 이 「트래픽 2포트 시험」 인
    것을 알아보는 데만 쓴다. 절차·판정은 이미 있는 것을 그대로 옮긴다 —
    지어낼 자리가 없으니 틀릴 자리도 없다.

    LLM 이 없거나 답을 못 줘도 글자 맞춤으로 찾아 준다. 찾는 일이 아예
    안 되는 것보다는 덜 똑똑해도 되는 편이 낫다.
    """
    want = str(payload.get("purpose") or "").strip()
    models = [str(x) for x in (payload.get("models") or []) if str(x).strip()]
    if not want:
        return {"ok": False, "error": "무엇을 시험하려는지 한 줄 적어 주세요"}

    try:
        metas = await db.tc_list_meta()
    except Exception:
        metas = []
    if not metas:
        return {"ok": True, "items": []}

    # 글자 맞춤 — 이름·요구사항에 든 낱말이 몇 개나 겹치나
    words = [w for w in re.split(r"[\s·,/()]+", want) if len(w) > 1]

    def _score(t):
        name = f"{t.get('name') or ''} {t.get('req_id') or ''} {t.get('type') or ''}"
        low = name.lower()
        s = sum(2 for w in words if w.lower() in low)
        # 같은 계열 장비로 돌린 적이 있으면 크게 친다 — 그 랩에서 실제로 된 것이다
        for m in models:
            if m and m.lower() in low:
                s += 3
        # 여러 번 돌아간 것일수록 믿을 만하다
        s += min(3, int(t.get("run_count") or 0))
        if str(t.get("status") or "").upper() == "PASS":
            s += 1
        return s

    #
    # 고를 거리를 넉넉히 준다.
    #
    # 처음에는 글자 맞춤 상위 여덟 개만 LLM 에게 보였다. 그랬더니 「E4320
    # 성능」 을 물었을 때 이름에 E4320 이 든 시험만 올라오고, 정작 맞는
    # 「N2X 트래픽 2포트 시험」 은 후보에도 못 들었다 — 그 이름에는 E4320 이
    # 없기 때문이다. 뜻으로 고르라고 시켜 놓고 글자로 미리 걸러 버린 셈이다.
    ranked = sorted(metas, key=_score, reverse=True)
    top = ranked[:40]

    # LLM 이 있으면 그중에서 고르게 한다 — 낱말이 안 겹쳐도 뜻이 닿는 것이 있다
    llm = _llm_pick("similar")
    picked = []
    if llm and len(top) > 1:
        brief = [
            {"tcid": t.get("tcid"), "name": t.get("name"), "type": t.get("type"), "req": t.get("req_id")}
            for t in top
        ]
        schema = {
            "type": "object",
            "properties": {"tcids": {"type": "array", "items": {"type": "string"}}},
            "required": ["tcids"],
        }
        sys_p = _prompt_of("similar")["system"]
        user_p = (
            "시험 목록:\n" + json.dumps(brief, ensure_ascii=False) +
            "\n\n사람이 하려는 것:\n" + want +
            "\n\n가장 가까운 것부터 최대 3개의 tcid 만 {\"tcids\":[...]} 로 출력하라."
        )
        try:
            got = await _llm_json(llm, sys_p, user_p, schema, timeout=60)
            ids = [str(x) for x in (got.get("tcids") or [])]
            byid = {str(t.get("tcid")): t for t in top}
            picked = [byid[i] for i in ids if i in byid]
        except Exception:
            picked = []

    order = picked + [t for t in top if t not in picked]
    return {
        "ok": True,
        "items": [
            {
                "tcid": t.get("tcid"),
                "name": t.get("name"),
                "type": t.get("type"),
                "req_id": t.get("req_id"),
                "runs": int(t.get("run_count") or 0),
                "status": t.get("status") or "",
                "why": "AI 가 고름" if t in picked else "이름이 닮음",
            }
            for t in order[:5]
        ],
    }


@app.post("/api/llm/wiring")
async def llm_wiring(payload: dict):
    """
    말로 적은 배선을 줄로 옮긴다.

    「E5724RL 1번 2번 포트를 N2X 4106/3, 4106/4 에 물렸어」 같은 문장을
    받아 배선 줄을 만든다.

    **지어낸 이름은 버린다.** 장비·포트 목록은 이미 자료로 있으므로, 그
    안에 없는 것은 서버에서 걸러 내고 무엇을 버렸는지 함께 알린다.
    로컬 모델은 그럴듯한 포트 이름을 곧잘 지어내는데, 그것이 그대로
    저장되면 실행할 때까지 아무도 모른다 — 조용히 틀리는 것이 제일 나쁘다.

    저장은 하지 않는다. 화면이 그림으로 보여 주고 사람이 정한다.
    """
    say = str(payload.get("text") or "").strip()
    if not say:
        return {"ok": False, "error": "무엇을 어떻게 물렸는지 적어 주세요"}
    devs = payload.get("devices") or []      # [{id,label,ports:[...]}]
    meters = payload.get("meters") or []     # [{id,label,ports:[...]}]
    if not devs or not meters:
        return {"ok": False, "error": "장비와 계측기가 있어야 배선을 그립니다"}

    llm = _llm_pick("wiring", str(payload.get("llm") or ""))
    if not llm:
        return {"ok": False, "error": "등록된 로컬 LLM 이 없습니다 — 설정 › LLM 설정에서 켜세요"}
    sys_p = _prompt_of("wiring")["system"]

    def _one(x):
        return {
            "id": str(x.get("id") or ""),
            "label": str(x.get("label") or x.get("id") or ""),
            "ports": [str(p) for p in (x.get("ports") or [])][:200],
        }

    D = [_one(x) for x in devs]
    M = [_one(x) for x in meters]
    # 계측기 포트를 안 불러왔으면 여기서 멈춘다. 빈 목록으로 LLM 에 보내면
    # 「목록에 없으면 만들지 마라」 규칙 때문에 빈손으로 돌아오는데, 화면은
    # 그걸 「못 알아들었다」 로 보여 줬다 — 이유를 말해야 사람이 고친다.
    if not any(m["ports"] for m in M):
        return {
            "ok": False,
            "error": "계측기 포트 목록이 비어 있습니다 — 결선 줄에서 계측기를 고르고 「불러오기」 를 먼저 누르세요.",
        }
    schema = {
        "type": "object",
        "properties": {
            "wires": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "dev": {"type": "string"},
                        "port": {"type": "string"},
                        "meter": {"type": "string"},
                        "meterPort": {"type": "string"},
                    },
                    "required": ["dev", "port", "meter", "meterPort"],
                },
            }
        },
        "required": ["wires"],
    }
    user_p = (
        "장비:\n" + json.dumps(D, ensure_ascii=False) +
        "\n계측기:\n" + json.dumps(M, ensure_ascii=False) +
        "\n\n사람이 말한 배선:\n" + say +
        "\n\n{\"wires\":[...]} 로만 출력하라."
    )
    try:
        got = await _llm_json(llm, sys_p, user_p, schema)
    except Exception as e:
        return {"ok": False, "error": str(e)[:300]}

    dmap = {x["id"]: set(x["ports"]) for x in D}
    mmap = {x["id"]: set(x["ports"]) for x in M}
    # 이름으로 부른 것도 받아 준다 — 사람도 모델도 id 보다 이름을 쓴다
    dbyname = {x["label"]: x["id"] for x in D}
    mbyname = {x["label"]: x["id"] for x in M}
    out, dropped = [], []
    seen = set()
    for w in (got.get("wires") or []):
        dv = str(w.get("dev") or "")
        mt = str(w.get("meter") or "")
        dv = dv if dv in dmap else dbyname.get(dv, dv)
        mt = mt if mt in mmap else mbyname.get(mt, mt)
        pt = str(w.get("port") or "")
        mp = str(w.get("meterPort") or "")
        why = ""
        if dv not in dmap:
            why = f"{w.get('dev')} 라는 장비가 없습니다"
        elif pt not in dmap[dv]:
            why = f"{dv} 에 {pt} 포트가 없습니다"
        elif mt not in mmap:
            why = f"{w.get('meter')} 라는 계측기가 없습니다"
        elif mp not in mmap[mt]:
            why = f"{mt} 에 {mp} 포트가 없습니다"
        elif (dv, pt) in seen or (mt, mp) in seen:
            why = "이미 쓴 포트입니다"
        if why:
            dropped.append(why)
            continue
        seen.add((dv, pt))
        seen.add((mt, mp))
        out.append({"dev": dv, "port": pt, "meter": mt, "meterPort": mp})
    if not out and not dropped:
        # LLM 이 규칙대로 「목록에 없으면 안 만든다」 를 지켜 빈손으로 온 것 —
        # 대개 포트 이름이 등록 목록과 다르다. 무엇이 있는지 알려 준다.
        avail = " · ".join(
            f"{x['label']}: {', '.join(x['ports'][:6])}{'…' if len(x['ports']) > 6 else ''}"
            for x in (D + M)[:4] if x["ports"]
        )
        return {
            "ok": False,
            "error": f"문장의 포트 이름이 등록 목록에 없는 것 같습니다. 등록된 포트 — {avail}",
        }
    return {"ok": True, "wires": out, "dropped": dropped}


@app.post("/api/llm/generate")
async def llm_generate(payload: dict, token: str = ""):
    """자연어 시험 목적 → 시험 절차(steps) 생성. 등록 LLM(vLLM) + 학습 예시 few-shot + JSON 강제."""
    import httpx, re as _re
    purpose = str(payload.get("purpose") or "").strip()
    dev_model = str(payload.get("model") or "").strip()
    role = str(payload.get("role") or "").strip()
    if not purpose:
        return {"ok": False, "error": "시험 목적을 입력하세요"}
    # 등록 LLM 선택 — vLLM/OpenAI 호환 로컬 LLM만 지원(/chat/completions 규격). Claude 는 /v1/messages 라 여기선 제외.
    init_llms_file()
    llms = (load_json(LLMS_FILE).get("llms") or [])
    def _ok(l):
        if not (l.get("status", "active") == "active" and l.get("endpoint")): return False
        t = str(l.get("type") or "").lower()
        # openai 호환 계열만 통과 (local/vllm/openai). claude/anthropic/gemini/bedrock 등 비호환은 제외.
        return t in ("local", "vllm", "openai", "openai-compatible", "") and t not in ("claude", "anthropic")
    llm = next((l for l in llms if _ok(l)), None)
    if not llm:
        return {"ok": False, "error": "등록된 로컬(OpenAI 호환) LLM이 없습니다. AI Assistant에서 vLLM/제마를 활성화하세요."}
    # few-shot: 학습 예시 중 모델/제품군 유사 상위 3건
    items = _load_learned().get("items", [])
    def _rel(it):
        s = 0
        if dev_model and dev_model in (it.get("models") or []): s += 2
        if role and it.get("role") == role: s += 1
        return s
    examples = sorted(items, key=_rel, reverse=True)[:3]
    schema = {"type": "object", "properties": {"steps": {"type": "array", "items": {
        "type": "object",
        "properties": {"desc": {"type": "string"}, "cli": {"type": "string"}, "type": {"type": "string", "enum": ["contains", "contains_all", "notcontains", "line"]}, "criteria": {"type": "string"}},
        "required": ["desc", "cli", "type"]}}}, "required": ["steps"]}
    sys_p = ("당신은 네트워크 장비 시험 절차 설계 전문가다. 사용자의 시험 목적에 맞는 CLI 시험 절차를 만든다. "
             "각 스텝은 desc(이 스텝이 무엇을 확인/시험하는지 한국어 한 줄 설명), cli(실행 명령), "
             "type(contains=출력에 포함되어야 정상 / notcontains=없어야 정상 / line=특정 라인 확인), "
             "criteria(판정에 쓸 문자열)로 구성한다. desc는 반드시 채운다. JSON만 출력하고 그 외 설명/주석은 쓰지 않는다.\n"
             "\n"
             "[Ubiquoss 장비 CLI 정정 규칙 — 반드시 준수]\n"
             "1) QoS Class/Queue: Cisco식 표기 금지. 다음 명령을 그대로 사용하라.\n"
             "   - 조회: 'show class-map' (X: 'show qos class-map'), 'show policy-map' (X: 'show qos policy-map <name>')\n"
             "   - 매치: 'match ip-dscp <0-63>' (X: 'match ip dscp ef'). 예) EF=46, AF11=10\n"
             "   - 큐 지정: 'set queueing <0-7>' (X: 'set queue <N>')\n"
             "   - 판정 문자열 예: 'Set Queueing : 1' (대문자 Q + 콜론), 'Match IP DSCP: 46'\n"
             "   - 원복: 'no service-policy input <NAME>' → 'no policy-map <NAME>' → 'no class-map <NAME>' 순\n"
             "2) IGMP Snooping — CLI 모드 주의:\n"
             "   - 'ip igmp snooping' 은 반드시 'interface vlan 1' 진입 후 설정. global config 에서 하면 '% Incomplete command'.\n"
             "   - 'show ip igmp snooping' 출력은 Global 섹션 + Vlan N 섹션으로 나뉜다. 판정 대상은 'Vlan 1' 섹션.\n"
             "3) IGMP proxy 는 DUT 설정이 아님:\n"
             "   - 'ip igmp proxy-service priority 200' 등은 상위 연동 OLT 설정이다. L2 DUT 스텝에 넣지 말 것.\n"
             "4) VLAN 표기(Ubiquoss):\n"
             "   - VLAN 생성: 'vlan database' 진입 후 'vlan N'. 'configure terminal' 로 바로 생성 금지.\n"
             "   - 포트: 'interface range GigabitEthernet 0/1-8' 또는 'interface GigabitEthernet 0/x' 사용.\n"
             "   - Cisco 표기(예: 'GigabitEthernet 1/0/1', 'gigabitethernet 1/0/x') 금지. Ubiquoss 엔 '1/0/' 계층 없음.\n"
             "   - 포트 모드: 'switchport mode access' / 'switchport mode trunk' / 'switchport mode hybrid'\n"
             "\n"
             "[판정기준 작성 규칙 — 거짓 합격 방지]\n"
             "- 판정에 쓸 값이 출력에 유일하면 값만 쓴다 (예: 'VLAN0010', '1.0.1', 'E5724RL').\n"
             "- 같은 값이 여러 줄에 등장하면 라벨 토큰 + 값을 함께 쓴다. 단 콜론·정렬 공백(2칸 이상)은 넣지 마라(기종별로 폭이 달라 깨진다).\n"
             "- 여러 줄 criteria 는 반드시 type='contains_all' 을 쓴다. 'contains' 는 여러 줄 검색 시 무조건 Fail.\n"
             "- desc 가 가리키는 항목과 criteria 는 반드시 같은 항목이어야 한다 (Main Memory 스텝에 Flash 값 금지).\n"
             "- 근거 없는 CLI/값은 창작하지 말고 criteria 를 '[확인필요]' 로 둔다.")
    ex_text = ""
    for ex in examples:
        ex_text += f"\n[예시] 목적: {ex.get('title','')} (모델 {','.join(ex.get('models') or [])})\n" + json.dumps({"steps": ex.get("steps") or []}, ensure_ascii=False)
    user_p = f"대상 모델: {dev_model or '공통'}\n시험 목적: {purpose}\n"
    if ex_text:
        user_p += "\n아래는 검증된 정상 절차 예시다. 명령 체계/스타일을 참고하라:" + ex_text
    user_p += "\n\n위 목적에 맞는 시험 절차를 {\"steps\":[...]} JSON으로만 출력하라."
    body = {"model": llm.get("model") or "", "messages": [{"role": "system", "content": sys_p}, {"role": "user", "content": user_p}],
            "temperature": 0.2, "max_tokens": 2048, "guided_json": schema}
    headers = {"Content-Type": "application/json"}
    if llm.get("apikey"):
        headers["Authorization"] = f"Bearer {llm['apikey']}"
    url = str(llm["endpoint"]).rstrip("/") + "/chat/completions"
    try:
        async with httpx.AsyncClient(timeout=180) as client:
            r = await client.post(url, headers=headers, json=body)
            if r.status_code != 200:   # guided_json 미지원 구버전 → json_object 폴백
                body.pop("guided_json", None)
                body["response_format"] = {"type": "json_object"}
                r = await client.post(url, headers=headers, json=body)
            if r.status_code != 200:
                return {"ok": False, "error": f"LLM {r.status_code}: {r.text[:300]}"}
            data = r.json()
            content = (((data.get("choices") or [{}])[0].get("message") or {}).get("content") or "").strip()
    except Exception as e:
        return {"ok": False, "error": str(e)[:300]}
    try:
        m = _re.search(r"\{.*\}", content, _re.DOTALL)
        obj = json.loads(m.group(0) if m else content)
        raw_steps = obj.get("steps") if isinstance(obj, dict) else []
    except Exception:
        return {"ok": False, "error": "응답 JSON 파싱 실패", "raw": content[:500]}
    clean = []
    for s in (raw_steps or []):
        if not isinstance(s, dict):
            continue
        cli = str(s.get("cli") or "").strip()
        if not cli:
            continue
        clean.append({"desc": str(s.get("desc") or ""), "cli": cli, "type": s.get("type") or "contains", "criteria": str(s.get("criteria") or "")})
    return {"ok": True, "steps": clean, "used_examples": len(examples), "model": llm.get("model"), "raw": content[:1200]}

@app.post("/api/llm/ask")
async def llm_ask(payload: dict, token: str = ""):
    """자연어 질문 → 학습 데이터 검색 + LLM(gemma) 요약 답변 (근거 포함)."""
    import httpx, re as _re
    query = str(payload.get("query") or "").strip()
    if not query:
        return {"ok": False, "error": "질문을 입력하세요"}
    items = _load_learned().get("items", [])
    terms = [t for t in _re.split(r"\s+", query.lower()) if t]
    def _hay(it):
        return (str(it.get("title", "")) + " " + " ".join(it.get("models") or []) + " " + str(it.get("role", "")) + " " + str(it.get("vendor", "")) + " "
                + " ".join((str(s.get("desc", "")) + " " + str(s.get("cli", "")) + " " + str(s.get("imageText", ""))) for s in (it.get("steps") or []))).lower()
    def _score(it):
        h = _hay(it); return sum(1 for t in terms if t in h)
    ranked = sorted(items, key=_score, reverse=True)
    matched = [it for it in ranked if _score(it) > 0][:5] or ranked[:3]
    # llm_generate 와 동일 규칙 — vLLM/OpenAI 호환 로컬 LLM만. Claude 등은 /chat/completions 미지원이라 404 원인.
    llms = (load_json(LLMS_FILE).get("llms") or [])
    def _ok_ask(l):
        if not (l.get("status", "active") == "active" and l.get("endpoint")): return False
        t = str(l.get("type") or "").lower()
        return t in ("local", "vllm", "openai", "openai-compatible", "") and t not in ("claude", "anthropic")
    llm = next((l for l in llms if _ok_ask(l)), None)
    if not llm:
        return {"ok": False, "error": "등록된 로컬(OpenAI 호환) LLM이 없습니다. AI Assistant에서 vLLM/제마를 활성화하세요.", "matched": matched}
    ctx = ""
    for it in matched:
        ctx += f"\n■ 시험항목: {it.get('title','')} (모델 {','.join(it.get('models') or [])} / 제품군 {it.get('role','')})\n"
        for s in (it.get("steps") or []):
            ctx += f"   - {s.get('desc','')}: `{s.get('cli','')}` [판정 {s.get('type','')} \"{s.get('criteria','')}\"]"
            if s.get("imageText"):
                ctx += f" / 이미지인식: {str(s.get('imageText'))[:150]}"
            ctx += "\n"
    sys_p = ("너는 사내 네트워크 장비 시험 절차 지식 어시스턴트다. 아래 '검색된 학습 데이터'만 근거로 사용자 질문에 한국어로 간결하고 명확하게 답한다. "
             "관련 시험항목과 절차(명령/판정)를 정리해 설명한다. 검색 데이터에 없는 내용은 추측하지 말고 '학습된 데이터에 없습니다'라고 답한다.")
    user_p = f"[검색된 학습 데이터]{ctx if ctx.strip() else ' (없음)'}\n\n[질문] {query}"
    body = {"model": llm.get("model") or "", "messages": [{"role": "system", "content": sys_p}, {"role": "user", "content": user_p}],
            "temperature": 0.3, "max_tokens": 1500}
    headers = {"Content-Type": "application/json"}
    if llm.get("apikey"):
        headers["Authorization"] = f"Bearer {llm['apikey']}"
    url = str(llm["endpoint"]).rstrip("/") + "/chat/completions"
    try:
        async with httpx.AsyncClient(timeout=180) as client:
            r = await client.post(url, headers=headers, json=body)
            if r.status_code != 200:
                return {"ok": False, "error": f"LLM {r.status_code}: {r.text[:300]}", "matched": matched}
            answer = (((r.json().get("choices") or [{}])[0].get("message") or {}).get("content") or "").strip()
    except Exception as e:
        return {"ok": False, "error": str(e)[:300], "matched": matched}
    return {"ok": True, "answer": answer, "matched": matched}

# ── AI 사용 로그 / 피드백 / 통계 ──
AI_USAGE_FILE = DATA_DIR / "state" / "ai_usage.json"
FEEDBACK_FILE = DATA_DIR / "state" / "ai_feedback.json"

_ITEMS_STORE_KV_MAP = {}   # str(path) → kv key
def _load_items_store(path):
    """path 가 KV 로 이전된 파일이면 DB(app_kv) 캐시 반환. 아니면 기존 파일 로드."""
    _key = _ITEMS_STORE_KV_MAP.get(str(path))
    if _key:
        d = _kv_load_sync(_key, {"items": []})
        if isinstance(d, dict) and isinstance(d.get("items"), list):
            return d
        return {"items": []}
    try:
        if path.exists():
            d = load_json(path)
            if isinstance(d, dict) and isinstance(d.get("items"), list):
                return d
    except Exception:
        pass
    return {"items": []}

def _save_items_store(path, data):
    """path 가 KV 로 이전된 파일이면 DB(app_kv) 저장. 아니면 파일 저장 (레거시)."""
    _key = _ITEMS_STORE_KV_MAP.get(str(path))
    if _key:
        _kv_save_sync(_key, data)
        return
    save_json(path, data)

def _user_of(token):
    s = SESSIONS.get(token or "")
    if s:
        return s.get("name") or s.get("username") or ""
    return ""

def _est_tokens(text):
    return max(0, round(len(str(text or "")) / 3))   # 한글·혼합 대략 3자/토큰

def _log_ai_usage(token, model, kind, question, answer, usage=None):
    try:
        store = _load_items_store(AI_USAGE_FILE)
        pin = (usage or {}).get("prompt_tokens")
        pout = (usage or {}).get("completion_tokens")
        if pin is None:
            pin = _est_tokens(question)
        if pout is None:
            pout = _est_tokens(answer)
        store["items"].insert(0, {
            "at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "user": _user_of(token) or "(미상)",
            "model": str(model or ""),
            "kind": str(kind or "chat"),
            "question": str(question or "")[:500],
            "tokens_in": int(pin or 0),
            "tokens_out": int(pout or 0),
            "tokens": int((pin or 0) + (pout or 0)),
            "estimated": usage is None,
        })
        store["items"] = store["items"][:20000]
        _save_items_store(AI_USAGE_FILE, store)
    except Exception:
        pass

@app.post("/api/ai/usage")
async def ai_usage_post(payload: dict, token: str = ""):
    _log_ai_usage(token, payload.get("model"), payload.get("kind"),
                  payload.get("question"), payload.get("answer"), payload.get("usage"))
    return {"ok": True}

@app.post("/api/ai/feedback")
async def ai_feedback_save(payload: dict, token: str = ""):
    store = _load_items_store(FEEDBACK_FILE)
    item = {
        "id": "fb-" + str(int(datetime.now().timestamp() * 1000)),
        "at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "by": _user_of(token) or str(payload.get("by") or "") or "(미상)",
        "model": str(payload.get("model") or ""),
        "thumb": int(payload.get("thumb") or 0),
        "score": int(payload.get("score") or 0),
        "reasons": [str(r)[:80] for r in (payload.get("reasons") or [])][:12] if isinstance(payload.get("reasons"), list) else [],
        "comment": str(payload.get("comment") or "")[:2000],
        "question": str(payload.get("question") or "")[:2000],
        "answer": str(payload.get("answer") or "")[:4000],
    }
    store["items"].insert(0, item)
    _save_items_store(FEEDBACK_FILE, store)
    return {"ok": True, "id": item["id"], "count": len(store["items"])}

@app.get("/api/ai/feedback")
async def ai_feedback_list(limit: int = 300):
    return {"items": _load_items_store(FEEDBACK_FILE).get("items", [])[:max(1, min(limit, 2000))]}

@app.delete("/api/ai/feedback/{fid}")
async def ai_feedback_del(fid: str, token: str = ""):
    store = _load_items_store(FEEDBACK_FILE)
    store["items"] = [it for it in store["items"] if it.get("id") != fid]
    _save_items_store(FEEDBACK_FILE, store)
    return {"ok": True}

@app.get("/api/ai/stats")
async def ai_stats(days: int = 30):
    from datetime import timedelta
    items = _load_items_store(AI_USAGE_FILE).get("items", [])
    fb = _load_items_store(FEEDBACK_FILE).get("items", [])
    if days and days > 0:
        cutoff = (datetime.now() - timedelta(days=days)).strftime("%Y-%m-%d")
        items = [it for it in items if str(it.get("at", ""))[:10] >= cutoff]
    # user → 조직 매핑
    umap = {}
    try:
        for u in _users_load_sync().get("users", []):
            for key in (u.get("name"), u.get("username")):
                if key:
                    umap[key] = u
    except Exception:
        pass
    by_user, by_model, by_org, by_day = {}, {}, {}, {}
    for it in items:
        u = it.get("user") or "(미상)"
        m = it.get("model") or "(미상)"
        day = str(it.get("at", ""))[:10]
        tk = int(it.get("tokens") or 0); ti = int(it.get("tokens_in") or 0); to = int(it.get("tokens_out") or 0)
        du = by_user.setdefault(u, {"user": u, "questions": 0, "tokens_in": 0, "tokens_out": 0, "tokens": 0})
        du["questions"] += 1; du["tokens_in"] += ti; du["tokens_out"] += to; du["tokens"] += tk
        dm = by_model.setdefault(m, {"model": m, "questions": 0, "tokens": 0})
        dm["questions"] += 1; dm["tokens"] += tk
        urec = umap.get(u)
        org = "(미지정)"
        if urec:
            org = " ▸ ".join([x for x in [urec.get("company"), urec.get("dept"), urec.get("team")] if x]) or "(미지정)"
        do = by_org.setdefault(org, {"org": org, "questions": 0, "tokens": 0, "_users": set()})
        do["questions"] += 1; do["tokens"] += tk; do["_users"].add(u)
        if day:
            dd = by_day.setdefault(day, {"day": day, "messages": 0, "tokens": 0})
            dd["messages"] += 1; dd["tokens"] += tk
    users = sorted(by_user.values(), key=lambda x: x["tokens"], reverse=True)
    models = sorted(by_model.values(), key=lambda x: x["tokens"], reverse=True)
    orgs = sorted([{"org": o["org"], "questions": o["questions"], "tokens": o["tokens"], "users": len(o["_users"])} for o in by_org.values()],
                  key=lambda x: x["tokens"], reverse=True)
    daily = sorted(by_day.values(), key=lambda x: x["day"])
    scored = [int(f.get("score") or 0) for f in fb if f.get("score")]
    fb_avg = round(sum(scored) / len(scored), 2) if scored else 0
    return {"users": users, "models": models, "orgs": orgs, "daily": daily,
            "total_questions": len(items), "total_tokens": sum(int(it.get("tokens") or 0) for it in items),
            "distinct_users": len(by_user), "days": days,
            "feedback_count": len(fb), "feedback_avg": fb_avg, "recent": items[:50]}

NOTIF_FILE = DATA_DIR / "state" / "notifications.json"

def _load_notifs() -> dict:
    try:
        if NOTIF_FILE.exists():
            d = load_json(NOTIF_FILE)
            if isinstance(d, dict) and isinstance(d.get("items"), list):
                return d
    except Exception:
        pass
    return {"items": [], "seq": 0}

def _save_notifs(d: dict):
    NOTIF_FILE.parent.mkdir(parents=True, exist_ok=True)
    save_json(NOTIF_FILE, d)

def _add_notif(to_user: str, from_user: str, from_name: str, text: str, link: str = ""):
    d = _load_notifs()
    d["seq"] = int(d.get("seq", 0)) + 1
    item = {"id": "n" + str(d["seq"]), "to": to_user, "from": from_user, "from_name": from_name,
            "text": text, "link": link, "read": False,
            "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
    d["items"].append(item)
    if len(d["items"]) > 500:
        d["items"] = d["items"][-500:]
    _save_notifs(d)
    return item

@app.get("/api/users/mentionable")
async def api_users_mentionable(token: str = ""):
    u = _user_from_token(token)
    if not u:
        raise HTTPException(401, "로그인이 필요합니다")
    return {"users": [{"username": x.get("username"), "name": x.get("name"), "email": x.get("email", ""),
                       "dept": x.get("dept", ""), "team": x.get("team", "")}
                      for x in _users_load_sync()["users"] if x.get("active", True)]}

@app.get("/api/notifications")
async def api_notifs_get(token: str = ""):
    u = _user_from_token(token)
    if not u:
        raise HTTPException(401, "로그인이 필요합니다")
    me = u.get("username")
    items = [n for n in _load_notifs()["items"] if n.get("to") == me]
    items = list(reversed(items))[:100]
    unread = sum(1 for n in items if not n.get("read"))
    return {"items": items, "unread": unread}

@app.post("/api/notifications/read")
async def api_notifs_read(payload: dict, token: str = ""):
    u = _user_from_token(token)
    if not u:
        raise HTTPException(401, "로그인이 필요합니다")
    me = u.get("username")
    nid = payload.get("id")
    d = _load_notifs()
    for n in d["items"]:
        if n.get("to") == me and (nid is None or n.get("id") == nid):
            n["read"] = True
    _save_notifs(d)
    return {"ok": True}

@app.post("/api/mention")
async def api_mention(payload: dict, token: str = ""):
    u = _user_from_token(token)
    if not u:
        raise HTTPException(401, "로그인이 필요합니다")
    mentions = payload.get("mentions") or []
    if isinstance(mentions, str):
        mentions = [mentions]
    text = str(payload.get("text", "")).strip()
    link = str(payload.get("link", "")).strip()
    ctx = str(payload.get("context", "")).strip()
    by_name = {x.get("username"): x for x in _users_load_sync()["users"]}
    me = u.get("username"); my_name = u.get("name") or me
    mail_cfg = _load_mail_cfg()
    notified = []
    for mname in mentions:
        tu = by_name.get(mname)
        if not tu or mname == me:
            continue
        body_text = my_name + "님이 회원님을 멘션했습니다" + ((" · " + ctx) if ctx else "") + ((": " + text) if text else "")
        _add_notif(mname, me, my_name, body_text, link)
        notified.append(mname)
        try:
            if mail_cfg.get("enabled") and tu.get("email"):
                _send_mail(tu["email"], "[ubiQuoss-TOP] " + my_name + "님의 멘션",
                           body_text + (("\n\n바로가기: " + link) if link else ""))
        except Exception:
            pass
    return {"ok": True, "notified": notified}

# ───────────────────────────────────────────
# 라우터 - REQ/TC 파일별 관리
# ───────────────────────────────────────────
REQ_DIR      = DATA_DIR / "req"
TC_DIR       = DATA_DIR / "tc"
CYCLE_DIR    = DATA_DIR / "cycle"
TRASH_DIR    = DATA_DIR / "trash"   # REQ/TC 소프트 삭제(휴지통) — 복원 가능
FOLDERS_FILE = DATA_DIR / "state" / "folders.json"

# 서버 시작 시 디렉토리 초기화
for _d in [REQ_DIR, TC_DIR, CYCLE_DIR, TRASH_DIR]:
    _d.mkdir(exist_ok=True)
if not FOLDERS_FILE.exists():
    save_json(FOLDERS_FILE, {"folders": []})

def init_req_dirs():
    REQ_DIR.mkdir(exist_ok=True)
    TC_DIR.mkdir(exist_ok=True)
    CYCLE_DIR.mkdir(exist_ok=True)
    TRASH_DIR.mkdir(exist_ok=True)
    if not FOLDERS_FILE.exists():
        save_json(FOLDERS_FILE, {"folders": []})

def _trash_put(kind, item_id, data, bundle=None):
    """REQ/TC 삭제 시 휴지통에 보관(복원 가능). bundle=REQ 삭제 시 딸린 TC 데이터 목록."""
    import datetime as _dt
    TRASH_DIR.mkdir(exist_ok=True)
    _now = _dt.datetime.now()
    tid = _now.strftime("%Y%m%d_%H%M%S_") + str(_now.microsecond) + "__" + str(kind) + "__" + str(item_id)
    try:
        nm = str((data or {}).get("name") or (data or {}).get("title") or (data or {}).get("summary") or item_id)
    except Exception:
        nm = str(item_id)
    rec = {"trash_id": tid, "kind": kind, "id": item_id, "name": nm,
           "deleted_at": _now.isoformat(timespec="seconds"), "data": data, "bundle": bundle or []}
    try:
        save_json(TRASH_DIR / (tid + ".json"), rec)
    except Exception:
        pass
    return tid
# 폴더 구조
@app.get("/api/folders")
async def get_folders():
    init_req_dirs()
    return load_json(FOLDERS_FILE)

@app.post("/api/folders")
async def save_folders(data: dict):
    init_req_dirs()
    save_json(FOLDERS_FILE, data)
    return {"success": True}


# ───────────────────────────────────────────
# 요구사항 분류 (3단 고정: 대분류 > 중분류 > 소분류)
#
# 옛 폴더 트리는 깊이 제한이 없어 프로토콜·계층·기능이 한 경로에 섞였다
# (IPV4_L2 > VLAN). 그래서 같은 기능이 여러 가지에 중복 등록됐다.
# 상한을 두되 3단까지는 허용한다. 이 규칙을 서버에서 강제한다 —
# DB 제약으로는 재귀 깊이를 막을 수 없다.
# ───────────────────────────────────────────
# 원래 폴더 구조에 이미 4단짜리가 있었다
# (U-REQ-PA1T-TC > 1. 부품 변경 > 1-1. 메인 메모리/WDT TC > 1-1-1. 부팅 1000회).
# 3단으로 묶어두니 E43·E57·LG 처럼 이미 3단을 쓴 가지가 어디로도 못 갔다.
MAX_CAT_DEPTH = 4
CAT_DEPTH_MSG = "분류는 4단까지만 만들 수 있습니다"
PRJ_ROOT_MSG = "프로젝트는 트리 맨 위에만 둘 수 있습니다"


async def _is_project_cat(cid: str) -> bool:
    """이 분류가 프로젝트(최상위 전용)인가 — 이동 검증이 쓴다."""
    async with db.pool().acquire() as c:
        return await c.fetchval("SELECT 1 FROM project WHERE cat_id=$1", cid) is not None


async def _req_chain_resync() -> int:
    """요구사항의 분류 사슬(cat1~4)을 트리 기준으로 다시 쓴다.

    사슬은 트리의 사본이라 폴더가 이동하면 낡는다. 낡은 사본은 옛 폴더에
    요구사항을 계속 매달아 두어 「폴더를 옮겼는데 요구사항이 안 따라왔다」
    로 보인다(실사고: 폴더를 프로젝트 밑으로 옮겼을 때). 놓인 칸(살아
    있는 가장 깊은 칸)이 사실이고, 그 조상 사슬로 위 칸들을 다시 채운다.
    폴더 이동 때마다·기동 때 1회 부른다 — 한 바퀴 훑기라 수천 건도 싸다.
    """
    async with db.pool().acquire() as c:
        cats = {
            r["id"]: r["parent_id"]
            for r in await c.fetch("SELECT id, parent_id FROM req_category")
        }
        # 화면(req_list_full)은 data(JSONB) 를 서빙하므로 사실도 data 에서
        # 읽고, 고칠 때도 컬럼과 data 를 함께 쓴다 — 컬럼만 고치면 목록이
        # 옛값을 계속 보인다 (TC 메타에서 겪은 함정).
        rows = await c.fetch(
            """
            SELECT id, data->>'cat1' AS cat1, data->>'cat2' AS cat2,
                   data->>'cat3' AS cat3, data->>'cat4' AS cat4
            FROM req
            """
        )
        n = 0
        for r in rows:
            deep = next(
                (r[k] for k in ("cat4", "cat3", "cat2", "cat1") if r[k] and r[k] in cats),
                None,
            )
            if not deep:
                continue
            chain: list = []
            cur = deep
            while cur and cur in cats and cur not in chain:
                chain.insert(0, cur)
                cur = cats[cur]
            chain = (chain + ["", "", "", ""])[:4]
            if [r["cat1"] or "", r["cat2"] or "", r["cat3"] or "", r["cat4"] or ""] != chain:
                await c.execute(
                    """
                    UPDATE req SET
                      cat1=NULLIF($1,''), cat2=NULLIF($2,''),
                      cat3=NULLIF($3,''), cat4=NULLIF($4,''),
                      data = data || jsonb_build_object(
                        'cat1', $1::text, 'cat2', $2::text,
                        'cat3', $3::text, 'cat4', $4::text)
                    WHERE id=$5
                    """,
                    *chain, r["id"],
                )
                n += 1
        return n


class ReqCategoryIn(BaseModel):
    name: str
    parent_id: Optional[str] = None
    sort_order: int = 0


async def _cat_children_map() -> dict:
    m: dict = {}
    for c in await db.cat_list():
        m.setdefault(c.get("parent_id"), []).append(c["id"])
    return m


async def _cat_descendants(cid: str, include_self: bool = True) -> set:
    kids = await _cat_children_map()
    out, stack = (set([cid]) if include_self else set()), list(kids.get(cid, []))
    while stack:
        x = stack.pop()
        if x in out:
            continue
        out.add(x)
        stack.extend(kids.get(x, []))
    return out


async def _is_descendant(node: str, ancestor: str) -> bool:
    return node in await _cat_descendants(ancestor, include_self=False)


async def _subtree_height(cid: str) -> int:
    """자기만 있으면 1, 자식이 있으면 2, 손자까지면 3."""
    kids = await _cat_children_map()

    def h(x: str) -> int:
        ch = kids.get(x, [])
        return 1 if not ch else 1 + max(h(k) for k in ch)

    return h(cid)


@app.get("/api/req-categories")
async def list_req_categories():
    return {"categories": await db.cat_list()}


@app.post("/api/req-categories")
async def create_req_category(body: ReqCategoryIn):
    name = (body.name or "").strip()
    if not name:
        raise HTTPException(400, "분류 이름을 입력하세요")

    parent_id = (body.parent_id or "").strip() or None
    if parent_id:
        parent = await db.cat_get(parent_id)
        if parent is None:
            raise HTTPException(404, "상위 분류를 찾을 수 없습니다")
        if await db.cat_depth(parent_id) >= MAX_CAT_DEPTH:
            raise HTTPException(400, CAT_DEPTH_MSG)

    cid = f"cat-{int(datetime.now().timestamp() * 1000)}"
    try:
        await db.cat_upsert(cid, name, parent_id, body.sort_order)
    except Exception as e:
        # 유니크 인덱스 위반 = 같은 상위 아래 같은 이름
        if "uq_req_category" in str(e):
            raise HTTPException(409, f"'{name}' 은 이미 있습니다") from e
        raise
    return {"success": True, "id": cid}


# ★ 이 라우트는 반드시 /{cat_id} 라우트보다 위에 있어야 한다.
#   아래에 두면 'reorder' 가 cat_id 로 잡혀 405 가 난다.
class ReqCategoryOrderIn(BaseModel):
    """한 상위 아래 형제들의 새 순서. ids 에 적힌 차례대로 sort_order 를 매긴다."""
    parent_id: Optional[str] = None
    ids: list[str]


@app.post("/api/req-categories/reorder")
async def reorder_req_categories(body: ReqCategoryOrderIn):
    """형제 순서 재배치 + 필요하면 상위 이동까지 한 번에.

    화면에서 '폴더와 폴더 사이' 에 놓으면 여기로 온다. 한 건씩 PUT 하면
    중간 상태가 보이고, 실패했을 때 절반만 적용된 채로 남는다.
    그래서 한 트랜잭션에서 형제 전체를 다시 매긴다.
    """
    parent = (body.parent_id or "").strip() or None
    if parent and await db.cat_get(parent) is None:
        raise HTTPException(404, "상위 분류를 찾을 수 없습니다")

    base = await db.cat_depth(parent) if parent else 0
    for cid in body.ids:
        if parent and (cid == parent or await _is_descendant(parent, cid)):
            raise HTTPException(400, "자기 하위 분류 밑으로는 옮길 수 없습니다")
        if parent and await _is_project_cat(cid):
            raise HTTPException(400, PRJ_ROOT_MSG)
        if base + await _subtree_height(cid) > MAX_CAT_DEPTH:
            raise HTTPException(400, CAT_DEPTH_MSG)

    async with db.pool().acquire() as c:
        async with c.transaction():
            for i, cid in enumerate(body.ids):
                await c.execute(
                    "UPDATE req_category SET parent_id=$1, sort_order=$2, updated_at=now() WHERE id=$3",
                    parent, i * 10, cid,
                )
    # 사이에 끼우기로도 상위가 바뀐다 — 여기서도 사슬을 맞춘다.
    try:
        await _req_chain_resync()
    except Exception as e:
        print(f"[reorder] 사슬 재작성 실패: {e}", flush=True)
    return {"success": True, "count": len(body.ids)}


@app.put("/api/req-categories/{cat_id}")
async def update_req_category(cat_id: str, body: ReqCategoryIn):
    cur = await db.cat_get(cat_id)
    if cur is None:
        raise HTTPException(404, "분류를 찾을 수 없습니다")

    name = (body.name or "").strip()
    if not name:
        raise HTTPException(400, "분류 이름을 입력하세요")

    parent_id = (body.parent_id or "").strip() or None
    if parent_id == cat_id:
        raise HTTPException(400, "자기 자신을 상위로 지정할 수 없습니다")
    if parent_id:
        # 프로젝트는 트리 맨 위가 자리다 — 폴더 밑으로 들어가면
        # 최상위=프로젝트 층 자체가 무너진다.
        if await _is_project_cat(cat_id):
            raise HTTPException(400, PRJ_ROOT_MSG)
        parent = await db.cat_get(parent_id)
        if parent is None:
            raise HTTPException(404, "상위 분류를 찾을 수 없습니다")
        # 자기 자손 밑으로 옮기면 순환이 된다.
        if cat_id in await _cat_descendants(cat_id, include_self=False) or await _is_descendant(
            parent_id, cat_id
        ):
            raise HTTPException(400, "자기 하위 분류 밑으로는 옮길 수 없습니다")
        # 옮긴 뒤 (상위 깊이 + 이 가지의 높이) 가 상한을 넘으면 안 된다.
        if await db.cat_depth(parent_id) + await _subtree_height(cat_id) > MAX_CAT_DEPTH:
            raise HTTPException(400, CAT_DEPTH_MSG)
    try:
        await db.cat_upsert(cat_id, name, parent_id, body.sort_order)
    except Exception as e:
        if "uq_req_category" in str(e):
            raise HTTPException(409, f"'{name}' 은 이미 있습니다") from e
        raise
    # 상위가 바뀌었으면(이동) 요구사항 사슬을 트리에 맞춘다 — 안 하면
    # 옛 폴더가 그 요구사항들을 계속 잡고 있다.
    if (cur.get("parent_id") or None) != parent_id:
        try:
            await _req_chain_resync()
        except Exception as e:
            print(f"[cat] 사슬 재작성 실패: {e}", flush=True)
    return {"success": True}


@app.delete("/api/req-categories/{cat_id}")
async def delete_req_category(cat_id: str):
    """하위 분류까지 함께 지운다. 요구사항은 지우지 않고 '미분류'가 된다."""
    if not await db.cat_delete(cat_id):
        raise HTTPException(404, "분류를 찾을 수 없습니다")
    return {"success": True}


# ───────────────────────────────────────────
# 프로젝트 — 요구사항 트리의 최상위 폴더가 곧 프로젝트다(itest 방식).
# 이름의 정본은 폴더(req_category.name)라서 폴더 이름 변경이 곧 프로젝트명
# 변경이고, 폴더 삭제가 곧 프로젝트 삭제다(FK CASCADE). 여기는 고객사·
# 모델 같은 메타만 맡는다.
# ───────────────────────────────────────────
class ProjectIn(BaseModel):
    name: str
    customer: str = ""
    model_group: str = ""
    model: str = ""
    description: str = ""


# ───────────────────────────────────────────
# ID 옮기기 — 모델그룹 기준(E61xx-R0001)으로
#
# **화면에서 눌러서 한다.** 서버에 들어가 명령을 치는 방식이면 253 처럼
# 제가 손 못 대는 곳은 사람이 거기까지 가서 쳐야 한다. 미리 보기로
# 무엇이 무엇으로 바뀌는지 먼저 보이고, 그다음에 누르게 한다 —
# 되돌릴 수는 있지만(id_alias), 안 보고 누르게 두면 안 된다.
# ───────────────────────────────────────────
@app.get("/api/id-alias")
async def id_alias_lookup(old: str = ""):
    """옛 ID 로 물으면 새 ID 를 준다.

    주소·위키·메일에 붙여 둔 옛 ID 는 우리가 못 고친다. 화면이 못 찾았을 때
    여기 한 번 물어보고 넘어가면, 옛 링크가 계속 살아 있다. 로그인만 되면
    쓸 수 있게 두었다 — 이 표는 「무엇이 무엇이 되었나」 뿐이라 숨길 것이 없고,
    막아 두면 링크가 끊기는 쪽 손해가 크다.
    """
    v = (old or "").strip()
    if not v:
        return {"new_id": ""}
    async with db.pool().acquire() as c:
        row = await c.fetchrow("SELECT new_id, kind FROM id_alias WHERE old_id = $1", v)
    return {"new_id": row["new_id"] if row else "", "kind": row["kind"] if row else ""}


@app.get("/api/id-migrate/plan")
async def id_migrate_plan(token: str = ""):
    _require_admin(token)
    async with db.pool().acquire() as c:
        return await id_migrate.plan(c)


@app.post("/api/id-migrate/apply")
async def id_migrate_apply(token: str = ""):
    _require_admin(token)
    async with db.pool().acquire() as c:
        p = await id_migrate.plan(c)
        if not p["moves"]:
            return {"ok": True, "counts": {}, "note": "옮길 것이 없습니다"}
        counts = await id_migrate.apply(c, p)
    print(f"[id] 옮김 {counts}", flush=True)
    return {"ok": True, "counts": counts, "skipped": len(p["skipped"])}


@app.get("/api/projects")
async def list_projects():
    async with db.pool().acquire() as c:
        rows = await c.fetch(
            """
            SELECT p.id, p.cat_id, c.name, p.customer, p.model_group, p.model,
                   p.description, p.created_at
            FROM project p JOIN req_category c ON c.id = p.cat_id
            ORDER BY c.name
            """
        )
    return {"projects": [dict(r) for r in rows]}


@app.post("/api/projects")
async def create_project(body: ProjectIn):
    name = (body.name or "").strip()
    if not name:
        raise HTTPException(400, "프로젝트 이름을 입력하세요")
    now_ms = int(datetime.now().timestamp() * 1000)
    cid, pid = f"cat-{now_ms}", f"prj-{now_ms}"
    try:
        await db.cat_upsert(cid, name, None, 0)
    except Exception as e:
        if "uq_req_category" in str(e):
            raise HTTPException(409, f"'{name}' 은 이미 있습니다") from e
        raise
    async with db.pool().acquire() as c:
        await c.execute(
            """
            INSERT INTO project (id, cat_id, customer, model_group, model, description)
            VALUES ($1, $2, $3, $4, $5, $6)
            """,
            pid, cid, (body.customer or "").strip(), (body.model_group or "").strip(),
            (body.model or "").strip(), (body.description or "").strip(),
        )
    return {"success": True, "id": pid, "cat_id": cid}


# ───────────────────────────────────────────
# 문서 → 마크다운 변환
#
# 워드(.docx)·PDF 는 브라우저가 제대로 읽지 못한다. 서버에서 바꿔서 돌려준다.
# 결과를 마크다운으로 두는 이유는 그게 이 시스템의 정본이기 때문이다 —
# 구현내용으로 들어가고, 벡터 DB 에 실리고, 시험항목 생성의 입력이 된다.
# ───────────────────────────────────────────
# ───────────────────────────────────────────
# 구현내용에 붙이는 이미지
#
# 파일은 data/req_images/ 에 둔다. 이 폴더는 도커 볼륨(app-data)이라
# 소스 트리에 쌓이지 않고, 볼륨 하나만 챙기면 함께 백업된다.
# 마크다운에는 ![](/api/req-images/<파일명>) 으로 들어간다 — 원문이
# 정본이므로 경로도 원문 안에 남아야 한다.
# ───────────────────────────────────────────
REQ_IMG_DIR = DATA_DIR / "req_images"

_IMG_EXT = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg", ".bmp"}


@app.post("/api/upload/image")
async def upload_image(file: UploadFile = File(...)):
    ext = Path(file.filename or "").suffix.lower()
    if ext not in _IMG_EXT:
        raise HTTPException(400, f"이미지 파일만 올릴 수 있습니다 ({', '.join(sorted(_IMG_EXT))})")
    raw = await file.read()
    if not raw:
        raise HTTPException(400, "빈 파일입니다")
    if len(raw) > 10 * 1024 * 1024:
        raise HTTPException(413, "10MB 이하만 올릴 수 있습니다")

    REQ_IMG_DIR.mkdir(parents=True, exist_ok=True)
    # 이름은 서버가 정한다. 사용자가 준 이름을 그대로 쓰면 경로 조작과
    # 덮어쓰기가 열린다.
    import secrets
    name = f"{int(datetime.now().timestamp() * 1000)}-{secrets.token_hex(4)}{ext}"
    (REQ_IMG_DIR / name).write_bytes(raw)
    return {"url": f"/api/req-images/{name}", "name": name, "size": len(raw)}


@app.get("/api/req-images/{name}")
async def get_req_image(name: str):
    # 이름만 받는다. 경로가 섞여 들어오면 거부 — 상위 폴더 탈출 방지.
    if "/" in name or "\\" in name or name.startswith("."):
        raise HTTPException(400, "잘못된 파일명입니다")
    f = REQ_IMG_DIR / name
    if not f.is_file():
        raise HTTPException(404, "이미지를 찾을 수 없습니다")
    return FileResponse(str(f), headers={"Cache-Control": "public, max-age=31536000, immutable"})


# ───────────────────────────────────────────
# 자원 점유 (장비 · 계측기)
#
# 같은 장비를 두 사람이 동시에 잡으면 시험이 통째로 망가진다. 50명이
# 함께 쓰면 반드시 필요하다.
#
# 락은 사이클이 끝날 때까지 유지한다(시간 만료 없음). 자동으로 풀면
# 실제 시험 중인 장비를 남이 뺏을 수 있다. 대신 살아있음 신호를 남겨
# '응답 없음' 을 보여주고, 푸는 것은 사람이 판단한다.
# ───────────────────────────────────────────
class LockIn(BaseModel):
    resource_id: str
    kind: str = "device"          # 'device' | 'instrument'
    cycle_id: Optional[str] = None
    note: Optional[str] = None


def _me(request) -> dict:
    """미들웨어가 넣어둔 세션. 없으면 401 (미들웨어가 이미 막지만 방어적으로)."""
    s = getattr(request.state, "user", None)
    if not s:
        raise HTTPException(401, "로그인이 필요합니다")
    return s


# ───────────────────────────────────────────
# 장비 (PG)
#
# 키는 IP. 같은 모델이 여러 대여도 접속 대상은 IP 로 갈린다.
# 제품군(L2·L3·OLT·ONT·CPE·HGW)은 role, 제조사는 vendor 에 담는다.
#
# 옛 /api/devices (devices.json) 와 경로를 나눠 둔다 — 옛 화면이 아직
# 그걸 쓰고 있어서, 한 번에 갈아치우면 옛 화면이 멈춘다.
# ───────────────────────────────────────────
DEVICE_ROLES = ["L2", "L3", "OLT", "ONT", "CPE", "HGW", "계측기", "기타"]


# 기존 앱의 /api/device-catalog(app_kv 기반) 과 경로가 겹친다.
# 먼저 선언된 쪽이 이기므로 그대로 두면 옛 화면이 조용히 망가진다.
# devices2 와 같은 규칙으로 2 를 붙인다.
@app.get("/api/codes")
async def codes_list(kind: str = ""):
    """드롭다운에 들어가는 값 목록. 화면은 여기서만 읽는다."""
    items = await db.code_list(kind)
    for it in items:
        it["used"] = await db.code_usage(it["kind"], it["value"])
    # 탭 이름 덮어쓰기 — 기본 이름(상태 등)을 사람이 바꿀 수 있다
    kinds = dict(db.CODE_KINDS)
    try:
        ov = _kv_load_sync("code_kind_labels", {}) or {}
        for k, v in ov.items():
            if k in kinds and str(v).strip():
                kinds[k] = str(v).strip()
        # 숨긴 기본 칸 — 탭·입력칸에서 빠진다 (값 자료는 남는다)
        hid = _kv_load_sync("code_kind_hidden", []) or []
        for k in hid:
            kinds.pop(k, None)
    except Exception:
        pass
    return {"items": items, "kinds": kinds}


@app.get("/api/codes/orphans")
async def codes_orphans(kind: str = ""):
    """쓰이고 있는데 목록에 없는 값 — 설정 화면이 「목록에 넣기」를 띄운다."""
    if not kind:
        return {"items": []}
    return {"items": await db.code_orphans(kind)}


@app.get("/api/cycle-desc-template")
async def cycle_desc_template_get():
    """사이클 설명 틀 — 보고서 패턴을 맞추려고 사람이 정의해 둔다."""
    d = _kv_load_sync("cycle_desc_template", {}) or {}
    return {"text": str(d.get("text") or "")}


@app.post("/api/cycle-desc-template")
async def cycle_desc_template_set(payload: dict):
    _kv_save_sync("cycle_desc_template", {"text": str(payload.get("text") or "")})
    return {"success": True}


@app.post("/api/codes/kind-hidden")
async def codes_kind_hidden(payload: dict, token: str = ""):
    """기본 칸(탭) 숨기기/되살리기 — 값 자료는 지우지 않는다.

    **관리자만**(지시). 세 화면이 함께 쓰는 설정이라, 한 사람이 고치면 모두의
    화면이 바뀐다."""
    _require_admin(token)
    kind = str(payload.get("kind") or "").strip()
    hidden = bool(payload.get("hidden"))
    if kind not in db.CODE_KINDS:
        raise HTTPException(400, f"알 수 없는 종류입니다: {kind}")
    hid = set(_kv_load_sync("code_kind_hidden", []) or [])
    if hidden:
        hid.add(kind)
    else:
        hid.discard(kind)
    _kv_save_sync("code_kind_hidden", sorted(hid))
    return {"success": True}


@app.post("/api/codes/kind-label")
async def codes_kind_label(payload: dict, token: str = ""):
    """기본 칸(탭)의 표시 이름 바꾸기 — 빈 이름이면 원래대로. **관리자만**(지시)."""
    _require_admin(token)
    kind = str(payload.get("kind") or "").strip()
    label = str(payload.get("label") or "").strip()
    if kind not in db.CODE_KINDS:
        raise HTTPException(400, f"알 수 없는 종류입니다: {kind}")
    ov = _kv_load_sync("code_kind_labels", {}) or {}
    if label:
        ov[kind] = label
    else:
        ov.pop(kind, None)
    _kv_save_sync("code_kind_labels", ov)
    return {"success": True}


@app.get("/api/codes/kind-style")
async def codes_kind_style_get():
    """필드(탭) 단위 모양 — 폭·모양·정렬.

    값마다의 색은 code.note 에 산다. 이건 **그 필드 전체**의 생김새다:
    목록에서 몇 px 를 차지하고, 값을 셀 채움으로 그릴지 알약으로 그릴지.
    여태 코드에 박혀 있어 폭 하나 고치는 데도 배포를 해야 했다(지시).
    """
    return {"styles": _kv_load_sync("code_kind_style", {}) or {}}


@app.post("/api/codes/kind-style")
async def codes_kind_style_set(payload: dict, token: str = ""):
    """{kind, w, shape, align, weight, size, font, caps} — 빈 값은 지운다.

    **관리자만**(지시). 열 폭 하나가 모두의 목록을 바꾼다."""
    _require_admin(token)
    kind = str(payload.get("kind") or "").strip()
    if not kind:
        raise HTTPException(400, "어느 필드인지 알려 주세요")
    cur = _kv_load_sync("code_kind_style", {}) or {}
    one = dict(cur.get(kind) or {})
    for k in ("w", "shape", "align", "weight", "size", "font", "caps"):
        v = payload.get(k)
        if v is None or str(v).strip() == "":
            one.pop(k, None)
        else:
            one[k] = str(v).strip()
    if one:
        cur[kind] = one
    else:
        cur.pop(kind, None)
    _kv_save_sync("code_kind_style", cur)
    return {"success": True, "styles": cur}


@app.post("/api/codes")
async def codes_save(payload: dict, token: str = ""):
    """드롭다운에 들어가는 값 추가·수정 — **관리자만**(지시)."""
    _require_admin(token)
    try:
        await db.code_upsert(payload)
    except ValueError as e:
        raise HTTPException(400, str(e)) from e
    return {"success": True}


@app.delete("/api/codes/{kind}/{value}")
async def codes_delete(kind: str, value: str, token: str = ""):
    _require_admin(token)
    # 쓰는 건수가 있어도 막지 않는다(피드백) — 지우는 것은 고르기 목록의
    # 항목뿐이고, 기록(data)에 저장된 값 문자열은 그대로 남는다.
    # 몇 건이 쓰는지는 화면이 확인창에서 미리 알린다.
    if not await db.code_delete(kind, value):
        raise HTTPException(404, "없는 항목입니다")
    return {"success": True}


# ───────────────────────────────────────────
# 커스텀 필드
#
# 팀마다 TC·요구사항에 적어두고 싶은 항목이 다르다. 컬럼을 늘리는 대신
# 여기서 정의만 관리하고, 값은 data->'custom' 에 담는다.
# ───────────────────────────────────────────
@app.get("/api/custom-fields")
async def custom_fields_list(target: str = ""):
    items = await db.cf_list(target)
    for it in items:
        it["used"] = await db.cf_usage(it["target"], it["key"])
    return {"items": items, "targets": db.CF_TARGETS, "types": db.CF_TYPES}


@app.post("/api/custom-fields")
async def custom_fields_save(payload: dict):
    try:
        cf_id = await db.cf_upsert(payload)
    except ValueError as e:
        raise HTTPException(400, str(e)) from e
    return {"success": True, "id": cf_id}


@app.delete("/api/custom-fields/{cf_id}")
async def custom_fields_delete(cf_id: int):
    cur = await db.cf_get(cf_id)
    if cur is None:
        raise HTTPException(404, "없는 필드입니다")
    # 값은 data->'custom' 에 그대로 남는다. 지우는 것은 정의뿐이라
    # 되돌리려면 같은 키로 다시 만들면 값이 도로 보인다. 그래서 쓰는 건수가
    # 있어도 막지 않고, 몇 건인지만 화면이 미리 물어보게 한다.
    if not await db.cf_delete(cf_id):
        raise HTTPException(404, "없는 필드입니다")
    return {"success": True}


@app.get("/api/device-catalog2")
async def device_catalog_list(kind: str = ""):
    items = await db.catalog_list(kind)
    # 지울 수 있는지 화면이 알 수 있게 쓰는 장비 수를 함께 준다
    for it in items:
        it["used"] = await db.catalog_usage(it["kind"], it["name"])
    return {"items": items, "kinds": list(db.CATALOG_KINDS)}


@app.post("/api/device-catalog2")
async def device_catalog_save(payload: dict):
    try:
        await db.catalog_upsert(payload)
    except ValueError as e:
        raise HTTPException(400, str(e)) from e
    return {"success": True}


@app.post("/api/device-catalog2/classify")
async def device_catalog_classify(payload: dict):
    """모델을 벤더·제품군·모델그룹으로 옮긴다 — 준 칸만 고친다."""
    try:
        await db.catalog_classify(str(payload.get("name") or ""), payload)
    except ValueError as e:
        raise HTTPException(400, str(e)) from e
    return {"success": True}


@app.post("/api/device-catalog2/rename")
async def device_catalog_rename(payload: dict):
    kind = str(payload.get("kind") or "").strip()
    old = str(payload.get("old") or "").strip()
    new = str(payload.get("new") or "").strip()
    if kind == "model":
        raise HTTPException(400, "모델명은 사이클·시험이 물려 있어 여기서 못 바꿉니다")
    if kind not in db.CATALOG_KINDS or not old or not new:
        raise HTTPException(400, "kind·old·new 가 필요합니다")
    try:
        await db.catalog_rename(kind, old, new)
    except ValueError as e:
        raise HTTPException(404, str(e)) from e
    return {"success": True}


@app.delete("/api/device-catalog2/{kind}/{name}")
async def device_catalog_delete(kind: str, name: str):
    used = await db.catalog_usage(kind, name)
    if used:
        # 몇 대인지만 말하면 다음에 할 일을 모른다 — 어느 장비인지 찍어 준다
        who = await db.catalog_users(kind, name)
        tail = " …" if used > len(who) else ""
        raise HTTPException(
            400,
            f"{used}대가 쓰고 있어 지울 수 없습니다 — {', '.join(who)}{tail}\n"
            "장비 화면에서 이 장비의 모델을 바꾸거나 장비를 지운 뒤 다시 시도하세요",
        )
    if not await db.catalog_delete(kind, name):
        raise HTTPException(404, "없는 항목입니다")
    return {"success": True}


@app.get("/api/device-roles")
async def device_roles():
    # 카탈로그에 등록된 것을 먼저 쓴다. 비어 있으면 기본 목록으로 시작한다.
    cat = await db.catalog_list()
    by: dict = {}
    for it in cat:
        by.setdefault(it["kind"], []).append(it["name"])

    # 이미 쓰고 있는 값도 함께 준다. 카탈로그에 아직 안 올린 것이 목록에서
    # 빠지면 그 장비를 편집할 때 값이 사라진 것처럼 보인다.
    async def distinct(col: str) -> list[str]:
        async with db.pool().acquire() as c:
            rows = await c.fetch(
                f"SELECT DISTINCT {col} AS v FROM device "
                f"WHERE {col} IS NOT NULL AND {col} <> '' ORDER BY 1"
            )
        return [r["v"] for r in rows]

    def merge(kind: str, used: list[str], fallback: list[str] | None = None) -> list[str]:
        out = list(by.get(kind) or fallback or [])
        for v in used:
            if v not in out:
                out.append(v)
        return out

    return {
        "roles": merge("family", await distinct("role"), DEVICE_ROLES),
        "labs": merge("lab", await distinct("lab")),
        "vendors": merge("vendor", await distinct("vendor")),
        "models": merge("model", await distinct("model")),
        "groups": by.get("group") or [],
        "usernames": await distinct("username"),
        # 모델을 고르면 제조사·제품군·기본 인터페이스를 채운다
        "model_info": {
            it["name"]: {
                "vendor": it.get("vendor"),
                "model_group": it.get("model_group"),
                "family": it.get("family"),
                "interfaces": it.get("interfaces"),
            }
            for it in cat
            if it["kind"] == "model"
        },
        "protocols": list(db.PROTOCOLS),
        "cli_protocols": list(db.CLI_PROTOCOLS),
    }


# 접속 확인은 장비에 실제로 붙어 본다. 동기 라이브러리(paramiko/telnetlib)라
# 그대로 await 하면 이벤트 루프가 멈춰 50명 전원이 같이 멈춘다. 스레드로 보낸다.
_PROBE_TIMEOUT = 6


def _probe_sync(proto: str, host: str, port: int) -> tuple[bool, str]:
    """포트가 열려 있는지만 본다.

    로그인까지 해보면 확실하지만 시간이 오래 걸리고, 잘못된 계정으로 여러 번
    시도하면 장비가 계정을 잠근다. 목록의 '연결상태' 는 '길이 열려 있나' 로 충분하다.
    """
    import socket
    if not host or not port:
        return False, "주소 또는 포트가 비어 있습니다"
    try:
        with socket.create_connection((host, int(port)), timeout=_PROBE_TIMEOUT):
            return True, ""
    except OSError as e:
        return False, str(e)


def _snmp_probe_sync(host: str, port: int, community: str) -> tuple[bool, str]:
    """SNMPv2c 로 sysDescr.0 을 실제로 읽어 본다 — 의존성 없이 최소 BER.

    UDP 라 TCP 처럼 포트 열림을 볼 수 없고, community 가 틀리면 장비가
    아예 응답하지 않는 것이 보통이다. 그래서 「응답 없음」 은 주소·포트·
    community 셋 중 하나가 틀렸다는 뜻이다.
    """
    import os
    import socket

    if not host:
        return False, "주소가 비어 있습니다"

    def tlv(t: int, v: bytes) -> bytes:
        n = len(v)
        if n < 0x80:
            return bytes([t, n]) + v
        eb = n.to_bytes((n.bit_length() + 7) // 8, "big")
        return bytes([t, 0x80 | len(eb)]) + eb + v

    def ber_int(n: int) -> bytes:
        b = n.to_bytes((max(n.bit_length(), 1) + 8) // 8, "big", signed=True)
        return tlv(0x02, b)

    oid = bytes([0x2B, 6, 1, 2, 1, 1, 1, 0])  # 1.3.6.1.2.1.1.1.0 = sysDescr.0
    rid = int.from_bytes(os.urandom(2), "big") & 0x7FFF
    vb = tlv(0x30, tlv(0x06, oid) + b"\x05\x00")
    pdu = tlv(0xA0, ber_int(rid) + ber_int(0) + ber_int(0) + tlv(0x30, vb))
    msg = tlv(0x30, ber_int(1) + tlv(0x04, (community or "public").encode()) + pdu)

    s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
    s.settimeout(_PROBE_TIMEOUT)
    try:
        s.sendto(msg, (host, int(port or 161)))
        data, _ = s.recvfrom(65535)
    except socket.timeout:
        return False, "응답 없음 — 주소·포트(161)·community 확인 (SNMP 줄의 계정 칸이 community, 비우면 public)"
    except OSError as e:
        return False, str(e)
    finally:
        s.close()

    # 관대하게 판다 — error-status 만 읽고, 나머지가 이상해도 응답이 온
    # 것 자체가 SNMP 가 살아 있다는 뜻이다.
    def read_tlv(b: bytes, i: int):
        t = b[i]
        ln = b[i + 1]
        i += 2
        if ln & 0x80:
            k = ln & 0x7F
            ln = int.from_bytes(b[i : i + k], "big")
            i += k
        return t, b[i : i + ln], i + ln

    try:
        _, body, _ = read_tlv(data, 0)          # SEQUENCE
        _, _, i = read_tlv(body, 0)             # version
        _, _, i = read_tlv(body, i)             # community
        t, pdu_b, _ = read_tlv(body, i)         # GetResponse(0xA2)
        _, _, j = read_tlv(pdu_b, 0)            # request-id
        _, est, j = read_tlv(pdu_b, j)          # error-status
        if int.from_bytes(est or b"\x00", "big"):
            return False, f"SNMP 오류 (error-status {int.from_bytes(est, 'big')}) — community 권한을 확인하세요"
        _, _, j = read_tlv(pdu_b, j)            # error-index
        _, vbl, _ = read_tlv(pdu_b, j)          # varbind list
        _, vb1, _ = read_tlv(vbl, 0)
        _, _, k = read_tlv(vb1, 0)              # oid
        vt, _, _ = read_tlv(vb1, k)             # value
        if vt in (0x80, 0x81, 0x82):            # noSuchObject 류
            return False, "장비가 sysDescr 를 주지 않습니다"
        return True, ""
    except Exception:
        return True, ""


def _snmp_bulk_sync(host: str, port: int, community: str,
                    roots: list, max_rep: int = 60) -> dict:
    """SNMPv2c GetBulk 로 서브트리를 읽는다 — 포트 상태(ifOperStatus) 몫.

    의존성 없이 최소 BER. roots 의 각 서브트리에 대해 {끝자리 index: 값} 을
    돌려준다. 값은 INTEGER 면 int, 아니면 bytes 그대로.
    """
    import os
    import socket

    def tlv(t: int, v: bytes) -> bytes:
        n = len(v)
        if n < 0x80:
            return bytes([t, n]) + v
        eb = n.to_bytes((n.bit_length() + 7) // 8, "big")
        return bytes([t, 0x80 | len(eb)]) + eb + v

    def ber_int(n: int) -> bytes:
        b = n.to_bytes((max(n.bit_length(), 1) + 8) // 8, "big", signed=True)
        return tlv(0x02, b)

    def oid_enc(parts: tuple) -> bytes:
        out = [40 * parts[0] + parts[1]]
        for x in parts[2:]:
            if x < 0x80:
                out.append(x)
            else:
                stack = [x & 0x7F]
                x >>= 7
                while x:
                    stack.append((x & 0x7F) | 0x80)
                    x >>= 7
                out.extend(reversed(stack))
        return tlv(0x06, bytes(out))

    def oid_dec(b: bytes) -> tuple:
        if not b:
            return ()
        out = [b[0] // 40, b[0] % 40]
        val = 0
        for c in b[1:]:
            val = (val << 7) | (c & 0x7F)
            if not c & 0x80:
                out.append(val)
                val = 0
        return tuple(out)

    def read_tlv(b: bytes, i: int):
        t = b[i]
        ln = b[i + 1]
        i += 2
        if ln & 0x80:
            k = ln & 0x7F
            ln = int.from_bytes(b[i : i + k], "big")
            i += k
        return t, b[i : i + ln], i + ln

    result: dict = {tuple(r): {} for r in roots}
    cursors = [tuple(r) for r in roots]
    sock = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
    sock.settimeout(_PROBE_TIMEOUT)
    try:
        for _round in range(20):  # 48포트 × 2열이면 한두 번에 끝난다 — 폭주 방지 상한
            rid = int.from_bytes(os.urandom(2), "big") & 0x7FFF
            vbs = b"".join(tlv(0x30, oid_enc(c) + b"\x05\x00") for c in cursors)
            pdu = tlv(0xA5, ber_int(rid) + ber_int(0) + ber_int(max_rep) + tlv(0x30, vbs))
            msg = tlv(0x30, ber_int(1) + tlv(0x04, community.encode()) + pdu)
            sock.sendto(msg, (host, int(port or 161)))
            data, _ = sock.recvfrom(65535)
            _, body, _ = read_tlv(data, 0)
            _, _, i = read_tlv(body, 0)
            _, _, i = read_tlv(body, i)
            _, pdu_b, _ = read_tlv(body, i)
            _, _, j = read_tlv(pdu_b, 0)
            _, est, j = read_tlv(pdu_b, j)
            if int.from_bytes(est or b"\x00", "big"):
                break
            _, _, j = read_tlv(pdu_b, j)
            _, vbl, _ = read_tlv(pdu_b, j)
            k = 0
            n_roots = len(cursors)
            col = 0
            done = [False] * n_roots
            last = list(cursors)
            while k < len(vbl):
                _, vb1, k = read_tlv(vbl, k)
                to, ob, m = read_tlv(vb1, 0)
                vt, vv, _ = read_tlv(vb1, m)
                oid = oid_dec(ob)
                root = tuple(roots[col % n_roots])
                if oid[: len(root)] == root and vt not in (0x82,):  # endOfMibView 제외
                    idx = oid[len(root) :]
                    val = int.from_bytes(vv, "big", signed=True) if vt == 0x02 else vv
                    result[root][idx[-1] if len(idx) == 1 else idx] = val
                    last[col % n_roots] = oid
                else:
                    done[col % n_roots] = True
                col += 1
            cursors = last
            if all(done) or col == 0:
                break
    except OSError:
        pass
    finally:
        sock.close()
    return result


# 장비마다 20초 캐시 — 랙뷰 카드가 뜰 때마다 장비를 두드리지 않게
_SNMP_PORTS_CACHE: dict = {}


@app.get("/api/devices2/{dev_id}/snmp-ports")
async def devices2_snmp_ports(dev_id: str):
    """포트 형상 실측 — SNMP(ifDescr·ifOperStatus)로 링크 up/down 을 읽는다."""
    d = await db.device_get(dev_id)
    if d is None:
        raise HTTPException(404, "장비를 찾을 수 없습니다")
    a = next((x for x in (d.get("access") or [])
              if x.get("protocol") == "snmp" and x.get("enabled", True)), None)
    if a is None:
        return {"ok": False, "reason": "SNMP 미등록"}
    import time as _time
    ent = _SNMP_PORTS_CACHE.get(d["id"])
    if ent and _time.time() - ent[0] < 20:
        return ent[1]
    host = (a.get("host") or d.get("ip") or "").strip()
    comm = (a.get("username") or "").strip() or "public"
    import asyncio as _aio
    loop = _aio.get_running_loop()
    res = await loop.run_in_executor(
        None, _snmp_bulk_sync, host, a.get("port") or 161, comm,
        [(1, 3, 6, 1, 2, 1, 2, 2, 1, 2), (1, 3, 6, 1, 2, 1, 2, 2, 1, 8)],
    )
    names = res.get((1, 3, 6, 1, 2, 1, 2, 2, 1, 2), {})
    stats = res.get((1, 3, 6, 1, 2, 1, 2, 2, 1, 8), {})
    # 물리 포트와 VLAN 을 가른다 — ifDescr 에는 mgmt·port-channel·CPU 도
    # 섞여 온다. 실물 포트는 예외 없이 슬롯/포트(Giga0/1) 꼴이라 '/' 가
    # 곧 물리의 표식이다. vlan 은 VLAN 정보로, 나머지 논리들은 뺀다.
    ports, vlans, others = [], [], 0
    for idx in names.keys():
        nm = names[idx]
        nm = nm.decode("utf-8", "replace") if isinstance(nm, (bytes, bytearray)) else str(nm)
        low = nm.lower()
        st = stats.get(idx)
        row = {"name": nm, "up": st == 1}
        if "vlan" in low or low.startswith("br"):
            vlans.append(row)
        elif "/" in nm:
            ports.append(row)
        else:
            others += 1
    # '/' 없는 장비(드물다)면 물리 표식이 안 통한 것 — 다 보여주는 쪽이 낫다
    if not ports and others:
        for idx in names.keys():
            nm = names[idx]
            nm = nm.decode("utf-8", "replace") if isinstance(nm, (bytes, bytearray)) else str(nm)
            if "vlan" in nm.lower() or nm.lower().startswith("br"):
                continue
            ports.append({"name": nm, "up": stats.get(idx) == 1})
    # ifIndex 차례는 포트 번호 차례가 아니다 — 이름을 자연 정렬한다
    # (Giga0/2 < Giga0/10 이 되게 숫자 덩어리는 숫자로 비교)
    import re as _re

    def _natkey(nm: str):
        return [(0, int(t)) if t.isdigit() else (1, t.lower())
                for t in _re.split(r"(\d+)", nm) if t]

    ports.sort(key=lambda x: _natkey(x["name"]))
    vlans.sort(key=lambda x: _natkey(x["name"]))
    out = {"ok": len(ports) + len(vlans) > 0, "ports": ports, "vlans": vlans,
           "reason": "" if ports or vlans else "SNMP 응답 없음"}
    _SNMP_PORTS_CACHE[d["id"]] = (_time.time(), out)
    return out


@app.post("/api/devices2/{dev_id}/check")
async def devices2_check(dev_id: str, protocol: str = ""):
    """접속해 보고 결과를 남긴다.

    protocol 을 주면 그것 하나만 — 목록에서 Telnet 칸만 눌러 확인하는 경우다.
    비우면 등록된 방식 전부."""
    d = await db.device_get(dev_id)
    if d is None:
        raise HTTPException(404, "장비를 찾을 수 없습니다")

    import asyncio
    loop = asyncio.get_running_loop()
    out = []
    want = (protocol or "").strip().lower()
    for a in d.get("access") or []:
        if not a.get("enabled", True):
            continue
        proto = a["protocol"]
        if want and proto != want:
            continue
        host = (a.get("host") or d.get("ip") or "").strip()
        # N2X 는 TCP 포트가 없다. 소켓으로 찔러 보는 대신, 중계로 ping 을
        # 보내 실제 섀시 세션이 열리는지 본다. STC 는 REST 라 그 쪽으로.
        if proto == "n2x":
            # 검사는 세션 하나를 **공유**한다(label 고정). 계측기마다 새
            # 세션을 열면 N2X 세션 한도("maximum sessions running")를 금방
            # 넘긴다 — 섀시는 하나여도 UTOP 이 7대를 각각 열려 하기 때문.
            r = await loop.run_in_executor(None, _n2x_send, host, "utop", "ping")
            ok = bool(isinstance(r, dict) and r.get("ok"))
            err = "" if ok else str((r or {}).get("error") or "N2X 응답 없음")
        elif proto == "stc":
            # STC 는 REST 서버(host)에 붙고, 그 서버가 섀시(장비 ip)로 연결한다.
            # 소켓만 찔러 보면 REST 서버가 살아있는지만 알지 섀시까지는 모른다.
            # 실제 섀시 인벤토리를 읽어 본다.
            r = await stc_conncheck({
                "chassis": (d.get("ip") or "").strip(),
                "restIp": host or "localhost",
                "restPort": a.get("port") or 8888,
            })
            ok = bool(isinstance(r, dict) and r.get("ok"))
            err = "" if ok else str((r or {}).get("error") or "STC 응답 없음")
        elif proto == "snmp":
            # community 는 SNMP 줄의 계정 칸 — 비우면 public.
            # 공용 계정(root 따위)으로 폴백하면 안 된다 — 그건 CLI 로그인
            # 계정이지 community 가 아니라서, 장비가 침묵해 「응답 없음」 이
            # 됐다(겪었다).
            comm = (a.get("username") or "").strip() or "public"
            ok, err = await loop.run_in_executor(
                None, _snmp_probe_sync, host, a.get("port") or 161, comm
            )
        else:
            ok, err = await loop.run_in_executor(
                None, _probe_sync, proto, host, a.get("port") or 0
            )
        await db.device_access_mark(d["id"], proto, ok, err)
        out.append({"protocol": proto, "host": host, "port": a.get("port"),
                    "ok": ok, "error": err})
    return {"success": True, "results": out}


@app.get("/api/devices2")
async def devices2_list(ifs: int = 1):
    """
    `ifs=0` 이면 **인터페이스 줄을 싣지 않고 개수만** 준다(성능).

    목록 화면은 인터페이스를 「48」 처럼 수로만 쓰는데, 장비 92대 × 48줄이면
    4천 줄이 브라우저로 넘어가 첫 화면이 무겁다(지적).
    """
    devs = await db.device_list(with_ifs=bool(ifs))
    if not ifs:
        async with db.pool().acquire() as c:
            rows = await c.fetch(
                "SELECT device_id, name FROM device_interface ORDER BY device_id, sort_order, name"
            )
        by: dict = {}
        for r in rows:
            by.setdefault(r["device_id"], []).append(r["name"])
        for d in devs:
            names = by.get(d["id"], [])
            d["if_count"] = len(names)
            # 「gi1/0/1-48, te1/1-4」 — 표가 구성을 그대로 보여 준다(지시).
            # 48줄을 다 실으면 목록이 무거워지므로 여기서 접어 보낸다.
            d["if_brief"] = _compress_ifs(names) if names else ""
    return {"devices": devs}


# 이 라우트는 /api/devices2/{dev_id} 보다 먼저 선언되어야 한다.
# 뒤에 두면 'export.csv' 가 dev_id 로 잡혀 404 가 난다.
@app.get("/api/devices2/export.csv")
async def devices2_export(with_secrets: int = 0):
    """장비 목록을 CSV 로. 비밀번호는 기본적으로 비운다.

    평문 저장은 결정된 사항이지만, CSV 는 메일과 메신저로 쉽게 돌아다닌다.
    파일 하나가 사내 장비 전체의 비밀번호가 되는 것은 저장과 다른 문제다.
    가져올 때 빈 칸은 '기존 값 유지' 로 처리하므로 이대로도 왕복이 된다.
    """
    import csv, io
    devs = await db.device_list()
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(DEV_CSV_COLS)
    for d in devs:
        tel, ssh = _acc_of(d, "telnet"), _acc_of(d, "ssh")
        con, snmp = _acc_of(d, "console"), _acc_of(d, "snmp")
        w.writerow([
            d.get("lab") or "",
            d.get("ip") or "", d.get("operator") or "", d.get("vendor") or "",
            d.get("role") or "", d.get("model") or "",
            tel.get("port") or "", ssh.get("port") or "",
            con.get("host") or "", con.get("port") or "",
            (snmp.get("username") or snmp.get("community") or ""),
            ((snmp.get("params") or {}).get("community_rw") or ""),
            d.get("username") or "",
            (d.get("password") or "") if with_secrets else "",
            _compress_ifs([i["name"] for i in d.get("interfaces") or []]),
        ])
    # 엑셀이 UTF-8 을 알아보게 BOM 을 붙인다. 없으면 한글이 깨져서 열린다.
    data = "﻿" + buf.getvalue()
    from fastapi.responses import Response
    return Response(
        content=data.encode("utf-8"),
        media_type="text/csv; charset=utf-8",
        headers={"Content-Disposition": 'attachment; filename="devices.csv"'},
    )


@app.get("/api/devices2/{dev_id}")
async def devices2_get(dev_id: str):
    d = await db.device_get(dev_id)
    if d is None:
        raise HTTPException(404, "장비를 찾을 수 없습니다")
    return d


@app.post("/api/devices2")
async def devices2_save(payload: dict):
    ip = str(payload.get("ip") or "").strip()
    if not ip:
        raise HTTPException(400, "IP 를 입력하세요")
    # 같은 IP 가 이미 있으면 그 장비를 고치는 것으로 본다. IP 가 키다.
    cur = await db.device_get(ip)
    if cur and not payload.get("id"):
        payload["id"] = cur["id"]
    try:
        dev_id = await db.device_upsert(payload)
    except ValueError as e:
        raise HTTPException(400, str(e)) from e
    return {"success": True, "id": dev_id}


@app.post("/api/devices2/{dev_id}/rack")
async def devices2_set_rack(dev_id: str, payload: dict):
    """랙 자리 지정/해제 — 랙뷰에서 끌어다 놓거나 뺀다. rack_id 비우면 해제.

    겹침은 서버가 최종 판정한다(409). 화면 검사만 믿으면 두 사람이 같은
    칸에 동시에 끌어다 놓았을 때 늦게 온 쪽이 조용히 겹쳐 앉는다.
    """
    rid = str(payload.get("rack_id") or "").strip()
    if rid:
        try:
            pos = int(payload.get("rack_pos") or 0)
            units = max(1, int(payload.get("rack_units") or 1))
        except (TypeError, ValueError):
            raise HTTPException(400, "자리(U)가 숫자가 아닙니다")
        if pos < 1:
            raise HTTPException(400, "자리(U)가 필요합니다")
        kv = _kv_load_sync("racks", {}) or {}
        rk = next((r for r in (kv.get("racks") or []) if str(r.get("id")) == rid), None)
        top = int((rk or {}).get("units") or 45)
        if pos + units - 1 > top:
            raise HTTPException(409, f"{top}U 랙 위를 벗어납니다")
        used: set = set()
        for d in await db.device_list(with_ifs=False):
            if str(d.get("rack_id") or "") != rid or not d.get("rack_pos"):
                continue
            if d["id"] == dev_id or d.get("ip") == dev_id:
                continue  # 자기 자신은 빼고 센다 — 제자리 이동·크기 변경 몫
            used.update(range(d["rack_pos"], d["rack_pos"] + (d.get("rack_units") or 1)))
        for b in kv.get("blanks") or []:
            brid = str(b.get("rack_id") or "")
            same = brid == rid or (not brid and rk and b.get("rack_name") == rk.get("name"))
            if not same:
                continue
            try:
                bp, bu = int(b.get("pos") or 0), int(b.get("units") or 1)
            except (TypeError, ValueError):
                continue
            used.update(range(bp, bp + bu))
        bad = sorted(u for u in range(pos, pos + units) if u in used)
        if bad:
            raise HTTPException(409, f"{bad[0]}U 가 이미 차 있습니다 — 다른 자리에 놓으세요")
    ok = await db.device_set_rack(
        dev_id, payload.get("rack_id"), payload.get("rack_pos"), payload.get("rack_units")
    )
    if not ok:
        raise HTTPException(404, "장비를 찾을 수 없습니다")
    return {"ok": True}


@app.post("/api/devices2/{dev_id}/default-protocol")
async def devices2_set_default_protocol(dev_id: str, payload: dict):
    """
    이 장비가 무엇으로 붙는지 바꾼다.

    전에는 장비 화면까지 가야 했다. 시험을 짜다가 「telnet 인데 왜 22번으로
    나가지」 를 알아차리는 자리는 세션 줄인데, 고치는 자리는 딴 데였다.
    """
    proto = str((payload or {}).get("protocol") or "").strip().lower()
    try:
        ok = await db.device_access_set_default(dev_id, proto)
    except ValueError as e:
        raise HTTPException(400, str(e)) from e
    if not ok:
        raise HTTPException(404, f"이 장비에 {proto} 접속이 등록되어 있지 않습니다")
    return {"success": True, "protocol": proto}


@app.delete("/api/devices2/{dev_id}")
async def devices2_delete(dev_id: str):
    if not await db.device_delete(dev_id):
        raise HTTPException(404, "장비를 찾을 수 없습니다")
    return {"success": True}


@app.post("/api/devices2/import-legacy")
async def devices2_import(request: Request):
    """옛 devices.json 을 PG 로 옮긴다. 여러 번 눌러도 안전하다(IP 기준 upsert)."""
    _me(request)
    src = DEVICES_FILE
    if not src.exists():
        raise HTTPException(404, f"{src.name} 이 없습니다")
    data = load_json(src) or {}
    # 랙 이름만 적힌 옛 장비 몫 — 이름을 KV 랙 id 로 풀어 준다
    _rk = _kv_load_sync("racks", {}) or {}
    rack_by_name = {str(r.get("name") or ""): str(r.get("id") or "") for r in (_rk.get("racks") or [])}
    n = 0
    for d in data.get("devices", []) or []:
        if not str(d.get("ip") or "").strip():
            continue
        payload = {
            "id": d.get("id") or d.get("ip"),
            "ip": d.get("ip"),
            "name": d.get("id") or d.get("ip"),
            "model": d.get("model"),
            "device_group": d.get("group"),
            "lab": d.get("lab"),
            "protocol": d.get("protocol"),
            "port": d.get("port"),
            "username": d.get("username"),
            "password": d.get("password"),
            "description": d.get("description"),
            "status": d.get("status"),
        }
        rid = d.get("rack_id") or rack_by_name.get(str(d.get("rack_name") or ""))
        if rid and d.get("rack_pos"):
            payload["rack_id"] = rid
            payload["rack_pos"] = d.get("rack_pos")
            payload["rack_units"] = d.get("rack_units") or 1
        await db.device_upsert(payload)
        n += 1
    return {"success": True, "imported": n}


# ── 장비 CSV 일괄 처리 ────────────────────────────────────────────
# 장비 30대를 창 하나씩 열어 등록하는 것은 현실적이지 않다. 내보내고,
# 엑셀에서 고치고, 다시 넣는 왕복 하나로 일괄등록·수정을 함께 해결한다.
DEV_CSV_COLS = [
    "LAB", "IP", "사업자", "제조사", "제품군", "모델명",
    "telnet포트", "ssh포트", "console주소", "console포트", "snmp", "snmp_rw",
    "계정", "비밀번호", "인터페이스",
]


def _expand_ifs(text: str) -> list[str]:
    """gi1/0/1-48 을 펼친다. 화면의 입력 규칙과 같아야 한다."""
    import re
    out: list[str] = []
    for raw in re.split(r"[,\n]", text or ""):
        s = raw.strip()
        if not s:
            continue
        # te6/1~te6/8 처럼 앞자리를 되풀이해 적은 것도 받는다. 기존 자료가
        # 물결로 들어와 있어서 '-' 만 알면 48포트가 1개로 세어진다.
        m = re.match(r"^(.*?)(\d+)\s*[-~]\s*(?:)?(\d+)$", s)
        if not m:
            out.append(s)
            continue
        prefix, a, b = m.group(1), int(m.group(2)), int(m.group(3))
        if b < a or b - a > 512:
            out.append(s)
            continue
        out.extend(f"{prefix}{i}" for i in range(a, b + 1))
    return out


def _compress_ifs(names: list[str]) -> str:
    """펼친 포트를 범위로 다시 접는다.

    48포트를 그대로 내보내면 한 칸이 화면을 넘어가 엑셀에서 손댈 수가 없다.
    _expand_ifs 의 역이라 내보내고 다시 넣어도 같은 결과가 나온다.
    """
    import re
    out: list[str] = []
    st: dict = {"pre": None, "from": 0, "to": 0, "width": 1}

    def flush():
        if st["pre"] is None:
            return
        if st["from"] == st["to"]:
            out.append(f"{st['pre']}{str(st['from']).zfill(st['width'])}")
        else:
            out.append(f"{st['pre']}{st['from']}-{st['to']}")

    for nm in names:
        m = re.match(r"^(.*?)(\d+)$", nm)
        if not m:
            flush()
            st["pre"] = None
            out.append(nm)
            continue
        pre, digits = m.group(1), m.group(2)
        num, width = int(digits), len(digits)
        # 0 으로 채운 이름(gi1/0/01)은 접으면 자릿수가 사라진다. 그대로 둔다.
        padded = width > 1 and digits[0] == "0"
        if st["pre"] == pre and num == st["to"] + 1 and not padded:
            st["to"] = num
            continue
        flush()
        st.update({"pre": pre, "from": num, "to": num, "width": width})
    flush()
    return ",".join(out)


def _acc_of(d: dict, proto: str) -> dict:
    for a in d.get("access") or []:
        if a.get("protocol") == proto:
            return a
    return {}


@app.post("/api/devices2/import-csv")
async def devices2_import_csv(payload: dict):
    """CSV 로 일괄 등록·수정. IP 가 키라 같은 IP 는 덮어쓴다.

    dry_run=true 면 저장하지 않고 무엇이 바뀌는지만 돌려준다 — 30줄을 넣기
    전에 확인할 수 있어야 한다.
    """
    import csv, io
    text = str(payload.get("csv") or "").lstrip("﻿")
    if not text.strip():
        raise HTTPException(400, "CSV 내용이 비어 있습니다")
    dry = bool(payload.get("dry_run"))

    rows = list(csv.DictReader(io.StringIO(text)))
    if not rows:
        raise HTTPException(400, "머리글만 있고 자료가 없습니다")

    missing = [c for c in ("IP",) if c not in (rows[0].keys() or [])]
    if missing:
        raise HTTPException(400, f"필수 열이 없습니다: {', '.join(missing)}. 내보내기 파일의 머리글을 그대로 쓰세요")

    created, updated, errors = [], [], []
    for n, r in enumerate(rows, start=2):   # 2 = 머리글 다음 줄
        ip = (r.get("IP") or "").strip()
        if not ip:
            errors.append(f"{n}행: IP 가 비어 있습니다")
            continue
        cur = await db.device_get(ip)

        def pick(key: str, old):
            """빈 칸은 기존 값 유지. 비밀번호를 안 내보내도 왕복이 되게 한다."""
            v = (r.get(key) or "").strip()
            return v if v else (old or None)

        def num(key: str, old, dflt=None):
            v = (r.get(key) or "").strip()
            if not v:
                return old if old is not None else dflt
            try:
                return int(v)
            except ValueError:
                errors.append(f"{n}행: {key} 가 숫자가 아닙니다 ({v})")
                return old if old is not None else dflt

        access = []
        tel_old, ssh_old = _acc_of(cur or {}, "telnet"), _acc_of(cur or {}, "ssh")
        con_old, snmp_old = _acc_of(cur or {}, "console"), _acc_of(cur or {}, "snmp")

        tp = num("telnet포트", tel_old.get("port"))
        if tp:
            access.append({"protocol": "telnet", "port": tp, "enabled": True,
                           "is_default": True})
        sp = num("ssh포트", ssh_old.get("port"))
        if sp:
            access.append({"protocol": "ssh", "port": sp, "enabled": True,
                           "is_default": not tp})
        ch, cp = pick("console주소", con_old.get("host")), num("console포트", con_old.get("port"))
        if ch or cp:
            access.append({"protocol": "console", "host": ch, "port": cp, "enabled": True})
        comm = pick("snmp", snmp_old.get("username") or snmp_old.get("community"))
        rw = pick("snmp_rw", (snmp_old.get("params") or {}).get("community_rw"))
        if comm or rw:
            access.append({"protocol": "snmp", "port": snmp_old.get("port") or 161,
                           # RO 는 읽는 쪽이 보는 칸(username)에 적는다
                           "username": comm or None, "community": comm or None,
                           "params": {**(snmp_old.get("params") or {}),
                                      **({"community_rw": rw, "rw": True} if rw else {})},
                           "enabled": True})

        if_text = (r.get("인터페이스") or "").strip()
        payload_dev = {
            "id": (cur or {}).get("id") or ip,
            "ip": ip,
            "lab": pick("LAB", (cur or {}).get("lab")),
            "operator": pick("사업자", (cur or {}).get("operator")),
            "vendor": pick("제조사", (cur or {}).get("vendor")),
            "role": pick("제품군", (cur or {}).get("role")),
            "model": pick("모델명", (cur or {}).get("model")),
            "username": pick("계정", (cur or {}).get("username")),
            "password": pick("비밀번호", (cur or {}).get("password")),
            "access": access,
        }
        # 인터페이스 칸이 비면 건드리지 않는다. 빈 칸을 '전부 삭제' 로 읽으면
        # 내보내기에서 지우고 올린 사람이 48포트를 통째로 잃는다.
        if if_text:
            payload_dev["interfaces"] = [
                {"name": nm, "kind": "general"} for nm in _expand_ifs(if_text)
            ]

        if dry:
            (updated if cur else created).append(
                {"ip": ip, "name": payload_dev.get("model"),
                 "interfaces": len(payload_dev.get("interfaces") or []) or None,
                 "access": [a["protocol"] for a in access]}
            )
            continue
        try:
            await db.device_upsert(payload_dev)
            (updated if cur else created).append({"ip": ip, "name": payload_dev.get("model")})
        except Exception as e:
            errors.append(f"{n}행 ({ip}): {e}")

    return {"success": not errors, "dry_run": dry,
            "created": created, "updated": updated, "errors": errors}


@app.get("/api/locks")
async def list_locks():
    """지금 잡혀 있는 자원 전부. 화면이 '누가 언제부터' 를 보여줄 수 있게
    신호가 끊긴 지 얼마나 됐는지도 함께 준다."""
    async with db.pool().acquire() as c:
        rows = await c.fetch(
            """
            SELECT resource_id, kind, locked_by, locked_name, cycle_id, note,
                   locked_at, heartbeat_at,
                   EXTRACT(EPOCH FROM (now() - heartbeat_at))::int AS stale_sec
            FROM resource_lock ORDER BY locked_at
            """
        )
        # 「어느 사이클에서 쓰는 중인가」 — id 만으로는 사람이 못 읽는다(지시)
        ids = [r["cycle_id"] for r in rows if r["cycle_id"]]
        nm: dict = {}
        if ids:
            cn = await c.fetch(
                "SELECT id, name, data_summary FROM cycle WHERE id = ANY($1::text[])", ids
            )
            for r2 in cn:
                d2 = dict(r2["data_summary"] or {})
                nm[r2["id"]] = {"name": r2["name"], "cid": d2.get("cid") or ""}
    out = []
    for r in rows:
        d = dict(r)
        info = nm.get(d.get("cycle_id") or "")
        d["cycle_name"] = (info or {}).get("name") or ""
        d["cycle_cid"] = (info or {}).get("cid") or ""
        out.append(d)
    return {"locks": out}


class LockBulkIn(BaseModel):
    resource_ids: list[str] = []
    kind: str = "device"
    cycle_id: Optional[str] = None
    note: Optional[str] = None


@app.post("/api/locks/bulk")
async def acquire_locks_bulk(body: LockBulkIn, request: Request):
    """
    사이클 실행이 거는 자동 점유(지시). 걸려는 장비 가운데 **남이 잡은 것이
    하나라도 있으면 아무것도 잡지 않고 물러난다** — 반쯤 잡힌 채로 실패하면
    남의 자리만 붙들고 있게 된다.

    막힌 자리는 「누가 · 어느 사이클에서」 까지 돌려준다.
    """
    me = _me(request)
    ids = [str(x).strip() for x in (body.resource_ids or []) if str(x).strip()]
    if not ids:
        return {"success": True, "locked": [], "blocked": []}

    async with db.pool().acquire() as c:
        cur = await c.fetch(
            "SELECT * FROM resource_lock WHERE resource_id = ANY($1::text[])", ids
        )
        mine_name = me.get("username")
        blocked = [dict(r) for r in cur if r["locked_by"] != mine_name]
        if blocked:
            cids = [b["cycle_id"] for b in blocked if b["cycle_id"]]
            nm: dict = {}
            if cids:
                cn = await c.fetch("SELECT id, name FROM cycle WHERE id = ANY($1::text[])", cids)
                nm = {r["id"]: r["name"] for r in cn}
            for b in blocked:
                b["cycle_name"] = nm.get(b.get("cycle_id") or "") or ""
                b["locked_at"] = b["locked_at"].isoformat() if b.get("locked_at") else None
                b["heartbeat_at"] = b["heartbeat_at"].isoformat() if b.get("heartbeat_at") else None
            return {"success": False, "locked": [], "blocked": blocked}

        held = {r["resource_id"] for r in cur}
        for rid in ids:
            if rid in held:
                await c.execute(
                    "UPDATE resource_lock SET heartbeat_at=now(), cycle_id=COALESCE($2, cycle_id) "
                    "WHERE resource_id=$1",
                    rid, body.cycle_id,
                )
                continue
            await c.execute(
                """INSERT INTO resource_lock
                   (resource_id, kind, locked_by, locked_name, cycle_id, note)
                   VALUES ($1,$2,$3,$4,$5,$6)""",
                rid, body.kind, mine_name, me.get("name") or mine_name,
                body.cycle_id, body.note,
            )
    return {"success": True, "locked": ids, "blocked": []}


@app.post("/api/locks")
async def acquire_lock(body: LockIn, request: Request):
    rid = (body.resource_id or "").strip()
    if not rid:
        raise HTTPException(400, "자원 id 가 필요합니다")
    me = _me(request)

    async with db.pool().acquire() as c:
        cur = await c.fetchrow("SELECT * FROM resource_lock WHERE resource_id=$1", rid)
        if cur:
            if cur["locked_by"] != me.get("username"):
                who = cur["locked_name"] or cur["locked_by"]
                since = cur["locked_at"].strftime("%m-%d %H:%M") if cur["locked_at"] else ""
                raise HTTPException(
                    409, f"이미 {who} 님이 잡고 있습니다 (시작 {since}). 확인 후 진행하세요."
                )
            # 내가 이미 잡고 있으면 신호만 갱신한다
            await c.execute(
                "UPDATE resource_lock SET heartbeat_at=now(), cycle_id=COALESCE($2, cycle_id) WHERE resource_id=$1",
                rid, body.cycle_id,
            )
            return {"success": True, "renewed": True}

        await c.execute(
            """INSERT INTO resource_lock
               (resource_id, kind, locked_by, locked_name, cycle_id, note)
               VALUES ($1,$2,$3,$4,$5,$6)""",
            rid, body.kind, me.get("username"), me.get("name") or me.get("username"),
            body.cycle_id, body.note,
        )
    return {"success": True, "renewed": False}


@app.post("/api/locks/{resource_id}/heartbeat")
async def heartbeat_lock(resource_id: str, request: Request):
    me = _me(request)
    async with db.pool().acquire() as c:
        r = await c.execute(
            "UPDATE resource_lock SET heartbeat_at=now() WHERE resource_id=$1 AND locked_by=$2",
            resource_id, me.get("username"),
        )
    if not r.endswith(" 1"):
        raise HTTPException(404, "내가 잡고 있는 자원이 아닙니다")
    return {"success": True}


@app.delete("/api/locks/{resource_id}")
async def release_lock(resource_id: str, request: Request):
    """해제는 잡은 본인과 관리자만. 남의 시험을 아무나 끊을 수 없어야 한다."""
    me = _me(request)
    async with db.pool().acquire() as c:
        cur = await c.fetchrow("SELECT * FROM resource_lock WHERE resource_id=$1", resource_id)
        if not cur:
            raise HTTPException(404, "잡혀 있지 않습니다")
        is_admin = me.get("role") == "관리자"
        if cur["locked_by"] != me.get("username") and not is_admin:
            who = cur["locked_name"] or cur["locked_by"]
            raise HTTPException(403, f"{who} 님이 잡은 자원입니다. 본인 또는 관리자만 해제할 수 있습니다.")
        await c.execute("DELETE FROM resource_lock WHERE resource_id=$1", resource_id)
    return {"success": True, "forced": cur["locked_by"] != me.get("username")}


@app.delete("/api/locks/by-cycle/{cycle_id}")
async def release_locks_of_cycle(cycle_id: str, request: Request):
    """사이클이 끝나면 그 사이클이 잡은 것을 한꺼번에 푼다."""
    _me(request)
    async with db.pool().acquire() as c:
        r = await c.execute("DELETE FROM resource_lock WHERE cycle_id=$1", cycle_id)
    return {"success": True, "released": int(r.rsplit(" ", 1)[-1] or 0)}


@app.post("/api/convert/markdown")
async def convert_to_markdown(file: UploadFile = File(...)):
    name = (file.filename or "").strip()
    raw = await file.read()
    if not raw:
        raise HTTPException(400, "빈 파일입니다")
    # 20MB 제한 — 이보다 큰 규격서는 통째로 넣기보다 나눠 올리는 게 낫다
    if len(raw) > 20 * 1024 * 1024:
        raise HTTPException(413, "20MB 이하 파일만 변환합니다")

    ext = Path(name).suffix.lower()
    if ext in (".md", ".markdown", ".txt"):
        return {"markdown": raw.decode("utf-8", "replace"), "source": name}

    try:
        from markitdown import MarkItDown
    except ImportError:
        raise HTTPException(
            501,
            "문서 변환 기능이 설치되지 않았습니다. requirements.txt 의 markitdown 을 "
            "설치한 뒤 서버를 다시 띄우세요.",
        )

    # markitdown 은 파일 경로를 받는다. 임시 파일로 떨군 뒤 지운다.
    import tempfile
    tmp = None
    try:
        with tempfile.NamedTemporaryFile(suffix=ext or ".bin", delete=False) as fh:
            fh.write(raw)
            tmp = fh.name
        md = MarkItDown(enable_plugins=False).convert(tmp).text_content or ""
    except Exception as e:
        raise HTTPException(422, f"변환하지 못했습니다: {e}") from e
    finally:
        if tmp:
            try:
                os.unlink(tmp)
            except OSError:
                pass

    return {"markdown": md, "source": name}


# ───────────────────────────────────────────
# 요구사항 구현내용 → 벡터 저장
#
# 마크다운을 제목(##) 단위로 자른다. 문단 길이로 자르면 '## 판정 기준' 의
# 표가 반토막 나서, 나중에 이 조각으로 시험항목을 만들 때 기준을 놓친다.
# ───────────────────────────────────────────
def _split_markdown(md: str, max_chars: int = 1200) -> list[str]:
    """제목 단위로 자르되, 한 절이 너무 길면 줄 단위로 더 쪼갠다."""
    lines = (md or "").splitlines()
    blocks: list[list[str]] = [[]]
    for ln in lines:
        if ln.lstrip().startswith("#") and blocks[-1]:
            blocks.append([])
        blocks[-1].append(ln)

    out: list[str] = []
    for b in blocks:
        text = "\n".join(b).strip()
        if not text:
            continue
        if len(text) <= max_chars:
            out.append(text)
            continue
        cur: list[str] = []
        size = 0
        for ln in b:
            if size + len(ln) > max_chars and cur:
                out.append("\n".join(cur).strip())
                cur, size = [], 0
            cur.append(ln)
            size += len(ln) + 1
        if cur:
            out.append("\n".join(cur).strip())
    return [c for c in out if c]


class ReqEmbedIn(BaseModel):
    text: str


@app.post("/api/req/{req_id}/embed")
async def embed_requirement(req_id: str, body: ReqEmbedIn):
    text = (body.text or "").strip()
    if not text:
        raise HTTPException(400, "구현내용이 비어 있습니다")

    chunks = _split_markdown(text)
    if not chunks:
        raise HTTPException(400, "저장할 내용이 없습니다")

    vecs = await _embed_texts(chunks)
    if not vecs:
        raise HTTPException(
            503,
            "임베딩 서버가 설정되지 않았습니다. 시스템 > RAG 설정에서 embed_url 을 "
            "지정하세요. (구현내용은 요구사항에 이미 저장되어 있습니다)",
        )

    import numpy as _np
    async with db.pool().acquire() as c:
        async with c.transaction():
            # 같은 요구사항의 옛 조각은 지우고 새로 넣는다 — 수정본과 옛 본이
            # 같이 검색되면 어느 것이 맞는지 알 수 없다.
            await c.execute("DELETE FROM rag_embed WHERE key LIKE $1", f"req:{req_id}#%")
            for i, (chunk, vec) in enumerate(zip(chunks, vecs)):
                await c.execute(
                    "INSERT INTO rag_embed (key, embed, meta) VALUES ($1,$2,$3::jsonb)",
                    f"req:{req_id}#{i}",
                    _np.asarray(vec, dtype="float32").tobytes(),
                    {"req_id": req_id, "chunk": i, "text": chunk},
                )
    return {"success": True, "chunks": len(chunks)}


# REQ 목록 (전체)
@app.get("/api/req")
async def get_all_req():
    # 분류는 req_category 테이블로 넘어갔다. 옛 화면은 /api/folders 를
    # 따로 부르므로 여기서 folders 를 함께 실어 보낼 이유가 없다.
    return {"reqs": await db.req_list_full()}

# REQ 단건 조회
@app.get("/api/req/{req_id}")
async def get_req(req_id: str):
    r = await db.req_get(req_id)
    if r is None:
        raise HTTPException(404, "REQ를 찾을 수 없습니다")
    # TC 는 참조(tcid/name/status) 만 반환 — 예전엔 각 TC 를 풀 데이터로 확장해 붙였으나
    # REQ 하나에 큰 TC 여러 개면 응답 수 MB 로 팽창, 프론트가 그대로 saveOneREQ 로 다시 POST 시
    # 서버가 각 tc 를 풀 데이터로 오판정해 tc_upsert 를 반복 실행 → 10초+ 지연 발생.
    # 개별 TC 상세는 /api/tc/{tcid} 로 lazy 로드 (loadTCFull) 로 이미 처리됨.
    refs = []
    for tc in r.get("tc", []) or []:
        if isinstance(tc, dict) and tc.get("tcid"):
            refs.append({"tcid": tc.get("tcid",""), "name": tc.get("name",""), "status": tc.get("status","대기")})
    r["tc"] = refs
    return r

# REQ 저장 (생성/수정)
# ══════════════ 폴더·요구사항·시험 통째로 복사 ══════════════
#
# 「+ Copy」 창 하나가 이 일을 다 한다(승인 2026-08-22). 파일로 내보냈다
# 가져오는 길은 「이 줄이 어느 요구사항에 붙나」 를 사람이 다시 정해 줘야
# 했다 — 붙일 자리를 먼저 고르고 옮기면 그 물음이 아예 사라진다.
#
# 규칙(승인): 새 ID 로 발번 · 요구사항↔시험 연결은 새 ID 끼리 유지 ·
# 대상 프로젝트의 모델그룹·모델명으로 갈아 끼움(끌 수 있음) · 같은 이름이면
# 「(복제)」 를 붙여 새로 만든다(덮어쓰기 없음) · 실행 이력은 안 가져온다.


def _cat_path(cats: dict, cid: str) -> list:
    """뿌리부터 이 폴더까지 — req 의 cat1..cat4 가 이 길을 담는다."""
    out, cur, guard = [], cid, 0
    while cur and guard < 12:
        out.append(cur)
        cur = (cats.get(cur) or {}).get("parent_id")
        guard += 1
    return list(reversed(out))


def _leaf_cat(r: dict) -> str:
    for k in ("cat4", "cat3", "cat2", "cat1"):
        v = str((r.get(k) or "")).strip()
        if v:
            return v
    return ""


async def _next_req_id(c) -> str:
    from datetime import datetime as _dt
    import re as _r
    iso = _dt.now().isocalendar()
    prefix = "REQ-%02d%02d-" % (iso[0] % 100, iso[1])
    rows = await c.fetch("SELECT data->>'reqid' AS reqid FROM req WHERE data->>'reqid' LIKE $1", prefix + "%")
    mx = 0
    for r in rows:
        m = _r.match("^" + _r.escape(prefix) + r"(\d+)$", r["reqid"] or "")
        if m:
            mx = max(mx, int(m.group(1)))
    return prefix + str(mx + 1).zfill(4)


async def _next_tc_id(c) -> str:
    from datetime import datetime as _dt
    import re as _r
    iso = _dt.now().isocalendar()
    prefix = "TC-%02d%02d-" % (iso[0] % 100, iso[1])
    rows = await c.fetch("SELECT tcid FROM tc WHERE tcid LIKE $1", prefix + "%")
    mx = 0
    for r in rows:
        m = _r.match("^" + _r.escape(prefix) + r"(\d+)$", r["tcid"] or "")
        if m:
            mx = max(mx, int(m.group(1)))
    return prefix + str(mx + 1).zfill(4)


@app.post("/api/copy-tree")
async def copy_tree(body: dict, token: str = ""):
    """Source 에서 고른 것들을 Destination 아래로 **복사**한다.

    items: [{kind: 'cat'|'req'|'tc', id}]   — 여럿 가능
    dst  : {kind: 'cat'|'req', id}          — 폴더나 요구사항
    swap_model: 대상 프로젝트의 모델그룹·모델명으로 갈아 끼울까(기본 켜짐)
    """
    _user_from_token(token)
    items = [x for x in (body.get("items") or []) if isinstance(x, dict)]
    dst = body.get("dst") or {}
    dst_kind = str(dst.get("kind") or "")
    dst_id = str(dst.get("id") or "")
    swap = body.get("swap_model") is not False
    # 무엇을 복사하나(승인): all=요구사항+시험 · req=요구사항만 · tc=시험만
    mode = str(body.get("mode") or "all")
    # 폴더·요구사항을 통째로 고르되 **뺄 시험**은 따로 온다(체크 해제한 것)
    skip_tc = {str(x) for x in (body.get("skip_tcs") or [])}
    if not items or not dst_id:
        raise HTTPException(400, "무엇을 어디로 복사할지 골라 주세요")

    cats = {c["id"]: dict(c) for c in await db.cat_list()}
    reqs = {r["id"]: dict(r) for r in await db.req_list_full()}
    tcs_meta = await db.tc_list_meta()
    tc_by_req: dict = {}
    for t in tcs_meta:
        tc_by_req.setdefault(str(t.get("req_id") or ""), []).append(str(t.get("tcid")))

    # 대상 프로젝트(뿌리 폴더)의 모델 정보 — 갈아 끼울 값
    base_cat = dst_id if dst_kind == "cat" else _leaf_cat(reqs.get(dst_id) or {})
    root = (_cat_path(cats, base_cat) or [base_cat])[0]
    async with db.pool().acquire() as c:
        prow = await c.fetchrow(
            "SELECT customer, model_group, model FROM project WHERE cat_id = $1", root
        )
    dst_mg = str((prow or {}).get("model_group") or "") if prow else ""
    dst_md = str((prow or {}).get("model") or "") if prow else ""

    made = {"cats": 0, "reqs": 0, "tcs": 0}
    now_ms = int(datetime.now().timestamp() * 1000)
    seq = {"n": 0}

    def _uid(p: str) -> str:
        seq["n"] += 1
        return f"{p}-{now_ms}-{seq['n']}"

    async def copy_tc(c, tcid: str, req_id: str) -> None:
        if mode == "req":
            return          # 「요구사항만」 — 시험은 따라가지 않는다
        if tcid in skip_tc:
            return          # 시험 칸에서 체크를 뺀 것
        src = await db.tc_get(tcid)
        if not src:
            return
        nid = await _next_tc_id(c)
        d = dict(src)
        d["tcid"] = nid
        d["req_id"] = req_id
        if swap:
            if dst_mg:
                d["model_group"] = dst_mg
            if dst_md:
                d["model"] = dst_md
        # 실행 흔적은 안 가져온다 — 복사본은 「아직 안 돌린 것」 이다(승인)
        for k in ("result_history", "issue_list", "last_run", "cycles"):
            d.pop(k, None)
        for st in d.get("checks") or []:
            for k in ("output", "status", "reason", "repeatResult", "executed_at", "took_ms", "rounds", "response"):
                st.pop(k, None)
        # 세션(장비 자리)은 **그대로 둔다**(지시). 장비가 여러 대면 어느 것을
        # 어느 자리에 앉힐지 기계가 정할 수 없다 — 사람이 고를 일이다.
        await db.tc_upsert(nid, d)
        made["tcs"] += 1

    async def copy_req(c, rid: str, cat_id: str) -> None:
        src = reqs.get(rid)
        if not src:
            return
        path = _cat_path(cats, cat_id)
        nid = _uid("rq")
        d = dict(src.get("data") or src)
        d["id"] = nid
        d["reqid"] = await _next_req_id(c)
        d["tc"] = []
        for i in range(4):
            d[f"cat{i + 1}"] = path[i] if i < len(path) else None
        if swap:
            if dst_mg:
                d["model_group"] = dst_mg
            if dst_md:
                d["model"] = dst_md
        await db.req_upsert(nid, d)
        made["reqs"] += 1
        for t in tc_by_req.get(rid, []):
            await copy_tc(c, t, nid)

    async def copy_cat(c, cid: str, parent: str) -> None:
        src = cats.get(cid)
        if not src:
            return
        name = str(src.get("name") or "")
        # 같은 이름이 이미 있으면 「(복제)」 — 덮어쓰지 않는다(승인)
        sibs = {str(v.get("name") or "") for v in cats.values() if str(v.get("parent_id") or "") == str(parent or "")}
        if name in sibs:
            name = f"{name} (복제)"
        nid = _uid("cat")
        await db.cat_upsert(nid, name, parent or None, int(src.get("sort_order") or 0))
        cats[nid] = {"id": nid, "name": name, "parent_id": parent}
        made["cats"] += 1
        for r in list(reqs.values()):
            if _leaf_cat(r) == cid:
                await copy_req(c, str(r["id"]), nid)
        for k, v in list(cats.items()):
            if str(v.get("parent_id") or "") == cid and not k.startswith("cat-" + str(now_ms)):
                await copy_cat(c, k, nid)

    async with db.pool().acquire() as c:
        for it in items:
            kind = str(it.get("kind") or "")
            sid = str(it.get("id") or "")
            if kind == "cat":
                if dst_kind != "cat":
                    raise HTTPException(400, "폴더는 폴더 아래로만 복사합니다")
                await copy_cat(c, sid, dst_id)
            elif kind == "req":
                if dst_kind != "cat":
                    raise HTTPException(400, "요구사항은 폴더 아래로만 복사합니다")
                await copy_req(c, sid, dst_id)
            elif kind == "tc":
                if dst_kind != "req":
                    raise HTTPException(400, "시험 항목은 요구사항 아래로만 복사합니다")
                await copy_tc(c, sid, dst_id)
    return {"ok": True, **made}


@app.get("/api/req-next-id")
async def req_next_id():
    """다음 요구사항 ID — REQ-<연2><ISO주차2>-<주차별 순번4>.

    2026년 ISO 15주차의 첫 건이면 REQ-2615-0001. 순번은 그 주차 안에서만
    센다(다음 주면 다시 0001). 그 프리픽스의 현재 최대 순번 +1 을 서버에서
    매겨, 두 사람이 같은 순간에 만들어도 안 겹친다(겹치면 save_req 가 한 번
    더 올려 준다)."""
    from datetime import datetime as _dt
    import re as _re_r
    iso = _dt.now().isocalendar()
    yy = iso[0] % 100
    ww = iso[1]
    prefix = "REQ-%02d%02d-" % (yy, ww)
    async with db.pool().acquire() as c:
        rows = await c.fetch(
            "SELECT data->>'reqid' AS reqid FROM req WHERE data->>'reqid' LIKE $1",
            prefix + "%",
        )
    mx = 0
    for r in rows:
        m = _re_r.match("^" + _re_r.escape(prefix) + r"(\d+)$", r["reqid"] or "")
        if m:
            mx = max(mx, int(m.group(1)))
    return {"reqid": prefix + str(mx + 1).zfill(4), "prefix": prefix, "seq": mx + 1}


@app.get("/api/tc-next-id")
async def tc_next_id():
    """다음 시험 ID — TC-<연2><ISO주차2>-<주차별 순번4>.

    요구사항(REQ-2632-0001)과 같은 규칙이라 나란히 읽힌다. tcid 는 곧 PK 라
    겹치면 남의 시험을 덮어쓴다 — 그래서 그 프리픽스의 현재 최대 순번 +1 을
    서버가 매긴다."""
    from datetime import datetime as _dt
    import re as _re_t
    iso = _dt.now().isocalendar()
    prefix = "TC-%02d%02d-" % (iso[0] % 100, iso[1])
    async with db.pool().acquire() as c:
        rows = await c.fetch("SELECT tcid FROM tc WHERE tcid LIKE $1", prefix + "%")
    mx = 0
    for r in rows:
        m = _re_t.match("^" + _re_t.escape(prefix) + r"(\d+)$", r["tcid"] or "")
        if m:
            mx = max(mx, int(m.group(1)))
    return {"tcid": prefix + str(mx + 1).zfill(4), "prefix": prefix, "seq": mx + 1}


@app.post("/api/req/{req_id}")
async def save_req(req_id: str, data: dict, response: Response):
    import time as _tm
    _t0 = _tm.perf_counter()
    # ★ URL id 와 data.id 를 통일 — 다르면 두 row 로 갈라져 REQ 중복 표시 버그 발생.
    #    URL 이 항상 진짜 PK. body 안 id 는 강제로 URL 값으로 덮어씀.
    data = {**data, "id": req_id}
    # ★ reqid 중복 방지 — 같은 폴더 내에 같은 reqid 를 가진 REQ 가 이미 있으면 자동으로 다음 번호 부여.
    _reqid = str(data.get("reqid") or "").strip()
    if _reqid:
        async with db.pool().acquire() as c:
            _dup = await c.fetch(
                "SELECT id FROM req WHERE data->>'reqid'=$1 AND id<>$2",
                _reqid, req_id,
            )
        if _dup:
            import re as _re_r
            _m = _re_r.match(r'^(.+-)(\d+)$', _reqid)
            if _m:
                _prefix = _m.group(1)
                # 이 prefix 로 존재하는 최대 번호 조회 후 +1
                async with db.pool().acquire() as c:
                    _rows = await c.fetch(
                        "SELECT data->>'reqid' AS reqid FROM req WHERE data->>'reqid' LIKE $1",
                        _prefix + '%',
                    )
                _max = 0
                for _r in _rows:
                    _mm = _re_r.match('^' + _re_r.escape(_prefix) + r'(\d+)$', _r['reqid'] or '')
                    if _mm:
                        _v = int(_mm.group(1))
                        if _v > _max:
                            _max = _v
                _new = _prefix + str(_max + 1).zfill(len(_m.group(2)))
                print(f"[save_req] reqid 중복 회피: {_reqid!r} -> {_new!r} (PK={req_id})")
                data["reqid"] = _new
    # TC 분리 저장 — 성능 최적화: 참조(ref only)만 온 것들은 존재 여부만 EXISTS 로 배치 확인
    tcs = data.get("tc", [])
    tc_refs = []
    REF_KEYS = {"tcid", "name", "status", "req_id"}
    # 1) full 데이터/ref only 분리
    full_items = []   # (tcid, tc)
    ref_items = []    # (tcid, tc)
    for tc in tcs:
        tcid = tc.get("tcid", "")
        if not tcid:
            tc_refs.append(tc)
            continue
        is_full = any(k not in REF_KEYS for k in tc.keys())
        if is_full:
            full_items.append((tcid, tc))
        else:
            ref_items.append((tcid, tc))
    # 2) full 데이터는 그대로 upsert
    for tcid, tc in full_items:
        await db.tc_upsert(tcid, {**tc, "req_id": data.get("id", "")})
        tc_refs.append({"tcid": tcid, "name": tc.get("name", ""), "status": tc.get("status", "대기")})
    # 3) ref only 는 EXISTS 로 한 번에 확인 (개당 tc_get 하면 REQ 아래 TC 수만큼 왕복)
    if ref_items:
        ref_ids = [tcid for tcid, _ in ref_items]
        async with db.pool().acquire() as c:
            rows = await c.fetch("SELECT tcid FROM tc WHERE tcid = ANY($1::text[])", ref_ids)
            existing_set = {r["tcid"] for r in rows}
        # 이름/상태만 바뀌는 경우 잘 안 일어남 → 존재하는 것만 참조 유지 (name/status 변경은 별도 저장 경로에서)
        for tcid, tc in ref_items:
            if tcid in existing_set:
                tc_refs.append({"tcid": tcid, "name": tc.get("name", ""), "status": tc.get("status", "대기")})
            # else: 잔여 참조 (DB 에 없는 TC) → 무시 (자동 재생성 방지)
    # REQ에는 TC 참조만 저장
    _t1 = _tm.perf_counter()
    req_data = {**data, "tc": tc_refs}
    await db.req_upsert(req_id, req_data)
    _t2 = _tm.perf_counter()
    # RAG 자동 색인(source=req) — 백그라운드, 실패해도 저장에 영향 없음
    try:
        if _ai_settings().get("auto_index_req"):
            asyncio.create_task(_rag_index_req(req_id, req_data))
    except Exception:
        pass
    try: asyncio.create_task(broadcast({"type": "req_updated", "req_id": req_id}))
    except Exception: pass
    _t3 = _tm.perf_counter()
    _msg = f"total={_t3-_t0:.3f}s tc_check={_t1-_t0:.3f}s req_upsert={_t2-_t1:.3f}s tail={_t3-_t2:.3f}s"
    try:
        response.headers["X-Save-Time"] = _msg
    except Exception: pass
    try:
        import sys as _sys
        print(f"[save_req] {req_id} {_msg}", flush=True)
        _sys.stdout.flush()
    except Exception: pass
    return {"success": True}

# REQ 삭제
@app.delete("/api/req/{req_id}")
async def delete_req(req_id: str):
    r = await db.req_get(req_id)
    _deleted_tcids = []
    if r:
        bundle = []
        for tc in r.get("tc", []):
            _tid = tc.get('tcid','')
            if not _tid: continue
            tcfull = await db.tc_get(_tid)
            if tcfull:
                bundle.append(tcfull)
                await db.tc_delete(_tid)
                _deleted_tcids.append(_tid)
        try: _trash_put("req", req_id, r, bundle)
        except Exception: pass
        await db.req_delete(req_id)
    try: asyncio.create_task(broadcast({"type": "req_deleted", "req_id": req_id, "tcids": _deleted_tcids}))
    except Exception: pass
    return {"success": True}

# TC 전체 목록 (REQ별)
TC_META_KEYS = ("tcid","id","name","req_id","folder","severity","priority","status","assignee","reporter","created_at","updated_at","created_by","updated_by","tags","custom","issue_list","result_history")

# 파일 단위 mtime 기반 캐시 — 변경된 파일만 다시 읽고 나머지는 in-memory 값 재사용.
# 이렇게 하면 요청마다 88개 파일을 다 파싱하지 않아 API 응답 시간이 크게 줄어듦.
_tc_cache = {}   # path_str → {"mtime":..., "full":dict, "meta":dict}
_cycle_cache = {}

def _cached_load(cache: dict, path: Path, meta_keys: tuple = None, meta_extra_fn=None):
    """파일 하나를 캐시에서 가져오거나 mtime 이 바뀌었으면 다시 읽음."""
    key = str(path)
    try:
        mt = path.stat().st_mtime
    except Exception:
        return None
    hit = cache.get(key)
    if hit and hit.get("mtime") == mt:
        return hit
    try:
        d = load_json(path)
    except Exception:
        return None
    entry = {"mtime": mt, "full": d}
    if meta_keys:
        meta = {k: d.get(k) for k in meta_keys if k in d}
        if meta_extra_fn:
            try: meta_extra_fn(meta, d)
            except Exception: pass
        entry["meta"] = meta
    cache[key] = entry
    return entry

def _tc_meta_extra(meta: dict, d: dict):
    _checks = d.get("checks") or []
    meta["_checks_count"] = len(_checks)
    # CLI 스텝(실제 시험 절차) 개수 — ⚠(세션 없음) 게이트가 쓴다
    meta["_cli_count"] = sum(1 for c in _checks if (c.get("kind") or "cli") == "cli")
    # 판정(PASS/FAIL)이 나오는 스텝 수(합의) — 목록 배지·Automation 탭이 같이 쓴다.
    # 주석·메시지·대기·치환·If/Loop/Switch 는 판정 단위가 아니라 안 센다.
    _judge = {"cli", "ping", "snmp_get", "snmp_set", "snmp_trap", "diff",
              "instrument", "connect", "disconnect", "auto"}
    meta["_step_count"] = sum(
        1 for c in _checks if (c.get("kind") or "cli") in _judge
    )

@app.get("/api/tc-last-result")
async def tc_last_result():
    """시험마다 **가장 최근 시험 결과** — 목록의 한 열(지시).

    결과는 사이클 안에 산다. 다만 `result` 칸만 보면 안 된다 — 자동 실행은
    항목 칸을 비워 두고 **스텝에만** 결과를 남긴다(사람이 손으로 찍을 때만
    항목 칸이 찬다). 그래서 돌려 놓고도 목록이 「–」 였다(지적: 제대로
    반영되고 있는 건가). 판정은 실행 화면과 같은 규칙(_item_verdict)으로
    스텝에서 유도한다.

    **가장 나중에 돌린 것**이 이긴다 — 사이클을 고친 시각이 아니라 항목을
    실행한 시각(`last_run`)으로 고른다.
    """
    try:
        async with db.pool().acquire() as c:
            rows = await c.fetch(
                """
                SELECT id, name, updated_at, data->'items' AS items
                FROM cycle
                ORDER BY updated_at DESC NULLS LAST
                """
            )
    except Exception as e:
        return {"items": {}, "error": str(e)[:200]}
    def _when(v) -> str:
        """시각 문자열을 견줄 수 있는 한 가지 꼴로. 자료에 `2026-06-29 14:53:44`
        와 `2026-06-29T14:53:44` 가 섞여 있어, 그대로 견주면 같은 순간인데도
        T 쪽이 늘 나중으로 읽힌다."""
        return str(v or "")[:19].replace("T", " ")

    best: dict = {}
    for r in rows:
        items = r["items"] if isinstance(r["items"], list) else []
        cyc_at = r["updated_at"].isoformat() if r["updated_at"] else ""
        for it in items:
            if not isinstance(it, dict):
                continue
            t = str(it.get("tcid") or "").strip()
            if not t:
                continue
            v = _item_verdict(it)
            if not v:
                continue        # 아직 안 돌린 항목은 「최근 결과」 가 아니다
            # **언제 돌렸나**로 고른다(지적: 실행했을 때 Pass 인데 안 바뀐다).
            # 여태는 「사이클을 마지막으로 고친 시각」 순으로 앞엣것을 썼다.
            # 그래서 같은 시험이 두 사이클에 들어 있으면, 어제 Pass 로 돌린
            # 것이 아니라 오늘 이름만 고친 사이클의 옛 Fail 이 이겼다. 말풍선의
            # 시각도 실행 시각이 아니라 사이클을 고친 시각이었다.
            at = _when(it.get("last_run")) or _when(cyc_at)
            prev = best.get(t)
            if prev and prev[0] >= at:
                continue
            # 화면 딱지는 설정(실행 판정 기준)의 값 이름을 쓴다 — PASS/FAIL
            # 대문자를 그 이름으로 되돌린다
            label = {"PASS": "Pass", "FAIL": "Fail", "N/A": "진행불가", "BLOCKED": "Blocked"}.get(v, v)
            best[t] = (at, {
                "result": label,
                "cycle_id": str(r["id"] or ""),
                "cycle_name": str(r["name"] or ""),
                "at": str(it.get("last_run") or cyc_at or ""),
            })
    return {"items": {k: val for k, (_, val) in best.items()}}


@app.get("/api/tc")
async def get_all_tc(meta: int = 0):
    """
    meta=1 : 목록 렌더에 필요한 메타만 반환 (checks/steps 등 큰 필드 제외 → 초기 로딩 대폭 단축)
    meta=0 (기본): 기존 동작 유지 — 전체 반환
    """
    if meta:
        return {"tcs": await db.tc_list_meta()}
    return {"tcs": await db.tc_list_full()}

# tcid 안에 슬래시(/)·백슬래시(\) 가 있으면 프론트가 __U2F__/__U5C__ sentinel 로 치환해 보냄.
# 서버 라우팅은 {tc_id:path} 로 받아 여기서 원본 tcid 로 복원. FastAPI 가 이미 URL 디코딩한 상태.
def _tc_id_norm(tc_id: str) -> str:
    return (tc_id or "").replace("__U5C__", "\\").replace("__U2F__", "/")

# TC 단건 조회
@app.get("/api/tc/{tc_id}")
async def get_tc(tc_id: str):
    tc_id = _tc_id_norm(tc_id)
    d = await db.tc_get(tc_id)
    if d is None:
        raise HTTPException(404, "TC를 찾을 수 없습니다")
    # 프론트가 checks 를 Array 로 기대(없으면 "시험 절차 로딩 중" 무한대기) — 빈 배열 보장.
    if not isinstance(d.get("checks"), list):
        d["checks"] = []
    return d

@app.get("/api/tc/{tc_id}/path")
async def get_tc_path(tc_id: str):
    """이 시험이 Coverage 트리의 **어디에 있나** — 사업자·폴더·요구사항 차례.

    AI 화면 머리줄이 이 길을 그대로 보여 준다(지시). 트리를 통째로 내려받아
    거슬러 올라가면 화면이 무거워지므로, 여기서 한 번에 짚어 준다.
    """
    tc_id = _tc_id_norm(tc_id)
    async with db.pool().acquire() as c:
        row = await c.fetchrow("SELECT tcid, name, req_id FROM tc WHERE tcid=$1", tc_id)
        if row is None:
            raise HTTPException(404, "TC를 찾을 수 없습니다")
        cats: list[dict] = []
        req = None
        if row["req_id"]:
            r = await c.fetchrow(
                "SELECT id, reqid, title, cat1, cat2, cat3, cat4 FROM req WHERE id=$1",
                row["req_id"],
            )
            if r is not None:
                req = {"id": r["id"], "reqid": r["reqid"] or "", "title": r["title"] or ""}
                at = r["cat4"] or r["cat3"] or r["cat2"] or r["cat1"]
                seen: set[str] = set()
                chain: list[dict] = []
                while at and at not in seen:
                    seen.add(at)
                    cr = await c.fetchrow(
                        "SELECT id, name, parent_id FROM req_category WHERE id=$1", at
                    )
                    if cr is None:
                        break
                    chain.append({"id": cr["id"], "name": cr["name"]})
                    at = cr["parent_id"]
                cats = list(reversed(chain))
    return {"tcid": row["tcid"], "name": row["name"] or "", "cats": cats, "req": req}


# TC 저장
# TC 스텝 스냅샷(자동 백업) — 스텝 수가 급감/이전 값 유실 방지, 최근 20개 유지
TC_SNAP_DIR = DATA_DIR / "tc_snapshots"
TC_SNAP_DIR.mkdir(parents=True, exist_ok=True)

def _tc_snap_dir(tc_id: str) -> Path:
    safe = "".join(ch if (ch.isalnum() or ch in "-_.") else "_" for ch in tc_id)
    p = TC_SNAP_DIR / safe
    p.mkdir(parents=True, exist_ok=True)
    return p

def _tc_snap_save(tc_id: str, prev: dict) -> None:
    try:
        d = _tc_snap_dir(tc_id)
        ts = datetime.now().strftime("%Y%m%d-%H%M%S-%f")[:-3]
        (d / f"{ts}.json").write_text(json.dumps(prev, ensure_ascii=False), encoding="utf-8")
        # 오래된 스냅샷 정리 (최근 20개만 유지)
        files = sorted(d.glob("*.json"))
        for old in files[:-20]:
            try: old.unlink()
            except Exception: pass
    except Exception:
        pass


@app.post("/api/tc/{tc_id}")
async def save_tc(tc_id: str, data: dict):
    tc_id = _tc_id_norm(tc_id)
    """
    저장.

    여러 사람이 같은 시험을 열어 두는 일이 잦다. 잠그지는 않는다 — 대개는
    한 사람이 보기만 하고, 잠가 버리면 보려던 사람이 못 들어오고 잠근
    사람이 자리를 뜨면 아무도 못 고친다.

    대신 **내가 읽은 뒤에 남이 저장했으면 그때 알린다.** 조용히 덮는 것보다
    낫다. 화면이 `_rev`(읽을 때 받은 값)를 같이 보내면 여기서 견준다.
    """
    _base = str(data.pop("_rev", "") or "")
    if _base:
        _now = await db.tc_rev(tc_id)
        if _now and _now != _base:
            who = str((data.get("updated_by") or "")).strip()
            raise HTTPException(
                409,
                f"이 시험을 다른 사람이 먼저 저장했습니다 ({_now[:16].replace('T', ' ')}). "
                "새로 읽어 확인한 뒤 다시 저장하세요."
                + (f" (내 이름: {who})" if who else ""),
            )
    # ★ 안전장치: payload 에 checks 필드가 없거나 배열이 아닌데 DB 에 기존 checks 가 있으면
    #    lazy load 미로드 상태의 tc 를 그대로 저장 시도한 것 → checks 보존해서 스텝 유실 방지.
    #    (프론트 lazy loading 제거했지만 캐시된 옛 코드/외부 API 호출 등에 대한 서버측 마지막 방어선)
    try:
        _need_checks = not isinstance(data.get("checks"), list)
        # ★ 세션도 같은 방어선(지적: 장비 세션 증발 재발) — 필드가 아예 없으면
        #   지운 게 아니라 안 보낸 것이다. [] 는 의도적 삭제로 인정.
        _need_sess = not isinstance(data.get("sessions"), list)
        if _need_checks or _need_sess:
            _prev_full = await db.tc_get(tc_id)
            if isinstance(_prev_full, dict):
                if _need_checks and isinstance(_prev_full.get("checks"), list) and _prev_full["checks"]:
                    data["checks"] = _prev_full["checks"]
                    print(f"[save_tc] checks 누락 감지 → 기존 값 {len(_prev_full['checks'])}건 보존 (tcid={tc_id})", flush=True)
                if _need_sess and isinstance(_prev_full.get("sessions"), list) and _prev_full["sessions"]:
                    data["sessions"] = _prev_full["sessions"]
                    print(f"[save_tc] sessions 누락 감지 → 기존 값 {len(_prev_full['sessions'])}건 보존 (tcid={tc_id})", flush=True)
    except Exception:
        pass
    # 저장 직전 이전 값 스냅샷 — 스텝(checks) 이 있고 새 값과 스텝 수가 다르면 백업.
    # 전체 데이터를 SELECT/비교하면 크기가 커지면 매우 느려짐 → step_count 만 조회해서
    # 스텝 수가 줄어드는 경우(사라짐 위험)에만 전체 데이터 조회+백업 (신규 생성/증가 시엔 스킵).
    try:
        async with db.pool().acquire() as _c:
            prev_step_count = await _c.fetchval("SELECT step_count FROM tc WHERE tcid=$1", tc_id)
        new_checks = data.get("checks") if isinstance(data.get("checks"), list) else []
        new_cli_count = sum(1 for x in new_checks if isinstance(x, dict) and (x.get("kind") or "cli") == "cli")
        if prev_step_count is not None and prev_step_count > 0 and new_cli_count < prev_step_count:
            # 스텝 수가 실제로 줄어들 때만 전체 데이터 조회해서 스냅샷 (비용 큰 경로 최소화)
            prev = await db.tc_get(tc_id)
            if isinstance(prev, dict):
                asyncio.get_event_loop().run_in_executor(None, _tc_snap_save, tc_id, prev)
    except Exception:
        pass
    await db.tc_upsert(tc_id, data)
    # broadcast 는 fire-and-forget (다수 접속 시 순차 send 대기로 응답 느려짐)
    #
    # 누가 저장했는지 함께 싣는다 — 받는 쪽이 「내가 방금 저장한 것」 을
    # 걸러내야 하고, 남이 저장한 것이면 이름을 말해 줘야 한다.
    _by = str(data.get("updated_by") or "").strip()
    try: asyncio.create_task(broadcast({"type": "tc_updated", "tcid": tc_id, "user": _by}))
    except Exception: pass
    return {"success": True}


@app.get("/api/tc/{tc_id}/snapshots")
async def list_tc_snapshots(tc_id: str):
    tc_id = _tc_id_norm(tc_id)
    """TC 자동 백업(스텝 스냅샷) 목록 — 파일명(=시각)만 반환."""
    d = _tc_snap_dir(tc_id)
    items = []
    for f in sorted(d.glob("*.json"), reverse=True):
        try:
            st = f.stat()
            items.append({"name": f.stem, "size": st.st_size, "mtime": st.st_mtime})
        except Exception: pass
    return {"ok": True, "items": items}


@app.get("/api/tc/{tc_id}/snapshots/{name}")
async def get_tc_snapshot(tc_id: str, name: str):
    tc_id = _tc_id_norm(tc_id)
    """스냅샷 상세(원본 TC 데이터) 반환."""
    d = _tc_snap_dir(tc_id)
    safe_name = "".join(ch if (ch.isalnum() or ch in "-_.") else "_" for ch in name)
    p = d / f"{safe_name}.json"
    if not p.exists():
        raise HTTPException(404, "스냅샷 없음")
    try:
        return {"ok": True, "data": json.loads(p.read_text(encoding="utf-8"))}
    except Exception:
        raise HTTPException(500, "스냅샷 읽기 실패")


@app.post("/api/tc/{tc_id}/snapshots/{name}/restore")
async def restore_tc_snapshot(tc_id: str, name: str):
    tc_id = _tc_id_norm(tc_id)
    """스냅샷으로 현재 TC 를 복원 (복원 전 현재 값도 스냅샷)."""
    d = _tc_snap_dir(tc_id)
    safe_name = "".join(ch if (ch.isalnum() or ch in "-_.") else "_" for ch in name)
    p = d / f"{safe_name}.json"
    if not p.exists():
        raise HTTPException(404, "스냅샷 없음")
    try:
        snap = json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        raise HTTPException(500, "스냅샷 읽기 실패")
    # 복원 전 현재 값도 자동 백업
    try:
        prev = await db.tc_get(tc_id)
        if isinstance(prev, dict):
            _tc_snap_save(tc_id, prev)
    except Exception: pass
    # tc_id 는 유지 (구조상 다를 수 있으니 강제)
    if isinstance(snap, dict):
        snap["tcid"] = tc_id
    await db.tc_upsert(tc_id, snap)
    try: asyncio.create_task(broadcast({"type": "tc_updated", "tcid": tc_id}))
    except Exception: pass
    return {"ok": True, "data": snap}

# TC 실행 History 저장 폴더 (tcid 별 하나의 파일)
TC_RUNHIST_DIR = DATA_DIR / "tc_run_history"
TC_RUNHIST_DIR.mkdir(parents=True, exist_ok=True)

def _tc_hist_path(tc_id: str) -> Path:
    safe = "".join(ch if (ch.isalnum() or ch in "-_.") else "_" for ch in tc_id)
    return TC_RUNHIST_DIR / f"{safe}.json"

@app.get("/api/tc/{tc_id}/run-history")
async def get_tc_run_history(tc_id: str):
    tc_id = _tc_id_norm(tc_id)
    """TC 실행 이력(모든 사용자 통합) 조회."""
    p = _tc_hist_path(tc_id)
    if not p.exists():
        return {"ok": True, "history": []}
    try:
        d = load_json(p)
        return {"ok": True, "history": d.get("history", [])}
    except Exception:
        return {"ok": True, "history": []}

@app.post("/api/tc/{tc_id}/run-history")
async def append_tc_run_history(tc_id: str, payload: dict, token: str = ""):
    tc_id = _tc_id_norm(tc_id)
    """새 실행 이력 1건을 추가. 사용자 이름 자동 태깅. WS로 전 접속자에게 알림."""
    u = _user_from_token(token) if token else None
    who = (u.get("name") or u.get("username")) if u else ""
    p = _tc_hist_path(tc_id)
    try:
        d = load_json(p) if p.exists() else {"history": []}
    except Exception:
        d = {"history": []}
    hist = d.get("history") or []
    entry = payload or {}
    entry["user"] = who or entry.get("user") or ""
    # 각 항목 크기 상한 — log 배열이 너무 크면 최근 5000줄만
    if isinstance(entry.get("log"), list) and len(entry["log"]) > 5000:
        entry["log"] = entry["log"][-5000:]
    hist.insert(0, entry)
    if len(hist) > 100:
        hist = hist[:100]
    d["history"] = hist
    save_json(p, d)
    # 다른 접속자에게 새 이력 알림
    try:
        await broadcast({"type": "tc_run_history_new", "tcid": tc_id, "at": entry.get("at",""), "user": entry.get("user",""), "pass": entry.get("pass",0), "fail": entry.get("fail",0), "sec": entry.get("sec",0)})
    except Exception:
        pass
    return {"ok": True, "count": len(hist)}

@app.delete("/api/tc/{tc_id}/run-history")
async def delete_tc_run_history(tc_id: str, idx: int = -1):
    tc_id = _tc_id_norm(tc_id)
    """이력 개별 삭제(idx>=0) 또는 전체 삭제(idx=-1)."""
    p = _tc_hist_path(tc_id)
    if not p.exists():
        return {"ok": True}
    try:
        d = load_json(p)
        hist = d.get("history") or []
        if idx < 0:
            hist = []
        elif 0 <= idx < len(hist):
            hist.pop(idx)
        d["history"] = hist
        save_json(p, d)
        try:
            await broadcast({"type": "tc_run_history_delete", "tcid": tc_id, "idx": idx})
        except Exception:
            pass
        return {"ok": True, "count": len(hist)}
    except Exception as e:
        raise HTTPException(500, str(e))

# 이 TC 를 참조하는 모든 사이클 items 에서 해당 항목 제거 (백그라운드용 헬퍼).
# ⚠️ 반드시 @app.delete 데코레이터 없이 순수 async 함수여야 함 — 데코레이터가 붙으면
#    DELETE /api/tc/{tc_id} 라우팅이 이 헬퍼로 가버려 실제 tc_delete 호출이 안 됨 (버그).
async def _clean_cycle_refs(tc_id: str):
    try:
        async with db.pool().acquire() as c:
            rows = await c.fetch(
                """
                SELECT id FROM cycle
                WHERE data->'items' @? ('$[*] ? (@.tcid == "' || $1::text || '")')::jsonpath
                """,
                tc_id,
            )
            for r in rows:
                cid = r["id"]
                cy = await db.cycle_get(cid)
                if not cy: continue
                items = cy.get("items") or []
                cleaned = [it for it in items if (it or {}).get("tcid") != tc_id]
                if len(cleaned) != len(items):
                    cy["items"] = cleaned
                    await db.cycle_upsert(cid, cy)
                    try: asyncio.create_task(broadcast({"type": "cycle_updated", "cycle_id": cid}))
                    except Exception: pass
    except Exception:
        pass


# TC 삭제 (실제 라우팅 대상)
@app.delete("/api/tc/{tc_id}")
async def delete_tc(tc_id: str):
    tc_id = _tc_id_norm(tc_id)
    # 존재 여부만 EXISTS 로 가볍게 확인 (전체 데이터 SELECT 안 함 — 큰 checks 로드 회피).
    # 휴지통 백업이 필요하면 별도 백그라운드 task 에서 SELECT+파일저장.
    async with db.pool().acquire() as _c:
        exists = await _c.fetchval("SELECT 1 FROM tc WHERE tcid=$1", tc_id)
    if exists:
        async def _trash_bg():
            try:
                existing = await db.tc_get(tc_id)
                if existing:
                    await asyncio.get_event_loop().run_in_executor(None, _trash_put, "tc", tc_id, existing)
            except Exception: pass
        try: asyncio.create_task(_trash_bg())
        except Exception: pass
        await db.tc_delete(tc_id)
        # 이 TC 를 참조하는 사이클 items 도 정리 (백그라운드) — 안 하면 프론트가 404 반복 조회
        try: asyncio.create_task(_clean_cycle_refs(tc_id))
        except Exception: pass
    # WS 브로드캐스트도 백그라운드 (다수 접속 시 순차 send 로 응답 지연)
    try: asyncio.create_task(broadcast({"type": "tc_deleted", "tcid": tc_id}))
    except Exception: pass
    return {"success": True}

# ── 휴지통(삭제 복원) ──
@app.get("/api/trash")
async def list_trash():
    TRASH_DIR.mkdir(exist_ok=True)
    items = []
    for f in sorted(TRASH_DIR.glob("*.json"), reverse=True):
        try:
            rec = load_json(f)
            items.append({"trash_id": rec.get("trash_id"), "kind": rec.get("kind"),
                          "id": rec.get("id"), "name": rec.get("name"),
                          "deleted_at": rec.get("deleted_at"),
                          "tc_count": len(rec.get("bundle") or [])})
        except Exception:
            pass
    return {"items": items}

@app.post("/api/trash/restore/{trash_id}")
async def restore_trash(trash_id: str):
    tf = TRASH_DIR / f"{trash_id}.json"
    if not tf.exists():
        raise HTTPException(404, "휴지통 항목을 찾을 수 없습니다")
    rec = load_json(tf)
    kind = rec.get("kind"); data = rec.get("data") or {}
    restored = {"kind": kind, "id": rec.get("id"), "tc": []}
    if kind == "req":
        rid = rec.get("id")
        if rid:
            await db.req_upsert(rid, data)
        for tcd in (rec.get("bundle") or []):
            tcid = tcd.get("tcid") or tcd.get("id")
            if tcid:
                await db.tc_upsert(tcid, tcd)
                restored["tc"].append(tcid)
    elif kind == "tc":
        tcid = rec.get("id")
        if tcid:
            await db.tc_upsert(tcid, data)
    tf.unlink()
    return {"success": True, "restored": restored}

@app.delete("/api/trash/{trash_id}")
async def purge_trash(trash_id: str):
    tf = TRASH_DIR / f"{trash_id}.json"
    if tf.exists():
        tf.unlink()
    return {"success": True}

# ───────────────────────────────────────────
# 라우터 - 매뉴얼 학습 (AI 참고 문서)
# ───────────────────────────────────────────
MANUALS_DIR = DATA_DIR / "manuals"
def init_manuals_dir():
    MANUALS_DIR.mkdir(parents=True, exist_ok=True)

MANUAL_FOLDERS_FILE = DATA_DIR / "state" / "manual_folders.json"

@app.get("/api/manual-folders")
async def manual_folders_get():
    try:
        if MANUAL_FOLDERS_FILE.exists():
            d = load_json(MANUAL_FOLDERS_FILE)
            if isinstance(d, dict) and isinstance(d.get("folders"), list):
                return {"folders": d["folders"]}
    except Exception:
        pass
    return {"folders": []}

@app.post("/api/manual-folders")
async def manual_folders_set(payload: dict, token: str = ""):
    folders = []
    for f in (payload.get("folders") or []):
        nm = str(f or "").strip()[:60]
        if nm and nm not in folders:
            folders.append(nm)
    save_json(MANUAL_FOLDERS_FILE, {"folders": folders})
    return {"ok": True, "folders": folders}

@app.get("/api/manuals")
async def get_manuals():
    out = []
    for d in await db.manuals_list_full():
        out.append({
            "id": d.get("id"), "name": d.get("name", ""),
            "chars": d.get("chars", len(d.get("text", ""))),
            "active": d.get("active", True),
            "source": d.get("source", ""),
            "created_at": d.get("created_at", ""),
            "folder": d.get("folder", ""),
            "image_count": len(d.get("images") or []),
        })
    return {"manuals": out}

@app.get("/api/manual/{mid}")
async def get_manual(mid: str, images: bool = True):
    d = await db.manuals_get(mid)
    if d is None:
        raise HTTPException(404, "매뉴얼을 찾을 수 없습니다")
    # 이미지 제외 옵션 — 큰 base64 없이 메타·텍스트만 (기본은 기존 호환)
    if not images and isinstance(d, dict):
        d = {k: v for k, v in d.items() if k != "images"}
    return d

@app.get("/api/manual/{mid}/images")
async def get_manual_images(mid: str):
    """이미지만 별도 fetch — 청크 화면에서 지연 로드용"""
    d = await db.manuals_get(mid)
    if d is None:
        raise HTTPException(404, "매뉴얼을 찾을 수 없습니다")
    return {"images": d.get("images") or []}

@app.post("/api/manual/{mid}")
async def save_manual(mid: str, data: dict):
    data["id"] = mid
    if "chars" not in data:
        data["chars"] = len(data.get("text", ""))
    await db.manuals_upsert(mid, data)
    return {"success": True}

@app.delete("/api/manual/{mid}")
async def delete_manual(mid: str):
    await db.manuals_delete(mid)
    return {"success": True}

# ───────────────────────────────────────────
# 라우터 - TC Cycle 관리
# ───────────────────────────────────────────
def init_cycle_dir():
    CYCLE_DIR.mkdir(exist_ok=True)

@app.get("/api/prompts")
async def get_prompts():
    if PROMPTS_FILE.exists():
        return load_json(PROMPTS_FILE)
    return {}

@app.post("/api/prompts")
async def save_prompts(data: dict):
    PROMPTS_FILE.parent.mkdir(parents=True, exist_ok=True)
    save_json(PROMPTS_FILE, data)
    return {"ok": True}

@app.get("/api/custom-fields")
async def get_custom_fields():
    if CUSTOM_FIELDS_FILE.exists():
        return load_json(CUSTOM_FIELDS_FILE)
    return {"req": [], "tc": [], "cycle": []}

@app.post("/api/custom-fields")
async def save_custom_fields(data: dict):
    CUSTOM_FIELDS_FILE.parent.mkdir(parents=True, exist_ok=True)
    save_json(CUSTOM_FIELDS_FILE, data)
    return {"ok": True}

@app.get("/api/help")
async def get_help():
    if HELP_FILE.exists():
        return load_json(HELP_FILE)
    return {"sections": []}

@app.post("/api/help")
async def save_help(data: dict):
    HELP_FILE.parent.mkdir(parents=True, exist_ok=True)
    save_json(HELP_FILE, data)
    return {"ok": True}

# AI 채팅 기록(세션) — 전체 공유(서버 단일 저장). 다른 사용자 질문도 보임.
# id 기준 업서트/삭제로 병합(전체 덮어쓰기 방지 → 동시 작성 안전).
def _chat_load_all():
    d = _kv_load_sync("chat_sessions", {"sessions": []})
    if isinstance(d, dict) and isinstance(d.get("sessions"), list):
        return d["sessions"]
    return []

def _chat_save_all(sessions):
    _kv_save_sync("chat_sessions", {"sessions": sessions})

@app.get("/api/chat-sessions")
async def get_chat_sessions(token: str = ""):
    # 일반 사용자는 본인 채팅만, 관리자는 전원 채팅을 모두 열람
    alls = _chat_load_all()
    u = _user_from_token(token)
    if not u:
        return {"sessions": []}
    if u.get("role") == "관리자":
        return {"sessions": alls}
    uname = u.get("username")
    return {"sessions": [s for s in alls if isinstance(s, dict) and s.get("user") == uname]}

@app.post("/api/chat-sessions")
async def upsert_chat_sessions(data: dict, token: str = ""):
    u = _user_from_token(token)
    uname = u.get("username") if u else None
    is_admin = bool(u and u.get("role") == "관리자")
    sessions = _chat_load_all()
    idx = {}
    for i, s in enumerate(sessions):
        if isinstance(s, dict) and s.get("id"):
            idx[s["id"]] = i
    incoming = []
    if isinstance(data.get("session"), dict):
        incoming = [data["session"]]
    elif isinstance(data.get("sessions"), list):
        incoming = [s for s in data["sessions"] if isinstance(s, dict)]
    for s in incoming:
        sid = s.get("id")
        if not sid:
            continue
        if not s.get("user"):           # 소유자 보장: 누락 시 현재 사용자로
            s["user"] = uname or "default"
        if sid in idx:
            old = sessions[idx[sid]]
            owner = old.get("user") if isinstance(old, dict) else None
            if not is_admin and owner and owner != uname:
                continue                # 일반 사용자는 남의 세션을 덮어쓸 수 없음
            sessions[idx[sid]] = s
        else:
            idx[sid] = len(sessions)
            sessions.append(s)
    _chat_save_all(sessions)
    try:
        await broadcast({"type": "chat_update"})
    except Exception:
        pass
    return {"ok": True, "count": len(sessions)}

@app.post("/api/chat-sessions/delete")
async def delete_chat_session(data: dict, token: str = ""):
    u = _user_from_token(token)
    uname = u.get("username") if u else None
    is_admin = bool(u and u.get("role") == "관리자")
    sid = data.get("id")
    kept = []
    for s in _chat_load_all():
        if isinstance(s, dict) and s.get("id") == sid and (is_admin or s.get("user") == uname or not s.get("user")):
            continue                    # 소유자/관리자/주인없음 → 삭제, 그 외엔 보존
        kept.append(s)
    _chat_save_all(kept)
    try:
        await broadcast({"type": "chat_update"})
    except Exception:
        pass
    return {"ok": True}

@app.get("/api/device-catalog")
async def get_device_catalog():
    d = _kv_load_sync("device_catalog", {"devices": []})
    return d if isinstance(d, dict) else {"devices": []}

DEVICE_CATALOG_BACKUP_DIR = DATA_DIR / "backups" / "device_catalog"

def _device_catalog_backup(prev: dict) -> None:
    """저장 직전 이전 device_catalog 를 타임스탬프 파일로 백업. 최근 30개 유지."""
    try:
        DEVICE_CATALOG_BACKUP_DIR.mkdir(parents=True, exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d-%H%M%S-%f")[:-3]
        (DEVICE_CATALOG_BACKUP_DIR / f"{ts}.json").write_text(
            json.dumps(prev, ensure_ascii=False), encoding="utf-8"
        )
        files = sorted(DEVICE_CATALOG_BACKUP_DIR.glob("*.json"))
        for old in files[:-30]:
            try: old.unlink()
            except Exception: pass
    except Exception:
        pass


@app.post("/api/device-catalog")
async def save_device_catalog(data: dict, force: bool = False):
    # 저장 직전 이전 값 자동 백업 + 안전장치 (급격한 감소 거부)
    try:
        prev = _kv_load_sync("device_catalog", None)
        if isinstance(prev, dict):
            prev_n = len((prev.get("devices") or []))
            new_n = len((data.get("devices") or []) if isinstance(data, dict) else [])
            # ★ 안전장치: 개수가 급격히 줄면(50% 이상 or 5대 이상 유실) 거부하고 백업만 남김.
            if not force and prev_n >= 5 and (new_n <= prev_n * 0.5 or (prev_n - new_n) >= 5):
                asyncio.get_event_loop().run_in_executor(None, _device_catalog_backup, prev)
                return {
                    "ok": False,
                    "error": f"장비 개수 급감 감지 (기존 {prev_n}대 → 요청 {new_n}대) — 유실 방지를 위해 저장 거부. 실제 삭제라면 force=true 옵션으로 재요청하거나 페이지 새로고침 후 다시 시도하세요.",
                    "prev_count": prev_n,
                    "new_count": new_n,
                }
            if prev_n > 0 and new_n < prev_n:
                asyncio.get_event_loop().run_in_executor(None, _device_catalog_backup, prev)
            elif prev_n > 0 and (new_n // 10) != (prev_n // 10):
                asyncio.get_event_loop().run_in_executor(None, _device_catalog_backup, prev)
    except Exception:
        pass
    _kv_save_sync("device_catalog", data)
    return {"ok": True}


@app.get("/api/device-catalog/backups")
async def list_device_catalog_backups():
    """자동 백업 목록 (최근 순)."""
    if not DEVICE_CATALOG_BACKUP_DIR.exists():
        return {"ok": True, "items": []}
    items = []
    for f in sorted(DEVICE_CATALOG_BACKUP_DIR.glob("*.json"), reverse=True):
        try:
            st = f.stat()
            # 파일 내 장비 개수 미리보기
            try:
                d = json.loads(f.read_text(encoding="utf-8"))
                cnt = len(d.get("devices") or [])
            except Exception:
                cnt = 0
            items.append({"name": f.stem, "size": st.st_size, "mtime": st.st_mtime, "count": cnt})
        except Exception: pass
    return {"ok": True, "items": items}


@app.get("/api/device-catalog/backups/{name}")
async def get_device_catalog_backup(name: str):
    safe = "".join(ch if (ch.isalnum() or ch in "-_.") else "_" for ch in name)
    p = DEVICE_CATALOG_BACKUP_DIR / f"{safe}.json"
    if not p.exists():
        raise HTTPException(404, "백업 없음")
    try:
        return {"ok": True, "data": json.loads(p.read_text(encoding="utf-8"))}
    except Exception:
        raise HTTPException(500, "백업 읽기 실패")

# ── Netmiko: 실장비 telnet/ssh 접속 ──
def _netmiko_params(p: dict) -> dict:
    protocol = (p.get("protocol") or "telnet").lower()
    try:
        port = int(p.get("port")) if p.get("port") else (22 if protocol == "ssh" else 23)
    except Exception:
        port = 22 if protocol == "ssh" else 23
    device_type = p.get("device_type") or ("cisco_ios" if protocol == "ssh" else "cisco_ios_telnet")
    return {
        "device_type": device_type,
        "host": (p.get("host") or p.get("ip") or "").strip(),
        "port": port,
        "username": p.get("username") or "",
        "password": p.get("password") or "",
        "secret": p.get("secret") or "",
        "timeout": int(p.get("timeout") or 12),
        "fast_cli": True,            # netmiko 내부 지연 최소화 (명령당 ~1초 → ~0.1초)
        "global_delay_factor": 0.5,  # 출력 안정성 (0.1은 출력 잘림 발생)
    }

def _conn_fail_msg(params: dict, err: Exception) -> str:
    """
    접속 실패를 사람이 쓸 수 있게 적는다.

    netmiko 는 「어디에 못 붙었나」 를 맨 끝에 적는다 —
    `Device settings: cisco_ios 220.1.12.3:22`. 그 문장을 앞에서 200자로
    자르고 있었더니 하필 그 줄이 잘려 「cisco_ios 220.1」 만 남았다.
    반토막 주소는 오해를 부른다 — 등록이 잘못된 줄 알고 장비를 뒤진다.

    그래서 **주소를 맨 앞으로 끌어온다.** 뒤엣말은 잘려도 되지만 어디에
    못 붙었는지는 잘리면 안 된다.
    """
    who = f"{params.get('host', '')}:{params.get('port', '')}"
    proto = "telnet" if "telnet" in str(params.get("device_type", "")) else "ssh"
    body = " ".join(str(err).split())
    if len(body) > 300:
        # 가운데를 접는다. 끝에도 쓸 말이 있다.
        body = body[:200] + " … " + body[-80:]
    return f"{proto} {who} 에 붙지 못했습니다 — {body}"


@app.post("/api/lab-test")
def lab_test(payload: dict):
    # tcl(IXIA N2X 계측기): telnet/ssh가 아니라 Tcl 데몬(9001)으로 연결 — N2X 데몬에 ping
    if str(payload.get("protocol", "")).upper() == "TCL":
        server = str(payload.get("host", "") or payload.get("ip", "")).strip()
        label = str(payload.get("username", "")).strip() or "2"   # N2X 계정 = 등록 ID
        if not server:
            return {"ok": False, "status": "실패", "error": "N2X 서버 IP가 없습니다"}
        try:
            res = _n2x_send(server, label, "ping")
        except Exception as e:
            return {"ok": False, "status": "실패", "error": "N2X 데몬 오류: " + str(e)[:300]}
        if res and res.get("ok"):
            return {"ok": True, "status": "연결됨", "prompt": "N2X session " + str(res.get("session", "")), "enabled": True}
        return {"ok": False, "status": "실패", "error": "N2X 연결 실패: " + str((res or {}).get("error", ""))[:300]}
    params = _netmiko_params(payload)
    if not params["host"]:
        return {"ok": False, "status": "실패", "error": "IP가 없습니다"}
    try:
        from netmiko import ConnectHandler
        conn = ConnectHandler(**params)
        enabled = False
        try:
            if params.get("secret"):
                conn.enable(); enabled = True
        except Exception:
            pass
        prompt = ""
        try:
            prompt = conn.find_prompt()
        except Exception:
            pass
        conn.disconnect()
        return {"ok": True, "status": "연결됨", "prompt": prompt, "enabled": enabled}
    except Exception as e:
        return {"ok": False, "status": "실패", "error": str(e)[:400]}

import threading as _threading
import time as _t
# 장비 연결 캐시 (세션 재사용 → 스텝마다 재접속 방지로 성능 향상)
_conn_cache = {}
_conn_cache_lock = _threading.Lock()
_CONN_IDLE_SEC = 180  # 이 시간 이상 idle이면 생존 확인 후 필요 시 재접속

def _conn_key(p):
    """접속 하나를 가리키는 열쇠.

    ★ `sess`(세션 자리 번호)를 넣는다 — 같은 장비에 세션을 열 개 앉히고
      **동시에** 명령을 넣는 시험이 있다(지시). 자리까지 넣지 않으면 접속
      하나를 열이 나눠 쓰느라 차례로 줄을 선다. 안 보내면 예전처럼 하나다.
    """
    return "{}|{}|{}|{}|{}".format(
        p.get("host"), p.get("port"), p.get("device_type"), p.get("username"), p.get("sess"),
    )

def _get_conn_entry(params):
    key = _conn_key(params)
    with _conn_cache_lock:
        ent = _conn_cache.get(key)
        if ent is None:
            ent = {"conn": None, "ts": 0.0, "lock": _threading.Lock()}
            _conn_cache[key] = ent
    return ent

# netmiko가 접속 시 자동으로 보내는 명령들을 끔 — 'terminal width 511'(미지원 % Invalid input) + 'terminal length 0'(disable_paging, _force_enable에서 1회만 보내도록 중복 제거 + 2초 대기 제거)
try:
    from netmiko.base_connection import BaseConnection as _NMBC
    _NMBC.set_terminal_width = lambda self, *a, **k: ""
    _NMBC.disable_paging = lambda self, *a, **k: ""
    # 접속 지연의 핵심: netmiko silence 대기(last_read 기본 2초)를 0.3초로 — find_prompt·초기 banner 읽기·session 준비가 모두 빨라짐
    _orig_rct = _NMBC.read_channel_timing
    def _fast_rct(self, last_read=0.3, read_timeout=120.0, **k):
        return _orig_rct(self, last_read=last_read, read_timeout=read_timeout, **k)
    _NMBC.read_channel_timing = _fast_rct
    # telnet_login: netmiko가 fast_cli면 delay_factor를 강제로 1로 되돌려 time.sleep(1)×여러번(≈2.5초+) → 최초 접속 지연의 주범.
    # 동일 로직을 작은 고정 지연(0.15초)으로 교체 → 텔넷 로그인 대폭 단축.
    import re as _nmr
    try:
        pass
    except Exception:
        pass
    def _fast_telnet_login(self, pri_prompt_terminator=r"#\s*$", alt_prompt_terminator=r">\s*$",
                           username_pattern=r"(?:user:|username|login|user name)", pwd_pattern=r"assword",
                           delay_factor=1.0, max_loops=20):
        # 고정 sleep/빈 RETURN(race) 대신 '실제 프롬프트가 올 때까지' 기다림 → 빠르면서 안정적(Login incorrect 방지)
        _msg = ""
        try:
            _msg += self.read_until_pattern(pattern=username_pattern, read_timeout=8, re_flags=_nmr.I)
            self.write_channel(self.username + "\r")
        except Exception:
            pass
        try:
            _msg += self.read_until_pattern(pattern=pwd_pattern, read_timeout=8, re_flags=_nmr.I)
            self.write_channel(str(self.password) + "\r")
        except Exception:
            pass
        try:
            _msg += self.read_until_pattern(pattern=r"[>#]\s*$", read_timeout=8, re_flags=_nmr.M)
        except Exception:
            pass
        return _msg
    _NMBC.telnet_login = _fast_telnet_login
    try:
        from netmiko.cisco_base_connection import CiscoBaseConnection as _NMCisco
        _NMCisco.telnet_login = _fast_telnet_login   # cisco_ios_telnet 은 여기서 override 하므로 꼭 패치해야 함
    except Exception:
        pass
except Exception:
    pass

def _force_enable(conn, params, ent=None):
    """접속 직후 User EXEC(>)면 무조건 enable(#) 진입 + 페이징 끄기.
    netmiko conn.enable()이 안 먹는 장비(Ericsson-LG/유비쿼스 등) 대비 직접 'enable' 전송.
    ent(커넥션 엔트리) 를 넘기면 'terminal length 0' 은 세션당 1회만 전송(중복 로그 방지)."""
    import re as _re3
    _P = r"[>#]\s*$"
    at_enable = False
    try:
        conn.write_channel("\n")
        try: cur = conn.read_until_pattern(pattern=_P, read_timeout=6, re_flags=_re3.M)
        except Exception: cur = ""
        if cur.rstrip().endswith("#"):
            at_enable = True
        else:
            conn.write_channel("enable\n")
            try: out = conn.read_until_pattern(pattern=r"(assword|[>#]\s*$)", read_timeout=6, re_flags=_re3.M | _re3.I)
            except Exception: out = ""
            if out.rstrip().endswith("#"):
                at_enable = True
            elif _re3.search(r"password|passwd|암호|비밀번호|secret", out, _re3.I):
                pw = str(params.get("secret") or params.get("password") or "")
                conn.write_channel(pw + "\n")
                try: out2 = conn.read_until_pattern(pattern=_P, read_timeout=6, re_flags=_re3.M)
                except Exception: out2 = ""
                if out2.rstrip().endswith("#"):
                    at_enable = True
    except Exception:
        pass
    # 'terminal length 0' 은 enable(#) 진입 후, 세션당 1회만.
    # (스텝마다 _force_enable 이 호출되므로 flag 로 중복 전송 방지 — 로그 스팸 해소 + 성능 개선)
    if at_enable:
        if ent is not None and ent.get("paging_off"):
            return
        try:
            conn.write_channel("terminal length 0\n")
            conn.read_until_pattern(pattern=r"#\s*$", read_timeout=6, re_flags=_re3.M)
            if ent is not None:
                ent["paging_off"] = True
        except Exception:
            pass



# ── 설정 모드 문맥 유지 ────────────────────────────────────────────
#   「한 스텝 = CLI 하나」 로 나누면 `configure terminal` 과 그다음 명령이
#   다른 호출로 갈린다. 그 사이에 장비가 설정 모드에서 빠져나오면(유휴로
#   빠지는 장비가 있다) 다음 명령이 privileged 프롬프트로 나가 `% Invalid
#   input` 이 난다 — 사용자가 실제로 겪었다(2026-08-19).
#
#   그래서 **세션이 설정 문맥을 기억**한다. 보낸 명령을 보고 문맥을 쌓거나
#   비우고, 보내기 직전에 지금 프롬프트가 설정 모드가 아니면 쌓아 둔 문맥을
#   조용히 다시 밟아 준다. 사람이 적은 절차는 그대로 두고, 잃어버린 상태만
#   되돌리는 방식이다.
_CFG_ENTER = re.compile(r"^\s*(do\s+)?(conf(ig(ure)?)?(\s+t(erminal)?)?|vlan\s+database)\s*$", re.I)
_CFG_LEAVE = re.compile(r"^\s*(end|exit|quit)\s*$", re.I)


def _cfg_ctx_keep(conn, ent, cmd):
    """이 명령을 보내기 전에 — 설정 문맥이 풀렸으면 다시 밟는다."""
    ctx = (ent or {}).get("cfg_ctx") or []
    if not ctx:
        return
    try:
        pr = conn.find_prompt() or ""
    except Exception:
        return
    if "(" in pr:        # 이미 (config)# · (config-if)# 안이다
        return
    for c in ctx:        # 잃어버렸다 — 조용히 되밟는다
        try:
            conn.write_channel(c + "\n")
            conn.read_until_pattern(pattern=r"[>#]\s*$", read_timeout=8, re_flags=re.M)
        except Exception:
            return


def _cfg_ctx_note(ent, cmd):
    """보낸 뒤 — 문맥을 쌓거나 비운다."""
    if ent is None:
        return
    c = str(cmd or "").strip()
    if not c:
        return
    ctx = list(ent.get("cfg_ctx") or [])
    if _CFG_LEAVE.match(c):
        ctx = [] if c.lower().startswith("end") else ctx[:-1]
    elif _CFG_ENTER.match(c):
        ctx = [c]
    elif ctx:
        # 설정 모드 안에서 문맥을 더 파고드는 명령(interface·vlan …)만 쌓는다.
        # 값을 바꾸는 명령(shutdown·ip address …)은 쌓지 않는다 — 되밟으면 두 번 걸린다.
        if re.match(r"^(interface|vlan|line|router|policy-map|class-map)\b", c, re.I):
            ctx = ctx + [c]
    ent["cfg_ctx"] = ctx


def _ensure_conn(ent, params):
    from netmiko import ConnectHandler
    now = _t.time()
    conn = ent.get("conn")
    if conn is not None:
        if now - ent.get("ts", 0.0) < _CONN_IDLE_SEC:
            try:
                conn.find_prompt()  # 생존 확인 (가벼움)
                ent["ts"] = now
                return conn
            except Exception:
                pass
        try:
            conn.disconnect()
        except Exception:
            pass
        ent["conn"] = None
        ent["paging_off"] = False   # 재접속 → 새 세션은 paging 다시 꺼야 함
    conn = ConnectHandler(**params)
    ent["paging_off"] = False   # 새 커넥션도 초기화
    ent["cfg_ctx"] = []         # 새 세션은 설정 문맥도 없다
    _force_enable(conn, params, ent)
    ent["conn"] = conn
    ent["ts"] = now
    return conn

# ── 라이브 스트리밍: SessionLog.write() 를 monkey-patch 해서 netmiko 가 read 한 데이터를
#    SessionLog 에 기록하는 순간 그대로 WebSocket 으로 push. netmiko 는 read_channel 결과를
#    session_log.write(str) 로 즉시 기록하므로, write 훅이 send_command 내부에서 실시간으로
#    호출됨 → 폴러/flush 없이 즉각 스트리밍.
_MAIN_LOOP = None   # startup 훅에서 asyncio.get_running_loop() 로 잡아둠 (워커 스레드용)

def _start_live_pusher(ent, conn, live_key):
    if not live_key: return None
    _slog = getattr(conn, "session_log", None)
    if _slog is None: return None
    # 메인 이벤트 루프 참조 (워커 스레드에서 broadcast 예약용). startup 훅에서 저장한 값 사용.
    _loop = _MAIN_LOOP
    if _loop is None:
        try: _loop = asyncio.get_event_loop()
        except Exception: _loop = None
    _orig_write = _slog.write
    def _push_bytes(_b):
        if not _b or _loop is None: return
        try:
            if isinstance(_b, (bytes, bytearray)):
                _txt = bytes(_b).decode("utf-8", "replace")
            else:
                _txt = str(_b)
            if not _txt: return
            asyncio.run_coroutine_threadsafe(
                broadcast({"type": "cli-live", "live_key": live_key, "chunk": _txt}),
                _loop
            )
        except Exception:
            pass
    def _patched(_data):
        try: _push_bytes(_data)
        except Exception: pass
        return _orig_write(_data)
    try:
        _slog.write = _patched
    except Exception:
        return None
    return (_slog, _orig_write)

def _stop_live_pusher(handle):
    if not handle: return
    try:
        _slog, _orig_write = handle
        try: _slog.write = _orig_write
        except Exception: pass
    except Exception:
        pass

@app.post("/api/run-cli")
def run_cli(payload: dict):
    params = _netmiko_params(payload)
    commands = payload.get("commands") or ([payload["command"]] if payload.get("command") else [])
    if not params["host"]:
        return {"ok": False, "error": "IP가 없습니다", "outputs": []}
    ent = _get_conn_entry(params)
    with ent["lock"]:  # 같은 장비는 순차 (netmiko 비스레드세이프), 다른 장비는 병렬
        try:
            if payload.get("require_session"):
                # iTest 모델: Session Open 으로 열린 세션이 있을 때만 실행 (자동접속 금지)
                # 다만 유휴 정리(_CONN_IDLE_SEC 초과)로 서버 측 conn이 사라진 경우, 프론트의
                # _procSessOpen 은 여전히 세션 있다고 판단하고 있으므로 자동 재접속을 시도.
                # 재접속 실패 시에만 no_session 반환.
                conn = ent.get("conn")
                _auto_reconn = False
                if conn is None:
                    try:
                        conn = _ensure_conn(ent, params)
                        _auto_reconn = True
                    except Exception as _re0:
                        return {"ok": False, "error": "세션이 열려 있지 않습니다 — 먼저 Session Open 스텝을 실행하세요 · 자동 재접속 실패: " + _conn_fail_msg(params, _re0), "no_session": True, "outputs": []}
                ent["ts"] = _t.time()
                _force_enable(conn, params, ent)  # 안전망: 세션이 아직 User EXEC(>)면 enable 재시도 (paging 은 세션당 1회)
                if _auto_reconn:
                    payload["_auto_reconn_notice"] = True
            else:
                conn = _ensure_conn(ent, params)
            repeat = max(1, int(payload.get("repeat", 1) or 1))
            interval = float(payload.get("interval", 1) or 1)
            try:
                cmd_delay = max(0.0, float(payload.get("cmd_delay", 100) or 0) / 1000.0)  # ms→s (iTest식 명령 간 딜레이)
            except Exception:
                cmd_delay = 0.1
            try:
                _tail_max = max(0.0, float(payload.get("tail_wait", 2.0) or 0))  # SecureCRT식: 명령 직후 비동기 로그 수집 상한(초)
            except Exception:
                _tail_max = 2.0
            _live_key = payload.get("live_key") or ""   # 있으면 send_command 대신 _exec_streaming 사용 → WS 로 chunk push
            # Completion Wait (스텝 옵션): >0 이면 이 스텝의 모든 명령은 send_command(=프롬프트 대기)를
            # 쓰지 않고 write_channel 로 명령을 쓰고 지정 초 동안 응답만 수집한다. 세션은 [y/n] 등
            # 대기 상태를 유지한 채 반환되며, 다음 스텝 명령(y/n)이 그 위치에 바로 이어져 입력된다.
            try:
                _wait_only_sec = max(0.0, float(payload.get("wait_only_sec", 0) or 0))
            except Exception:
                _wait_only_sec = 0.0
            # 프롬프트까지 기다려 읽기 (빠른 연속 전송에도 출력 온전히 — 출력 누락 방지)
            import re as _re
            try:
                # 프롬프트는 캐시(연결 시점) 대신 항상 재탐지 — 시험 중 hostname 변경 시
                # 옛 프롬프트를 기다리다 read_timeout(약 25초)을 까먹는 지연을 방지
                # 호스트명만 추출(config 모드 괄호 제거) → enable/config/config-if 어느 레벨이든 매칭되게
                _bp = ""
                try: _bp = (conn.find_prompt() or "").strip().rstrip("#>$ ").split("(")[0].strip()
                except Exception: _bp = ""
                if not _bp:
                    _bp = (getattr(conn, "base_prompt", "") or "").strip().rstrip("#>$ ").split("(")[0].strip()
            except Exception:
                _bp = ""
            # 확인 프롬프트(continue ...? [y/n]: 등)도 함께 매칭 — 장비가 y/n 응답을 기다리며
            # #/> 로 끝나지 않는 줄에서 멈추는 경우, 원래 프롬프트를 못 만나 read_timeout까지
            # 불필요하게 기다리거나 다음 명령(yes/no)이 엉뚱한 타이밍에 꼬여 들어가는 문제 방지.
            _confirm_pat = r"\(y/n\)|\[y/n\]|\(yes/no\)|\[yes/no\]"
            _expect = (r"(?:" + _re.escape(_bp) + r"\S*[#>]\s*$" + r"|" + _confirm_pat + r")") if _bp else None
            _yn_re = _re.compile(_confirm_pat, _re.IGNORECASE)
            _lb = ent.get("log_buf"); _lstart = len(_lb.getvalue()) if _lb else 0
            outputs = []
            _skip_next = False   # 확인 프롬프트에 자동응답한 다음 명령(yes/no 그 자체)은 건너뜀 — 중복 전송 방지
            for _ci, cmd in enumerate(commands):
                # 스텝을 나눠 보내면 그 사이 설정 모드가 풀릴 수 있다 — 되밟는다(지시)
                _cfg_ctx_keep(conn, ent, cmd)
                _cfg_ctx_note(ent, cmd)
                if _skip_next:
                    _skip_next = False
                    continue
                if _ci > 0 and cmd_delay > 0:
                    _t.sleep(min(cmd_delay, 5))
                iters = []
                for k in range(repeat):
                    if k > 0 and interval > 0:
                        _t.sleep(min(interval, 60))
                    # ① 전송 직전: 직전 idle/대기 동안 장비가 밀어낸 잔여 출력(syslog·재표시 프롬프트)을 회수(보존).
                    #    안 비우면 그 끝 프롬프트에 send_command가 즉시 매칭돼 "첫 명령 빈 결과(off-by-one)"가 됨.
                    _pre = ""
                    try:
                        _pre = conn.read_channel() or ""
                    except Exception:
                        _pre = ""
                    # ② 명령 전송 + 프롬프트까지 읽기
                    # Completion Wait On(_wait_only_sec>0): 프롬프트를 기다리지 않는다. write_channel로
                    # 명령만 보내고 지정 초 동안 응답 수집. [y/n] 프롬프트가 뜬 채 반환되고 세션은 그
                    # 상태로 유지 → 다음 스텝 명령(y/n)이 [y/n]: 위치에 바로 이어져 입력된다.
                    # 단일 y/yes/n/no 명령은 직전에 [y/n]이 있었을 가능성 → send_command 대신 write_channel + 3초 대기.
                    _cmd_stripped_pre = (cmd or "").strip().lower()
                    _is_yn_only = _cmd_stripped_pre in ("y", "yes", "n", "no")
                    try:
                        if _wait_only_sec > 0:
                            try:
                                conn.write_channel(cmd + "\n")
                            except Exception as ce_wr:
                                raise ce_wr
                            _wo_buf = ""
                            _dl_wo = _t.time() + min(_wait_only_sec, 600.0)
                            while _t.time() < _dl_wo:
                                try: _ch_wo = conn.read_channel() or ""
                                except Exception: _ch_wo = ""
                                if _ch_wo:
                                    _wo_buf += _ch_wo
                                _t.sleep(0.1)
                            out = _wo_buf.rstrip("\r\n")
                            if not out: out = "[정보] " + cmd + " 전송됨 (Completion Wait " + str(int(_wait_only_sec)) + "s 대기, 응답 없음)"
                        elif _is_yn_only:
                            try:
                                conn.write_channel(cmd + "\n")
                            except Exception as ce_wr:
                                raise ce_wr
                            # y/n 후 응답 수집: 총 상한 6초 내에서, 새 프롬프트(#/>/[y/n]:) 감지되면 조기 종료.
                            # config 저장 확인 등 순차 확인 프롬프트(Save to ...cfg? [y/n]:) 도 수집해야 함.
                            _yn_wait = ""
                            _yn_max = 6.0
                            _yn_quiet = 0.6   # 마지막 데이터 후 조용한 시간 이내에 새 프롬프트 못 만나면 종료
                            _yn_start = _t.time()
                            _yn_last_data = _yn_start
                            _yn_prompt_re = _re.compile(r"(\(y/n\)|\[y/n\]|\(yes/no\)|\[yes/no\])\s*:?\s*$", _re.IGNORECASE|_re.MULTILINE)
                            _yn_end_prompt_re = _re.compile(r"[#>]\s*$", _re.MULTILINE)
                            while (_t.time() - _yn_start) < _yn_max:
                                try: _ch_yw = conn.read_channel() or ""
                                except Exception: _ch_yw = ""
                                if _ch_yw:
                                    _yn_wait += _ch_yw
                                    _yn_last_data = _t.time()
                                    # 새 확인 프롬프트나 일반 프롬프트 감지 → 잠깐 여유(0.3s) 두고 종료
                                    _tail_seg = _yn_wait[-80:]
                                    if _yn_prompt_re.search(_tail_seg) or _yn_end_prompt_re.search(_tail_seg):
                                        _t.sleep(0.3)
                                        try: _extra = conn.read_channel() or ""
                                        except Exception: _extra = ""
                                        if _extra: _yn_wait += _extra
                                        break
                                else:
                                    # 조용한 시간이 지속되면 종료
                                    if (_t.time() - _yn_last_data) >= _yn_quiet and _yn_wait.strip():
                                        break
                                    _t.sleep(0.1)
                            out = _yn_wait.strip()
                            if not out: out = "[정보] " + cmd + " 전송됨 (응답 수집 없음 — 재부팅/즉시명령 등)"
                        elif _expect:
                            _lph = _start_live_pusher(ent, conn, _live_key) if _live_key else None
                            try:
                                out = conn.send_command(cmd, expect_string=_expect, read_timeout=25, cmd_verify=False)
                            finally:
                                _stop_live_pusher(_lph)
                        else:
                            _lph = _start_live_pusher(ent, conn, _live_key) if _live_key else None
                            try:
                                out = conn.send_command(cmd, read_timeout=25, cmd_verify=False)
                            finally:
                                _stop_live_pusher(_lph)
                    except Exception as ce:
                        _emsg = str(ce); out = None
                        # 연결 자체가 끊긴 경우(소켓 강제종료 등) — 명령이 장비에 도달했는지 알 수 없으므로
                        # 재연결 후 동일 명령을 1회 재시도. 재연결도 실패하면 이후 스텝은 실행해도 전부
                        # 실패하므로(죽은 세션에 계속 송신하며 거짓 "완료"가 찍히는 문제) 여기서 확정 실패 처리.
                        _emsg_low = _emsg.lower()
                        _connlost = (
                            isinstance(ce, (OSError, EOFError))
                            or "winerror" in _emsg_low
                            or "not connected" in _emsg_low
                            or ("connection" in _emsg_low and "pattern" not in _emsg_low)
                        )
                        if _connlost:
                            try:
                                try: conn.disconnect()
                                except Exception: pass
                                ent["conn"] = None
                                conn = _ensure_conn(ent, params)
                                if payload.get("require_session"):
                                    _force_enable(conn, params, ent)
                                try:
                                    _bp3 = (conn.find_prompt() or "").strip().rstrip("#>$ ").split("(")[0].strip()
                                    if _bp3: _bp = _bp3; _expect = _re.escape(_bp3) + r"\S*[#>]\s*$"
                                except Exception: pass
                                if _expect:
                                    _lph2 = _start_live_pusher(ent, conn, _live_key) if _live_key else None
                                    try:
                                        out = conn.send_command(cmd, expect_string=_expect, read_timeout=25, cmd_verify=False)
                                    finally:
                                        _stop_live_pusher(_lph2)
                                else:
                                    _lph2 = _start_live_pusher(ent, conn, _live_key) if _live_key else None
                                    try:
                                        out = conn.send_command(cmd, read_timeout=25, cmd_verify=False)
                                    finally:
                                        _stop_live_pusher(_lph2)
                                # 재접속 후 재실행 성공 — 사용자 로그에 알림 노이즈 없이 결과만 그대로 반환
                            except Exception as ce3:
                                out = "[실패] 연결이 끊어졌고 재접속도 실패했습니다: " + str(ce3)
                        # 프롬프트 미검출이어도 명령은 이미 실행됨 → 재전송 금지!
                        # (exit 등 상태변경 명령을 중복 전송하면 모드 이탈·로그아웃 위험)
                        # 버퍼에 남은 출력만 회수하고, 다음 명령용 프롬프트(호스트명)만 갱신.
                        elif "Pattern not detected" in _emsg or "pattern not" in _emsg.lower():
                            try:
                                _acc = ""
                                for _rr in range(10):
                                    _ch = ""
                                    try: _ch = conn.read_channel() or ""
                                    except Exception: _ch = ""
                                    if _ch: _acc += _ch
                                    else: _t.sleep(0.2)
                                out = _acc.strip()
                                try:
                                    _bp3 = (conn.find_prompt() or "").strip().rstrip("#>$ ").split("(")[0].strip()
                                    if _bp3: _bp = _bp3; _expect = _re.escape(_bp3) + r"\S*[#>]\s*$"
                                except Exception: pass
                                if not out: out = "[경고] 프롬프트 미검출 — 출력 일부만 수집됨(재전송 안 함)"
                            except Exception as ce2:
                                out = "[ERROR] " + str(ce2)
                        if out is None:
                            out = "[ERROR] " + _emsg
                    # ②-b 확인 프롬프트(continue...? [y/n]:) 정책:
                    # 백엔드는 자동응답 하지 않는다. [y/n]은 _expect 정규식에 포함되어 있어
                    # send_command가 [y/n]을 만나면 즉시 정상 반환한다. 다음 응답(y/n)은 사용자가 UI에서
                    # 명시적으로 다음 Step에 넣어야 한다. reload는 "Save to ...cfg? [y/n]" 등 여러
                    # 프롬프트가 나올 수 있어 자동 y가 위험(구성 저장 등 부작용) → 사용자가 명령으로 통제.
                    # 스텝 사이 지연으로 장비가 [y/n] 취소하는 문제는 스텝별 "Completion Wait"
                    # (프롬프트 대기) 기능으로 해결한다.
                    # ③ SecureCRT식: 명령 직후 장비가 비동기로 흘리는 로그(reboot의 dying-gasp 등)를
                    #    조용해질 때까지(연속 무출력) 추가 수집 → 콘솔에서 직접 친 것과 동일하게 결과에 포함.
                    _tail = ""
                    if _tail_max > 0:
                        try:
                            _idle = 0
                            _dl = _t.time() + _tail_max
                            while _t.time() < _dl:
                                _t.sleep(0.12)
                                _ch = conn.read_channel() or ""
                                if _ch:
                                    _tail += _ch
                                    _idle = 0
                                else:
                                    _idle += 1
                                    if _idle >= 2:  # 약 0.24s 무출력이면 종료
                                        break
                        except Exception:
                            _tail = ""
                    # 비동기 출력 보존: 직전 잔여(_pre)는 앞, 직후 로그(_tail)는 뒤에 붙여 콘솔처럼 그대로 노출
                    if _pre.strip():
                        out = _pre.rstrip("\r\n") + "\n" + out
                    if _tail.strip():
                        out = out.rstrip("\r\n") + "\n" + _tail.rstrip("\r\n")
                    # 명령 에코 중복 제거: 같은 명령줄이 2번 이상 나오면(텔넷 이중 에코·_pre 잔여) 첫 1줄만 남김
                    _cs = (cmd or "").strip()
                    if _cs and out:
                        _ol = out.split("\n"); _dd2 = []; _ne = 0
                        for _ln in _ol:
                            if _ln.strip() == _cs:
                                _ne += 1
                                if _ne > 1:
                                    continue
                            _dd2.append(_ln)
                        out = "\n".join(_dd2)
                    iters.append({"output": out, "at": datetime.now().strftime("%Y-%m-%d %H:%M:%S")})
                entry = {"command": cmd, "output": iters[-1]["output"], "at": iters[-1]["at"]}
                if repeat > 1:
                    entry["iterations"] = iters
                outputs.append(entry)
            # 자동 재접속(유휴 정리로 conn 손실 후 복구)해도 사용자 로그는 깨끗하게 — 알림 노이즈 없이 결과만 반환
            ent["ts"] = _t.time()  # 세션 유지(재사용)
            _transcript = ""
            if _lb is not None:
                try:
                    _transcript = _lb.getvalue()[_lstart:].decode("utf-8", "replace")
                    _pw0 = params.get("password") or ""
                    if _pw0:
                        _transcript = _transcript.replace(_pw0, "********")
                    _lb.seek(0); _lb.truncate(0)  # 버퍼 비움 → 세션 길어도 메모리 누적 없음
                except Exception:
                    _transcript = ""
            return {"ok": True, "outputs": outputs, "transcript": _transcript}
        except Exception as e:
            try:
                if ent.get("conn"):
                    ent["conn"].disconnect()
            except Exception:
                pass
            ent["conn"] = None
            return {"ok": False, "error": str(e)[:400], "outputs": []}

@app.post("/api/run-cli-stream")
async def run_cli_stream(payload: dict):
    """명령 출력을 줄 단위로 실시간(SSE) 전송 — 블로킹 대신 스트리밍. 명령 에코·끝 프롬프트는 빼고 출력만.
    async generator 로 매 yield 마다 즉시 client 로 flush 되도록 함 (sync generator 는 threadpool 배치 이슈).
    netmiko 호출은 blocking → 짧은 chunk 사이마다 asyncio.sleep(0) 으로 이벤트 루프 양보."""
    from fastapi.responses import StreamingResponse
    import json as _jstr, re as _restr, time as _tstr
    params = _netmiko_params(payload)
    commands = payload.get("commands") or ([payload["command"]] if payload.get("command") else [])
    host_ok = bool(params["host"])
    ent = _get_conn_entry(params) if host_ok else None
    # 프롬프트가 온 뒤에도 얼마나 더 기다릴 것인가.
    #
    # 프롬프트 뒤에 늦게 올라오는 syslog 를 놓치지 않으려는 대기다. 예전에는
    # 2.0 이 코드에 박혀 있어서 **명령마다 2초**가 그냥 나갔다 — 명령 10개짜리
    # 스텝이면 순수 대기만 20초다. 기본을 낮추고 부르는 쪽이 정하게 한다.
    # 뭔가 오면 그때부터 다시 이 시간만큼 연장하므로, 실제로 늦게 오는
    # 출력이 있으면 짧게 잡아도 놓치지 않는다.
    try:
        _quiet_wait = min(30.0, max(0.0, float(payload.get("tail_wait", 0.3) or 0)))
    except Exception:
        _quiet_wait = 0.3
    def _sse(obj):
        return "data: " + _jstr.dumps(obj, ensure_ascii=False) + "\n\n"
    async def _gen():
        if not host_ok:
            yield _sse({"err": "IP가 없습니다"}); yield _sse({"done": True}); return
        # 이 장비 세션 락을 **이벤트 루프를 막지 않고** 잡는다. 동기 `with` 로
        # 잡으면, 앞 스트림이 안 풀어 둔 락을 다음 호출이 기다리며 **API 전체가
        # 언다** — health 까지 죽는다(실사고: run-cli-stream 하나가 걸리자 서버가
        # 통째로 「로딩중」). 논블로킹 시도 + await sleep 로 잡고, 오래 못 잡으면
        # 비켜 준다. finally 로 반드시 푼다(스트림이 중간에 끊겨도).
        _lk_got = False
        _lk_dl = _t.time() + 20
        while not ent["lock"].acquire(blocking=False):
            await asyncio.sleep(0.05)
            if _t.time() > _lk_dl:
                yield _sse({"err": "이 장비 세션이 다른 실행에 잡혀 있습니다 — 잠시 뒤 다시 시도하세요"}); yield _sse({"done": True}); return
        _lk_got = True
        try:
            try:
                # netmiko 는 **블로킹**이다. 접속·enable·프롬프트 찾기를 async 안에서
                # 그대로 부르면, 그 몇 초 동안 **이벤트 루프 전체가 선다** — 같은
                # 반복의 SNMP 도, 남의 요청도, health 까지 멎는다(지적: 반복이 많은
                # 시험이 엄청 느리거나 멈춘다). 48회 반복이면 그 멈춤이 48번 쌓인다.
                # 스레드로 밀어내면 기다리는 동안에도 서버는 계속 돈다.
                if payload.get("require_session"):
                    conn = ent.get("conn")
                    if conn is None:
                        try:
                            conn = await asyncio.to_thread(_ensure_conn, ent, params)
                        except Exception as _re0s:
                            yield _sse({"err": "세션이 열려 있지 않습니다 — 자동 재접속 실패: " + _conn_fail_msg(params, _re0s)}); yield _sse({"done": True}); return
                    ent["ts"] = _t.time()
                    await asyncio.to_thread(_force_enable, conn, params, ent)
                else:
                    conn = await asyncio.to_thread(_ensure_conn, ent, params)
                ent["ts"] = _t.time()
                bp = (getattr(conn, "base_prompt", "") or "").strip().rstrip("#>$ ").split("(")[0].strip()
                if not bp:
                    try:
                        _fp = await asyncio.to_thread(conn.find_prompt)
                        bp = (_fp or "").strip().rstrip("#>$ ").split("(")[0].strip()
                    except Exception: bp = ""
                pr = (_restr.escape(bp) + r"\S*[#>]\s*$") if bp else None
                for cmd in commands:
                    yield _sse({"cmd": cmd})       # 명령 입력 표시(라이브 터미널에 '$ cmd')
                    # 장비가 **무엇을 돌려줬는지** 남긴다.
                    #
                    # 여태 이 자리에 기록이 없어서, 「명령은 나갔는데 화면에
                    # 아무것도 안 나온다」 를 만났을 때 접속이 죽은 것인지
                    # 장비가 침묵한 것인지 가릴 방법이 없었다(지적). 받은
                    # 바이트와 걸린 시간만 남겨도 그 둘이 갈린다 — 0바이트면
                    # 장비가 안 보낸 것이고, 오래 걸렸으면 기다리다 끝난 것이다.
                    _cli_t0 = _tstr.time(); _cli_n = 0
                    await asyncio.sleep(0)
                    # 스텝을 나눠 보내면 그 사이 설정 모드가 풀릴 수 있다 —
                    # 풀렸으면 쌓아 둔 문맥을 조용히 되밟는다(지시: 프롬프트 유지)
                    await asyncio.to_thread(_cfg_ctx_keep, conn, ent, cmd)
                    try: await asyncio.to_thread(conn.read_channel)
                    except Exception: pass
                    await asyncio.to_thread(conn.write_channel, cmd + "\n")
                    _cfg_ctx_note(ent, cmd)
                    echo_done = False; pending = ""; idle = 0; dl = _tstr.time() + 30
                    while _tstr.time() < dl:
                        ch = ""
                        try: ch = conn.read_channel() or ""
                        except Exception: ch = ""
                        if ch:
                            idle = 0
                            _cli_n += len(ch)
                            if not echo_done:
                                pending += ch
                                _nl = pending.find("\n")
                                if _nl < 0:
                                    await asyncio.sleep(0)
                                    continue
                                echo_done = True
                                # 명령 에코(첫 줄)에는 **그때의 진짜 프롬프트**가 들어 있다 —
                                # `R3(config)#interface …`. 여태 통째로 버려서 화면이
                                # 늘 `E6100#` 로 굳어 있었다(지적). 프롬프트만 떼어 보낸다.
                                _echo = pending[:_nl].replace("\r", "").rstrip()
                                _mpr = _restr.match(r"^(\S.*?[#>])\s*(?:" + _restr.escape(cmd) + r")?\s*$", _echo)
                                if _mpr:
                                    yield _sse({"pr": _mpr.group(1)})
                                pending = pending[_nl + 1:]   # 에코 줄 자체는 버린다
                            else:
                                pending += ch
                            # chunk 즉시 push (줄 단위 대기 X) — echo 처리된 이후 모든 데이터 실시간
                            if pending:
                                _emit = pending
                                # 마지막 줄이 완전한 줄이 아니면(개행 없음) 프롬프트 확인 위해 남겨둠
                                if pending.endswith("\n"):
                                    pending = ""
                                else:
                                    _last_nl = pending.rfind("\n")
                                    if _last_nl >= 0:
                                        _emit = pending[:_last_nl + 1]
                                        pending = pending[_last_nl + 1:]
                                    else:
                                        _emit = ""   # 개행 없는 짧은 tail 은 프롬프트 후보 → 보류
                                if _emit:
                                    yield _sse({"o": _emit})
                                    await asyncio.sleep(0)
                            if pr and pending.strip() and _restr.search(pr, pending.strip()):
                                _quiet_dl = _tstr.time() + _quiet_wait
                                _saw_more = False
                                while _tstr.time() < _quiet_dl:
                                    try: _ch2 = conn.read_channel() or ""
                                    except Exception: _ch2 = ""
                                    if _ch2:
                                        pending += _ch2; _saw_more = True
                                        if pending.endswith("\n"):
                                            yield _sse({"o": pending}); pending = ""
                                        else:
                                            _last_nl = pending.rfind("\n")
                                            if _last_nl >= 0:
                                                yield _sse({"o": pending[:_last_nl + 1]})
                                                pending = pending[_last_nl + 1:]
                                        await asyncio.sleep(0)
                                        _quiet_dl = _tstr.time() + _quiet_wait
                                    else:
                                        await asyncio.sleep(0.05)
                                if _saw_more and pending.strip() and not _restr.search(pr, pending.strip()):
                                    continue
                                pending = ""; break
                        else:
                            if echo_done and pending.strip():
                                _ps = pending.strip()
                                if (pr and _restr.search(pr, _ps)) or _restr.search(r"\S+[#>]\s*$", _ps):
                                    try:
                                        _np = _ps.split("\n")[-1].strip().rstrip("#>$ ")
                                        if _np: pr = _restr.escape(_np) + r"\S*[#>]\s*$"
                                    except Exception: pass
                                    pending = ""; break
                            idle += 1
                            if idle > 240:
                                if pending: yield _sse({"o": pending})
                                pending = ""; break
                            await asyncio.sleep(0.05)
                    if pending.strip() and not (pr and _restr.search(pr, pending.strip())):
                        yield _sse({"o": pending})
                        await asyncio.sleep(0)
                    _cli_ms = int((_tstr.time() - _cli_t0) * 1000)
                    print(
                        f"[cli] {params.get('host')} sess={payload.get('sess')} "
                        f"{_cli_n}B {_cli_ms}ms"
                        f"{' 응답없음' if _cli_n == 0 else ''} :: {cmd[:80]}",
                        flush=True,
                    )
                ent["ts"] = _t.time()
                yield _sse({"done": True})
            except Exception as e:
                yield _sse({"err": str(e)[:200]}); yield _sse({"done": True})
        finally:
            if _lk_got:
                try: ent["lock"].release()
                except Exception: pass
    return StreamingResponse(_gen(), media_type="text/event-stream", headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache", "Content-Encoding": "identity"})

@app.post("/api/cli-complete")
def cli_complete(payload: dict):
    """터미널 Tab 자동완성: 영속 세션에 '부분명령+Tab'을 보내 장비가 완성한 명령을 읽어 반환."""
    import re as _re
    params = _netmiko_params(payload)
    partial = str(payload.get("partial", "") or "")
    help_mode = bool(payload.get("help"))   # '?' 도움말: 입력 가능한 명령 목록(실행 안 함)
    if not params["host"]:
        return {"ok": False, "error": "IP가 없습니다"}
    ent = _get_conn_entry(params)
    with ent["lock"]:
        conn = ent.get("conn")
        if conn is None:
            return {"ok": False, "error": "세션이 열려 있지 않습니다", "no_session": True}
        try:
            try: conn.read_channel()   # 잔여 비우기
            except Exception: pass
            conn.write_channel(partial + ("?" if help_mode else "\t"))
            _t.sleep(0.25)
            data = ""; _idle = 0; _dl = _t.time() + (2.5 if help_mode else 1.2)
            while _t.time() < _dl:
                ch = ""
                try: ch = conn.read_channel() or ""
                except Exception: ch = ""
                if ch:
                    data += ch; _idle = 0
                else:
                    _idle += 1
                    if _idle >= 2: break
                    _t.sleep(0.1)
            try:  # 장비 입력 라인 비우기(Ctrl+U) → 다음 명령 오염 방지
                conn.write_channel("\x15"); _t.sleep(0.05); conn.read_channel()
            except Exception: pass
            ent["ts"] = _t.time()
            clean = _re.sub(r"\x1b\[[0-9;?]*[a-zA-Z]", "", data).replace("\x07", "")
            buf = []
            for chh in clean:
                if chh == "\x08":
                    if buf: buf.pop()
                else:
                    buf.append(chh)
            clean = "".join(buf)
            seg = _re.split(r"[\r\n]", clean)
            if help_mode:
                # '?' 도움말: 프롬프트(에코 'host# show ?' / 재표시) 줄을 빼고 명령 목록만 반환
                _isprompt = _re.compile(r"^\s*\S+[>#]")
                opts = []
                for s in seg:
                    st = s.rstrip()
                    if not st.strip():
                        continue
                    if _isprompt.match(st):
                        continue
                    opts.append(st)
                return {"ok": True, "help": "\n".join(opts), "options": opts, "raw": clean[-3000:]}
            last = ""
            for s in seg:
                if s.strip(): last = s
            m = _re.search(r"[>#]\s*(.*)$", last)
            if m and m.group(1).strip():
                completed = m.group(1).strip()
            else:
                completed = last.strip() or partial
            options = [s.rstrip() for s in seg if s.strip() and not _re.search(r"[>#]\s*$", s)]
            return {"ok": True, "completed": completed, "options": options, "raw": clean[-2000:]}
        except Exception as e:
            return {"ok": False, "error": str(e)[:300]}

@app.post("/api/session-open")
def session_open(payload: dict):
    params = _netmiko_params(payload)
    if not params["host"]:
        return {"ok": False, "error": "IP가 없습니다"}
    ent = _get_conn_entry(params)
    with ent["lock"]:
        try:
            # Session Open: 기존 세션 닫고 새로 열어 로그인 과정(배너/Username/Password)을 캡처
            if ent.get("conn"):
                try:
                    ent["conn"].disconnect()
                except Exception:
                    pass
                ent["conn"] = None
            # 빠른 접속(터미널용): 로그인 전문 캡처를 생략 → 접속 지연 최소화. 프롬프트만 반환(로그인 로그 노출 안 함).
            if payload.get("fast"):
                import io as _iof, re as _ref
                _ts0 = _t.time()
                buf = _iof.BytesIO()
                from netmiko import ConnectHandler as _CH
                p2 = dict(params); p2["session_log"] = buf; p2["session_log_record_writes"] = True
                p2["global_delay_factor"] = 0.1   # 텔넷 로그인 루프의 sleep(0.5s×~20)이 최초 접속 지연의 주범 → 낮춰 단축 (장비 응답은 그대로)
                p2["conn_timeout"] = 8            # TCP 접속 타임아웃 단축
                conn = _CH(**p2)                # 빠른 읽기 + 로그인 전문 캡처
                _ts_conn = _t.time()            # ① 접속 + telnet 로그인 완료 시점
                ent["paging_off"] = False       # Session Open: 새 세션이므로 paging flag 초기화
                _force_enable(conn, params, ent)  # 무조건 enable(#) + 페이징 끄기 (세션당 1회)
                _ts_enable = _t.time()          # ② enable/페이징 완료
                prompt = ""
                try:
                    prompt = conn.find_prompt()
                except Exception:
                    pass
                _ts_prompt = _t.time()          # ③ 프롬프트 확인
                try:
                    conn.session_log.flush()
                except Exception:
                    pass
                ent["conn"] = conn
                ent["ts"] = _t.time()
                ent["log_buf"] = buf
                login_log = ""
                try:
                    login_log = buf.getvalue().decode("utf-8", "replace")
                    for _pw in (params.get("password") or "", params.get("secret") or ""):
                        if _pw:
                            login_log = login_log.replace(_pw, "********")
                    login_log = _ref.sub(r"(?i)(password[^\r\n:]*:[ \t]*)\S+", r"\1********", login_log)
                except Exception:
                    pass
                return {"ok": True, "prompt": prompt, "login_log": login_log,
                        "elapsed": round(_t.time() - _ts0, 1),
                        "t_connect": round(_ts_conn - _ts0, 1),
                        "t_enable": round(_ts_enable - _ts_conn, 1),
                        "t_prompt": round(_ts_prompt - _ts_enable, 1)}
            import io as _io, re as _re2
            buf = _io.BytesIO()
            from netmiko import ConnectHandler
            p2 = dict(params); p2["session_log"] = buf; p2["session_log_record_writes"] = True
            p2["fast_cli"] = False          # 로그인 배너/Username/Password 교환을 session_log에 온전히 담기 위해 천천히 읽기
            p2["global_delay_factor"] = 1
            conn = ConnectHandler(**p2)
            ent["paging_off"] = False   # Session Open: 새 세션이므로 paging flag 초기화
            _force_enable(conn, params, ent)  # 접속 직후 무조건 enable(#) 진입 + terminal length 0 (세션당 1회)
            prompt = ""
            try:
                prompt = conn.find_prompt()
            except Exception:
                pass
            # netmiko(SessionLog)는 로그를 메모리(slog_buffer)에 모았다가 flush 시점에만 실제 버퍼(buf)에 기록한다.
            # 접속 직후엔 flush 전이라 buf가 비어 로그인 과정이 누락됨 → 명시적으로 flush 해서 로그인 전문을 buf에 내린다.
            try:
                conn.session_log.flush()
            except Exception:
                pass
            ent["conn"] = conn
            ent["ts"] = _t.time()
            ent["log_buf"] = buf  # 이후 run-cli 명령들의 raw 입출력도 이 버퍼에 누적 → 실시간 터미널 로그
            try:
                login_log = buf.getvalue().decode("utf-8", "replace")
            except Exception:
                login_log = ""
            # 비밀번호 마스킹 (실제 값 치환 + Password: 뒤 입력 방어적 마스킹)
            for _pw in (params.get("password") or "", params.get("secret") or ""):
                if _pw:
                    login_log = login_log.replace(_pw, "********")
            login_log = _re2.sub(r"(?i)(password[^\r\n:]*:[ \t]*)\S+", r"\1********", login_log)
            # 셋업 명령(페이징 'terminal length 0' / 'terminal width')만 로그인 로그에서 숨긴다.
            # 그로 인한 '프롬프트만' 연속 중복 줄만 접고(로그인 배너·login:·Password: 등 실제 내용 줄은 절대 제거 X),
            # 만약 결과가 비면 원본을 유지해 로그인 과정이 사라지지 않게 한다(전송 동작엔 영향 없음 — 표시만 정리).
            _ll = []
            _prompt_re = _re2.compile(r"^\S+[>#]\s*$")   # 프롬프트만 있는 줄
            for _ln in login_log.split("\n"):
                _s = _ln.strip()
                if ("terminal length 0" in _s) or ("terminal width" in _s):
                    continue
                if ("% Invalid input" in _s) or (_s == "^"):   # terminal width 511 미지원 장비의 에러 표시 제거
                    continue
                if _ll and _ll[-1].strip() == _s and _prompt_re.match(_s):
                    continue
                _ll.append(_ln)
            _filtered = "\n".join(_ll)
            if _filtered.strip():
                login_log = _filtered
            return {"ok": True, "prompt": prompt, "login_log": login_log}
        except Exception as e:
            ent["conn"] = None
            return {"ok": False, "error": str(e)[:400]}

@app.post("/api/session-close")
def session_close(payload: dict):
    params = _netmiko_params(payload)
    ent = _get_conn_entry(params)
    with ent["lock"]:
        try:
            if ent.get("conn"):
                ent["conn"].disconnect()
        except Exception:
            pass
        ent["conn"] = None
        ent["ts"] = 0.0
        return {"ok": True}


@app.post("/api/session-key")
def session_key(payload: dict):
    """탭 완성·`?` 도움말 — 터미널의 진짜 CLI 손맛(지적: 캡쳐 터미널에서 안 됨).

    지금 세션 채널에 **엔터 없이 글자만** 흘리고, 장비가 돌려주는 것
    (완성된 낱말·도움말 목록)을 읽어 온다. 명령은 실행되지 않는다.

    읽고 나면 Ctrl-U 로 장비 쪽 입력줄을 비운다 — 안 비우면 장비 버퍼에
    친 글자가 남아, 다음에 보내는 진짜 명령 앞에 붙어 엉뚱한 명령이 된다.
    (netmiko 의 send_command 는 보내기 전에 수신 버퍼를 비우므로, 남는
    재출력 프롬프트는 다음 명령에 안 섞인다.)
    """
    params = _netmiko_params(payload)
    text = str(payload.get("text") or "")
    # 개행은 금지 — 이 길은 완성/도움말용이지 실행이 아니다
    text = text.replace("\r", "").replace("\n", "")
    if not text:
        return {"ok": False, "error": "보낼 글자가 없습니다"}
    ent = _get_conn_entry(params)
    with ent["lock"]:
        conn = ent.get("conn")
        if not conn:
            return {"ok": False, "error": "세션이 없습니다 — 먼저 접속하세요"}
        try:
            conn.write_channel(text)
            out = ""
            quiet = 0.0
            t0 = _t.time()
            # 장비가 조용해질 때까지 모은다 — 도움말이 길면 여러 조각으로 온다
            while _t.time() - t0 < 2.5:
                _t.sleep(0.12)
                chunk = conn.read_channel()
                if chunk:
                    out += chunk
                    quiet = 0.0
                else:
                    quiet += 0.12
                    if out and quiet >= 0.4:
                        break
            # 장비 입력줄 비우기 (Ctrl-U) — 대부분의 네트워크 OS 가 받는다
            try:
                conn.write_channel("\x15")
                _t.sleep(0.2)
                conn.read_channel()
            except Exception:
                pass
            return {"ok": True, "out": out}
        except Exception as exc:
            return {"ok": False, "error": str(exc)}

@app.post("/api/ping-stream")
async def ping_stream(payload: dict):
    """ping 을 줄 단위로 흘려보낸다 (SSE).

    `/api/ping` 은 다 끝난 뒤 한 번에 준다. 재부팅 시험에서는 그게 쓸모가
    없다 — '안 되고… 안 되고… 됐다' 가 실시간으로 보여야 언제 살아났는지
    안다. 4번을 다 기다린 뒤 결과만 보면 그 순간을 놓친다.

    subprocess 를 asyncio 로 띄워 stdout 을 한 줄씩 읽어 보낸다.
    """
    from fastapi.responses import StreamingResponse
    import json as _jp

    host = (payload.get("host") or "").strip()
    try:
        count = max(1, min(60, int(payload.get("count", 4) or 4)))
    except Exception:
        count = 4

    def _sse(obj):
        return "data: " + _jp.dumps(obj, ensure_ascii=False) + "\n\n"

    async def _gen():
        if not host:
            yield _sse({"err": "대상 IP 가 없습니다"})
            yield _sse({"done": True, "alive": False})
            return
        # -c 는 리눅스. 컨테이너 안에서만 도므로 윈도우 분기는 두지 않는다.
        proc = None
        try:
            proc = await asyncio.create_subprocess_exec(
                "ping", "-c", str(count), "-W", "2", host,
                stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.STDOUT,
            )
        except FileNotFoundError:
            yield _sse({"err": "이미지에 ping 이 없습니다"})
            yield _sse({"done": True, "alive": False})
            return
        except Exception as e:
            yield _sse({"err": str(e)[:300]})
            yield _sse({"done": True, "alive": False})
            return

        try:
            assert proc.stdout is not None
            while True:
                # ping -c N 은 대략 N초가 걸린다. 넉넉히 잡되 영영 매달리지는
                # 않게 한 줄마다 상한을 둔다.
                try:
                    line = await asyncio.wait_for(proc.stdout.readline(), timeout=count + 20)
                except asyncio.TimeoutError:
                    yield _sse({"err": "응답이 너무 늦습니다"})
                    break
                if not line:
                    break
                yield _sse({"o": line.decode("utf-8", "replace")})
            rc = await proc.wait()
        except Exception as e:
            print(f"[ping_stream] 읽기 실패: {e}", flush=True)
            rc = 1
        finally:
            try:
                if proc and proc.returncode is None:
                    proc.kill()
            except Exception as e:
                print(f"[ping_stream] 정리 실패: {e}", flush=True)
        yield _sse({"done": True, "alive": rc == 0})

    return StreamingResponse(
        _gen(),
        media_type="text/event-stream",
        headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache", "Content-Encoding": "identity"},
    )


@app.post("/api/ping")
def ping_host(payload: dict):
    host = (payload.get("host") or "").strip()
    if not host:
        return {"ok": False, "error": "host(IP)가 없습니다", "output": ""}
    import subprocess, platform
    is_win = platform.system().lower().startswith("win")
    try:
        count = max(1, min(20, int(payload.get("count", 4) or 4)))
    except Exception:
        count = 4
    cmd = ["ping", "-n" if is_win else "-c", str(count), host]
    try:
        p = subprocess.run(cmd, capture_output=True, timeout=40)
        enc = "cp949" if is_win else "utf-8"
        out = (p.stdout or b"").decode(enc, errors="replace")
        err = (p.stderr or b"").decode(enc, errors="replace")
        full = out + (("\n" + err) if err.strip() else "")
        return {"ok": True, "output": full.strip(), "returncode": p.returncode, "alive": p.returncode == 0}
    except Exception as e:
        return {"ok": False, "error": str(e)[:300], "output": ""}

def _snmp_is_noinstance(s):
    s = str(s)
    return any(k in s for k in ("No Such Instance", "No Such Object", "noSuchInstance",
                                "noSuchObject", "endOfMibView", "No more variables"))

def _fmt_timeticks(ticks):
    # SNMP TimeTicks(1/100초) → "N일 H시간 M분 S초" 사람이 읽는 형식
    try:
        t = int(ticks)
    except Exception:
        return None
    if t < 0:
        return None
    sec = t // 100
    d = sec // 86400; sec %= 86400
    h = sec // 3600; sec %= 3600
    m = sec // 60; s = sec % 60
    out = ("%d days, " % d) if d else ""
    return out + ("%d hours, %d mins, %d secs" % (h, m, s))

def _snmp_val_str(vobj):
    # varbind 값 문자열 — TimeTicks면 "원시값 (N일 H시간 M분 S초)" 로 변환
    try:
        s = vobj.prettyPrint()
    except Exception:
        return str(vobj)
    try:
        if vobj.__class__.__name__ == "TimeTicks":
            ft = _fmt_timeticks(int(vobj))
            if ft and ("days" not in s and "hours" not in s):
                s = s + " (" + ft + ")"
    except Exception:
        pass
    # 값이 OBJECT IDENTIFIER(또는 숫자 OID 형태)면 이름으로 해석 (예: 1.3.6.1.4.1.7800.1.238 → .iso.org...E5724RL)
    try:
        cn = vobj.__class__.__name__
        st = s.strip()
        _looks_oid = st.lstrip(".").replace(".", "").isdigit() and st.lstrip(".").count(".") >= 4
        if cn in ("ObjectIdentifier", "ObjectIdentity") or _looks_oid:
            nm = _snmp_oid_name(st)
            if nm and nm.lstrip(".") != st.lstrip("."):
                s = nm + " (" + st.lstrip(".") + ")"
    except Exception:
        pass
    return s

# SNMP enum 이름 매핑 — OID(인스턴스 제외 베이스) → {정수값(str): 이름}.
# 대량 매핑은 data/MIB/ 에서 추출한 data/snmp/snmp_enums.json 에서 자동 로드 (tools/mib_enums.py 로 재생성).
_SNMP_MANUAL = {
    # 수동 보정(MIB 없거나 덮어쓸 OID만 여기 추가). MIB 추출값보다 우선.
    "1.3.6.1.4.1.7800.100.1.1.3.6": {"1": "fiveSec", "2": "oneMin", "3": "fiveMin"},
}
SNMP_ENUM_MAP = dict(_SNMP_MANUAL)
# SNMP SET 타입 약어 → 표시명 (snmp_set_api 응답/에러힌트용). 동작 불변.
_SNMP_TYPE_NAMES = {"i": "Integer", "u": "Unsigned", "c": "Counter32", "g": "Gauge32", "t": "TimeTicks", "s": "String", "a": "IpAddress", "x": "Hex"}
SNMP_OID_NAMES = {}            # 숫자 OID → 이름(전체 경로). 값이 OBJECT IDENTIFIER 인 응답·OID 표시 해석용 (mib_enums.py 추출)
_SNMP_ENUM_MTIME = [None]   # data/snmp/snmp_enums.json 의 마지막 로드 mtime — 변경 시 자동 재로드
def _load_snmp_enums(force=False):
    # data/snmp/snmp_enums.json(MIB 추출) 로드 후 SNMP_ENUM_MAP 재구성(수동 항목 우선). 파일 mtime이 바뀐 경우만 다시 읽음 → 재시작 없이 반영.
    try:
        f = DATA_DIR / "snmp" / "snmp_enums.json"
        mt = f.stat().st_mtime if f.exists() else 0.0
        if (not force) and (_SNMP_ENUM_MTIME[0] == mt):
            return
        _SNMP_ENUM_MTIME[0] = mt
        newmap = dict(_SNMP_MANUAL); cnt = 0
        if f.exists():
            import json as _json
            data = _json.loads(f.read_text(encoding="utf-8"))
            for oid, m in data.items():
                if oid not in newmap and isinstance(m, dict):
                    newmap[oid] = {str(k): str(v) for k, v in m.items()}; cnt += 1
        SNMP_ENUM_MAP.clear(); SNMP_ENUM_MAP.update(newmap)   # 같은 객체 갱신(참조 유지)
        # OID → 이름 맵도 같이 로드(숫자 OID 해석)
        nf = DATA_DIR / "snmp" / "snmp_names.json"; ncnt = 0
        if nf.exists():
            import json as _json
            nd = _json.loads(nf.read_text(encoding="utf-8"))
            if isinstance(nd, dict):
                SNMP_OID_NAMES.clear(); SNMP_OID_NAMES.update({str(k): str(v) for k, v in nd.items()}); ncnt = len(SNMP_OID_NAMES)
        print(f"[SNMP] MIB enum {cnt}개 / OID 이름 {ncnt}개 로드 (총 enum {len(SNMP_ENUM_MAP)})")
    except Exception as e:
        print(f"[SNMP] enum/name JSON 로드 실패: {e}")
def _snmp_oid_name(num):
    # 숫자 OID(앞 점 무관) → 이름. 인스턴스 접미(.0, .N)도 떼고 재시도. 없으면 None.
    try:
        s = str(num).lstrip(".")
        if s in SNMP_OID_NAMES: return SNMP_OID_NAMES[s]
        ps = s.split(".")
        for cut in (1, 2):
            if cut < len(ps):
                cand = ".".join(ps[: len(ps) - cut])
                if cand in SNMP_OID_NAMES: return SNMP_OID_NAMES[cand]
        return None
    except Exception:
        return None
_load_snmp_enums(force=True)
def _oid_to_num(nm):
    try:
        s = str(nm)
        s = s.replace("SNMPv2-SMI::enterprises", "1.3.6.1.4.1").replace("SNMPv2-SMI::", "")
        s = s.replace("iso.org.dod.internet.private.enterprises", "1.3.6.1.4.1").replace("iso.org.dod.internet", "1.3.6.1")
        return s.lstrip(".")
    except Exception:
        return str(nm)
def _snmp_enum(oid_num, vobj, val):
    # 정수형 값이고 매핑에 OID(인스턴스 제외)가 있으면 이름으로 표시
    try:
        if vobj.__class__.__name__ not in ("Integer", "Integer32", "Unsigned32", "Gauge32", "Counter32"):
            return val
        ps = oid_num.split(".")
        for cut in (1, 2, 0):
            if cut > len(ps):
                continue
            cand = ".".join(ps[: len(ps) - cut]) if cut else oid_num
            m = SNMP_ENUM_MAP.get(cand)
            if m:
                k = str(int(vobj))
                if k in m:
                    return m[k]
        return val
    except Exception:
        return val

def _snmp_set_enum(oid, value):
    # SET 값이 enum '이름'(예: fiveMin)이고 그 OID에 enum 맵이 있으면 정수(예: 3)로 변환. 이미 숫자/없음이면 None.
    try:
        v = str(value).strip()
        if v == "" or v.lstrip("-").isdigit():
            return None
        ps = str(oid).lstrip(".").split(".")
        for cut in (1, 2, 0):
            if cut > len(ps):
                continue
            cand = ".".join(ps[: len(ps) - cut]) if cut else ".".join(ps)
            m = SNMP_ENUM_MAP.get(cand)
            if m:
                for k, nm in m.items():
                    if str(nm).lower() == v.lower():
                        return str(k)
        return None
    except Exception:
        return None

def _snmp_close(eng):
    # SNMP 엔진/디스패처(UDP 소켓) 정리 — 안 닫으면 반복 시 FD 누적 → Windows select() 512 한계 초과로 크래시
    if eng is None:
        return
    for m in ("close_dispatcher", "closeDispatcher"):
        try:
            getattr(eng, m)()
            return
        except Exception:
            pass
    try:
        td = getattr(eng, "transport_dispatcher", None) or getattr(eng, "transportDispatcher", None)
        if td:
            for m in ("close_dispatcher", "closeDispatcher"):
                try:
                    getattr(td, m)()
                    return
                except Exception:
                    pass
    except Exception:
        pass

@app.get("/api/snmp-oids")
async def snmp_oids(q: str = "", limit: int = 50):
    """MIB 에서 뽑아 둔 OID 이름표를 찾는다.

    화면에서 OID 를 손으로 치게 두면 `1.3.6.1.2.1.1.3.0` 를 외우거나 문서를
    뒤져야 한다. 이름으로 찾아 눌러 넣게 한다.

    자료는 `data/snmp/snmp_names.json` 이고 `tools/mib_enums.py` 가 만든다.
    없으면 빈 목록을 주되 어디서 만드는지 함께 알려 준다 — 빈 화면만 보면
    기능이 고장난 줄 안다.
    """
    _load_snmp_enums()
    n = (q or "").strip().lower()
    lim = max(1, min(300, int(limit or 50)))
    out = []
    for oid, name in SNMP_OID_NAMES.items():
        if n and n not in oid.lower() and n not in str(name).lower():
            continue
        out.append({"oid": oid, "name": name})
        if len(out) >= lim:
            break
    # 이름 순이 사람이 찾기 좋다. OID 숫자순은 트리 구조라 눈에 안 들어온다.
    out.sort(key=lambda x: str(x["name"]))
    return {
        "oids": out,
        "total": len(SNMP_OID_NAMES),
        "source": "data/snmp/snmp_names.json",
        "hint": "" if SNMP_OID_NAMES else "MIB 를 아직 안 뽑았습니다 — data/MIB/ 에 파일을 넣고 tools/mib_enums.py 를 돌리세요",
    }


_SNMP_WCOMM_CACHE: dict = {}   # host -> 최근에 SET 이 통한 쓰기 커뮤니티


async def _snmp_instances(host: str, comm: str, mp: int, col_oid: str, limit: int = 12):
    """그 열(column)에 **실제로 있는 인스턴스 번호**를 몇 개 — 진단용.

    장비의 ifIndex 는 1·2·3 이 아니라 101·102·112·1003… 처럼 띄엄띄엄한 경우가
    많다(지적: MIB 브라우저로 보면 그렇다). 그때 `.8.2` 로 SET 하면 noSuchName
    인데, 「없다」 만으로는 **무슨 번호를 써야 하는지** 알 수 없다. 있는 번호를
    같이 보여 주면 반복을 그 값으로 맞출 수 있다.
    """
    out = []
    try:
        from pysnmp.hlapi.v3arch.asyncio import (
            SnmpEngine, CommunityData, UdpTransportTarget, ContextData,
            ObjectType, ObjectIdentity, walk_cmd)
        try:
            from pysnmp.hlapi.v3arch.asyncio import bulk_walk_cmd as _bw
        except Exception:
            _bw = None
        eng = SnmpEngine()
        tr = await UdpTransportTarget.create((host, 161), timeout=1.2, retries=0)
        auth = CommunityData(comm, mpModel=mp)
        base = col_oid.strip().lstrip(".")
        walker = (_bw(eng, auth, tr, ContextData(), 0, 25, ObjectType(ObjectIdentity(base)), lexicographicMode=False)
                  if _bw is not None and mp == 1 else
                  walk_cmd(eng, auth, tr, ContextData(), ObjectType(ObjectIdentity(base)), lexicographicMode=False))
        async for (ei, es2, ex, vbs) in walker:
            if ei or es2:
                break
            for vb in vbs:
                nm = _oid_to_num(vb[0].prettyPrint()).lstrip(".")
                if nm.startswith(base + "."):
                    out.append(nm[len(base) + 1:])
            if len(out) >= limit:
                break
        _snmp_close(eng)
    except Exception:
        pass
    return out[:limit]



async def _snmp_write_comms(host: str) -> list:
    """SET 에 시도할 **쓰기 커뮤니티 후보**를 차례로.

    장비마다 관례가 다르다: 어떤 건 읽기 public·쓰기 private, 어떤 건
    public 하나로 읽기·쓰기 다 된다(public-RW). 하나만 골라 보내면 한쪽
    장비에서 늘 막힌다(지적: 수동 .8.2 는 되는데 도구는 noSuchName). 그래서
    **여러 개를 차례로** 시도한다 — 실패한 SET 은 장비를 바꾸지 않으니 안전하다.

    차례: 등록된 쓰기 커뮤니티 → 등록된 읽기 커뮤니티 → private → public.
    """
    out = []
    try:
        host = (host or "").strip()
        ro = wo = ""
        for d in await db.device_list(with_ifs=False):
            if str(d.get("ip") or "").strip() != host:
                continue
            snmp = _acc_of(d, "snmp")
            params = snmp.get("params") or {}
            for _ in range(2):
                if isinstance(params, str):
                    try:
                        params = json.loads(params)
                    except Exception:
                        params = {}
            if not isinstance(params, dict):
                params = {}
            ro = snmp.get("username") or snmp.get("community") or ""
            wo = params.get("community_rw") or ""
            break
        _cached = _SNMP_WCOMM_CACHE.get(host)
        for c in (_cached, wo, ro, "private", "public"):
            if c and c not in out:
                out.append(c)
    except Exception:
        out = ["private", "public"]
    return out or ["private", "public"]


async def _snmp_comm_for(host: str, rw: bool):
    """이 IP 장비에 **저장된 커뮤니티**를 찾는다.

    읽기(public)는 되고 쓰기만 noAccess 로 막히던 까닭이 여기 있었다(지적):
    SNMP Set 이 늘 기본값 'private' 로 나갔는데, 장비의 쓰기 커뮤니티는 따로
    등록돼 있다(장비 SNMP 줄의 `community_rw`). 그 값을 꺼내 쓴다.

    rw=True 면 쓰기 커뮤니티를 먼저, 없으면 읽기 커뮤니티, 그것도 없으면
    None(부르는 쪽이 기본값을 쓴다).
    """
    try:
        host = (host or "").strip()
        if not host:
            return None
        for d in await db.device_list(with_ifs=False):
            if str(d.get("ip") or "").strip() != host:
                continue
            snmp = _acc_of(d, "snmp")
            params = snmp.get("params") or {}
            # params 가 문자열(때로 이중 인코딩)로 저장된 자료가 있다 — 풀어 준다
            for _ in range(2):
                if isinstance(params, str):
                    try:
                        params = json.loads(params)
                    except Exception:
                        params = {}
            if not isinstance(params, dict):
                params = {}
            ro = snmp.get("username") or snmp.get("community") or ""
            wo = params.get("community_rw") or ""
            if rw:
                # 쓰기는 **쓰기 커뮤니티만** 쓴다. 읽기 커뮤니티로 되돌아가면
                # 안 된다(지적: 수동은 private 로 되는데 도구는 안 된다) —
                # 읽기 community 가 public 이면 SET 을 public 으로 보내 noAccess
                # 가 난다. 없으면 None → 부르는 쪽이 관례값 'private' 를 쓴다.
                return (wo or None)
            return (ro or None)
    except Exception:
        pass
    return None


@app.post("/api/snmp-get")
async def snmp_get_api(payload: dict):
    _load_snmp_enums()   # JSON(MIB 추출) 변경 시 자동 재로드 → mib_enums.py 재실행만으로 반영(서버 재시작 불필요)
    eng = None
    host = (payload.get("host") or "").strip()
    oid = (payload.get("oid") or "").strip()
    community = payload.get("community") or await _snmp_comm_for(host, rw=False) or "public"
    ver = (payload.get("version") or "v2c").lower()
    try:
        port = int(payload.get("port", 161) or 161)
    except Exception:
        port = 161
    if not host:
        return {"ok": False, "error": "host(IP)가 없습니다", "output": ""}
    if not oid:
        return {"ok": False, "error": "OID가 없습니다", "output": ""}
    mp = 0 if ver == "v1" else 1
    mode = (payload.get("mode") or "auto").lower()   # auto(GET→없으면 WALK) | get | walk
    try:
        from pysnmp.hlapi.v3arch.asyncio import (
            SnmpEngine, CommunityData, UdpTransportTarget, ContextData,
            ObjectType, ObjectIdentity, get_cmd, walk_cmd)
        try:
            from pysnmp.hlapi.v3arch.asyncio import bulk_walk_cmd as _bulk_walk_cmd
        except Exception:
            _bulk_walk_cmd = None
        eng = SnmpEngine()
        # 첫 패킷이 늦으면 이 시간을 다 기다린다 — 가끔 3s 씩 튀던 까닭이다
        # (지적). 짧게 잡고 재시도 1 로 유실만 메꾼다.
        transport = await UdpTransportTarget.create((host, port), timeout=1.2, retries=1)
        auth = CommunityData(community, mpModel=mp)

        async def _do_walk():
            rows = []
            try:
                # v2c(mp==1) 면 GETBULK — 한 번에 여러 행을 받아 왕복·유실을 줄인다.
                # v1 이나 미지원이면 GETNEXT(walk_cmd)로 떨어진다.
                if _bulk_walk_cmd is not None and mp == 1:
                    _walker = _bulk_walk_cmd(
                        eng, auth, transport, ContextData(),
                        0, 25,
                        ObjectType(ObjectIdentity(oid)), lexicographicMode=False)
                else:
                    _walker = walk_cmd(
                        eng, auth, transport, ContextData(),
                        ObjectType(ObjectIdentity(oid)), lexicographicMode=False)
                async for (eInd, eStat, eIdx, vbs) in _walker:
                    if eInd or eStat:
                        break
                    for vb in vbs:
                        try:
                            nm = vb[0].prettyPrint(); val = _snmp_enum(_oid_to_num(nm), vb[1], _snmp_val_str(vb[1]))
                        except Exception:
                            nm = str(vb); val = ""
                        if _snmp_is_noinstance(val):
                            continue
                        rows.append(nm + " = " + val)
                    if len(rows) >= 500:
                        break
            except Exception:
                pass
            return rows

        if mode == "walk":
            rows = await _do_walk()
            if rows:
                return {"ok": True, "output": "\n".join(rows), "count": len(rows), "mode": "walk"}
            return {"ok": False, "error": "WALK 결과 없음", "output": "[SNMP] " + oid + " 하위에 데이터가 없습니다"}

        # GET 먼저
        errInd, errStat, errIdx, varBinds = await get_cmd(
            eng, auth, transport, ContextData(), ObjectType(ObjectIdentity(oid)))
        if errInd:
            return {"ok": False, "error": str(errInd), "output": "[SNMP] " + str(errInd)}
        if errStat:
            return {"ok": False, "error": str(errStat.prettyPrint()), "output": "[SNMP] " + str(errStat.prettyPrint())}
        lines = []; noinst = False
        for vb in varBinds:
            try:
                nm = vb[0].prettyPrint(); val = _snmp_enum(_oid_to_num(nm), vb[1], _snmp_val_str(vb[1]))
                lines.append(nm + " = " + val)
                if _snmp_is_noinstance(val):
                    noinst = True
            except Exception:
                lines.append(str(vb))
        # GET이 No Such Instance/Object → 테이블 컬럼일 가능성 → WALK 자동 폴백
        if noinst and mode == "auto":
            rows = await _do_walk()
            if rows:
                return {"ok": True, "output": "\n".join(rows), "count": len(rows), "mode": "walk"}
            return {"ok": False, "error": "No Such Instance — 이 OID에 인스턴스가 없습니다 (스칼라는 끝에 .0, 테이블은 인덱스 필요)",
                    "output": "[SNMP] " + ("\n".join(lines) if lines else oid + " : No Such Instance")}
        if noinst:
            return {"ok": False, "error": "No Such Instance/Object",
                    "output": "[SNMP] " + ("\n".join(lines) if lines else "(빈 응답)")}
        return {"ok": True, "output": "\n".join(lines) if lines else "(빈 응답)", "mode": "get"}
    except Exception as e:
        _msg = str(e)
        if "No module named" in _msg and ("pysnmp" in _msg or "pyasn1" in _msg):
            _msg = "pysnmp 미설치 — 백엔드에서 'python -m pip install pysnmp' 실행 후 서버 재시작"
        return {"ok": False, "error": _msg[:300], "output": "[SNMP 오류] " + _msg[:220]}
    finally:
        _snmp_close(eng)   # 소켓 정리(FD 누수 방지)

@app.post("/api/snmp-set")
async def snmp_set_api(payload: dict):
    host = (payload.get("host") or "").strip()
    oid = (payload.get("oid") or "").strip().lstrip(".")
    value = payload.get("value")
    value = "" if value is None else str(value)
    # 쓰기 커뮤니티는 장비에 등록된 것을 먼저 쓴다(지적: noAccess). 없으면 private.
    # 명시했으면 그것만. 아니면 여러 쓰기 커뮤니티를 차례로 시도한다(지적).
    _explicit = payload.get("community")
    comm_cands = [_explicit] if _explicit else (await _snmp_write_comms(host))
    community = comm_cands[0]
    ver = (payload.get("version") or "v2c").lower()
    vtype = (payload.get("type") or "").strip().lower()   # 선택: i/s/u/a … 없으면 자동(숫자→정수, 그 외→문자열)
    _load_snmp_enums()                                     # enum 맵 최신화
    _ev = _snmp_set_enum(oid, value)                       # enum 이름(fiveMin)이면 정수(3)로 변환 → 숫자 흐름으로
    if _ev is not None:
        value = _ev
    eng = None
    try:
        port = int(payload.get("port", 161) or 161)
    except Exception:
        port = 161
    if not host:
        return {"ok": False, "error": "host(IP)가 없습니다", "output": ""}
    if not oid:
        return {"ok": False, "error": "OID가 없습니다", "output": ""}
    mp = 0 if ver == "v1" else 1
    try:
        from pysnmp.hlapi.v3arch.asyncio import (
            SnmpEngine, CommunityData, UdpTransportTarget, ContextData,
            ObjectType, ObjectIdentity, set_cmd)
        from pysnmp.proto.rfc1902 import Integer32, OctetString, Unsigned32, IpAddress, Counter32, Gauge32, TimeTicks
        def _mkval(tt, val):
            if tt in ("i", "int", "integer"):
                return Integer32(int(val))
            if tt in ("u", "uint", "unsigned"):
                return Unsigned32(int(val))
            if tt in ("g", "gauge", "gauge32"):
                return Gauge32(int(val))
            if tt in ("c", "counter", "counter32"):
                return Counter32(int(val))
            if tt in ("t", "ticks", "timeticks"):
                return TimeTicks(int(val))
            if tt in ("a", "ip", "ipaddress"):
                return IpAddress(val)
            if tt in ("x", "hex"):
                return OctetString(hexValue=val.replace(" ", "").replace("0x", ""))
            return OctetString(val)
        # 후보 타입: 명시 타입 있으면 그것만. 없으면 숫자→[Integer32→Unsigned32→Counter32→Gauge32] wrongType 시 자동 재시도, 그 외→문자열
        if vtype:
            cands = [vtype]
        else:
            _digits = value.lstrip("-")
            cands = ["i", "u", "c", "g"] if (_digits.isdigit() and value not in ("", "-")) else ["s"]
        eng = SnmpEngine()
        # 안 맞는 커뮤니티는 응답이 없어 타임아웃까지 매달린다 — 짧게(지적:
        # 시험 진행 중 갑자기 느려진다). 되는 커뮤니티는 캐시로 첫 시도에 맞는다.
        transport = await UdpTransportTarget.create((host, port), timeout=1.5, retries=0)
        last_err = None
        # OID 후보 — 스칼라는 인스턴스 `.0` 을 찍어야 SET 이 먹는다.
        # `…3.6`(객체)으로 SET 하면 딱 noAccess 가 난다(지적: 커뮤니티도 RW 인데
        # noAccess). 수동 snmpset 은 `…3.6.0` 을 쓴다. 인스턴스가 없어 보이면
        # `.0` 을 붙인 것도 후보에 넣어, noAccess/noSuchName 이면 그것으로 다시 건다.
        _bare = oid.rstrip(".")
        oid_cands = [oid]
        _last_arc = _bare.rsplit(".", 1)[-1] if "." in _bare else ""
        if _last_arc != "0":
            oid_cands.append(_bare + ".0")
        # 시도 목록을 **한 겹으로 펼친다** — 버전 × 커뮤니티 × OID × 타입.
        # 중첩을 쌓으면 어디서 빠져나왔는지 알 수 없고, 보고하는 오류도 마지막
        # 시도 것이 되어 엉뚱한 곳을 가리킨다(지적: private RW 인데 거부).
        #
        # 버전: 명시했으면 그것만. 아니면 v2c 뒤에 v1 도 본다 — 수동 snmpset 은
        # 버전을 안 주면 흔히 v1 로 나가고, 그것만 쓰기를 받는 장비가 있다.
        ver_cands = [ver] if payload.get("version") else ([ver, "v1"] if ver != "v1" else ["v1"])
        attempts = []
        for _vr in ver_cands:
            for _cm in comm_cands:
                for _od in oid_cands:
                    for _tt in cands:
                        attempts.append((_vr, _cm, _od, _tt))
        _used_oid, _used_comm, _used_ver = oid, community, ver
        _first_err = None       # 원래 OID·첫 커뮤니티의 오류 — 보고는 이걸로
        _tried = []             # 무엇을 어떻게 보냈는지 (진단용)
        _seen_err = {}          # (ver,comm) -> 마지막 오류. 같은 짝을 헛돌지 않게
        _need_zero = False      # noSuchName 을 봤을 때만 `.0` 을 시도한다
        _oid_class = False      # 오류가 **OID 문제**로 보이나 (noSuchName …)
        _acc_class = False      # 오류가 **권한 문제**로 보이나 (noAccess …)
        for _vr, _cm, _od, tt in attempts:
            # **필요할 때만 넓힌다** — 조합을 다 돌면 응답 없는 커뮤니티에서
            # 1.5초씩 먹어 시험이 느려진다(앞서 겪은 것). 규칙:
            #  · `.0` 곁가지는 noSuchName(인스턴스 없음)을 봤을 때만
            #  · 같은 (버전·커뮤니티)에서 wrongType 이 아니면 다른 타입은 무의미
            _prev = _seen_err.get((_vr, _cm))
            # 오류 **갈래대로만** 넓힌다(지적: SET 동작이 이상하다 — 8조합을
            # 헛돌았다). noSuchName 은 OID 문제라 커뮤니티·버전을 바꿔도 소용이
            # 없고, noAccess 는 권한 문제라 OID 를 바꿔도 소용이 없다.
            if _oid_class and (_vr, _cm) != (ver_cands[0], comm_cands[0]):
                continue        # OID 문제 — 커뮤니티·버전은 건드릴 이유가 없다
            if _acc_class and _od != oid:
                continue        # 권한 문제 — `.0` 곁가지는 뜻이 없다
            if _od != oid and not _need_zero:
                continue
            if _prev is not None and "wrongType" not in _prev and tt != cands[0]:
                continue
            _used_oid, _used_comm, _used_ver = _od, _cm, _vr
            try:
                pv = _mkval(tt, value)
            except Exception as ex:
                last_err = str(ex); continue
            auth = CommunityData(_cm, mpModel=(0 if _vr == "v1" else 1))
            errInd, errStat, errIdx, varBinds = await set_cmd(
                eng, auth, transport, ContextData(), ObjectType(ObjectIdentity(_od), pv))
            if errInd:
                return {"ok": False, "error": str(errInd), "output": "[SNMP SET] " + str(errInd)}
            if errStat:
                es = str(errStat.prettyPrint()); last_err = es
                _sig = _vr + "/" + str(_cm) + " " + _od + " (" + _SNMP_TYPE_NAMES.get(tt, tt) + ")"
                if _sig not in _tried:
                    _tried.append(_sig)
                if _first_err is None:
                    _first_err = (es, _od, _cm, _vr, tt)
                _seen_err[(_vr, _cm)] = es
                if "noSuchName" in es or "noSuchInstance" in es:
                    _need_zero = True     # 인스턴스 문제 — `.0` 을 붙여 볼 값이 있다
                    _oid_class = True
                elif any(x in es for x in ("noAccess", "authorizationError", "notWritable", "readOnly")):
                    _acc_class = True
                continue        # 다음 시도로 — 다 해 보고 아래에서 보고한다
            # ── 성공 ──
            lines2 = []
            for vb in varBinds:
                try:
                    lines2.append(vb[0].prettyPrint() + " = " + _snmp_val_str(vb[1]))
                except Exception:
                    lines2.append(str(vb))
            _tn = _SNMP_TYPE_NAMES.get(tt, tt)
            _note = ""
            if _used_oid != oid:
                _note += "\n\u2192 \uc778\uc2a4\ud134\uc2a4 `.0` \uc744 \ubd99\uc5ec \uc131\uacf5\ud588\uc2b5\ub2c8\ub2e4 (" + _used_oid + ")."
            if not _explicit and _used_comm != comm_cands[0]:
                _note += "\n\u2192 \uc4f0\uae30 \ucee4\ubba4\ub2c8\ud2f0 '" + str(_used_comm) + "' \ub85c \ub410\uc2b5\ub2c8\ub2e4. Devices \uc5d0 \ub123\uc5b4 \ub450\uba74 \ub2e4\uc74c\ubd80\ud134 \ubc14\ub85c \ub429\ub2c8\ub2e4."
            if _used_ver != ver:
                _note += "\n\u2192 SNMP " + _used_ver + " \ub85c \ub410\uc2b5\ub2c8\ub2e4 (v2c \ub294 \uac70\ubd80). \uc7a5\ube44 SNMP \uc124\uc815\uc758 \ubc84\uc804\uc744 " + _used_ver + " \ub85c \ub450\uc138\uc694."
            if not _explicit:
                _SNMP_WCOMM_CACHE[host] = _used_comm
            return {"ok": True, "output": "[SNMP SET OK] (type=" + _tn + ")\n" + ("\n".join(lines2) if lines2 else (_used_oid + " = " + value)) + _note, "mode": "set"}

        # ── 다 실패 ── 원래 OID·첫 시도의 오류로 보고한다(마지막 것은 `.0` 등 곁가지다)
        if _first_err:
            es, _eo, _ec, _ev2, _et = _first_err
            _tnn = _SNMP_TYPE_NAMES.get(_et, _et)
            _hint = "\n\u2192 \ubcf4\ub0b8 \uac12: [" + str(value) + "] \ud0c0\uc785: " + _tnn
            _hint += "\n\u2192 \ubcf4\ub0b8 OID: " + _eo + " · \ucee4\ubba4\ub2c8\ud2f0: " + str(_ec) + " · \ubc84\uc804: " + _ev2
            try:
                _ps = oid.split("."); _em2 = None
                for _cut in (1, 2, 0):
                    _cand = ".".join(_ps[: len(_ps) - _cut]) if _cut else oid
                    if SNMP_ENUM_MAP.get(_cand):
                        _em2 = SNMP_ENUM_MAP[_cand]; break
                if _em2:
                    _hint += "\n\u2192 \uc720\ud6a8\uac12: " + ", ".join(nm + "(" + k + ")" for k, nm in sorted(_em2.items(), key=lambda x: int(x[0])))
            except Exception:
                pass
            if "wrongType" in es:
                _hint += "\n\u2192 \ud0c0\uc785 \ubd88\uc77c\uce58. [i:" + value + "]\u00b7[u:" + value + "]\u00b7[s:..]\u00b7[x:HEX] \ub85c \uc9c0\uc815 \uac00\ub2a5"
            elif "noSuchName" in es or "noSuchInstance" in es:
                # 있는 번호를 실제로 물어봐서 알려 준다 — 「없다」 만으로는
                # 무슨 번호를 써야 할지 알 수 없다(지적: 인덱스가 101·1003…)
                try:
                    _col = oid.rstrip(".").rsplit(".", 1)[0]
                    _have = await _snmp_instances(host, _ec, (0 if _ev2 == "v1" else 1), _col)
                    if _have:
                        _hint += ("\n\u2192 \uc774 \uc5f4\uc5d0 **\uc2e4\uc81c\ub85c \uc788\ub294 \ubc88\ud638**: "
                                  + " \u00b7 ".join(_have) + " \u2026  \ubc18\ubcf5\uc744 \uc774 \uac12\uc73c\ub85c \ub450\uc138\uc694(\ubaa9\ub85d \ubc29\uc2dd).")
                except Exception:
                    pass
                _hint += ("\n\u2192 \uc774 \uc7a5\ube44\ub294 **\uadf8 \ubc88\ud638\uc5d0 \uc4f0\uae30\ub97c \uc548 \ubc1b\uc2b5\ub2c8\ub2e4**(" + es + "). "
                          "\uc77d\uae30(GET)\ub294 \ub418\ub294\ub370 SET \ub9cc \uc774\ub7ec\uba74, \uadf8 \ubc88\ud638\uac00 \uc4f0\uae30 \ub300\uc0c1\uc774 \uc544\ub2c8\uac70\ub098 "
                          "CLI \ud3ec\ud2b8 \ubc88\ud638\uc640 SNMP ifIndex \uac00 \ub2e4\ub978 \uacbd\uc6b0\uc785\ub2c8\ub2e4 \u2014 \uc218\ub3d9\uc73c\ub85c \ub41c \ubc88\ud638\uc640 \uacac\uc918 \ubcf4\uc138\uc694.")
            elif any(x in es for x in ("noAccess", "authorizationError", "notWritable", "readOnly")):
                _hint += ("\n\u2192 \uc7a5\ube44\uac00 \uc4f0\uae30\ub97c \uac70\ubd80\ud588\uc2b5\ub2c8\ub2e4(" + es + "). \uc218\ub3d9 snmpset \uc774 \ub41c\ub2e4\uba74 \uadf8\ub54c\uc758 "
                          "**OID\u00b7\ucee4\ubba4\ub2c8\ud2f0\u00b7\ubc84\uc804**\uc744 \uc704 \uc904\uacfc \uacac\uc918 \ubcf4\uc138\uc694 \u2014 \ud558\ub098\ub77c\ub3c4 \ub2e4\ub974\uba74 \uadf8\uac83\uc774 \uae30\uc900\uc785\ub2c8\ub2e4. "
                          "\uc7a5\ube44 \ucabd ACL(\ud5c8\uc6a9 IP)\uc774 \uc788\uc73c\uba74 \uc774 \uc11c\ubc84 IP \ub3c4 \ub123\uc5b4\uc57c \ud569\ub2c8\ub2e4.")
            elif "genErr" in es:
                _hint += "\n\u2192 genErr = \uc7a5\ube44\uac00 SET \uac70\ubd80. \ud574\ub2f9 \ud3ec\ud2b8 \uc0c1\ud0dc\u00b7\uc4f0\uae30\uad8c\ud55c\uc744 \ud655\uc778\ud558\uc138\uc694"
            if len(_tried) > 1:
                _hint += "\n\u2192 \uc2dc\ub3c4: " + " · ".join(_tried[:8]) + ("  \u2026" if len(_tried) > 8 else "") + " (\ub2e4 \uac70\ubd80)"
            return {"ok": False, "error": es, "output": "[SNMP SET \uc2e4\ud328] " + es + _hint}
        return {"ok": False, "error": last_err or "SET 실패", "output": "[SNMP SET 오류] " + str(last_err or "")}
    except Exception as e:
        _msg = str(e)
        if "No module named" in _msg and ("pysnmp" in _msg or "pyasn1" in _msg):
            _msg = "pysnmp 미설치 — 'python -m pip install pysnmp' 후 서버 재시작"
        return {"ok": False, "error": _msg[:300], "output": "[SNMP SET 오류] " + _msg[:220]}
    finally:
        _snmp_close(eng)   # 소켓 정리(FD 누수 방지)

# ── SNMP Trap 수신기 (장비가 보내는 Notification 수신·판정용) ──
_TRAP_BUF = []           # [{ts, from, oid, varbinds:[{oid,value}]}]
_TRAP_LOCK = _threading.Lock()
_trap_state = {"started": False, "error": "", "port": 162}

def _trap_listener_thread(port):
    try:
        import asyncio
        from pysnmp.entity import engine, config
        from pysnmp.carrier.asyncio.dgram import udp
        from pysnmp.entity.rfc3413 import ntfrcv
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        snmpEngine = engine.SnmpEngine()
        config.add_transport(snmpEngine, udp.DOMAIN_NAME,
                             udp.UdpTransport().open_server_mode(('0.0.0.0', port)))
        # v1/v2c community (수신은 community 검증 느슨하게 — 흔한 public/private 등록)
        for comm in ("public", "private", "ubiquoss"):
            try:
                config.add_v1_system(snmpEngine, "utop-" + comm, comm)
            except Exception:
                pass

        def cbFun(snmpEngine, stateReference, contextEngineId, contextName, varBinds, cbCtx):
            src = ""
            try:
                td, ta = snmpEngine.message_dispatcher.get_transport_info(stateReference)
                src = str(ta[0])
            except Exception:
                pass
            vbs, trap_oid = [], ""
            for oid, val in varBinds:
                try:
                    o = oid.prettyPrint(); v = val.prettyPrint()
                except Exception:
                    o, v = str(oid), str(val)
                vbs.append({"oid": o, "value": v})
                if o.endswith("1.3.6.1.6.3.1.1.4.1.0"):  # snmpTrapOID.0
                    trap_oid = v
            with _TRAP_LOCK:
                _TRAP_BUF.append({"ts": _t.time(), "from": src, "oid": trap_oid, "varbinds": vbs})
                if len(_TRAP_BUF) > 500:
                    del _TRAP_BUF[:len(_TRAP_BUF) - 500]

        ntfrcv.NotificationReceiver(snmpEngine, cbFun)
        snmpEngine.transport_dispatcher.job_started(1)
        _trap_state["started"] = True
        _trap_state["error"] = ""
        snmpEngine.transport_dispatcher.run_dispatcher()
    except Exception as e:
        _trap_state["error"] = str(e)[:300]
        _trap_state["started"] = False

def _ensure_trap_listener(port):
    if _trap_state.get("started"):
        return True
    _trap_state["port"] = port
    th = _threading.Thread(target=_trap_listener_thread, args=(port,), daemon=True)
    th.start()
    for _ in range(25):
        if _trap_state.get("started"):
            return True
        if _trap_state.get("error"):
            return False
        _t.sleep(0.1)
    return _trap_state.get("started", False)





@app.post("/api/snmp-trap/wait")
async def snmp_trap_wait(payload: dict):
    import asyncio
    oid = (payload.get("oid") or "").strip()
    try:
        timeout = float(payload.get("timeout", 15) or 15)
    except Exception:
        timeout = 15
    try:
        port = int(payload.get("port", 162) or 162)
    except Exception:
        port = 162
    ok = _ensure_trap_listener(port)
    if not ok:
        return {"ok": False, "error": "Trap 수신기 시작 실패: " + (_trap_state.get("error") or ("UDP " + str(port) + " 바인드 불가 — 관리자 권한 또는 포트 사용중 확인"))}
    with _TRAP_LOCK:
        base = len(_TRAP_BUF)
    start = _t.time()
    while _t.time() - start < timeout:
        with _TRAP_LOCK:
            for tr in _TRAP_BUF[base:]:
                if (not oid) or (oid in (tr.get("oid") or "")) or any(oid in (v.get("oid") or "") for v in tr.get("varbinds", [])):
                    return {"ok": True, "trap": tr}
        await asyncio.sleep(0.3)
    return {"ok": True, "trap": None, "error": "timeout"}

CYCLE_META_KEYS = ("id","name","model","version","version_group","folder_id","assignee","start_date","end_date","mail_send","created_at","updated_at")

def _cycle_item_meta_lite(it: dict) -> dict:
    """UI 목록·판정 집계에 필요한 최소 필드만.
    result 값 카운트로만 집계 → steps 배열 자체를 안 보내고 pass/fail/total 3개 숫자만 반환."""
    m = {k: it.get(k) for k in ("tcid","name","req_id","severity","priority","assignee","devId","devName","executed_by","executed_at","executed_auto","issues")}
    _stp = it.get("steps") or []
    _p = 0; _f = 0; _o = 0
    for s in _stp:
        r = s.get("result") or ""
        if r == "Pass": _p += 1
        elif r == "Fail": _f += 1
        elif r: _o += 1
    m["_steps_count"] = len(_stp)
    m["_steps_pass"] = _p
    m["_steps_fail"] = _f
    m["_steps_other"] = _o
    # 각 step 은 판정 집계에 필요한 최소 필드만 (result, action, manual, cli 첫줄).
    # output/verdictMsg 등 무거운 필드는 완전 제외.
    def _lite(s):
        return {
            "result": s.get("result",""),
            "action": s.get("action",""),
            "manual": bool(s.get("manual")),
        }
    m["steps"] = [_lite(s) for s in _stp]
    return m

def _cycle_meta_extra(meta: dict, d: dict):
    _items = d.get("items") or []
    meta["items"] = [_cycle_item_meta_lite(it) for it in _items]

@app.get("/api/cycle")
async def get_all_cycles(meta: int = 0):
    """
    meta=1 : 목록·판정 집계에 필요한 필드만 반환 (각 item.steps 에서 output/verdictMsg 등 큰 필드 제거)
    meta=0 (기본): 전체 반환
    """
    if meta:
        # meta 모드: cycle.data - items (items 는 통째 제외). 프론트가 다시 loadCycleFull 로 개별 로드.
        return {"cycles": await db.cycle_list_meta()}
    return {"cycles": await db.cycle_list_full()}

# ───────────────────────────────────────────
# 사이클 트리 집계 — 폴더 한 층의 현황을 **서버가** 센다.
#
# 여태 화면이 회차·항목을 다 받아 브라우저에서 셌다. 사업자 층은 수천
# 건이라 트리 위로 갈수록 느려진다. 여기서 한 번에 세어 내려준다.
#
# 트리 경로는 화면(pathOfCycle)과 **같은 규칙**이다:
#   Root/사업자/제품군/모델그룹/모델명/버전그룹
# 폴더를 손으로 정해 둔 회차는 Root/<그 경로> 에 그대로 붙는다.
# ───────────────────────────────────────────
_RU_ROOT = "Root"
_RU_NO_CUST = "(사업자 없음)"
_RU_NO_CAT = "(카탈로그에 없는 모델)"
_RU_NO_MGROUP = "(모델그룹 없음)"
_RU_NO_GROUP = "(버전그룹 없음)"
_RU_LEVELS = ["root", "operator", "family", "model_group", "model", "version_group", "cycle"]


def _ru_path(c: dict, fam: dict, mgrp: dict) -> str:
    own = str(c.get("folder") or "").strip().strip("/")
    if own:
        return f"{_RU_ROOT}/{own}"
    model = str(c.get("model") or "").strip() or "(모델 없음)"
    cust = str(c.get("customer") or "").strip() or _RU_NO_CUST
    f = (fam.get(model) or "(제품군 없음)") if model in fam else _RU_NO_CAT
    mg = str(c.get("model_group") or "").strip() or mgrp.get(model) or _RU_NO_MGROUP
    vg = str(c.get("version_group") or "").strip() or _RU_NO_GROUP
    return f"{_RU_ROOT}/{cust}/{f}/{mg}/{model}/{vg}"


def _ru_day(s) -> str:
    """항목이 남긴 시각에서 날짜만 — 「2026-08-19T14:21」 → 「2026-08-19」"""
    t = str(s or "")[:10]
    return t if len(t) == 10 and t[4] == "-" else ""


async def _ru_groups() -> dict:
    """판정 → 집계 계열. 설정 「실행 판정 기준」 이 정본이다."""
    g = {"Pass": "pass", "Fail": "fail", "": "none"}
    try:
        for it in await db.code_list("cycle_result"):
            v = str(it.get("value") or "")
            try:
                meta = json.loads(it.get("note") or "{}")
            except Exception:
                meta = {}
            grp = meta.get("group")
            g[v] = grp if grp in ("pass", "fail") else g.get(v, "neutral")
    except Exception:
        pass
    return g


async def _rollup(path: str = _RU_ROOT, date_from: str = "", date_to: str = "",
                  axis: str = "") -> dict:
    """
    폴더 한 층의 현황. 프리뷰의 KPI·막대·표·추이가 모두 이 하나를 쓴다.

    · `path`      — 「Root/LGUPLUS/L3」 처럼 트리 경로
    · `date_from` / `date_to` — 항목이 실행된 날(YYYY-MM-DD) 로 자른다.
                    비우면 전부. 잘라도 **회차 수·항목 수는 그대로**고,
                    판정 집계와 추이만 그 기간 것으로 센다.

    합격률은 **합격 ÷ (합격+실패)** 다 — 실행한 것 중 합격 비율.
    미실행이 얼마나 남았는지는 진척률(실행/전체)이 따로 말한다.
    """
    base = str(path or _RU_ROOT).strip().strip("/") or _RU_ROOT
    metas = await db.cycle_list_meta()
    cat = await db.catalog_list("model")
    fam = {str(m.get("name") or ""): str(m.get("family") or "").strip() for m in cat}
    mgrp = {str(m.get("name") or ""): str(m.get("model_group") or "").strip() for m in cat}
    grp_of = await _ru_groups()

    depth = len(base.split("/"))
    level = _RU_LEVELS[depth] if depth < len(_RU_LEVELS) else "cycle"

    def tally() -> dict:
        return {"n": 0, "pass": 0, "fail": 0, "other": 0, "none": 0, "cycles": 0,
                "last_run": "", "open_defects": 0}

    total = tally()
    kids: dict[str, dict] = {}
    axes: dict[str, dict] = {}
    trend: dict[str, dict] = {}
    rows: list[dict] = []

    for c in metas:
        p = _ru_path(c, fam, mgrp)
        if p != base and not p.startswith(base + "/"):
            continue
        rest = p[len(base):].strip("/")
        key = rest.split("/")[0] if rest else (str(c.get("cid") or c.get("id") or ""))
        kid = kids.setdefault(key, {**tally(), "key": key, "leaf": not rest})
        kid["cycles"] += 1
        total["cycles"] += 1

        # 축 열쇠 — 「무엇으로 나눠 볼까」. 회차에서 오는 것과 항목에서
        # 오는 것이 있어 둘 다 받는다.
        cyc_key = {
            "cycle": str(c.get("cid") or c.get("name") or c.get("id") or "–"),
            "version_group": str(c.get("version_group") or "(버전그룹 없음)"),
            "model": str(c.get("model") or "(모델 없음)"),
            "customer": str(c.get("customer") or "(고객 없음)"),
            "status": str(c.get("status") or "(상태 없음)"),
        }.get(axis, "")

        cy = tally()
        for it in (c.get("items") or []):
            if not isinstance(it, dict):
                continue
            day = _ru_day(it.get("executed_at"))
            if date_from and day and day < date_from:
                continue
            if date_to and day and day > date_to:
                continue
            v = str(it.get("_verdict") or it.get("result") or "")
            g = grp_of.get(v, "neutral" if v else "none")
            cy["n"] += 1
            cy[g if g in ("pass", "fail", "none") else "other"] += 1
            if day and day > cy["last_run"]:
                cy["last_run"] = day
            if g in ("pass", "fail") and day:
                wk = trend.setdefault(day[:7] + "-" + str((int(day[8:10]) - 1) // 7 + 1),
                                      {"k": "", "pass": 0, "fail": 0})
                wk["k"] = day
                wk[g] += 1
            if it.get("issues"):
                cy["open_defects"] += len(it.get("issues") or [])

            if axis:
                if cyc_key:
                    ak = cyc_key
                else:
                    raw = it.get("severity") if axis == "severity" else it.get("assignee")
                    ak = str(raw or "").strip() or "(없음)"
                ax = axes.setdefault(ak, {**tally(), "key": ak})
                ax["n"] += 1
                ax[g if g in ("pass", "fail", "none") else "other"] += 1
                if day and day > ax["last_run"]:
                    ax["last_run"] = day

        for f in ("n", "pass", "fail", "other", "none", "open_defects"):
            kid[f] += cy[f]
            total[f] += cy[f]
        if cy["last_run"] > kid["last_run"]:
            kid["last_run"] = cy["last_run"]
        if cy["last_run"] > total["last_run"]:
            total["last_run"] = cy["last_run"]

        if not rest:  # 이 층이 곧 회차 목록이다(버전그룹 아래)
            rows.append({
                "id": c.get("id"), "cid": c.get("cid"), "name": c.get("name"),
                "version": c.get("version"), "version_group": c.get("version_group"),
                "model": c.get("model"), "status": c.get("status"),
                "assignee": c.get("assignee"), "end_date": c.get("end_date"),
                **{k: cy[k] for k in ("n", "pass", "fail", "other", "none", "last_run")},
            })

    def pct(t: dict) -> dict:
        done = t["pass"] + t["fail"]
        t["pass_rate"] = round(t["pass"] / done * 100) if done else 0
        t["progress"] = round((t["n"] - t["none"]) / t["n"] * 100) if t["n"] else 0
        return t

    return {
        "path": base,
        "level": level,
        "totals": pct(total),
        "children": sorted((pct(k) for k in kids.values()), key=lambda x: (x["pass_rate"], -x["n"])),
        "axis": axis,
        "groups": sorted((pct(a) for a in axes.values()), key=lambda x: (x["pass_rate"], -x["n"])),
        "cycles": rows,
        "trend": [
            {"at": v["k"], "pass": v["pass"], "fail": v["fail"],
             "pass_rate": round(v["pass"] / (v["pass"] + v["fail"]) * 100) if (v["pass"] + v["fail"]) else 0}
            for _k, v in sorted(trend.items())
        ],
    }


@app.get("/api/cycle/rollup")
async def cycle_rollup_get(path: str = _RU_ROOT, date_from: str = "", date_to: str = "",
                           axis: str = ""):
    """
    폴더 한 층의 현황 — 화면(KPI·막대·추이·표)이 이 하나를 쓴다.

    `axis` 를 주면 **하위 폴더 대신 그것으로 나눈** 막대를 함께 내려준다:
    cycle · version_group · model · customer · status · severity · assignee.
    (옛 Reports 의 「축 갈아끼우기」 가 이 자리로 왔다)
    """
    return await _rollup(path, date_from, date_to, axis)


@app.get("/api/cycle/rollup/items")
async def cycle_rollup_items(
    path: str = _RU_ROOT,
    date_from: str = "",
    date_to: str = "",
    q: str = "",
    kind: str = "",
    severity: str = "",
    cycle: str = "",
    verdict: str = "",
    limit: int = 20,
    offset: int = 0,
):
    """
    결과 상세 — 이 폴더에 걸린 **항목 한 줄씩**. 옛 Reports 의 아래 표다.
    거르개: 찾기 · 타입(auto·manual) · 심각도 · 사이클 · 판정 · 기간.
    """
    metas = await db.cycle_list_meta()
    cat = await db.catalog_list("model")
    fam = {str(m.get("name") or ""): str(m.get("family") or "").strip() for m in cat}
    mgrp = {str(m.get("name") or ""): str(m.get("model_group") or "").strip() for m in cat}
    grp_of = await _ru_groups()
    base = str(path or _RU_ROOT).strip().strip("/") or _RU_ROOT
    ql = q.strip().lower()

    out: list[dict] = []
    cycles_seen: list[dict] = []
    for c in metas:
        p = _ru_path(c, fam, mgrp)
        if p != base and not p.startswith(base + "/"):
            continue
        cid = str(c.get("cid") or c.get("id") or "")
        cnm = str(c.get("name") or cid)
        cycles_seen.append({"id": cid, "name": cnm})
        if cycle and cycle not in (cid, cnm):
            continue
        for it in (c.get("items") or []):
            if not isinstance(it, dict):
                continue
            day = _ru_day(it.get("executed_at"))
            if date_from and day and day < date_from:
                continue
            if date_to and day and day > date_to:
                continue
            v = str(it.get("_verdict") or it.get("result") or "")
            g = grp_of.get(v, "neutral" if v else "none")
            if verdict and verdict != g:
                continue
            if severity and str(it.get("severity") or "") != severity:
                continue
            if kind:
                steps = it.get("steps") or []
                man = sum(1 for x in steps if isinstance(x, dict) and x.get("manual"))
                aut = len(steps) - man
                k = "manual" if man and not aut else ("auto" if aut and not man else "mixed")
                if k != kind:
                    continue
            if ql:
                hay = f"{it.get('tcid') or ''} {it.get('name') or ''} {it.get('req_id') or ''}".lower()
                if ql not in hay:
                    continue
            out.append({
                "tcid": it.get("tcid"), "name": it.get("name"), "verdict": v, "group": g,
                "severity": it.get("severity"), "req_id": it.get("req_id"),
                "cycle": cnm, "cycle_id": cid,
                "executed_at": it.get("executed_at"),
                "fails": int(it.get("_steps_fail") or 0),
            })

    out.sort(key=lambda x: str(x.get("executed_at") or ""), reverse=True)
    lim = max(1, min(int(limit or 20), 500))
    off = max(0, int(offset or 0))
    seen: dict[str, str] = {}
    for c in cycles_seen:
        seen[c["id"]] = c["name"]
    return {"total": len(out), "rows": out[off:off + lim],
            "cycles": [{"id": k, "name": v} for k, v in seen.items()]}


# ───────────────────────────────────────────
# 보고서 — 같은 집계를 **한 장**으로. PDF 는 화면이 인쇄로 뽑고(같은 그림),
# 메일은 여기서 HTML 로 지어 보낸다. 두 갈래 다 같은 자료다.
# ───────────────────────────────────────────
def _rp_headline(r: dict) -> str:
    """메일 첫 문단 — 그대로 읽어도 말이 되게(회차 요약 헤드라인과 같은 결)"""
    t = r["totals"]
    nm = r["path"].split("/")[-1] or "Root"
    left = t["n"] - t["none"]
    s2 = (f"{nm} 시험 현황입니다. 항목 {t['n']}건 가운데 {left}건을 실행해 "
          f"진척 {t['progress']}%, 합격률 {t['pass_rate']}% (합격 {t['pass']} · 실패 {t['fail']}) 입니다.")
    if t["none"]:
        s2 += f" 아직 {t['none']}건이 남아 있습니다."
    if t["open_defects"]:
        s2 += f" 열린 결함은 {t['open_defects']}건입니다."
    if t["last_run"]:
        s2 += f" 마지막 실행은 {t['last_run']} 입니다."
    return s2


def _rp_html(r: dict, note: str = "") -> str:
    t = r["totals"]
    nm = r["path"].split("/")[-1] or "Root"
    kid_lb = {"root": "사업자", "operator": "제품군", "family": "모델그룹",
              "model_group": "모델명", "model": "버전그룹",
              "version_group": "회차"}.get(r["level"], "하위")

    def bar(x: dict) -> str:
        n = max(1, x["n"])
        seg = [("#16a34a", x["pass"]), ("#dc2626", x["fail"]),
               ("#f0b429", x["other"]), ("#c3cad4", x["none"])]
        cells = "".join(
            f'<td width="{round(v / n * 100)}%" bgcolor="{c}" style="height:8px;font-size:0;line-height:0">&nbsp;</td>'
            for c, v in seg if v
        )
        return f'<table width="150" cellpadding="0" cellspacing="0" style="border-radius:4px;overflow:hidden"><tr>{cells}</tr></table>'

    if r["level"] == "version_group":
        head = ["회차", "버전", "항목", "진행", "합격률"]
        body = "".join(
            f"<tr><td>{c.get('cid') or c.get('id')}</td><td>{c.get('version') or c.get('name') or '–'}</td>"
            f"<td align=right>{c['n']}</td><td>{bar(c)}</td>"
            f"<td align=right><b>{round(c['pass'] / (c['pass'] + c['fail']) * 100) if (c['pass'] + c['fail']) else 0}%</b></td></tr>"
            for c in r["cycles"]
        )
    else:
        head = [kid_lb, "회차", "항목", "진행", "합격률"]
        body = "".join(
            f"<tr><td>{k['key']}</td><td align=right>{k['cycles']}</td><td align=right>{k['n']}</td>"
            f"<td>{bar(k)}</td><td align=right><b>{k['pass_rate']}%</b></td></tr>"
            for k in r["children"]
        )

    kpi = "".join(
        f'<td style="padding:8px 12px;border:1px solid #e3e8ef;border-radius:8px">'
        f'<div style="font-size:11px;color:#9ca3af">{lb}</div>'
        f'<div style="font-size:19px;font-weight:700;color:{col}">{val}</div></td>'
        for lb, val, col in [
            ("합격률", f"{t['pass_rate']}%", "#16a34a" if t["pass_rate"] >= 80 else ("#dc2626" if t["pass_rate"] < 50 else "#1f2937")),
            ("진척률", f"{t['progress']}%", "#1f2937"),
            ("시험 항목", t["n"], "#1f2937"),
            ("열린 결함", t["open_defects"], "#dc2626" if t["open_defects"] else "#1f2937"),
            ("마지막 실행", t["last_run"] or "–", "#1f2937"),
        ]
    )
    note_html = f'<p style="margin:0 0 14px;padding:10px 12px;background:#fffbea;border:1px solid #fde68a;border-radius:8px">{note}</p>' if note else ""
    return f"""<div style="font-family:-apple-system,'Segoe UI',Roboto,'Noto Sans KR',sans-serif;color:#1f2937;font-size:13px;line-height:1.6;max-width:760px">
  <h2 style="margin:0 0 4px;font-size:18px">{nm} 시험 현황</h2>
  <div style="color:#6b7280;font-size:12px;margin-bottom:14px">{r['path']}</div>
  {note_html}
  <p style="margin:0 0 14px">{_rp_headline(r)}</p>
  <table cellspacing="6" cellpadding="0" style="margin:0 0 16px"><tr>{kpi}</tr></table>
  <table cellpadding="6" cellspacing="0" width="100%" style="border-collapse:collapse;font-size:12px">
    <tr style="background:#f5f7fa">{''.join(f'<th align=left style="border-bottom:1px solid #e3e8ef;color:#6b7280;font-size:11px">{h}</th>' for h in head)}</tr>
    {body}
  </table>
  <p style="color:#9ca3af;font-size:11px;margin-top:14px">
    합격률 = 합격 ÷ (합격+실패) · 진척률 = 실행 ÷ 전체 · 색: 합격 초록 · 실패 빨강 · 그 밖 노랑 · 미실행 회색<br>
    ubiQuoss-TOP 이 보낸 자동 요약입니다.
  </p>
</div>"""


@app.get("/api/cycle/rollup/csv")
async def cycle_rollup_csv(
    path: str = _RU_ROOT,
    date_from: str = "",
    date_to: str = "",
    q: str = "",
    kind: str = "",
    severity: str = "",
    cycle: str = "",
    verdict: str = "",
):
    """결과 상세를 원자료 그대로 — 지금 걸린 폴더·기간·거르개가 그대로 나간다."""
    import csv as _csv
    import io as _io
    from urllib.parse import quote

    d = await cycle_rollup_items(path, date_from, date_to, q, kind, severity, cycle,
                                 verdict, limit=100000, offset=0)
    buf = _io.StringIO()
    w = _csv.writer(buf)
    w.writerow(["결과", "TC ID", "시험항목", "부적합", "심각도", "요구사항", "사이클", "실행일"])
    for r in d["rows"]:
        w.writerow([
            {"pass": "합격", "fail": "불합격", "none": "미실행"}.get(r["group"], r["verdict"] or ""),
            r.get("tcid") or "", r.get("name") or "", r.get("fails") or 0,
            r.get("severity") or "", r.get("req_id") or "", r.get("cycle") or "",
            str(r.get("executed_at") or "")[:16].replace("T", " "),
        ])
    nm = (path.split("/")[-1] or "Root").replace(" ", "_")
    # 엑셀이 UTF-8 을 알아보게 BOM 을 붙인다 — 없으면 한글이 깨져 열린다
    body = ("\ufeff" + buf.getvalue()).encode("utf-8")
    return Response(
        content=body,
        media_type="text/csv; charset=utf-8",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(nm)}_result.csv"},
    )


@app.get("/api/cycle/rollup/preview")
async def cycle_rollup_preview(path: str = _RU_ROOT, date_from: str = "", date_to: str = "", note: str = ""):
    """메일로 나갈 그 모습 그대로 — 보내기 전에 눈으로 본다."""
    r = await _rollup(path, date_from, date_to)
    return {"subject": f"[UTOP] {r['path'].split('/')[-1]} 시험 현황",
            "headline": _rp_headline(r), "html": _rp_html(r, note)}


@app.post("/api/cycle/rollup/mail")
async def cycle_rollup_mail(payload: dict):
    """이 폴더의 현황을 메일로 보낸다 — 화면에서 보는 것과 같은 자료다."""
    path = str(payload.get("path") or _RU_ROOT)
    to = payload.get("to") or ""
    if not to:
        raise HTTPException(400, "받는 사람을 적어 주세요")
    r = await _rollup(path, str(payload.get("date_from") or ""), str(payload.get("date_to") or ""))
    subject = str(payload.get("subject") or "").strip() or f"[UTOP] {path.split('/')[-1]} 시험 현황"
    html = _rp_html(r, str(payload.get("note") or "").strip())
    try:
        sent = _send_mail(to, subject, html, html=True)
    except Exception as e:
        raise HTTPException(400, f"보내지 못했습니다 — {e}")
    return {"success": True, "to": sent, "subject": subject}


# 버전그룹 폴더 — `{ "<모델명>": ["R200", "R300"] }`
#
# 모델그룹·모델명은 장비 카탈로그가 master 다. 자유 입력으로 두었더니
# `E4320-24P_2` 같은 것이 생겼다. 버전그룹만 사람이 만든다 — R200, R300
# 은 카탈로그가 알 수 없는, 이 회차 묶음의 이름이라서다.
#
# **파일이 아니라 DB 에 둔다.** 옛 사이클 폴더는
# `data/state/cycle_folders.json` 이었고, 자료를 옮길 때 딸려오지 않아
# 사이클 23건이 전부 이름 없는 폴더를 가리키게 됐다.
_VGROUP_KV = "cycle_version_groups"


@app.get("/api/cycle-version-groups")
async def get_cycle_version_groups():
    return {"groups": await db.kv_get(_VGROUP_KV) or {}}


# ── 자연어로 시험 짜기 ──────────────────────────────────────
#
# 「E5724RL 시스템 정보 시험해줘」 한 줄로 돌아가게 하는 자리.
#
# **LLM 이 시험을 지어내게 하지 않는다.** 있는 TC 중에서 고르게만 한다.
# 스텝을 자유롭게 만들게 하면 그럴듯한데 틀린 시험이 나오고, 그건 사람이
# 검토하는 데 더 오래 걸린다. 고르게 하면 결과가 「이 3건」 이라 눈으로
# 바로 확인된다.
#
# 고른 뒤에 무엇을 할지는 화면이 정한다 — 여기서는 계획만 돌려준다.

_NL_SCHEMA = {
    "type": "object",
    "properties": {
        "model": {"type": "string"},
        "tcids": {"type": "array", "items": {"type": "string"}},
        "device_ip": {"type": "string"},
        "why": {"type": "string"},
    },
    "required": ["tcids", "why"],
}


@app.post("/api/nl/plan")
async def nl_plan(payload: dict):
    """말 한 줄 → 돌릴 시험 목록. 고르기만 하고 실행하지는 않는다."""
    text = str(payload.get("text") or "").strip()
    if not text:
        raise HTTPException(400, "무엇을 시험할지 적어 주세요")

    tcs = await db.tc_list_meta()
    models = [c["name"] for c in await db.catalog_list("model")]
    devices = await db.device_list()
    reqs = await db.req_list_full()

    # 근거를 넓게 준다.
    #
    # TC 이름만 보고 고르면 「시스템 정보」 같은 말에는 맞지만 「Gi0/1 링크
    # 시험」 처럼 장비·포트를 가리키는 말에는 못 맞춘다. 어떤 장비가 있고
    # 무슨 포트가 달렸는지, 그 시험이 어느 요구사항 아래인지까지 함께
    # 준다 — 위에서부터 쌓인 것이 다 근거다.
    req_by_id = {}
    for r in reqs:
        rid = str(r.get("reqid") or r.get("id") or "")
        if rid:
            req_by_id[rid] = str(r.get("title") or "")

    lines = []
    for t in tcs:
        nm = str(t.get("name") or "").strip()
        if not nm:
            continue
        rt = req_by_id.get(str(t.get("req_id") or ""), "")
        lines.append(f"{t.get('tcid')}\t{nm}\t{rt}")

    dev_lines = []
    for d in devices:
        if d.get("role") == "계측기":
            continue
        ifs = [str(i.get("name")) for i in (d.get("interfaces") or [])][:12]
        dev_lines.append(
            f"{d.get('ip')}\t{d.get('model') or ''}\t{d.get('role') or ''}"
            f"\t포트: {', '.join(ifs) or '(등록 안 됨)'}"
        )

    sys_p = (
        "너는 네트워크 장비 시험 담당자를 돕는다. 사람이 한 말에 맞는 시험을 "
        "**아래 목록에서 고르기만** 한다. 목록에 없는 tcid 는 절대 만들지 않는다. "
        "맞는 것이 없으면 tcids 를 빈 배열로 두고 why 에 그렇게 적는다. "
        "말에 장비나 모델이 나오면 model 과 device_ip 에 **등록된 것 중에서** 골라 적는다. "
        "why 는 왜 이것들을 골랐는지 한국어 한두 문장."
    )
    user_p = (
        f"사람이 한 말: {text}\n\n"
        f"등록된 모델: {', '.join(models) or '(없음)'}\n\n"
        "등록된 장비 (IP<TAB>모델<TAB>역할<TAB>포트):\n" + "\n".join(dev_lines) + "\n\n"
        "시험 목록 (tcid<TAB>이름<TAB>요구사항):\n" + "\n".join(lines)
    )

    ans, err = await _ai_chat(
        [{"role": "system", "content": sys_p}, {"role": "user", "content": user_p}],
        max_tokens=800,
        json_schema=_NL_SCHEMA,
    )
    if err:
        raise HTTPException(502, err)

    plan, perr = _json_from_llm(ans)
    if plan is None:
        raise HTTPException(502, f"AI 응답을 읽지 못했습니다 — {perr}")
    if isinstance(plan, list):
        plan = {"tcids": [x for x in plan if isinstance(x, str)], "why": ""}
    if not isinstance(plan, dict):
        raise HTTPException(502, f"AI 응답이 예상 밖입니다 — {str(ans)[:200]}")

    # 지어낸 tcid 를 걸러낸다. 없는 것을 돌리려다 실패하면 왜인지 알기 어렵다
    known = {str(t.get("tcid")) for t in tcs}
    picked = [x for x in (plan.get("tcids") or []) if str(x) in known]
    dropped = [x for x in (plan.get("tcids") or []) if str(x) not in known]
    by_id = {str(t.get("tcid")): t for t in tcs}

    # 장비도 실제로 있는 것만 남긴다
    dev_ips = {str(d.get("ip")) for d in devices}
    dev_ip = str(plan.get("device_ip") or "")
    return {
        "model": plan.get("model") or "",
        "device_ip": dev_ip if dev_ip in dev_ips else "",
        "why": plan.get("why") or "",
        "tcs": [
            {"tcid": x, "name": by_id[str(x)].get("name") or "", "req_id": by_id[str(x)].get("req_id") or ""}
            for x in picked
        ],
        # 지어낸 것이 있었다는 사실도 알려 준다 — 조용히 지우면 왜 빠졌는지 모른다
        "dropped": dropped,
    }


def _json_from_llm(text):
    """LLM 이 돌려준 글에서 JSON 을 꺼낸다.

    `guided_json` 을 줘도 모델이 ```json 울타리를 씌워 보내는 일이 흔하다
    (gemma 가 그렇다). 앞뒤에 설명을 한 줄 붙이기도 한다. 그대로
    `json.loads` 하면 「AI 응답을 읽지 못했습니다」 만 뜨고, 무엇이 왔는지
    알 수가 없다.

    울타리를 벗기고, 그래도 안 되면 첫 `{` 부터 짝이 맞는 `}` 까지를 잘라
    본다. 끝내 못 읽으면 받은 글을 함께 돌려줘 화면이 보여 줄 수 있게 한다.
    """
    s = str(text or "").strip()
    if not s:
        return None, "빈 응답"
    # ```json … ``` 벗기기
    if s.startswith("```"):
        s = s.split("\n", 1)[-1]
        if s.rstrip().endswith("```"):
            s = s.rstrip()[:-3]
        s = s.strip()
    try:
        return json.loads(s), None
    except Exception:
        pass
    # 첫 { 부터 짝이 맞는 } 까지
    start = s.find("{")
    if start >= 0:
        depth, in_str, esc = 0, False, False
        for i in range(start, len(s)):
            ch = s[i]
            if in_str:
                if esc:
                    esc = False
                elif ch == "\\":
                    esc = True
                elif ch == '"':
                    in_str = False
                continue
            if ch == '"':
                in_str = True
            elif ch == "{":
                depth += 1
            elif ch == "}":
                depth -= 1
                if depth == 0:
                    try:
                        return json.loads(s[start:i + 1]), None
                    except Exception:
                        break
    return None, "받은 글: " + str(text or "")[:300]


# ── 자연어로 **새 시험 만들기** ──────────────────────────────
#
# 있는 시험을 찾아 주는 것이 아니라, 있는 것을 **참고해서** 새 시험을
# 짜고 돌리고 결과를 알려 주는 것이 목적이다.
#
# 다만 1차는 **조회 명령만** 짓게 한다. 설정을 바꾸는 명령을 AI 가 지어내
# 장비로 보내면 되돌릴 수가 없다. 조회는 틀려도 「출력이 없다」 로 끝난다.

# 이 랩에서 실제로 통한 조회 명령의 머리말. 여기 없는 것은 안 내보낸다.
_READ_HEADS = (
    "show", "display", "get", "dir", "more", "cat", "ping", "traceroute",
    "do show", "showtech", "who", "history",
)
# 한 글자라도 걸리면 자른다 — 조회처럼 보여도 뒤에 붙는 경우가 있다
_WRITE_WORDS = (
    "configure", "conf t", "config t", "write", "wr ", "reload", "erase",
    "delete", "format", "copy ", "clear ", "no ", "set ", "shutdown",
    "reset", "reboot", "boot ", "upgrade", "install", "restore", "factory",
)


# 설정 시험에서만 푸는 명령. 「무엇을 풀지」 를 여기 한 곳에 적어 둔다 —
# 여러 군데 흩어 두면 한쪽만 고쳐 놓고 풀린 줄 안다.
_CONFIG_HEADS = (
    "configure terminal", "conf t", "config t", "end", "exit",
    "interface", "int ", "no shutdown", "shutdown",
)
# 무엇을 풀어 주든 이것만은 못 지나간다. 되돌릴 수 없거나 장비가 죽는다.
_NEVER_WORDS = (
    "reload", "reboot", "erase", "format", "factory", "upgrade", "firmware",
    "install", "restore", "write ", "wr ", "copy ", "delete", "rmdir",
    "clear config", "halt", "boot ",
)


def _config_allowed(cli: str) -> bool:
    """설정 시험에서 이 명령을 써도 되나 (allow_config 일 때만 부른다).

    허용 목록 방식이다 — 「막을 것을 적는」 방식은 새 명령이 생길 때마다
    구멍이 난다. 링크를 내리고 올리는 시험에 필요한 것만 연다.
    """
    s0 = str(cli or "").strip().lower()
    if not s0:
        return False
    for line in s0.splitlines():
        ln = line.strip()
        if not ln:
            continue
        if any(w in ln for w in _NEVER_WORDS):
            return False
        if any(ch.isdigit() or ch == "." for ch in ln) and all(
            ch.isdigit() or ch == "." for ch in ln
        ):
            continue                      # OID
        if any(ln.startswith(h) for h in _READ_HEADS):
            continue                      # 조회는 언제나 된다
        if any(ln.startswith(h) for h in _CONFIG_HEADS):
            continue
        return False
    return True


def _is_read_only(cli: str) -> bool:
    """조회 명령인가. SNMP OID(숫자와 점)도 조회로 본다."""
    s = str(cli or "").strip().lower()
    if not s:
        return False
    for line in s.splitlines():
        ln = line.strip()
        if not ln:
            continue
        if all(ch.isdigit() or ch == "." for ch in ln):   # OID
            continue
        if any(w in ln for w in _WRITE_WORDS):
            return False
        if not any(ln.startswith(h) for h in _READ_HEADS):
            return False
    return True


_TC_SCHEMA = {
    "type": "object",
    "properties": {
        "name": {"type": "string"},
        "object": {"type": "string"},
        "device_ip": {"type": "string"},
        "device_ips": {"type": "array", "items": {"type": "string"}},
        "steps": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "desc": {"type": "string"},
                    "kind": {"type": "string"},
                    "cli": {"type": "string"},
                    "type": {"type": "string"},
                    "criteria": {"type": "string"},
                    "session": {"type": "integer"},
                    "loopCount": {"type": "integer"},
                    "waitSec": {"type": "number"},
                },
                "required": ["desc"],
            },
        },
    },
    "required": ["name", "steps"],
}


@app.post("/api/nl/tc")
async def nl_make_tc(payload: dict):
    """말 한 줄 → **새 시험 초안**. 만들기만 하고 저장·실행은 화면이 한다."""
    text = str(payload.get("text") or "").strip()
    if not text:
        raise HTTPException(400, "무엇을 시험할지 적어 주세요")
    # 설정 시험(링크 down/up 같은 것)은 사람이 켤 때만 만든다. 기본은 조회다.
    allow_config = bool(payload.get("allow_config"))

    devices = [d for d in await db.device_list() if d.get("role") != "계측기"]
    tcs = await db.tc_list_full()

    # 이 랩에서 실제로 통한 명령을 모은다.
    #
    # 이것이 이 기능의 근거다. 일반적인 네트워크 지식으로 명령을 지으면
    # 이 장비에서 안 통하는 것이 나온다. 여기서 실제로 오간 것을 주면
    # 「이 장비가 알아듣는 말」 안에서 고르게 된다.
    seen, cmds = set(), []
    for tc in tcs:
        for c in (tc.get("checks") or []):
            cli = str((c or {}).get("cli") or "").strip()
            if not cli or not _is_read_only(cli):
                continue
            head = cli.splitlines()[0].strip()
            if head in seen:
                continue
            seen.add(head)
            cmds.append(head)
    cmds = cmds[:120]

    # 비슷한 시험의 판정기준을 예로 준다 — 무엇을 어떻게 보는지의 본
    samples = []
    for tc in tcs[:40]:
        for c in (tc.get("checks") or [])[:3]:
            cli = str((c or {}).get("cli") or "").strip()
            cr = str((c or {}).get("criteria") or "").strip()
            if cli and cr and _is_read_only(cli):
                samples.append(f"{cli.splitlines()[0]} → [{c.get('type') or 'contains'}] {cr}")
        if len(samples) >= 25:
            break

    dev_lines = [
        f"{d.get('ip')}\t{d.get('model') or ''}\t{d.get('role') or ''}"
        for d in devices
    ]

    sys_p = (
        "너는 네트워크 장비 시험 절차를 짠다. 사람이 말한 것을 확인할 수 있는 "
        "**조회 시험**을 만든다.\n"
        "규칙:\n"
        "1. 명령은 아래 「이 랩에서 통한 명령」 에 있는 것을 그대로 쓰거나 그 꼴을 따른다. "
        "일반적인 지식으로 새 명령을 지어내지 않는다.\n"
        + ("2. **설정 시험이다.** 아래만 쓸 수 있다 — configure terminal · interface <이름> · "
           "shutdown · no shutdown · end · exit, 그리고 조회 명령(show…). "
           "reload·write·copy·erase·factory 는 절대 쓰지 않는다.\n"
           "2-1. 링크를 내렸으면 **반드시 다시 올린다**(no shutdown). 내려 둔 채 끝내면 "
           "다음 시험이 전부 깨진다.\n"
           "2-2. 상태가 반영되기까지 시간이 걸린다 — 내리고/올린 **뒤에 wait 스텝**을 둔다.\n"
           if allow_config else
           "2. **조회 명령만.** configure·write·reload·no·set·clear 같은 것은 절대 쓰지 않는다.\n")
        + "3. **스텝마다 판정기준(criteria)을 반드시 적는다.** 비워 두면 그 스텝은 "
        "돌기만 하고 아무것도 확인하지 못한다. 출력에 늘 나오는 **항목 이름**을 "
        "기준으로 삼으면 안전하다 — 값은 장비마다 다르지만 이름은 같다. "
        "예: `show cpu usage` → `Average CPU load`, `show memory usage` → `Total`. "
        "확인할 문구를 딱 집기 애매하면 **type 을 ok** 로 둔다 — 「명령이 오류 없이 "
        "응답하면 합격」 이라는 뜻이고, 조회 시험은 대개 이것이면 된다. criteria 는 "
        "비워 둔다. type 은 contains(문구 포함) · contains_all(콤마로 여럿, 모두 포함) · "
        "ok(오류만 없으면) 중 하나. none 은 쓰지 않는다.\n"
        "4. 스텝은 2~10개. 많을수록 좋은 것이 아니다.\n"
        "5. desc 는 그 스텝이 무엇을 확인하는지 한국어 한 줄.\n"
        "6. 스텝 종류(kind): cli(명령·기본) · wait(기다리기, waitSec 초) · "
        "loop(여기부터 아래를 loopCount 번 되풀이). loop 는 되풀이할 묶음 **앞**에 한 번 둔다.\n"
        "7. 장비가 둘이면 session 으로 가른다 — 0 이 첫 장비, 1 이 둘째. "
        "device_ips 에 쓴 차례와 같다. 장비가 하나면 session 은 적지 않는다.\n\n"
        "아래 꼴 그대로, **다른 말 없이 JSON 만** 답한다:\n"
        '{"name":"E5724RL 시스템 정보 확인","object":"모델명과 메모리를 확인한다",'
        '"device_ip":"210.1.1.254","steps":['
        '{"desc":"모델명을 확인한다","cli":"show system","type":"contains","criteria":"E5724RL"},'
        '{"desc":"CPU 사용량이 조회되는지 확인한다","cli":"show cpu usage","type":"ok","criteria":""}'
        ']}'
        + ("\n설정 시험 예 (링크를 내렸다 올리며 상대 장비에서 확인):\n"
           '{"name":"gi0/1 링크 down/up 100회","object":"링크를 내렸다 올리며 상대에서 상태를 본다",'
           '"device_ips":["210.1.1.254","210.1.1.253"],"steps":['
           '{"desc":"100회 되풀이","kind":"loop","loopCount":100},'
           '{"desc":"A 장비 gi0/1 을 내린다","kind":"cli","session":0,'
           '"cli":"configure terminal\\ninterface gi0/1\\nshutdown","type":"ok","criteria":""},'
           '{"desc":"상태가 반영되기를 기다린다","kind":"wait","waitSec":2},'
           '{"desc":"B 장비에서 링크가 내려갔는지 본다","kind":"cli","session":1,'
           '"cli":"show interface gi0/2","type":"contains","criteria":"down"},'
           '{"desc":"A 장비 gi0/1 을 올린다","kind":"cli","session":0,'
           '"cli":"configure terminal\\ninterface gi0/1\\nno shutdown","type":"ok","criteria":""},'
           '{"desc":"상태가 반영되기를 기다린다","kind":"wait","waitSec":2},'
           '{"desc":"B 장비에서 링크가 올라왔는지 본다","kind":"cli","session":1,'
           '"cli":"show interface gi0/2","type":"contains","criteria":"up"}'
           ']}'
           if allow_config else "")
    )
    user_p = (
        f"사람이 한 말: {text}\n\n"
        "등록된 장비 (IP<TAB>모델<TAB>역할):\n" + "\n".join(dev_lines) + "\n\n"
        "이 랩에서 통한 명령:\n" + "\n".join(cmds) + "\n\n"
        "판정기준 예:\n" + "\n".join(samples)
    )

    ans, err = await _ai_chat(
        [{"role": "system", "content": sys_p}, {"role": "user", "content": user_p}],
        max_tokens=1400,
        json_schema=_TC_SCHEMA,
    )
    if err:
        raise HTTPException(502, err)
    draft, perr = _json_from_llm(ans)
    if draft is None:
        raise HTTPException(502, f"AI 응답을 읽지 못했습니다 — {perr}")

    # 최상위를 배열로 보내는 일이 있다. 스키마를 줘도 그렇다.
    #   · [{...steps...}]        → 스텝 목록 그 자체
    #   · [{"name":…, "steps":…}] → 감싼 것이 하나뿐
    # 여기서 받아 주지 않으면 500 이 나고, 화면에는 이유가 안 보인다.
    if isinstance(draft, list):
        if len(draft) == 1 and isinstance(draft[0], dict) and "steps" in draft[0]:
            draft = draft[0]
        else:
            draft = {"name": text[:40], "steps": [x for x in draft if isinstance(x, dict)]}
    if not isinstance(draft, dict):
        raise HTTPException(502, f"AI 응답이 예상 밖입니다 — {str(ans)[:200]}")

    # 조회가 아닌 명령은 잘라낸다. 조용히 지우지 않고 무엇을 왜 뺐는지 알린다
    # 모델이 스키마를 무시하고 제 나름의 이름을 쓴다.
    #
    #   {"cmd": "show version", "criteria": "contains", "value": "E4300"}
    #
    # `cmd` 가 명령이고 `criteria` 자리에 **판정 종류**가, `value` 에 기준이
    # 들어 있다. 이름 하나 다르다고 빈 화면을 보여 줄 이유가 없다.
    _TYPES = {"contains", "contains_all", "notcontains", "line", "ok", "none", "expr", "table"}

    def _pick(s, *names):
        for n in names:
            v = s.get(n)
            if v not in (None, ""):
                return str(v).strip()
        return ""

    def _num(v, dflt=0):
        try:
            return type(dflt)(v)
        except Exception:
            return dflt

    keep, cut = [], []
    for s in (draft.get("steps") or []):
        if not isinstance(s, dict):
            continue

        # 명령이 아닌 스텝 — 되풀이(loop)·기다리기(wait). cli 가 없어도 산다.
        skind = _pick(s, "kind", "step_kind").lower()
        if skind in ("loop", "wait"):
            row = {"desc": _pick(s, "desc", "description", "purpose"), "kind": skind}
            if skind == "loop":
                row["loopCount"] = max(1, _num(s.get("loopCount") or s.get("count"), 1))
            else:
                row["waitSec"] = max(0.1, _num(s.get("waitSec") or s.get("sec"), 1.0))
            keep.append(row)
            continue

        cli = _pick(s, "cli", "cmd", "command", "input")
        if not cli:
            continue
        # 설정 시험이면 허용 목록까지, 아니면 조회만. 어느 쪽이든 되돌릴 수
        # 없는 명령은 못 지나간다.
        ok_cmd = _config_allowed(cli) if allow_config else _is_read_only(cli)
        if not ok_cmd:
            cut.append(cli.splitlines()[0])
            continue
        kind = _pick(s, "type", "judge", "mode")
        crit = _pick(s, "criteria", "value", "expected", "expect")
        # `criteria` 자리에 종류가 들어온 경우 — 서로 바꿔 놓는다
        if not kind and crit in _TYPES:
            kind, crit = crit, _pick(s, "value", "expected", "expect")
        if kind not in _TYPES:
            kind = "contains" if crit else "ok"
        # 기준이 비었는데 문구를 보라고 온 것은 판정을 못 한다.
        # 「오류만 없으면 합격」 으로 돌린다 — 조회 시험의 기본값이다.
        if kind in ("contains", "contains_all", "notcontains", "line") and not crit:
            kind = "ok"
        row = {
            "desc": _pick(s, "desc", "description", "purpose"),
            "kind": "cli",
            "cli": cli,
            "type": kind,
            "criteria": crit,
        }
        # 장비가 둘 이상일 때만 세션을 싣는다 — 하나뿐이면 0 이 당연해서 군더더기다
        if s.get("session") is not None:
            row["session"] = max(0, _num(s.get("session"), 0))
        keep.append(row)

    dev_ips = {str(d.get("ip")) for d in devices}
    ip = str(draft.get("device_ip") or "")
    ips = [str(x) for x in (draft.get("device_ips") or []) if str(x) in dev_ips]
    if not ips and ip in dev_ips:
        ips = [ip]
    return {
        "name": str(draft.get("name") or draft.get("title") or "").strip() or text[:40],
        "object": str(draft.get("object") or "").strip(),
        "device_ip": ips[0] if ips else "",
        "device_ips": ips,
        "steps": keep,
        "cut": cut,
        "allow_config": allow_config,
    }


@app.post("/api/cycle-version-groups")
async def save_cycle_version_groups(payload: dict):
    groups = payload.get("groups")
    if not isinstance(groups, dict):
        raise HTTPException(400, "groups 는 { 모델명: [버전그룹…] } 이어야 합니다")
    clean = {}
    for model, arr in groups.items():
        m = str(model).strip()
        if not m or not isinstance(arr, list):
            continue
        seen = []
        for g in arr:
            g = str(g).strip()
            if g and g not in seen:
                seen.append(g)
        clean[m] = seen
    await db.kv_set(_VGROUP_KV, clean)
    return {"ok": True, "groups": clean}


CYCLE_FOLDERS_FILE = DATA_DIR / "state" / "cycle_folders.json"

@app.get("/api/cycle-folders")
async def get_cycle_folders():
    if not CYCLE_FOLDERS_FILE.exists():
        return {"folders": []}
    return load_json(CYCLE_FOLDERS_FILE)

@app.post("/api/cycle-folders")
async def save_cycle_folders(data: dict):
    save_json(CYCLE_FOLDERS_FILE, data)
    return {"success": True}

RACKS_FILE = DATA_DIR / "state" / "racks.json"

def _ensure_blank_ids(d):
    """옛 자료로 들어온 부품에 id 가 없으면 채워서 저장해 둔다.
    부품 바꾸기·빼기·옮기기가 id 로 찾기 때문에, 없으면 저장이 헛돌거나
    (자기 자신 제외가 안 걸려) 겹침 검사에 걸려 저장 버튼이 죽는다."""
    try:
        if not isinstance(d, dict):
            return d
        dirty = False
        for i, b in enumerate(d.get("blanks") or []):
            if isinstance(b, dict) and not b.get("id"):
                b["id"] = f"blk-fix-{i}-{b.get('pos', 0)}"
                dirty = True
        if dirty:
            _kv_save_sync("racks", d)
    except Exception:
        pass
    return d

@app.get("/api/racks")
async def get_racks():
    d = _ensure_blank_ids(_kv_load_sync("racks", {"racks": []}))
    return d if isinstance(d, dict) else {"racks": []}

@app.post("/api/racks")
async def save_racks(data: dict):
    # 안전장치: 기존 데이터가 있는데 요청이 완전히 비어있으면 거부 (빈 저장으로 랙 배치 유실 방지)
    prev = _kv_load_sync("racks", {})
    if isinstance(prev, dict) and (prev.get("labs") or prev.get("racks")) and not (data.get("labs") or data.get("racks")):
        return {"success": False, "error": "빈 데이터로 기존 랙 배치를 덮어쓸 수 없습니다"}
    _kv_save_sync("racks", data)
    return {"success": True}


# ══════════════════════════════════════════════════════════════════════
# 데이터 이사 — 묶음 단위 내보내기/가져오기 (설정 → 데이터)
#
# 랩마다 UTOP 이 따로 서 있어 자료를 통째로 옮기는 일이 잦다. DB 를 그대로
# 복사하면 장비 비밀번호·LLM 키까지 따라가므로, 묶음을 골라 JSON 하나로
# 뜨고, 받는 쪽은 ID 기준 합치기(upsert)로 넣는다.
# ══════════════════════════════════════════════════════════════════════

_TRANSFER_PARTS = ("req", "tc", "cycle", "defect", "device", "catalog", "settings")


def _strip_derived(d: dict) -> dict:
    """내보낼 때 붙인 파생 키(_created_at 등)를 걷는다 — 원본에 없던 것이다."""
    return {k: v for k, v in d.items() if not str(k).startswith("_")}


@app.get("/api/tc/{tc_id}/revisions")
async def tc_revisions_api(tc_id: str):
    """이 시험의 지난 판들 — 최신이 앞."""
    return {"items": await db.tc_revisions(tc_id)}


@app.post("/api/tc/{tc_id}/revisions/{rev_id}/restore")
async def tc_revision_restore(tc_id: str, rev_id: int, request: Request):
    """그 판으로 되돌린다. 지금 판은 되돌리기 직전에 자동으로 이력에 남는다."""
    data = await db.tc_revision_get(tc_id, rev_id)
    if data is None:
        raise HTTPException(404, "그 판이 없습니다")
    _by = ""
    try:
        _by = _user_of(_token_from(request)) or ""
    except Exception:
        pass
    if isinstance(data, dict):
        data = dict(data)
        data["updated_by"] = _by
    await db.tc_upsert(tc_id, data)
    try: asyncio.create_task(broadcast({"type": "tc_updated", "tcid": tc_id, "user": _by}))
    except Exception: pass
    return {"ok": True}


@app.get("/api/tc-running")
async def tc_running_now():
    """지금 자동 실행 중인 시험들 — 방금 접속한 사람이 현황을 받는다.
    브로드캐스트는 이미 붙어 있는 사람에게만 가므로, 이 GET 이 초기값이다."""
    return {"items": {k: v for k, v in _tc_running.items()}}


@app.get("/api/presence")
async def presence_roster(prefix: str = ""):
    """지금 접속해 있는 사람들 — prefix 로 화면을 좁힌다 (cycle → cycle:*)."""
    seen = []
    for st in ws_state.values():
        u = st.get("user")
        pg = str(st.get("page") or "")
        if u and (not prefix or pg == prefix or pg.startswith(prefix + ":")) and u not in seen:
            seen.append(u)
    return {"users": seen}


# ══════════════════════════════════════════════════════════════════
# 위키 — 프로젝트마다 갖는 문서
#
# 본문(body)은 편집기가 읽고 쓰는 **블록**이 정본이고, 찾기용 민글(plain)을
# 함께 담는다. 블록을 뒤져 찾을 수는 없다.
# ══════════════════════════════════════════════════════════════════
def _wiki_plain(body) -> str:
    """블록에서 글자만 훑어 낸다 — 찾기가 읽을 것.

    블록 꼴은 편집기가 정한다. 우리가 아는 것은 「어딘가에 text 가 있다」
    뿐이라, 모양을 따지지 않고 재귀로 긁는다. 모양이 바뀌어도 안 깨진다.
    """
    out = []

    def walk(v):
        if isinstance(v, dict):
            t = v.get("text")
            if isinstance(t, str):
                out.append(t)
            # 「살아 있는 표」 는 글자가 없다 — 담긴 것은 질의뿐이라, 찾기가
            # 훑을 것이 하나도 없어 문서에서 통째로 사라진다. 무엇을 가리키는
            # 블록인지만 남긴다: 사이클 ID 로 문서를 찾는 일이 실제로 있다.
            if v.get("type") == "utopView":
                p = v.get("props") or {}
                out.append(" ".join(
                    str(x) for x in ("UTOP 표", p.get("view"), p.get("cycle"), p.get("project")) if x
                ))
            for x in v.values():
                walk(x)
        elif isinstance(v, list):
            for x in v:
                walk(x)

    walk(body)
    return " ".join(out)[:200000]


@app.get("/api/wiki")
async def wiki_list(project: str = ""):
    """문서 트리 — 본문은 안 준다. 목록에 본문까지 실으면 수백 KB 가 된다."""
    async with db.pool().acquire() as c:
        rows = await c.fetch(
            # 프로젝트를 골랐어도 **프로젝트 없는 문서는 늘 보인다.**
            #
            # 「전체 프로젝트」 로 두고 쓴 문서는 project 가 빈 값으로 저장된다.
            # 그런데 나중에 프로젝트를 하나 고르면 그 문서들이 목록에서 통째로
            # 사라져, 쓴 사람은 **지워진 줄 안다**(지적: 문서가 다 날아갔다).
            # 빈 값은 「이 프로젝트 것이 아니다」 가 아니라 「어느 프로젝트에도
            # 매이지 않았다」 — 공용 문서다. 공용은 어디서 보든 보여야 한다.
            "SELECT id, project, parent_id, title, ord, updated_by, updated_at "
            "FROM wiki_page WHERE ($1='' OR project=$1 OR coalesce(project,'')='') "
            "ORDER BY ord, title",
            project,
        )
    return {
        "pages": [
            {**dict(r), "updated_at": r["updated_at"].isoformat() if r["updated_at"] else None}
            for r in rows
        ]
    }


@app.get("/api/wiki/search")
async def wiki_search(q: str, project: str = "", limit: int = 40):
    """**본문까지** 찾는다 — 이름만으로는 「그 말이 어느 문서에 있더라」 를 못 찾는다.

    민글(plain)을 그대로 훑는다. 전문검색 색인을 쓰지 않는 것은 문서가 수천
    장이 아니기 때문이다 — 지금 크기에서 ILIKE 로 충분하고, 색인은 한국어
    형태소를 걸어야 제구실을 해서 값이 크다.
    """
    n = (q or "").strip()
    if not n:
        return {"hits": []}
    async with db.pool().acquire() as c:
        rows = await c.fetch(
            "SELECT id, title, plain FROM wiki_page "
            # 목록과 **같은 규칙** — 프로젝트 없는 문서는 늘 걸린다.
            # 목록에는 보이는데 찾기에는 안 걸리면 그건 더 헷갈린다.
            "WHERE ($2='' OR project=$2 OR coalesce(project,'')='') "
            "AND (title ILIKE $1 OR plain ILIKE $1) "
            "ORDER BY updated_at DESC LIMIT $3",
            f"%{n}%", project, max(1, min(200, limit)),
        )
    out = []
    for r in rows:
        p = r["plain"] or ""
        i = p.lower().find(n.lower())
        # 걸린 자리 앞뒤를 잘라 보여 준다 — 「어디에 있나」 를 열지 않고 알게
        snip = p[max(0, i - 40) : i + 80] if i >= 0 else ""
        out.append({"id": r["id"], "title": r["title"], "snippet": snip})
    return {"hits": out}


@app.get("/api/wiki/{pid}")
async def wiki_get(pid: str):
    async with db.pool().acquire() as c:
        r = await c.fetchrow("SELECT * FROM wiki_page WHERE id=$1", pid)
    if not r:
        raise HTTPException(404, "문서를 찾을 수 없습니다")
    d = dict(r)
    for k in ("created_at", "updated_at"):
        if d.get(k):
            d[k] = d[k].isoformat()
    if isinstance(d.get("body"), str):
        try:
            d["body"] = json.loads(d["body"])
        except Exception:
            d["body"] = []
    return {"page": d}


# ── 고정 주소는 {pid} 보다 **먼저** 등록한다 ─────────────────────────
#
# FastAPI 는 먼저 등록된 길부터 맞춰 본다. `/api/wiki/{pid}` 가 위에 있으면
# `/api/wiki/pdf` 요청이 pid="pdf" 로 걸려 **「pdf 라는 이름의 문서를 저장」**
# 이 된다. 200 이 돌아오니 화면은 성공으로 보이는데 정작 PDF 는 없다 —
# 그리고 wiki_page 에 쓰레기 문서가 하나 생긴다. 워드 가져오기도 같은 일을
# 겪었다(지적: PDF·워드 둘 다 안 된다).
@app.post("/api/wiki/pdf")
async def wiki_pdf(payload: dict):
    """문서를 **PDF 파일로 구워서** 돌려준다.

    여태는 브라우저 인쇄 창을 띄웠다 — 미리보기가 뜨고, 대상을 고르고, 저장을
    눌러야 했다. 게다가 종이가 화면과 자꾸 갈렸다: 인쇄 창이 앱 CSS 를 못
    불러오거나 옛 판을 들고 갔다.

    화면이 보내 준 **그 HTML 그대로** 크로미움으로 찍는다. 화면을 그리는
    엔진과 종이를 찍는 엔진이 하나라, 갈릴 자리가 없다.
    """
    html = str(payload.get("html") or "")
    title = str(payload.get("title") or "문서")
    # 들어온 것부터 남긴다 — 요청이 여기까지 왔는지가 첫 갈림길이다.
    print(f"[pdf] 요청 도착 — {title} / html {len(html)}글자", flush=True)
    if not html.strip():
        return {"ok": False, "error": "찍을 내용이 없습니다"}
    try:
        from playwright.async_api import async_playwright
    except Exception as e:
        return {"ok": False, "error": "PDF 엔진이 없습니다: " + str(e)[:120]}

    import base64 as _b64
    try:
        async with async_playwright() as pw:
            br = await pw.chromium.launch(args=["--no-sandbox", "--disable-dev-shm-usage"])
            pg = await br.new_page()
            # 글꼴·그림이 다 앉은 뒤에 찍는다. 바로 찍으면 글자가 자리를 잡기
            # 전이라 줄이 어긋나고 그림 자리가 빈다.
            await pg.set_content(html, wait_until="networkidle")
            await pg.emulate_media(media="print")
            # 배율 90%(지시).
            #
            # 100% 로 찍으면 CSS 에 적은 크기가 그대로 나간다 — 틀리진 않지만
            # 종이에서는 빡빡하다. 90% 로 한 번 줄이면 한 장에 더 담기고,
            # **글자와 표와 그림이 같은 비율로** 줄어 균형이 안 깨진다.
            # 글자 크기만 따로 줄이면 표·그림만 커 보인다.
            pdf = await pg.pdf(
                format="A4",
                margin={"top": "14mm", "right": "14mm", "bottom": "14mm", "left": "14mm"},
                print_background=True,
                scale=0.9,
            )
            await br.close()
    except Exception as e:
        return {"ok": False, "error": "PDF 를 만들지 못했습니다: " + str(e)[:300]}

    print(f"[pdf] {title} — {len(pdf)}B", flush=True)
    return {"ok": True, "name": f"{title}.pdf", "data": _b64.b64encode(pdf).decode()}


@app.post("/api/wiki/import-docx")
async def wiki_import_docx(payload: dict):
    """워드(.docx) 를 위키가 읽을 수 있는 HTML 로 푼다.

    브라우저는 .docx 를 못 읽는다 — 압축 파일이라 풀어야 하고, 그림은 그 안에
    따로 들어 있다. 그래서 서버가 푼다.

    **그림은 문서에 함께 담는다**(data URI). 파일 서버를 따로 두지 않아도 되고,
    문서를 옮기거나 내보내도 그림이 안 깨진다. 다만 무한정 담지는 않는다 —
    큰 사진이 몇 장만 있어도 문서가 수십 MB 가 되어 여는 것부터 느려진다.

    **표 안의 표**가 까다롭다. 편집기의 표는 칸 안에 표를 담지 못한다. 그렇다고
    버리면 내용이 사라지므로, 안쪽 표를 **바깥 표 뒤로 떼어** 내고 원래 자리에는
    「여기에 표가 있었다」 는 표시를 남긴다. 모양은 펴지지만 잃는 글자는 없다.
    """
    import base64 as _b64, io as _io
    raw = str(payload.get("data") or "")
    if raw.strip().startswith("data:") and "," in raw:
        raw = raw.split(",", 1)[1]
    try:
        blob = _b64.b64decode(raw)
    except Exception as e:
        return {"ok": False, "error": "파일을 읽지 못했습니다: " + str(e)[:120]}
    if not blob:
        return {"ok": False, "error": "빈 파일입니다"}

    try:
        import mammoth
        from bs4 import BeautifulSoup
    except Exception as e:
        return {"ok": False, "error": "변환기를 불러오지 못했습니다: " + str(e)[:120]}

    # 그림 — 너무 큰 것은 줄여 담는다. 원본 그대로면 문서가 못 열 만큼 무거워진다.
    def _img(image):
        with image.open() as f:
            data = f.read()
        ctype = (image.content_type or "image/png")
        try:
            from PIL import Image
            im = Image.open(_io.BytesIO(data))
            if im.width > 1600:
                im = im.convert("RGB") if im.mode in ("P", "CMYK") else im
                h = max(1, round(im.height * 1600 / im.width))
                im = im.resize((1600, h))
                buf = _io.BytesIO()
                im.save(buf, format="PNG")
                data, ctype = buf.getvalue(), "image/png"
        except Exception:
            pass  # 못 줄이면 원본 그대로 — 그림 하나 때문에 가져오기를 막지 않는다
        return {"src": f"data:{ctype};base64,{_b64.b64encode(data).decode()}"}

    try:
        res = mammoth.convert_to_html(_io.BytesIO(blob), convert_image=mammoth.images.img_element(_img))
        html = res.value or ""
    except Exception as e:
        return {"ok": False, "error": "워드 문서를 푸는 데 실패했습니다: " + str(e)[:200]}

    # mammoth 가 빈손이면 **직접 뜯는다.**
    #
    # mammoth 는 스타일·번호 정의가 온전한 문서를 전제한다. 다른 도구가 만든
    # docx, 옛 문서를 변환한 docx 는 그 전제를 깨서 빈 결과가 나온다. 그럴 때
    # 「못 가져왔다」 로 끝내면 사람은 방법이 없다 — 글자는 분명히 문서 안에
    # 있는데도.
    #
    # word/document.xml 을 열어 문단(w:p)과 표(w:tbl)만 곧이곧대로 옮긴다.
    # 서식은 잃지만 **글과 표 구조는 살아 남는다.** 아무것도 못 가져오는 것보다
    # 낫고, 사람이 문서에서 이어 고칠 수 있다.
    if not html.strip():
        try:
            import zipfile as _zip, re as _re2, html as _h
            with _zip.ZipFile(_io.BytesIO(blob)) as z:
                xml = z.read("word/document.xml").decode("utf-8", "ignore")

            def _text(node: str) -> str:
                return _h.escape("".join(_re2.findall(r"<w:t[^>]*>(.*?)</w:t>", node, _re2.S)))

            out = []
            # 표와 문단을 **나온 차례대로** 훑는다 — 문단만 먼저 모으면 표가
            # 문서 끝으로 밀려 읽는 차례가 바뀐다.
            for m in _re2.finditer(r"<w:tbl>.*?</w:tbl>|<w:p[ >].*?</w:p>", xml, _re2.S):
                blk = m.group(0)
                if blk.startswith("<w:tbl"):
                    rows = []
                    for tr in _re2.findall(r"<w:tr[ >].*?</w:tr>", blk, _re2.S):
                        cells = _re2.findall(r"<w:tc[ >].*?</w:tc>", tr, _re2.S)
                        rows.append("<tr>" + "".join(f"<td>{_text(c)}</td>" for c in cells) + "</tr>")
                    if rows:
                        out.append("<table>" + "".join(rows) + "</table>")
                else:
                    t = _text(blk)
                    if not t.strip():
                        continue
                    lvl = _re2.search(r'w:pStyle w:val="Heading(\d)"', blk)
                    out.append(f"<h{lvl.group(1)}>{t}</h{lvl.group(1)}>" if lvl else f"<p>{t}</p>")
            if out:
                html = "".join(out)
                print(f"[docx] mammoth 빈손 → 직접 뜯음: {len(out)}덩이", flush=True)
        except Exception as e:
            print(f"[docx] 직접 뜯기도 실패: {str(e)[:120]}", flush=True)

    # 푼 결과가 비었으면 **왜 비었는지** 말한다.
    #
    # 200 으로 답했는데 화면은 「못 가져왔다」 만 띄우면, 서버 탓인지 문서 탓인지
    # 알 길이 없다(지적: 200인데 못 가져왔다고 한다). 흔한 까닭은 옛 .doc 를
    # 이름만 .docx 로 바꾼 경우다 — 속이 전혀 다른 형식이라 풀리지 않는다.
    if not html.strip():
        why = "문서에서 옮길 내용을 찾지 못했습니다"
        if blob[:2] == b"\xd0\xcf":
            why = "옛 워드(.doc) 형식입니다 — 워드에서 「다른 이름으로 저장 → .docx」 한 뒤 다시 해 주세요"
        elif blob[:2] != b"PK":
            why = "워드 문서가 아닙니다(.docx 가 아님)"
        elif not (res.messages or []):
            why = "문서가 비어 있습니다"
        print(f"[docx] 빈 결과 — {len(blob)}B 머리={blob[:4]!r} :: {why}", flush=True)
        return {
            "ok": False,
            "error": why,
            "bytes": len(blob),
            # 파일 머리 네 글자 — 사람이 화면에서 바로 읽고 갈릴 수 있게
            "head": blob[:4].hex(" "),
            "messages": [str(m) for m in (res.messages or [])][:10],
        }

    # 표 안의 표를 바깥으로 떼어 내고, **글자 크기는 버린다.**
    moved = 0
    try:
        soup = BeautifulSoup(html, "html.parser")

        # 워드에서 온 글자 크기를 걷어낸다.
        #
        # 워드 문서는 같은 「제목 2」 인데도 12pt 인 것과 10pt 인 것이 섞여
        # 있다(실제 원본이 그랬다). mammoth 는 그 크기를 글자마다 그대로
        # 옮기므로, 편집기에서도 제목이 어떤 건 크고 어떤 건 작다(지적).
        #
        # 크기가 아니라 **제목 몇 단인가**만 가져온다. 그러면 워드가 어떻든
        # 이 문서의 눈금으로 정돈된다 — 위키는 문서마다 제목 크기가 달라지면
        # 목차로 읽히지 않는다. 색·굵기 같은 다른 꾸밈은 건드리지 않는다.
        import re as _restyle
        for el in soup.find_all(style=True):
            st = _restyle.sub(r"font-size\s*:[^;]*;?", "", el["style"]).strip()
            if st:
                el["style"] = st
            else:
                del el["style"]
        for outer in list(soup.find_all("table")):
            inners = [t for t in outer.find_all("table") if t is not outer]
            for inner in inners:
                mark = soup.new_tag("p")
                moved += 1
                mark.string = f"[표 {moved}] — 아래에 이어집니다"
                inner.replace_with(mark)
                cap = soup.new_tag("p")
                cap.string = f"[표 {moved}] 위 표 안에 있던 표"
                outer.insert_after(inner)
                outer.insert_after(cap)
        html = str(soup)
    except Exception:
        pass  # 못 펴도 가져오기는 계속한다

    print(
        f"[docx] {len(blob)}B 머리={blob[:4]!r} html={len(html)}글자 "
        f"표속표={moved} 알림={len(res.messages or [])}",
        flush=True,
    )
    return {
        "ok": True,
        "html": html,
        "nested_tables": moved,
        # 무엇이 안 넘어왔는지 사람이 알아야 한다 — 조용히 빠지면 나중에 찾는다
        "messages": [str(m) for m in (res.messages or [])][:20],
    }


@app.post("/api/wiki/{pid}")
async def wiki_save(pid: str, payload: dict, request: Request):
    """만들기·고치기 공통.

    저장할 때마다 **지난 판을 한 줄 남긴다**(wiki_rev). 되돌릴 수 있어야 사람이
    마음 놓고 고친다 — 못 되돌리면 지우기가 무서워 문서가 안 정리된다.
    """
    s = getattr(request.state, "user", None)
    who = (s or {}).get("username") or ""
    title = str(payload.get("title") or "")
    body = payload.get("body")
    if body is None:
        body = []
    plain = _wiki_plain(body)
    async with db.pool().acquire() as c:
        old = await c.fetchrow("SELECT title, body FROM wiki_page WHERE id=$1", pid)
        if old:
            await c.execute(
                "INSERT INTO wiki_rev (page_id, title, body, who) VALUES ($1,$2,$3::jsonb,$4)",
                pid, old["title"], old["body"] if isinstance(old["body"], str) else json.dumps(old["body"], ensure_ascii=False), who,
            )
            await c.execute(
                "UPDATE wiki_page SET title=$2, body=$3::jsonb, plain=$4, updated_by=$5, "
                "updated_at=now() WHERE id=$1",
                pid, title, json.dumps(body, ensure_ascii=False), plain, who,
            )
        else:
            await c.execute(
                "INSERT INTO wiki_page (id, project, parent_id, title, body, plain, ord, created_by, updated_by) "
                "VALUES ($1,$2,$3,$4,$5::jsonb,$6,$7,$8,$8)",
                pid, str(payload.get("project") or ""), payload.get("parent_id"),
                title, json.dumps(body, ensure_ascii=False), plain,
                int(payload.get("ord") or 0), who,
            )
    try: asyncio.create_task(broadcast({"type": "wiki_updated", "id": pid}))
    except Exception: pass
    return {"ok": True, "id": pid}


@app.patch("/api/wiki/{pid}")
async def wiki_patch(pid: str, payload: dict):
    """자리 옮기기·이름 바꾸기 — 본문은 안 건드린다(지난 판도 안 남긴다)."""
    sets, args = [], []
    for k in ("title", "parent_id", "project", "ord"):
        if k in payload:
            args.append(payload[k])
            sets.append(f"{k}=${len(args)}")
    if not sets:
        return {"ok": True}
    sets.append("updated_at=now()")
    args.append(pid)
    async with db.pool().acquire() as c:
        await c.execute(f"UPDATE wiki_page SET {', '.join(sets)} WHERE id=${len(args)}", *args)
    return {"ok": True}


@app.delete("/api/wiki/{pid}")
async def wiki_delete(pid: str):
    """지운다. **아래 문서가 있으면 안 지운다** — 통째로 사라지면 되돌릴 수 없다."""
    async with db.pool().acquire() as c:
        kid = await c.fetchval("SELECT count(*) FROM wiki_page WHERE parent_id=$1", pid)
        if kid:
            raise HTTPException(400, f"아래 문서가 {kid}개 있습니다 — 먼저 옮기거나 지우세요")
        await c.execute("DELETE FROM wiki_page WHERE id=$1", pid)
    return {"ok": True}


@app.get("/api/wiki/{pid}/revs")
async def wiki_revs(pid: str, limit: int = 30):
    async with db.pool().acquire() as c:
        rows = await c.fetch(
            "SELECT id, title, who, at FROM wiki_rev WHERE page_id=$1 ORDER BY at DESC LIMIT $2",
            pid, max(1, min(200, limit)),
        )
    return {"revs": [{**dict(r), "at": r["at"].isoformat()} for r in rows]}


@app.get("/api/wiki/rev/{rev_id}")
async def wiki_rev_get(rev_id: int):
    """지난 판 하나. 되돌리려면 그때의 본문이 있어야 한다."""
    async with db.pool().acquire() as c:
        r = await c.fetchrow("SELECT * FROM wiki_rev WHERE id=$1", rev_id)
    if not r:
        raise HTTPException(404, "그 판을 찾을 수 없습니다")
    d = dict(r)
    d["at"] = d["at"].isoformat()
    if isinstance(d.get("body"), str):
        try:
            d["body"] = json.loads(d["body"])
        except Exception:
            d["body"] = []
    return {"rev": d}


@app.get("/api/audit")
async def audit_list_api(limit: int = 300):
    """수정 이력 — 알림 종이 읽는다. 최신이 앞."""
    return {"items": await db.audit_list(max(1, min(1000, limit)))}


@app.get("/api/transfer/export")
async def transfer_export(parts: str = "", secrets: int = 0):
    _require_admin("")  # 미들웨어 세션이 컨텍스트에 있다
    want = {x.strip() for x in parts.split(",") if x.strip()} or set(_TRANSFER_PARTS)
    out = {"app": "utop", "version": 1,
           "exported_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"), "parts": {}}
    P = out["parts"]
    if "req" in want:
        P["req"] = {"categories": await db.cat_list(), "reqs": await db.req_list_full()}
    if "tc" in want:
        P["tc"] = {"tcs": await db.tc_list_full()}
    if "cycle" in want:
        P["cycle"] = {"cycles": await db.cycle_list_full()}
    if "defect" in want:
        P["defect"] = {"defects": await db.defect_list(limit=100000)}
    if "device" in want:
        devs = await db.device_list()
        if not secrets:
            # 비밀번호는 기본 제외 — 파일이 어디로 돌지 모른다
            for d in devs:
                d.pop("password", None)
                d.pop("enable_password", None)
                for a in d.get("access") or []:
                    a.pop("password", None)
                    a.pop("enable_password", None)
        P["device"] = {"devices": devs, "secrets": bool(secrets)}
    if "catalog" in want:
        P["catalog"] = {"items": await db.catalog_list(),
                        "racks": _kv_load_sync("racks", {}) or {}}
    if "settings" in want:
        cfs = await db.cf_list("")
        P["settings"] = {
            "custom_fields": cfs,
            "global_params": _load_global_params(),
            "branding": _load_branding(),
        }
    return out


@app.post("/api/transfer/import")
async def transfer_import(payload: dict):
    """합치기(upsert) — 같은 ID 는 덮고 없는 것은 만든다. 지우지는 않는다."""
    _require_admin("")
    parts = payload.get("parts") or {}
    done: dict = {}

    if "req" in parts:
        n = 0
        for c in parts["req"].get("categories") or []:
            if c.get("id") and c.get("name"):
                await db.cat_upsert(str(c["id"]), str(c["name"]), c.get("parent_id"),
                                    int(c.get("sort_order") or 0))
        for r in parts["req"].get("reqs") or []:
            rid = str(r.get("id") or "").strip()
            if not rid:
                continue
            try:
                await db.req_upsert(rid, _strip_derived(r))
                n += 1
            except Exception as e:
                done.setdefault("_errors", []).append(f"요구사항 {rid}: {e}")
        done["req"] = n

    if "tc" in parts:
        n = 0
        for t in parts["tc"].get("tcs") or []:
            tid = str(t.get("tcid") or t.get("id") or "").strip()
            if not tid:
                continue
            try:
                await db.tc_upsert(tid, _strip_derived(t))
                n += 1
            except Exception as e:
                done.setdefault("_errors", []).append(f"시험 {tid}: {e}")
        done["tc"] = n

    if "cycle" in parts:
        n = 0
        for cyc in parts["cycle"].get("cycles") or []:
            cid = str(cyc.get("id") or "").strip()
            if not cid:
                continue
            try:
                await db.cycle_upsert(cid, _strip_derived(cyc))
                n += 1
            except Exception as e:
                done.setdefault("_errors", []).append(f"사이클 {cid}: {e}")
        done["cycle"] = n

    if "defect" in parts:
        n = 0
        for d in parts["defect"].get("defects") or []:
            did = str(d.get("id") or "").strip()
            if not did:
                continue
            try:
                if await db.defect_get(did):
                    await db.defect_update(did, d)
                else:
                    await db.defect_create(d)
                n += 1
            except Exception as e:
                done.setdefault("_errors", []).append(f"결함 {did}: {e}")
        done["defect"] = n

    if "device" in parts:
        n = 0
        errs: list = []
        for d in parts["device"].get("devices") or []:
            ip = str(d.get("ip") or "").strip()
            if not ip:
                continue
            try:
                # 장비의 실질 키는 IP 다. 서버마다 id 를 다르게 만들어 둬서,
                # 239의 id 로 넣으면 id 충돌은 안 나고 ip UNIQUE 에 걸려
                # 통째로 500 이 났다 — IP 로 찾은 기존 장비의 id 를 입힌다.
                cur = await db.device_get(ip)
                if cur:
                    d["id"] = cur["id"]
                    # 비밀번호 없이 온 파일이면 기존 비밀번호를 지킨다 —
                    # upsert 가 전 칸을 쓰므로 그냥 넣으면 빈 값으로 덮인다
                    for k in ("password", "enable_password"):
                        if not d.get(k):
                            d[k] = cur.get(k)
                    accs = {a.get("protocol"): a for a in (cur.get("access") or [])}
                    for a in d.get("access") or []:
                        old = accs.get(a.get("protocol")) or {}
                        for k in ("password", "enable_password"):
                            if not a.get(k):
                                a[k] = old.get(k)
                await db.device_upsert(_strip_derived(d))
                n += 1
            except Exception as e:
                if len(errs) < 10:
                    errs.append(f"장비 {ip}: {e}")
        done["device"] = n
        if errs:
            done.setdefault("_errors", []).extend(errs)

    if "catalog" in parts:
        n = 0
        for it in parts["catalog"].get("items") or []:
            if it.get("kind") and it.get("name"):
                await db.catalog_upsert(_strip_derived(it))
                n += 1
        racks = parts["catalog"].get("racks")
        if isinstance(racks, dict) and (racks.get("racks") or racks.get("labs")):
            _kv_save_sync("racks", racks)
        done["catalog"] = n

    if "settings" in parts:
        st = parts["settings"]
        n = 0
        for cf in st.get("custom_fields") or []:
            try:
                await db.cf_upsert({k: v for k, v in cf.items() if k not in ("used",)})
                n += 1
            except Exception:
                pass
        gp = st.get("global_params")
        if isinstance(gp, dict) and gp:
            GLOBAL_PARAMS_FILE.write_text(
                json.dumps(gp, ensure_ascii=False, indent=2), encoding="utf-8")
        br = st.get("branding")
        if isinstance(br, dict) and br:
            save_json(BRANDING_FILE, br)
        done["settings"] = n
    errs = done.pop("_errors", [])
    return {"ok": True, "done": done, "errors": errs[:10],
            "error_count": len(errs)}


@app.get("/api/rackview")
async def rackview():
    """랙뷰 한 판 — 랙 틀(KV 'racks') + PG 장비 배치 + 아직 안 옮긴 옛 배치.

    옛 devices.json 의 배치는 IP 로 겹침을 가른다: 같은 IP 가 PG 에 있으면
    PG 가 정본이고, 없으면 회색 유령으로 보여 준다 — 랙에 꽂혀 있는 것은
    사실이니 숨기지 않는다(숨김 금지 원칙).
    """
    kv = _ensure_blank_ids(_kv_load_sync("racks", {}) or {})
    racks = kv.get("racks") or []
    # 인터페이스 이름까지 싣는다 — 호버 카드의 포트 형상 몫
    devs = await db.device_list(with_ifs=True)
    placed, unplaced, pg_ips = [], [], set()
    for d in devs:
        ip = str(d.get("ip") or "").strip()
        if ip:
            pg_ips.add(ip)
        slim = {
            "id": d["id"], "ip": ip, "name": d.get("name"), "model": d.get("model"),
            "lab": d.get("lab"), "role": d.get("role"), "vendor": d.get("vendor"),
            "rack_units": d.get("rack_units"), "power_w": d.get("power_w"),
            "ifs": [str(i.get("name") or "") for i in (d.get("interfaces") or [])],
        }
        if d.get("rack_id") and d.get("rack_pos"):
            placed.append({
                **slim, "source": "pg",
                "rack_id": d["rack_id"], "rack_pos": d["rack_pos"],
                "rack_units": d.get("rack_units") or 1,
                "access": [
                    {"protocol": a.get("protocol"), "status": a.get("last_status"),
                     "enabled": a.get("enabled")}
                    for a in (d.get("access") or [])
                ],
            })
        else:
            unplaced.append(slim)
    legacy = []
    try:
        old = load_json(DEVICES_FILE) or {}
        by_name = {str(r.get("name") or ""): str(r.get("id") or "") for r in racks}
        for d in old.get("devices", []) or []:
            ip = str(d.get("ip") or "").strip()
            rid = d.get("rack_id") or by_name.get(str(d.get("rack_name") or ""))
            pos = d.get("rack_pos")
            if not rid or not pos or (ip and ip in pg_ips):
                continue
            legacy.append({
                "ip": ip, "name": d.get("id") or d.get("name") or ip,
                "model": d.get("model"), "lab": d.get("lab"), "source": "legacy",
                "rack_id": rid, "rack_pos": int(pos),
                "rack_units": int(d.get("rack_units") or 1),
            })
    except Exception:
        pass
    return {
        "labs": kv.get("labs") or [],
        "racks": racks,
        "blanks": kv.get("blanks") or [],
        "devices": placed + legacy,
        "unplaced": unplaced,
        # 부품 견본 — 비어 있으면 화면이 기본 팔레트를 쓴다
        "part_presets": kv.get("part_presets") or [],
    }

# ───────────────────────────────────────────
# 라우터 - STC (Spirent) 트래픽 실행 (py2.7 스크립트 subprocess)
# ───────────────────────────────────────────
_stc_proc = {"p": None}

@app.post("/api/stc/run")
async def stc_run(data: dict):
    import os, shlex
    script = (data.get("script_path") or "").strip()
    py = (data.get("py_cmd") or "py -2.7").strip()
    if not script or not os.path.exists(script):
        return {"ok": False, "error": "STC 스크립트를 찾을 수 없습니다: " + (script or "(빈 경로)")}
    p = _stc_proc.get("p")
    if p is not None and p.poll() is None:
        return {"ok": False, "error": "이미 실행 중입니다. 먼저 중지하세요."}
    try:
        cmd = shlex.split(py, posix=False) + [script]
    except Exception:
        cmd = py.split() + [script]
    env = dict(os.environ)
    env["STC_GUI_CHILD"] = "1"
    # U-TOP 폼 파라미터 → 스크립트 오버라이드 (params 가 있으면 STC_OVERRIDE=1)
    params = data.get("params") or {}
    if params:
        env["STC_OVERRIDE"] = "1"
        def _setenv(k, v):
            if v is not None and str(v) != "":
                env[k] = str(v)
        _setenv("STC_REST_IP", params.get("restIp"))
        _setenv("STC_REST_PORT", params.get("restPort"))
        _setenv("STC_USER", params.get("user"))
        _setenv("STC_SESSION", params.get("session"))
        _setenv("STC_CHASSIS", params.get("chassis"))
        _setenv("STC_PORT_A", params.get("portA"))
        _setenv("STC_PORT_B", params.get("portB"))
        _setenv("STC_FRAME", params.get("frame"))
        _setenv("STC_LOAD", params.get("load"))
        _setenv("STC_LOAD_UNIT", params.get("loadUnit"))
        _setenv("STC_PROTO", params.get("proto"))
        _setenv("STC_DST_PORT", params.get("dstPort"))
        _setenv("STC_SRC_PORT", params.get("srcPort"))
        _setenv("STC_DURATION", params.get("duration"))
        _setenv("STC_INTERVAL", params.get("interval"))
        _setenv("STC_BIDIR", "1" if params.get("bidir") else "0")
        _setenv("STC_DEVA_IP", params.get("devAip"))
        _setenv("STC_DEVA_GW", params.get("devAgw"))
        _setenv("STC_DEVA_MAC", params.get("devAmac"))
        _setenv("STC_DEVB_IP", params.get("devBip"))
        _setenv("STC_DEVB_GW", params.get("devBgw"))
        _setenv("STC_DEVB_MAC", params.get("devBmac"))
    # REST 대상이 로컬이면 stcweb.exe REST 서버를 자동 기동.
    # 안 알려 줬으면 계측기 등록에서 찾는다 — 이 서버는 리눅스라 localhost 엔 없다
    _rip0, rest_port = await _stc_rest_for(
        str((params.get("chassis") if params else "") or ""),
        (params.get("restIp") if params else ""),
        (params.get("restPort") if params else None),
    )
    rest_ip = str(_rip0).strip().lower()
    if params is not None:
        params["restIp"], params["restPort"] = _rip0, rest_port
        _setenv("STC_REST_IP", _rip0)
        _setenv("STC_REST_PORT", rest_port)
    if rest_ip in ("localhost", "127.0.0.1", "") and not _port_listening(rest_port):
        await broadcast({"type": "stc_line", "line": "[REST] localhost:" + str(rest_port) + " 미기동 → 서버 자동 시작"})
        srv = await stc_server_start({"port": rest_port})
        if not srv.get("ok"):
            return {"ok": False, "error": "REST 서버 시작 실패: " + str(srv.get("error"))}
    workdir = os.path.dirname(script) or "."
    # Windows + uvicorn(SelectorEventLoop) 환경에서 asyncio.create_subprocess_exec 는
    # NotImplementedError 가 나므로, 동기 Popen + 스레드로 stdout 을 읽어 WebSocket 으로 흘린다.
    try:
        proc = subprocess.Popen(
            cmd, cwd=workdir,
            stdout=subprocess.PIPE, stderr=subprocess.STDOUT, env=env, bufsize=1)
    except Exception as e:
        return {"ok": False, "error": "실행 실패: " + str(e) + " (py/stcrestclient 설치 확인)"}
    _stc_proc["p"] = proc
    loop = asyncio.get_event_loop()
    await broadcast({"type": "stc_start", "cmd": " ".join(cmd)})

    def _emit(msg):
        try:
            asyncio.run_coroutine_threadsafe(broadcast(msg), loop)
        except Exception:
            pass

    def _pump():
        try:
            for raw in iter(proc.stdout.readline, b""):
                txt = raw.decode("utf-8", "replace").rstrip("\r\n")
                _emit({"type": "stc_line", "line": txt})
            rc = proc.wait()
            _emit({"type": "stc_done", "code": rc})
        except Exception as e:
            _emit({"type": "stc_line", "line": "[backend error] " + str(e)})
            _emit({"type": "stc_done", "code": -1})

    threading.Thread(target=_pump, daemon=True).start()
    return {"ok": True, "cmd": " ".join(cmd)}

@app.post("/api/stc/stop")
async def stc_stop():
    p = _stc_proc.get("p")
    if p is not None and p.poll() is None:
        try:
            p.terminate()
        except Exception:
            pass
        return {"ok": True, "stopped": True}
    return {"ok": True, "stopped": False}


# ── 위저드 트래픽 전송: 수집한 설정(JSON)으로 stc_traffic.py 통빌드+전송, 결과 스트리밍 ──
_stc_traffic_proc = {"p": None, "stop": None}

@app.post("/api/stc/traffic/run")
async def stc_traffic_run(data: dict):
    import os
    import json as _json
    p = _stc_traffic_proc.get("p")
    if p is not None and p.poll() is None:
        return {"ok": False, "error": "이미 전송 중입니다. 먼저 정지하세요."}
    _rip1, rest_port = await _stc_rest_for(str(data.get("chassis") or ""), data.get("restIp"), data.get("restPort"))
    rest_ip = str(_rip1).strip().lower()
    data["restIp"], data["restPort"] = _rip1, rest_port   # 전송 스크립트도 같은 곳을 본다
    if rest_ip in ("localhost", "127.0.0.1", "") and not _port_listening(rest_port):
        srv = await stc_server_start({"port": rest_port})
        if not srv.get("ok"):
            return {"ok": False, "error": "REST 서버 시작 실패: " + str(srv.get("error"))}
    cfgdir = str(Path(__file__).parent)
    cfgpath = os.path.join(cfgdir, "stc", "_stc_traffic_cfg.json")
    stoppath = os.path.join(cfgdir, "_stc_traffic_stop")
    try:
        if os.path.exists(stoppath):
            os.remove(stoppath)
    except Exception:
        pass
    try:
        with open(cfgpath, "w") as f:
            _json.dump(data, f)
    except Exception as e:
        return {"ok": False, "error": "config 쓰기 실패: " + repr(e)}
    # 전송 후 위저드 예약(U_TOP_op) 복원/사전해제용 정보
    _ports_csv = ",".join([str(x).split("/")[-2] + "/" + str(x).split("/")[-1] for x in (data.get("ports") or [])])
    _user = str(data.get("user") or "admin")
    _chassis = str(data.get("chassis") or "192.168.5.100").strip()
    _rip = data.get("restIp") or rest_ip or "localhost"
    _rport = rest_port
    # 전송 직전: 위저드 예약(U_TOP_op)을 빠르게 해제 → 포트 free → U_TOP_tx 가 RevokeOwner(~50초) 없이 즉시 예약.
    #  (전송 스크립트가 예약 후 레지스트리에 user 를 다시 써서 위저드 '내 예약' 표시를 유지함)
    if _ports_csv:
        try:
            await asyncio.to_thread(_run_stc_helper, "releaseports", _chassis, _rip, _rport,
                                    {"user": _user, "ports": _ports_csv, "fast": True})
        except Exception:
            pass
    script = os.path.join(cfgdir, "stc", "stc_traffic.py")
    cmd = [sys.executable, "-u", script, cfgpath, stoppath]   # -u: 버퍼링 없이 실시간 stdout
    env = dict(os.environ)
    env["STC_GUI_CHILD"] = "1"
    env["PYTHONUNBUFFERED"] = "1"
    try:
        proc = subprocess.Popen(cmd, cwd=cfgdir, stdout=subprocess.PIPE,
                                stderr=subprocess.STDOUT, env=env, bufsize=1)
    except Exception as e:
        return {"ok": False, "error": "실행 실패: " + repr(e)}
    _stc_traffic_proc["p"] = proc
    _stc_traffic_proc["stop"] = stoppath
    loop = asyncio.get_event_loop()
    await broadcast({"type": "stc_start", "cmd": "traffic"})

    def _emit(msg):
        try:
            asyncio.run_coroutine_threadsafe(broadcast(msg), loop)
        except Exception:
            pass

    def _pump():
        try:
            for raw in iter(proc.stdout.readline, b""):
                _emit({"type": "stc_line", "line": raw.decode("utf-8", "replace").rstrip("\r\n")})
            rc = proc.wait()
            _emit({"type": "stc_line", "line": "[전송 종료] code=" + str(rc)})
            # U_TOP_op 를 그대로 써서 포트 예약이 유지되므로 복원 불필요. 상태 캐시만 무효화.
            try:
                _stc_status_cache.pop(_chassis, None)
            except Exception:
                pass
            _emit({"type": "stc_done", "code": rc})
        except Exception as e:
            _emit({"type": "stc_line", "line": "[backend error] " + repr(e)})
            _emit({"type": "stc_done", "code": -1})

    threading.Thread(target=_pump, daemon=True).start()
    return {"ok": True}

@app.post("/api/stc/traffic/stop")
async def stc_traffic_stop(data: dict = None):
    # 정지 신호 파일 생성 → 스크립트가 트래픽 정지 후 포트 해제/세션 종료(정리 보장).
    sp = _stc_traffic_proc.get("stop")
    if sp:
        try:
            with open(sp, "w") as f:
                f.write("stop")
        except Exception:
            pass
    return {"ok": True, "stopped": True}

# ───────────────────────────────────────────
# 라우터 - STC REST API 서버 (stcweb.exe) 관리
# stcrestclient 는 이 REST 서버에 붙고, 서버가 섀시로 연결한다.
# 서버 옆 stcweb.yaml 에 service addr(:8888)·대상 섀시가 설정돼 있다.
# ───────────────────────────────────────────
_stcweb_proc = {"p": None}

# stcweb.exe 후보 경로 (설치본 위치). 환경변수 STCWEB_PATH 로 덮어쓸 수 있다.
_STCWEB_CANDIDATES = [
    os.environ.get("STCWEB_PATH", ""),
    r"C:\Program Files\Spirent Communications\Spirent TestCenter 5.23\Spirent TestCenter Application\stcweb.exe",
    r"D:\Spirent Communications\Spirent TestCenter 5.23\Spirent TestCenter Application\stcweb.exe",
    r"C:\Spirent Communications\Spirent TestCenter 5.23\Spirent TestCenter Application\stcweb.exe",
]

# 버전(5.23 등)이 달라도 찾도록 설치 폴더를 와일드카드로 탐색 (최신 버전 우선)
_STCWEB_GLOBS = [
    r"C:\Program Files\Spirent Communications\Spirent TestCenter *\Spirent TestCenter Application\stcweb.exe",
    r"C:\Program Files (x86)\Spirent Communications\Spirent TestCenter *\Spirent TestCenter Application\stcweb.exe",
    r"D:\Spirent Communications\Spirent TestCenter *\Spirent TestCenter Application\stcweb.exe",
    r"C:\Spirent Communications\Spirent TestCenter *\Spirent TestCenter Application\stcweb.exe",
]

def _find_stcweb():
    for c in _STCWEB_CANDIDATES:
        if c and os.path.exists(c):
            return c
    import glob as _glob
    for pat in _STCWEB_GLOBS:
        hits = sorted(_glob.glob(pat))
        if hits:
            return hits[-1]  # 정렬상 마지막 = 최신 버전
    return None

def _port_listening(port, host="127.0.0.1"):
    s = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    s.settimeout(0.5)
    try:
        return s.connect_ex((host, int(port))) == 0
    except Exception:
        return False
    finally:
        s.close()

@app.get("/api/stc/server/status")
async def stc_server_status(port: int = 8888):
    p = _stcweb_proc.get("p")
    managed = bool(p is not None and p.poll() is None)
    return {
        "listening": _port_listening(port),
        "managed": managed,
        "exe": _find_stcweb(),
    }

# ── IXIA N2X 트래픽 시험 (상주 데몬: n2xtclsh85 가 세션을 계속 유지) ──
# 실행 PC 가 리눅스로 바뀌면 이 경로는 없다. 코드를 고치지 않고 바꿀 수
# 있게 환경변수로 뺀다 — 리눅스용 N2X Tcl 이 있으면 그 경로를 넣으면 된다.
N2X_TCLSH = os.environ.get("N2X_TCLSH") or r"C:\N2xTcl85\bin\n2xtclsh85.exe"

# ── N2X 중계 ────────────────────────────────────────────────
#
# N2X 는 STC 와 처지가 다르다. STC 는 REST 서버를 **네트워크 너머로**
# 가리킬 수 있어서 리눅스 백엔드에서도 붙지만, N2X 는 백엔드가 있는 그
# 기계에서 Tcl 프로세스를 직접 띄운다. 실행 PC 를 리눅스로 옮기는 순간
# 이 길이 끊긴다.
#
# 그래서 윈도우 PC 한 대에 백엔드를 하나 더 띄우고 N2X 명령만 그리로
# 넘긴다. 시험은 리눅스가 돌리고, Tcl 만 건너간다.
#
#   리눅스 백엔드  ──HTTP──▶  윈도우 백엔드  ──▶  N2X Tcl  ──▶  섀시
#     N2X_RELAY_URL              N2X_RELAY_KEY
#
# 비워 두면 예전처럼 이 기계에서 직접 띄운다 — 실행 PC 가 윈도우면 아무
# 설정도 필요 없다. 즉 어느 쪽을 골라도 이 코드 하나로 된다.
N2X_RELAY_URL = (os.environ.get("N2X_RELAY_URL") or "").rstrip("/")
# 중계는 로그인 세션이 아니라 이 열쇠로 연다. 계측기를 아무나 못 돌리게
# 하려면 양쪽에 같은 값을 넣어야 한다.
N2X_RELAY_KEY = os.environ.get("N2X_RELAY_KEY") or ""
# 중계 전용으로 뜬다 — DB 도 RAG 도 잡지 않는다.
#
# 중계는 N2X 앱 서버(윈도우) 위에 올라간다. 그 기계는 계측기를 돌리는 것이
# 일이지 시험 자료를 들고 있지 않다. 거기에 PostgreSQL 을 물리게 하면
# 랩 네트워크에 구멍을 하나 더 내는 셈이고, DB 가 잠깐 흔들리면 계측기까지
# 같이 멈춘다.
N2X_RELAY_ONLY = (os.environ.get("N2X_RELAY_ONLY") or "").strip() not in ("", "0", "false")
N2X_DAEMON = os.path.join(os.path.dirname(os.path.abspath(__file__)), "n2x", "n2x_daemon.tcl")
_n2x_daemons = {}            # key "server|label" -> {proc, lock, ready}
_n2x_reg_lock = threading.Lock()

def _n2x_log(msg):
    """콘솔 인코딩(Windows cp949 등)이 유니코드 문자를 못 담아 크래시하는 것을 방지.
    print 실패 시 ASCII-safe 로 대체 후 재출력. HTTP 응답과 무관하게 서버 로그 전용."""
    try:
        print(msg)
    except (UnicodeEncodeError, UnicodeError):
        try:
            print(str(msg).encode("ascii", "replace").decode("ascii"))
        except Exception:
            pass
    except Exception:
        pass

def _n2x_readline(proc, timeout):
    """타임아웃 있는 readline. 시간 초과 시 None 반환(데몬이 멈춘 것으로 간주)."""
    box = {}
    def _r():
        try:
            box["line"] = proc.stdout.readline()
        except Exception as e:
            box["err"] = e
    t = threading.Thread(target=_r, daemon=True)
    t.start()
    t.join(timeout)
    if t.is_alive():
        return None
    if "err" in box:
        raise box["err"]
    return box.get("line", "")

def _n2x_drain_stderr(proc):
    """데몬 stderr 남은 내용 흡수(디버그용). 최대 4KB."""
    try:
        if proc and proc.stderr:
            data = proc.stderr.read(4096) if not proc.stderr.closed else ""
            return (data or "").strip()
    except Exception:
        pass
    return ""

def _n2x_start_daemon(server, label):
    """새 데몬 프로세스 기동 + ready 대기. 성공 시 {proc, lock, ready} 반환, 실패 시 {error}."""
    # 무엇이 없는지 말해 준다. "N2X Tcl 환경 없음" 만 보면 섀시가 안 켜진
    # 건지, 이 서버에 뭘 깔아야 하는 건지 알 수가 없다.
    if not os.path.exists(N2X_TCLSH):
        return {"error": f"N2X Tcl 이 이 서버에 없습니다 — 찾은 곳: {N2X_TCLSH} "
                         f"(N2X_TCLSH 환경변수로 경로를 지정하세요. "
                         f"지금 이 서버는 {platform.system()} 입니다)"}
    if not os.path.exists(N2X_DAEMON):
        return {"error": f"N2X 데몬 스크립트가 없습니다 — {N2X_DAEMON}"}
    try:
        proc = subprocess.Popen(
            [N2X_TCLSH, N2X_DAEMON, server, label],
            stdin=subprocess.PIPE, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            text=True, bufsize=1, encoding="utf-8", errors="replace")
        rl = _n2x_readline(proc, 40)
        if rl is None:
            err = _n2x_drain_stderr(proc)
            try: proc.kill()
            except Exception: pass
            _n2x_log(f"[N2X] daemon start timeout server={server} label={label} stderr={err[:400]}")
            return {"error": "N2X 연결 시간 초과 — 서버 상태 확인 후 다시 시도"}
        ready = (rl or "").strip()
        if not ready:
            err = _n2x_drain_stderr(proc)
            try: proc.kill()
            except Exception: pass
            _n2x_log(f"[N2X] daemon ready empty server={server} label={label} stderr={err[:400]}")
            return {"error": "N2X 데몬 초기화 실패(빈 응답) — 서버 확인"}
        # ready 응답에 error 필드가 있으면(연결 자체 실패) 그것도 전달
        try:
            rj = json.loads(ready)
            if isinstance(rj, dict) and rj.get("ready") is False:
                err_msg = str(rj.get("error", "알 수 없음"))
                try: proc.kill()
                except Exception: pass
                _n2x_log(f"[N2X] daemon ready=false server={server} label={label} err={err_msg}")
                return {"error": "N2X 세션 연결 실패: " + err_msg}
        except Exception:
            pass
        return {"proc": proc, "lock": threading.Lock(), "ready": ready}
    except Exception as e:
        _n2x_log(f"[N2X] daemon spawn exception: {e}")
        return {"error": str(e)}

def _n2x_get_daemon(server, label):
    key = server + "|" + label
    with _n2x_reg_lock:
        d = _n2x_daemons.get(key)
        if d and d["proc"].poll() is None:
            return d
        # 죽은 데몬은 제거하고 새로 기동
        if d:
            _n2x_daemons.pop(key, None)
        nd = _n2x_start_daemon(server, label)
        if "error" in nd:
            return nd
        _n2x_daemons[key] = nd
        return nd

# 데몬이 내주는 잘못 코드 → 사람이 읽는 말.
#
# 문구를 Tcl 안에 한국어로 적어 두었더니 그대로 깨져 나왔다
# (`∞£åφÜ¿φÒ£ ∞èñφè¬ …`). Windows 의 Tcl 은 .tcl 파일을 시스템 인코딩으로
# 읽는데 파일은 UTF-8 이라, 한글이 그 자리에서 어긋난다. 데몬은 ASCII 만
# 말하고 한국어는 여기서 붙인다.
_N2X_ERRS = {
    "no_valid_stream":
        "스트림을 하나도 만들지 못했습니다 — 그 포트를 이 세션이 잡고 있지 않습니다",
}


def _n2x_humanize(res):
    """데몬 응답의 error 를 사람이 읽는 말로 바꾼다. 모르는 것은 그대로 둔다."""
    if not isinstance(res, dict):
        return res
    code = str(res.get("error") or "")
    ko = _N2X_ERRS.get(code)
    if not ko:
        return res
    bad = str(res.get("badPorts") or "").strip()
    if bad:
        ko += f" (못 잡은 포트: {bad}). 「포트 확인」 으로 실제 포트 번호를 보고 Traffic 탭의 시험 포트와 맞추세요"
    out = dict(res)
    out["error"] = ko
    out["code"] = code
    return out


def _n2x_send(server, label, cmd, _retry=True):
    """N2X 명령 한 줄. 중계가 설정돼 있으면 그리로, 아니면 이 기계에서 직접.

    모든 N2X 기능(ports · reserve · release · traffic · ping)이 이 함수
    하나를 지난다. 그래서 여기 한 곳만 갈라 두면 기능마다 따로 손볼 것이
    없다.
    """
    if N2X_RELAY_URL:
        try:
            # N2X 명령은 길다 — ports 스캔이 45초, traffic 이 60초까지 간다.
            # 중계 타임아웃은 그보다 넉넉해야 중간에 끊기지 않는다.
            r = httpx.post(
                N2X_RELAY_URL + "/api/n2x/send",
                json={"server": server, "label": label, "cmd": cmd, "key": N2X_RELAY_KEY},
                timeout=150,
            )
            if r.status_code != 200:
                return {"ok": False,
                        "error": f"N2X 중계 오류 {r.status_code} — {r.text[:200]}"}
            return r.json()
        except Exception as e:
            _n2x_log(f"[N2X] relay failed url={N2X_RELAY_URL} err={e}")
            return {"ok": False,
                    "error": f"N2X 중계에 못 붙었습니다 ({N2X_RELAY_URL}) — {e}"}
    return _n2x_send_local(server, label, cmd, _retry)


def _n2x_send_local(server, label, cmd, _retry=True):
    """상주 데몬에 한 줄 명령 전송 → JSON 응답. 데몬 EOF(=크래시) 감지 시 1회 자동 재기동+재시도.
    NOTE: reserve/release 는 이미 서버에 반영됐을 가능성이 있어 auto-retry 하지 않는다(중복 명령 방지)."""
    # reserve/release 는 부작용 있는 명령 — retry 금지
    _is_side_effect = cmd.startswith("reserve") or cmd.startswith("release")
    if _is_side_effect:
        _retry = False
    d = _n2x_get_daemon(server, label)
    if "error" in d:
        return {"ok": False, "error": d["error"]}
    if d["proc"].poll() is not None:
        with _n2x_reg_lock:
            _n2x_daemons.pop(server + "|" + label, None)
        d = _n2x_get_daemon(server, label)
        if "error" in d:
            return {"ok": False, "error": d["error"]}
    proc = d["proc"]
    # 명령별 timeout — ping 은 짧게, reserve/release 는 넉넉히(AddPort 는 N2X 서버 부하 시 15-25s 소요),
    # ports 는 전체 스캔이라 길게, traffic 은 중간
    if cmd.startswith("ports"):
        tmo = 45
    elif cmd.startswith("ping"):
        tmo = 10
    elif cmd.startswith("reserve") or cmd.startswith("release"):
        tmo = 15   # reserve/release 는 짧게 잡음. 초과 시 데몬 kill 하지 않고 background 로 결과 확인 → 다른 조회 진행 가능
    else:
        tmo = 60
    # 데몬 파이프는 단일이라 lock 으로 명령 직렬화. 다만 앞선 명령이 hang 이면 뒤 요청이 무한 대기 →
    # lock 획득 자체에도 타임아웃을 걸어, 대기 초과 시 즉시 실패로 반환 (사용자가 계속 응답 안 오는 상태 방지).
    # ports 는 캐시(5초)로 대체 가능하니 짧게(3s), 나머지는 명령 timeout+5s 로 설정.
    if cmd.startswith("ports"):
        lock_wait = 3.0
    elif cmd.startswith("ping"):
        lock_wait = 5.0
    else:
        lock_wait = tmo + 5.0
    if not d["lock"].acquire(timeout=lock_wait):
        _n2x_log(f"[N2X] lock acquire timeout({lock_wait}s) cmd={cmd[:40]} -- daemon busy with previous command")
        return {"ok": False, "error": f"데몬 사용 중 — 앞선 명령 대기 시간 초과({lock_wait}s). 잠시 후 다시 시도"}
    _lock_transferred = [False]   # async 반환 시 background 로 lock 소유권 넘기고 여기 finally 에서 release 안 하도록
    try:
        try:
            # 전송 직전 프로세스 죽음 감지 → 자동 재기동
            if proc.poll() is not None:
                err = _n2x_drain_stderr(proc)
                _n2x_log(f"[N2X] daemon dead before send cmd={cmd[:40]} stderr={err[:400]}")
                with _n2x_reg_lock:
                    _n2x_daemons.pop(server + "|" + label, None)
                if _retry:
                    return _n2x_send_local(server, label, cmd, _retry=False)
                return {"ok": False, "error": "데몬 재기동 실패 — 백엔드 로그 확인"}
            proc.stdin.write(cmd + "\n")
            proc.stdin.flush()
            rl = _n2x_readline(proc, tmo)
            if rl is None:
                _n2x_log(f"[N2X] readline timeout cmd={cmd[:40]} tmo={tmo}s -- returning early (background will drain)")
                # reserve/release: 데몬은 살려두고, 백그라운드 스레드가 남은 응답을 읽어 캐시 무효화만 처리
                # → 다른 조회(ports 등) 는 lock 이 즉시 풀리므로 바로 응답 가능
                if cmd.startswith("reserve") or cmd.startswith("release"):
                    def _bg_drain(_d, _s, _l, _c):
                        try:
                            extra = _n2x_readline(_d["proc"], 60)   # 최대 60초 더 대기
                            if extra:
                                _n2x_log(f"[N2X] bg drain done cmd={_c[:40]} resp={extra[:80]}")
                                # 결과에 따라 캐시 무효화 (실제 상태 변화 반영)
                                _n2x_ports_cache_invalidate(_s, _l)
                            else:
                                _n2x_log(f"[N2X] bg drain empty (daemon crashed) cmd={_c[:40]}")
                                with _n2x_reg_lock:
                                    _n2x_daemons.pop(_s + "|" + _l, None)
                        except Exception as _e:
                            _n2x_log(f"[N2X] bg drain exception cmd={_c[:40]} err={_e}")
                        finally:
                            try: _d["lock"].release()
                            except Exception: pass
                    _lock_transferred[0] = True   # finally 에서 release 하지 않도록 표시
                    threading.Thread(target=_bg_drain, args=(d, server, label, cmd), daemon=True).start()
                    return {"ok": False, "async": True,
                            "error": f"N2X 처리 시간이 길어 background 로 확인 중 ({tmo}s+). 잠시 후 재조회 시 반영됩니다.",
                            "hint": "reserve/release 는 실제 서버에 반영됐을 가능성이 큼 — 재조회 시 상태 확인"}
                # ports 등 조회 명령은 기존대로 kill/재기동
                err = _n2x_drain_stderr(proc)
                _n2x_log(f"[N2X] readline timeout cmd={cmd[:40]} tmo={tmo}s stderr={err[:400]}")
                try: proc.kill()
                except Exception: pass
                with _n2x_reg_lock:
                    _n2x_daemons.pop(server + "|" + label, None)
                return {"ok": False, "error": f"데몬 응답 시간 초과({tmo}s) — 다시 시도하세요"}
            line = (rl or "").strip()
            if not line:
                # EOF = 데몬 프로세스 종료. stderr 확인 후 1회 재시도
                err = _n2x_drain_stderr(proc)
                exit_code = proc.poll()
                _n2x_log(f"[N2X] EOF (daemon died) cmd={cmd[:40]} exit={exit_code} stderr={err[:400]}")
                with _n2x_reg_lock:
                    _n2x_daemons.pop(server + "|" + label, None)
                if _retry:
                    return _n2x_send_local(server, label, cmd, _retry=False)
                return {"ok": False, "error": "데몬 응답 없음(연결 끊김): " + (err[:200] if err else "원인 불명 — 백엔드 로그 확인")}
            return json.loads(line)
        except (BrokenPipeError, OSError) as e:
            err = _n2x_drain_stderr(proc)
            _n2x_log(f"[N2X] pipe broken cmd={cmd[:40]} err={e} stderr={err[:400]}")
            with _n2x_reg_lock:
                _n2x_daemons.pop(server + "|" + label, None)
            if _retry:
                return _n2x_send_local(server, label, cmd, _retry=False)
            return {"ok": False, "error": "데몬 파이프 끊김 — 다시 시도"}
        except Exception as e:
            return {"ok": False, "error": str(e)}
    finally:
        if not _lock_transferred[0]:
            try: d["lock"].release()
            except Exception: pass

@app.post("/api/n2x/send")
def n2x_send_relay(data: dict):
    """중계 창구 — 다른 UTOP 백엔드가 보낸 N2X 명령을 이 기계에서 실행한다.

    이 기계에 N2X Tcl 이 깔려 있어야 한다. 리눅스 백엔드가 여기로 넘긴다.

    로그인 세션이 아니라 열쇠로 연다 — 서버끼리 부르는 자리라 사람의
    세션이 없다. 열쇠를 안 정해 두면 아무나 계측기를 돌릴 수 있으므로
    비어 있으면 아예 막는다.
    """
    if not N2X_RELAY_KEY:
        return {"ok": False, "error": "이 서버는 중계로 열려 있지 않습니다 (N2X_RELAY_KEY 없음)"}
    if str(data.get("key") or "") != N2X_RELAY_KEY:
        return {"ok": False, "error": "중계 열쇠가 다릅니다"}
    cmd = str(data.get("cmd") or "").strip()
    if not cmd:
        return {"ok": False, "error": "cmd 필요"}
    return _n2x_send_local(str(data.get("server") or ""), str(data.get("label") or "utop"), cmd)


@app.get("/api/n2x/ping")
def n2x_ping(server: str = "210.1.2.248", label: str = "utop"):
    return _n2x_send(server, label, "ping")

@app.get("/api/n2x/diag")
def n2x_diag(server: str = "210.1.2.248", label: str = "utop"):
    """데몬 상태 진단 — 등록된 데몬 프로세스 목록, alive 여부, stderr 잔여 등."""
    out = {"ok": True, "target": f"{server}|{label}", "daemons": []}
    with _n2x_reg_lock:
        for key, d in _n2x_daemons.items():
            proc = d.get("proc")
            alive = proc and proc.poll() is None
            info = {"key": key, "alive": bool(alive), "pid": proc.pid if proc else None,
                    "ready": d.get("ready", ""), "exit_code": proc.poll() if proc else None}
            out["daemons"].append(info)
    out["target_alive"] = any(d["key"] == f"{server}|{label}" and d["alive"] for d in out["daemons"])
    return out

def _n2x_local_ver() -> int:
    """저장소에 있는 n2x_daemon.tcl 의 판. 없으면 0."""
    try:
        import re as _re
        with open(N2X_DAEMON, encoding="utf-8") as f:
            mm = _re.search(r"^set DAEMON_VER (\d+)", f.read(), _re.M)
        return int(mm.group(1)) if mm else 0
    except Exception:
        return 0


@app.get("/api/n2x/daemon.tcl")
def n2x_daemon_file():
    """
    지금 서버가 갖고 있는 데몬 스크립트를 그대로 내려준다.

    이 파일은 N2X 기계(윈도우)의 사본이 도는데, 그것을 어디서 받아야 하는지가
    어디에도 없었다. 저장소를 뒤지거나 사람에게 물어야 했다. 서버가 제
    사본을 내주면 그 기계에서 브라우저로 열어 받으면 끝난다 — 판이 어긋날
    자리도 그만큼 줄어든다.
    """
    if not os.path.exists(N2X_DAEMON):
        raise HTTPException(404, f"데몬 스크립트가 없습니다 — {N2X_DAEMON}")
    return FileResponse(
        N2X_DAEMON,
        media_type="text/plain; charset=utf-8",
        filename="n2x_daemon.tcl",
        headers={"Cache-Control": "no-store"},
    )


@app.get("/api/n2x/relay.py")
def n2x_relay_file():
    """중계 스크립트도 같은 자리에서. 처음 깔 때 이것부터 필요하다."""
    p = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "tools", "n2x_relay.py")
    if not os.path.exists(p):
        raise HTTPException(404, f"중계 스크립트가 없습니다 — {p}")
    return FileResponse(
        p,
        media_type="text/plain; charset=utf-8",
        filename="n2x_relay.py",
        headers={"Cache-Control": "no-store"},
    )


@app.get("/api/n2x/ver")
def n2x_ver(server: str = "210.1.2.248", label: str = "utop"):
    """
    윈도우에서 도는 데몬이 몇 번째 판인가.

    이 스크립트는 리눅스가 아니라 N2X 기계의 사본이 돈다. 저장소만 고치고
    컨테이너를 다시 올려도 실제로 도는 것은 안 바뀐다 — 그런데 화면에는
    그 사실이 어디에도 안 나와서, 고친 것이 왜 안 먹는지 알 수가 없었다.
    """
    want = _n2x_local_ver()
    res = _n2x_send(server, label, "ver")
    got = 0
    if isinstance(res, dict):
        try:
            got = int(res.get("ver") or 0)
        except Exception:
            got = 0

    # 「닿았는데 옛 판」 과 「아예 안 닿는다」 는 다른 일이다.
    #
    # 전에는 둘을 하나로 묶어 「옛 판입니다 (거기 알 수 없음)」 이라고 했다.
    # 주소를 잘못 적어 중계가 없는 자리를 고른 사람에게 이 말은 거짓이고,
    # 시키는 대로 파일을 복사해도 아무것도 안 바뀐다 — 고칠 곳은 주소다.
    reachable = got > 0
    stale = reachable and want > got
    if not reachable:
        note = (
            f"N2X 중계에 닿지 않습니다 ({server}). 주소가 맞는지, 그 기계에서 "
            "n2x_relay.py 가 떠 있는지 보세요."
        )
    elif stale:
        note = (
            f"N2X 기계의 n2x_daemon.tcl 이 옛 판입니다 (거기 {got} · 여기 {want}). "
            "backend/n2x/n2x_daemon.tcl 을 그 기계의 중계 폴더로 복사하고 n2x_relay.py 를 다시 띄우세요."
        )
    else:
        note = ""
    return {
        "ok": True,
        "local": want,
        "remote": got,
        "reachable": reachable,
        "stale": stale,
        "note": note,
    }


@app.post("/api/n2x/reset")
def n2x_reset(data: dict):
    """
    데몬 강제 재기동 — 섀시가 붙잡고 있는 세션을 놓게 한다.

    N2X 섀시는 동시에 열 수 있는 세션 수가 정해져 있다. 그것이 차면
    「The system already has maximum sessions running」 로 트래픽 시작이
    막힌다. 세션은 우리가 띄운 Tcl 데몬이 하나씩 쥐고 있으므로, 데몬을
    죽이는 것이 곧 세션을 놓는 것이다.

    **라벨을 안 주면 그 섀시로 띄운 데몬을 전부** 정리한다. 라벨 하나만
    죽이게 두었더니, 화면마다 다른 라벨로 띄운 것들이 남아 아무리 눌러도
    자리가 안 났다.
    """
    server = str(data.get("server", "210.1.2.248"))
    label = str(data.get("label") or "").strip()
    killed = []
    with _n2x_reg_lock:
        if label:
            keys = [k for k in list(_n2x_daemons) if k == server + "|" + label]
        else:
            keys = [k for k in list(_n2x_daemons) if k.startswith(server + "|")]
        for k in keys:
            d = _n2x_daemons.pop(k, None)
            if d and d.get("proc"):
                try:
                    d["proc"].kill()
                    killed.append(k.split("|", 1)[1])
                except Exception:
                    pass
    for lb in (killed or [label or "utop"]):
        _n2x_ports_cache_invalidate(server, lb)
    return {"ok": True, "killed": killed, "count": len(killed),
            "note": ("정리한 세션 " + ", ".join(killed)) if killed
                    else "우리가 띄운 세션은 없었습니다 — 남은 세션은 다른 PC 나 N2X GUI 가 쥐고 있습니다"}

# ── ports 응답 캐시 (server|label → {ts, data}) ─────────────────────────
# N2X 서버는 매 조회마다 모든 모듈·포트를 순차 스캔해서 부하가 크고 느림.
# 짧은 창(N2X_PORTS_CACHE_TTL초) 내 반복 조회는 캐시로 응답 → 사용자 여럿이 페이지를 열어도
# 실제 N2X 호출은 창당 1회. 예약/해제 성공 시 즉시 무효화(_n2x_ports_cache_invalidate).
_n2x_ports_cache = {}
_n2x_ports_cache_lock = threading.Lock()
N2X_PORTS_CACHE_TTL = 5.0   # 초

def _n2x_ports_cache_key(server, label):
    return str(server) + "|" + str(label)

def _n2x_ports_cache_invalidate(server, label):
    with _n2x_ports_cache_lock:
        _n2x_ports_cache.pop(_n2x_ports_cache_key(server, label), None)

def _n2x_ports_cached(server, label, force=False):
    """캐시 우선 조회. force=True 면 캐시 무시하고 새로 조회.
    데몬이 busy(lock timeout) 이면 만료된 캐시라도 반환 → 예약 진행 중에도 조회 응답 유지 (stale-while-busy 전략)."""
    import time as _t
    key = _n2x_ports_cache_key(server, label)
    if not force:
        with _n2x_ports_cache_lock:
            hit = _n2x_ports_cache.get(key)
            if hit and (_t.time() - hit["ts"]) < N2X_PORTS_CACHE_TTL:
                d = dict(hit["data"]); d["cached"] = True; d["cache_age"] = round(_t.time() - hit["ts"], 2)
                return d
    data = _n2x_send(server, label, "ports")
    if isinstance(data, dict) and data.get("ok"):
        with _n2x_ports_cache_lock:
            _n2x_ports_cache[key] = {"ts": _t.time(), "data": data}
    elif isinstance(data, dict) and not data.get("ok"):
        # 데몬 busy / lock timeout 등으로 실패 → 만료된 캐시라도 있으면 그걸로 대체 (사용자가 계속 이전 상태는 보게)
        with _n2x_ports_cache_lock:
            hit = _n2x_ports_cache.get(key)
        if hit:
            age = round(_t.time() - hit["ts"], 1)
            d = dict(hit["data"])
            d["cached"] = True; d["stale"] = True; d["cache_age"] = age
            d["stale_reason"] = data.get("error", "데몬 사용 중")
            return d
    return data

@app.get("/api/n2x/probe")
def n2x_probe(server: str = "210.1.2.248", label: str = "utop", force: int = 0):
    return _n2x_ports_cached(server, label, force=bool(force))

@app.get("/api/n2x/ports")
def n2x_ports(server: str = "210.1.2.248", label: str = "utop", force: int = 0):
    return _n2x_ports_cached(server, label, force=bool(force))

@app.post("/api/n2x/reserve-batch")
def n2x_reserve_batch(data: dict):
    """여러 (module, port) 예약을 한 번의 요청으로. 데몬 파이프는 단일이라 서버 단에서 순차 처리 →
    프론트가 병렬로 개별 요청 보낼 때 발생하던 파이프 race/timeout 문제 회피.
    사전 헬스체크(ping) + 포트 상태 조회(ports)로 이미 다른 세션이 잡은 포트는 시도 없이 명확한 에러 반환.
    payload: {server, label, targets: [{module, port}, ...]} → {ok:true, results:[{module,port,ok,error?}]}"""
    server = str(data.get("server", "210.1.2.248"))
    label = str(data.get("label", "utop"))
    targets = data.get("targets") or []
    if not targets:
        return {"ok": False, "error": "targets 필요"}
    # 데몬 사전 헬스체크: ping 실패 시 좀비 데몬 강제 정리 후 재기동 유도
    _hc = _n2x_send(server, label, "ping")
    if not (isinstance(_hc, dict) and _hc.get("ok")):
        with _n2x_reg_lock:
            _z = _n2x_daemons.pop(server + "|" + label, None)
        if _z:
            try: _z["proc"].kill()
            except Exception: pass
        _hc2 = _n2x_send(server, label, "ping")
        if not (isinstance(_hc2, dict) and _hc2.get("ok")):
            return {"ok": False, "error": "N2X 데몬 응답 없음 — 데몬 재기동 실패 · 서버 관리자 확인 필요"}
    # 포트 상태 조회 — reserve 전에 다른 세션이 잡고 있는지 미리 확인 (잡혀있으면 hang 방지 위해 시도 스킵)
    port_state = {}   # "module/port" -> {"lock": "sessionId", "label": "누구", "mine": bool}
    try:
        pj = _n2x_ports_cached(server, label)
        if isinstance(pj, dict) and pj.get("ok"):
            for m in (pj.get("modules") or []):
                mid = str(m.get("id"))
                for p in (m.get("portList") or []):
                    port_state[mid + "/" + str(p.get("port"))] = {
                        "lock": str(p.get("lock", "0")),
                        "label": p.get("label") or "",
                        "mine": bool(p.get("mine")),
                    }
    except Exception:
        pass
    results = []
    for t in targets:
        m = ""; p = ""
        try:
            m = str((t or {}).get("module", "")).strip()
            p = str((t or {}).get("port", "")).strip()
            if not m or not p:
                results.append({"module": m, "port": p, "ok": False, "error": "module/port 누락"})
                continue
            key = m + "/" + p
            _st = port_state.get(key)
            # 다른 세션이 잡고 있으면 시도 스킵 (hang 방지) — 강제 예약을 원하면 force 플래그 사용해야 함
            if _st and not _st["mine"] and _st["lock"] != "0":
                _who = _st["label"] or ("세션 " + _st["lock"])
                results.append({"module": m, "port": p, "ok": False,
                                "error": "이미 다른 세션이 사용 중 (label: " + _who + ")",
                                "locked_by": _who})
                continue
            # 내가 이미 잡은 포트면 성공으로 (재예약 불필요)
            if _st and _st["mine"]:
                results.append({"module": m, "port": p, "ok": True, "already_mine": True})
                continue
            r = _n2x_send(server, label, "reserve " + m + " " + p)
            if isinstance(r, dict) and r.get("ok"):
                results.append({"module": m, "port": p, "ok": True})
            else:
                results.append({"module": m, "port": p, "ok": False, "error": (r or {}).get("error", "실패")})
        except Exception as e:
            results.append({"module": m, "port": p, "ok": False, "error": str(e)})
    _n2x_ports_cache_invalidate(server, label)
    return {"ok": True, "results": results}

def _n2x_verify_reserved(server, label, module, port):
    """포트가 실제로 label 세션에 예약됐는지 서버에 물어봐 확인. 상태 변경 후이므로 캐시 무시."""
    pj = _n2x_ports_cached(server, label, force=True)
    if not isinstance(pj, dict) or not pj.get("ok"):
        return None   # 확인 불가
    for m in (pj.get("modules") or []):
        if str(m.get("id")) == str(module):
            for p in (m.get("portList") or []):
                if str(p.get("port")) == str(port):
                    return bool(p.get("mine"))
    return False

@app.post("/api/n2x/reserve")
def n2x_reserve(data: dict):
    module = str(data.get("module", ""))
    ports = data.get("ports", [])
    if not module or not ports:
        return {"ok": False, "error": "module/ports 누락"}
    server = str(data.get("server", "210.1.2.248"))
    label = str(data.get("label", "utop"))
    port = str(ports[0])
    if data.get("force"):
        # 강제: 다른 세션이 이 포트를 잠갔으면 그 세션(label)에서 해당 포트만 release 후 점유
        pj = _n2x_ports_cached(server, label)
        lk = None
        if isinstance(pj, dict) and pj.get("ok"):
            for m in (pj.get("modules") or []):
                if str(m.get("id")) == module:
                    for p in (m.get("portList") or []):
                        if str(p.get("port")) == port and not p.get("mine") and str(p.get("lock", "0")) != "0":
                            lk = p.get("label")
        if lk and str(lk) != label:
            _n2x_send(server, str(lk), "release " + module + " " + port)
            _n2x_ports_cache_invalidate(server, str(lk))
            _n2x_ports_cache_invalidate(server, label)
            import time
            time.sleep(0.5)
    res = _n2x_send(server, label, "reserve " + module + " " + port)
    _n2x_ports_cache_invalidate(server, label)
    # reserve 실패 응답이 왔더라도 서버 상태 재확인 — 파이프 끊김 사이 이미 예약됐을 수 있음
    if isinstance(res, dict) and not res.get("ok"):
        try:
            import time
            time.sleep(0.3)   # N2X 서버 상태 반영 대기
            ok_actual = _n2x_verify_reserved(server, label, module, port)
            if ok_actual is True:
                _n2x_log(f"[N2X] reserve reported fail but port actually reserved -- recovering module={module} port={port}")
                return {"ok": True, "reserved": module + "/" + port, "recovered": True, "note": res.get("error", "")}
        except Exception as e:
            _n2x_log(f"[N2X] verify after reserve fail failed: {e}")
    return res

@app.post("/api/n2x/release")
def n2x_release(data: dict):
    module = str(data.get("module", ""))
    port = str(data.get("port", ""))
    if not module or not port:
        return {"ok": False, "error": "module/port 필요"}
    server = str(data.get("server", "210.1.2.248"))
    label = str(data.get("label", "utop"))
    res = _n2x_send(server, label, "release " + module + " " + port)
    _n2x_ports_cache_invalidate(server, label)
    # release 실패 응답이 왔더라도 실제로 해제됐는지 재확인
    if isinstance(res, dict) and not res.get("ok"):
        try:
            import time
            time.sleep(0.3)
            ok_actual = _n2x_verify_reserved(server, label, module, port)
            if ok_actual is False:   # 내 예약에 없음 = 해제된 것
                _n2x_log(f"[N2X] release reported fail but port actually released -- recovering module={module} port={port}")
                return {"ok": True, "released": module + "/" + port, "recovered": True, "note": res.get("error", "")}
        except Exception as e:
            _n2x_log(f"[N2X] verify after release fail failed: {e}")
    return res

def _n2x_streams_from(data: dict):
    streams = data.get("streams") or []
    # 단일(구버전 폼) 호환: module/txPort/rxPort 가 오면 1개 스트림으로 변환
    if not streams and data.get("module") and data.get("txPort") and data.get("rxPort"):
        streams = [{"txMod": data.get("module"), "txPort": data.get("txPort"),
                    "rxMod": data.get("module"), "rxPort": data.get("rxPort"),
                    "proto": "udp", "frame": data.get("frame", 64),
                    "pps": data.get("pps", 1000), "npkt": data.get("npkt", 0)}]
    return streams


# 화면의 말 → 데몬의 말.
#
# 데몬은 윈도우에서 도는 Tcl 이라 한글이 그대로 가면 깨진다(전에 오류
# 메시지가 그렇게 깨져 읽을 수가 없었다). 여기서 ASCII 로 바꿔 보낸다.
_N2X_MODS = {"증가": "inc", "감소": "dec", "무작위": "rand", "Yes": "inc", "No": "fix"}


def _n2x_specs(streams):
    def _clean(v):
        t = str(v if v is not None else "").strip()
        t = _N2X_MODS.get(t, t)
        return t.replace(",", "").replace(" ", "")
    # unit 은 맨 뒤에 붙인다 — 자리로 읽는 형식이라, 중간에 끼우면 옛 spec 이
    # 통째로 어긋난다. 없으면 데몬이 pps 로 본다.
    #
    # 뒤의 일곱은 **주소를 여럿으로 뿌리기** 위한 것이다.
    #
    # 여태 계측기에는 값이 하나씩만 갔다(`SetFieldFixedValue`). 화면에서는
    # 「01 부터 열 개」 로 적어 두고 선로에는 01 하나만 나갔는데, 화면
    # 어디에도 그 말이 없었다 — 시험은 돌고 결과도 나오는데 잰 것이 딴것이다.
    # 시작 · 개수 · 모드를 함께 보내고, 데몬이 목록을 만들어
    # `SetFieldValueList` 로 넣는다.
    keys = [
        "txMod", "txPort", "rxMod", "rxPort", "proto", "frame", "pps", "npkt",
        "srcMac", "dstMac", "srcIp", "dstIp", "unit", "frameMax",
        "cnt", "srcMacMod", "dstMacMod", "srcIpMod", "dstIpMod", "vlan", "vlanMod",
    ]
    specs = []
    for s in streams:
        # 송신 모듈만 주면 수신 모듈도 동일하게
        if not s.get("rxMod"):
            s["rxMod"] = s.get("txMod") or s.get("module") or ""
        if not s.get("txMod"):
            s["txMod"] = s.get("module") or ""
        specs.append(",".join(_clean(s.get(k, "")) for k in keys))
    return specs




@app.post("/api/n2x/traffic/start")
def n2x_traffic_start(data: dict):
    # 비동기 시작 — 즉시 리턴(대기 X). dur 0/미지정 = 연속(데몬에서 1시간), 이후 /stat 폴링
    streams = _n2x_streams_from(data)
    if not streams:
        return {"ok": False, "error": "streams(또는 module/txPort/rxPort) 필요"}
    # 보내는 줄을 그대로 남긴다.
    #
    # 「10 갈래로 잡았는데 두 줄만 나온다」 를 쫓는 데 한참 걸렸다. 화면 ·
    # 서버 · 데몬 셋 중 어디서 값이 빠지는지 볼 데가 없었기 때문이다.
    # 이 한 줄이면 무엇이 실제로 나갔는지 바로 보인다.
    cmd = "tstart " + str(data.get("dur") or 0) + " " + " ".join(_n2x_specs(streams))
    _n2x_log("[N2X] " + cmd)
    return _n2x_humanize(_n2x_send(
        str(data.get("server", "210.1.2.248")), str(data.get("label", "utop")), cmd))


@app.post("/api/n2x/traffic/stat")
def n2x_traffic_stat(data: dict):
    # 실시간 통계 폴링 (전송 중에도 조회)
    return _n2x_send(str(data.get("server", "210.1.2.248")), str(data.get("label", "utop")), "tstat")


@app.post("/api/n2x/traffic/stop")
def n2x_traffic_stop(data: dict):
    return _n2x_send(str(data.get("server", "210.1.2.248")), str(data.get("label", "utop")), "tstop")


@app.post("/api/n2x/arp")
def n2x_arp(data: dict):
    """
    GW 에게 ARP 를 보내 MAC 을 받아 온다.

    L3 로 쏘려면 프레임의 목적지 MAC 이 첫 홉(=GW)의 MAC 이어야 한다.
    지금까지 그 값은 사람이 장비에서 `show arp` 로 읽어 손으로 옮겨 적었다.
    한 자만 틀려도 프레임이 장비로 안 가고 손실 100% 로 나오는데, 화면에는
    「안 받았다」 만 뜬다.

    데몬이 그 명령을 아직 모르면 `unknown` 이 돌아온다. 그때는 **거짓으로
    성공을 만들지 않는다** — 무엇이 없어서 못 하는지 그대로 말한다.
    """
    server = str(data.get("server") or "210.1.2.248")
    label = str(data.get("label") or "utop")
    port = str(data.get("port") or "").strip()
    gw = str(data.get("gw") or "").strip()
    if not gw:
        return {"ok": False, "error": "GW 가 비어 있습니다"}
    if not port:
        return {"ok": False, "error": "이 스트림의 보내는 포트가 비어 있습니다"}
    mod, _, pnum = port.partition("/")
    res = _n2x_send(
        server,
        label,
        "arp %s %s %s %s %s"
        % (
            mod or "-",
            pnum or "-",
            gw,
            str(data.get("srcIp") or "-").strip() or "-",
            str(data.get("srcMac") or "-").strip() or "-",
        ),
    )
    if isinstance(res, dict) and str(res.get("error") or "") == "unknown":
        return {
            "ok": False,
            "error": (
                "이 N2X 데몬은 아직 ARP 를 모릅니다. n2x_daemon.tcl 을 새 판으로 "
                "바꾸고 n2x_relay.py 를 다시 띄우세요."
            ),
        }
    return res


@app.post("/api/n2x/traffic/clear")
def n2x_traffic_clear(data: dict):
    return _n2x_send(str(data.get("server", "210.1.2.248")), str(data.get("label", "utop")), "tclear")

@app.post("/api/stc/meter/{action}")
async def stc_meter_action(action: str, data: dict):
    # TC 계측기 스텝 실행: stc_meter.py 로 한 액션 수행(영속 U_TOP_meter 세션). stdout 텍스트 반환.
    import json as _json
    cfg = (data or {}).get("cfg") or {}
    if action not in ("build", "arp", "start", "stop", "query", "close", "disconnect"):
        return {"ok": False, "error": "알 수 없는 action: " + str(action)}
    _rip, _rport = await _stc_rest_for(str(cfg.get("chassis") or ""), cfg.get("restIp"), cfg.get("restPort"))
    cfg["restIp"], cfg["restPort"] = _rip, _rport
    rest_ip = str(_rip).strip().lower()
    rest_port = int(_rport)
    if rest_ip in ("localhost", "127.0.0.1", "") and not _port_listening(rest_port):
        srv = await stc_server_start({"port": rest_port})
        if not srv.get("ok"):
            return {"ok": False, "error": "REST 서버 시작 실패: " + str(srv.get("error"))}
    cfgdir = str(Path(__file__).parent)
    cfgpath = os.path.join(cfgdir, "stc", "_stc_meter_cfg.json")
    try:
        with open(cfgpath, "w") as f:
            _json.dump(cfg, f)
    except Exception as e:
        return {"ok": False, "error": "cfg 쓰기 실패: " + repr(e)}
    script = os.path.join(cfgdir, "stc", "stc_meter.py")

    def _run():
        cmd = [sys.executable, "-u", script, action, cfgpath]
        env = dict(os.environ)
        env["PYTHONUNBUFFERED"] = "1"
        try:
            cp = subprocess.run(cmd, cwd=cfgdir, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, env=env, timeout=150)
            return (cp.stdout.decode("utf-8", "replace") if cp.stdout else ""), cp.returncode
        except subprocess.TimeoutExpired:
            return "[시간 초과 150초]", 1
        except Exception as e:
            return "[실행 실패] " + repr(e), 1

    text, rc = await asyncio.to_thread(_run)
    ok = (rc == 0) and ("[ERROR]" not in text)
    return {"ok": ok, "text": (text or "").strip(), "code": rc}

@app.post("/api/stc/server/start")
async def _stc_rest_for(chassis: str, rest_ip: str = "", rest_port=None):
    """이 섀시의 **REST 서버 주소**. 화면이 안 알려 주면 등록에서 찾는다.

    STC 는 두 자리가 있다: 섀시(장비 IP)와 REST 서버(윈도우 PC). 화면 몇
    군데가 REST 서버를 `localhost` 로 박아 두었는데, 이 서버는 리눅스라
    거기엔 아무도 없다 — 그래서 「stcweb.exe 를 찾을 수 없습니다」 로
    끝났다(지적: 계측기는 붙는데 시험 탭에서만).

    계측기 등록의 stc 접속 줄에 그 주소가 이미 있다(host·port). 그것이
    정본이다. 못 찾으면 받은 값을 그대로 쓴다 — 지어내지 않는다.
    """
    ip = str(rest_ip or "").strip()
    port = int(rest_port or 0) or 0
    local = ip.lower() in ("", "localhost", "127.0.0.1")
    if not local and port:
        return ip, port
    try:
        for d in await db.device_list(with_ifs=False):
            if str(d.get("ip") or "").strip() != str(chassis or "").strip():
                continue
            for a in (d.get("access") or []):
                if str(a.get("protocol") or "").lower() != "stc":
                    continue
                h = str(a.get("host") or "").strip()
                p = int(a.get("port") or 0) or 0
                if local and h:
                    ip = h
                if not port and p:
                    port = p
                break
    except Exception as e:
        print(f"[stc] 등록에서 REST 주소를 못 읽었습니다: {e}", flush=True)
    return (ip or "localhost"), (port or 8888)


async def stc_server_start(data: dict = None):
    data = data or {}
    port = int(data.get("port") or 8888)
    if _port_listening(port):
        return {"ok": True, "already": True, "listening": True}
    exe = _find_stcweb()
    if not exe:
        # 이 서버가 리눅스면 stcweb 은 여기서 뜰 수 있는 물건이 아니다.
        # 「환경변수를 지정하라」 만 적어 두면 없는 파일을 찾아 헤매게 된다
        # (지적: 계측기는 붙는데 시험 탭에서만 이 말이 뜬다).
        if os.name != "nt":
            return {"ok": False, "error":
                    f"STC REST 서버(포트 {port})에 닿지 못했습니다 — 이 서버는 리눅스라 "
                    "stcweb 을 여기서 띄울 수 없습니다. STC PC 에서 REST 서버를 켜고, "
                    "계측기 등록의 포트가 그 서버 포트와 같은지 확인하세요."}
        return {"ok": False, "error": "stcweb.exe 를 찾을 수 없습니다. STCWEB_PATH 환경변수로 경로를 지정하세요."}
    workdir = os.path.dirname(exe)
    # Popen 으로 detached 기동 (asyncio subprocess 는 Windows uvicorn 루프에서 미지원)
    try:
        proc = subprocess.Popen([exe], cwd=workdir,
                                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except Exception as e:
        return {"ok": False, "error": "REST 서버 실행 실패: " + str(e)}
    _stcweb_proc["p"] = proc
    await broadcast({"type": "stc_line", "line": "[REST] stcweb.exe 시작 — 포트 " + str(port) + " 대기..."})
    # 포트가 열릴 때까지 최대 ~25초 대기
    for _ in range(50):
        await asyncio.sleep(0.5)
        if proc.poll() is not None:
            return {"ok": False, "error": "REST 서버가 즉시 종료됨 (exit " + str(proc.poll()) + ")"}
        if _port_listening(port):
            await broadcast({"type": "stc_line", "line": "[REST] 서버 준비됨 (localhost:" + str(port) + ")"})
            return {"ok": True, "listening": True, "exe": exe}
    return {"ok": False, "error": "REST 서버가 시간 내 포트 " + str(port) + " 를 열지 못함"}


# STC 인벤토리 캐시 — 매번 ChassisConnect 하면 수십 초가 걸린다.
# 담아 두고, 그 안의 재조회는 즉시 돌려준다. 「새로고침」(force) 이면 무시한다.
#
# 60초는 짧았다. 이 값이 말하는 것은 **어떤 슬롯에 몇 포트가 꽂혀 있나**
# 이고, 그것은 사람이 카드를 뽑았다 꽂을 때나 바뀐다 — 시험 하나 만드는
# 동안 바뀔 일이 없다(지적: 느리다). 예약 상태는 이 캐시로 보지 않는다.
_STC_CC_CACHE = {}   # "chassis|rest_ip:rest_port" -> {"ts": t, "data": {...}}
_STC_CC_TTL = 600


@app.post("/api/stc/conncheck")
async def stc_conncheck(data: dict = None):
    """실제 섀시 연결 확인 (트래픽 생성 없음). 섀시/모듈 인벤토리를 반환.

    매번 새 세션을 열고 섀시에 접속하므로 수십 초 걸린다. 그래서 결과를
    잠깐 캐시한다 — 포트 현황을 다시 열거나 다른 사람이 같은 섀시를 봐도
    바로 뜬다. force=1 이면 캐시를 건너뛴다.
    """
    data = data or {}
    _ck = str(data.get("chassis") or "") + "|" + str(data.get("restIp") or "") + ":" + str(data.get("restPort") or "")
    if not data.get("force"):
        hit = _STC_CC_CACHE.get(_ck)
        if hit and (_t.time() - hit["ts"]) < _STC_CC_TTL:
            d = dict(hit["data"]); d["cached"] = True
            d["cache_age"] = round(_t.time() - hit["ts"])
            return d
    chassis = (data.get("chassis") or "192.168.5.100").strip()
    # 화면이 안 알려 줬으면 계측기 등록에서 찾는다
    rest_ip, rest_port = await _stc_rest_for(chassis, data.get("restIp"), data.get("restPort"))
    # 로컬 REST 서버 자동 기동
    if rest_ip.lower() in ("localhost", "127.0.0.1", "") and not _port_listening(rest_port):
        srv = await stc_server_start({"port": rest_port})
        if not srv.get("ok"):
            return {"ok": False, "error": "REST 서버 시작 실패: " + str(srv.get("error"))}
    helper = str(Path(__file__).parent / "stc" / "stc_conncheck.py")
    # 백엔드와 동일한 python(sys.executable)을 직접 호출한다.
    # 'py -3.12' 런처를 쓰면 손자 python 프로세스가 stdout 파이프를 잡아 timeout 이 안 풀린다.
    # 또한 이 인터프리터에는 stcrestclient 가 설치돼 있다(백엔드가 그 위에서 동작).
    cmd = [sys.executable, helper, chassis, rest_ip, str(rest_port)]
    # 동기 subprocess 를 스레드에서 실행 (Windows uvicorn 루프 호환)
    def _run():
        return subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, timeout=90)
    try:
        cp = await asyncio.to_thread(_run)
        raw = cp.stdout
    except subprocess.TimeoutExpired:
        return {"ok": False, "error": "연결 확인 시간 초과 (90초)"}
    except Exception as e:
        return {"ok": False, "error": "실행 실패: " + repr(e)}
    text = raw.decode("utf-8", "replace") if raw else ""
    # 마지막 JSON 라인 파싱
    result = None
    for line in reversed(text.strip().splitlines()):
        line = line.strip()
        if line.startswith("{") and line.endswith("}"):
            try:
                result = json.loads(line)
                break
            except Exception:
                continue
    if result is None:
        return {"ok": False, "error": "결과 파싱 실패", "raw": text[-500:]}
    if isinstance(result, dict) and result.get("ok"):
        _STC_CC_CACHE[_ck] = {"ts": _t.time(), "data": result}
    return result

# ───────────────────────────────────────────
# 라우터 - STC 실제 포트 예약/해제 (영속 세션 U_TOP_reserve)
# 예약은 세션에 묶이므로 reserve 후 세션을 유지 → Spirent STC 프로그램에 'Reserved by utop' 로 보임.
# ───────────────────────────────────────────
async def _run_stc_reserve(action, ports, chassis, rest_ip, rest_port):
    # 로컬 REST 서버 자동 기동
    if str(rest_ip).lower() in ("localhost", "127.0.0.1", "") and not _port_listening(rest_port):
        srv = await stc_server_start({"port": rest_port})
        if not srv.get("ok"):
            return {"ok": False, "error": "REST 서버 시작 실패: " + str(srv.get("error"))}
    helper = str(Path(__file__).parent / "stc" / "stc_reserve.py")
    cmd = [sys.executable, helper, action, (ports or "-"), chassis, rest_ip, str(rest_port)]
    def _run():
        return subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, timeout=90)
    try:
        cp = await asyncio.to_thread(_run)
        raw = cp.stdout
    except subprocess.TimeoutExpired:
        return {"ok": False, "error": "시간 초과 (90초)"}
    except Exception as e:
        return {"ok": False, "error": "실행 실패: " + repr(e)}
    text = raw.decode("utf-8", "replace") if raw else ""
    for line in reversed(text.strip().splitlines()):
        line = line.strip()
        if line.startswith("{") and line.endswith("}"):
            try:
                return json.loads(line)
            except Exception:
                continue
    return {"ok": False, "error": "결과 파싱 실패", "raw": text[-500:]}

@app.post("/api/stc/reserve")
async def stc_reserve(data: dict = None):
    data = data or {}
    chassis = (data.get("chassis") or "192.168.5.100").strip()
    rest_ip, rest_port = await _stc_rest_for(chassis, data.get("restIp"), data.get("restPort"))
    ports = (data.get("ports") or "").strip()  # "1/15,1/16" — 예약할 전체 집합
    return await _run_stc_reserve("reserve", ports, chassis, rest_ip, rest_port)

@app.post("/api/stc/reserve/status")
async def stc_reserve_status(data: dict = None):
    data = data or {}
    chassis = (data.get("chassis") or "192.168.5.100").strip()
    rest_ip, rest_port = await _stc_rest_for(chassis, data.get("restIp"), data.get("restPort"))
    return await _run_stc_reserve("status", "-", chassis, rest_ip, rest_port)

@app.post("/api/stc/release")
async def stc_release(data: dict = None):
    data = data or {}
    chassis = (data.get("chassis") or "192.168.5.100").strip()
    rest_ip, rest_port = await _stc_rest_for(chassis, data.get("restIp"), data.get("restPort"))
    return await _run_stc_reserve("release", "-", chassis, rest_ip, rest_port)

# ───────────────────────────────────────────
# 라우터 - STC 트래픽 시험 단계 (영속 세션 U_TOP_traffic)
# connect/reserve/devices/streams/start/stop/counters/status/end 를 한 세션에서 순차 수행
# ───────────────────────────────────────────
_STC_SESS_ACTIONS = {"connect", "reserve", "releaseports", "forcereset", "devices", "streams", "start", "stop", "counters", "status", "portstatus", "end"}

# ── STC: 영속 in-process 세션 + 단일 락 + 상태 캐시 ──────────
#  핵심: STC 예약은 '연속 연결 + 살아있는 포트 핸들'을 전제로 한다. 액션마다 subprocess 를
#  새로 띄우면 핸들이 끊겨 해제가 불가능 → 하나의 영속 세션(StcLive)을 백엔드 안에 유지한다.
#   · connect/reserve/releaseports/forcereset/portstatus → StcLive(공유 작업 세션 U_TOP_work)
#   · 단일 락(_stc_live_lock)으로 직렬화(계정이 몇이든 안전). '누가 예약'은 레지스트리가 추적.
#   · 상태는 백엔드 단일 폴러가 캐시 → 모든 브라우저는 캐시 + 계정별 mine 만 덧칠(계측기 부하 일정).
#   · devices/streams/start/stop/counters/end → 아직 subprocess(U_TOP_traffic, 트래픽 단계에서 이관).
from stc.stc_live import StcLive
_stc_live = StcLive()
_stc_live_lock = asyncio.Lock()       # 예약/해제/연결(U_TOP_work) 직렬화
_stc_status_lock = asyncio.Lock()     # 상태 서브프로세스(U_TOP_status) 중복 방지
_stc_status_cache = {}         # chassis -> {"ts": float, "rows": [...]}
_stc_status_targets = {}       # chassis -> (rest_ip, rest_port)
# 마지막으로 「포트 상태 좀」 하고 물어본 시각. chassis -> epoch
#
# 여태 한 번 물어본 섀시는 **영영** 목록에 남았고, 폴러가 3초마다 TCL
# 서브프로세스를 띄웠다. 아무도 안 보고 있어도 그랬다 — 253 의 CPU 가
# 종일 붙어 있던 것이 이것이다(지적). 보는 사람이 없으면 멈춘다.
_stc_status_seen = {}
_STC_WATCH_TTL = 90            # 이 초 동안 아무도 안 물어보면 폴러에서 뺀다
_stc_poller_started = False
_STC_LIVE_ACTIONS = {"connect", "reserve", "releaseports", "forcereset"}

def _stc_err(msg: str) -> str:
    """STC 가 준 실패를 사람 말로. 원문은 뒤에 남긴다."""
    t = str(msg or "")
    low = t.lower()
    if "timed out waiting for session" in low:
        return ("STC REST 서버가 세션을 못 띄웠습니다 — 계측기 안의 BLL 이 뜨는 데 "
                "오래 걸리는 중입니다. 20초쯤 뒤 다시 누르거나, 그래도 같으면 "
                "REST 서버를 다시 시작하세요. (원문: " + t[:120] + ")")
    if "connection refused" in low or "failed to establish" in low:
        return ("STC REST 서버에 닿지 못했습니다 — 서버가 떠 있는지(기본 8888) 확인하세요. "
                "(원문: " + t[:120] + ")")
    if "session not found" in low:
        return ("STC 세션이 끊겼습니다 — 다시 누르면 새로 붙습니다. (원문: " + t[:120] + ")")
    return t


def _run_stc_helper(action, chassis, rest_ip, rest_port, params):
    """stc_session.py 서브프로세스(트래픽 단계용). JSON dict 반환(예외도 dict)."""
    helper = str(Path(__file__).parent / "stc" / "stc_session.py")
    cmd = [sys.executable, helper, action, chassis, rest_ip, str(rest_port), json.dumps(params)]
    try:
        cp = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, timeout=120)
        raw = cp.stdout
    except subprocess.TimeoutExpired:
        return {"ok": False, "error": "시간 초과 (120초)"}
    except Exception as e:
        return {"ok": False, "error": "실행 실패: " + repr(e)}
    text = raw.decode("utf-8", "replace") if raw else ""
    for line in reversed(text.strip().splitlines()):
        line = line.strip()
        if line.startswith("{") and line.endswith("}"):
            try:
                return json.loads(line)
            except Exception:
                continue
    return {"ok": False, "error": "결과 파싱 실패", "raw": text[-500:]}

async def _stc_live_status(chassis, rest_ip, rest_port):
    # 상태는 읽기전용 서브프로세스(U_TOP_status)로 — 영속 세션 핸들 손상 없이 안정적.
    #  예약 세션(StcLive/U_TOP_work)과 독립이라 락 불필요.
    # light: 링크·속도(포트마다 REST 왕복)는 건너뛴다. 목록과 예약 상태가
    # 이 화면이 쓰는 전부다 — 왕복 수가 절반 아래로 준다(지적: 느리다).
    res = await asyncio.to_thread(_run_stc_helper, "portstatus", chassis, rest_ip, rest_port,
                                  {"user": "_utop_status_", "light": True})
    if isinstance(res, dict) and res.get("ok"):
        _stc_status_cache[chassis] = {"ts": _t.time(), "rows": res.get("ports", [])}
    return res

async def _stc_poller_loop():
    while True:
        try:
            now = _t.time()
            for ch, (rip, rport) in list(_stc_status_targets.items()):
                # 보고 있는 사람이 없으면 뺀다. 서브프로세스 한 번이 가볍지
                # 않다 — 섀시 하나만 남아 있어도 코어 하나를 문다.
                if now - _stc_status_seen.get(ch, 0) > _STC_WATCH_TTL:
                    _stc_status_targets.pop(ch, None)
                    _stc_status_seen.pop(ch, None)
                    print(f"[stc] 보는 사람이 없어 상태 폴링을 멈춥니다 — {ch}", flush=True)
                    continue
                async with _stc_status_lock:   # sync-poll 과 같은 U_TOP_status 세션 → 직렬화 필수
                    await _stc_live_status(ch, rip, rport)
        except Exception:
            pass
        # 볼 것이 없으면 느리게 — 빈 채로 3초마다 깨울 이유가 없다
        await asyncio.sleep(3 if _stc_status_targets else 15)

def _ensure_poller():
    global _stc_poller_started
    if not _stc_poller_started:
        _stc_poller_started = True
        try:
            asyncio.create_task(_stc_poller_loop())
        except Exception:
            _stc_poller_started = False

def _overlay_mine(rows, user):
    """캐시된 raw 상태(예약은 전부 other)에 요청 계정의 '내 예약'을 덧칠."""
    out = []
    for r in rows:
        st = r.get("status"); who = r.get("who")
        nr = {"slot": r.get("slot"), "port": r.get("port"), "status": st}
        if r.get("link"):
            nr["link"] = r.get("link")
        if r.get("speed"):
            nr["speed"] = r.get("speed")
        if st in ("other", "mine"):
            if who and who == user:
                nr["status"] = "mine"; nr["who"] = user
            else:
                nr["status"] = "other"; nr["who"] = who or "외부"
        out.append(nr)
    return out

@app.post("/api/stc/sess/{action}")
async def stc_sess(action: str, data: dict = None):
    if action not in _STC_SESS_ACTIONS:
        return {"ok": False, "error": "알 수 없는 단계: " + action}
    data = data or {}
    chassis = (data.get("chassis") or "192.168.5.100").strip()
    rest_ip, rest_port = await _stc_rest_for(chassis, data.get("restIp"), data.get("restPort"))
    params = data.get("params") or {}
    user = str(params.get("user") or "admin")
    # 로컬 REST 서버 자동 기동
    if rest_ip.lower() in ("localhost", "127.0.0.1", "") and not _port_listening(rest_port):
        srv = await stc_server_start({"port": rest_port})
        if not srv.get("ok"):
            return {"ok": False, "error": "REST 서버 시작 실패: " + str(srv.get("error"))}

    # 상태조회: 단일 폴러 캐시 + 계정별 mine 덧칠
    if action == "portstatus":
        _stc_status_targets[chassis] = (rest_ip, rest_port)
        _stc_status_seen[chassis] = _t.time()      # 지금 보고 있다
        _ensure_poller()
        cache = _stc_status_cache.get(chassis)
        now = _t.time()
        # 있는 값을 **먼저 준다**. 섀시에 묻는 일은 포트마다 REST 왕복이라
        # 몇 초가 든다 — 그동안 화면이 멎어 있으면 「느리다」 가 된다(지적).
        # 묵은 값이면 뒤에서 새로 읽어 두고, 다음 번에 새 값이 나간다.
        if cache and (now - cache["ts"]) > 12:
            async def _bg():
                async with _stc_status_lock:
                    c2 = _stc_status_cache.get(chassis)
                    if not c2 or (_t.time() - c2["ts"]) > 12:
                        await _stc_live_status(chassis, rest_ip, rest_port)
            try:
                asyncio.create_task(_bg())
            except Exception as e:
                print(f"[stc] 뒷일로 못 넘겼습니다: {e}", flush=True)
        elif not cache:
            # 처음 한 번은 어쩔 수 없이 기다린다 — 줄 것이 없다
            async with _stc_status_lock:
                if not _stc_status_cache.get(chassis):
                    await _stc_live_status(chassis, rest_ip, rest_port)
            cache = _stc_status_cache.get(chassis)
        rows = cache["rows"] if cache else []
        return {"ok": True, "action": "portstatus", "user": user,
                "ports": _overlay_mine(rows, user),
                "cached": True, "ts": (cache["ts"] if cache else 0),
                # 몇 초 전 값인가 — 화면이 「지금 것」 인 척하지 않게
                "age": round(now - cache["ts"], 1) if cache else 0}

    # 그 외 모든 액션(예약/해제/강제리셋/연결/트래픽): 서브프로세스.
    #  예약/해제는 command-only(ReservePortCommand/RevokeOwner)라 포트 오브젝트를 안 만들어
    #  세션이 손상되지 않음(검증됨). 작업 세션(U_TOP_work) 동시접근 방지 위해 직렬화.
    async with _stc_live_lock:
        res = await asyncio.to_thread(_run_stc_helper, action, chassis, rest_ip, rest_port, params)
    if action in ("reserve", "releaseports", "forcereset", "connect"):
        _stc_status_cache.pop(chassis, None)   # 점유 변화 → 다음 조회에서 실상태 반영
    # 영어 한 덩어리를 그대로 던지지 않는다 — 무엇을 해야 하는지까지 적는다
    if isinstance(res, dict) and res.get("error"):
        res["error"] = _stc_err(res["error"])
    return res

@app.get("/api/cycle/{cycle_id}")
async def get_cycle(cycle_id: str):
    d = await db.cycle_get(cycle_id)
    if d is None:
        raise HTTPException(404, "Cycle을 찾을 수 없습니다")
    return d

@app.get("/api/pptx-templates")
def pptx_templates():
    """고를 수 있는 고객사 양식. 파일이 없는 것은 빼고 준다."""
    import pptx_tpl
    return {"templates": pptx_tpl.list_templates()}


@app.post("/api/pptx-render")
async def pptx_render(payload: dict):
    """
    고객사 양식에 값을 채워 결과서를 만든다.

    **내용은 화면이 조립해서 보낸다.** 미리보기가 쓰는 것과 같은 자료·같은
    쪽 나누기를 그대로 쓰기 위해서다 — 서버에서 따로 조립하면 두 벌이 되고,
    한쪽만 고치는 순간 화면에서 본 장수와 파일의 장수가 어긋난다.

    서버가 하는 일은 하나다: 고객사가 준 pptx 를 열어 **그 안의 장을 복제해
    값만 갈아 끼운다.** 글꼴·표선·색·머리글은 손대지 않으므로 받는 쪽 눈에는
    자기네 양식 그대로다.
    """
    import pptx_tpl
    from pptx import Presentation

    tid = str((payload or {}).get("template") or "lguplus")
    tpl = pptx_tpl.TEMPLATES.get(tid)
    if not tpl:
        raise HTTPException(400, f"모르는 양식입니다: {tid}")
    path = pptx_tpl.TPL_DIR / str(tpl["file"])
    if not path.exists():
        raise HTTPException(404, f"양식 파일이 없습니다: {path.name}")

    slides = (payload or {}).get("slides") or []
    if not slides:
        raise HTTPException(400, "채울 내용이 없습니다")

    prs = Presentation(str(path))
    src_first = prs.slides[int(tpl["first"])]
    src_more = prs.slides[int(tpl["more"])]
    n_tpl = len(list(prs.slides))

    for sl in slides:
        kind = str((sl or {}).get("kind") or "first")
        vals = {k: str(v if v is not None else "") for k, v in ((sl or {}).get("values") or {}).items()}
        if kind == "more":
            new = pptx_tpl.clone_slide(prs, src_more)
            pptx_tpl.fill(new, tpl["more_spots"], vals)
        else:
            new = pptx_tpl.clone_slide(prs, src_first)
            pptx_tpl.fill(new, tpl["spots"], vals)
        # 양식에 얹혀 있던 **예시 그림·글상자**를 걷는다. 안 걷으면 만든
        # 결과서 모든 쪽에 남의 시험 화면이 실려 진짜 결과를 덮는다.
        pptx_tpl.strip_samples(new)
        # 그 위에 이 시험의 그림을 얹는다 — 구성도, 그리고 CLI 캡쳐
        pics = tpl.get("more_pics" if kind == "more" else "pics") or {}
        for name, spec in pics.items():
            blob = pptx_tpl.decode_img(str((sl or {}).get(name) or ""))
            if not blob:
                continue
            try:
                rect = pptx_tpl.pic_rect(new, spec)
                if rect:
                    pptx_tpl.place_pic(new, blob, rect)
            except Exception:
                # 그림 하나 때문에 결과서 전체가 안 나오면 안 된다
                pass

    # 본보기 장은 결과에 남기지 않는다. 뒤에서부터 빼야 번호가 안 밀린다.
    for i in range(n_tpl - 1, -1, -1):
        pptx_tpl.drop_slide(prs, i)

    import io
    from fastapi.responses import StreamingResponse
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    name = pptx_tpl.safe(str((payload or {}).get("name") or "결과서")) + ".pptx"
    from urllib.parse import quote
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(name)}"},
    )


@app.post("/api/cycle/{cycle_id}")
async def save_cycle(cycle_id: str, data: dict):
    # 부여 ID(C-<연2><주차2>-<순번3>) — 없을 때만 새로 매긴다. 한 번 박히면 영원하다
    if not str((data or {}).get("cid") or "").strip():
        from datetime import datetime as _dt
        try:
            data["cid"] = await db.cycle_next_cid(db._cid_prefix_of(_dt.now()))
        except Exception:
            pass
    await db.cycle_upsert(cycle_id, data)
    # 누가 고쳤는지 함께 싣는다. 받는 쪽이 「내가 방금 저장한 것」 을 걸러야
    # 하고, 남이 한 것이면 이름을 말해 줘야 한다 — 사이클은 여럿이 나눠
    # 돌리는 자리라 「누가 3번을 Fail 로 바꿨나」 가 곧 알아야 할 일이다.
    _by = str((data or {}).get("updated_by") or "").strip()
    try:
        asyncio.create_task(broadcast({"type": "cycle_updated", "cycle_id": cycle_id, "user": _by}))
    except Exception:
        pass
    return {"success": True}

@app.post("/api/cycle/{cycle_id}/exec-ids")
async def cycle_exec_ids(cycle_id: str):
    """실행 ID 부여 — 사이클에 포함되는 값이다.

    CE 는 사이클 ID 에서 파생한다 (C-2633-002 → CE-2633-002). 사이클:실행이
    1:1 이라는 결정 그대로 — 재시험은 Clone 이 새 사이클을 만드니 새 CE 다.
    항목은 CETC-<파생>-NN. 멱등이라 실행 화면에 들어올 때마다 불러도 되고,
    나중에 항목을 더 넣으면 빈 번호만 채운다."""
    data = await db.cycle_get(cycle_id)
    if not data:
        raise HTTPException(404, "사이클이 없습니다")
    cid = str(data.get("cid") or "").strip()
    if not cid:
        from datetime import datetime as _dt
        try:
            cid = await db.cycle_next_cid(db._cid_prefix_of(_dt.now()))
            data["cid"] = cid
        except Exception:
            cid = ""
    base = cid[2:] if cid.startswith("C-") else cid
    changed = False
    if base:
        if not str(data.get("ce") or "").strip():
            data["ce"] = f"CE-{base}"
            changed = True
        items = data.get("items") or []
        used = set()
        for it in items:
            if isinstance(it, dict):
                m = str(it.get("ceid") or "")
                if m.startswith(f"CETC-{base}-"):
                    try:
                        used.add(int(m.rsplit("-", 1)[1]))
                    except ValueError:
                        pass
        nxt = 1
        for it in items:
            if not isinstance(it, dict) or str(it.get("ceid") or "").strip():
                continue
            while nxt in used:
                nxt += 1
            it["ceid"] = f"CETC-{base}-{nxt:02d}"
            used.add(nxt)
            changed = True
    if changed:
        await db.cycle_upsert(cycle_id, data)
    return {"ce": str(data.get("ce") or ""), "changed": changed}

@app.delete("/api/cycle/{cycle_id}")
async def delete_cycle(cycle_id: str):
    await db.cycle_delete(cycle_id)
    try: asyncio.create_task(broadcast({"type": "cycle_deleted", "cycle_id": cycle_id}))
    except Exception: pass
    return {"success": True}

# ───────────────────────────────────────────
# 라우터 - 장비
# ───────────────────────────────────────────
# ── 페이지·모듈 권한 ─────────────────────────────────────────────────
# 참고한 것: QMetry(모듈 × 권리 격자) · TestRail(역할 = 권한 묶음, 이름을
# 바꾸고 새로 만들 수 있다) · Zephyr Scale(맨 위 켬/끔) · Xray(제 체계를 안
# 만들고 Jira 권한에 얹는다).
#
# 어느 툴도 「메뉴 보임」 을 따로 관리하지 않는다 — 격자 하나에서 파생시킨다.
# 표를 두 벌 두면 반드시 어긋나기 때문이다. 여기도 그 방식이다: 「보기」 가
# 없으면 메뉴에 안 뜬다.
#
# **꺼진 채로 나간다.** 켜는 순간 아무도 아무것도 못 하는 사고를 막는다 —
# 표를 다 채운 뒤 사람이 켠다.
#
# 역할에 `jira` 칸을 비워 둔다. 계정 연동이 정리되면 Jira 그룹·프로젝트
# 역할이 여기 들어와 정본이 된다(지시). 그때 표를 다시 짜지 않아도 되게.
_PERM_RIGHTS = ("view", "create", "edit", "delete", "run", "folder")

_PERM_DEFAULT_ROLES = [
    {"key": "admin", "label": "관리자", "builtin": True, "jira": []},
    {"key": "lead", "label": "팀장", "builtin": True, "jira": []},
    {"key": "owner", "label": "담당", "builtin": True, "jira": []},
    {"key": "member", "label": "팀원", "builtin": True, "jira": []},
]


def _perm_doc() -> dict:
    """저장된 권한 문서. 없으면 「사용 안 함」 기본값."""
    d = _kv_load_sync("permissions", None)
    if not isinstance(d, dict):
        d = {}
    roles = d.get("roles")
    if not isinstance(roles, list) or not roles:
        roles = [dict(r) for r in _PERM_DEFAULT_ROLES]
    grid = d.get("grid")
    if not isinstance(grid, dict):
        grid = {}
    return {
        "enabled": bool(d.get("enabled")),
        "roles": roles,
        "grid": grid,
        # 옛 화면이 읽던 칸 — 아직 살려 둔다
        "perms": d.get("perms") if isinstance(d.get("perms"), dict) else {},
    }


@app.get("/api/permissions")
async def get_permissions():
    """모든 화면이 메뉴를 그리기 전에 읽는다 — 로그인만 하면 볼 수 있다."""
    return _perm_doc()


@app.post("/api/permissions")
async def save_permissions(data: dict = None, token: str = ""):
    """**관리자만**(지시) — 여기서 잘못 저장하면 아무도 못 들어온다."""
    _require_admin(token)
    data = data or {}
    cur = _perm_doc()

    roles = data.get("roles")
    if isinstance(roles, list) and roles:
        cur["roles"] = [
            {
                "key": str(r.get("key") or "").strip(),
                "label": str(r.get("label") or "").strip(),
                "builtin": bool(r.get("builtin")),
                "jira": [str(x) for x in (r.get("jira") or []) if str(x).strip()],
            }
            for r in roles
            if isinstance(r, dict) and str(r.get("key") or "").strip()
        ]
    grid = data.get("grid")
    if isinstance(grid, dict):
        cur["grid"] = {
            str(m): {
                str(rk): [x for x in (rv or []) if x in _PERM_RIGHTS]
                for rk, rv in (mv or {}).items()
            }
            for m, mv in grid.items()
        }
    if "enabled" in data:
        cur["enabled"] = bool(data.get("enabled"))
    if isinstance(data.get("perms"), dict):
        cur["perms"] = data["perms"]

    # **관리자를 0명으로 만들 수 없다.** 관리자 역할에서 SETUP 접근을 빼면
    # 아무도 이 화면에 다시 못 들어온다 — 잠긴 방에 열쇠를 두고 나오는 꼴이다.
    admin_key = next((r["key"] for r in cur["roles"] if r.get("builtin") and r["key"] == "admin"), "admin")
    cur["grid"].setdefault("settings", {})
    cur["grid"]["settings"][admin_key] = list(_PERM_RIGHTS)

    _kv_save_sync("permissions", cur)
    return {"ok": True, **cur}


# ───────────────────────────────────────────
@app.get("/api/devices")
async def get_devices():
    return load_json(DEVICES_FILE)

class DeviceCreate(BaseModel):
    group: str
    model: str
    ip: str
    protocol: str
    port: int
    username: str = ""
    password: str = ""
    description: str = ""

@app.post("/api/devices")
async def add_device(device: DeviceCreate):
    data = load_json(DEVICES_FILE)
    new_id = f"dev{len(data['devices'])+1:03d}"
    new_dev = {"id": new_id, "status": "unknown", **device.model_dump()}
    data["devices"].append(new_dev)
    save_json(DEVICES_FILE, data)
    return {"success": True, "device": new_dev}

@app.put("/api/devices/{device_id}")
async def update_device(device_id: str, device: DeviceCreate):
    data = load_json(DEVICES_FILE)
    for i, d in enumerate(data["devices"]):
        if d["id"] == device_id:
            status = d.get("status", "unknown")
            data["devices"][i] = {"id": device_id, "status": status, **device.model_dump()}
            save_json(DEVICES_FILE, data)
            return {"success": True}
    raise HTTPException(404, "장비를 찾을 수 없습니다")

@app.delete("/api/devices/{device_id}")
async def delete_device(device_id: str):
    data = load_json(DEVICES_FILE)
    data["devices"] = [d for d in data["devices"] if d["id"] != device_id]
    save_json(DEVICES_FILE, data)
    return {"success": True}

@app.post("/api/devices/{device_id}/connect")
async def connect_device(device_id: str):
    data = load_json(DEVICES_FILE)
    device = next((d for d in data["devices"] if d["id"] == device_id), None)
    if not device:
        raise HTTPException(404, "장비를 찾을 수 없습니다")

    proto = device.get("protocol", "").upper()
    ip = device["ip"]
    port = device["port"]

    if proto == "SSH":
        ok = check_ssh(ip, port, device.get("username",""), device.get("password",""))
    elif proto == "TELNET":
        ok = check_telnet(ip, port)
    elif proto in ("TCL", "API", "REST"):
        ok = check_tcp(ip, port)
    else:
        ok = check_tcp(ip, port)

    status = "connected" if ok else "disconnected"
    for d in data["devices"]:
        if d["id"] == device_id:
            d["status"] = status
    save_json(DEVICES_FILE, data)
    await broadcast({"type": "device_status", "id": device_id, "status": status})
    return {"success": True, "status": status}

@app.post("/api/devices/{device_id}/command")
async def run_command(device_id: str, body: dict):
    data = load_json(DEVICES_FILE)
    device = next((d for d in data["devices"] if d["id"] == device_id), None)
    if not device:
        raise HTTPException(404, "장비를 찾을 수 없습니다")
    command = body.get("command", "")
    proto = device.get("protocol", "").upper()
    try:
        if proto == "SSH":
            output = ssh_exec(device["ip"], device["port"], device["username"], device["password"], command)
        elif proto == "TCL":
            output = tcl_exec(command)
        else:
            output = f"[{proto}] 직접 명령 실행은 SSH/TCL만 지원합니다."
    except Exception as e:
        output = f"[오류] {e}"
    return {"output": output}

# ───────────────────────────────────────────
# 라우터 - 시험 절차
# ───────────────────────────────────────────
@app.get("/api/procedures")
async def get_procedures():
    return load_json(PROCEDURES_FILE)

class ProcedureCreate(BaseModel):
    group: str
    model: str
    name: str
    description: str = ""
    steps: list = []

@app.post("/api/procedures")
async def add_procedure(proc: ProcedureCreate):
    data = load_json(PROCEDURES_FILE)
    new_id = f"proc{len(data['procedures'])+1:03d}"
    new_proc = {"id": new_id, **proc.model_dump()}
    data["procedures"].append(new_proc)
    save_json(PROCEDURES_FILE, data)
    return {"success": True, "procedure": new_proc}

@app.put("/api/procedures/{proc_id}")
async def update_procedure(proc_id: str, proc: ProcedureCreate):
    data = load_json(PROCEDURES_FILE)
    for i, p in enumerate(data["procedures"]):
        if p["id"] == proc_id:
            data["procedures"][i] = {"id": proc_id, **proc.model_dump()}
            save_json(PROCEDURES_FILE, data)
            return {"success": True}
    raise HTTPException(404, "절차를 찾을 수 없습니다")

@app.delete("/api/procedures/{proc_id}")
async def delete_procedure(proc_id: str):
    data = load_json(PROCEDURES_FILE)
    data["procedures"] = [p for p in data["procedures"] if p["id"] != proc_id]
    save_json(PROCEDURES_FILE, data)
    return {"success": True}

# ───────────────────────────────────────────
# 라우터 - 시험 실행
# ───────────────────────────────────────────
@app.post("/api/run/{proc_id}")
async def run_procedure(proc_id: str):
    proc_data = load_json(PROCEDURES_FILE)
    proc = next((p for p in proc_data["procedures"] if p["id"] == proc_id), None)
    if not proc:
        raise HTTPException(404, "절차를 찾을 수 없습니다")

    dev_data = load_json(DEVICES_FILE)
    device = next((d for d in dev_data["devices"] if d["model"] == proc["model"]), None)

    result_id = f"result_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    results = []

    await broadcast({"type": "run_start", "proc_id": proc_id, "proc_name": proc["name"]})

    for step in proc.get("steps", []):
        await broadcast({"type": "step_start", "seq": step["seq"], "name": step["name"]})
        output = ""
        status = "PASS"

        try:
            if step["type"] == "CLI" and device:
                proto = device.get("protocol", "").upper()
                if proto == "SSH":
                    output = ssh_exec(device["ip"], device["port"], device["username"], device["password"], step["command"])
                elif proto == "TELNET":
                    output = await asyncio.to_thread(engine.netmiko_exec, device, step["command"])
                else:
                    output = f"[{proto}] CLI 명령 실행"
            elif step["type"] == "TCL":
                output = f"[TCL] {step['command']} 실행 시뮬레이션\n트래픽 생성 완료"
            elif step["type"] == "API":
                output = f"[API] {step['command']} 실행 완료"
            elif step["type"] == "검증":
                output = f"[검증] {step['command']}\n결과: 정상"
            elif step["type"] == "리포트":
                output = "[리포트] 결과 저장 완료"
            else:
                output = f"[{step['type']}] {step['command']}"
        except Exception as e:
            output = str(e)
            status = "FAIL"

        step_result = {
            "seq": step["seq"],
            "name": step["name"],
            "type": step["type"],
            "output": output,
            "status": status
        }
        results.append(step_result)
        await broadcast({"type": "step_done", **step_result})
        await asyncio.sleep(0.3)

    final = {
        "id": result_id,
        "proc_id": proc_id,
        "proc_name": proc["name"],
        "model": proc["model"],
        "timestamp": datetime.now().isoformat(),
        "steps": results,
        "overall": "PASS" if all(r["status"] == "PASS" for r in results) else "FAIL"
    }
    result_path = RESULTS_DIR / f"{result_id}.json"
    save_json(result_path, final)
    await broadcast({"type": "run_done", "result": final})
    return final

@app.get("/api/results")
async def get_results():
    results = []
    for f in sorted(RESULTS_DIR.glob("*.json"), reverse=True)[:20]:
        results.append(load_json(f))
    return {"results": results}

# ───────────────────────────────────────────
# 라우터 - 로컬 LLM 프록시
# ───────────────────────────────────────────
class LocalLLMRequest(BaseModel):
    endpoint: str
    model: str
    messages: list
    max_tokens: int = 4096
    context_size: int = 262144
    temperature: float = 0.7
    apikey: str = ""

@app.get("/api/chat/local/models")
async def get_local_models(endpoint: str):
    import httpx
    url = endpoint.rstrip("/") + "/models"
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            r = await client.get(url)
            return r.json()
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/chat/local/stream")
async def chat_local_stream(req: LocalLLMRequest):
    import httpx
    from fastapi.responses import StreamingResponse
    url = req.endpoint.rstrip("/") + "/chat/completions"
    headers = {"Content-Type": "application/json"}
    if req.apikey:
        headers["Authorization"] = f"Bearer {req.apikey}"
    safe_max_tokens = min(req.max_tokens, req.context_size, 4096)
    body = {
        "model": req.model,
        "messages": req.messages,
        "max_tokens": safe_max_tokens,
        "temperature": req.temperature,
        "stream": True,
    }
    async def generate():
        import json as _json
        # timeout: read=None(무제한 스트림). aiter_bytes 로 즉시 오는 바이트를 라인 단위로 파싱해 각 delta 를 곧바로 yield
        # (aiter_lines 는 라인이 완성될 때까지 잡아두는 구현이 있어 첫 청크가 지연 도착하면 전체가 한꺼번에 오는 것처럼 보일 수 있음)
        async with httpx.AsyncClient(timeout=httpx.Timeout(connect=10.0, read=None, write=10.0, pool=10.0)) as client:
            async with client.stream("POST", url, headers=headers, json=body) as r:
                _buf = b""
                async for _chunk in r.aiter_bytes():
                    if not _chunk:
                        continue
                    _buf += _chunk
                    while True:
                        _nl = _buf.find(b"\n")
                        if _nl < 0:
                            break
                        _line = _buf[:_nl].decode("utf-8", errors="ignore").rstrip("\r")
                        _buf = _buf[_nl+1:]
                        if not _line.startswith("data: "):
                            continue
                        data = _line[6:]
                        if data.strip() == "[DONE]":
                            yield "data: [DONE]\n\n"
                            return
                        try:
                            ch = _json.loads(data)
                            delta = ch["choices"][0]["delta"].get("content", "")
                            if delta:
                                yield f"data: {_json.dumps({'text': delta})}\n\n"
                        except Exception:
                            pass
    # X-Accel-Buffering: no — Nginx/역방향 프록시 앞단이 있을 때 SSE 청크 버퍼링 방지(즉시 flush)
    # Content-Encoding: identity — 상위 GZipMiddleware 가 이미 인코딩된 응답으로 인식해 재압축을 스킵함
    return StreamingResponse(generate(), media_type="text/event-stream",
                             headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache",
                                      "Connection": "keep-alive", "Content-Encoding": "identity"})

# ───────────────────────────────────────────
# 라우터 - Dify ChatFlow 지식 어시스턴트 (OpenWebUI Functions 이식)
# API 키/URL 은 서버에만 보관(브라우저로 노출하지 않음). 환경변수로 덮어쓸 수 있다.
# 프론트는 assistant id(specs/qag/trouble) 만 보낸다.
# ───────────────────────────────────────────
DIFY_BASE_URL = os.environ.get("DIFY_BASE_URL", "http://10.10.30.219:3897/v1")
# file_var: 첨부 파일을 받을 ChatFlow 입력변수 이름(Dify Start 노드의 File 변수). 빈 값이면 message files 로 전달.
_DIFY_FILE_VAR = os.environ.get("DIFY_FILE_VAR", "log_file")
# 지식 어시스턴트는 하드코딩하지 않는다 — LLM 설정 화면에서 관리(dify_assistants.json). API 키는 그 파일(서버)에만 존재한다.

# ── 지식 어시스턴트(Dify) 동적 관리 — 서버측 저장(dify_assistants.json). API 키는 프론트로 노출하지 않음 ──
DIFY_FILE = DATA_DIR / "integrations" / "dify_assistants.json"

def _dify_load():
    """저장된 Dify 어시스턴트 목록 — LLM 설정(dify_assistants.json)에서만 관리. 하드코딩/자동 시드 없음."""
    try:
        data = load_json(DIFY_FILE)
    except Exception:
        data = None
    lst = data.get("assistants") if isinstance(data, dict) else None
    return lst or []

def _dify_save(lst):
    save_json(DIFY_FILE, {"assistants": lst})

def _dify_get(aid):
    aid = (aid or "").strip()
    for a in _dify_load():
        if a.get("id") == aid:
            return a
    return None

def _dify_slug(name, lst):
    base = "".join(ch if (ord(ch) < 128 and ch.isalnum()) else "-" for ch in (name or "").lower()).strip("-") or ("dify" + str(len(lst) + 1))
    aid, n = base, 2
    while any(a.get("id") == aid for a in lst):
        aid = base + "-" + str(n); n += 1
    return aid

@app.get("/api/dify/assistants")
async def dify_assistants_list():
    # 키는 빼고 has_key 만 반환 — 브라우저로 API 키를 보내지 않는다
    out = []
    for a in _dify_load():
        out.append({"id": a.get("id"), "name": a.get("name"), "endpoint": a.get("endpoint", ""),
                    "file_var": a.get("file_var", ""), "icon": a.get("icon", ""), "has_key": bool(a.get("key")),
                    "greeting": a.get("greeting", ""), "placeholder": a.get("placeholder", ""),
                    "public": a.get("public", True),
                    "type": a.get("type", "dify"), "llm_id": a.get("llm_id", ""),
                    "prompt": a.get("prompt", ""), "rag": bool(a.get("rag", False)),
                    "kb_group": a.get("kb_group", "external")})
    return {"assistants": out}

@app.get("/api/dify/assistants/{aid}")
async def dify_assistant_detail(aid: str):
    # 관리 화면(LLM 설정) 전용 — API 키 포함 반환(편집 필드에 실제 값 표시용). 채팅 목록(GET /assistants)은 여전히 키 마스킹.
    a = _dify_get(aid)
    if not a:
        return {"ok": False, "error": "어시스턴트를 찾을 수 없습니다."}
    return {"id": a.get("id"), "name": a.get("name"), "endpoint": a.get("endpoint", ""),
            "file_var": a.get("file_var", ""), "icon": a.get("icon", ""), "key": a.get("key", ""),
            "greeting": a.get("greeting", ""), "placeholder": a.get("placeholder", ""),
            "public": a.get("public", True),
            "type": a.get("type", "dify"), "llm_id": a.get("llm_id", ""),
            "prompt": a.get("prompt", ""), "rag": bool(a.get("rag", False)),
            "rag_sources": a.get("rag_sources") or [],   # 저장된 소스별 활성화+우선순위 (편집 폼 복원)
            "kb_group": a.get("kb_group", "external")}

@app.post("/api/dify/assistants")
async def dify_assistants_add(data: dict):
    name = str(data.get("name", "")).strip()
    if not name:
        return {"ok": False, "error": "이름은 필수입니다."}
    lst = _dify_load()
    aid = str(data.get("id", "")).strip() or _dify_slug(name, lst)
    if any(a.get("id") == aid for a in lst):
        aid = _dify_slug(aid, lst)
    _type = str(data.get("type", "dify")).strip() or "dify"
    lst.append({"id": aid, "name": name,
                "type": _type,
                "llm_id": str(data.get("llm_id", "")).strip(),
                "prompt": str(data.get("prompt", "")),
                "rag": bool(data.get("rag", False)),
                "kb_group": str(data.get("kb_group", "external")).strip() or "external",
                "endpoint": str(data.get("endpoint", "")).strip() or (DIFY_BASE_URL if _type == "dify" else ""),
                "key": str(data.get("key", "")).strip(),
                "file_var": str(data.get("file_var", _DIFY_FILE_VAR)).strip(),
                "icon": str(data.get("icon", "")).strip(),
                "greeting": str(data.get("greeting", "")),
                "placeholder": str(data.get("placeholder", "")),
                "public": bool(data.get("public", True))})
    _dify_save(lst)
    return {"ok": True, "id": aid}

@app.put("/api/dify/assistants/{aid}")
async def dify_assistants_update(aid: str, data: dict):
    lst = _dify_load()
    for a in lst:
        if a.get("id") == aid:
            if str(data.get("name", "")).strip():
                a["name"] = str(data.get("name")).strip()
            if "type" in data:
                a["type"] = str(data.get("type", "")).strip() or a.get("type", "dify")
            if "llm_id" in data:
                a["llm_id"] = str(data.get("llm_id", "")).strip()
            if "prompt" in data:
                a["prompt"] = str(data.get("prompt", ""))
            if "rag" in data:
                a["rag"] = bool(data.get("rag"))
            if "rag_sources" in data:
                # 소스별 활성화+우선순위: [{"source":"tc"|"manual"|"confluence","enabled":bool,"priority":int}, ...]
                rs = data.get("rag_sources")
                if isinstance(rs, list):
                    clean = []
                    for it in rs:
                        if not isinstance(it, dict):
                            continue
                        src = str(it.get("source", "")).strip()
                        if src not in ("tc", "manual", "confluence"):
                            continue
                        clean.append({
                            "source": src,
                            "enabled": bool(it.get("enabled")),
                            "priority": int(it.get("priority") or 0),
                        })
                    a["rag_sources"] = clean
            if "kb_group" in data:
                a["kb_group"] = str(data.get("kb_group", "")).strip()
            if "endpoint" in data:
                a["endpoint"] = str(data.get("endpoint", "")).strip() or a.get("endpoint", "")
            if "file_var" in data:
                a["file_var"] = str(data.get("file_var", "")).strip()
            if "icon" in data:
                a["icon"] = str(data.get("icon", "")).strip()
            if "public" in data:
                a["public"] = bool(data.get("public"))
            if "greeting" in data:
                a["greeting"] = str(data.get("greeting", ""))
            if "placeholder" in data:
                a["placeholder"] = str(data.get("placeholder", ""))
            k = str(data.get("key", "")).strip()
            if k:  # 키는 새로 입력했을 때만 교체(빈값이면 기존 키 유지)
                a["key"] = k
            _dify_save(lst)
            return {"ok": True}
    return {"ok": False, "error": "어시스턴트를 찾을 수 없습니다."}

@app.delete("/api/dify/assistants/{aid}")
async def dify_assistants_delete(aid: str):
    lst = _dify_load()
    new = [a for a in lst if a.get("id") != aid]
    if len(new) == len(lst):
        return {"ok": False, "error": "어시스턴트를 찾을 수 없습니다."}
    _dify_save(new)
    return {"ok": True}

@app.post("/api/dify/chat")
async def dify_chat(data: dict):
    import httpx
    from fastapi.responses import StreamingResponse
    assistant = (data.get("assistant") or "").strip()
    query = (data.get("query") or "").strip()
    conv_id = (data.get("conversation_id") or "").strip()
    user = (data.get("user") or "utop-user").strip() or "utop-user"
    cfg = _dify_get(assistant)

    async def generate():
        def sse(obj):
            return "data: " + json.dumps(obj, ensure_ascii=False) + "\n\n"
        if not cfg:
            yield sse({"text": "알 수 없는 어시스턴트입니다."}); yield "data: [DONE]\n\n"; return
        if not query:
            yield sse({"text": "빈 메시지입니다."}); yield "data: [DONE]\n\n"; return
        headers = {"Authorization": "Bearer " + cfg["key"], "Content-Type": "application/json"}
        payload = {"inputs": {}, "query": query, "response_mode": "streaming", "user": user}
        if conv_id:
            payload["conversation_id"] = conv_id
        # 첨부 파일(미리 /api/dify/upload 로 올려 받은 upload_file_id)
        dify_files = []
        for f in (data.get("files") or []):
            fid = ((f.get("upload_file_id") or f.get("id") or "").strip()) if isinstance(f, dict) else ""
            if fid:
                dify_files.append({"type": (f.get("type") or "document"), "transfer_method": "local_file", "upload_file_id": fid})
        if dify_files:
            fvar = (cfg.get("file_var") or "").strip()
            if fvar:
                # ChatFlow 입력변수(예: log_file)로 전달 — Start 노드의 File 변수
                payload["inputs"][fvar] = dify_files[0]
            else:
                # 입력변수 미지정 시 message files(vision/sys.files)로 전달
                payload["files"] = dify_files
        url = (cfg.get("endpoint") or DIFY_BASE_URL).rstrip("/") + "/chat-messages"
        timeout = httpx.Timeout(connect=10.0, read=300.0, write=10.0, pool=10.0)
        accumulated = ""
        has_output = False
        new_conv = conv_id
        try:
            async with httpx.AsyncClient(timeout=timeout) as client:
                async with client.stream("POST", url, headers=headers, json=payload) as resp:
                    if resp.status_code != 200:
                        b = await resp.aread()
                        yield sse({"text": "Dify API 오류 (%d): %s" % (resp.status_code, b.decode(errors="replace")[:400])})
                        yield "data: [DONE]\n\n"
                        return
                    async for line in resp.aiter_lines():
                        if not line or not line.startswith("data: "):
                            continue
                        try:
                            d = json.loads(line[6:])
                        except Exception:
                            continue
                        event = d.get("event", "")
                        if event in ("message", "agent_message"):
                            answer = d.get("answer", "")
                            if not answer:
                                continue
                            cid = d.get("conversation_id", "")
                            if cid:
                                new_conv = cid
                            if accumulated and answer.startswith(accumulated):
                                delta = answer[len(accumulated):]
                                accumulated = answer
                            else:
                                delta = answer
                                accumulated += answer
                            if delta:
                                has_output = True
                                yield sse({"text": delta})
                        elif event == "message_end":
                            cid = d.get("conversation_id", "")
                            if cid:
                                new_conv = cid
                        elif event == "message_replace":
                            answer = d.get("answer", "")
                            if answer:
                                has_output = True
                                yield sse({"text": answer})
                        elif event == "text_chunk":
                            chunk = (d.get("data") or {}).get("text", "")
                            if chunk:
                                has_output = True
                                yield sse({"text": chunk})
                        elif event == "workflow_finished":
                            wf = d.get("data", {}) or {}
                            if wf.get("status") == "failed":
                                yield sse({"text": "\n\n❌ 워크플로우 실패: " + str(wf.get("error", ""))})
                                break
                            if not has_output:
                                outputs = wf.get("outputs", {})
                                if isinstance(outputs, dict):
                                    text = outputs.get("answer", "") or outputs.get("result", "") or outputs.get("text", "")
                                    if text:
                                        has_output = True
                                        yield sse({"text": text})
                        elif event == "node_finished":
                            nd = d.get("data", {}) or {}
                            if nd.get("status") == "failed":
                                yield sse({"text": "\n\n❌ 노드 실패 [%s]: %s" % (nd.get("title", ""), nd.get("error", ""))})
                                break
                        elif event == "error":
                            emsg = d.get("message", "") or d.get("msg", "")
                            code = d.get("code", "")
                            yield sse({"text": "\n\n❌ Dify 오류 [%s]: %s" % (code, emsg)})
                            break
        except httpx.TimeoutException:
            yield sse({"text": "\n\n⏰ 타임아웃(300초)"})
        except httpx.ConnectError:
            yield sse({"text": "\n\n❌ 연결 실패: " + DIFY_BASE_URL})
        except Exception as e:
            yield sse({"text": "\n\n❌ %s: %s" % (type(e).__name__, e)})
        if new_conv:
            yield sse({"conv": new_conv})
        yield "data: [DONE]\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")

@app.post("/api/dify/upload")
async def dify_upload(assistant: str = Form(...), user: str = Form("utop-user"), file: UploadFile = File(...)):
    # 파일을 Dify(/v1/files/upload)에 올려 upload_file_id 를 받아 프론트에 돌려준다. 키는 서버 보관.
    import httpx
    cfg = _dify_get(assistant)
    if not cfg:
        return {"ok": False, "error": "알 수 없는 어시스턴트입니다."}
    try:
        content = await file.read()
    except Exception as e:
        return {"ok": False, "error": "파일 읽기 실패: %s" % e}
    files = {"file": (file.filename or "upload", content, file.content_type or "application/octet-stream")}
    form = {"user": (user or "utop-user")}
    headers = {"Authorization": "Bearer " + cfg["key"]}
    url = (cfg.get("endpoint") or DIFY_BASE_URL).rstrip("/") + "/files/upload"
    try:
        async with httpx.AsyncClient(timeout=httpx.Timeout(connect=10.0, read=120.0, write=120.0, pool=10.0)) as client:
            r = await client.post(url, headers=headers, files=files, data=form)
        if r.status_code not in (200, 201):
            low = (r.text or "").lower()
            if r.status_code == 413 or "file_too_large" in low or "too large" in low:
                return {"ok": False, "error": "파일이 Dify 허용 크기를 초과했습니다(413). 더 작은 파일을 쓰거나, 이미지는 축소/캡처 후 첨부하세요. (제한은 Dify 서버 설정)"}
            return {"ok": False, "error": "Dify 업로드 오류 (%d): %s" % (r.status_code, r.text[:300])}
        j = r.json()
        fid = j.get("id")
        if not fid:
            return {"ok": False, "error": "Dify 응답에 파일 id 가 없습니다."}
        mime = (j.get("mime_type") or file.content_type or "").lower()
        ext = (j.get("extension") or "").lower().lstrip(".")
        is_img = mime.startswith("image/") or ext in ("jpg", "jpeg", "png", "gif", "webp", "bmp", "svg")
        return {"ok": True, "id": fid, "type": ("image" if is_img else "document"), "name": j.get("name") or file.filename}
    except httpx.ConnectError:
        return {"ok": False, "error": "연결 실패: " + DIFY_BASE_URL}
    except httpx.TimeoutException:
        return {"ok": False, "error": "업로드 타임아웃"}
    except Exception as e:
        return {"ok": False, "error": "%s: %s" % (type(e).__name__, e)}

@app.post("/api/chat/local")
async def chat_local(req: LocalLLMRequest):
    import httpx
    url = req.endpoint.rstrip("/") + "/chat/completions"
    print(f"[LocalLLM] URL: {url}")
    print(f"[LocalLLM] Model: {req.model}")
    headers = {"Content-Type": "application/json"}
    if req.apikey:
        headers["Authorization"] = f"Bearer {req.apikey}"
    # max_tokens는 context_size를 초과할 수 없음 (Dify 방식)
    safe_max_tokens = min(req.max_tokens, req.context_size, 4096)
    body = {
        "model": req.model,
        "messages": req.messages,
        "max_tokens": safe_max_tokens,
        "temperature": req.temperature,
    }
    try:
        async with httpx.AsyncClient(timeout=120) as client:
            r = await client.post(url, headers=headers, json=body)
            print(f"[LocalLLM] Status: {r.status_code}")
            print(f"[LocalLLM] Response: {r.text[:500]}")
            if not r.is_success:
                return {"reply": f"[로컬 LLM 오류] {r.status_code}\n\n{r.text}"}
            data = r.json()
            reply = data["choices"][0]["message"]["content"]
            return {"reply": reply}
    except Exception as e:
        print(f"[LocalLLM] Error: {e}")
        return {"reply": f"[로컬 LLM 오류] {e}"}

# ───────────────────────────────────────────
# 라우터 - Claude 채팅
# ───────────────────────────────────────────
class ChatRequest(BaseModel):
    message: str
    history: list = []
    max_tokens: int = 2048

# ══════════ RAG: 청킹 + BM25 검색 (Phase1 희소검색, 임베딩은 후속) ══════════
RAG_CHUNK_SIZE = 500
RAG_CHUNK_OVERLAP = 80
_RAG_CACHE = {"sig": None, "corpus": None}
# ── 청크 임베딩 캐시 (numpy .npy 바이너리) — 재시작/파일변경 시 전체 재임베딩 방지 + 빠른 로드/검색 ──
EMBED_NPY_FILE = DATA_DIR / "state" / "rag_embed.npy"
EMBED_KEYS_FILE = DATA_DIR / "state" / "rag_embed_keys.json"
EMBED_CACHE_FILE = DATA_DIR / "state" / "rag_embed_cache.json"   # 구버전(JSON) — 자동 변환
_EMBED = {"keys": None, "rows": None, "mat": None, "norm": None, "dirty": False}
def _embed_key(text, model):
    import hashlib
    return hashlib.md5((str(model) + "|" + str(text)).encode("utf-8")).hexdigest()
def _embed_renorm():
    import numpy as np
    m = _EMBED["mat"]
    if m is None or not len(m):
        _EMBED["norm"] = None; return
    n = np.linalg.norm(m, axis=1, keepdims=True); n[n == 0] = 1.0
    _EMBED["norm"] = (m / n).astype("float32")
def _embed_load():
    if _EMBED["keys"] is not None:
        return
    import numpy as np
    keys = []; mat = None
    try:
        if EMBED_NPY_FILE.exists() and EMBED_KEYS_FILE.exists():
            mat = np.load(str(EMBED_NPY_FILE)); keys = load_json(EMBED_KEYS_FILE) or []
        elif EMBED_CACHE_FILE.exists():   # 구버전 {key:vec} JSON → 행렬로 변환
            old = load_json(EMBED_CACHE_FILE) or {}; keys = list(old.keys())
            if keys:
                mat = np.asarray([old[k] for k in keys], dtype="float32")
    except Exception:
        keys = []; mat = None
    if mat is None or not len(keys):
        _EMBED["keys"] = {}; _EMBED["rows"] = []; _EMBED["mat"] = None; _EMBED["norm"] = None
    else:
        _EMBED["mat"] = np.asarray(mat, dtype="float32"); _EMBED["rows"] = list(keys)
        _EMBED["keys"] = {k: i for i, k in enumerate(keys)}; _embed_renorm()
        if not EMBED_NPY_FILE.exists():   # JSON→npy 1회 변환 저장 후 구파일 제거
            _EMBED["dirty"] = True; _embed_save()
            try: EMBED_CACHE_FILE.unlink()
            except Exception: pass
def _embed_add(keys, vecs):
    import numpy as np
    _embed_load()
    newk = []; newv = []
    for k, v in zip(keys, vecs):
        if k in _EMBED["keys"]: continue
        _EMBED["keys"][k] = len(_EMBED["rows"]) + len(newk); newk.append(k); newv.append(v)
    if not newk: return
    nv = np.asarray(newv, dtype="float32")
    _EMBED["mat"] = nv if _EMBED["mat"] is None else np.vstack([_EMBED["mat"], nv])
    _EMBED["rows"].extend(newk); _EMBED["dirty"] = True; _embed_renorm()
def _embed_save():
    import numpy as np
    if not _EMBED.get("dirty"): return
    try:
        if _EMBED["mat"] is not None:
            np.save(str(EMBED_NPY_FILE), _EMBED["mat"]); save_json(EMBED_KEYS_FILE, _EMBED["rows"])
        _EMBED["dirty"] = False
    except Exception:
        pass
@app.on_event("startup")
async def _db_init():
    """PostgreSQL 커넥션 풀 초기화 (필수 — 이후 모든 데이터 접근이 db.* 로 감).
    정리·백필 작업은 서버 기동을 막지 않도록 백그라운드로 미룸."""
    if N2X_RELAY_ONLY:
        print("[startup] N2X 중계 전용 — DB 를 잡지 않습니다", flush=True)
        return
    await db.init_pool()
    # 스키마를 기동할 때마다 적용한다. 도커 initdb 훅은 볼륨이 빌 때 한 번만
    # 돌아서, 이미 쓰고 있는 설치처에는 나중에 추가한 컬럼이 반영되지 않는다.
    try:
        await db.apply_schema()
    except Exception as e:
        print(f'[startup] 스키마 적용 실패: {e}', flush=True)
    # 겹싸여 저장된 jsonb(장비 data · 접속 params) 벗기기 — 한 번만 돌면 끝난다
    try:
        await db._repair_double_json()
    except Exception as e:
        print(f'[startup] jsonb 손질 실패: {e}', flush=True)
    # 워커 스레드(run_cli 등)에서 asyncio.run_coroutine_threadsafe 호출용 메인 루프 참조 저장.
    # 스레드 안에서는 asyncio.get_event_loop() 가 새 루프를 만들거나 실패하므로,
    # 요청 처리 스레드에서 broadcast() 를 예약하려면 반드시 여기서 잡은 루프를 써야 한다.
    global _MAIN_LOOP
    try: _MAIN_LOOP = asyncio.get_running_loop()
    except Exception: _MAIN_LOOP = None

    # 세션: 레거시 sessions.json → DB 이전 (1회) 후 DB → in-memory 로드
    try: await _migrate_sessions_file_to_db()
    except Exception as e: print(f"[startup] session migration failed: {e}", flush=True)
    try: await _load_sessions_from_db()
    except Exception as e: print(f"[startup] session load failed: {e}", flush=True)

    # 사용자: 레거시 users.json → DB(app_kv 'users') 로 이전 (1회)
    try: await _migrate_users_file_to_db()
    except Exception as e: print(f"[startup] users migration failed: {e}", flush=True)
    # 사용자 캐시 초기 로드 + 기본 admin 계정 보장 (DB 비었으면 생성). async 로 호출해야 loop deadlock 없음.
    try: await _init_users_async()
    except Exception as e: print(f"[startup] init_users_file failed: {e}", flush=True)

    # 결함 ID — 옛 무작위 꼬리(DEF-a1b2c3…)를 DEF-<프로젝트키>-<순번> 으로 이전 (멱등)
    try:
        _dn = await db.defect_renumber_legacy()
        if _dn: print(f"[startup] 결함 ID {_dn}건을 새 체계로 이전", flush=True)
    except Exception as e:
        print(f"[startup] defect renumber failed: {e}", flush=True)

    # 사이클 부여 ID — cid 없는 회차에 C-<연2><주차2>-<순번3> 을 채운다 (멱등)
    try:
        _cn = await db.cycle_backfill_cids()
        if _cn: print(f"[startup] 사이클 ID {_cn}건 부여 (C-연주차-순번)", flush=True)
    except Exception as e:
        print(f"[startup] cycle cid backfill failed: {e}", flush=True)

    # 실행 타입 「혼합」 은 뺐다(합의) — 기동 때 지워 두면 253 도
    # update.sh 만으로 같아진다. 없으면 그냥 지나간다(멱등).
    try:
        if await db.code_delete("tc_run_type", "혼합"):
            print("[startup] 실행 타입 「혼합」 제거", flush=True)
    except Exception as e:
        print(f"[startup] 혼합 제거 실패: {e}", flush=True)

    # 폴더 이동으로 낡은 요구사항 분류 사슬(cat1~4)을 트리 기준으로
    # 재작성한다 — 이동 API 가 그때그때 맞추지만, 그 전에 낡은 자료가
    # 이미 있고 253 도 update.sh 만으로 같아져야 한다 (멱등).
    try:
        _rn = await _req_chain_resync()
        if _rn:
            print(f"[startup] 요구사항 분류 사슬 {_rn}건 재작성", flush=True)
    except Exception as e:
        print(f"[startup] 사슬 재작성 실패: {e}", flush=True)

    # 파일 → DB(app_kv) 이전 (파일이 정본이면 DB 덮어씀). ai_usage/ai_feedback 는 _load_items_store 매핑도 등록.
    _KV_MIGRATIONS = [
        ("chat_sessions", CHAT_SESS_FILE),
        ("ai_usage", AI_USAGE_FILE),
        ("ai_feedback", FEEDBACK_FILE),
        ("learned_procedures", LEARNED_FILE),
        ("release_summary", RELEASE_SUMMARY_FILE),
        ("manpower", RSC_MANPOWER_FILE),
        ("device_catalog", DEVICE_CATALOG_FILE),
        ("racks", RACKS_FILE),
        # ★ 여기 빠지면 저장(DB)은 되는데 재시작 후 안 읽힌다 —
        #   _kv_load_sync 는 등록된 키만 기동 때 DB 에서 캐시로 채운다.
        #   실사고: 253 에서 update.sh(재시작) 뒤 탭 이름이 초기값으로 복귀.
        #   파일은 원래 없던 키라 경로는 자리표시용이다(파일 없음 → DB 만).
        ("code_kind_labels", DATA_DIR / "code_kind_labels.json"),
        ("code_kind_hidden", DATA_DIR / "code_kind_hidden.json"),
        # INFO 필드의 폭·모양·정렬·글꼴. 바로 위 두 형제는 등록해 두고 이것만
        # 빠져 있었다 — 저장은 DB 에 되는데 다시 올리면 `_kv_load_sync` 가
        # 빈 값을 캐시에 박고, 다음 저장이 그 빈 값으로 DB 를 덮어썼다.
        # 실사고: 「폭이 자꾸 변경돼」 — 배포할 때마다 열 폭이 기본값으로 복귀.
        ("code_kind_style", DATA_DIR / "code_kind_style.json"),
        # 페이지·모듈 권한. 옛 파일(config/permissions.json)이 정본이면 그것을
        # DB 로 옮긴다. 등록을 빼면 재시작 때 빈 격자가 캐시에 박히고 다음
        # 저장이 DB 를 덮어써 **권한이 통째로 날아간다** — 바로 위에서 겪은 것.
        ("permissions", PERMISSIONS_FILE),
        ("cycle_desc_template", DATA_DIR / "cycle_desc_template.json"),
        # 자연어 시험 첫 화면의 질문 보기 — 등록 안 하면 재시작 때 빈 값이
        # 캐시에 박히고 다음 저장이 DB 를 덮어쓴다(원본 앱에서 겪은 덫).
        ("ai_examples", DATA_DIR / "ai_examples.json"),
        # 자연어 시험 기록. 등록 안 하면 재시작 때 _kv_load_sync 가 빈 값을
        # 캐시에 박고(등록된 키만 DB 에서 채운다), 다음 저장이 그 빈 값으로
        # DB 를 덮어써 **기록이 통째로 날아간다**. 실사고: 재시작 뒤 시험
        # 기록이 사라졌다.
        ("nl_chats", DATA_DIR / "nl_chats.json"),
        # 조직도(회사 → 그룹 → 담당 → 팀 → 사람). 계정 화면이 이걸로 묶어 본다.
        # 등록을 빼면 재시작 때 빈 조직도가 캐시에 박히고 다음 저장이 DB 를
        # 덮어써 통째로 날아간다 — 위 형제들이 겪은 그 덫이다.
        ("org_tree", DATA_DIR / "org_tree.json"),
    ]
    for _key, _fp in _KV_MIGRATIONS:
        _kv_register_fallback(_key, _fp)
        try: await _kv_init_async(_key, _fp, sizeguard=True)
        except Exception as _me: print(f"[startup] KV migrate '{_key}' failed: {_me}", flush=True)
    # 조직도 씨앗 — **비어 있을 때만** 채운다.
    #
    # 조직도는 app_kv(DB) 에 산다. 그래서 코드만 받은 서버(253)는 계정 화면이
    # 예전 납작한 목록 그대로였다(지적). DATA_DIR 은 도커 볼륨이라 자료를
    # 거기 두면 이미지를 따라가지 못한다 — 그래서 씨앗은 backend/ 안에 둔다.
    #
    # 이미 조직도가 있으면 **손대지 않는다.** 사람이 옮겨 놓은 것을 배포할
    # 때마다 되돌리면, 고쳐도 소용없는 화면이 된다.
    try:
        if not _kv_load_sync("org_tree", None):
            _seed = Path(__file__).parent / "seed" / "org_tree.json"
            if _seed.exists():
                _kv_save_sync("org_tree", json.loads(_seed.read_text(encoding="utf-8")))
                print("[startup] 조직도 씨앗 심음", flush=True)
    except Exception as _se:
        print(f"[startup] 조직도 씨앗 실패: {_se}", flush=True)

    # 역할 씨앗 — **딱 한 번만** 심는다.
    #
    # 팀장 24 · 담당 18 은 사람이 손으로 정한 값이다(role_by 가 비어 있다).
    # 조직도처럼 다시 만들어 낼 수 없어, 아이디→역할 명단을 씨앗으로 싣는다.
    #
    # 두 가지를 지킨다.
    #  · **팀원인 사람만** 올린다 — 이미 정해 둔 역할을 배포가 덮으면 안 된다.
    #    관리자는 손도 안 댄다(내리면 그 사람이 화면을 못 쓴다).
    #  · 한 번 심고 표식을 남긴다. 안 그러면 253 에서 누군가를 일부러 팀원으로
    #    되돌려도 다음 재시작이 도로 올려, 고쳐도 소용없는 화면이 된다.
    try:
        _rd = _users_load_sync()
        if not _rd.get("role_seed_v1"):
            _rs = Path(__file__).parent / "seed" / "roles.json"
            _seeded = 0
            if _rs.exists():
                _want = json.loads(_rs.read_text(encoding="utf-8"))
                for _u in _rd["users"]:
                    _r = _want.get(_u.get("username"))
                    if _r and _u.get("role") == "팀원":
                        _u["role"] = _r
                        _seeded += 1
            _rd["role_seed_v1"] = True
            _users_save_sync(_rd)
            if _seeded:
                print(f"[startup] 역할 씨앗 {_seeded}명 심음", flush=True)
    except Exception as _re2:
        print(f"[startup] 역할 씨앗 실패: {_re2}", flush=True)

    # 조직도의 장 → 계정 역할 「담당」 을 **한 번 맞춘다**. 저장할 때만 맞추면
    # 이미 들어 있는 조직도는 아무도 다시 저장하기 전까지 표(담당)와 편집판
    # (팀원)이 어긋난 채로 남는다(지적). 관리자는 안 내린다.
    try:
        _org = _kv_load_sync("org_tree", None)
        if isinstance(_org, dict) and _org.get("name"):
            _r = _apply_org_roles(_org)
            if _r.get("role_up") or _r.get("role_down"):
                print(f"[startup] 조직 역할 맞춤: {_r}", flush=True)
    except Exception as _oe:
        print(f"[startup] 조직 역할 맞춤 실패: {_oe}", flush=True)

    # ai_usage/ai_feedback 는 _load_items_store(path) 우회 매핑 등록
    _ITEMS_STORE_KV_MAP[str(AI_USAGE_FILE)] = "ai_usage"
    _ITEMS_STORE_KV_MAP[str(FEEDBACK_FILE)] = "ai_feedback"

    async def _bg_maintenance():
        # cycle.data_summary 백필
        try:
            n = await db.cycle_backfill_summary()
            if n > 0:
                print(f"[startup-bg] cycle data_summary backfilled: {n} rows")
        except Exception as e:
            print(f"[startup-bg] cycle backfill failed: {e}")
        # REQ.tc 안 stale 참조 정리
        try:
            n = await _cleanup_stale_req_tc_refs()
            if n > 0:
                print(f"[startup-bg] REQ 안 stale TC 참조 정리: {n} 건")
        except Exception as e:
            print(f"[startup-bg] stale ref cleanup failed: {e}")
        # REQ 중복 row 정리
        try:
            n = await _cleanup_duplicate_reqs()
            if n > 0:
                print(f"[startup-bg] 중복 REQ row 정리: {n} 건")
        except Exception as e:
            print(f"[startup-bg] duplicate REQ cleanup failed: {e}")

    asyncio.create_task(_bg_maintenance())


async def _cleanup_stale_req_tc_refs() -> int:
    """모든 REQ 를 순회하며 tc 배열에서 DB 에 없는 tcid 참조를 제거."""
    async with db.pool().acquire() as c:
        alive_rows = await c.fetch("SELECT tcid FROM tc")
        alive = {r["tcid"] for r in alive_rows}
    reqs = await db.req_list_full()
    fixed = 0
    for r in reqs:
        if not isinstance(r, dict):   # data 가 dict 아니면 (예: 문자열) 스킵
            continue
        refs = r.get("tc") or []
        if not isinstance(refs, list):
            continue
        cleaned = [ref for ref in refs if isinstance(ref, dict) and (ref.get("tcid") in alive)]
        if len(cleaned) != len(refs):
            r["tc"] = cleaned
            await db.req_upsert(r.get("id") or r.get("reqid"), r)
            fixed += (len(refs) - len(cleaned))
    return fixed


async def _cleanup_duplicate_reqs() -> int:
    """
    예전 saveOneREQ 버그로 PK 와 data.id 가 다른 두 REQ row 가 생긴 경우 병합.
    같은 data.id 를 가진 여러 PK 발견 시:
      - 하나만 남기고 나머지 삭제
      - 남길 것 = data.id 와 PK 가 일치하는 row (가장 정통)
      - 없으면 updated_at 최신 것
    """
    async with db.pool().acquire() as c:
        rows = await c.fetch("SELECT id AS pk, data, updated_at FROM req")
    # group by data.id
    groups = {}
    for r in rows:
        d = r["data"] or {}
        if not isinstance(d, dict):   # 손상된 row (data 가 dict 아님) 스킵
            continue
        did = d.get("id") or r["pk"]
        groups.setdefault(did, []).append({"pk": r["pk"], "data": d, "updated_at": r["updated_at"]})
    total_removed = 0
    for did, arr in groups.items():
        if len(arr) < 2:
            continue
        # 남길 것 선택
        canonical = next((x for x in arr if x["pk"] == did), None)
        if canonical is None:
            arr.sort(key=lambda x: x["updated_at"] or "", reverse=True)
            canonical = arr[0]
            # PK 를 data.id 로 통일 위해 canonical 을 did (data.id) 로 upsert 하고 옛 PK row 삭제
            await db.req_upsert(did, {**canonical["data"], "id": did})
        # 나머지 삭제
        for x in arr:
            if x["pk"] == canonical["pk"] or (x["pk"] == did and canonical["pk"] != did):
                continue
            async with db.pool().acquire() as c:
                await c.execute("DELETE FROM req WHERE id=$1", x["pk"])
            total_removed += 1
    return total_removed

@app.on_event("shutdown")
async def _db_close():
    try:
        await db.close_pool()
    except Exception:
        pass

@app.on_event("startup")
async def _rag_warmup():
    if N2X_RELAY_ONLY:
        return
    """시작 시 임베딩 캐시(.npy)·코퍼스를 백그라운드로 미리 로드 → 첫 RAG 질의도 빠름."""
    import threading
    def _warm():
        try:
            _embed_load(); _manual_chunk_corpus()
        except Exception:
            pass
    threading.Thread(target=_warm, daemon=True).start()

def _rag_tokenize(s):
    import re as _re
    s = str(s or "").lower()
    toks = _re.findall(r"[a-z0-9]+|[가-힣]+", s)
    out = []
    for t in toks:
        out.append(t)
        if len(t) >= 2 and ('가' <= t[0] <= '힣'):
            for i in range(len(t) - 1):
                out.append(t[i:i + 2])   # 한글 2-gram → 부분일치
    return out

def _chunk_text(text, size=RAG_CHUNK_SIZE, overlap=RAG_CHUNK_OVERLAP):
    import re as _re
    text = str(text or "").strip()
    if not text:
        return []
    paras = [p.strip() for p in _re.split(r"\n\s*\n", text) if p.strip()]
    chunks, buf = [], ""
    for p in paras:
        if len(buf) + len(p) + 1 <= size:
            buf = (buf + "\n" + p).strip()
        else:
            if buf:
                chunks.append(buf)
            if len(p) <= size:
                buf = p
            else:
                step = max(1, size - overlap)
                for i in range(0, len(p), step):
                    chunks.append(p[i:i + size])
                buf = ""
    if buf:
        chunks.append(buf)
    return chunks

async def _manual_chunk_corpus_async():
    """DB 에서 manuals 를 가져와 청크 코퍼스 생성 — async 네이티브 버전."""
    # 캐시 히트 fast path — 시그니처만 조회
    async with db.pool().acquire() as c:
        sig_rows = await c.fetch("SELECT id, updated_at FROM manuals WHERE active=true ORDER BY id")
    sig_only = tuple((r["id"], r["updated_at"].timestamp() if r["updated_at"] else 0) for r in sig_rows)
    if _RAG_CACHE.get("sig") == sig_only and _RAG_CACHE.get("corpus") is not None:
        return _RAG_CACHE["corpus"]
    # 캐시 미스 → 전체 fetch
    async with db.pool().acquire() as c:
        rows_full = await c.fetch(
            "SELECT id, data, updated_at FROM manuals WHERE active=true ORDER BY id"
        )
        rows = [(r["id"], r["data"], r["updated_at"]) for r in rows_full]
    # 코퍼스 조립 (CPU 작업 — 이벤트 루프 잠깐 잡음. 매뉴얼 100개 정도면 문제 없음)
    return _build_corpus_from_rows(rows, sig_only)


def _build_corpus_from_rows(rows, sig):
    """rows 를 corpus 리스트로 변환 + 캐시 저장. CPU-only 라 sync."""
    corpus = []
    _emodel = str((_rag_cfg().get("embed_model") if callable(globals().get("_rag_cfg")) else None) or "bge-m3")
    for _rid, d, _ua in rows:
        try:
            name = d.get("name", "")
            imgs = d.get("images") or []
            _src = str(d.get("source") or "").strip().lower()
            stag = "confluence" if _src.startswith("conf") else (_src if _src in ("tc", "req", "jira", "manual") else "manual")
            for ch in _chunk_text(d.get("text", "")):
                corpus.append({"name": name, "text": ch, "tokens": _rag_tokenize(ch), "key": _embed_key(ch, _emodel), "images_ref": imgs, "source": stag})
        except Exception:
            pass
    _RAG_CACHE["sig"] = sig
    _RAG_CACHE["corpus"] = corpus
    return corpus


def _manual_chunk_corpus():
    """레거시 sync 진입점 — 워밍업 스레드 등에서만 호출. async 컨텍스트에서는 _manual_chunk_corpus_async 사용."""
    # 새 이벤트 루프 + 별도 asyncpg 연결로 조회 (풀 재사용 불가 — 다른 루프에 바인딩됨)
    async def _work():
        conn = await __import__('asyncpg').connect(dsn=db.DSN)
        try:
            sig_rows = await conn.fetch("SELECT id, updated_at FROM manuals WHERE active=true ORDER BY id")
            sig_only = tuple((r["id"], r["updated_at"].timestamp() if r["updated_at"] else 0) for r in sig_rows)
            if _RAG_CACHE.get("sig") == sig_only and _RAG_CACHE.get("corpus") is not None:
                return _RAG_CACHE["corpus"]
            rows_full = await conn.fetch("SELECT id, data, updated_at FROM manuals WHERE active=true ORDER BY id")
            rows = [(r["id"], r["data"], r["updated_at"]) for r in rows_full]
        finally:
            try: await conn.close()
            except Exception: pass
        return _build_corpus_from_rows(rows, sig_only)
    loop = asyncio.new_event_loop()
    try:
        return loop.run_until_complete(_work())
    finally:
        try: loop.close()
        except Exception: pass


def _finalize_corpus(rows, sig):   # 하위호환용 shim (사용 안 함)
    return _build_corpus_from_rows(rows, sig)
    # updated_at 시그니처로 캐시 히트 판정
    sig = tuple((rid, ua.timestamp() if ua else 0) for rid, _d, ua in rows)
    if _RAG_CACHE["sig"] == sig and _RAG_CACHE["corpus"] is not None:
        return _RAG_CACHE["corpus"]
    corpus = []
    _emodel = str((_rag_cfg().get("embed_model") if callable(globals().get("_rag_cfg")) else None) or "bge-m3")
    for _rid, d, _ua in rows:
        try:
            name = d.get("name", "")
            imgs = d.get("images") or []
            _src = str(d.get("source") or "").strip().lower()
            stag = "confluence" if _src.startswith("conf") else (_src if _src in ("tc", "req", "jira", "manual") else "manual")
            for ch in _chunk_text(d.get("text", "")):
                corpus.append({"name": name, "text": ch, "tokens": _rag_tokenize(ch), "key": _embed_key(ch, _emodel), "images_ref": imgs, "source": stag})
        except Exception:
            pass
    _RAG_CACHE["sig"] = sig
    _RAG_CACHE["corpus"] = corpus
    return corpus

def _bm25_search(query, corpus, top_k=6, k1=1.5, b=0.75):
    import math as _m
    if not corpus:
        return []
    q = set(_rag_tokenize(query))
    if not q:
        return []
    N = len(corpus)
    avgdl = sum(len(c["tokens"]) for c in corpus) / max(1, N)
    df = {}
    for c in corpus:
        for t in set(c["tokens"]):
            if t in q:
                df[t] = df.get(t, 0) + 1
    scored = []
    for c in corpus:
        dl = len(c["tokens"]) or 1
        tf = {}
        for t in c["tokens"]:
            if t in q:
                tf[t] = tf.get(t, 0) + 1
        if not tf:
            continue
        s = 0.0
        for t, f in tf.items():
            idf = _m.log(1 + (N - df.get(t, 0) + 0.5) / (df.get(t, 0) + 0.5))
            s += idf * (f * (k1 + 1)) / (f + k1 * (1 - b + b * dl / avgdl))
        scored.append((s, c))
    scored.sort(key=lambda x: x[0], reverse=True)
    return [c for s, c in scored[:top_k] if s > 0]

# ── 임베딩(bge-m3) + 리랭커(bge-reranker) 하이브리드 ──
RAG_CONFIG_FILE = DATA_DIR / "config" / "rag_config.json"
_RAG_DEFAULT_CFG = {"embed_url": "", "embed_model": "bge-m3", "rerank_url": "", "rerank_model": "bge-reranker-v2-m3", "use_embed": True, "use_rerank": True, "min_score": 0.0}

def _rag_cfg():
    try:
        if RAG_CONFIG_FILE.exists():
            d = load_json(RAG_CONFIG_FILE)
            if isinstance(d, dict):
                return {**_RAG_DEFAULT_CFG, **d}
    except Exception:
        pass
    return dict(_RAG_DEFAULT_CFG)


async def _embed_texts_raw(texts):
    """(벡터, 실패이유) 를 함께 돌려준다.

    검색 경로는 실패하면 그냥 넘어가면 되지만, 연결 테스트는 '왜' 안 되는지를
    답해야 한다. 예전에는 둘 다 None 만 받아서 화면에 '실패' 한 단어만 떴고,
    주소가 틀린 건지 서버가 죽은 건지 모델 이름이 틀린 건지 알 수 없었다.
    """
    cfg = _rag_cfg(); url = str(cfg.get("embed_url") or "").rstrip("/")
    if not url:
        return None, "서버 주소가 비어 있습니다"
    if not texts:
        return None, "보낼 내용이 없습니다"
    import httpx
    model = cfg.get("embed_model") or "bge-m3"
    out = []
    try:
        async with httpx.AsyncClient(timeout=120) as client:
            for i in range(0, len(texts), 64):
                batch = texts[i:i + 64]
                r = await client.post(url + "/v1/embeddings",
                                      json={"model": model, "input": batch})
                if r.status_code != 200:
                    return None, await _why_http(client, url, model, r)
                data = sorted(r.json().get("data") or [], key=lambda x: x.get("index", 0))
                out.extend([d.get("embedding") for d in data])
        return out, ""
    except Exception as e:
        return None, f"연결하지 못했습니다 — {e}"


async def _embed_texts(texts):
    v, _ = await _embed_texts_raw(texts)
    return v


async def _rerank_raw(query, docs, top_k):
    """(결과, 실패이유). _embed_texts_raw 와 같은 이유로 나눠 둔다."""
    cfg = _rag_cfg(); url = str(cfg.get("rerank_url") or "").rstrip("/")
    if not url:
        return None, "서버 주소가 비어 있습니다"
    if not docs:
        return None, "보낼 내용이 없습니다"
    import httpx
    model = cfg.get("rerank_model") or "bge-reranker-v2-m3"
    payload = {"model": model, "query": query, "documents": docs, "top_n": top_k}
    try:
        async with httpx.AsyncClient(timeout=120) as client:
            # 서버마다 경로가 갈린다(vLLM 은 /v1/rerank, TEI 는 /rerank).
            # 먼저 것이 아니면 두 번째로 한 번 더 본다.
            r = await client.post(url + "/v1/rerank", json=payload)
            if r.status_code != 200:
                r2 = await client.post(url + "/rerank", json=payload)
                # 둘 다 실패면 앞의 응답으로 이유를 말한다 — 보통 그쪽이
                # 진짜 경로라 메시지가 더 쓸모 있다.
                if r2.status_code == 200:
                    r = r2
                else:
                    return None, await _why_http(client, url, model, r)
            res = r.json().get("results") or []
            return [(it.get("index"), it.get("relevance_score", it.get("score", 0))) for it in res], ""
    except Exception as e:
        return None, f"연결하지 못했습니다 — {e}"


async def _rerank(query, docs, top_k):
    v, _ = await _rerank_raw(query, docs, top_k)
    return v


async def _why_http(client, url: str, model: str, r) -> str:
    """200 이 아닐 때 사람이 읽을 이유를 만든다.

    가장 흔한 실수가 모델 이름이라, 서버가 살아 있으면 /v1/models 를 물어
    '주소는 맞는데 그 모델이 없다' 를 따로 짚어 준다. 이것을 안 하면
    주소부터 다시 의심하느라 시간을 버린다.
    """
    if r.status_code in (401, 403):
        return f"인증 실패 ({r.status_code}) — 키를 확인하세요"
    try:
        m = await client.get(url + "/v1/models")
        if m.status_code == 200:
            names = [x.get("id") for x in (m.json().get("data") or []) if x.get("id")]
            if names and model not in names:
                return f"서버는 응답하는데 '{model}' 모델이 없습니다 (있는 것: {', '.join(names[:5])})"
    except Exception:
        pass
    body = (r.text or "").strip().replace("\n", " ")[:200]
    return f"응답 코드 {r.status_code}{' — ' + body if body else ''}"

async def _ensure_embeddings(corpus):
    """캐시에 없는 청크만 임베딩 후 .npy 캐시에 추가 (다음부터 재사용)."""
    _embed_load()
    todo = [i for i, c in enumerate(corpus) if c.get("key") and c["key"] not in _EMBED["keys"]]
    if not todo:
        return True
    # 캐시에 이미 있는 동일 텍스트(키)는 한 번만 임베딩
    seen = {}; uniq = []
    for i in todo:
        k = corpus[i]["key"]
        if k not in seen and k not in _EMBED["keys"]:
            seen[k] = True; uniq.append(i)
    if not uniq:
        return True
    vecs = await _embed_texts([corpus[i]["text"] for i in uniq])
    if not vecs or len(vecs) != len(uniq):
        return False
    _embed_add([corpus[i]["key"] for i in uniq], vecs)
    _embed_save()
    return True

async def _hybrid_search(query, top_k=6, min_score=None, sources=None):
    cfg = _rag_cfg()
    corpus = await _manual_chunk_corpus_async()
    if sources:   # 소스 필터: tc/req/manual/confluence/jira
        _ss = {str(s).strip().lower() for s in sources if str(s).strip()}
        if _ss:
            corpus = [c for c in corpus if c.get("source", "manual") in _ss]
    if not corpus:
        return [], "none"
    # 1) BM25 후보
    bm = _bm25_search(query, corpus, top_k=max(top_k * 4, 20))
    mode = "bm25"
    cand = list(bm)
    # 2) 임베딩 의미검색 후보 합치기
    if cfg.get("use_embed") and cfg.get("embed_url"):
        ok = await _ensure_embeddings(corpus)
        qv = (await _embed_texts([query])) if ok else None
        if ok and qv and _EMBED.get("norm") is not None:
            try:
                import numpy as np
                q = np.asarray(qv[0], dtype="float32"); qn = float(np.linalg.norm(q)) or 1.0; q = q / qn
                km = _EMBED["keys"]; rows = []; ci = []
                for i, c in enumerate(corpus):
                    r = km.get(c.get("key"))
                    if r is not None and r < len(_EMBED["norm"]):
                        rows.append(r); ci.append(i)
                if rows:
                    sims = _EMBED["norm"][rows] @ q          # numpy 벡터화 (22K도 수십 ms)
                    kk = min(max(top_k * 4, 20), len(sims))
                    top = np.argpartition(-sims, kk - 1)[:kk]; top = top[np.argsort(-sims[top])]
                    seen = set(id(c) for c in cand)
                    for t in top:
                        c = corpus[ci[int(t)]]
                        if id(c) not in seen:
                            cand.append(c); seen.add(id(c))
                    mode = "hybrid"
            except Exception:
                pass
    if not cand:
        cand = corpus[:max(top_k * 4, 20)]
    cand = cand[:max(top_k * 4, 24)]   # 리랭커 부하 제한 (후보 과다 방지)
    # 3) 리랭킹
    if cfg.get("use_rerank") and cfg.get("rerank_url") and cand:
        rr = await _rerank(query, [c["text"] for c in cand], top_k * 2)
        if rr:
            try:
                ms = float(min_score if min_score is not None else (cfg.get("min_score") or 0))
            except Exception:
                ms = 0.0
            ordered = [(cand[i], (float(sc) if sc is not None else None)) for i, sc in rr
                       if 0 <= i < len(cand) and (sc is None or float(sc) >= ms)]
            return ordered[:top_k], mode + "+rerank"
    return [(c, None) for c in cand[:top_k]], mode

@app.post("/api/rag/search")
async def rag_search(payload: dict):
    import re as _re
    q = str(payload.get("query") or "")
    k = max(1, min(int(payload.get("top_k") or 6), 20))
    _ms = payload.get("min_score")
    srcs = payload.get("sources") or ([payload.get("source")] if payload.get("source") else None)
    hits, mode = await _hybrid_search(q, top_k=k, min_score=(float(_ms) if _ms is not None else None), sources=srcs)
    out = []
    for h, sc in hits:
        txt = str(h.get("text") or "")
        imgs = []
        for mk in _re.findall(r"\[\[IMG:(\d+)\]\]", txt):
            try:
                idx = int(mk); ref = h.get("images_ref") or []
                if 0 <= idx < len(ref):
                    imgs.append(ref[idx])
            except Exception:
                pass
        clean = _re.sub(r"\[\[IMG:\d+\]\]", "", txt).strip()
        out.append({"name": h["name"], "text": clean, "images": imgs[:6], "score": (round(sc, 3) if sc is not None else None), "source": h.get("source", "manual")})
    # 라이브 Confluence 검색 결과를 앞에 합침 (Dify식 — import 불필요). confluence=False면 생략(단계별 표시용)
    if payload.get("confluence", True):
        conf = await _confluence_live_search(q, 3)
        if conf:
            out = [{"name": c["name"], "text": str(c.get("text") or "")[:600], "images": (c.get("images") or [])[:3], "url": c.get("url", "")} for c in conf] + out
            mode = mode + "+confluence"
    _corpus_tc = await _manual_chunk_corpus_async()
    return {"hits": out, "total_chunks": len(_corpus_tc), "mode": mode}

@app.post("/api/confluence/search")
async def conf_search(payload: dict):
    """라이브 Confluence 검색만 (FAB 단계별 표시용). live_query 꺼져있으면 빈 결과."""
    q = str((payload or {}).get("query") or "")
    lim = max(1, min(int((payload or {}).get("limit") or 4), 8))
    hits = await _confluence_live_search(q, lim)
    return {"hits": hits, "count": len(hits)}

_CONF_MODELS_CACHE = {"models": None}

@app.get("/api/confluence/models")
async def conf_models():
    """11.Feature List 하위 스펙 페이지에서 실제 모델명 추출 (HITL 모델 칩용, 캐시)."""
    if _CONF_MODELS_CACHE.get("models") is not None:
        return {"models": _CONF_MODELS_CACHE["models"]}
    cfg = _conf_cfg(); base = str(cfg.get("base_url") or "").rstrip("/")
    models = []
    if base:
        import httpx, re as _ri
        headers = _conf_headers(cfg); auth = _conf_auth(cfg)
        try:
            async with httpx.AsyncClient(timeout=40, verify=False) as client:
                children = await _conf_children(client, base, headers, auth, "11.Feature List")
                for c in children:
                    try:
                        hit = await client.get(base + f"/rest/api/content/{c.get('id')}", headers=headers, auth=auth, params={"expand": "body.storage"})
                        if hit.status_code != 200:
                            continue
                        body = (((hit.json().get("body") or {}).get("storage") or {}).get("value")) or ""
                        txt = _html_to_text(body)
                        m = _ri.search(r"Model\s*Name\s*[:：]\s*([^\n]+)", txt, _ri.I)
                        added = False
                        if m:
                            for tok in _ri.findall(r"[A-Za-z]{1,4}\d{3,}[A-Za-z0-9]*", m.group(1)):
                                if tok not in models:
                                    models.append(tok); added = True
                        if not added:
                            ttl = str(c.get("title") or "").replace("_Series_Spec", "")
                            if ttl and ttl not in models:
                                models.append(ttl)
                    except Exception:
                        continue
        except Exception:
            pass
    _CONF_MODELS_CACHE["models"] = models
    return {"models": models}

@app.get("/api/rag/config")
async def rag_config_get(token: str = ""):
    return _rag_cfg()

@app.post("/api/rag/config")
async def rag_config_set(payload: dict, token: str = ""):
    _require_admin(token)
    cfg = _rag_cfg()
    for k in ("embed_url", "embed_model", "rerank_url", "rerank_model"):
        if k in payload:
            cfg[k] = str(payload.get(k) or "").strip()
    for k in ("use_embed", "use_rerank"):
        if k in payload:
            cfg[k] = bool(payload.get(k))
    if "min_score" in payload:
        try:
            cfg["min_score"] = max(0.0, float(payload.get("min_score") or 0))
        except Exception:
            pass
    save_json(RAG_CONFIG_FILE, cfg)
    _RAG_CACHE["corpus"] = None  # 임베딩 재계산 유도
    return {"ok": True, **cfg}

@app.post("/api/rag/test")
async def rag_test(payload: dict, token: str = ""):
    """임베딩·리랭커 연결 테스트. 실패하면 왜 안 되는지까지 돌려준다."""
    out = {}
    ev, ew = await _embed_texts_raw(["연결 테스트", "embedding test"])
    out["embed"] = {"ok": bool(ev), "dim": (len(ev[0]) if ev else 0), "detail": ew}
    rr, rw = await _rerank_raw("테스트 질문", ["문서1 테스트", "관계없는 문서"], 2)
    out["rerank"] = {"ok": bool(rr), "results": len(rr or []), "detail": rw}
    return out

# ══════════ AI 통합: RAG 색인 API · Cycle 요약 · Fail→Jira · 통합 검색 · 자연어 실행 ══════════
AI_SETTINGS_FILE = DATA_DIR / "config" / "ai_settings.json"
RAG_INDEX_FOLDER = "AI 자동 색인"

def _ai_settings():
    """AI 자동화 스위치. auto_jira 만 기본 꺼짐(이슈 대량생성 방지) — 나머지는 로컬 LLM이라 부담 없음."""
    base = {"auto_index_tc": True, "auto_index_req": True, "auto_summary": True, "auto_jira": False}
    try:
        if AI_SETTINGS_FILE.exists():
            d = load_json(AI_SETTINGS_FILE)
            if isinstance(d, dict):
                base.update({k: bool(d[k]) for k in base if k in d})
    except Exception:
        pass
    return base

@app.get("/api/ai/settings")
async def ai_settings_get():
    return _ai_settings()

@app.post("/api/ai/settings")
async def ai_settings_set(payload: dict, token: str = ""):
    cur = _ai_settings()
    for k in cur:
        if k in (payload or {}):
            cur[k] = bool(payload[k])
    save_json(AI_SETTINGS_FILE, cur)
    return {"ok": True, **cur}

# ── 페이지별 AI(fab) 설정: Tests/Cycle/Reports 각각의 LLM + 시스템 프롬프트 (전 계정 공유) ──
PAGE_AI_FILE = DATA_DIR / "config" / "page_ai.json"
_PAGE_AI_KEYS = ("tests", "cycle", "report", "jira_ai")

_PAGE_AI_FIELDS = ("llm_id", "prompt", "greeting", "placeholder")

def _clean_rag_sources(rs):
    """지식 소스 3종(TC 절차/매뉴얼/Confluence) 활성화+우선순위 목록 정제 — dify 어시스턴트와 동일 스키마."""
    if not isinstance(rs, list):
        return []
    out = []
    for it in rs:
        if not isinstance(it, dict):
            continue
        src = str(it.get("source", "")).strip()
        if src not in ("tc", "manual", "confluence"):
            continue
        out.append({"source": src, "enabled": bool(it.get("enabled")), "priority": int(it.get("priority") or 0)})
    return out

def _clean_quick(qs):
    """추천 질문(퀵 질문 칩) 목록 정제 — 문자열 배열, 빈 값/중복 제거, 최대 20개."""
    if not isinstance(qs, list):
        return []
    out = []
    for q in qs:
        s = str(q or "").strip()
        if s and s not in out:
            out.append(s)
    return out[:20]

def _page_ai_settings():
    base = {k: {f: "" for f in _PAGE_AI_FIELDS} for k in _PAGE_AI_KEYS}
    for k in _PAGE_AI_KEYS:
        base[k]["rag_sources"] = []
        base[k]["quick"] = []
    try:
        if PAGE_AI_FILE.exists():
            d = load_json(PAGE_AI_FILE)
            if isinstance(d, dict):
                for k in _PAGE_AI_KEYS:
                    v = d.get(k)
                    if isinstance(v, dict):
                        for f in _PAGE_AI_FIELDS:
                            base[k][f] = str(v.get(f) or "")
                        base[k]["rag_sources"] = _clean_rag_sources(v.get("rag_sources"))
                        base[k]["quick"] = _clean_quick(v.get("quick"))
    except Exception:
        pass
    return base

@app.get("/api/page-ai")
async def page_ai_get():
    return _page_ai_settings()

@app.post("/api/page-ai")
async def page_ai_set(payload: dict, token: str = ""):
    cur = _page_ai_settings()
    p = payload or {}
    for k in _PAGE_AI_KEYS:
        v = p.get(k)
        if isinstance(v, dict):
            for f in _PAGE_AI_FIELDS:
                if f in v:
                    cur[k][f] = str(v.get(f) or "")
            if "rag_sources" in v:
                cur[k]["rag_sources"] = _clean_rag_sources(v.get("rag_sources"))
            if "quick" in v:
                cur[k]["quick"] = _clean_quick(v.get("quick"))
    save_json(PAGE_AI_FILE, cur)
    return {"ok": True, **cur}

def _rag_doc_id(raw):
    import re as _re
    s = _re.sub(r"[^0-9A-Za-z._-]", "_", str(raw or ""))[:120]
    return s or "doc"

def _rag_index_doc(doc_id, name, text, source, folder=RAG_INDEX_FOLDER, url=""):
    """RAG 문서 업서트 — 매뉴얼 저장소(DB.manuals) 재활용: 기존 청킹·임베딩·검색 파이프라인 그대로 탄다."""
    text = str(text or "").strip()
    if not text:
        return False
    mid = _rag_doc_id(doc_id)
    payload = {
        "id": mid, "name": str(name or mid), "text": text, "chars": len(text),
        "source": str(source or "manual"), "active": True, "folder": folder, "url": url,
        "created_at": datetime.now().strftime("%Y-%m-%d"),
    }
    # sync → async 브리지
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            asyncio.run_coroutine_threadsafe(db.manuals_upsert(mid, payload), loop).result(timeout=10)
        else:
            asyncio.run(db.manuals_upsert(mid, payload))
    except RuntimeError:
        asyncio.run(db.manuals_upsert(mid, payload))
    _RAG_CACHE["corpus"] = None
    return True

async def _rag_index_warm():
    """색인 후 새 청크 임베딩을 백그라운드로 미리 계산 (첫 검색 지연 방지)."""
    try:
        _corpus_w = await _manual_chunk_corpus_async()
        await _ensure_embeddings(_corpus_w)
    except Exception:
        pass

@app.post("/api/rag/index")
async def rag_index(payload: dict, token: str = ""):
    """RAG 색인 추가/갱신 — {id?, name, text, source?, folder?, url?}. 같은 id 재전송 = 덮어쓰기(업서트)."""
    name = str((payload or {}).get("name") or "").strip()
    text = str((payload or {}).get("text") or "").strip()
    if not text:
        return {"ok": False, "error": "text가 비어 있습니다"}
    source = str(payload.get("source") or "manual").strip().lower()
    doc_id = str(payload.get("id") or "").strip() or ("rag-" + _embed_key(source + "|" + (name or text[:80]), "id")[:16])
    ok = _rag_index_doc(doc_id, name or doc_id, text, source, folder=str(payload.get("folder") or RAG_INDEX_FOLDER), url=str(payload.get("url") or ""))
    if ok:
        asyncio.create_task(_rag_index_warm())
    return {"ok": ok, "id": doc_id, "source": source}

# ── 로컬 LLM(gemma) 공용 호출 ──
def _ai_llm(llm_id: str = ""):
    """AI 통합 기능용 LLM 선택 — 로컬(vLLM/OpenAI 호환) 우선. claude 타입은 스키마가 달라 제외.

    `llm_id` 는 화면에서 사람이 고른 것이다(지시: 탭마다 드롭바). 고른 것이
    있으면 그것을 쓴다 — 목록에서 사라졌으면 여느 때처럼 고른다.
    """
    init_llms_file()
    llms = (load_json(LLMS_FILE).get("llms") or [])
    if llm_id:
        got = next((l for l in llms if str(l.get("id") or "") == llm_id), None)
        if got and got.get("endpoint"):
            return got
    act = [l for l in llms if l.get("status", "active") == "active" and l.get("endpoint") and str(l.get("type") or "").lower() != "claude"]
    loc = [l for l in act if str(l.get("type") or "").lower() == "local"]
    return (loc[0] if loc else (act[0] if act else None))

async def _ai_chat(messages, max_tokens=1800, temperature=0.3, json_schema=None, timeout=180,
                   llm_id: str = ""):
    """OpenAI 호환 chat/completions 1회 호출 → (content, error). json_schema 지정 시 vLLM guided_json."""
    llm = _ai_llm(llm_id)
    if not llm:
        return None, "등록된 로컬 LLM이 없습니다 — AI Assistant ▸ LLM 설정에서 등록하세요."
    import httpx
    body = {"model": llm.get("model") or "", "messages": messages, "temperature": temperature, "max_tokens": max_tokens}
    if json_schema:
        body["guided_json"] = json_schema
    headers = {"Content-Type": "application/json"}
    _key = str(llm.get("apikey") or "")
    if _key and not _key.startswith("http"):   # 일부 등록건은 apikey 칸에 URL이 들어있음 — 방어
        headers["Authorization"] = f"Bearer {_key}"
    url = str(llm["endpoint"]).rstrip("/") + "/chat/completions"
    try:
        async with httpx.AsyncClient(timeout=timeout) as client:
            r = await client.post(url, headers=headers, json=body)
            if r.status_code != 200 and json_schema:   # guided_json 미지원 → json_object 폴백
                body.pop("guided_json", None)
                body["response_format"] = {"type": "json_object"}
                r = await client.post(url, headers=headers, json=body)
            if r.status_code != 200:
                return None, f"LLM {r.status_code}: {r.text[:200]}"
            content = (((r.json().get("choices") or [{}])[0].get("message") or {}).get("content") or "").strip()
            return content, None
    except Exception as e:
        return None, str(e)[:200]

def _ai_json(content):
    """LLM 응답에서 JSON 오브젝트 추출."""
    import re as _re
    try:
        m = _re.search(r"\{.*\}", str(content or ""), _re.DOTALL)
        return json.loads(m.group(0) if m else content)
    except Exception:
        return None

# ── Cycle → Markdown 변환 · RAG 색인 · AI 요약 ──
def _step_verdict(s):
    """스텝 하나의 판정 — 화면(types.ts stepVerdict)과 **글자 하나까지 같게**.

    실행기는 `status`(PASS/FAIL)·`repeatResult`(Pass/Fail) 에 적고, 옛 자료와
    손 입력은 `result` 에 있다.

    아는 판정만 판정으로 친다. `repeatResult` 에는 「실행완료」 처럼 판정이
    아닌 말도 들어 있어서, 그것을 그대로 판정으로 읽으면 셈이 어긋난다.
    """
    legacy = str(s.get("result") or "").strip()
    if legacy:
        return _norm_verdict(legacy)
    # 화면은 `status ?? repeatResult` — status 가 아예 없을 때만 옛 칸을 본다
    v = s.get("status")
    if v is None:
        v = s.get("repeatResult")
    u = str(v or "").strip().upper()
    return u if u in ("PASS", "FAIL", "WIP", "BLOCKED") else ""


def _norm_verdict(v):
    """같은 판정을 부르는 이름이 여럿이다 — 한 가지로 모은다."""
    u = str(v or "").strip().upper()
    if u in ("PASS", "합격"):
        return "PASS"
    if u in ("FAIL", "불합격"):
        return "FAIL"
    if u in ("N/A", "NA", "진행불가"):
        return "N/A"
    return str(v or "").strip()


def _item_verdict(item):
    """항목 판정 — **화면(Cycles.tsx itemVerdict)과 같은 규칙**.

    여태 이 함수만 항목의 `status` 칸을 함께 읽었다. 화면은 `result` 만 본다.
    그래서 한 번 깨진 뒤 다시 돌린 항목처럼 두 칸이 어긋난 자료에서, **상세는
    Pass 인데 목록의 「최근 결과」는 Fail** 이 됐다(지적). 판정을 두 규칙으로
    내리면 둘 중 하나는 반드시 틀린다 — 화면 쪽으로 맞춘다.

    `미실행` 은 값이 아니라 표식이다. 사람이 「이건 안 돌린 것으로 둬라」 고
    적어 둔 것이라, 스텝에 결과가 남아 있어도 판정으로 올리지 않는다.
    """
    r = str(item.get("result") or "").strip()
    if r == "미실행":
        return ""
    if r:
        return _norm_verdict(r)
    steps = item.get("steps") or []
    # 수동 스텝은 자동 판정에서 뺀다 — 사람이 보는 것은 사람이 따로 적는다
    auto = [
        s for s in steps
        if isinstance(s, dict) and not (s.get("manual") or s.get("action") == "수동")
    ]
    if not auto:
        return "N/A" if steps else ""
    rs = [_step_verdict(s) for s in auto]
    if len(auto) == 1:
        return rs[0]
    if any(x == "FAIL" for x in rs):          # Fail 하나라도 → Fail
        return "FAIL"
    if any(x == "PASS" for x in rs):          # Fail 없고 Pass 있으면 → Pass
        return "PASS"
    # 메시지·주석처럼 판정이 없는 줄뿐이면 그 줄의 말을 그대로 (없으면 미실행)
    return next((x for x in rs if x), "")


def _tc_item_md(item, cycle):
    """Cycle 항목(스텝 포함)을 검색용 Markdown으로 직렬화."""
    lines = [f"# [TC] {item.get('tcid','')} {item.get('name','')}",
             f"- 모델: {cycle.get('model','')} / 버전: {cycle.get('version','')} ({cycle.get('version_group','')})",
             f"- 결과: {item.get('result') or item.get('status') or '-'} · 실행 {str(item.get('last_run') or '')[:16]}", ""]
    for i, s in enumerate(item.get("steps") or [], 1):
        if not isinstance(s, dict):
            continue
        lines.append(f"## Step {i}. {s.get('desc') or ''}")
        if s.get("cli"):
            lines.append(f"- CLI: `{s['cli']}`")
        if s.get("criteria"):
            lines.append(f"- 판정({s.get('type','')}): {s.get('criteria')}")
        out = str(s.get("output") or "").strip()
        if out:
            lines.append("```\n" + out[:800] + "\n```")
    return "\n".join(lines)

async def _rag_index_cycle(cycle):
    """Cycle 완료 → Pass 항목 절차를 RAG(source=tc)로 자동 색인. 같은 TC는 최신 실행으로 덮어씀."""
    n = 0
    for it in (cycle.get("items") or []):
        if _item_verdict(it) != "PASS":
            continue
        tcid = str(it.get("tcid") or "").strip()
        if not tcid:
            continue
        if _rag_index_doc("ragtc-" + tcid, f"[TC] {tcid} {it.get('name','')}", _tc_item_md(it, cycle), "tc"):
            n += 1
    if n:
        await _rag_index_warm()
    return n

def _req_to_md(req):
    """REQ 문서를 검색용 Markdown으로 직렬화."""
    lines = [f"# [REQ] {req.get('reqid') or req.get('id','')} {req.get('title') or req.get('name','')}"]
    meta = []
    for k, lb in (("folder", "폴더"), ("status", "상태"), ("priority", "우선순위")):
        if req.get(k):
            meta.append(f"{lb}: {req[k]}")
    prods = req.get("products")
    if isinstance(prods, list) and prods:
        meta.append("제품군: " + ", ".join(str(p) for p in prods))
    if meta:
        lines.append("- " + " · ".join(meta))
    for k, lb in (("overview", "개요"), ("object", "목적"), ("desc", "설명")):
        v = str(req.get(k) or "").strip()
        if v:
            lines.append(f"\n## {lb}\n{v}")
    scen = req.get("scenarios")
    if isinstance(scen, list) and scen:
        lines.append("\n## 시험 시나리오")
        for i, s in enumerate(scen, 1):
            if isinstance(s, dict):
                lines.append(f"{i}. " + " / ".join(str(v) for v in s.values() if v and isinstance(v, (str, int, float))))
            elif s:
                lines.append(f"{i}. {s}")
    cf = req.get("custom_fields")
    if isinstance(cf, dict) and cf:
        lines.append("\n## 추가 필드")
        for k, v in cf.items():
            if v:
                lines.append(f"- {k}: {v}")
    tcs = [t for t in (req.get("tc") or []) if isinstance(t, dict) and t.get("tcid")]
    if tcs:
        lines.append("\n## 연결 TC")
        for t in tcs:
            lines.append(f"- {t.get('tcid')} {t.get('name','')} ({t.get('status','')})")
    return "\n".join(lines)

async def _rag_index_req(req_id, req_data):
    # _rag_index_doc 은 sync 함수인데 그 안에서 run_coroutine_threadsafe(...).result(timeout=10) 을
    # 호출한다. 이 호출을 현재 async 이벤트 루프 안에서 그대로 실행하면 자기 자신 blocking 대기가
    # 발생해 다른 API 요청까지 몇 초 이상 지연된다. → 별도 스레드로 위임해서 이벤트 루프에서 분리.
    try:
        doc_id = "ragreq-" + req_id
        name = f"[REQ] {req_data.get('reqid') or req_id} {req_data.get('title','')}"
        text = _req_to_md(req_data)
        # sync 호출을 별도 스레드로 실행 → 이벤트 루프 blocking 방지
        ok = await asyncio.get_event_loop().run_in_executor(
            None, _rag_index_doc, doc_id, name, text, "req"
        )
        if ok:
            await _rag_index_warm()
    except Exception:
        pass

def _cycle_result_ctx(cycle, fails_detail=True):
    """Cycle 실행 결과를 LLM 컨텍스트 문자열로 요약."""
    items = cycle.get("items") or []
    p = f = na = un = 0
    rows = []
    for it in items:
        v = _item_verdict(it)
        if v == "PASS":
            p += 1
        elif v == "FAIL":
            f += 1
        elif v in ("N/A", "NA"):
            na += 1
        else:
            un += 1
        rows.append(f"- {it.get('tcid','')} {it.get('name','')}: {v or '미실행'}" + (f" — {it.get('memo')}" if it.get("memo") else ""))
    head = (f"Cycle: {cycle.get('id','')} / 모델 {cycle.get('model','')} / 버전 {cycle.get('version','')}\n"
            f"실행일시: {cycle.get('executed_at','')}\n"
            f"전체 {len(items)} · Pass {p} · Fail {f} · N/A {na} · 미실행 {un}\n\n[항목별 결과]\n" + "\n".join(rows))
    if fails_detail:
        fd = []
        for it in items:
            if _item_verdict(it) != "FAIL":
                continue
            fd.append(f"\n■ FAIL: {it.get('tcid','')} {it.get('name','')}")
            for i, s in enumerate(it.get("steps") or [], 1):
                if not isinstance(s, dict) or str(s.get("result") or "").strip().upper() not in ("FAIL", "불합격"):
                    continue
                fd.append(f"  - Step{i} {s.get('desc','')} / CLI `{s.get('cli','')}` / 기대({s.get('type','')}): {s.get('criteria','')}")
                out = str(s.get("output") or "").strip()
                if out:
                    fd.append("    출력(끝부분): " + out[-400:].replace("\n", " ⏎ "))
        if fd:
            head += "\n\n[Fail 상세]" + "\n".join(fd)
    return head

async def _cycle_ai_summary(cycle_id, llm_id: str = ""):
    """Gemma로 Cycle 요약 생성 → cycle 의 ai_summary 에 저장."""
    cycle = await db.cycle_get(cycle_id)
    if cycle is None:
        return None, "Cycle을 찾을 수 없습니다"
    ctx = _cycle_result_ctx(cycle)
    # 전체·수동·자동 집계를 앞에 실어 준다 — 화면 요약 바와 같은 축으로 분석하게
    def _grp_of(it):
        st2 = [x for x in (it.get("steps") or []) if isinstance(x, dict)]
        auto2 = [x for x in st2 if not (x.get("manual") or x.get("action") == "수동")]
        return "자동" if auto2 else "수동"
    _tly = {"전체": {}, "수동": {}, "자동": {}}
    for _it in (cycle.get("items") or []):
        if not isinstance(_it, dict):
            continue
        _v = _item_verdict(_it) or "미실행"
        for _k in ("전체", _grp_of(_it)):
            _tly[_k][_v] = _tly[_k].get(_v, 0) + 1
    def _tline(k):
        d2 = _tly[k]
        n2 = sum(d2.values())
        body = " · ".join(f"{a} {b}건" for a, b in sorted(d2.items(), key=lambda x: -x[1]))
        return f"{k} {n2}건 — {body or '없음'}"
    # 제품 정보 — 실행 화면 정보 상자와 같은 축 (카탈로그의 제조사·제품군 보강)
    _pi_vendor = _pi_family = ""
    _pi_mg = str(cycle.get("model_group") or "").strip()
    try:
        for _c in await db.catalog_list("model"):
            if _c.get("name") == cycle.get("model"):
                _pi_vendor = str(_c.get("vendor") or "").strip()
                _pi_family = str(_c.get("family") or "").strip()
                if not _pi_mg:
                    _pi_mg = str(_c.get("model_group") or "").strip()
                break
    except Exception:
        pass
    _pinfo = " / ".join(
        f"{k} {v}" for k, v in [
            ("제조사", _pi_vendor), ("제품군", _pi_family), ("모델그룹", _pi_mg),
            ("제품명", str(cycle.get("model") or "").strip()),
            ("버전그룹", str(cycle.get("version_group") or "").strip()),
            ("버전명", str(cycle.get("version") or "").strip()),
            ("사이클", (str(cycle.get("cid") or "") + " " + str(cycle.get("name") or "")).strip()),
        ] if v
    )
    ctx = ("[제품 정보] " + _pinfo + "\n"
           + "[전체·수동·자동 집계]\n" + "\n".join(_tline(k) for k in ("전체", "수동", "자동")) + "\n\n" + ctx)
    # 이 글의 말투·차례는 설정(용도별 프롬프트 · Cycle-Test Summary)이 정한다.
    # 비워 두면 아래 기본 규칙을 그대로 쓴다.
    _slot = _prompt_of("cycle_summary").get("system") or ""
    sys_p = _slot + "\n" if _slot else ""
    sys_p += ("너는 네트워크 장비 시험(QA) 결과 분석 전문가다. 아래 Test Cycle 실행 결과를 근거로 한국어 Markdown 분석 보고서를 작성한다. "
             "보고서 첫 줄에 제품 정보(제조사·제품군·제품명·버전)를 한 줄로 쓰고, "
             "구성: ## 총평(2~3문장 — 이 제품·버전 회차의 품질 판단) / ## 전체·수동·자동 현황(집계를 표로 정리하고 눈에 띄는 점 1~2문장) / "
             "## Fail 분석(항목별 원인 추정과 근거 — 출력·판정기준 인용) / ## 권고사항(재시험·설정확인 등 구체적으로). "
             "결과에 없는 내용은 추측하지 말고, Fail이 없으면 Fail 분석은 '해당 없음'으로 쓴다.")
    ans, err = await _ai_chat([{"role": "system", "content": sys_p}, {"role": "user", "content": ctx}],
                              max_tokens=1600, llm_id=llm_id)
    if err:
        return None, err
    llm = _ai_llm(llm_id) or {}
    cycle = await db.cycle_get(cycle_id)   # 재로드(요약 생성 동안의 변경 보존)
    cycle["ai_summary"] = {"text": ans, "at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"), "model": llm.get("model", "")}
    await db.cycle_upsert(cycle_id, cycle)
    try:
        await broadcast({"type": "cycle_ai_summary", "cycle_id": cycle_id})
    except Exception:
        pass
    return cycle["ai_summary"], None

_cb_run_state = None   # 마지막 실행 진행 상태 — 새로고침으로 재접속한 브라우저가 GET으로 복원

@app.post("/api/cycle-run-progress")
async def cycle_run_progress(payload: dict):
    """클라이언트 주도 Cycle 자동 실행 진행 상태를 전 접속자에게 중계 (WebSocket broadcast).
    다른 계정 브라우저도 실행 배너·진행 중 오버레이·결과 갱신을 실시간으로 받는다.
    주의: /api/cycle/{cycle_id} POST(먼저 등록)가 /api/cycle/run-progress 를 가로채므로 경로를 분리했다."""
    global _cb_run_state
    msg = {"type": "cb_run_progress"}
    for k in ("evt", "ids", "key", "name", "done", "total", "user", "stepIdx", "stepCnt", "stepName", "stepAction", "stepOutput", "stepResult"):
        if k in (payload or {}):
            msg[k] = payload[k]
    if msg.get("evt") == "done":
        _cb_run_state = None
    else:
        st = {k: v for k, v in msg.items() if k != "type"}
        st["_at"] = _t.time()
        _cb_run_state = st
    try:
        await broadcast(msg)
    except Exception:
        pass
    return {"ok": True}

@app.get("/api/dashboard")
async def dashboard_data():
    """대시보드 집계 — 위젯 전부를 한 번에. 사이클은 요약본(data_summary)만 읽어 가볍다."""
    from datetime import datetime as _dt, timedelta as _td
    devices = await db.device_list()
    meters = [d for d in devices if str(d.get("role") or "") == "계측기"]
    dev_groups = {}
    for d in devices:
        role = str(d.get("role") or "")
        if role == "계측기":
            continue
        g = role or "기타"
        dev_groups[g] = dev_groups.get(g, 0) + 1
    try:
        defects = await db.defect_list()
    except Exception:
        defects = []
    _closed = ("closed", "resolved", "done", "완료", "해결", "닫힘")
    opened = [x for x in defects if str(x.get("status") or "").strip().lower() not in _closed]
    wk = (_dt.now() - _td(days=7)).strftime("%Y-%m-%d")
    week_new = sum(1 for x in opened if str(x.get("created_at") or "")[:10] >= wk)
    metas = await db.cycle_list_meta()
    today = _dt.now().strftime("%Y-%m-%d")
    yday = (_dt.now() - _td(days=1)).strftime("%Y-%m-%d")
    days = [(_dt.now() - _td(days=i)).strftime("%Y-%m-%d") for i in range(13, -1, -1)]
    daily = {d2: {"runs": 0, "ok": 0, "bad": 0} for d2 in days}
    versions = []

    def _vd(it):
        """항목 결과를 화면 말로 — Blocked·WIP·커스텀도 살린다 (_item_verdict 는 PASS/FAIL 만)"""
        raw = str(it.get("result") or "").strip()
        if raw and raw != "미실행":
            return raw
        v0 = _item_verdict(it)
        return {"PASS": "Pass", "FAIL": "Fail", "N/A": "진행불가"}.get(v0, "")

    overall = {}
    attention = []
    latest_by_tc = {}
    for c in metas:
        items = [x for x in (c.get("items") or []) if isinstance(x, dict)]
        ok2 = bad = done = 0
        for it in items:
            v = _item_verdict(it)
            lb = _vd(it)
            overall[lb or "미실행"] = overall.get(lb or "미실행", 0) + 1
            at2 = str(it.get("executed_at") or "")
            tid = str(it.get("tcid") or "")
            if tid and lb:
                cur0 = latest_by_tc.get(tid)
                if not cur0 or at2 >= cur0[0]:
                    latest_by_tc[tid] = (at2, lb)
            if lb in ("Fail", "Blocked", "진행불가"):
                attention.append({
                    "tcid": tid, "name": str(it.get("name") or ""), "label": lb,
                    "cycle_id": str(c.get("id") or ""), "version": str(c.get("version") or c.get("cid") or ""),
                    "at": at2,
                })
            if v == "PASS":
                ok2 += 1
            elif v == "FAIL":
                bad += 1
            if v:
                done += 1
            d3 = at2[:10]
            if d3 in daily:
                daily[d3]["runs"] += 1
                if v == "PASS":
                    daily[d3]["ok"] += 1
                elif v == "FAIL":
                    daily[d3]["bad"] += 1
        versions.append({
            "id": c.get("id"), "cid": str(c.get("cid") or ""),
            "version": str(c.get("version") or ""), "name": str(c.get("name") or ""),
            "updated": str(c.get("_updated_at_pg") or ""),
            "total": len(items), "ok": ok2, "bad": bad, "done": done,
        })
    versions.sort(key=lambda x: x.get("updated") or "", reverse=True)
    run = None
    st2 = _cb_run_state
    if st2 and _t.time() - st2.get("_at", 0) <= 1800:
        run = {k: st2.get(k) for k in ("key", "name", "done", "total", "user")}
    # 자산 현황 + 자동화율 — TC 메타에서
    try:
        tcs = await db.tc_list_meta()
    except Exception:
        tcs = []
    try:
        reqs_n = len(await db.req_list_full())
    except Exception:
        reqs_n = 0
    auto_n = sum(1 for t2 in tcs if str(t2.get("kind") or t2.get("run_type") or "").strip() == "자동")
    # 자주 깨지는 TC — 모든 회차의 Fail 을 tcid 로 센다
    fail_by = {}
    for c in metas:
        for it in (c.get("items") or []):
            if not isinstance(it, dict):
                continue
            v = _item_verdict(it)
            tid = str(it.get("tcid") or "")
            if not tid or not v:
                continue
            rec = fail_by.setdefault(tid, {"tcid": tid, "name": str(it.get("name") or ""), "fails": 0, "runs": 0})
            rec["runs"] += 1
            if v == "FAIL":
                rec["fails"] += 1
    top_fail = sorted((x for x in fail_by.values() if x["fails"] > 0), key=lambda x: -x["fails"])[:5]
    attention.sort(key=lambda x: x.get("at") or "", reverse=True)
    # 요구사항 커버리지 — TC 가 가리키는 요구사항 / 전체
    try:
        reqs_all = await db.req_list_full()
        req_ids = set()
        for r2 in reqs_all:
            for k2 in ("id", "reqid"):
                v2 = str(r2.get(k2) or "").strip()
                if v2:
                    req_ids.add(v2)
        covered = {str(t2.get("req_id") or "").strip() for t2 in tcs} & req_ids
        coverage = {"total": len(reqs_all), "covered": len({str(r2.get("id") or r2.get("reqid") or "") for r2 in reqs_all if (str(r2.get("id") or "").strip() in covered) or (str(r2.get("reqid") or "").strip() in covered)})}
    except Exception:
        coverage = {"total": 0, "covered": 0}
    # TC 실행 현황 — 최근 결과 기준(항목당 마지막 회차)
    tc_ids_all = {str(t2.get("tcid") or "") for t2 in tcs if t2.get("tcid")}
    exec_pass = sum(1 for tid, (_, lb) in latest_by_tc.items() if tid in tc_ids_all and lb == "Pass")
    exec_fail = sum(1 for tid, (_, lb) in latest_by_tc.items() if tid in tc_ids_all and lb == "Fail")
    exec_n = sum(1 for tid in latest_by_tc if tid in tc_ids_all)
    tcexec = {"total": len(tc_ids_all), "executed": exec_n, "passed": exec_pass, "failed": exec_fail}
    recent_defects = [
        {k: str(x.get(k) or "") for k in ("id", "title", "severity", "status", "created_at", "cycle_id", "tcid")}
        for x in opened[:5]
    ]
    return {
        "devices": {"total": len(devices) - len(meters), "groups": dev_groups},
        "meters": {"total": len(meters)},
        "defects": {"open": len(opened), "week_new": week_new},
        "today": daily.get(today, {"runs": 0, "ok": 0}),
        "daily": [{"date": d2, **daily[d2]} for d2 in days],
        "versions": versions[:8],
        "running": run,
        "assets": {"reqs": reqs_n, "tcs": len(tcs), "cycles": len(metas)},
        "overall": overall,
        "yday_runs": daily.get(yday, {}).get("runs", 0) if yday in daily else 0,
        "attention": attention[:6],
        "coverage": coverage,
        "tcexec": tcexec,
        "automation": {"auto": auto_n, "manual": max(len(tcs) - auto_n, 0)},
        "top_fail": top_fail,
        "recent_defects": recent_defects,
    }

@app.get("/api/cycle-run-progress")
async def cycle_run_progress_get():
    """진행 중인 실행 상태 조회 — 새로고침한 접속자가 배너·오버레이를 복원할 때 사용."""
    st = _cb_run_state
    if st and _t.time() - st.get("_at", 0) > 1800:   # 30분 무갱신 → 중단으로 간주
        st = None
    return {"ok": True, "state": st}

@app.post("/api/cycle-run-stop")
async def cycle_run_stop(payload: dict = None, token: str = ""):
    """다른 사용자가 실행 중인 Cycle 자동 실행을 원격으로 중지 요청.
    실행자(브라우저)에게 WebSocket으로 stop 신호 전송 → 실행자 프론트가 tcCheckRunStop 호출.
    누가 요청했는지도 함께 브로드캐스트되어 로그·배너에 표시된다."""
    global _cb_run_state
    u = _user_from_token(token) if token else None
    requester = (u.get("name") or u.get("username")) if u else ""
    reason = str((payload or {}).get("reason") or "")
    # 실행 상태 즉시 완료 처리 (다른 시청자도 오버레이 해제)
    _cb_run_state = None
    msg = {"type": "cb_run_stop_request", "user": requester, "reason": reason}
    try:
        await broadcast(msg)
    except Exception:
        pass
    # done 이벤트도 함께 브로드캐스트 → 오버레이/배너 즉시 정리
    done_msg = {"type": "cb_run_progress", "evt": "done", "user": requester}
    try:
        await broadcast(done_msg)
    except Exception:
        pass
    return {"ok": True}

@app.post("/api/cycle/{cycle_id}/summarize")
async def cycle_summarize(cycle_id: str, payload: dict = None):
    # 누구에게 맡길지 화면이 고른다(지시) — 안 고르면 여느 때처럼
    summ, err = await _cycle_ai_summary(cycle_id, str((payload or {}).get("llm") or ""))
    if err:
        return {"ok": False, "error": err}
    return {"ok": True, "summary": summ}

async def _on_cycle_complete(cycle_id, cycle):
    """engine.run_cycle 완료 훅 — RAG 자동 색인(S: TC절차) + 자동 요약(S4). 실패해도 본 실행에 영향 없음."""
    st = _ai_settings()
    if st.get("auto_index_tc"):
        try:
            await _rag_index_cycle(cycle)
        except Exception:
            pass
    if st.get("auto_summary"):
        try:
            await _cycle_ai_summary(cycle_id)
        except Exception:
            pass

engine.on_cycle_complete = _on_cycle_complete

@app.post("/api/cycle/{cycle_id}/auto-jira")
async def cycle_auto_jira(cycle_id: str, payload: dict):
    """Fail 항목 → Gemma로 이슈 제목/설명 생성 → Jira 일괄 등록. {project, issuetype, tcids?, dry_run?}"""
    cycle = await db.cycle_get(cycle_id)
    if cycle is None:
        return {"ok": False, "error": "Cycle을 찾을 수 없습니다"}
    project = str((payload or {}).get("project") or "").strip()
    itype = str((payload or {}).get("issuetype") or "Bug").strip()
    only = set((payload or {}).get("tcids") or [])
    dry = bool((payload or {}).get("dry_run"))
    if not dry and not project:
        return {"ok": False, "error": "Jira 프로젝트 키가 필요합니다"}
    fails = [it for it in (cycle.get("items") or []) if _item_verdict(it) == "FAIL" and (not only or str(it.get("tcid") or "") in only)]
    if not fails:
        return {"ok": True, "issues": [], "message": "Fail 항목이 없습니다"}
    schema = {"type": "object", "properties": {"summary": {"type": "string"}, "description": {"type": "string"}}, "required": ["summary", "description"]}
    sys_p = ("너는 QA 엔지니어다. 네트워크 장비 시험 Fail 결과로 Jira 이슈를 작성한다. "
             "summary는 [모델/버전] 증상 한 줄(80자 이내). description은 Jira wiki 마크업으로 "
             "h3.환경 / h3.재현 절차(실행 CLI 순서) / h3.기대 결과 / h3.실제 결과(출력 인용 {code}...{code}) 구성. "
             "결과에 없는 내용은 추측하지 않는다. JSON {\"summary\":...,\"description\":...} 만 출력한다.")
    issues = []
    for it in fails:
        fctx = (f"모델 {cycle.get('model','')} / 버전 {cycle.get('version','')}\nTC: {it.get('tcid','')} {it.get('name','')}\n메모: {it.get('memo','')}\n[스텝]\n")
        for i, s in enumerate(it.get("steps") or [], 1):
            if not isinstance(s, dict):
                continue
            fctx += f"{i}. {s.get('desc','')} / CLI `{s.get('cli','')}` / 기대({s.get('type','')}): {s.get('criteria','')} / 결과: {s.get('result','')}\n"
            out = str(s.get("output") or "").strip()
            if out and str(s.get("result") or "").strip().upper() in ("FAIL", "불합격"):
                fctx += "   출력:\n" + out[-600:] + "\n"
        content, err = await _ai_chat([{"role": "system", "content": sys_p}, {"role": "user", "content": fctx}], max_tokens=1400, json_schema=schema)
        obj = _ai_json(content) if content else None
        summary = (obj or {}).get("summary") or f"[{cycle.get('model','')}/{cycle.get('version','')}] {it.get('tcid','')} {it.get('name','')} 시험 Fail"
        desc = (obj or {}).get("description") or ("h3.환경\n" + f"모델 {cycle.get('model','')} / 버전 {cycle.get('version','')}\n\nh3.실패 내역\n" + fctx[:3000])
        desc += f"\n\n----\n자동 등록: utop Cycle {cycle_id}"
        rec = {"tcid": it.get("tcid", ""), "summary": summary}
        if dry:
            rec["description"] = desc
        else:
            res = await jira_create_issue({"project": project, "issuetype": itype, "summary": summary, "description": desc, "labels": ["utop-auto"]})
            rec.update({"ok": bool(res.get("ok")), "key": res.get("key", ""), "url": res.get("url", ""), "error": res.get("error", "")})
        issues.append(rec)
    if not dry:
        cycle = await db.cycle_get(cycle_id)
        cycle.setdefault("auto_jira", []).extend([{"tcid": r["tcid"], "key": r.get("key", ""), "at": datetime.now().strftime("%Y-%m-%d %H:%M")} for r in issues if r.get("ok")])
        await db.cycle_upsert(cycle_id, cycle)
    return {"ok": True, "issues": issues, "dry_run": dry}

# ── S7: 통합 검색 오케스트레이터 (RAG + Confluence + Jira → Rerank → 단일 답변) ──
@app.post("/api/ai/search-all")
async def ai_search_all(payload: dict):
    q = str((payload or {}).get("query") or "").strip()
    if not q:
        return {"ok": False, "error": "query가 비어 있습니다"}
    top_k = max(1, min(int(payload.get("top_k") or 8), 16))

    async def _jira_task():
        try:
            _q = q.replace('"', ' ').strip()
            r, err = await asyncio.to_thread(_jira_call, "GET", "/rest/api/2/search",
                                             params={"jql": f'text ~ "{_q}" ORDER BY updated DESC',
                                                     "maxResults": 5, "fields": "summary,status,issuetype,resolution"})
            if err or not r.is_success:
                return []
            out = []
            cfg = _jira_cfg()
            for iss in (r.json().get("issues") or []):
                fl = iss.get("fields") or {}
                out.append({"source": "jira", "name": iss.get("key", ""),
                            "text": f"[{((fl.get('issuetype') or {}).get('name',''))} · {((fl.get('status') or {}).get('name',''))}] {fl.get('summary','')}",
                            "url": (str(cfg.get('url') or '').rstrip('/') + "/browse/" + iss.get("key", ""))})
            return out
        except Exception:
            return []

    async def _rag_task():
        try:
            hits, _m = await _hybrid_search(q, top_k=top_k, sources=payload.get("sources"))
            return [{"source": h.get("source", "manual"), "name": h.get("name", ""), "text": str(h.get("text") or ""), "url": ""} for h, _sc in hits]
        except Exception:
            return []

    async def _conf_task():
        try:
            conf = await _confluence_live_search(q, 3)
            return [{"source": "confluence", "name": c.get("name", ""), "text": str(c.get("text") or "")[:800], "url": c.get("url", "")} for c in (conf or [])]
        except Exception:
            return []

    rag_hits, conf_hits, jira_hits = await asyncio.gather(_rag_task(), _conf_task(), _jira_task())
    cand = rag_hits + conf_hits + jira_hits
    if not cand:
        return {"ok": True, "answer": "관련 자료를 찾지 못했습니다. RAG 색인·Confluence·Jira 연동 상태를 확인하세요.", "sources": []}
    # 전체 후보 통합 리랭킹 (소스 간 순위 정렬)
    rr = await _rerank(q, [c["text"][:1200] for c in cand], min(top_k, len(cand)))
    if rr:
        cand = [cand[i] for i, _sc in rr if 0 <= i < len(cand)]
    cand = cand[:top_k]
    ctx = ""
    for i, c in enumerate(cand, 1):
        ctx += f"\n[{i}] ({c['source']}) {c['name']}\n{c['text'][:900]}\n"
    sys_p = ("너는 사내 시험 지식 통합 어시스턴트다. 아래 검색 자료([번호] (소스) 제목)만 근거로 한국어로 답한다. "
             "문장 근거에 [1][2] 형태로 출처 번호를 붙이고, 마지막 줄에 '출처: [1] 이름, [2] 이름' 목록을 쓴다. "
             "자료에 없는 내용은 추측하지 말고 없다고 답한다.")
    ans, err = await _ai_chat([{"role": "system", "content": sys_p}, {"role": "user", "content": f"[검색 자료]{ctx}\n\n[질문] {q}"}], max_tokens=1600)
    if err:
        return {"ok": False, "error": err, "sources": [{"n": i + 1, "source": c["source"], "name": c["name"], "url": c.get("url", "")} for i, c in enumerate(cand)]}
    return {"ok": True, "answer": ans, "sources": [{"n": i + 1, "source": c["source"], "name": c["name"], "url": c.get("url", "")} for i, c in enumerate(cand)]}

# ── S6: 자연어 시험 진행 (실험적) — 자연어 → CLI 변환 → 실행 → 결과 판정 ──
def _nl_cmd_allowed(cmd, allow_config=False):
    """생성된 CLI 안전 필터.

    allow_config 면 설정 허용 목록(_config_allowed)까지 열고, 아니면 조회만.
    어느 쪽이든 _NEVER_WORDS 는 못 지나간다 — TC 초안과 **같은 정책**을 쓴다.
    """
    c = str(cmd or "").strip().lower()
    if not c:
        return False
    if any(b in c for b in _NEVER_WORDS):
        return False
    if allow_config:
        return _config_allowed(cmd)
    head = c.split()[0]
    return head in ("show", "display", "ping", "traceroute", "dir", "more", "cat", "get", "status", "monitor")

@app.post("/api/ai/nl-exec")
async def ai_nl_exec(payload: dict):
    """자연어 지시 → CLI 생성(Gemma) → (execute=true면) 실행 → 출력 해석·판정.
    기본은 조회성 명령만 허용. 장비 접속 파라미터는 /api/run-cli 와 동일하게 전달."""
    text = str((payload or {}).get("text") or "").strip()
    if not text:
        return {"ok": False, "error": "지시문(text)이 비어 있습니다"}
    dev_model = str(payload.get("model") or "").strip()
    schema = {"type": "object", "properties": {
        "commands": {"type": "array", "items": {"type": "string"}},
        "purpose": {"type": "string"}, "criteria": {"type": "string"}}, "required": ["commands"]}
    sys_p = ("너는 네트워크 장비 CLI 전문가다. 사용자의 자연어 지시를 장비에서 실행할 CLI 명령 목록으로 변환한다. "
             "조회(show/display/ping 등) 명령만 생성하고 설정 변경·재부팅 명령은 절대 만들지 않는다. "
             "purpose에는 무엇을 확인하는지, criteria에는 출력에서 확인할 판정 포인트를 쓴다. JSON만 출력한다.")
    content, err = await _ai_chat([{"role": "system", "content": sys_p},
                                   {"role": "user", "content": (f"대상 모델: {dev_model or '공통'}\n지시: {text}")}],
                                  max_tokens=800, json_schema=schema)
    if err:
        return {"ok": False, "error": err}
    obj = _ai_json(content) or {}
    raw_cmds = [str(c).strip() for c in (obj.get("commands") or []) if str(c or "").strip()]
    allow_cfg = bool(payload.get("allow_config"))
    cmds = [c for c in raw_cmds if _nl_cmd_allowed(c, allow_cfg)]
    blocked = [c for c in raw_cmds if c not in cmds]
    result = {"ok": True, "commands": cmds, "blocked": blocked, "purpose": obj.get("purpose", ""), "criteria": obj.get("criteria", "")}
    if not cmds:
        result.update({"ok": False, "error": "실행 가능한(안전한) 명령이 생성되지 않았습니다"})
        return result
    if not payload.get("execute"):
        return result   # 프리뷰 모드: 명령만 반환 (기본)
    run = await asyncio.to_thread(run_cli, {**payload, "commands": cmds})
    outputs = run.get("outputs") or []
    result["outputs"] = outputs
    if not run.get("ok", True) and run.get("error"):
        result.update({"ok": False, "error": run.get("error")})
        return result
    # 출력 해석·판정
    jschema = {"type": "object", "properties": {"verdict": {"type": "string", "enum": ["Pass", "Fail", "Unknown"]}, "reason": {"type": "string"}}, "required": ["verdict", "reason"]}
    out_txt = json.dumps(outputs, ensure_ascii=False)[:6000]
    sys_j = ("너는 네트워크 장비 시험 판정 전문가다. 실행 출력이 판정 포인트를 만족하면 Pass, 아니면 Fail, 판단 불가는 Unknown. "
             "reason에 근거(출력 인용)를 한국어로 쓴다. JSON만 출력한다.")
    jcontent, jerr = await _ai_chat([{"role": "system", "content": sys_j},
                                     {"role": "user", "content": f"지시: {text}\n판정 포인트: {obj.get('criteria','')}\n[실행 출력]\n{out_txt}"}],
                                    max_tokens=700, json_schema=jschema)
    jobj = _ai_json(jcontent) if jcontent else None
    result["verdict"] = (jobj or {}).get("verdict", "Unknown")
    result["reason"] = (jobj or {}).get("reason", jerr or "")
    return result

# ══════════ 지식 소스 통합 설정 (① 일반 gemma ② UTOP 내부 지식(시험절차+매뉴얼) ③ Jira ④ Confluence — 전역 On/Off) ══════════
KNOWLEDGE_SRC_FILE = DATA_DIR / "state" / "knowledge_sources.json"
_KNOWLEDGE_SRC_DEFAULT = {"general": True, "internal": True, "jira": True, "confluence": True}

def _knowledge_src_cfg():
    try:
        if KNOWLEDGE_SRC_FILE.exists():
            d = load_json(KNOWLEDGE_SRC_FILE)
            if isinstance(d, dict):
                return {**_KNOWLEDGE_SRC_DEFAULT, **{k: bool(d.get(k, v)) for k, v in _KNOWLEDGE_SRC_DEFAULT.items()}}
    except Exception:
        pass
    return dict(_KNOWLEDGE_SRC_DEFAULT)

@app.get("/api/knowledge-sources")
async def knowledge_sources_get():
    cfg = _knowledge_src_cfg()
    cfg["confluence_scopes"] = _conf_cfg().get("scopes") or []
    return cfg

@app.post("/api/knowledge-sources")
async def knowledge_sources_set(payload: dict, token: str = ""):
    _require_admin(token)
    cur = _knowledge_src_cfg()
    p = payload or {}
    for k in _KNOWLEDGE_SRC_DEFAULT:
        if k in p:
            cur[k] = bool(p.get(k))
    save_json(KNOWLEDGE_SRC_FILE, cur)
    return {"ok": True, **cur}

# ══════════ Confluence 연동 (페이지→매뉴얼 동기화: 트리=폴더, 첨부=이미지) ══════════
CONFLUENCE_CONFIG_FILE = DATA_DIR / "integrations" / "confluence_config.json"
_CONF_DEFAULT = {"base_url": "", "auth_type": "bearer", "token": "", "username": "", "password": "", "space_key": "", "enabled": False, "live_query": False, "scopes": []}
# scopes: [{"id","label","space_key","parent_title","enabled"}, ...] — 검색 범위를 "스페이스 + 상위 페이지(하위 포함)"
# 여러 개 등록해 관리자가 각각 On/Off. 비어 있으면 기존 하드코딩 경로(11.Feature List/12.How to debuging, 전역 space_key)로 폴백.

def _clean_conf_scopes(scopes):
    if not isinstance(scopes, list):
        return []
    out = []
    for it in scopes:
        if not isinstance(it, dict):
            continue
        parent = str(it.get("parent_title", "")).strip()
        page_id = str(it.get("page_id", "")).strip()
        url = str(it.get("url", "")).strip()
        space = str(it.get("space_key", "")).strip()
        label = str(it.get("label", "")).strip()
        # 최소 하나의 식별자(page_id / parent_title / space_key / url) 라도 있어야 저장
        if not (parent or page_id or space or url):
            continue
        # id 자동 생성 — page_id 우선, 없으면 parent+space, 없으면 url
        _seed = page_id or (parent + space) or url
        try:
            _depth = int(it.get("depth")) if it.get("depth") not in (None, "") else None
        except Exception:
            _depth = None
        out.append({
            "id": str(it.get("id") or _rag_doc_id(_seed)),
            "label": label or parent or space or (("pageId=" + page_id) if page_id else url),
            "url": url,
            "space_key": space,
            "parent_title": parent,
            "page_id": page_id,
            "depth": _depth if (_depth is not None and 1 <= _depth <= 10) else None,
            "enabled": bool(it.get("enabled", True)),
        })
    return out

def _conf_cfg():
    try:
        if CONFLUENCE_CONFIG_FILE.exists():
            d = load_json(CONFLUENCE_CONFIG_FILE)
            if isinstance(d, dict):
                merged = {**_CONF_DEFAULT, **d}
                merged["scopes"] = _clean_conf_scopes(merged.get("scopes"))
                return merged
    except Exception:
        pass
    return dict(_CONF_DEFAULT)

def _conf_headers(cfg):
    h = {"Accept": "application/json"}
    if cfg.get("auth_type") == "bearer" and cfg.get("token"):
        h["Authorization"] = "Bearer " + str(cfg["token"])
    return h

def _conf_auth(cfg):
    if cfg.get("auth_type") == "basic" and cfg.get("username"):
        return (str(cfg["username"]), str(cfg.get("password") or cfg.get("token") or ""))
    return None

def _html_to_text(html):
    import re as _rh
    s = str(html or "")
    s = _rh.sub(r"(?is)<(script|style).*?</\1>", "", s)
    s = _rh.sub(r"(?i)<br\s*/?>", "\n", s)
    s = _rh.sub(r"(?i)</(p|div|li|tr|h[1-6]|td)>", "\n", s)
    s = _rh.sub(r"<[^>]+>", " ", s)
    for a, b in (("&nbsp;", " "), ("&amp;", "&"), ("&lt;", "<"), ("&gt;", ">"), ("&quot;", '"')):
        s = s.replace(a, b)
    s = _rh.sub(r"[ \t]+", " ", s)
    s = _rh.sub(r"\n\s*\n\s*\n+", "\n\n", s)
    return s.strip()

@app.get("/api/confluence/config")
async def conf_config_get():
    c = _conf_cfg()
    c2 = dict(c)
    c2["token"] = "***" if c.get("token") else ""
    c2["password"] = "***" if c.get("password") else ""
    return c2

@app.post("/api/confluence/config")
async def conf_config_set(payload: dict, token: str = ""):
    _require_admin(token)
    cfg = _conf_cfg()
    for k in ("base_url", "auth_type", "username", "space_key"):
        if k in payload:
            cfg[k] = str(payload.get(k) or "").strip()
    for k in ("token", "password"):
        if k in payload and payload.get(k) not in (None, "", "***"):
            cfg[k] = str(payload.get(k))
    for k in ("enabled", "live_query"):
        if k in payload:
            cfg[k] = bool(payload.get(k))
    if "default_depth" in payload:
        try:
            _dd = int(payload.get("default_depth"))
            if 1 <= _dd <= 10:
                cfg["default_depth"] = _dd
        except Exception:
            pass
    if "scopes" in payload:
        cfg["scopes"] = _clean_conf_scopes(payload.get("scopes"))
    save_json(CONFLUENCE_CONFIG_FILE, cfg)
    return {"ok": True}

# 로고·아이콘·장식 제외 + 큰 이미지(콘텐츠성) 우선해서 페이지 대표 이미지 추출
_CONF_IMG_SKIP = ("logo", "icon", "banner", "button", "arrow", "bullet", "emoticon",
                  "avatar", "badge", "favicon", "header", "footer", "divider", "spacer",
                  "thumb", "small", "line.", "dot.", "bg.", "background")

def _conf_mark_images(body):
    """storage 본문의 <ac:image>를 [[CIMG:파일명]] 마커로 치환 → 텍스트 추출 시 이미지 위치 보존."""
    import re as _ri
    def repl(m):
        block = m.group(0)
        fm = _ri.search(r'ri:filename="([^"]+)"', block)
        if fm:
            return f" [[CIMG:{fm.group(1)}]] "
        return " "
    return _ri.sub(r"(?is)<ac:image.*?</ac:image>", repl, str(body or ""))

def _conf_inline_image_urls(text_marked, base, att_map):
    """[[CIMG:파일명]] → 마크다운 이미지 ![image](공개URL). (위키 공개라 URL 직접 렌더 — Dify식)"""
    import re as _ri
    def repl(m):
        fn = m.group(1)
        info = att_map.get(fn)
        if not info or not info.get("dl") or "image" not in str(info.get("mt")):
            return ""
        fnl = fn.lower()
        if any(b in fnl for b in _CONF_IMG_SKIP) or fnl.endswith(".gif") or "svg" in str(info.get("mt")):
            return ""
        url = (base + info["dl"]) if info["dl"].startswith("/") else info["dl"]
        return f"\n![image]({url})\n"
    return _ri.sub(r"\[\[CIMG:([^\]]+)\]\]", repl, text_marked)

async def _conf_page_attachments(client, base, headers, auth, pid):
    """페이지 첨부 맵 {파일명:{mt,dl,fs}} + 비이미지 파일목록(PDF 등) 반환."""
    att_map = {}; files = []
    try:
        ar = await client.get(base + f"/rest/api/content/{pid}/child/attachment", headers=headers, auth=auth, params={"limit": 50})
        if ar.status_code == 200:
            for att in (ar.json().get("results") or []):
                fn = att.get("title") or ""
                mt = ((att.get("metadata") or {}).get("mediaType")) or ""
                dl = ((att.get("_links") or {}).get("download")) or ""
                try:
                    fs = int(((att.get("extensions") or {}).get("fileSize")) or 0)
                except Exception:
                    fs = 0
                if not fn or not dl:
                    continue
                att_map[fn] = {"mt": mt, "dl": dl, "fs": fs}
                fnl = fn.lower()
                if "image" not in str(mt) and any(fnl.endswith(e) for e in (".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".hwp", ".txt", ".csv")):
                    files.append({"name": fn, "url": base + dl})
    except Exception:
        pass
    return att_map, files

# ── Dify식 트리 스코핑 + 모델 정밀 매칭 ──
_CONF_TREE_CACHE = {}  # 부모 제목 -> [{"id","title"}] (메모리 캐시)

_CONF_PID_CACHE = {}
async def _conf_parent_id(client, base, headers, auth, parent_title, space_key=None):
    """페이지 제목 → pageId 조회. 결과 캐시. space 미지정이면 전역(모든 space)에서 정확 title 매치.
    (CQL 대신 페이지 API 로만 조회 — space 제약 없이 동작)."""
    sp = space_key or _conf_cfg().get("space_key") or ""
    _key = f"{sp or '*'}::{parent_title}"
    if _key in _CONF_PID_CACHE:
        return _CONF_PID_CACHE[_key]
    _pid = None
    try:
        _params = {"title": parent_title, "limit": 10}
        if sp: _params["spaceKey"] = sp
        r = await client.get(base + "/rest/api/content", headers=headers, auth=auth, params=_params)
        if r.status_code == 200:
            for pg in (r.json().get("results") or []):
                if (pg.get("title") or "").strip() == parent_title:
                    _pid = pg.get("id"); break
        # 지정 space 에서 못 찾았으면 전역 재조회
        if not _pid and sp:
            r2 = await client.get(base + "/rest/api/content", headers=headers, auth=auth, params={"title": parent_title, "limit": 10})
            if r2.status_code == 200:
                for pg in (r2.json().get("results") or []):
                    if (pg.get("title") or "").strip() == parent_title:
                        _pid = pg.get("id"); break
    except Exception:
        pass
    _CONF_PID_CACHE[_key] = _pid
    return _pid

async def _conf_children(client, base, headers, auth, parent_title, space_key=None, depth=1):
    """부모 페이지의 직계 자식 목록 (모델코드 정밀 매칭용). depth 는 사용 안 함 — 하위 전체 검색은 CQL ancestor 로.
    (예전 동작 유지 — 재귀 순회는 매우 느려서 depth 파라미터를 프론트 UX 로만 남기고 여기선 무시)."""
    cache_key = f"{space_key or ''}::{parent_title}"
    if cache_key in _CONF_TREE_CACHE:
        return _CONF_TREE_CACHE[cache_key]
    out = []
    pid = await _conf_parent_id(client, base, headers, auth, parent_title, space_key)
    if pid:
        try:
            r = await client.get(base + f"/rest/api/content/{pid}/child/page", headers=headers, auth=auth, params={"limit": 100})
            if r.status_code == 200:
                out = [{"id": p.get("id"), "title": p.get("title") or ""} for p in (r.json().get("results") or [])]
        except Exception:
            pass
    _CONF_TREE_CACHE[cache_key] = out
    return out

def _spec_title_regex(title):
    """'E61xx_Series_Spec' → ^E61\\d\\d$ (소문자 x = 숫자 와일드카드). 모델코드 정밀 매칭용."""
    import re as _rs
    b = str(title or "").replace("_Series_Spec", "").replace("_Series", "")
    pat = "".join((r"\d" if ch == "x" else _rs.escape(ch)) for ch in b)
    try:
        return _rs.compile("^" + pat + "$", _rs.I)
    except Exception:
        return None

# 페이지 본문 hit 캐시 (pageId → hit dict). TTL 로 오래된 것 자동 폐기 (변경 감지 어려워 짧게).
_CONF_HIT_CACHE = {}
_CONF_HIT_TTL = 600   # 10분

async def _conf_fetch_page_hit(client, base, headers, auth, pid, title, toks=None):
    """페이지 1개를 풀 조회 → 본문(이미지 인라인 URL)+첨부파일 hit 구성. 결과 10분 캐시."""
    import time as _tm
    _now = _tm.time()
    _cached = _CONF_HIT_CACHE.get(pid)
    if _cached and (_now - _cached[0]) < _CONF_HIT_TTL:
        return _cached[1]
    try:
        r = await client.get(base + f"/rest/api/content/{pid}", headers=headers, auth=auth, params={"expand": "body.storage"})
        if r.status_code != 200:
            return None
        pg = r.json()
        body = (((pg.get("body") or {}).get("storage") or {}).get("value")) or ""
        att_map, files = await _conf_page_attachments(client, base, headers, auth, pid)
        files = [f for f in files if f.get("name") and (f["name"] in body or (toks and any(t.lower() in f["name"].lower() for t in toks)))]
        text = _conf_inline_image_urls(_html_to_text(_conf_mark_images(body))[:7000], base, att_map)
        _hit = {"name": "Confluence · " + (title or pg.get("title") or ""), "text": text, "images": [],
                "files": files, "url": base + (((pg.get("_links") or {}).get("webui")) or "")}
        _CONF_HIT_CACHE[pid] = (_now, _hit)
        return _hit
    except Exception:
        return None

def _spec_digit_regex(title):
    """제목의 숫자 시그니처 → 정규식 ('E71xx'→71\\d\\d, 'E4300'→4300). 숫자만 입력(7100·4300) 매칭용."""
    import re as _rs
    b = str(title or "").replace("_Series_Spec", "").replace("_Series", "")
    sig = "".join(ch for ch in b if ch.isdigit() or ch == "x")
    if not sig or not any(c.isdigit() for c in sig):
        return None
    pat = "".join((r"\d" if ch == "x" else ch) for ch in sig)
    try:
        return _rs.compile("^" + pat + "$")
    except Exception:
        return None

_CONF_MODEL_INDEX = {"map": None}
async def _conf_model_index(client, base, headers, auth):
    """본문 'Model Name' 라인에서 실제 모델명 추출 → {모델명(소문자): 페이지} 매핑 (E4320·E7148T 등 형제모델 포함, 캐시)."""
    if _CONF_MODEL_INDEX.get("map") is not None:
        return _CONF_MODEL_INDEX["map"]
    import re as _ri
    mp = {}
    children = await _conf_children(client, base, headers, auth, "11.Feature List")
    for c in children:
        try:
            r = await client.get(base + f"/rest/api/content/{c.get('id')}", headers=headers, auth=auth, params={"expand": "body.storage"})
            if r.status_code != 200:
                continue
            body = (((r.json().get("body") or {}).get("storage") or {}).get("value")) or ""
            m = _ri.search(r"Model\s*Name\s*[:：]\s*([^\n]+)", _html_to_text(body), _ri.I)
            if not m:
                continue
            for mod in _ri.findall(r"[A-Za-z]{1,4}\d{3,}[A-Za-z0-9]*", m.group(1)):
                ref = {"id": c.get("id"), "title": c.get("title")}
                mp[mod.lower()] = ref
                dg = _ri.sub(r"\D", "", mod)
                if dg:
                    mp.setdefault("#" + dg, ref)   # 숫자만(4320) 키
        except Exception:
            continue
    _CONF_MODEL_INDEX["map"] = mp
    return mp

async def _conf_children_tree(client, base, headers, auth, page_id, depth=3, cache=None):
    """pageId 하위 페이지 트리 (자기 자신 포함) 재귀 수집 — child/page API 직접 호출.
    cache: (page_id, depth) → 결과 리스트 캐시 (프로세스 생존 동안)."""
    if cache is None:
        cache = _CONF_TREE_ID_CACHE
    _key = f"{page_id}::d{depth}"
    if _key in cache:
        return cache[_key]
    out = []
    # 자기 자신 (title 조회)
    try:
        r0 = await client.get(base + f"/rest/api/content/{page_id}", headers=headers, auth=auth)
        if r0.status_code == 200:
            j = r0.json()
            out.append({"id": j.get("id"), "title": j.get("title") or ""})
    except Exception:
        cache[_key] = out
        return out
    frontier = [page_id]
    for _ in range(max(1, int(depth or 1))):
        next_ids = []
        _tasks = [client.get(base + f"/rest/api/content/{pid}/child/page", headers=headers, auth=auth, params={"limit": 200}) for pid in frontier]
        _rs = await asyncio.gather(*_tasks, return_exceptions=True)
        for _r in _rs:
            if not _r or isinstance(_r, Exception) or _r.status_code != 200:
                continue
            try:
                for p in (_r.json().get("results") or []):
                    _id = p.get("id"); _t = p.get("title") or ""
                    if _id:
                        out.append({"id": _id, "title": _t})
                        next_ids.append(_id)
            except Exception:
                continue
        if not next_ids: break
        frontier = next_ids
    cache[_key] = out
    return out

_CONF_TREE_ID_CACHE = {}

async def _conf_match_under_id(client, base, headers, auth, query, page_id, depth=3):
    """pageId 기반 매칭 — child/page API 로 트리 확보 후 제목/본문 매칭.
    CQL 완전 걷어냄 — space 제약 없이 페이지 트리 자체가 검색 대상."""
    import re as _rq
    pages = await _conf_children_tree(client, base, headers, auth, page_id, depth)
    if not pages:
        return None
    raw = str(query).replace('"', " ").strip()
    if not raw:
        return None
    _stop = {"알려", "정보", "대해", "무엇", "어떻", "해줘", "해주", "있나", "뭐야", "어디", "그리고",
             "주세요", "알려줘", "확인", "내용", "추가", "질문", "정확", "어떤", "찾으", "중에", "대한", "관련"}
    _allt = []
    for t in _rq.findall(r"[A-Za-z0-9]{2,}|[가-힣]{2,}", raw):
        if t not in _stop and t not in _allt:
            _allt.append(t)
    _alnum = [t for t in _allt if _rq.match(r"^[A-Za-z0-9]+$", t)]
    toks = (_alnum + [t for t in _allt if t not in _alnum])[:6]
    if not toks:
        # 키워드 없으면 자기 자신(루트) 본문만 반환
        top = pages[:1]
    else:
        # 제목 매칭 스코어 + 날짜 최신순 tiebreaker (주간 업무처럼 날짜가 title 에 있으면 최근 우선)
        _date_rx = _rq.compile(r'(\d{2})[년._-]\s*(\d{1,2})[월._-]\s*(\d{1,2})|(\d{4})[-._]?(\d{2})[-._]?(\d{2})')
        def _title_date(pg):
            t = str(pg.get("title") or "")
            m = _date_rx.search(t)
            if not m: return 0
            try:
                g = m.groups()
                if g[3]:  # YYYYMMDD
                    return int(g[3])*10000 + int(g[4])*100 + int(g[5])
                # YY년MM월DD (26년05월21일 → 20260521)
                return (2000+int(g[0]))*10000 + int(g[1])*100 + int(g[2])
            except Exception:
                return 0
        def _title_score(pg):
            t = str(pg.get("title") or "").lower()
            sc = 0
            for tk in _alnum:
                if tk.lower() in t: sc += 5
            for tk in toks:
                if tk in t: sc += 1
            return sc
        # 1차: 매칭 스코어 내림차순, 2차: 날짜 최신 우선
        pages_sorted = sorted(pages, key=lambda p: (_title_score(p), _title_date(p)), reverse=True)
        top = pages_sorted[:3]
        # 매칭 스코어가 전부 0 이면 날짜 최신 우선 + 자기 자신
        if all(_title_score(p) == 0 for p in top):
            top = sorted(pages, key=_title_date, reverse=True)[:3]
    # 상위 페이지 본문 병렬 fetch
    fetch_tasks = [_conf_fetch_page_hit(client, base, headers, auth, pg.get("id"), pg.get("title") or "", toks) for pg in top]
    hits_top = await asyncio.gather(*fetch_tasks, return_exceptions=True)
    merged_parts = []
    combined_files = []
    first_hit = None
    pages_meta = []
    for hit in hits_top:
        if not hit or isinstance(hit, Exception): continue
        if first_hit is None: first_hit = hit
        merged_parts.append(f"=== {hit['name']} ===\n{hit['text']}")
        combined_files.extend(hit.get("files") or [])
        pages_meta.append({"name": hit.get("name") or "", "url": hit.get("url") or ""})
    if not merged_parts:
        return None
    return {
        "name": (first_hit or {}).get("name") if first_hit else "Confluence · 여러 페이지",
        "text": "\n\n".join(merged_parts),
        "images": [],
        "files": combined_files,
        "url": (first_hit or {}).get("url", "") if first_hit else "",
        "pages": pages_meta,
    }

async def _conf_match_under_id_OLD_CQL(client, base, headers, auth, query, page_id, depth=3):
    """(사용 안 함) 옛 CQL ancestor 방식 — 참고용 보존."""
    import re as _rq
    raw = str(query).replace('"', " ").strip()
    if not raw:
        return None
    _stop = {"알려", "정보", "대해", "무엇", "어떻", "해줘", "해주", "있나", "뭐야", "어디", "그리고",
             "주세요", "알려줘", "확인", "내용", "추가", "질문", "정확", "어떤", "찾으", "중에", "대한", "관련"}
    _allt = []
    for t in _rq.findall(r"[A-Za-z0-9]{2,}|[가-힣]{2,}", raw):
        if t not in _stop and t not in _allt:
            _allt.append(t)
    _alnum = [t for t in _allt if _rq.match(r"^[A-Za-z0-9]+$", t)]
    toks = (_alnum + [t for t in _allt if t not in _alnum])[:6]
    if not toks:
        return None
    terms = " OR ".join(f'text ~ "{t}"' for t in toks)
    # ancestor 는 pageId 하위 전체 트리 (자기 자신 제외) → 자기 자신 포함하려면 OR 추가
    cql = f'({terms}) AND (ancestor = "{page_id}" OR id = {page_id}) AND type = page'
    try:
        r = await client.get(base + "/rest/api/content/search", headers=headers, auth=auth,
                              params={"cql": cql, "limit": 6, "expand": "body.storage"})
        if r.status_code != 200:
            return None
        results = r.json().get("results") or []
        if not results:
            return None
        # 스코어링: 제목/본문에 alnum 토큰 매칭 우선
        def _score(pg):
            title = str(pg.get("title") or ""); body = str((((pg.get("body") or {}).get("storage") or {}).get("value")) or "")
            sc = 0
            for t in _alnum:
                if t.lower() in title.lower(): sc += 5
                if t in body: sc += 1
            for t in toks:
                if t in body: sc += 1
            return sc
        results.sort(key=_score, reverse=True)
        top = results[:3]
        fetch_tasks = [_conf_fetch_page_hit(client, base, headers, auth, pg.get("id"), pg.get("title") or "", toks) for pg in top]
        hits_top = await asyncio.gather(*fetch_tasks, return_exceptions=True)
        # 상위 페이지들을 병합해 컨텍스트로 보내되, 각 hit 마다 자기 페이지 URL 을 유지 (프론트가 출처 링크로 표시)
        merged_parts = []
        combined_files = []
        first_hit = None
        pages_meta = []
        for hit in hits_top:
            if not hit or isinstance(hit, Exception):
                continue
            if first_hit is None:
                first_hit = hit
            merged_parts.append(f"=== {hit['name']} ===\n{hit['text']}")
            combined_files.extend(hit.get("files") or [])
            pages_meta.append({"name": hit.get("name") or "", "url": hit.get("url") or ""})
        if not merged_parts:
            return None
        return {
            "name": (first_hit or {}).get("name", "Confluence · 여러 페이지"),
            "text": "\n\n".join(merged_parts),
            "images": [],
            "files": combined_files,
            "url": (first_hit or {}).get("url", ""),
            "pages": pages_meta,   # 각 페이지 개별 URL — 프론트 출처 표시용
        }
    except Exception:
        return None

async def _conf_match_under(client, base, headers, auth, query, parent_title, use_index=False, space_key=None, depth=1):
    """parent_title → pageId → 트리 순회 방식으로 매칭 (CQL 안 씀)."""
    import re as _ri
    pid = await _conf_parent_id(client, base, headers, auth, parent_title, space_key)
    if pid:
        _hit = await _conf_match_under_id(client, base, headers, auth, query, pid, depth)
        if _hit:
            return _hit
    # pageId 못 찾은 경우: 기존 직계 자식 조회 (모델코드 정밀 매칭용)
    children = await _conf_children(client, base, headers, auth, parent_title, space_key)
    if not children:
        return None
    toks = _ri.findall(r"[A-Za-z]{0,4}\d{3,}[A-Za-z0-9]*", str(query))  # 숫자만(4300)도 허용
    idx = (await _conf_model_index(client, base, headers, auth)) if (use_index and toks) else {}
    for tok in toks:
        tl = tok.lower(); dg = _ri.sub(r"\D", "", tok)
        page = None
        # ① 실제 모델명 정확 일치 (E4320, E7148T …) — 스펙 인덱스 사용 시
        if use_index and tl in idx:
            page = idx[tl]
        # ② 제목 정규식 (E61xx→E61\d\d) — 풀 토큰 매칭
        if not page:
            for c in children:
                rgx = _spec_title_regex(c.get("title"))
                if rgx and rgx.match(tok):
                    page = {"id": c.get("id"), "title": c.get("title")}; break
        # ③ 숫자 시그니처 매칭 (7100→E71xx, 4300→E4300)
        if not page and dg:
            if use_index and ("#" + dg) in idx:
                page = idx["#" + dg]
            else:
                for c in children:
                    drx = _spec_digit_regex(c.get("title"))
                    if drx and drx.match(dg):
                        page = {"id": c.get("id"), "title": c.get("title")}; break
        if page:
            hit = await _conf_fetch_page_hit(client, base, headers, auth, page["id"], page["title"], toks)
            if hit:
                return hit
    # 폴백: 모델코드 매칭이 안 됐으면(토큰 없음 또는 전부 실패) 이 하위 페이지들 안에서 일반 CQL 텍스트검색
    return await _conf_search_within(client, base, headers, auth, query, children)

async def _conf_search_within(client, base, headers, auth, query, pages):
    """주어진 페이지 목록(하위 트리) 안에서 CQL 텍스트 검색.
    상위 1개만이 아니라 관련 페이지 여러 개를 합쳐 반환 — 세부 페이지가 있을 때 최상위 목록만 뽑히는 문제 해결."""
    import re as _rq
    raw = str(query).replace('"', " ").strip()
    if not raw or not pages:
        return None
    _stop = {"알려", "정보", "대해", "무엇", "어떻", "해줘", "해주", "있나", "뭐야", "어디", "그리고",
             "주세요", "알려줘", "확인", "내용", "추가", "질문", "정확", "어떤", "찾으", "중에", "대한", "관련"}
    toks = [t for t in _rq.findall(r"[A-Za-z0-9]{2,}|[가-힣]{2,}", raw) if t not in _stop][:6]
    if not toks:
        return None
    ids = [p.get("id") for p in pages if p.get("id")]
    if not ids:
        return None
    # 전체 페이지 목록을 조각(chunk) 으로 나눠 각 조각마다 CQL 실행 → 40개 하드 제한 회피
    _alnum = [t for t in toks if _rq.match(r"^[A-Za-z0-9]+$", t)]
    hits_all = []
    chunk_size = 35   # id_clause 길이 여유
    try:
        for _ci in range(0, len(ids), chunk_size):
            _chunk = ids[_ci:_ci+chunk_size]
            terms = " OR ".join(f'text ~ "{t}"' for t in toks)
            id_clause = " OR ".join(f'id = {pid}' for pid in _chunk)
            cql = f'({terms}) AND ({id_clause}) AND type = page'
            r = await client.get(base + "/rest/api/content/search", headers=headers, auth=auth,
                                  params={"cql": cql, "limit": 6, "expand": "body.storage"})
            if r.status_code != 200:
                continue
            for pg in (r.json().get("results") or []):
                hits_all.append(pg)
        if not hits_all:
            return None
        # 스코어링: 제목/본문에 영숫자 토큰(모델코드) 정확 매칭 우선. 그 다음 텍스트 매칭 개수.
        def _score(pg):
            title = str(pg.get("title") or ""); body = str((((pg.get("body") or {}).get("storage") or {}).get("value")) or "")
            sc = 0
            for t in _alnum:
                if t.lower() in title.lower(): sc += 5   # 제목 매치는 강한 신호
                if t in body: sc += 1
            for t in toks:
                if t in body: sc += 1
            return sc
        hits_all.sort(key=_score, reverse=True)
        # 상위 3개 페이지 hit 을 하나의 텍스트 블록으로 합쳐 반환 (LLM 컨텍스트 확장)
        top = hits_all[:3]
        merged_parts = []
        combined_files = []
        for pg in top:
            hit = await _conf_fetch_page_hit(client, base, headers, auth, pg.get("id"), pg.get("title") or "", toks)
            if not hit:
                continue
            merged_parts.append(f"=== {hit['name']} ===\n{hit['text']}")
            combined_files.extend(hit.get("files") or [])
        if not merged_parts:
            return None
        # 상위 페이지들 병렬 fetch → 페이지별 URL 개별 유지 (프론트 출처 하이퍼링크용)
        top_pages = hits_all[:2]
        fetch_tasks = [_conf_fetch_page_hit(client, base, headers, auth, pg.get("id"), pg.get("title") or "", toks) for pg in top_pages]
        top_hits = await asyncio.gather(*fetch_tasks, return_exceptions=True)
        pages_meta = []
        first_hit = None
        for _h in top_hits:
            if not _h or isinstance(_h, Exception): continue
            if first_hit is None: first_hit = _h
            pages_meta.append({"name": _h.get("name") or "", "url": _h.get("url") or ""})
        return {
            "name": (first_hit or {}).get("name") if first_hit else "Confluence · 여러 페이지",
            "text": "\n\n".join(merged_parts),
            "images": [],
            "files": combined_files,
            "url": (first_hit or {}).get("url", "") if first_hit else "",
            "pages": pages_meta,
        }
    except Exception:
        return None

async def _conf_match_spec(client, base, headers, auth, query):
    return await _conf_match_under(client, base, headers, auth, query, "11.Feature List", use_index=True)

async def _confluence_live_search(query, limit=4):
    """질문 시 Confluence를 CQL로 라이브 검색 → 관련 페이지(본문·이미지) 반환 (Dify식)."""
    if not _knowledge_src_cfg().get("confluence", True):
        return []   # 지식 소스 통합 설정에서 Confluence 검색 자체를 껐으면 즉시 중단
    cfg = _conf_cfg()
    base = str(cfg.get("base_url") or "").rstrip("/")
    if not base or not cfg.get("live_query") or not str(query or "").strip():
        return []
    import httpx
    headers = _conf_headers(cfg); auth = _conf_auth(cfg)
    space = str(cfg.get("space_key") or "").strip()
    import re as _rq
    raw = str(query).replace('"', " ").strip()
    _stop = {"알려", "정보", "대해", "무엇", "어떻", "해줘", "해주", "있나", "뭐야", "어디", "그리고",
             "주세요", "알려줘", "확인", "내용", "추가", "질문", "정확", "어떤", "찾으", "중에", "대한", "관련"}
    _allt = []
    for t in _rq.findall(r"[A-Za-z0-9]{2,}|[가-힣]{2,}", raw):
        if t not in _stop and t not in _allt:
            _allt.append(t)
    # 영숫자(모델코드·식별자) 토큰을 앞에 배치 → CQL/부스트에서 누락 방지 (긴 되물음 답변에서도 모델코드 보존)
    _alnum = [t for t in _allt if _rq.match(r"^[A-Za-z0-9]+$", t)]
    toks = (_alnum + [t for t in _allt if t not in _alnum])[:6]
    if not toks and raw:
        toks = [raw[:40]]
    if not toks:
        return []
    scopes = [s for s in (cfg.get("scopes") or []) if s.get("enabled")]
    # 여러 스페이스를 등록했으면 CQL에서도 그 스페이스들만 대상으로(전역 space_key 하나만 걸던 것 확장)
    scope_spaces = sorted({s["space_key"] for s in scopes if s.get("space_key")})
    space_clause = ""
    if scope_spaces:
        space_clause = " AND (" + " OR ".join(f'space = "{sp}"' for sp in scope_spaces) + ")"
    elif space:
        space_clause = f' AND space = "{space}"'
    terms = " OR ".join(f'text ~ "{t}"' for t in toks)
    cql = f"({terms})" + space_clause + " AND type = page"
    out = []
    fetch_n = min(max(limit * 3, 10), 18)   # 후보를 넉넉히 받아 리랭크/스펙우선 정렬 후 상위만 사용
    try:
        async with httpx.AsyncClient(timeout=30, verify=False) as client:
            if scopes:
                # 등록된 모든 scope 병렬 검색 → 각 scope 결과를 개별 hit 으로 리턴 (병합 X)
                # 프론트가 hit 개수만큼 프롬프트에 담아 슬라이스가 균등 분배됨
                default_depth = int(cfg.get("default_depth") or 3)
                async def _one(sc):
                    _depth = int(sc.get("depth") or default_depth or 3)
                    try:
                        if sc.get("page_id"):
                            return await _conf_match_under_id(client, base, headers, auth, query, sc["page_id"], _depth)
                        if sc.get("parent_title"):
                            return await _conf_match_under(client, base, headers, auth, query, sc["parent_title"], use_index=False, space_key=(sc.get("space_key") or None), depth=_depth)
                    except Exception:
                        return None
                    return None
                _hits = await asyncio.gather(*[_one(sc) for sc in scopes], return_exceptions=True)
                # 각 scope 결과를 개별 hit 으로 (병합 안 함)
                out_hits = []
                seen_urls = set()
                for h in _hits:
                    if not h or isinstance(h, Exception): continue
                    _u = h.get("url") or ""
                    if _u and _u in seen_urls: continue
                    if _u: seen_urls.add(_u)
                    out_hits.append(h)
                if out_hits:
                    return out_hits[:8]   # 최대 8개 scope hit
            else:
                # scopes 미설정 시 기존 하드코딩 경로로 폴백(11.Feature List/12.How to debuging) — 기존 동작 유지
                _is_dbg = bool(_rq.search(r"디버깅|디버그|debug|트러블|장애|문제\s*해결|how\s*to", str(query), _rq.I))
                if _is_dbg:
                    _dbg = await _conf_match_under(client, base, headers, auth, query, "12.How to debuging")
                    if _dbg:
                        return [_dbg]
                else:
                    # Dify식 정밀 매칭: 모델코드 → 11.Feature List 하위 스펙 페이지 1개 (노이즈 0)
                    _spec = await _conf_match_spec(client, base, headers, auth, query)
                    if _spec:
                        return [_spec]
            # ② 폴백: 전체 CQL 전문검색 + 리랭크
            r = await client.get(base + "/rest/api/content/search", headers=headers, auth=auth,
                                  params={"cql": cql, "limit": fetch_n, "expand": "body.storage"})
            if r.status_code != 200:
                return []
            for pg in (r.json().get("results") or []):
                pid = pg.get("id"); title = pg.get("title") or ""
                body = (((pg.get("body") or {}).get("storage") or {}).get("value")) or ""
                att_map, files = await _conf_page_attachments(client, base, headers, auth, pid)
                # '관련 파일'은 본문에 링크됐거나(=ri:filename) 질의어와 일치하는 첨부만 (무관 첨부 제외)
                files = [f for f in files if (f.get("name") and (f["name"] in body or any(t.lower() in f["name"].lower() for t in toks)))]
                # 본문 이미지 위치에 마커 → 텍스트 추출 → [[IMG:n]]로 치환 + 인라인 이미지 다운로드
                text_marked = _html_to_text(_conf_mark_images(body))[:7000]
                text = _conf_inline_image_urls(text_marked, base, att_map)   # 이미지=공개URL 마크다운 인라인
                out.append({"name": "Confluence · " + title, "text": text, "images": [],
                            "files": files, "url": base + (((pg.get("_links") or {}).get("webui")) or "")})
    except Exception:
        return []
    # 의미 기반 리랭커로 적합도 재정렬 (제목/패턴 하드코딩 없이 모든 질문 유형에 일반 적용).
    # 리랭크 입력은 '제목 + 본문' — 제목 신호도 반영되어 인덱스/내비 페이지가 자연히 후순위로 밀림.
    if len(out) > 1:
        try:
            docs = [(str(h.get("name") or "").replace("Confluence · ", "") + "\n" + str(h.get("text") or ""))[:1800] for h in out]
            ranked = await _rerank(query, docs, len(out))
            if ranked:
                out = [out[i] for i, _sc in ranked if 0 <= i < len(out)]
        except Exception:
            pass
    # 식별 토큰(모델코드 등 영숫자 4자+) 포함 페이지 우선 — 일반 IR 부스트(특정 페이지 하드코딩 아님).
    # 안정 정렬이라 리랭크 의미순서는 그대로 두고, 식별어가 든 페이지만 앞으로 끌어올림.
    ents = [t for t in toks if _rq.match(r"^[A-Za-z0-9]{4,}$", t)]
    if ents:
        def _has_ent(h):
            blob = (str(h.get("name") or "") + " " + str(h.get("text") or "")).lower()
            return any(e.lower() in blob for e in ents)
        out.sort(key=lambda h: 0 if _has_ent(h) else 1)
    return out[:limit]

@app.post("/api/confluence/test")
async def conf_test(token: str = ""):
    cfg = _conf_cfg()
    base = str(cfg.get("base_url") or "").rstrip("/")
    if not base:
        return {"ok": False, "error": "base_url 미설정"}
    import httpx
    space = str(cfg.get("space_key") or "").strip()
    try:
        async with httpx.AsyncClient(timeout=30, verify=False) as client:
            if space:
                # 설정한 스페이스(kb)에 실제 접근 가능한지 직접 확인
                r = await client.get(base + f"/rest/api/space/{space}", headers=_conf_headers(cfg), auth=_conf_auth(cfg))
                if r.status_code == 200:
                    s = r.json()
                    # 페이지 수도 함께 확인
                    cnt = None
                    try:
                        pr = await client.get(base + "/rest/api/content", headers=_conf_headers(cfg), auth=_conf_auth(cfg), params={"spaceKey": space, "type": "page", "limit": 1})
                        if pr.status_code == 200:
                            cnt = (pr.json().get("size"))
                    except Exception:
                        pass
                    return {"ok": True, "space": {"key": s.get("key"), "name": s.get("name")}, "page_probe": cnt}
                # 접근 불가 → 접근 가능한 스페이스 목록 안내
                lr = await client.get(base + "/rest/api/space", headers=_conf_headers(cfg), auth=_conf_auth(cfg), params={"limit": 25})
                spaces = [{"key": x.get("key"), "name": x.get("name")} for x in (lr.json().get("results") or [])] if lr.status_code == 200 else []
                return {"ok": False, "error": f"설정한 스페이스 '{space}' 접근 불가 (HTTP {r.status_code}). 이 계정으로 접근 가능한 스페이스에서 골라 Space Key를 바꾸세요.", "spaces": spaces}
            # 스페이스 미지정 → 전체 목록
            r = await client.get(base + "/rest/api/space", headers=_conf_headers(cfg), auth=_conf_auth(cfg), params={"limit": 25})
            if r.status_code != 200:
                return {"ok": False, "error": f"HTTP {r.status_code}: {r.text[:200]}"}
            spaces = [{"key": s.get("key"), "name": s.get("name")} for s in (r.json().get("results") or [])]
            return {"ok": True, "spaces": spaces}
    except Exception as e:
        return {"ok": False, "error": str(e)[:200]}

@app.post("/api/confluence/sync")
async def conf_sync(payload: dict, token: str = ""):
    _require_admin(token)
    cfg = _conf_cfg()
    base = str(cfg.get("base_url") or "").rstrip("/")
    space = str((payload or {}).get("space_key") or cfg.get("space_key") or "").strip()
    if not base:
        return {"ok": False, "error": "base_url 미설정"}
    import httpx, base64 as _b64
    headers = _conf_headers(cfg); auth = _conf_auth(cfg)
    init_manuals_dir()
    saved = 0
    try:
        async with httpx.AsyncClient(timeout=90, verify=False) as client:
            start = 0; limit = 25; fetched = 0
            while fetched < 1000:
                params = {"type": "page", "limit": limit, "start": start, "expand": "body.storage,ancestors,space"}
                if space:
                    params["spaceKey"] = space
                r = await client.get(base + "/rest/api/content", headers=headers, auth=auth, params=params)
                if r.status_code != 200:
                    return {"ok": False, "error": f"목록 HTTP {r.status_code}: {r.text[:200]}", "saved": saved}
                results = r.json().get("results") or []
                if not results:
                    break
                for pg in results:
                    fetched += 1
                    pid = pg.get("id"); title = pg.get("title") or ""
                    body = (((pg.get("body") or {}).get("storage") or {}).get("value")) or ""
                    text = _html_to_text(body)
                    anc = pg.get("ancestors") or []
                    folder = (anc[0].get("title") if anc else None) or ((pg.get("space") or {}).get("name")) or space or "Confluence"
                    images = []
                    try:
                        ar = await client.get(base + f"/rest/api/content/{pid}/child/attachment", headers=headers, auth=auth, params={"limit": 25})
                        if ar.status_code == 200:
                            for att in (ar.json().get("results") or []):
                                mt = ((att.get("metadata") or {}).get("mediaType")) or ""
                                dl = ((att.get("_links") or {}).get("download")) or ""
                                if dl and "image" in str(mt):
                                    iu = (base + dl) if dl.startswith("/") else dl
                                    ir = await client.get(iu, headers=headers, auth=auth)
                                    if ir.status_code == 200 and len(ir.content) < 3_000_000:
                                        idx = len(images)
                                        images.append("data:" + str(mt) + ";base64," + _b64.b64encode(ir.content).decode())
                                        text += f"\n[[IMG:{idx}]]\n"
                                if len(images) >= 15:
                                    break
                    except Exception:
                        pass
                    if not text and not images:
                        continue
                    mid = "conf-" + str(pid)
                    await db.manuals_upsert(mid, {
                        "id": mid, "name": title, "text": text, "chars": len(text), "source": "Confluence",
                        "active": True, "folder": folder, "images": images,
                        "created_at": datetime.now().strftime("%Y-%m-%d"), "confluence_id": pid,
                        "url": base + (((pg.get("_links") or {}).get("webui")) or ""),
                    })
                    saved += 1
                start += limit
                if len(results) < limit:
                    break
    except Exception as e:
        return {"ok": False, "error": str(e)[:300], "saved": saved}
    _RAG_CACHE["corpus"] = None
    return {"ok": True, "saved": saved}

@app.get("/api/rag/info")
async def rag_info():
    corpus = await _manual_chunk_corpus_async()
    by = {}
    for c in corpus:
        by[c["name"]] = by.get(c["name"], 0) + 1
    return {"total_chunks": len(corpus), "by_manual": by, "chunk_size": RAG_CHUNK_SIZE, "overlap": RAG_CHUNK_OVERLAP}

@app.get("/api/rag/chunks")
async def rag_chunks(manual: str = "", limit: int = 800):
    # 특정 매뉴얼만 조회하는 경우: 전체 코퍼스 만들지 말고 그 매뉴얼만 SELECT
    # (전체 코퍼스는 매뉴얼 수십개의 data(JSONB, 이미지 base64 포함) 를 모두 가져와 무거움)
    lim = max(1, min(limit, 2000))
    if manual:
        async with db.pool().acquire() as c:
            row = await c.fetchrow(
                "SELECT data FROM manuals WHERE active=true AND data->>'name'=$1 LIMIT 1",
                manual,
            )
        if not row:
            return {"manual": manual, "count": 0, "chunks": []}
        d = row["data"] or {}
        # 청크는 캐시된 코퍼스 없이도 텍스트만 있으면 즉시 분할 가능
        chunks = _chunk_text(d.get("text", ""))
        return {"manual": manual, "count": len(chunks), "chunks": chunks[:lim]}
    # 전체 조회는 executor 로 (blocking 회피)
    corpus = await asyncio.get_event_loop().run_in_executor(None, _manual_chunk_corpus)
    sel = [c["text"] for c in corpus]
    return {"manual": manual, "count": len(sel), "chunks": sel[:lim]}

def _retrieve_knowledge(query: str, max_chars: int = 4500):
    """사용자 질문과 관련된 사내 지식을 검색해 컨텍스트 문자열로 반환 — 시험절차 학습 + 매뉴얼(RAG)."""
    import re as _re
    terms = [t for t in _re.split(r"\s+", (query or "").lower()) if len(t) >= 2]
    if not terms:
        return ""
    parts = []
    # 1) 시험절차 학습 데이터
    try:
        items = _load_learned().get("items", [])
        def _lscore(it):
            hay = (str(it.get("title", "")) + " " + " ".join(it.get("models") or []) + " " + str(it.get("role", "")) + " "
                   + " ".join((str(s.get("desc", "")) + " " + str(s.get("cli", "")) + " " + str(s.get("imageText", ""))) for s in (it.get("steps") or []))).lower()
            return sum(1 for t in terms if t in hay)
        lhits = sorted([it for it in items if _lscore(it) > 0], key=_lscore, reverse=True)[:4]
        if lhits:
            seg = "【시험절차 학습 데이터】\n"
            for it in lhits:
                seg += f"· 시험항목: {it.get('title','')} (모델 {','.join(it.get('models') or [])} / 제품군 {it.get('role','')})\n"
                for s in (it.get("steps") or []):
                    seg += f"   - {s.get('desc','')}: `{s.get('cli','')}` [판정 {s.get('type','')} \"{s.get('criteria','')}\"]"
                    if s.get("imageText"):
                        seg += f" / 이미지인식: {str(s.get('imageText'))[:150]}"
                    seg += "\n"
            parts.append(seg)
    except Exception:
        pass
    # 2) 매뉴얼/지식 문서 — 청킹 + BM25 검색 (RAG)
    try:
        hits = _bm25_search(query, _manual_chunk_corpus(), top_k=6)
        if hits:
            seg = "【매뉴얼/지식 문서 (RAG 발췌)】\n"
            for h in hits:
                _t = _re.sub(r"\[\[IMG:\d+\]\]", "", str(h.get("text") or ""))[:600]
                seg += f"· [{h['name']}] {_t}\n"
            parts.append(seg)
    except Exception:
        pass
    ctx = "\n\n".join(parts).strip()
    if len(ctx) > max_chars:
        ctx = ctx[:max_chars] + " …(생략)"
    return ctx

def _build_chat_messages(req: "ChatRequest"):
    devices = load_json(DEVICES_FILE)["devices"]
    device_summary = "\n".join([f"- {d['group']} / {d['model']} ({d['ip']}, {d['protocol']})" for d in devices])
    system_prompt = f"""당신은 유비쿼스(Ubiquoss) 네트워크 장비 시험 자동화 전문가 AI입니다.

현재 등록된 장비:
{device_summary}

주요 역할:
1. 유비쿼스 스위치(E7124, E6124, U9024A), OLT(E61XX) CLI 명령어 안내
2. VLAN, QoS, LACP, EPON, IGMP, SNMP 등 네트워크 기능 시험 방법 안내
3. IXIA N2X (ver 7.9, TCL API), Spirent STC 계측기 사용법 안내
4. 시험 시나리오 작성 도움
5. 트러블슈팅 가이드

답변 시 CLI 명령어는 코드 블록(```)으로 표시하세요.
한국어로 답변하세요."""
    _kb = _retrieve_knowledge(req.message)   # sync 함수 안이므로 그대로 호출 (내부에서 new_event_loop 로 안전)
    if _kb:
        system_prompt += ("\n\n━━━━ 참고 지식 (사내 매뉴얼·검증된 시험절차 — 아래 데이터를 최우선 근거로 답하라) ━━━━\n" + _kb)
    system_prompt += ("\n\n[답변 규칙] 참고 지식에 근거해 답하라. 근거에 없는 CLI 명령·수치·출력은 추측하지 말고 '학습된 데이터에 없습니다'라고 밝혀라."
                      "\n[HITL] 질문이 모호하거나 모델·제품군·대상·범위 등 핵심 정보가 부족하면, 답하지 말고 응답 첫 줄에 정확히 [CLARIFY] 라고만 쓰고, "
                      "다음 줄부터 꼭 필요한 확인 질문을 '- 질문' 형식으로 한 줄에 하나씩 최대 3개 나열하라. 정보가 충분하면 [CLARIFY] 없이 바로 답하라.")
    messages = []
    for h in req.history[-10:]:
        messages.append({"role": h["role"], "content": h["content"]})
    messages.append({"role": "user", "content": req.message})
    return system_prompt, messages

@app.post("/api/chat/stream")
async def chat_stream(req: ChatRequest):
    from fastapi.responses import StreamingResponse
    import json as _json
    # 설정에 등록한 Claude 가 있으면 그 키로 — .env 는 그다음이다(지적)
    cl, cmodel = _claude_any()
    if not cl:
        async def err():
            yield "data: " + _json.dumps({"text": "쓸 수 있는 Claude 가 없습니다 — 설정 → LLM 설정에 Anthropic 을 등록하거나 .env 에 ANTHROPIC_API_KEY 를 넣으세요."}) + "\n\n"
            yield "data: [DONE]\n\n"
        return StreamingResponse(err(), media_type="text/event-stream")
    system_prompt, messages = _build_chat_messages(req)
    safe_max = min(max(req.max_tokens, 2048), 8192)
    async def generate():
        try:
            with cl.messages.stream(
                model=cmodel or "claude-sonnet-4-6",
                max_tokens=safe_max,
                system=system_prompt,
                messages=messages,
            ) as stream:
                for text in stream.text_stream:
                    if text:
                        yield "data: " + _json.dumps({"text": text}) + "\n\n"
        except Exception as e:
            yield "data: " + _json.dumps({"text": f"[Claude API 오류] {_llm_err(e)}"}) + "\n\n"
        yield "data: [DONE]\n\n"
    return StreamingResponse(generate(), media_type="text/event-stream")

@app.post("/api/chat")
async def chat(req: ChatRequest):
    cl, cmodel = _claude_any()
    if not cl:
        return {"reply": "쓸 수 있는 Claude 가 없습니다.\n\n설정 → LLM 설정에서 Anthropic 을 등록하고 API 키를 넣거나,\nbackend 폴더와 같은 위치의 .env 에 ANTHROPIC_API_KEY=sk-ant-... 를 넣고 서버를 재시작하세요."}

    # 장비/절차 컨텍스트 로드
    devices = load_json(DEVICES_FILE)["devices"]
    device_summary = "\n".join([f"- {d['group']} / {d['model']} ({d['ip']}, {d['protocol']})" for d in devices])

    system_prompt = f"""당신은 유비쿼스(Ubiquoss) 네트워크 장비 시험 자동화 전문가 AI입니다.

현재 등록된 장비:
{device_summary}

주요 역할:
1. 유비쿼스 스위치(E7124, E6124, U9024A), OLT(E61XX) CLI 명령어 안내
2. VLAN, QoS, LACP, EPON, IGMP, SNMP 등 네트워크 기능 시험 방법 안내
3. IXIA N2X (ver 7.9, TCL API), Spirent STC 계측기 사용법 안내
4. 시험 시나리오 작성 도움
5. 트러블슈팅 가이드

답변 시 CLI 명령어는 코드 블록(```)으로 표시하세요.
한국어로 답변하세요."""

    messages = []
    for h in req.history[-10:]:
        messages.append({"role": h["role"], "content": h["content"]})
    messages.append({"role": "user", "content": req.message})

    try:
        response = cl.messages.create(
            model=cmodel or "claude-sonnet-4-6",
            max_tokens=min(max(req.max_tokens, 2048), 8192),
            system=system_prompt,
            messages=messages
        )
        reply = response.content[0].text
    except Exception as e:
        reply = f"[Claude API 오류] {_llm_err(e)}"

    return {"reply": reply}

# ───────────────────────────────────────────
# WebSocket
# ───────────────────────────────────────────
@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket):
    await websocket.accept()
    active_connections.append(websocket)
    ws_state[id(websocket)] = {"ws": websocket, "user": None, "page": None}
    try:
        while True:
            txt = await websocket.receive_text()
            try:
                msg = json.loads(txt)
            except Exception:
                continue
            if not isinstance(msg, dict):
                continue
            t = msg.get("type")
            st = ws_state.get(id(websocket))
            if t == "presence" and st is not None:
                old = st.get("page")
                if msg.get("user"):
                    st["user"] = str(msg.get("user"))
                st["page"] = msg.get("page")
                if old and old != st["page"]:
                    await _broadcast_presence(old)
                if st["page"]:
                    await _broadcast_presence(st["page"])
            elif t == "focus" and st is not None:
                # 같은 사이클 안에서 **어느 항목**을 보고 있나.
                #
                # 접속자(presence)는 화면 단위라, 사이클을 같이 보고 있다는
                # 것까지만 안다. 사이클은 항목을 나눠 돌리는 자리라 정작
                # 부딪히는 곳은 항목이다 — 둘이 같은 항목에 결과를 찍으면
                # 나중 사람이 앞사람 것을 덮는다.
                st["focus"] = msg.get("at")
                pg = msg.get("page") or st.get("page")
                if pg:
                    await _broadcast_focus(pg)
            elif t == "tc_running" and st is not None:
                # 실행 시작/끝을 모두에게 알린다. 연결이 끊기면(러너가 죽으면)
                # 아래 disconnect 에서 이 연결이 켠 것들을 자동으로 끈다.
                _tcid = str(msg.get("tcid") or "").strip()
                _on = bool(msg.get("on"))
                _u = str(msg.get("user") or st.get("user") or "").strip()
                if _tcid:
                    if _on:
                        _tc_running[_tcid] = {"user": _u, "at": _t.time()}
                        st.setdefault("running_tcs", set()).add(_tcid)
                    else:
                        _tc_running.pop(_tcid, None)
                        if isinstance(st.get("running_tcs"), set):
                            st["running_tcs"].discard(_tcid)
                    await _broadcast_tc_running(_tcid, _u, _on)
            elif t == "takeover":
                pg = msg.get("page"); u = msg.get("user")
                if pg and u:
                    page_controller[pg] = str(u)
                    await _broadcast_presence(pg)
    except WebSocketDisconnect:
        st = ws_state.pop(id(websocket), None)
        if websocket in active_connections:
            active_connections.remove(websocket)
        # 러너가 창을 닫거나 끊기면, 이 연결이 켠 「실행 중」을 자동으로 끈다 —
        # 아니면 유령 진행중이 남는다.
        for _rt in list((st or {}).get("running_tcs") or []):
            _tc_running.pop(_rt, None)
            try: await _broadcast_tc_running(_rt, "", False)
            except Exception: pass
        if st and st.get("page"):
            await _broadcast_presence(st["page"])
            await _broadcast_focus(st["page"])
    except Exception:
        ws_state.pop(id(websocket), None)
        if websocket in active_connections:
            try:
                active_connections.remove(websocket)
            except Exception:
                pass

# ───────────────────────────────────────────
# 전체 상태 체크 (백그라운드)
# ───────────────────────────────────────────
@app.get("/api/status")
async def get_status():
    data = load_json(DEVICES_FILE)
    summary = {"connected": 0, "disconnected": 0, "unknown": 0}
    for d in data["devices"]:
        s = d.get("status", "unknown")
        summary[s] = summary.get(s, 0) + 1
    return summary

# ===== regrafted: Confluence + Jira (UMS) =====
def _html_to_text(html_text: str) -> str:
    """HTML을 표 구조(탭/줄바꿈)를 살린 평문으로 변환 (bs4 없이 stdlib). Confluence는 main-content 영역만 추출."""
    import re as _re
    _m = _re.search(r'id=["\']main-content["\']', html_text)
    if _m:
        _gt = html_text.find('>', _m.start())
        if _gt >= 0:
            html_text = html_text[_gt + 1:]
    from html.parser import HTMLParser
    out = []

    class _P(HTMLParser):
        def __init__(self):
            super().__init__()
            self.skip = 0

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style", "noscript", "head"):
                self.skip += 1
            elif tag in ("br", "p", "div", "tr", "li", "h1", "h2", "h3", "h4", "h5", "h6", "table", "ul", "ol", "hr", "section"):
                out.append("\n")
            elif tag in ("td", "th"):
                out.append("\t")

        def handle_endtag(self, tag):
            if tag in ("script", "style", "noscript", "head") and self.skip > 0:
                self.skip -= 1

        def handle_data(self, data):
            if self.skip == 0:
                t = data.strip()
                if t:
                    out.append(t + " ")

    try:
        _P().feed(html_text)
    except Exception:
        pass
    lines = [ln.rstrip() for ln in "".join(out).splitlines()]
    cleaned, blank = [], 0
    for ln in lines:
        if ln.strip():
            cleaned.append(ln)
            blank = 0
        else:
            blank += 1
            if blank <= 1:
                cleaned.append("")
    return "\n".join(cleaned).strip()


@app.post("/api/confluence/fetch")
async def confluence_fetch(payload: dict):
    """REQ의 Confluence URL을 읽어 본문 텍스트를 반환 (TC 자동 생성용)."""
    import httpx
    url = str(payload.get("url", "")).strip()
    if not url.lower().startswith("http"):
        raise HTTPException(400, "올바른 URL이 아닙니다")
    cfg = {}
    if CONFLUENCE_FILE.exists():
        try:
            cfg = load_json(CONFLUENCE_FILE)
        except Exception:
            cfg = {}
    headers = {"User-Agent": "Mozilla/5.0 (U-TOP)", "Accept": "text/html,application/xhtml+xml"}
    auth = None
    if cfg.get("token"):
        headers["Authorization"] = "Bearer " + str(cfg["token"])
    elif cfg.get("username"):
        auth = (str(cfg["username"]), str(cfg.get("password", "")))
    try:
        async with httpx.AsyncClient(timeout=20, follow_redirects=True, verify=False) as client:
            resp = await client.get(url, headers=headers, auth=auth)
        if resp.status_code in (401, 403):
            raise HTTPException(502, f"위키 인증 필요 (HTTP {resp.status_code}). data/integrations/confluence.json에 username/password 또는 token을 설정하세요.")
        resp.raise_for_status()
        text = _html_to_text(resp.text)
        return {"success": True, "text": text[:18000], "length": len(text), "truncated": len(text) > 18000}
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(502, f"Confluence 읽기 실패: {exc}")


# ══════════════ Jira 연동 (Server 8.14, REST API v2) ══════════════
def _jira_cfg():
    if JIRA_FILE.exists():
        try:
            return load_json(JIRA_FILE)
        except Exception:
            return {}
    return {}

def _jira_headers(cfg):
    auth = (cfg.get("auth") or "basic").lower()
    tok = cfg.get("token") or ""
    user = cfg.get("user") or ""
    h = {"Content-Type": "application/json", "Accept": "application/json"}
    if auth == "bearer":
        h["Authorization"] = "Bearer " + tok        # Jira Server PAT
    else:
        import base64 as _b64
        h["Authorization"] = "Basic " + _b64.b64encode((user + ":" + tok).encode("utf-8")).decode("ascii")
    return h

def _jira_call(method, path, cfg=None, **kw):
    import httpx
    cfg = cfg or _jira_cfg()
    base = (cfg.get("url") or "").rstrip("/")
    if not base:
        return None, {"ok": False, "error": "Jira URL이 설정되지 않았습니다 (시스템 → Jira 연동 설정)"}
    try:
        with httpx.Client(timeout=25, verify=cfg.get("verify", True)) as c:
            r = c.request(method, base + path, headers=_jira_headers(cfg), **kw)
        return r, None
    except Exception as e:
        return None, {"ok": False, "error": str(e)[:300]}

@app.get("/api/jira/defect/schema")
async def api_defect_schema():
    """분류 스키마(드롭다운 옵션) 반환."""
    return {"ok": True, "device": _DEF_DEVICE, "category_field": _DEF_CAT_FIELD,
            "category_live": _DEF_CAT_LIVE, "item": _DEF_ITEM, "type3": _DEF_TYPE3}

@app.get("/api/jira/defect/class")
async def api_defect_class_get():
    """저장된 전체 분류 반환 {key: {...}}."""
    return {"ok": True, "classes": _load_defect_class()}

@app.post("/api/jira/defect/class")
async def api_defect_class_save(payload: dict):
    """수동 분류 저장/수정. payload={key, class:{source,device,category,item,type3}}."""
    key = str(payload.get("key") or "").strip()
    if not key:
        raise HTTPException(400, "이슈 키가 없습니다")
    store = _load_defect_class()
    cls = _defect_norm(payload.get("class") or {})
    cls["by"] = "manual"
    import datetime as _dt
    cls["at"] = _dt.datetime.now().strftime("%Y-%m-%d %H:%M")
    store[key] = cls
    _save_defect_class(store)
    return {"ok": True, "key": key, "class": cls}

@app.post("/api/jira/defect/classify")
async def api_defect_classify(payload: dict):
    """이슈키 목록을 LLM(제마)으로 자동 분류 → 저장. payload={keys:[...], overwrite:bool}."""
    keys = [str(k).strip() for k in (payload.get("keys") or []) if str(k).strip()]
    if not keys:
        raise HTTPException(400, "분류할 이슈가 없습니다")
    overwrite = bool(payload.get("overwrite"))
    store = _load_defect_class()
    todo = keys if overwrite else [k for k in keys if not (store.get(k) or {}).get("source")]
    if not todo:
        return {"ok": True, "classified": 0, "skipped": len(keys), "message": "이미 모두 분류됨(덮어쓰기 아님)"}
    # LLM 선택 (llms.json active + endpoint)
    llms = (load_json(LLMS_FILE).get("llms") or [])
    active = [l for l in llms if l.get("status", "active") == "active" and l.get("endpoint")]
    llm = next((l for l in active if str(l.get("type") or "").lower() not in ("claude", "anthropic") and "anthropic.com" not in str(l.get("endpoint", ""))), None) or (active[0] if active else None)
    if not llm:
        raise HTTPException(400, "사용 가능한 LLM이 없습니다 (시스템 → LLM 설정)")
    texts = _jira_issue_texts(todo)
    sys_p = (
        "너는 네트워크 장비 시험 이슈(Jira)를 defect로 분류하는 분류기다. 아래 스키마의 허용값 중에서만 골라 JSON으로만 답한다. 설명·코드펜스 금지.\n"
        + _defect_schema_text() +
        "판단 기준: 현장(운용망/고객사)에서 발생한 장애면 source=현장장애, 상용망 검증(BMT/사전검증)에서 발생하면 source=상용망검증. "
        "장비군은 L2/L3/FTTH(OLT/광가입자) 중 택1. 카테고리는 서비스/기능/운용/IPv6 등에서 가장 적합한 1개. "
        "상용망검증이면 item(RFP/표준Config/CR_Defect/UTS(SNMP)/부팅/반복Aging)과 type3도 함께 채운다. 모르면 빈 문자열.\n"
        '출력 형식: {"source":"","device":"","category":"","item":"","type3":""}'
    )
    import json as _json, datetime as _dt, re as _re
    now = _dt.datetime.now().strftime("%Y-%m-%d %H:%M")
    done = 0; fails = []
    for k in todo:
        txt = texts.get(k) or ""
        if not txt:
            fails.append(k); continue
        try:
            out = await _jira_llm_complete(llm, sys_p, "이슈:\n" + txt + "\n분류 JSON:", max_tokens=200, temp=0.0)
            m = _re.search(r"\{[\s\S]*\}", str(out or ""))
            cls = _defect_norm(_json.loads(m.group(0)) if m else {})
        except Exception:
            cls = _defect_norm({})
        if not cls.get("source"):
            fails.append(k)
        cls["by"] = "llm"; cls["at"] = now
        store[k] = cls; done += 1
    _save_defect_class(store)
    return {"ok": True, "classified": done, "failed": fails, "total": len(todo), "llm": llm.get("name") or llm.get("model") or ""}

@app.get("/api/jira/config")
async def jira_get_config():
    return _jira_cfg()

@app.post("/api/jira/config")
async def jira_save_config(data: dict):
    cur = _jira_cfg()
    for k in ["url", "user", "token", "auth", "default_project", "default_issuetype", "verify", "fav_projects", "ai", "panel_templates", "login_enabled", "login_auto_create", "login_url"]:
        if k in data:
            cur[k] = data[k]
    save_json(JIRA_FILE, cur)
    return {"ok": True}

def _jira_fetch_users(q: str = "", limit: int = 2000) -> tuple:
    """Jira 사용자 목록 — 활성·비활성을 함께 가져온다.

    조회 계정(연동 설정의 user/token)으로 부른다. 한 번에 다 안 오므로
    startAt 을 밀며 여러 번 받는다 — 200명 넘는 곳에서 첫 50명만 들어오면
    「왜 저 사람만 없나」 를 영원히 못 찾는다.
    """
    cfg = _jira_cfg()
    # 조회 계정이 없으면 Jira 는 **로그인 화면(HTML)** 을 401 로 돌려준다.
    # 그것을 그대로 화면에 뿌리면 「이건 뭐야」 가 된다(지적) — 먼저 막는다.
    if not str(cfg.get("user") or "").strip() or not str(cfg.get("token") or "").strip():
        return [], {"ok": False, "error": (
            "Jira 조회 계정이 없습니다 — SETUP → Jira 연동에서 아이디와 토큰(PAT)을 넣고 "
            "「연결 테스트」 가 통과한 뒤에 다시 누르세요. (로그인은 각자 비밀번호로 되지만, "
            "**명단을 통째로 읽는 것**은 조회 계정이 있어야 합니다)")}
    search = (q or cfg.get("user_search") or ALLOWED_EMAIL_DOMAIN or "ubiquoss.com").strip()
    out, seen, start = [], set(), 0
    while start < int(limit):
        r, err = _jira_call(
            "GET", "/rest/api/2/user/search",
            params={"username": search, "startAt": start, "maxResults": 200,
                    "includeActive": "true", "includeInactive": "true"},
        )
        if err:
            return [], err
        if not r.is_success:
            # Jira 는 실패를 **HTML 한 장**으로 준다. 사람이 읽을 한 줄로 바꾼다.
            why = {
                401: "Jira 가 조회 계정을 받지 않았습니다(401) — 아이디·토큰을 확인하세요. "
                     "Jira Server 라면 비밀번호 대신 PAT(개인 액세스 토큰)를 권합니다",
                403: "조회 계정에 사용자 조회 권한이 없습니다(403) — Jira 관리자에게 요청하세요",
                404: "Jira 주소가 잘못됐습니다(404) — REST 경로를 못 찾았습니다",
            }.get(r.status_code)
            if not why:
                body = re.sub(r"<[^>]+>", " ", r.text or "")
                body = " ".join(body.split())[:160]
                why = f"Jira 가 {r.status_code} 로 답했습니다 — {body}"
            return [], {"ok": False, "error": why}
        rows = r.json() or []
        if not rows:
            break
        for u in rows:
            name = str(u.get("name") or "").strip()
            if not name or name in seen:
                continue
            seen.add(name)
            _disp = str(u.get("displayName") or name).strip()
            # 이름 괄호에 부서가 든다: 「강경묵(생산)」·「권민수(검증)_중…」.
            # 소속 칸이 비어 있으면 이걸 채운다(직급/소속을 보고 싶다는 지적).
            _m = re.search(r"\(([^)]+)\)", _disp)
            out.append({
                "username": name,
                "jira_key": str(u.get("key") or u.get("accountId") or "").strip(),
                "name": _disp,
                "email": str(u.get("emailAddress") or "").strip(),
                "jira_active": bool(u.get("active")),
                "dept": (_m.group(1).strip() if _m else ""),
            })
        if len(rows) < 200:
            break
        start += 200
    return out, None


def _jira_leaders(cfg=None) -> set:
    """Jira 에서 **직급(리더)** 을 유추한다 — 팀장·그룹장.

    Jira 사용자 API 에는 직급 필드가 없다. 다만 「팀장」·「그룹장」 은 **그룹**
    으로 남아 있어(팀장·기술팀 팀장·연구소 팀장/그룹장·그룹장/담당 …), 그
    구성원을 읽으면 누가 리더인지 알 수 있다. 사원·선임·책임·수석은 그룹이
    없어 Jira 로는 가릴 수 없다 — 그건 못 채운다.

    비싼 짓(사람마다 조회)이 아니다: picker 로 리더 그룹 이름을 몇 개 찾고,
    그 그룹의 구성원만 읽는다(십수 번의 호출).
    """
    cfg = cfg or _jira_cfg()
    gnames = set()
    for q in ("팀장", "그룹장"):
        r, err = _jira_call("GET", "/rest/api/2/groups/picker",
                            cfg=cfg, params={"query": q, "maxResults": 50})
        if err or not r.is_success:
            continue
        for g in (r.json().get("groups") or []):
            nm = str(g.get("name") or "").strip()
            # 이름에 팀장/그룹장이 든 그룹만 — picker 는 느슨히 걸린다
            if nm and ("팀장" in nm or "그룹장" in nm):
                gnames.add(nm)
    leaders = set()
    for nm in gnames:
        start = 0
        while start < 2000:
            r, err = _jira_call("GET", "/rest/api/2/group/member", cfg=cfg,
                                params={"groupname": nm, "startAt": start,
                                        "maxResults": 200, "includeInactiveUsers": "true"})
            if err or not r.is_success:
                break
            j = r.json() or {}
            vals = j.get("values") or []
            for m in vals:
                un = str(m.get("name") or "").strip().lower()
                if un:
                    leaders.add(un)
            if j.get("isLast") or len(vals) < 200:
                break
            start += 200
    return leaders


@app.get("/api/users/jira-sync")
async def api_users_jira_sync_status(token: str = ""):
    """지난번 동기화가 언제·어떻게 됐나 (화면 머리줄)."""
    _require_admin(token)
    cfg = _jira_cfg()
    return {
        "ok": True,
        "url": str(cfg.get("url") or ""),
        "login_url": _jira_login_base(cfg),
        "user": str(cfg.get("user") or ""),
        "login_enabled": bool(cfg.get("login_enabled")),
        "auto_create": _jira_auto_create(),
        "last": cfg.get("last_user_sync") or None,
    }


@app.post("/api/users/jira-sync")
async def api_users_jira_sync(payload: dict = None, token: str = ""):
    """Jira 사용자를 UTOP 명단으로 **끌어온다.**

    ★ 관리자가 정한 것은 안 덮는다 — 역할(role)과 잠금(active)은 그대로 둔다.
      Jira 가 정본인 것은 **누가 있는가·이름·메일·Jira 활성**까지다.
    ★ 비밀번호는 여기서도 없다. Jira 가 갖고 있다.
    """
    _require_admin(token)
    q = str((payload or {}).get("search") or "").strip()
    rows, err = await asyncio.to_thread(_jira_fetch_users, q)
    if err:
        return {"ok": False, **err}
    # 직급(리더) 유추 — 팀장·그룹장 그룹의 구성원. 못 읽어도 동기화는 계속한다.
    try:
        leaders = await asyncio.to_thread(_jira_leaders)
    except Exception:
        leaders = set()
    data = _users_load_sync()
    by_name = {str(u.get("username") or "").lower(): u for u in data["users"]}
    by_key = {str(u.get("jira_key") or ""): u for u in data["users"] if u.get("jira_key")}
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    new_n = chg_n = off_n = back_n = lead_n = 0
    for j in rows:
        cur = by_key.get(j["jira_key"]) if j["jira_key"] else None
        if cur is None:
            cur = by_name.get(j["username"].lower())
        if cur is None:
            # **비활성(퇴사자) 신규는 안 만든다**(지시: 퇴사자 필요 없음).
            # 명단에 없고 이미 Jira 에서 나간 사람은 애초에 추가하지 않는다 —
            # 안 그러면 지워도 다음 동기화마다 되살아난다.
            if not j["jira_active"]:
                continue
            _is_lead = j["username"].lower() in leaders
            nu = {
                "id": j["username"], "username": j["username"], "name": j["name"],
                "role": "팀장" if _is_lead else "팀원", "email": j["email"], "active": True,
                "dept": j.get("dept") or "",
                "source": "jira",
                "jira_key": j["jira_key"], "jira_active": True,
                "created_at": now, "synced_at": now,
            }
            if _is_lead:
                nu["role_by"] = "jira"
                lead_n += 1
            data["users"].append(nu)
            new_n += 1
            continue
        before = (cur.get("name"), cur.get("email"), cur.get("jira_active"), cur.get("jira_key"))
        if j["name"]:
            cur["name"] = j["name"]
        if j["email"]:
            cur["email"] = j["email"]
        if j["jira_key"]:
            cur["jira_key"] = j["jira_key"]
        if j.get("dept") and not str(cur.get("dept") or "").strip():
            cur["dept"] = j["dept"]   # 비어 있을 때만 — 관리자가 정한 소속은 그대로
        # 직급(리더) — Jira 의 팀장/그룹장 그룹이 정본. **관리자가 손으로 정한
        # 역할(관리자·담당, 또는 손으로 준 팀장)은 안 건드린다** — jira 가 준
        # 것(role_by=jira)만 올리고 내린다. 이 규칙은 active 의 locked_by 와 같다.
        _is_lead = str(cur.get("username") or "").lower() in leaders
        if _is_lead and cur.get("role") == "팀원":
            cur["role"] = "팀장"; cur["role_by"] = "jira"; lead_n += 1
        elif not _is_lead and cur.get("role") == "팀장" and cur.get("role_by") == "jira":
            cur["role"] = "팀원"; cur.pop("role_by", None)
        cur["jira_active"] = j["jira_active"]
        cur["source"] = "jira"
        cur["synced_at"] = now
        """
        Jira 에서 나간 사람은 **여기서도 잠근다.**

        퇴사자가 명단에 활성으로 남아 있으면 「누가 들어올 수 있나」 가 틀린
        답을 준다. 다만 **우리가 잠근 것만** 되돌린다(locked_by) — 관리자가
        따로 잠근 사람을 Jira 가 살아났다고 풀어 주면 안 된다.
        """
        if not j["jira_active"] and cur.get("active", True):
            cur["active"] = False
            cur["locked_by"] = "jira"
            off_n += 1
        elif j["jira_active"] and cur.get("active") is False and cur.get("locked_by") == "jira":
            cur["active"] = True
            cur.pop("locked_by", None)
            back_n += 1
        if before != (cur.get("name"), cur.get("email"), cur.get("jira_active"), cur.get("jira_key")):
            chg_n += 1
    _users_save_sync(data)
    stat = {"at": now, "found": len(rows), "new": new_n, "changed": chg_n,
            "locked": off_n, "unlocked": back_n, "leads": lead_n,
            "active": len([x for x in rows if x["jira_active"]]),
            "inactive": len([x for x in rows if not x["jira_active"]])}
    cfg = _jira_cfg()
    cfg["last_user_sync"] = stat
    save_json(JIRA_FILE, cfg)
    return {"ok": True, **stat}


@app.post("/api/jira/login-test")
async def api_jira_login_test(payload: dict, token: str = ""):
    """**Jira 계정으로 로그인이 되는지** 관리자 자리에서 확인한다.

    화면에서 아이디·비밀번호를 넣어 눌러 보는 것 말고는 「왜 저 사람은 안
    되나」 를 알 길이 없었다. 비밀번호는 확인에만 쓰고 **어디에도 담지 않는다** —
    저장도, 로그도 안 한다. 여기서 성공해도 세션은 안 만든다.
    """
    _require_admin(token)
    uname = str(payload.get("username") or "").strip()
    pw = str(payload.get("password") or "")
    if not uname or not pw:
        return {"ok": False, "error": "아이디와 비밀번호를 넣으세요"}
    if not _jira_login_on():
        return {"ok": False, "error": "Jira 계정 로그인이 꺼져 있습니다 — 위에서 켜세요"}
    ju, why = await _jira_verify_login(uname, pw)
    if ju:
        known = _find_user(uname)
        return {"ok": True, "jira": {
            "username": str(ju.get("name") or uname),
            "key": str(ju.get("key") or ju.get("accountId") or ""),
            "name": str(ju.get("displayName") or ""),
            "email": str(ju.get("emailAddress") or ""),
            "active": ju.get("active"),
        }, "in_utop": bool(known), "auto_create": _jira_auto_create()}
    msg = {
        "denied": "Jira 가 아이디·비밀번호를 받지 않았습니다",
        "captcha": "Jira 가 CAPTCHA 를 걸었습니다 — 그 계정으로 Jira 웹에 한 번 로그인해 푸세요",
        "cert": "Jira 인증서 문제(만료 등) — 「TLS 인증서 검증」 을 끄거나 인증서를 갱신하세요",
        "unreachable": "Jira 에 닿지 못했습니다",
        "no-url": "Jira 주소가 없습니다",
    }.get(why, why or "확인하지 못했습니다")
    return {"ok": False, "why": why, "error": msg}


@app.get("/api/jira/login-check")
async def jira_login_check(token: str = ""):
    """**Jira 로그인이 지금 되는 상태인가** — 계정 관리 화면이 묻는다.

    「Jira 계정으로 로그인이 안 된다」 는 말은 셋 중 하나다: 꺼져 있거나,
    주소가 없거나, Jira 가 거절하거나. 셋을 갈라 보여 주지 않으면 어디를
    고쳐야 하는지 알 수 없다. 비밀번호는 여기에 없다.
    """
    _require_admin(token)
    cfg = _jira_cfg()
    url = _jira_login_base(cfg)
    out = {
        "enabled": bool(cfg.get("login_enabled")),
        "url": url,
        "issue_url": str(cfg.get("url") or "").strip().rstrip("/"),
        "separate": bool(str(cfg.get("login_url") or "").strip()),
        "auto_create": _jira_auto_create(),
        "last_fail": dict(_JIRA_LAST_FAIL) or None,
        # Jira Cloud 는 계정 비밀번호로 REST 인증이 안 된다 — 그것을 모르면
        # 「비밀번호가 맞는데 왜 안 되나」 를 끝없이 헤맨다
        "cloud": "atlassian.net" in url.lower(),
    }
    if not url:
        out["reachable"] = False
        out["reason"] = "Jira 주소가 없습니다 — 「Jira 연동」 에서 먼저 넣으세요"
        return out
    out["verify"] = _jira_verify_flag(cfg)
    import httpx as _hx
    try:
        async with _hx.AsyncClient(timeout=8, verify=out["verify"]) as c:
            r = await c.get(url + "/rest/api/2/serverInfo")
        out["reachable"] = r.status_code < 500
        out["status"] = r.status_code
        out["reason"] = "" if r.status_code < 500 else f"Jira 가 {r.status_code} 로 답했습니다"
    except Exception as exc:
        msg = str(exc)
        out["reachable"] = False
        if "CERTIFICATE" in msg.upper() or "SSL" in msg.upper():
            out["cert"] = True
            out["reason"] = (
                "Jira 인증서에 문제가 있습니다(만료 등) — 위 「TLS 인증서 검증」 을 끄거나 "
                "인증서를 갱신하세요"
            )
        else:
            out["reason"] = f"닿지 못했습니다 — {msg[:120]}"
    return out


@app.post("/api/jira/test")
async def jira_test(data: dict = None):
    cfg = data if (data and data.get("url")) else _jira_cfg()
    r, err = _jira_call("GET", "/rest/api/2/myself", cfg=cfg)
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:200]}"}
    j = r.json()
    return {"ok": True, "name": j.get("name"), "displayName": j.get("displayName"), "email": j.get("emailAddress")}

@app.get("/api/jira/user-search")
async def jira_user_search(q: str = "", project: str = "", limit: int = 100):
    """담당자 목록/검색. project 지정 시 그 프로젝트에 할당 가능한 사용자 전체(assignable),
    q 지정 시 이름/메일/ID 부분일치 필터 (Jira user/assignable/search · user/search 프록시)."""
    mx = max(1, min(int(limit or 100), 500))
    qs = str(q).strip()
    if project:
        params = {"project": project, "maxResults": mx}
        if qs:
            params["username"] = qs
        r, err = _jira_call("GET", "/rest/api/2/user/assignable/search", params=params)
        if (err or (r is not None and not r.is_success)) and not qs:
            # 일부 Jira 버전은 username 파라미터 필수 → '.'(대부분 메일에 포함)로 폴백
            r, err = _jira_call("GET", "/rest/api/2/user/assignable/search",
                                params={"project": project, "username": ".", "maxResults": mx})
    else:
        if not qs:
            return {"ok": True, "users": []}
        r, err = _jira_call("GET", "/rest/api/2/user/search", params={"username": qs, "maxResults": mx})
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:200]}"}
    return {"ok": True, "users": [{"name": u.get("name", ""), "displayName": u.get("displayName", ""),
                                   "email": u.get("emailAddress", "")} for u in (r.json() or [])]}

@app.get("/api/jira/components")
async def jira_components(project: str):
    """프로젝트 구성요소 + 컴포넌트 리드(기본 담당자) 목록 — 구성요소 선택 시 담당자 자동 지정용."""
    r, err = _jira_call("GET", f"/rest/api/2/project/{project}/components")
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:200]}"}
    out = []
    for c in (r.json() or []):
        lead = c.get("lead") or {}
        out.append({"id": str(c.get("id", "")), "name": c.get("name", ""),
                    "lead": lead.get("name", ""), "leadDisplay": lead.get("displayName", "")})
    return {"ok": True, "components": out}

@app.get("/api/jira/projects")
async def jira_projects(expand: str = ""):
    # expand=description → /project 리스트에 description 필드도 함께 조회 (Jira 8+)
    _url = "/rest/api/2/project" + ("?expand=description" if "description" in (expand or "") else "")
    r, err = _jira_call("GET", _url)
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:200]}"}
    _out = []
    for p in r.json():
        _out.append({"key": p.get("key"), "name": p.get("name"), "id": p.get("id"),
                     "description": (p.get("description") or "")})
    return {"ok": True, "projects": _out}

@app.get("/api/jira/issuetypes")
async def jira_issuetypes(project: str):
    r, err = _jira_call("GET", f"/rest/api/2/issue/createmeta?projectKeys={project}&expand=projects.issuetypes")
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:200]}"}
    types = []
    for p in (r.json().get("projects") or []):
        for t in (p.get("issuetypes") or []):
            types.append({"id": t.get("id"), "name": t.get("name"), "subtask": t.get("subtask", False)})
    return {"ok": True, "issuetypes": types}

@app.get("/api/jira/createmeta")
async def jira_createmeta(project: str, issuetype: str = None):
    r, err = _jira_call("GET", f"/rest/api/2/issue/createmeta?projectKeys={project}&expand=projects.issuetypes.fields")
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:200]}"}
    projs = r.json().get("projects") or []
    if not projs:
        return {"ok": False, "error": "프로젝트 접근 불가 또는 없음(권한 확인)"}
    out = []
    for it in projs[0].get("issuetypes", []):
        if issuetype and str(it.get("id")) != str(issuetype) and it.get("name") != issuetype:
            continue
        for fid, f in (it.get("fields") or {}).items():
            sch = f.get("schema") or {}
            av = f.get("allowedValues")
            opts = None
            if isinstance(av, list):
                opts = [{"id": o.get("id"), "name": (o.get("name") or o.get("value") or o.get("key") or str(o.get("id") or ""))} for o in av]
            out.append({
                "id": fid, "name": f.get("name"), "required": bool(f.get("required")),
                "type": sch.get("type"), "items": sch.get("items"), "custom": sch.get("custom"),
                "options": opts, "hasDefault": f.get("hasDefaultValue", False),
            })
        break
    return {"ok": True, "fields": out}

@app.post("/api/jira/issue")
async def jira_create_issue(data: dict):
    project = data.get("project")
    itype = data.get("issuetype")
    summary = (data.get("summary") or "(제목 없음)")[:250]
    if not project or not itype:
        return {"ok": False, "error": "프로젝트/이슈유형이 필요합니다"}
    it = {"id": str(itype)} if str(itype).isdigit() else {"name": str(itype)}
    fields = {"project": {"key": project}, "issuetype": it, "summary": summary, "description": data.get("description") or ""}
    if data.get("labels"):
        fields["labels"] = data["labels"]
    if data.get("priority"):
        fields["priority"] = {"name": data["priority"]}
    if isinstance(data.get("fields"), dict):
        fields.update(data["fields"])
    r, err = _jira_call("POST", "/rest/api/2/issue", json={"fields": fields})
    if err:
        return err
    dropped = []
    if not r.is_success:
        # 생성 화면에 없는/알 수 없는 필드(labels 등)는 자동 제거 후 1회 재시도
        try:
            bad = list((r.json().get("errors") or {}).keys())
        except Exception:
            bad = []
        removable = [k for k in bad if k in fields and k not in ("project", "issuetype", "summary")]
        if removable:
            for k in removable:
                fields.pop(k, None)
            dropped = removable
            r, err = _jira_call("POST", "/rest/api/2/issue", json={"fields": fields})
            if err:
                return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:400]}"}
    j = r.json()
    cfg = _jira_cfg()
    key = j.get("key", "")
    return {"ok": True, "key": key, "url": (cfg.get("url", "").rstrip("/") + "/browse/" + key), "dropped": dropped}

@app.post("/api/jira/issue/{key}/attach")
async def jira_attach(key: str, data: dict):
    import base64 as _b64, httpx
    cfg = _jira_cfg()
    base = (cfg.get("url") or "").rstrip("/")
    if not base:
        return {"ok": False, "error": "Jira URL 미설정"}
    raw = data.get("data") or ""
    if raw.strip().startswith("data:") and "," in raw:
        raw = raw.split(",", 1)[1]
    try:
        content = _b64.b64decode(raw)
    except Exception as e:
        return {"ok": False, "error": "이미지 디코드 실패: " + str(e)[:120]}
    fn = data.get("filename") or "구성도.png"
    mime = data.get("mime") or "image/png"   # txt 첨부(running-config 등)도 지원
    h = _jira_headers(cfg)
    h.pop("Content-Type", None)        # multipart 경계는 httpx가 설정
    h["X-Atlassian-Token"] = "no-check"
    try:
        with httpx.Client(timeout=40, verify=cfg.get("verify", True)) as c:
            r = c.post(base + f"/rest/api/2/issue/{key}/attachments", headers=h,
                       files={"file": (fn, content, mime)})
    except Exception as e:
        return {"ok": False, "error": str(e)[:300]}
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:300]}"}
    return {"ok": True, "attachments": [a.get("filename") for a in r.json()]}


# Release Summary — 이슈 상세(설명+댓글+첨부) lazy load
@app.get("/api/jira/issue/{key}")
async def jira_issue_detail(key: str):
    r, err = _jira_call("GET", f"/rest/api/2/issue/{key}",
                        params={"fields": "summary,description,comment,status,assignee,reporter,issuetype,priority,fixVersions,updated,attachment",
                                "expand": "renderedFields"})   # renderedFields = Jira가 렌더한 HTML(설명·댓글) → UI 동일 표현
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:300]}"}
    return {"ok": True, **r.json()}

@app.post("/api/jira/issue/{key}/comment")
async def jira_add_comment(key: str, data: dict):
    body = data.get("body", "")
    if not body:
        return {"ok": False, "error": "body is empty"}
    # 계정 오버라이드: user/pw가 오면 그 계정(basic 인증)으로 등록 — 없으면 기존 설정 계정
    cfg = None
    _u = str(data.get("user") or "").strip()
    _p = str(data.get("pw") or data.get("password") or "")
    if _u and _p:
        cfg = {**_jira_cfg(), "auth": "basic", "user": _u, "token": _p}
    r, err = _jira_call("POST", f"/rest/api/2/issue/{key}/comment", cfg=cfg, json={"body": body})
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:300]}"}
    return {"ok": True, "comment": r.json()}

# Release Summary — Jira 첨부(이미지) 인증 프록시 (브라우저가 직접 못 받으므로 백엔드가 인증해서 중계)
@app.get("/api/jira/attachment")
async def jira_attachment(url: str):
    import httpx
    from fastapi import Response
    cfg = _jira_cfg()
    base = (cfg.get("url") or "").rstrip("/")
    if not base:
        return Response(content=b"", status_code=400)
    if url.startswith("/"):
        url = base + url
    if not url.startswith(base):           # 보안: 설정된 Jira 호스트만 프록시 허용
        return Response(content=b"", status_code=403)
    try:
        with httpx.Client(timeout=30, verify=cfg.get("verify", True), follow_redirects=True) as c:
            rr = c.get(url, headers=_jira_headers(cfg))
        ct = rr.headers.get("content-type", "application/octet-stream")
        return Response(content=rr.content, media_type=ct)
    except Exception as e:
        return Response(content=str(e)[:200].encode(), status_code=502)

# Release Summary — 트랙2 데이터 저장(이슈·TC·스텝·판정)
RELEASE_SUMMARY_FILE = DATA_DIR / "state" / "release_summary.json"
@app.get("/api/release-summary")
async def release_summary_get():
    d = _kv_load_sync("release_summary", {"releases": []})
    return d if isinstance(d, dict) else {"releases": []}

@app.post("/api/release-summary")
async def release_summary_save(data: dict):
    _kv_save_sync("release_summary", data or {"releases": []})
    return {"ok": True}

async def _jira_llm_complete(llm, sys_p, user_p, max_tokens=400, temp=0.0):
    """LLM 1회 호출 → 텍스트 반환 (JQL 자동생성 등 보조용). 실패 시 ''."""
    import httpx as _hx
    if not llm: return ""
    ltype = str(llm.get("type") or "").lower(); ep = str(llm.get("endpoint") or "")
    try:
        if ltype in ("claude", "anthropic") or "anthropic.com" in ep:
            import anthropic as _ah
            m = _ah.Anthropic(api_key=llm.get("apikey") or "").messages.create(
                model=llm.get("model") or "claude-sonnet-4-6", max_tokens=max_tokens, system=sys_p,
                messages=[{"role": "user", "content": user_p}])
            return "".join(getattr(b, "text", "") for b in m.content).strip()
        body = {"model": llm.get("model") or "", "messages": [{"role": "system", "content": sys_p}, {"role": "user", "content": user_p}], "temperature": temp, "max_tokens": max_tokens}
        headers = {"Content-Type": "application/json"}; ak = llm.get("apikey")
        if ak and not str(ak).lower().startswith("http"): headers["Authorization"] = f"Bearer {ak}"
        async with _hx.AsyncClient(timeout=120) as client:
            rr = await client.post(ep.rstrip("/") + "/chat/completions", headers=headers, json=body)
            if rr.status_code == 200:
                return (((rr.json().get("choices") or [{}])[0].get("message") or {}).get("content") or "").strip()
    except Exception:
        return ""
    return ""

# ══════════════ Issue Sync · Defect 분류 (현장장애 / 상용망검증) ══════════════
# 발생상황: 현장장애 | 상용망검증  (LLM이 이슈 내용으로 판단)
#  · 현장장애 → device(L2/L3/FTTH) + category(서비스·기능·운용·IPv6)
#  · 상용망검증 → device(L2/L3/FTTH) + category(서비스·기능·운용·IPv6·신규기능·신규기능Side)
#              AND item(RFP·표준Config·CR_Defect·UTS(SNMP)·부팅·반복Aging) + type3(서비스·기능·운용·IPv6·BMS·신규기능Side)
_DEF_DEVICE = ["L2", "L3", "FTTH"]
_DEF_CAT_FIELD = ["서비스", "기능", "운용", "IPv6"]          # 현장장애 카테고리
_DEF_CAT_LIVE = ["서비스", "기능", "운용", "IPv6", "신규기능", "신규기능Side"]  # 상용망 카테고리
_DEF_ITEM = ["RFP", "표준Config", "CR_Defect", "UTS(SNMP)", "부팅", "반복Aging"]
_DEF_TYPE3 = ["서비스", "기능", "운용", "IPv6", "BMS", "신규기능Side"]

def _load_defect_class() -> dict:
    try:
        if DEFECT_CLASS_FILE.exists():
            return load_json(DEFECT_CLASS_FILE) or {}
    except Exception:
        pass
    return {}

def _save_defect_class(d: dict):
    save_json(DEFECT_CLASS_FILE, d or {})

def _defect_schema_text():
    return (
        "분류 체계(값은 아래 목록 중에서만 선택):\n"
        "- source(발생상황): 현장장애 | 상용망검증\n"
        "- device(유형/장비): " + " | ".join(_DEF_DEVICE) + "\n"
        "- category(카테고리): " + " | ".join(_DEF_CAT_LIVE) + "  (단, 현장장애는 " + " | ".join(_DEF_CAT_FIELD) + " 중에서만)\n"
        "- item(상용망 항목, 상용망검증일 때만): " + " | ".join(_DEF_ITEM) + "\n"
        "- type3(상용망 유형, 상용망검증일 때만): " + " | ".join(_DEF_TYPE3) + "\n"
    )

def _defect_norm(cls: dict) -> dict:
    """LLM/수동 분류값을 스키마에 맞게 정규화(허용값 외/누락은 '')."""
    def pick(v, allow):
        v = str(v or "").strip()
        for a in allow:
            if v == a or v.lower() == a.lower():
                return a
        return ""
    src = pick((cls or {}).get("source"), ["현장장애", "상용망검증"])
    dev = pick((cls or {}).get("device"), _DEF_DEVICE)
    cat = pick((cls or {}).get("category"), _DEF_CAT_LIVE if src == "상용망검증" else _DEF_CAT_FIELD)
    out = {"source": src, "device": dev, "category": cat}
    if src == "상용망검증":
        out["item"] = pick((cls or {}).get("item"), _DEF_ITEM)
        out["type3"] = pick((cls or {}).get("type3"), _DEF_TYPE3)
    else:
        out["item"] = ""; out["type3"] = ""
    return out

def _jira_issue_texts(keys):
    """이슈키 목록 → {key: '제목 + 설명 + 댓글요약'} (LLM 분류 입력용). Jira에서 fetch."""
    out = {}
    for k in keys:
        try:
            r, err = _jira_call("GET", "/rest/api/2/issue/" + str(k) + "?fields=summary,description,issuetype,labels,components,comment")
            if err or r is None or r.status_code != 200:
                out[k] = ""; continue
            f = (r.json().get("fields") or {})
            summ = str(f.get("summary") or "")
            desc = str(f.get("description") or "")[:1500]
            itype = str(((f.get("issuetype") or {}).get("name")) or "")
            labels = ", ".join([str(x) for x in (f.get("labels") or [])])
            comps = ", ".join([str((c or {}).get("name") or "") for c in (f.get("components") or [])])
            cmts = (f.get("comment") or {}).get("comments") or []
            cmt_txt = " / ".join([str((c or {}).get("body") or "")[:200] for c in cmts[:3]])
            out[k] = (f"[제목] {summ}\n[유형] {itype}\n[컴포넌트] {comps}\n[라벨] {labels}\n[설명] {desc}\n[댓글] {cmt_txt}").strip()
        except Exception:
            out[k] = ""
    return out

# ── Jira 메타데이터(프로젝트·이슈유형·상태·모델매핑) 조회·캐시 — JQL 생성 정확도용 (하드코딩 없음, 10분 TTL) ──
_JIRA_META_CACHE = {"ts": 0.0, "data": None}
def _jira_meta(cfg=None):
    import time as _t
    if _JIRA_META_CACHE["data"] is not None and (_t.time() - _JIRA_META_CACHE["ts"]) < 600:
        return _JIRA_META_CACHE["data"]
    meta = {"projects": [], "issuetypes": [], "statuses": [], "model_proj": {}}
    cfg = cfg or _jira_cfg()
    try:
        r, err = _jira_call("GET", "/rest/api/2/project", cfg=cfg)
        if (r is not None) and r.is_success:
            meta["projects"] = [{"key": str(p.get("key") or ""), "name": str(p.get("name") or "")} for p in (r.json() or []) if p.get("key")][:80]
    except Exception:
        pass
    try:
        r, err = _jira_call("GET", "/rest/api/2/issuetype", cfg=cfg)
        if (r is not None) and r.is_success:
            seen = []
            for it in (r.json() or []):
                n = str(it.get("name") or "")
                if n and n not in seen:
                    seen.append(n)
            meta["issuetypes"] = seen[:40]
    except Exception:
        pass
    try:
        r, err = _jira_call("GET", "/rest/api/2/status", cfg=cfg)
        if (r is not None) and r.is_success:
            seen = []
            for st in (r.json() or []):
                n = str(st.get("name") or "")
                if n and n not in seen:
                    seen.append(n)
            meta["statuses"] = seen[:60]
    except Exception:
        pass
    # 모델명 → 프로젝트 키 (Jira 프로젝트 패널 설정의 이슈 키 매핑 auto_models — 설정 데이터)
    try:
        pts = cfg.get("panel_templates") or {}
        for pk, t in pts.items():
            for m in ((t or {}).get("auto_models") or []):
                m = str(m).strip()
                if m:
                    meta["model_proj"][m.upper()] = str(pk)
    except Exception:
        pass
    _JIRA_META_CACHE["ts"] = _t.time(); _JIRA_META_CACHE["data"] = meta
    return meta

# 담당자/보고자 이름 → Jira username 리졸브 (user search API, 세션 캐시)
_JIRA_USER_CACHE = {}
def _jira_resolve_user(name, cfg=None):
    key = str(name or "").strip()
    if not key:
        return name
    if key in _JIRA_USER_CACHE:
        return _JIRA_USER_CACHE[key]
    out = key
    try:
        r, err = _jira_call("GET", "/rest/api/2/user/search", cfg=cfg, params={"username": key, "maxResults": 5})
        if (r is not None) and r.is_success:
            arr = r.json() or []
            if arr:
                out = str(arr[0].get("name") or key)
    except Exception:
        pass
    _JIRA_USER_CACHE[key] = out
    return out

def _jira_fix_assignees(jql, cfg=None):
    """assignee/reporter 값이 한글 표시명이면 username으로 치환 (ASCII id는 그대로)."""
    import re as _r
    def _need(v):
        return any(ord(ch) > 127 for ch in v) or (" " in v.strip())
    def _rep_in(m):
        parts = [p.strip().strip('"').strip("'") for p in m.group(2).split(",")]
        rs = [(_jira_resolve_user(p, cfg) if _need(p) else p) for p in parts if p]
        return m.group(1) + " in (" + ", ".join('"%s"' % x for x in rs) + ")"
    def _rep_eq(m):
        v = m.group(2)
        return m.group(1) + ' = "' + (_jira_resolve_user(v, cfg) if _need(v) else v) + '"'
    jql = _r.sub(r'\b(assignee|reporter)\s+in\s*\(([^)]*)\)', _rep_in, jql, flags=_r.I)
    jql = _r.sub(r'\b(assignee|reporter)\s*=\s*"([^"]+)"', _rep_eq, jql, flags=_r.I)
    return jql

def _jira_expand_text(jql, question):
    """text ~ "키워드"를 붙임/띄움 두 표기 OR로 확장 — 띄어쓰기에 따라 결과가 갈리는 문제 방지.
    LLM 이 이미 확장한 경우(같은 확장이 원본 JQL 에 이미 존재)엔 재확장하지 않아 중복 그룹 생성 방지."""
    import re as _r
    toks = [t for t in _r.split(r"\s+", str(question or "")) if t]
    # JQL 안에 이미 등장한 text ~ 값을 모아 중복 확장 판정에 쓴다
    _existing = set(_r.findall(r'text\s*~\s*"([^"]+)"', jql))
    def _rep(m):
        term = m.group(1).strip()
        var = {term}
        if " " in term:
            var.add(term.replace(" ", ""))
        else:
            # 질문에서 "A B"로 띄어 쓴 연속 토큰의 결합이 이 키워드와 같으면 띄운 표기도 추가
            for i in range(len(toks) - 1):
                if (toks[i] + toks[i + 1]) == term:
                    var.add(toks[i] + " " + toks[i + 1])
        if len(var) <= 1:
            return m.group(0)
        # 이미 다른 표기(붙임/띄움)가 JQL 안에 존재 → LLM 이 이미 확장한 상태이므로 재확장 스킵
        if any((v != term and v in _existing) for v in var):
            return m.group(0)
        return "(" + " OR ".join('text ~ "%s"' % v for v in sorted(var)) + ")"
    return _r.sub(r'text\s*~\s*"([^"]+)"', _rep, jql)

async def _jira_gen_jql(q, llm, proj, meta=None):
    """자연어 질문 → LLM이 Jira JQL 생성. 실제 Jira 값(프로젝트·이슈유형·상태·모델매핑)을 주입해 필터 정확도 확보. 실패 시 ''."""
    import re as _r, datetime as _dt
    today = _dt.date.today().isoformat()
    meta = meta or {}
    _mLines = ""
    if meta.get("projects"):
        _mLines += "\n[프로젝트 목록 (key:이름) — project 조건은 반드시 이 key만 사용]\n" + ", ".join((p["key"] + ":" + p["name"]) for p in meta["projects"])
    if meta.get("model_proj"):
        _mLines += "\n[모델→프로젝트 매핑 — 질문에 모델명이 있으면 해당 project 조건 포함]\n" + ", ".join((m + "→" + k) for m, k in sorted(meta["model_proj"].items()))
    if meta.get("issuetypes"):
        _mLines += "\n[이슈유형 목록 — issuetype 조건은 이 값만 사용]\n" + ", ".join(meta["issuetypes"])
    if meta.get("statuses"):
        _mLines += "\n[상태 목록 — status 조건은 이 값만 사용. '미처리/진행중' 같은 표현은 이 목록의 실제 상태명(여러 개면 status in (...))으로 매핑]\n" + ", ".join(meta["statuses"])
    sysp = ("너는 Jira(Server, JQL v2) 검색식 생성기다. 사용자 질문을 JQL '한 줄'로만 출력한다. 설명·코드펜스·접두어 금지, JQL만.\n"
            "필드: project, issuetype, status, priority, assignee, reporter, created, updated, text(제목·설명·댓글 전체검색), summary, labels, component.\n"
            "규칙:\n"
            "- 질문에서 추출 가능한 조건(프로젝트·이슈유형·상태·담당자·기간)은 해당 JQL 필터로 만들고, 나머지 키워드만 text ~ \"키워드\" 검색으로 남긴다.\n"
            "- 담당자/보고자 이름이 언급되면 assignee in (\"이름\") 형태 (이름 그대로, 시스템이 계정으로 변환).\n"
            "- 복합 명사 키워드는 붙임/띄움 두 표기를 OR로 포함: (text ~ \"현장이슈\" OR text ~ \"현장 이슈\").\n"
            "- 기간은 created(또는 updated) >= \"YYYY-MM-DD\" AND < \"YYYY-MM-DD\". 오늘=" + today + ".\n"
            "- 위 목록에 없는 프로젝트/이슈유형/상태 값을 지어내지 않는다. 확실하지 않은 조건은 넣지 않는다.\n"
            "- 끝에 ORDER BY created DESC.\n"
            + (('반드시 project = "%s" 포함.\n' % proj) if proj else "")
            + _mLines + "\n"
            + "예) '2025년 1월 U9500H Defect 미처리 목록' → " + (('project = "%s" AND ' % proj) if proj else "")
            + "issuetype = Defect AND status in (\"Open\", \"Assign\") AND text ~ \"U9500H\" AND created >= \"2025-01-01\" AND created < \"2025-02-01\" ORDER BY created DESC")
    out = await _jira_llm_complete(llm, sysp, "질문: " + str(q) + "\nJQL:", max_tokens=300, temp=0.0)
    s = _r.sub(r"```[a-zA-Z]*", "", str(out or "")).replace("`", "").replace("\r", " ").strip()
    for line in ([s] + s.split("\n")):
        line = _r.sub(r'^\s*(JQL|jql)\s*[:=]\s*', "", line.strip()).strip()
        if len(line) >= 5 and _r.search(r'(~|=|>|<|ORDER\s+BY)', line, _r.I):
            s = line; break
    s = s.strip()
    if not _r.search(r'(~|=|>|<|ORDER\s+BY|issuetype|project|text|created)', s, _r.I):
        return ""
    if proj and ('project' not in s.lower()):
        m = _r.search(r'\s+ORDER\s+BY\s', s, _r.I)
        s = (s[:m.start()] + (' AND project = "%s"' % proj) + s[m.start():]) if m else (s + (' AND project = "%s"' % proj))
    return s

async def _jira_ask_prep(payload: dict):
    """Jira 검색 + 컨텍스트·프롬프트·LLM 선택까지 공통 준비 (일반/스트리밍 공유)."""
    import re as _re5
    q = str(payload.get("question") or payload.get("query") or "").strip()
    img = str(payload.get("image") or "").strip()   # 첨부 캡처 이미지(data URL) — 비전 LLM에 전달
    if not q and img:
        q = "첨부한 이미지를 참고해 관련 이슈를 분석해줘"
    if not q:
        return {"error": "질문을 입력하세요"}
    def _jtext(v):
        if v is None: return ""
        if isinstance(v, str): return v
        if isinstance(v, (dict, list)):
            out = []
            def _walk(n):
                if isinstance(n, dict):
                    if n.get("type") == "text" and n.get("text"): out.append(n["text"])
                    for c in (n.get("content") or []): _walk(c)
                elif isinstance(n, list):
                    for c in n: _walk(c)
            _walk(v); return " ".join(out)
        return str(v)
    cfg = _jira_cfg(); ai = cfg.get("ai") if isinstance(cfg.get("ai"), dict) else {}
    _SAFE = 500   # 안전 상한(프롬프트 폭발·OOM 방지) — 사실상 무제한
    _rawmax = ai.get("max_issues")
    if _rawmax is None or str(_rawmax).strip() == "": _rawmax = payload.get("max")
    if _rawmax is None:
        _want = 30   # 미설정 기본
    else:
        try: _want = int(_rawmax)
        except Exception: _want = 30
        _want = _SAFE if _want <= 0 else min(_want, _SAFE)   # 0(=설정에서 빈칸 저장) → 개수 제한 없음(=안전상한 500)
    # desc_len / comment_n : 빈 값(None/"") → 무제한(전체) — max_issues 와 동일한 규칙
    _rawdesc = ai.get("desc_len"); _rawcmt = ai.get("comment_n")
    if _rawdesc is None or str(_rawdesc).strip() == "":
        _desclen = None   # None = 자르지 않음(전체)
    else:
        try: _desclen = int(_rawdesc)
        except Exception: _desclen = 2800
        if _desclen <= 0: _desclen = None
        elif _desclen < 200: _desclen = 200   # 너무 짧은 값은 최소 200 보호
    if _rawcmt is None or str(_rawcmt).strip() == "":
        _cmtn = None   # None = 모든 댓글
    else:
        try: _cmtn = int(_rawcmt)
        except Exception: _cmtn = 8
        if _cmtn < 0: _cmtn = 0
    try: _temp = float(ai.get("temperature")) if str(ai.get("temperature") or "").strip() != "" else 0.35
    except Exception: _temp = 0.35
    _maxtok = max(256, int(ai.get("max_tokens") or 3500)); _proj = str(ai.get("project") or "").strip()
    _aj = ai.get("auto_jql"); _autojql = True if _aj is None else (str(_aj).lower() not in ("false", "0", "off", "no", ""))
    # LLM 선택 (JQL 자동생성·답변 공용)
    llms = (load_json(LLMS_FILE).get("llms") or [])
    active = [l for l in llms if l.get("status", "active") == "active" and l.get("endpoint")]
    _lid = str(ai.get("llm_id") or "").strip()
    llm = (next((l for l in active if str(l.get("id") or "") == _lid), None) if _lid else None) \
          or next((l for l in active if str(l.get("type") or "").lower() not in ("claude", "anthropic") and "anthropic.com" not in str(l.get("endpoint", ""))), None) \
          or (active[0] if active else None)
    def _kw_jql():
        _stop = set("관련 이슈 이슈들 알려줘 알려 정리 정리해줘 정리해 해줘 현황 상태 보여줘 보여 대해 대한 뭐야 무엇 어떤 무슨 그리고 좀 해 줘 어디 누가 어떻게 있어 있나 인가 인지 대하여 모두 전부 list 리스트".split())
        toks = [t for t in _re5.split(r"[\s,./]+", _re5.sub(r'["\\]', " ", q)) if t and len(t) >= 2 and t not in _stop]
        tp = ("(" + " OR ".join('text ~ "' + t.replace('"', " ") + '"' for t in toks[:6]) + ")") if toks else ('text ~ "' + _re5.sub(r'["\\]', " ", q)[:60].strip() + '"')
        return ((('project = "%s" AND ' % _proj.replace('"', " ")) if _proj else "") + tp + " ORDER BY updated DESC")
    _fields = "summary,description,status,issuetype,priority,assignee,reporter,project,created,updated,comment,labels,components"
    def _page(_jql, _start, _cnt): return _jira_call("GET", "/rest/api/2/search", params={"jql": _jql, "startAt": _start, "maxResults": max(1, _cnt), "fields": _fields})
    # JQL 결정: ① 직접입력 → ② AI 자동생성(설정 on) → ③ 키워드 OR
    # (질문에 모델명이 있으면 아래에서 제목매칭 이슈를 추가 수집·병합하고 제목매칭/본문언급 태그를 붙인다)
    jql = str(payload.get("jql") or "").strip(); jql_mode = "직접"
    # 모델 토큰(영문 1~3자 + 숫자 3~5자 + 접미: U9532H, E7500, U9024A-10G 등) 감지
    _mdls = ([t for t in _re5.findall(r"\b[A-Za-z]{1,3}\d{3,5}[A-Za-z0-9-]*\b", q)][:3]) if not jql else []
    _meta = (_jira_meta(cfg) if not jql else {})   # 실제 Jira 값(프로젝트·유형·상태·모델매핑) — JQL 필터 정확도
    _preInj = ""   # 모델→프로젝트 주입 전 JQL (0건 시 전체 폴백용)
    if not jql:
        if _autojql and llm:
            try: _g = await _jira_gen_jql(q, llm, _proj, _meta)
            except Exception: _g = ""
            if _g: jql = _g; jql_mode = "AI생성"
        if not jql:
            jql = _kw_jql(); jql_mode = "키워드"
        # 후처리 ①: 키워드 붙임/띄움 두 표기 OR 확장 (띄어쓰기로 결과 갈리는 문제 방지)
        try: jql = _jira_expand_text(jql, q)
        except Exception: pass
        # 후처리 ②: 담당자/보고자 한글 이름 → Jira 계정 리졸브
        try: jql = _jira_fix_assignees(jql, cfg)
        except Exception: pass
        # 후처리 ③: 모델→프로젝트 매핑(패널 설정 auto_models) 결정적 주입 — project 조건이 없을 때만
        try:
            if _mdls and _meta.get("model_proj") and ("project" not in jql.lower()):
                _pk = next((_meta["model_proj"][m.upper()] for m in _mdls if m.upper() in _meta["model_proj"]), "")
                if _pk:
                    _m6 = _re5.search(r"\s+ORDER\s+BY\s", jql, _re5.I)
                    _body = (jql[:_m6.start()] if _m6 else jql).strip()
                    _tail = (jql[_m6.start():] if _m6 else "")
                    _preInj = jql
                    jql = ('project = "%s" AND (%s)%s' % (_pk, _body, _tail))
                    jql_mode += "+프로젝트매핑(" + _pk + ")"
        except Exception:
            pass
    r, err = _page(jql, 0, min(_want, 100))
    if err: return err
    if (r is not None) and (not r.is_success) and jql_mode == "AI생성":   # AI생성 JQL 문법오류 → 키워드로 1회 폴백
        jql = _kw_jql(); jql_mode = "키워드(AI생성 실패→폴백)"
        r, err = _page(jql, 0, min(_want, 100))
        if err: return err
    if not r.is_success:
        return {"error": f"Jira {r.status_code} · {r.text[:200]} (JQL: {jql})"}
    _j0 = r.json(); issues = list(_j0.get("issues") or []); _total = int(_j0.get("total") or len(issues))
    # 매핑 프로젝트로 제한했는데 0건이면 전체(주입 전 JQL)로 폴백 — 참고성 질문 커버
    if _total == 0 and _preInj:
        jql = _preInj; jql_mode += "→0건 전체폴백"
        _r0, _e0 = _page(jql, 0, min(_want, 100))
        if (not _e0) and (_r0 is not None) and _r0.is_success:
            _j0 = _r0.json(); issues = list(_j0.get("issues") or []); _total = int(_j0.get("total") or len(issues))
    while len(issues) < min(_want, _total) and len(issues) < _SAFE:   # 추가 페이지로 더 수집 (개수 제한 없음, 안전상한까지)
        _rr, _e = _page(jql, len(issues), min(100, _want - len(issues)))
        if _e or (_rr is None) or (not _rr.is_success): break
        _b = _rr.json().get("issues") or []
        if not _b: break
        issues.extend(_b)
    issues = issues[:_want]
    # 모델명 제목매칭 보강: 질문에 모델 토큰이 있으면 summary 매칭 이슈를 추가 수집해 병합(중복 제거, 최신순)
    if _mdls:
        try:
            _mq = "(" + " OR ".join('summary ~ "%s"' % m.replace('"', " ") for m in _mdls) + ")"
            _mjql = ((('project = "%s" AND ' % _proj.replace('"', " ")) if _proj else "") + _mq + " ORDER BY created DESC")
            _rm, _em = _page(_mjql, 0, min(_want, 100))
            if (not _em) and (_rm is not None) and _rm.is_success:
                _mis = list(_rm.json().get("issues") or [])
                _seen = {it.get("key") for it in issues}
                _added = [it for it in _mis if it.get("key") not in _seen]
                if _added:
                    issues = sorted(issues + _added, key=lambda it: str((it.get("fields") or {}).get("created") or ""), reverse=True)[: max(_want, len(issues))]
                    jql_mode = (jql_mode or "") + "+모델제목"
        except Exception:
            pass
    def _tmatch(it):
        if not _mdls: return ""
        _s = str((it.get("fields") or {}).get("summary") or "").lower()
        return "제목매칭" if any(m.lower() in _s for m in _mdls) else "본문언급"
    cited = []; ctx = []
    def _nm(d): return str((d or {}).get("displayName") or (d or {}).get("name") or "") if isinstance(d, dict) else ""
    _detn = min(len(issues), 25)   # 앞 25건까지 상세, 나머지는 요약줄 (개수 제한 없이 수집하되 프롬프트는 관리)
    for idx, it in enumerate(issues):
        key = it.get("key", ""); f = it.get("fields") or {}
        summ = str(f.get("summary") or ""); st = ((f.get("status") or {}).get("name") or "")
        itype = ((f.get("issuetype") or {}).get("name") or ""); prio = ((f.get("priority") or {}).get("name") or "")
        asgn = _nm(f.get("assignee")) or "미지정"; rep = _nm(f.get("reporter"))
        upd = str(f.get("updated") or "")[:10]; crt = str(f.get("created") or "")[:10]
        cited.append({"key": key, "summary": summ, "status": st})
        _tg = _tmatch(it); _tgs = (f" / {_tg}" if _tg else "")
        if idx < _detn:
            _d = _jtext(f.get("description")); desc = _d if _desclen is None else _d[:_desclen]
            cms = ((f.get("comment") or {}).get("comments") or [])
            # _cmtn None = 전체 댓글, 0 = 제외, N = 최근 N개
            if _cmtn is None: _pick = cms
            elif _cmtn == 0: _pick = []
            else: _pick = cms[-_cmtn:]
            cmt = "\n".join("  · [" + str(c.get("created") or "")[:10] + " " + _nm(c.get("author")) + "] " + _jtext(c.get("body"))[:700] for c in _pick)
            labels = ", ".join([str(x) for x in (f.get("labels") or [])][:8])
            comps = ", ".join([_nm(x) for x in (f.get("components") or [])][:6])
            meta = f"유형:{itype} / 상태:{st} / 우선순위:{prio} / 담당:{asgn} / 보고:{rep} / 생성:{crt} / 수정:{upd}{_tgs}"
            if labels: meta += f" / 라벨:{labels}"
            if comps: meta += f" / 컴포넌트:{comps}"
            ctx.append(f"### 이슈키 [{key}] 제목: {summ}\n  {meta}\n  [설명]\n{desc or '(설명 없음)'}" + (f"\n  [댓글 {len(cms)}개 중 최근]\n{cmt}" if cmt.strip() else "\n  [댓글 없음]"))
        else:
            ctx.append(f"### 이슈키 [{key}] 제목: {summ}\n  유형:{itype} / 상태:{st} / 우선순위:{prio} / 담당:{asgn} / 생성:{crt}{_tgs}  (요약)")
    # 프롬프트 문자 예산 — 이슈가 많을 때(수백 건) 컨텍스트 폭발 → LLM 컨텍스트 초과·빈 응답 방지
    try: _ctx_budget = int(ai.get("ctx_chars") or 60000)
    except Exception: _ctx_budget = 60000
    _acc, _used, _cut = [], 0, 0
    for _seg in ctx:
        if _acc and (_used + len(_seg) > _ctx_budget):
            _cut += 1; continue
        _acc.append(_seg); _used += len(_seg)
    if _cut:
        _acc.append(f"(프롬프트 길이 제한으로 {_cut}건 생략 — 검색 매칭 총 {len(issues)}건)")
    context = "\n\n══════════════\n\n".join(_acc) if _acc else "(검색 결과 없음)"
    sys_p = ("너는 사내 Jira 이슈를 분석해 주는 전문 어시스턴트다. 아래 '검색된 Jira 이슈'(제목·유형·상태·우선순위·담당·설명·댓글)만을 근거로, "
             "사용자 질문에 한국어로 **상세하고 구조적으로** 답한다. 다음 형식을 반드시 따른다:\n\n"
             "## 핵심 요약\n질문에 대한 답을 3~5문장으로 먼저 제시한다.\n\n"
             "## 이슈별 상세\n관련된 각 이슈마다 다음을 작성한다:\n"
             "- **[PROJ-123] 제목** — (상태/우선순위/담당)\n"
             "- 무엇이 문제/요청인지, 원인, 진행/조치 내용, 댓글에서 드러난 핵심 논의·결론을 2~5줄로 구체적으로.\n\n"
             "## 종합 결론\n공통 원인·패턴, 미해결(Open) 항목, 우선 처리 권고, 추가로 확인이 필요한 점을 정리한다.\n\n"
             "[규칙] 근거가 된 이슈는 반드시 이슈 키([P106-2436]처럼 '프로젝트코드-숫자' 형태, 각 이슈의 '이슈키 [...]' 값)로 인용한다. "
             "이슈 제목 안의 [U9532H]·[LGU]·[상용망] 같은 대괄호 태그는 이슈 키가 아니므로 절대 키 자리에 쓰지 않는다. "
             "이슈에 '제목매칭/본문언급' 표시가 있으면: 질문의 모델명이 제목에 있는 '제목매칭' 이슈를 그 모델의 이슈로 우선하고, '본문언급' 이슈는 참고로만 다룬다. "
             "설명·댓글이 없는 관리성 이슈(산출물·릴리즈 등)가 최신이면 그 사실을 명시하고, 실질 내용이 있는 최신 이슈도 함께 제시한다. "
             "설명·댓글에 실제로 있는 내용만 쓰고 추측·창작은 금지한다. "
             "검색식(JQL)을 임의로 만들어 답변에 표시하지 않는다 — 실제 사용된 JQL은 시스템이 하단에 별도 표시한다. "
             "정보가 부족하면 '해당 이슈의 설명/댓글에 정보가 부족함'이라고 명시한다. 검색 결과 자체가 없으면 '관련 이슈를 찾지 못했습니다'라고만 답한다. "
             "충분히 길고 빠짐없이, 마크다운(##, **, -)으로 가독성 있게 작성한다.")
    if str(ai.get("prompt") or "").strip(): sys_p = str(ai.get("prompt"))   # 설정에서 프롬프트 직접 지정 시 사용
    # Jira 프로젝트 Key ↔ 이름 매핑 — 설정된 게 있으면 시스템 프롬프트 끝에 붙여 LLM 이 프로젝트 판단에 활용
    _km = ai.get("key_mappings") if isinstance(ai.get("key_mappings"), list) else []
    _km = [x for x in _km if isinstance(x, dict) and str(x.get("key") or "").strip()]
    if _km:
        _kml = "\n\n[Jira 프로젝트 Key 매핑 — 질문의 프로젝트명을 이 표에서 Key 로 해석하라]"
        for _x in _km:
            _k = str(_x.get("key") or "").strip()
            _n = str(_x.get("name") or "").strip()
            _d = str(_x.get("desc") or "").strip()
            _kml += f"\n- {_k}"
            if _n: _kml += f" = {_n}"
            if _d: _kml += f" · {_d}"
        sys_p += _kml
    if img:
        sys_p += "\n[첨부 이미지] 사용자가 캡처 이미지를 첨부했다. 이미지 속 로그·CLI 출력·화면 내용을 판독해 질문·이슈 분석에 활용하고, 이미지에서 확인한 내용은 그 사실을 명시한다."
    _hdr = (f"[검색된 Jira 이슈 — 분석 {len(issues)}건" + (f" / 전체 매칭 {_total}건" if _total > len(issues) else "") + "]")
    user_p = f"{_hdr}\n\n{context}\n\n[사용자 질문]\n{q}\n\n위 이슈들을 근거로 형식에 맞춰 상세히 답하라."
    return {"cited": cited, "count": len(issues), "total": _total, "sys_p": sys_p, "user_p": user_p, "llm": llm, "temp": _temp, "maxtok": _maxtok, "jql": jql, "jql_mode": jql_mode, "img": img}

def _jira_mm_openai(user_p: str, img: str):
    """OpenAI 호환 user content — 이미지(data URL) 있으면 멀티모달 배열."""
    if not img:
        return user_p
    return [{"type": "text", "text": user_p}, {"type": "image_url", "image_url": {"url": img}}]

def _jira_mm_claude(user_p: str, img: str):
    """Anthropic user content — data URL → base64 이미지 블록."""
    if not img:
        return user_p
    try:
        head, b64 = img.split(",", 1)
        mt = head.split(":", 1)[1].split(";", 1)[0] if ":" in head else "image/png"
        return [{"type": "image", "source": {"type": "base64", "media_type": mt, "data": b64}},
                {"type": "text", "text": user_p}]
    except Exception:
        return user_p

@app.post("/api/jira/ask")
async def jira_ask(payload: dict):
    """질문 → Jira 이슈 검색 → LLM 답변 (비스트리밍)."""
    import httpx as _hx
    p = await _jira_ask_prep(payload)
    if p.get("error"): return {"ok": False, "error": p["error"]}
    llm = p["llm"]; cited = p["cited"]; count = p["count"]; sys_p = p["sys_p"]; user_p = p["user_p"]; _temp = p["temp"]; _maxtok = p["maxtok"]; jql = p["jql"]
    if not llm:
        return {"ok": False, "error": "등록된 LLM이 없습니다 (AI Assistant에서 LLM을 먼저 등록하세요)", "cited": cited, "count": count}
    answer = None; llm_err = None
    ltype = str(llm.get("type") or "").lower(); ep = str(llm.get("endpoint") or "")
    if ltype in ("claude", "anthropic") or "anthropic.com" in ep:
        try:
            import anthropic as _ah
            m = _ah.Anthropic(api_key=llm.get("apikey") or "").messages.create(
                model=llm.get("model") or "claude-sonnet-4-6", max_tokens=_maxtok, system=sys_p,
                messages=[{"role": "user", "content": _jira_mm_claude(user_p, p.get("img") or "")}])
            answer = "".join(getattr(b, "text", "") for b in m.content).strip()
        except Exception as e:
            llm_err = "Claude: " + str(e)[:200]
    else:
        body = {"model": llm.get("model") or "", "messages": [{"role": "system", "content": sys_p}, {"role": "user", "content": _jira_mm_openai(user_p, p.get("img") or "")}], "temperature": _temp, "max_tokens": _maxtok}
        headers = {"Content-Type": "application/json"}
        ak = llm.get("apikey")
        if ak and not str(ak).lower().startswith("http"):
            headers["Authorization"] = f"Bearer {ak}"
        url = ep.rstrip("/") + "/chat/completions"
        try:
            async with _hx.AsyncClient(timeout=180) as client:
                rr = await client.post(url, headers=headers, json=body)
                if rr.status_code == 200:
                    answer = (((rr.json().get("choices") or [{}])[0].get("message") or {}).get("content") or "").strip()
                else:
                    llm_err = f"LLM {rr.status_code}: {rr.text[:160]} (URL {url})"
        except Exception as e:
            llm_err = str(e)[:200]
    if not answer:
        return {"ok": False, "error": "LLM 호출 실패 — " + (llm_err or "빈 응답"), "cited": cited, "count": count, "llm": llm.get("name")}
    return {"ok": True, "answer": answer, "cited": cited, "count": count, "total": p.get("total"), "jql": jql, "llm": llm.get("name")}

@app.post("/api/jira/ask-stream")
async def jira_ask_stream(payload: dict):
    """질문 → Jira 검색 → LLM 답변 스트리밍 (SSE) — 느린 응답을 토큰 단위로 즉시 표시."""
    import httpx as _hx, json as _json
    from fastapi.responses import StreamingResponse
    p = await _jira_ask_prep(payload)
    async def gen():
        def sse(o): return "data: " + _json.dumps(o, ensure_ascii=False) + "\n\n"
        if p.get("error"):
            yield sse({"error": p["error"]}); return
        llm = p["llm"]
        yield sse({"meta": {"cited": p["cited"], "count": p["count"], "total": p.get("total"), "jql": p["jql"], "jql_mode": p.get("jql_mode"), "llm": (llm or {}).get("name")}})
        if not llm:
            yield sse({"error": "등록된 LLM이 없습니다 (AI Assistant에서 LLM을 먼저 등록하세요)"}); yield sse({"done": True}); return
        sys_p = p["sys_p"]; user_p = p["user_p"]; _temp = p["temp"]; _maxtok = p["maxtok"]; _img = p.get("img") or ""
        ltype = str(llm.get("type") or "").lower(); ep = str(llm.get("endpoint") or ""); got = False
        if ltype in ("claude", "anthropic") or "anthropic.com" in ep:
            try:
                import anthropic as _ah
                with _ah.Anthropic(api_key=llm.get("apikey") or "").messages.stream(
                        model=llm.get("model") or "claude-sonnet-4-6", max_tokens=_maxtok, system=sys_p,
                        messages=[{"role": "user", "content": _jira_mm_claude(user_p, _img)}]) as _st:
                    for _txt in _st.text_stream:
                        if _txt: got = True; yield sse({"delta": _txt})
            except Exception as e:
                yield sse({"error": "Claude: " + str(e)[:200]})
        else:
            body = {"model": llm.get("model") or "", "messages": [{"role": "system", "content": sys_p}, {"role": "user", "content": _jira_mm_openai(user_p, _img)}], "temperature": _temp, "max_tokens": _maxtok, "stream": True}
            headers = {"Content-Type": "application/json"}
            ak = llm.get("apikey")
            if ak and not str(ak).lower().startswith("http"):
                headers["Authorization"] = f"Bearer {ak}"
            url = ep.rstrip("/") + "/chat/completions"
            try:
                async with _hx.AsyncClient(timeout=300) as client:
                    async with client.stream("POST", url, headers=headers, json=body) as rr:
                        if rr.status_code != 200:
                            _tx = (await rr.aread()).decode("utf-8", "replace")[:200]
                            yield sse({"error": f"LLM {rr.status_code}: {_tx}"})
                        else:
                            async for line in rr.aiter_lines():
                                if not line or not line.startswith("data:"): continue
                                _d = line[5:].strip()
                                if _d == "[DONE]": break
                                try: _j = _json.loads(_d)
                                except Exception: continue
                                _delta = (((_j.get("choices") or [{}])[0].get("delta") or {}).get("content")) or ""
                                if _delta: got = True; yield sse({"delta": _delta})
            except Exception as e:
                yield sse({"error": str(e)[:200]})
        if not got:
            yield sse({"error": "빈 응답"})
        yield sse({"done": True})
    return StreamingResponse(gen(), media_type="text/event-stream", headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no", "Connection": "keep-alive", "Content-Encoding": "identity"})

@app.get("/api/jira/search-all")
async def jira_search_all(jql: str, fields: str = "summary,status,issuetype,reporter,project", cap: int = 10000):
    cfg = _jira_cfg()
    issues = []
    start = 0
    total = None
    for _ in range(200):   # 안전망: 최대 200페이지(=20000건)
        r, err = _jira_call("GET", "/rest/api/2/search", cfg=cfg,
                            params={"jql": jql, "startAt": start, "maxResults": 100, "fields": fields})
        if err:
            return err
        if not r.is_success:
            return {"ok": False, "error": f"{r.status_code} · {r.text[:300]}"}
        j = r.json()
        batch = j.get("issues", [])
        issues.extend(batch)
        total = j.get("total", len(issues))
        start += len(batch)
        if not batch or start >= total or len(issues) >= cap:
            break
    return {"ok": True, "issues": issues, "total": total if total is not None else len(issues)}

# ── Issue Sync: utop 서버 저장 + 증분 동기화 ──
def _issue_path(project: str):
    import re as _re
    safe = _re.sub(r"[^A-Za-z0-9_.-]", "_", str(project or "")) or "_"
    return DATA_DIR / "issue_sync" / (safe + ".json")

def _issue_load(project: str) -> dict:
    p = _issue_path(project)
    if p.exists():
        try:
            return json.loads(p.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}

def _issue_save(project: str, store: dict):
    p = _issue_path(project)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(store, ensure_ascii=False), encoding="utf-8")

@app.get("/api/issues/{project}")
async def issues_get(project: str):
    """utop에 저장된 이슈를 그대로 반환 (Jira 호출 없음)."""
    store = _issue_load(project)
    issues = store.get("issues", [])
    return {"ok": True, "project": project, "issues": issues,
            "count": len(issues), "last_synced_at": store.get("last_synced_at", "")}

@app.post("/api/issues/sync")
async def issues_sync(payload: dict):
    """프로젝트 이슈를 Jira에서 가져와 utop에 저장. 마지막 마커 이후 변경분만(증분), full=True면 전체."""
    from datetime import datetime as _dt, timedelta as _td
    project = str(payload.get("project") or "").strip()
    if not project:
        return {"ok": False, "error": "프로젝트가 없습니다"}
    fields = str(payload.get("fields") or "summary,status,issuetype,assignee,priority,updated")
    fset = [f.strip() for f in fields.split(",") if f.strip()]
    if "updated" not in fset:
        fset.append("updated")
    fields = ",".join(fset)
    full = bool(payload.get("full"))
    store = _issue_load(project)
    marker = store.get("updated_marker", "")
    existing = store.get("issues", []) or []
    if full or not marker or not existing:
        jql = f'project = "{project}" ORDER BY updated DESC'
        mode = "full"
    else:
        jql = f'project = "{project}" AND updated >= "{marker}" ORDER BY updated DESC'
        mode = "incremental"
    res = await jira_search_all(jql=jql, fields=fields, cap=20000)
    if not res.get("ok"):
        return {"ok": False, "error": res.get("error") or "Jira 조회 실패"}
    fetched = res.get("issues", []) or []
    by_key = {}
    if mode == "incremental":
        for it in existing:
            k = it.get("key")
            if k:
                by_key[k] = it
    added = 0; updated_n = 0
    for it in fetched:
        k = it.get("key")
        if not k:
            continue
        if k in by_key:
            updated_n += 1
        else:
            added += 1
        by_key[k] = it
    merged = list(by_key.values())
    def _upd(it):
        return ((it.get("fields") or {}).get("updated")) or ""
    merged.sort(key=_upd, reverse=True)
    # 다음 증분용 마커: Jira updated 최댓값 - 2분 버퍼 ("yyyy-MM-dd HH:mm")
    max_upd = ""
    for it in merged:
        u = _upd(it)
        if u > max_upd:
            max_upd = u
    new_marker = marker
    if max_upd:
        try:
            dt = _dt.strptime(max_upd[:19], "%Y-%m-%dT%H:%M:%S") - _td(minutes=2)
            new_marker = dt.strftime("%Y-%m-%d %H:%M")
        except Exception:
            new_marker = marker
    now_iso = _dt.now().strftime("%Y-%m-%d %H:%M:%S")
    _issue_save(project, {"project": project, "issues": merged, "fields": fields,
                          "updated_marker": new_marker, "last_synced_at": now_iso})
    return {"ok": True, "project": project, "mode": mode, "added": added,
            "updated": updated_n, "total": len(merged), "fetched": len(fetched),
            "last_synced_at": now_iso, "issues": merged}

@app.get("/api/jira/fields")
async def jira_fields():
    r, err = _jira_call("GET", "/rest/api/2/field")
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:300]}"}
    out = []
    for f in (r.json() or []):
        out.append({"id": f.get("id"), "name": f.get("name"), "custom": bool(f.get("custom"))})
    return {"ok": True, "fields": out}

@app.get("/api/jira/versions")
async def jira_versions(project: str):
    r, err = _jira_call("GET", f"/rest/api/2/project/{project}/versions")
    if err:
        return err
    if not r.is_success:
        return {"ok": False, "error": f"{r.status_code} · {r.text[:300]}"}
    vs = []
    for v in (r.json() or []):
        vs.append({"id": v.get("id"), "name": v.get("name"),
                   "released": bool(v.get("released")), "archived": bool(v.get("archived")),
                   "releaseDate": v.get("releaseDate") or "", "startDate": v.get("startDate") or "",
                   "description": v.get("description") or ""})
    return {"ok": True, "versions": vs}

# ===== Release Summary 적부 판정 (Jira 버전 이슈를 시험 항목으로) =====
RELEASE_JUDGE_FILE = DATA_DIR / "config" / "release_judge.json"

def _rj_load():
    if RELEASE_JUDGE_FILE.exists():
        try:
            return load_json(RELEASE_JUDGE_FILE)
        except Exception:
            return {}
    return {}



# ===== 게시판 (수정사항 요청) =====
BOARD_FILE = DATA_DIR / "state" / "board.json"
BOARD_FILES_DIR = DATA_DIR / "board_files"
BOARD_FILES_DIR.mkdir(parents=True, exist_ok=True)


def _board_load():
    if BOARD_FILE.exists():
        try:
            return load_json(BOARD_FILE)
        except Exception:
            return {"posts": []}
    return {"posts": []}


@app.get("/api/board")
async def board_list():
    return _board_load()


@app.post("/api/board")
async def board_add(payload: dict):
    data = _board_load()
    posts = data.get("posts", [])
    import time as _t
    post = {
        "id": str(int(_t.time() * 1000)),
        "title": (str(payload.get("title", "")).strip() or "(제목 없음)"),
        "body": str(payload.get("body", "")).strip(),
        "author": (str(payload.get("author", "")).strip() or "익명"),
        "status": "open",
        "created_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "done_at": "",
        "attachments": (payload.get("attachments") or []),
    }
    posts.insert(0, post)
    data["posts"] = posts
    save_json(BOARD_FILE, data)
    await broadcast({"type": "board_update"})
    return {"success": True, "post": post}


@app.post("/api/board/{pid}")
async def board_update(pid: str, payload: dict):
    data = _board_load()
    for p in data.get("posts", []):
        if p.get("id") == pid:
            if "status" in payload:
                p["status"] = str(payload.get("status") or "open").strip()  # open/approved/rejected/done 그대로 저장
                p["done_at"] = datetime.now().strftime("%Y-%m-%d %H:%M") if p["status"] == "done" else ""
            for k in ("title", "body", "author"):
                if k in payload:
                    p[k] = str(payload[k]).strip()
            if isinstance(payload.get("attachments"), list):
                import os as _os
                new_names = set(str((a or {}).get("name", "")) for a in payload["attachments"] if isinstance(a, dict))
                for old in (p.get("attachments") or []):
                    on = str((old or {}).get("name", ""))
                    if on and on not in new_names:
                        try:
                            fp = BOARD_FILES_DIR / _os.path.basename(on)
                            if fp.exists():
                                fp.unlink()
                        except Exception:
                            pass
                p["attachments"] = payload["attachments"]
            save_json(BOARD_FILE, data)
            try: asyncio.create_task(broadcast({"type": "board_update"}))
            except Exception: pass
            return {"success": True, "post": p, "_v": "v2", "_rs": payload.get("status"), "_haskey": ("status" in payload)}
    raise HTTPException(404, "글을 찾을 수 없습니다")


@app.delete("/api/board/{pid}")
async def board_remove(pid: str):
    import os as _os
    data = _board_load()
    keep = []
    for p in data.get("posts", []):
        if p.get("id") == pid:
            for a in (p.get("attachments") or []):
                try:
                    fp = BOARD_FILES_DIR / _os.path.basename(str(a.get("name", "")))
                    if fp.exists():
                        fp.unlink()
                except Exception:
                    pass
        else:
            keep.append(p)
    data["posts"] = keep
    save_json(BOARD_FILE, data)
    await broadcast({"type": "board_update"})
    return {"success": True}


@app.post("/api/board-upload")
async def board_upload(payload: dict):
    import time as _t
    import re as _re
    import base64 as _b64
    orig = str(payload.get("orig", "file")) or "file"
    data = str(payload.get("data", ""))
    if data.startswith("data:") and "," in data:
        data = data.split(",", 1)[1]
    try:
        raw = _b64.b64decode(data)
    except Exception:
        raise HTTPException(400, "파일 디코드 실패")
    if len(raw) > 25 * 1024 * 1024:
        raise HTTPException(413, "파일이 너무 큽니다 (25MB 이하)")
    safe = _re.sub(r"[^A-Za-z0-9._-]", "_", orig)
    fname = str(int(_t.time() * 1000)) + "_" + safe
    with open(BOARD_FILES_DIR / fname, "wb") as fp:
        fp.write(raw)
    ext = orig.rsplit(".", 1)[-1].lower() if "." in orig else ""
    is_image = ext in ("png", "jpg", "jpeg", "gif", "webp", "bmp", "svg")
    return {"success": True, "name": fname, "orig": orig, "url": "/api/board/file/" + fname, "size": len(raw), "is_image": is_image}


@app.get("/api/board/file/{fname}")
async def board_file(fname: str):
    import os as _os
    dest = BOARD_FILES_DIR / _os.path.basename(fname)
    if not dest.exists():
        raise HTTPException(404, "파일 없음")
    return FileResponse(str(dest))


@app.post("/api/board/{pid}/reply")
async def board_reply_add(pid: str, payload: dict):
    import time as _t
    data = _board_load()
    for p in data.get("posts", []):
        if p.get("id") == pid:
            if not isinstance(p.get("replies"), list):
                p["replies"] = []
            reply = {
                "id": "r" + str(int(_t.time() * 1000)),
                "author": (str(payload.get("author", "")).strip() or "익명"),
                "body": str(payload.get("body", "")).strip(),
                "attachments": (payload.get("attachments") or []),
                "created_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
            }
            p["replies"].append(reply)
            save_json(BOARD_FILE, data)
            await broadcast({"type": "board_update"})
            return {"success": True, "reply": reply}
    raise HTTPException(404, "글을 찾을 수 없습니다")


@app.post("/api/board/{pid}/reply/{rid}")
async def board_reply_update(pid: str, rid: str, payload: dict):
    import os as _os
    data = _board_load()
    for p in data.get("posts", []):
        if p.get("id") == pid:
            for rp in (p.get("replies") or []):
                if rp.get("id") == rid:
                    if "body" in payload:
                        rp["body"] = str(payload["body"]).strip()
                    if "author" in payload:
                        rp["author"] = str(payload["author"]).strip() or "익명"
                    if isinstance(payload.get("attachments"), list):
                        new_names = set(str((a or {}).get("name", "")) for a in payload["attachments"] if isinstance(a, dict))
                        for old in (rp.get("attachments") or []):
                            on = str((old or {}).get("name", ""))
                            if on and on not in new_names:
                                try:
                                    fp = BOARD_FILES_DIR / _os.path.basename(on)
                                    if fp.exists():
                                        fp.unlink()
                                except Exception:
                                    pass
                        rp["attachments"] = payload["attachments"]
                    rp["updated_at"] = datetime.now().strftime("%Y-%m-%d %H:%M")
                    save_json(BOARD_FILE, data)
                    await broadcast({"type": "board_update"})
                    return {"success": True, "reply": rp}
            raise HTTPException(404, "답글을 찾을 수 없습니다")
    raise HTTPException(404, "글을 찾을 수 없습니다")


@app.delete("/api/board/{pid}/reply/{rid}")
async def board_reply_remove(pid: str, rid: str):
    import os as _os
    data = _board_load()
    for p in data.get("posts", []):
        if p.get("id") == pid:
            keep = []
            for rp in (p.get("replies") or []):
                if rp.get("id") == rid:
                    for a in (rp.get("attachments") or []):
                        try:
                            fp = BOARD_FILES_DIR / _os.path.basename(str(a.get("name", "")))
                            if fp.exists():
                                fp.unlink()
                        except Exception:
                            pass
                else:
                    keep.append(rp)
            p["replies"] = keep
            save_json(BOARD_FILE, data)
            await broadcast({"type": "board_update"})
            return {"success": True}
    raise HTTPException(404, "글을 찾을 수 없습니다")




# ── 시스템 UI 옵션 (관리자 설정 → 전체 유저 공유) ──
UI_OPTIONS_FILE = DATA_DIR / "config" / "ui_options.json"
def _load_ui_options():
    if UI_OPTIONS_FILE.exists():
        try: return json.loads(UI_OPTIONS_FILE.read_text(encoding="utf-8"))
        except: pass
    return {"show_req_id": True, "show_tc_id": True}

# ── 스텝 종류별 쓰임새 (SETUP → TC Step Action) ─────────────────
#
# 종류마다 「실행 로그에 찍을까 · ＋스텝 목록에 내놓을까 · 결과서(PPTX)에
# 실을까」 가 현장마다 다르다. 여태 코드에 박혀 있어서 고치려면 배포를
# 해야 했다(지시: 설정 페이지로).
STEP_ACT_FILE = DATA_DIR / "state" / "step_actions.json"


def _load_step_actions() -> dict:
    try:
        if STEP_ACT_FILE.exists():
            return json.loads(STEP_ACT_FILE.read_text(encoding="utf-8")) or {}
    except Exception:
        pass
    return {}


@app.get("/api/step-actions")
async def step_actions_get():
    """{종류: {log, add, pptx}} — 안 적힌 종류는 셋 다 켠 것으로 본다."""
    return {"items": _load_step_actions()}


@app.post("/api/step-actions")
async def step_actions_save(payload: dict, token: str = ""):
    _require_admin(token)
    items = payload.get("items")
    if not isinstance(items, dict):
        raise HTTPException(400, "items 가 없습니다")
    clean = {}
    for k, v in items.items():
        if not isinstance(v, dict):
            continue
        clean[str(k)] = {
            "log": bool(v.get("log", True)),
            "add": bool(v.get("add", True)),
            "pptx": bool(v.get("pptx", True)),
        }
    STEP_ACT_FILE.parent.mkdir(parents=True, exist_ok=True)
    STEP_ACT_FILE.write_text(json.dumps(clean, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"ok": True, "items": clean}


@app.get("/api/ui-options")
async def ui_options_get():
    return _load_ui_options()

@app.post("/api/ui-options")
async def ui_options_save(payload: dict, token: str = ""):
    _require_admin(token)
    cur = _load_ui_options()
    if "show_req_id" in payload: cur["show_req_id"] = bool(payload["show_req_id"])
    if "show_tc_id" in payload: cur["show_tc_id"] = bool(payload["show_tc_id"])
    UI_OPTIONS_FILE.write_text(json.dumps(cur, ensure_ascii=False), encoding="utf-8")
    return {"ok": True, **cur}

# ── Global Parameters (TC 변수 치환용) ──
def _load_global_params():
    if GLOBAL_PARAMS_FILE.exists():
        try: return json.loads(GLOBAL_PARAMS_FILE.read_text(encoding="utf-8"))
        except: pass
    return {}

@app.get("/api/global-params")
async def global_params_get():
    return _load_global_params()

@app.post("/api/global-params")
async def global_params_save(payload: dict):
    GLOBAL_PARAMS_FILE.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"ok": True}


# ══════════════════════════════════════════════════════════════════════
# TO-DO (관리자 공용) — app_kv 에 저장. WebSocket 으로 실시간 동기화.
# 형식: {"items": [{text, status, at}, ...]}
# ══════════════════════════════════════════════════════════════════════
@app.get("/api/todo")
async def todo_list():
    d = await db.kv_get("admin_todo") or {}
    items = d.get("items") if isinstance(d, dict) else None
    return {"items": items if isinstance(items, list) else []}


@app.post("/api/todo")
async def todo_save(payload: dict):
    items = (payload or {}).get("items")
    if not isinstance(items, list):
        raise HTTPException(400, "items 는 배열이어야 합니다")
    # 정리 (문자열 필드만, 상태 화이트리스트)
    _valid_st = {"todo", "doing", "done"}
    _clean = []
    for it in items:
        if not isinstance(it, dict):
            continue
        _cmts_raw = it.get("comments") if isinstance(it.get("comments"), list) else []
        _cmts = []
        for c in _cmts_raw:
            if not isinstance(c, dict):
                continue
            _imgs_raw = c.get("images") if isinstance(c.get("images"), list) else []
            # 이미지는 dataURL 문자열만, 5MB 이하 각각·항목당 최대 10개
            _imgs = []
            for im in _imgs_raw[:10]:
                if isinstance(im, str) and im.startswith("data:") and len(im) < 5_000_000:
                    _imgs.append(im)
            _cmts.append({
                "id": str(c.get("id") or "")[:64],
                "text": str(c.get("text") or "")[:4000],
                "images": _imgs,
                "author": str(c.get("author") or "")[:100],
                "at": c.get("at") or 0,
            })
        _clean.append({
            "text": str(it.get("text") or "")[:2000],
            "status": (it.get("status") if it.get("status") in _valid_st else "todo"),
            "at": it.get("at") or 0,
            "comments": _cmts,
        })
    await db.kv_set("admin_todo", {"items": _clean})
    try: asyncio.create_task(broadcast({"type": "todo_updated"}))
    except Exception: pass
    return {"ok": True, "items": _clean}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)


# ══════════════════════════════════════════════════════════════════════
# 자연어로 시험 만들기
#
# "E6100 rate limit 시험 해줘" 한 줄에서 슬롯(장비·세션)과 스텝을 만든다.
# 이 기능의 성패는 '근거' 다. 모델에게 맨손으로 물으면 유비쿼스에 없는
# 명령을 그럴듯하게 지어낸다. 그래서 반드시 다음을 찾아 함께 넘긴다:
#   1. 등록된 장비와 그 인터페이스 (실제로 존재하는 포트만 쓰게)
#   2. 같은 모델로 이미 만든 TC 의 스텝 (사내에서 쓰는 실제 명령)
#   3. 요구사항·매뉴얼에서 찾은 조각 (임베딩 서버가 있을 때)
#
# 임베딩 서버가 없어도 1·2 만으로 동작한다 — 사내망 밖에서도 쓸 수 있어야
# 하고, 무엇보다 2번(우리가 쓰던 실제 명령)이 가장 정확한 근거다.
# ══════════════════════════════════════════════════════════════════════
async def _expand_series(hint: str) -> list[str]:
    """말에 나온 모델군을 그 군의 모델 이름으로 펼친다.

    "E6000 시리즈 rate limit 시험" 이라고 하면 E6100-48X · E6100-24X 로
    만든 기존 TC 도 근거로 잡혀야 한다. 시리즈 이름만으로 찾으면 아무것도
    안 나온다 — 스텝에는 모델명이 적혀 있기 때문이다.
    """
    async with db.pool().acquire() as c:
        rows = await c.fetch(
            "SELECT name, model_group FROM device_catalog "
            "WHERE kind='model' AND model_group IS NOT NULL"
        )
    low = hint.lower()
    out: list[str] = []
    for r in rows:
        g = (r["model_group"] or "").strip()
        if g and g.lower() in low:
            out.append(r["name"])
    return out


async def _grounding_devices(hint: str) -> list[dict]:
    """말에 나온 모델·모델군·IP 와 맞는 장비를 찾는다. 없으면 전부 조금씩."""
    devs = await db.device_list()
    words = [w for w in re.split(r"[\s,·]+", hint) if len(w) >= 2]
    hit = []
    for d in devs:
        hay = " ".join(
            str(d.get(k) or "")
            for k in ("ip", "model", "model_group", "vendor", "role", "lab")
        )
        if any(w.lower() in hay.lower() for w in words):
            hit.append(d)
    return (hit or devs)[:8]


async def _grounding_steps(hint: str, limit: int = 40) -> list[dict]:
    """같은 모델·주제로 이미 만든 TC 스텝. 사내에서 실제로 쓰는 명령이다.

    모델군으로 물어보면 그 군의 모델명까지 넓혀서 찾는다.
    """
    words = [w for w in re.split(r"[\s,·]+", hint) if len(w) >= 2][:6]
    words += await _expand_series(hint)
    if not words:
        return []
    async with db.pool().acquire() as c:
        rows = await c.fetch(
            """SELECT tcid, name, data FROM tc
               WHERE data::text ILIKE ANY($1::text[])
               ORDER BY updated_at DESC LIMIT 12""",
            [f"%{w}%" for w in words],
        )
    out = []
    for r in rows:
        d = r["data"] or {}
        for s in (d.get("checks") or [])[:8]:
            if not isinstance(s, dict):
                continue
            cmd = (s.get("cli") or s.get("data") or "").strip()
            if cmd:
                out.append({"tcid": r["tcid"], "cmd": cmd,
                            "expected": (s.get("expected") or s.get("criteria") or "")[:200]})
            if len(out) >= limit:
                return out
    return out


async def _grounding_docs(query: str, k: int = 6) -> list[dict]:
    """요구사항·매뉴얼에서 찾은 조각. 임베딩 서버가 없으면 빈 목록."""
    vecs = await _embed_texts([query])
    if not vecs:
        return []
    import numpy as _np
    qv = _np.asarray(vecs[0], dtype="float32")
    qn = float(_np.linalg.norm(qv)) or 1.0
    async with db.pool().acquire() as c:
        rows = await c.fetch("SELECT key, embed, meta FROM rag_embed")
    scored = []
    for r in rows:
        try:
            v = _np.frombuffer(r["embed"], dtype="float32")
            if v.shape != qv.shape:
                continue
            s = float(qv @ v) / (qn * (float(_np.linalg.norm(v)) or 1.0))
            scored.append((s, r))
        except Exception:
            continue
    scored.sort(key=lambda x: -x[0])
    cfg = _rag_cfg()
    lo = float(cfg.get("min_score") or 0)
    return [
        {"key": r["key"], "score": round(s, 3), "text": (r["meta"] or {}).get("text", "")[:800]}
        for s, r in scored[:k] if s >= lo
    ]


_GEN_SYSTEM = """당신은 유비쿼스 네트워크 장비 시험 자동화 도구의 시험 설계자다.
사용자의 한 줄 요청에서 시험 슬롯과 스텝을 만든다.

절대 규칙:
1. 아래 '근거' 에 없는 CLI 명령을 지어내지 마라. 근거의 명령을 그대로 쓰거나
   포트·값만 바꿔 쓴다. 근거에 없으면 그 스텝의 data 를 비우고 note 에
   "근거 없음 - 확인 필요" 라고 적는다.
2. 인터페이스는 '등록된 장비' 에 실제로 있는 이름만 쓴다.
3. 슬롯 key 는 s1, s2 … 순서대로. 스텝의 session 은 반드시 만든 슬롯의 key.
4. 계측기가 필요하면 슬롯을 따로 만들고 family 를 '계측기' 로 한다.
5. 판정 기준(criteria)은 응답에서 확인할 문자열이나 수치 조건으로 적는다.

반드시 JSON 만 출력한다. 설명 문장을 붙이지 마라. 모양:
{"slots":[{"key":"s1","label":"DUT","family":"L2","device_ip":"","protocol":"telnet"}],
 "steps":[{"kind":"auto","session":"s1","step":"...","data":"...","expected":"...","criteria":"...","rca":"...","note":""}],
 "summary":"무엇을 어떻게 시험하는지 두 문장",
 "unsure":["근거가 부족해 확인이 필요한 것"]}"""


@app.post("/api/tc/{tc_id}/generate")
async def tc_generate(tc_id: str, payload: dict):
    """자연어 한 줄 → 슬롯·스텝 제안. 저장하지 않고 돌려만 준다.

    바로 저장하지 않는 이유: 모델이 만든 스텝을 사람이 보기 전에 넣으면
    잘못된 명령이 장비로 나간다. 화면에서 확인하고 적용하게 한다.
    """
    prompt = str(payload.get("prompt") or "").strip()
    # 근거 찾기용 짧은 질의. 문서로 만들 때는 프롬프트가 길어서 그대로
    # 쓰면 "시험 제목:" 같은 껍데기 낱말로 장비를 찾게 된다.
    gquery = prompt

    if not prompt:
        # 프롬프트가 없으면 **요구사항 구현의도 + 시험 목적**으로 만든다.
        # 이쪽이 본류다 — 시험은 요구사항을 검증하려고 있는 것이라, 무엇을
        # 만들지는 그 두 글이 정한다. 한 줄 요청은 빠른 손을 위한 지름길이다.
        tc = payload.get("tc") if isinstance(payload.get("tc"), dict) else None
        if tc is None:
            tc = await db.tc_get(_tc_id_norm(tc_id))
        if not isinstance(tc, dict):
            raise HTTPException(404, "TC 를 찾을 수 없습니다")
        req = None
        rid = str(tc.get("req_id") or "").strip()
        if rid:
            try:
                # req_id 칸에는 PG 키(rq-…)와 부여 ID(U-REQ-…)가 섞여 있다
                req = await db.req_get(rid)
                if req is None:
                    async with db.pool().acquire() as c:
                        row = await c.fetchrow(
                            "SELECT data FROM req WHERE data->>'reqid'=$1 LIMIT 1", rid
                        )
                        req = dict(row["data"]) if row else None
            except Exception as e:
                print(f"[tc_generate] 요구사항 조회 실패({rid}): {e}", flush=True)
        intent = str((req or {}).get("desc") or "").strip()
        obj = str(tc.get("object_md") or "").strip()
        pre = str(tc.get("precondition_md") or "").strip()
        name = str(tc.get("name") or "").strip()
        if not intent and not obj:
            # 재료가 없으면 모델은 지어낼 수밖에 없다 — 만들지 않는 것이 맞다
            raise HTTPException(
                400,
                "요구사항 구현의도(Intent)와 시험 목적(Object)이 모두 비어 있습니다 — "
                "둘 중 하나는 있어야 스텝을 설계할 수 있습니다",
            )
        parts = [f"시험 제목: {name or tc_id}"]
        if req:
            parts.append(f"요구사항: {req.get('title') or req.get('reqid') or rid}")
        if intent:
            parts.append(f"=== 요구사항 구현의도 ===\n{intent[:4000]}")
        if obj:
            parts.append(f"=== 시험 목적 ===\n{obj[:2000]}")
        if pre:
            parts.append(f"=== 사전 준비 조건 ===\n{pre[:1000]}")
        parts.append("위 구현의도와 시험 목적을 검증하는 시험 스텝을 설계하라.")
        prompt = "\n\n".join(parts)
        gquery = " ".join(x for x in [name, str((req or {}).get("title") or ""), obj[:200]] if x)

    # 누구에게 맡길지 화면이 고른다(지시). 안 고르면 여태처럼 Claude 다.
    want_llm = str(payload.get("llm") or "").strip()
    cl, cmodel = _claude_any()
    if cl is None and not want_llm:
        raise HTTPException(
            503,
            "쓸 수 있는 Claude 가 없습니다 — 설정 → LLM 설정에 Anthropic 을 등록하거나 "
            ".env 에 ANTHROPIC_API_KEY 를 넣으세요",
        )

    devs = await _grounding_devices(gquery)
    prev = await _grounding_steps(gquery)
    docs = await _grounding_docs(gquery)

    dev_txt = "\n".join(
        f"- {d.get('ip')} · {d.get('model') or '?'}"
        f"{' (' + d['model_group'] + ')' if d.get('model_group') else ''}"
        f" · {d.get('role') or '?'} · {d.get('lab') or '?'}"
        f" · 접속 {','.join(a['protocol'] for a in (d.get('access') or []))}"
        f" · 포트 {','.join(i['name'] for i in (d.get('interfaces') or [])[:60]) or '없음'}"
        for d in devs
    ) or "(등록된 장비 없음)"
    prev_txt = "\n".join(f"- [{p['tcid']}] {p['cmd']}   → {p['expected']}" for p in prev) \
        or "(비슷한 시험 없음)"
    doc_txt = "\n\n".join(f"[{d['key']} {d['score']}]\n{d['text']}" for d in docs) \
        or "(임베딩 서버가 없어 문서 근거는 비어 있습니다)"

    user = (
        f"요청: {prompt}\n\n"
        f"=== 근거 1. 등록된 장비와 실제 포트 ===\n{dev_txt}\n\n"
        f"=== 근거 2. 우리가 이미 쓰는 명령 ===\n{prev_txt}\n\n"
        f"=== 근거 3. 요구사항·매뉴얼 ===\n{doc_txt}\n"
    )

    if want_llm:
        # 로컬이든 Claude 든 한 길로 — _llm_text 가 종류를 가려서 부른다
        _got, raw = await _ask_json("coverage_automation", _GEN_SYSTEM, user,
                                    max_tokens=4000, llm_id=want_llm)
        if isinstance(_got, dict):
            raw = json.dumps(_got, ensure_ascii=False)
    else:
        try:
            msg = cl.messages.create(
                model=cmodel or CLAUDE_FALLBACK_MODEL,
                max_tokens=4000,
                system=_GEN_SYSTEM,
                messages=[{"role": "user", "content": user}],
            )
            raw = "".join(b.text for b in msg.content if getattr(b, "type", "") == "text").strip()
        except Exception as e:
            raise HTTPException(502, f"모델 호출에 실패했습니다: {_llm_err(e)}") from e

    # 모델이 ```json 으로 감싸는 경우가 있다
    m = re.search(r"\{.*\}", raw, re.S)
    if not m:
        raise HTTPException(502, "모델이 JSON 을 돌려주지 않았습니다")
    try:
        out = json.loads(m.group(0))
    except json.JSONDecodeError as e:
        raise HTTPException(502, f"모델 응답을 읽지 못했습니다: {e}") from e

    return {
        "success": True,
        "proposal": out,
        "grounding": {
            "devices": len(devs),
            "prev_steps": len(prev),
            "docs": len(docs),
            "embed_ready": bool(docs) or bool(_rag_cfg().get("embed_url")),
        },
    }


def _llm_err(e: Exception) -> str:
    """모델이 준 실패를 **사람 말로** 바꾼다.

    Anthropic 은 영어 한 덩어리(JSON 통째)로 답한다. 「credit balance is
    too low」 를 그대로 화면에 던져 놓으면 무엇을 해야 하는지 알 수 없다
    (지적: 이건 뭐야). 흔한 셋은 풀어 쓰고, 나머지는 원문을 남긴다 —
    모르는 실패를 지어내 설명하는 것이 더 나쁘다.
    """
    t = str(e)
    low = t.lower()
    if "credit balance is too low" in low or ("insufficient" in low and "credit" in low):
        return ("Claude 계정에 크레딧이 없습니다 — console.anthropic.com 의 "
                "Plans & Billing 에서 충전하거나, 용도별 프롬프트에서 사용 LLM 을 "
                "랩 안의 로컬 LLM 으로 바꾸세요")
    if "authentication_error" in low or "invalid x-api-key" in low or "401" in t[:40]:
        return "API 키가 맞지 않습니다 — 설정 → LLM 설정에서 키를 다시 넣으세요"
    if "rate_limit" in low or "429" in t[:40]:
        return "잠시 뒤에 다시 하세요 — 짧은 사이에 너무 여러 번 불렀습니다(rate limit)"
    if "not_found_error" in low or ("model" in low and "not found" in low):
        return "모델 이름이 맞지 않습니다 — 설정 → LLM 설정에서 모델을 다시 고르세요"
    return t


def _json_from(raw: str):
    """모델이 준 글에서 **JSON 을 건져 낸다**.

    작은 모델은 ```json 으로 감싸거나, 앞에 「알겠습니다」 를 붙이거나,
    뒤에 설명을 단다. 한 번에 못 읽었다고 「모델이 JSON 을 돌려주지
    않았습니다」 로 끝내면 사람이 할 수 있는 일이 없다(지적).

    못 건지면 None 이다 — 지어내지 않는다.
    """
    t = str(raw or "").strip()
    if not t:
        return None
    # ```json … ``` 껍데기부터 벗긴다
    m = re.search(r"```(?:json)?\s*(.+?)```", t, re.S | re.I)
    if m:
        t = m.group(1).strip()
    for pat in (r"\{.*\}", r"\[.*\]"):
        m = re.search(pat, t, re.S)
        if not m:
            continue
        chunk = m.group(0)
        try:
            return json.loads(chunk)
        except json.JSONDecodeError:
            # 끝에 쉼표가 붙거나 홑따옴표를 쓰는 것 정도는 봐준다
            fixed = re.sub(r",\s*([}\]])", r"\1", chunk)
            try:
                return json.loads(fixed)
            except json.JSONDecodeError:
                continue
    return None


async def _ask_json(use: str, system: str, user: str, max_tokens: int = 1500,
                    llm_id: str = "", tries: int = 2):
    """JSON 을 받아 낼 때까지 (짧게) 다시 묻는다. (읽은 것, 마지막 원문)

    한 번 더 묻는 값이 사람이 다시 누르는 값보다 싸다. 두 번을 넘기지는
    않는다 — 안 되는 모델은 세 번도 안 된다.
    """
    raw = ""
    for i in range(max(1, tries)):
        u = user if i == 0 else (user + "\n\n앞의 답은 JSON 이 아니었다. **JSON 만** 출력하라. 설명·인사·코드펜스 금지.")
        raw = await _llm_text(use, system, u, max_tokens=max_tokens, llm_id=llm_id, want_json=True)
        got = _json_from(raw)
        if got is not None:
            return got, raw
    return None, raw


def _anthropic_from(llm: Optional[dict]):
    """등록해 둔 Claude 로 부르는 실물.

    설정 화면에서 키를 넣고 「연결 시험」 까지 통과했는데, 정작 일을 시킬
    때는 `.env` 의 키만 봤다 — 키가 없으면 「쓸 수 있는 LLM 이 없습니다」
    로 끝났다(지적: 등록하고 통신까지 확인했는데 왜 없다고 하나).
    """
    key = str((llm or {}).get("apikey") or "").strip()
    if not key:
        return claude_client
    try:
        ep = str((llm or {}).get("endpoint") or "").strip().rstrip("/")
        kw = {"api_key": key}
        # 기본 주소면 굳이 넘기지 않는다 — SDK 가 알아서 붙인다
        if ep and not ep.startswith("https://api.anthropic.com"):
            kw["base_url"] = ep
        return anthropic.Anthropic(**kw)
    except Exception as e:
        print(f"[_anthropic_from] 등록 Claude 를 세우지 못했습니다: {e}", flush=True)
        return claude_client


def _claude_any():
    """등록된 Claude 중 아무거나, 없으면 `.env` 의 것. (실물, 모델명)"""
    try:
        init_llms_file()
        llms = load_json(LLMS_FILE).get("llms") or []
    except Exception:
        llms = []
    for l in llms:
        if str(l.get("status", "active")) != "active":
            continue
        if str(l.get("type") or "").lower() in ("claude", "anthropic") and l.get("apikey"):
            c = _anthropic_from(l)
            if c is not None:
                return c, str(l.get("model") or "").strip() or CLAUDE_FALLBACK_MODEL
    return claude_client, CLAUDE_FALLBACK_MODEL


def _llm_for(use: str, llm_id: str = "") -> Optional[dict]:
    """이 일에 쓸 LLM 하나.

    `llm_id` 를 주면 그것을 쓴다 — 화면에서 사람이 고른 경우다. 랩 안에
    있는 로컬 LLM 과 Claude 는 잘하는 일이 달라서, 매번 고를 수 있어야 한다.

    안 주면 설정 → Chat LLM 에 등록한 것 중에서 고른다. `uses` 에 이 일의
    이름이 들어 있는 것이 먼저고, 없으면 활성인 아무 것. 그것도 없으면
    None 이고 부르는 쪽이 Claude 로 넘어간다.

    화면에 이미 `uses`·`field_prompts` 칸이 있는데 서버가 아무 데서도 안
    읽고 있었다. 새 설정 화면을 만드는 대신 그 칸을 쓴다.
    """
    try:
        init_llms_file()
        llms = load_json(LLMS_FILE).get("llms") or []
    except Exception as e:
        print(f"[_llm_for] LLM 목록을 읽지 못했습니다: {e}", flush=True)
        return None
    if llm_id:
        got = next((l for l in llms if str(l.get("id")) == llm_id), None)
        if got:
            return got
        # 골라 둔 것이 지워졌다 — 여기서 None 을 주면 「쓸 수 있는 LLM 이
        # 없습니다」 로 끝났다(지적). 산 것 중에서 다시 고른다.
        print(f"[_llm_for] 골라 둔 LLM({llm_id}) 이 목록에 없습니다 — 산 것으로 대신합니다", flush=True)
    # Anthropic 은 주소가 고정이라 칸이 비어 있을 수 있다 — 키가 있으면 산 것
    live = [
        l for l in llms
        if str(l.get("status", "active")) == "active"
        and (l.get("endpoint") or (str(l.get("type") or "").lower() in ("claude", "anthropic") and l.get("apikey")))
    ]
    return next((l for l in live if use in (l.get("uses") or [])), None) or (live[0] if live else None)


async def _llm_text(use: str, system: str, user: str, max_tokens: int = 1500,
                    llm_id: str = "", want_json: bool = False) -> str:
    """등록 LLM 으로 한 번 물어보고 글자만 돌려준다.

    OpenAI 호환(vLLM 등)과 Anthropic 을 둘 다 받는다. 등록된 것이 없으면
    `ANTHROPIC_API_KEY` 로 뜬 기본 Claude 를 쓴다 — 설정이 비어 있어도
    동작은 해야 한다.

    시스템 프롬프트는 설정에서 갈아끼울 수 있다. `field_prompts[use]` 가
    있으면 그것을 쓰고, 없으면 코드의 기본값을 쓴다. 장비 CLI 는 우리 것이
    특이해서 프롬프트를 배포 없이 고칠 수 있어야 한다.
    """
    # 'claude' 는 등록 목록에 없는 특별한 값 — .env 의 기본 Claude 를 뜻한다
    llm = None if llm_id == "claude" else _llm_for(use, llm_id)
    why = ""   # 등록 LLM 이 왜 안 됐는가 — 사람에게 그대로 알려 준다
    sys_p = str(((llm or {}).get("field_prompts") or {}).get(use) or "").strip() or system

    if llm and str(llm.get("type") or "").lower() not in ("claude", "anthropic", "bedrock"):
        import httpx
        base = str(llm.get("endpoint") or "").rstrip("/")
        url = base if base.endswith("/chat/completions") else base + "/chat/completions"
        headers = {"Content-Type": "application/json"}
        if llm.get("apikey"):
            headers["Authorization"] = "Bearer " + str(llm["apikey"])
        body = {
            "model": llm.get("model") or "",
            "max_tokens": int(llm.get("max_tokens") or max_tokens),
            "temperature": float(llm.get("temperature") or 0.7),
            "messages": [{"role": "system", "content": sys_p}, {"role": "user", "content": user}],
        }
        # JSON 이 필요하면 **규격으로** 부탁한다. 말로만 「JSON 만 출력하라」 고
        # 하면 작은 모델은 곧잘 설명을 앞에 붙인다(지적: 모델이 JSON 을 안 줬다).
        if want_json:
            body["response_format"] = {"type": "json_object"}
        try:
            async with httpx.AsyncClient(timeout=120) as c:
                r = await c.post(url, json=body, headers=headers)
                if r.status_code >= 400 and want_json:
                    # 이 규격을 모르는 서버가 있다 — 빼고 한 번 더
                    body.pop("response_format", None)
                    r = await c.post(url, json=body, headers=headers)
                r.raise_for_status()
                d = r.json()
            return str(((d.get("choices") or [{}])[0].get("message") or {}).get("content") or "").strip()
        except Exception as e:
            # 등록 LLM 이 죽어 있을 수 있다. 조용히 실패하지 않고 Claude 로 넘어간다.
            print(f"[_llm_text] 등록 LLM({llm.get('name')}) 호출 실패 → Claude 로 시도: {e}", flush=True)
            why = f"등록 LLM({llm.get('name') or llm.get('model') or '이름 없음'}) 호출 실패: {e}"

    if str((llm or {}).get("type") or "").lower() in ("claude", "anthropic"):
        # 등록해 둔 Claude — 그 키로 부른다(여태 .env 키만 봤다)
        cl, cmodel = _anthropic_from(llm), str((llm or {}).get("model") or "").strip()
    else:
        cl, cmodel = _claude_any()

    if cl is None:
        raise HTTPException(
            503,
            (why + " — 등록한 LLM 주소·모델명을 확인하세요")
            if why else
            "쓸 수 있는 LLM 이 없습니다 — 설정 → Chat LLM 에 등록하거나 "
            ".env 에 ANTHROPIC_API_KEY 를 넣으세요",
        )
    try:
        msg = cl.messages.create(
            model=cmodel or CLAUDE_FALLBACK_MODEL,
            max_tokens=max_tokens,
            system=sys_p,
            messages=[{"role": "user", "content": user}],
        )
        return "".join(b.text for b in msg.content if getattr(b, "type", "") == "text").strip()
    except Exception as e:
        raise HTTPException(502, f"모델 호출에 실패했습니다: {_llm_err(e)}") from e


_DESCRIBE_SYSTEM = """당신은 유비쿼스 네트워크 장비 시험 문서를 쓰는 사람이다.

이미 만들어진 시험 절차(스텝)를 읽고, 그 시험의 **목적**과 **사전 준비 조건**을
한국어로 쓴다. 스텝을 새로 만들거나 고치지 않는다.

규칙:
- 시험 목적은 '무엇을 확인하는 시험인가' 를 두세 문장으로. 명령을 나열하지
  말고, 그 명령들로 무엇을 확인하려는 것인지를 쓴다.
- 사전 준비 조건은 '시작 전에 되어 있어야 하는 것' 을 '- ' 로 시작하는
  목록으로. 스텝에서 읽어낼 수 있는 것만 쓴다(어떤 장비가 몇 대 필요한지,
  어떤 접속이 열려 있어야 하는지, 어떤 설정이 미리 있어야 하는지).
- 스텝에서 알 수 없는 것을 지어내지 않는다. 근거가 없으면 그 항목을 뺀다.
- 이미 적혀 있는 목적·사전조건이 함께 주어지면, 그것을 참고하되 스텝과
  어긋나는 부분은 스텝 쪽을 따른다.

JSON 만 출력한다. 형식:
{"object_md": "...", "precondition_md": "- ...\\n- ..."}
"""


@app.post("/api/tc/{tc_id}/describe")
async def tc_describe(tc_id: str, payload: dict):
    """스텝을 읽고 시험 목적·사전 준비 조건을 제안한다. 저장하지 않는다.

    `/api/tc/{id}/generate` 와 반대 방향이다. 저쪽은 '목적 → 스텝' 이고
    이쪽은 '스텝 → 목적' 이다. 명령어 캡쳐로 스텝을 먼저 만들게 되면서
    남는 일이 문서 쓰기라 이 방향이 필요해졌다.

    저장하지 않는 이유는 generate 와 같다 — 모델이 쓴 글을 사람이 보기 전에
    넣으면 틀린 설명이 그대로 문서가 된다.
    """
    tc_id = _tc_id_norm(tc_id)

    # 화면이 편집 중인 내용을 그대로 보낼 수 있게 payload 를 먼저 본다.
    # 저장하지 않은 스텝으로도 목적을 뽑을 수 있어야 한다 — 캡쳐 직후가
    # 바로 그 순간이다.
    tc = payload.get("tc") if isinstance(payload.get("tc"), dict) else None
    if tc is None:
        tc = await db.tc_get(tc_id)
    if not isinstance(tc, dict):
        raise HTTPException(404, "TC 를 찾을 수 없습니다")

    checks = tc.get("checks") if isinstance(tc.get("checks"), list) else []

    # 스텝이 없어도 쓴다(지적: 단추가 안 켜진다 — 생성 불가).
    #
    # 원래 이 길은 「스텝 → 목적」 이었다. 그런데 시험을 요구사항에서 먼저
    # 뽑아 만들면 이름만 있고 스텝은 아직 없다 — 정작 그때 목적이 필요하다.
    # 스텝이 없으면 **요구사항의 구현의도**를 재료로 쓴다. 그것마저 없으면
    # 이름뿐인데, 그때는 지어내지 말라고 일러 둔다.
    intent = ""
    try:
        rid = str(tc.get("req_id") or "").strip()
        if rid:
            r = await db.req_get(rid)
            if not isinstance(r, dict):
                for x in await db.req_list_full():
                    if str(x.get("reqid") or "") == rid or str(x.get("id") or "") == rid:
                        r = x
                        break
            if isinstance(r, dict):
                intent = str(r.get("desc") or "").strip()
    except Exception as e:
        print(f"[tc_describe] 요구사항을 못 읽었습니다: {e}", flush=True)

    sessions = tc.get("sessions") if isinstance(tc.get("sessions"), list) else []
    sess_txt = []
    for i, dev_id in enumerate(sessions):
        try:
            d = await db.device_get(str(dev_id))
        except Exception as e:
            print(f"[tc_describe] 장비 조회 실패({dev_id}): {e}", flush=True)
            d = None
        if d:
            sess_txt.append(f"- S{i+1}: {d.get('model') or '?'} · {d.get('role') or '?'} · {d.get('ip')}")
        else:
            sess_txt.append(f"- S{i+1}: (등록에 없는 장비 {dev_id})")

    lines = []
    for n, c in enumerate(checks[:200], start=1):
        if not isinstance(c, dict):
            continue
        kind = c.get("kind") or "cli"
        body = (c.get("cli") or c.get("data") or c.get("condition")
                or c.get("text") or c.get("oid") or c.get("step") or "")
        crit = c.get("criteria") or ""
        s = c.get("session")
        who = f"S{int(s)+1}" if isinstance(s, int) else ""
        lines.append(
            f"{n}. [{kind}]{(' ' + who) if who else ''} {str(body).strip()[:200]}"
            + (f"   → 기대: {str(crit).strip()[:120]}" if crit else "")
        )

    user = (
        f"시험 제목: {tc.get('name') or '(없음)'}\n"
        f"TC ID: {tc_id}\n\n"
        f"=== 쓰는 장비 ===\n" + ("\n".join(sess_txt) or "(지정 안 됨)") + "\n\n"
        f"=== 이미 적힌 목적 ===\n{tc.get('object_md') or '(비어 있음)'}\n\n"
        f"=== 이미 적힌 사전조건 ===\n{tc.get('precondition_md') or '(비어 있음)'}\n\n"
        + (
            f"=== 시험 절차 {len(lines)}스텝 ===\n" + "\n".join(lines) + "\n"
            if lines else
            "=== 시험 절차 ===\n(아직 없음 — 아래 요구사항과 시험 제목만 보고 쓴다. "
            "명령어나 수치를 지어내지 말고, 무엇을 확인하는 시험인지와 준비 조건만 적어라.)\n"
        )
        + (f"\n=== 요구사항 구현의도 ===\n{intent[:4000]}\n" if intent else "")
    )

    out, raw = await _ask_json(
        "tc_describe", _DESCRIBE_SYSTEM, user, max_tokens=1500,
        llm_id=str(payload.get("llm") or ""),
    )
    if not isinstance(out, dict):
        raise HTTPException(502, "모델이 JSON 을 돌려주지 않았습니다 — 받은 것: "
                                 + (str(raw)[:160].replace("\n", " ") or "(빈 응답)"))

    return {
        "success": True,
        "object_md": str(out.get("object_md") or ""),
        "precondition_md": str(out.get("precondition_md") or ""),
        "steps": len(lines),
    }


@app.post("/api/req/{req_id}/ai-intent")
async def req_ai_intent(req_id: str, payload: dict):
    """짧은 요청 → **구현의도** 초안. 저장하지 않는다.

    여태 이 글은 손으로만 썼다 — 용도별 프롬프트에 REQ-Intent 자리는 있는데
    화면에 부르는 자리가 없었다(지시: Intent 에도 LLM 드롭바를).

    저장하지 않는 까닭은 목적·스텝 쪽과 같다: 모델이 쓴 글이 사람 눈을
    거치지 않고 문서가 되면 안 된다. 화면이 보여 주고 사람이 「넣기」 를
    누른다.
    """
    req = payload.get("req") if isinstance(payload.get("req"), dict) else None
    if req is None:
        req = await db.req_get(req_id)
    if not isinstance(req, dict):
        raise HTTPException(404, "요구사항을 찾을 수 없습니다")

    say = str(payload.get("text") or "").strip()
    user = (
        f"요구사항 ID: {req.get('reqid') or req_id}\n"
        f"제목: {req.get('title') or '(없음)'}\n"
        f"자리: {' › '.join(str(req.get(k) or '') for k in ('cat1', 'cat2', 'cat3', 'cat4') if req.get(k))}\n\n"
        f"=== 사람이 적은 요청 ===\n{say or '(없음 — 제목과 이미 적힌 글로 다듬어라)'}\n\n"
        f"=== 이미 적힌 구현내용 ===\n{str(req.get('desc') or '')[:4000] or '(비어 있음)'}\n\n"
        "구현의도를 마크다운으로 써라. 제목 줄은 넣지 말고 본문만 쓴다."
    )
    text = await _llm_text(
        "req_intent", _prompt_of("req_intent")["system"], user,
        max_tokens=1800, llm_id=str(payload.get("llm") or ""),
    )
    return {"ok": True, "text": text}


def _tc_title(name: str, obj: str) -> str:
    """제안의 **제목**. 번호를 적어 오면 목적의 첫 줄로 바꾼다.

    모델이 name 자리에 「TC-1」 같은 순번을 적고 정작 제목은 object 에 쓰는
    일이 잦다(지적: 제목이 없어). 그대로 만들면 목록에 TC-1…TC-6 만 서서
    무슨 시험인지 알 수 없다.
    """
    nm = str(name or "").strip()
    ob = str(obj or "").strip()
    if nm and not re.match(r"^(tc|테스트|시험)?[\s\-_#.]*\d+$", nm, re.I):
        return nm
    if not ob:
        return nm
    first = re.split(r"[\n.]", ob)[0].strip()
    return (first or ob)[:160]


@app.post("/api/req/{req_id}/ai-coverage")
async def req_ai_coverage(req_id: str, payload: dict):
    """구현의도 → **덮을 시험 항목 목록** 초안. 저장하지 않는다.

    이름과 목적까지다 — 스텝은 여기서 만들지 않는다(용도 프롬프트의 규칙).
    사람이 골라서 시험항목으로 만든다.
    """
    req = payload.get("req") if isinstance(payload.get("req"), dict) else None
    if req is None:
        req = await db.req_get(req_id)
    if not isinstance(req, dict):
        raise HTTPException(404, "요구사항을 찾을 수 없습니다")

    intent = str(req.get("desc") or "").strip()
    if not intent:
        raise HTTPException(400, "구현내용이 비어 있습니다 — Intent 를 먼저 쓰세요")

    have = []
    try:
        for t in await db.tc_list_meta():
            if str(t.get("req_id") or "") == str(req.get("reqid") or req_id):
                have.append(str(t.get("name") or ""))
    except Exception as e:
        print(f"[req_ai_coverage] 이미 있는 시험을 읽지 못했습니다: {e}", flush=True)

    user = (
        f"요구사항: {req.get('title') or ''}\n\n"
        f"=== 구현의도 ===\n{intent[:6000]}\n\n"
        f"=== 이미 있는 시험 항목 ===\n" + ("\n".join(f"- {x}" for x in have) or "(없음)") + "\n\n"
        "이미 있는 것과 겹치지 않는 시험 항목만 제안하라.\n"
        "name 은 **시험 항목의 제목**이다 — 「무엇을 어떻게 확인하는가」 로 적고, "
        "TC-1 같은 번호나 순번을 적지 마라. object 는 그 시험의 목적이다.\n"
        'JSON 만 출력한다: {"items":[{"name":"...","object":"..."}]}'
    )
    out, raw = await _ask_json(
        "req_coverage", _prompt_of("req_coverage")["system"], user,
        max_tokens=2000, llm_id=str(payload.get("llm") or ""),
    )
    if not isinstance(out, dict):
        # 무엇을 받았는지 함께 보여 준다 — 「JSON 이 아니었다」 만으로는
        # 모델을 바꿔야 하는지 프롬프트를 고쳐야 하는지 알 수 없다(지적)
        raise HTTPException(502, "모델이 JSON 을 돌려주지 않았습니다 — 받은 것: "
                                 + (str(raw)[:160].replace("\n", " ") or "(빈 응답)"))
    items = []
    for x in (out.get("items") or []):
        ob = str(x.get("object") or "").strip()
        nm = _tc_title(x.get("name") or x.get("title") or "", ob)
        if nm:
            items.append({"name": nm, "object": ob})
    return {"ok": True, "items": items, "have": len(have)}


@app.post("/api/req/{req_id}/make-tcs")
async def req_make_tcs(req_id: str, payload: dict, token: str = ""):
    """제안한 시험 항목을 **진짜 시험항목으로 만든다**. 고른 것만.

    제안을 보여만 주고 끝냈더니 옮겨 적을 길이 없었다(지적: 저장을 할 수
    없다). 그렇다고 모델이 만든 것을 곧장 밀어 넣지도 않는다 — 사람이
    고른 것만 만들고, **스텝은 비워 둔다**. 무엇을 어떤 명령으로 볼지는
    Automation 탭에서 정한다.

    자리와 모델은 **요구사항에서 물려받는다** — 요구사항이 선 폴더가
    시험도 설 자리고, 모델그룹·모델명은 그 프로젝트의 것이다. 이 둘을
    사람이 다시 고르게 하면 제안을 받는 뜻이 없다.
    """
    u = _user_from_token(token)
    req = await db.req_get(req_id)
    if not isinstance(req, dict):
        raise HTTPException(404, "요구사항을 찾을 수 없습니다")
    items = [x for x in (payload.get("items") or []) if str((x or {}).get("name") or "").strip()]
    if not items:
        raise HTTPException(400, "만들 항목이 없습니다")

    # 프로젝트(뿌리 폴더)의 모델그룹·모델명 — 신규 시험은 이 둘이 있어야 한다
    mg = str(req.get("model_group") or "")
    md = str(req.get("model") or "")
    root = str(req.get("cat1") or "")
    if root and (not mg or not md):
        async with db.pool().acquire() as c:
            prow = await c.fetchrow(
                "SELECT model_group, model FROM project WHERE cat_id = $1", root
            )
        if prow:
            mg = mg or str(prow["model_group"] or "")
            md = md or str(prow["model"] or "")

    made = []
    async with db.pool().acquire() as c:
        for it in items:
            tcid = await _next_tc_id(c)
            d = {
                "tcid": tcid,
                "name": _tc_title(it.get("name") or "", it.get("object") or ""),
                # 제안이 말한 「무엇을 확인하는가」 는 시험 목적 자리에 그대로 앉힌다
                "object_md": str(it.get("object") or "").strip(),
                "req_id": str(req.get("id") or req.get("reqid") or req_id),
                "status": "작성중",
                "model_group": mg,
                "model": md,
                "checks": [],
                "created_by": (u or {}).get("name") or (u or {}).get("username") or "",
                "updated_by": (u or {}).get("name") or (u or {}).get("username") or "",
            }
            for i in range(4):
                k = f"cat{i + 1}"
                if req.get(k):
                    d[k] = req[k]
            await db.tc_upsert(tcid, d)
            made.append({"tcid": tcid, "name": d["name"]})

    # 요구사항 쪽 포인터에도 적어 둔다 — 연결은 두 곳에 산다
    try:
        cur = list(req.get("tc") or [])
        req["tc"] = cur + [m["tcid"] for m in made if m["tcid"] not in cur]
        await db.req_upsert(str(req.get("id") or req_id), req)
    except Exception as e:
        print(f"[req_make_tcs] 요구사항 쪽 연결을 적지 못했습니다: {e}", flush=True)

    return {"ok": True, "made": made, "model_group": mg, "model": md}


@app.post("/api/tc/{tc_id}/ai-manual")
async def tc_ai_manual(tc_id: str, payload: dict):
    """목적·자동 스텝을 읽고 **수동 시험서** 초안. 저장하지 않는다.

    수동 스텝은 사람이 읽고 따라 하는 글이라 셋으로 나뉜다 —
    무엇을 한다 / 무엇을 넣는다 / 무엇이 나와야 한다.
    """
    tc_id = _tc_id_norm(tc_id)
    tc = payload.get("tc") if isinstance(payload.get("tc"), dict) else None
    if tc is None:
        tc = await db.tc_get(tc_id)
    if not isinstance(tc, dict):
        raise HTTPException(404, "TC 를 찾을 수 없습니다")

    checks = tc.get("checks") if isinstance(tc.get("checks"), list) else []
    auto = []
    for c in checks[:200]:
        if not isinstance(c, dict) or c.get("kind") == "manual":
            continue
        auto.append(f"- {c.get('desc') or ''} / {c.get('cli') or c.get('oid') or ''} → {c.get('criteria') or ''}")

    user = (
        f"시험 제목: {tc.get('name') or '(없음)'}\n\n"
        f"=== 시험 목적 ===\n{str(tc.get('object_md') or '')[:3000] or '(비어 있음)'}\n\n"
        f"=== 사전 준비 조건 ===\n{str(tc.get('precondition_md') or '')[:1500] or '(비어 있음)'}\n\n"
        f"=== 자동 스텝(참고) ===\n" + ("\n".join(auto) or "(없음)") + "\n\n"
        '수동 시험서를 JSON 으로만 출력한다: '
        '{"steps":[{"step":"무엇을 한다","data":"무엇을 넣는다","expected":"무엇이 나와야 한다"}]}'
    )
    out, raw = await _ask_json(
        "coverage_manual", _prompt_of("coverage_manual")["system"], user,
        max_tokens=2500, llm_id=str(payload.get("llm") or ""),
    )
    if not isinstance(out, dict):
        raise HTTPException(502, "모델이 JSON 을 돌려주지 않았습니다 — 받은 것: "
                                 + (str(raw)[:160].replace("\n", " ") or "(빈 응답)"))
    steps = [
        {
            "step": str(x.get("step") or "").strip(),
            "data": str(x.get("data") or "").strip(),
            "expected": str(x.get("expected") or "").strip(),
        }
        for x in (out.get("steps") or []) if str(x.get("step") or "").strip()
    ]
    return {"ok": True, "steps": steps}


@app.get("/api/tc/{tc_id}/cycles")
async def tc_cycles(tc_id: str):
    """이 TC 가 어느 사이클에서 돌았고 결과가 어땠나.

    실행 이력(`/run-history`)과 다른 질문이다. 저쪽은 '이 화면에서 언제
    돌렸나' 고, 이쪽은 '어느 배포 검증에 들어갔나' 다. 자료도 다른 곳에
    있다 — 이력은 파일, 사이클은 DB 의 cycle 테이블이다.
    """
    tc_id = _tc_id_norm(tc_id)
    rows = await db.cycle_of_tc(tc_id)
    out = []
    for r in rows:
        fail = int(r.get("steps_fail") or 0)
        ok = int(r.get("steps_pass") or 0)
        # 옛 자료에는 item.status 가 없다. 스텝 집계로 만든다 —
        # 하나라도 실패면 FAIL, 하나도 안 돈 것은 미실행이다.
        status = (r.get("status") or "").upper()
        if status not in ("PASS", "FAIL"):
            status = "FAIL" if fail else ("PASS" if ok else "")
        at = r.get("executed_at") or r.get("start_date") or r.get("created_at") or ""
        if not at and r.get("updated_at") is not None:
            at = str(r["updated_at"])
        out.append({
            "cycle_id": r.get("cycle_id") or "",
            "model": r.get("model") or "",
            "version": r.get("version") or "",
            "at": str(at)[:19],
            "by": r.get("executed_by") or "",
            "auto": str(r.get("executed_auto") or "").lower() in ("1", "true", "y", "yes"),
            "device": r.get("device") or "",
            "status": status,
            "pass": ok,
            "fail": fail,
            "steps": int(r.get("steps_count") or 0),
            "issues": int(r.get("issues") or 0),
        })
    return {"ok": True, "cycles": out}


@app.get("/api/llm-choices")
async def llm_choices():
    """글을 맡길 수 있는 것들. 화면의 고르는 칸이 이것을 읽는다.

    등록 LLM 과 기본 Claude 를 한 목록으로 준다 — 사람 눈에는 둘 다 그냥
    '누가 쓸 것인가' 이고, 어디에 등록돼 있는지는 사정이다.
    """
    try:
        init_llms_file()
        llms = load_json(LLMS_FILE).get("llms") or []
    except Exception as e:
        print(f"[llm_choices] LLM 목록을 읽지 못했습니다: {e}", flush=True)
        llms = []
    out = [
        {"id": str(l.get("id")), "name": l.get("name") or l.get("model") or "(이름 없음)",
         "model": l.get("model") or "", "local": True}
        for l in llms
        if str(l.get("status", "active")) == "active" and l.get("endpoint")
    ]
    _cl, _cm = _claude_any()
    if _cl is not None:
        out.append({"id": "claude", "name": "Claude", "model": _cm or "claude-sonnet-4-5", "local": False})
    return {"choices": out}


# ══════════════════════════════════════════════════════════════════════
# 사이클 서버 실행
#
# 전에는 브라우저가 실행을 붙들고 있었다. 64건을 걸어 놓고 탭을 닫으면
# 거기서 멈췄고, 자리를 뜰 수가 없었다. 253 을 실행 서버로 둔 의미도
# 여기서 생긴다.
#
# 화면은 **일감을 줄에 걸어 놓고 손을 뗀다.** 실행기(runner 컨테이너)가
# 집어서 돌리고 진행을 여기로 올린다. 서버가 그것을 WebSocket 으로 뿌리면
# 보고 있던 사람들 화면이 같이 움직이고, 안 보고 있었어도 나중에 열어
# 처음부터 다 볼 수 있다.
#
# 판정기는 한 벌이다. 실행기는 화면과 **같은 runner.ts·judge.ts** 를 Node 로
# 묶어 돈다 — 두 벌이면 한 화면에서 적합인 것이 다른 화면에서 부적합이 된다.
# ══════════════════════════════════════════════════════════════════════

# 실행기가 자기를 밝히는 열쇠. 사람의 세션이 없는 자리라 이것으로 가른다.
RUNNER_KEY = os.environ.get("RUNNER_KEY") or ""


def _runner_guard(key: str):
    if not RUNNER_KEY:
        raise HTTPException(503, "RUNNER_KEY 가 정해져 있지 않습니다")
    if str(key or "") != RUNNER_KEY:
        raise HTTPException(403, "실행기 열쇠가 맞지 않습니다")


async def _run_push(run: dict, logs: list = None):
    """진행을 보고 있는 사람들에게 그대로 넘긴다."""
    try:
        msg = {"type": "run_progress", "run": run}
        if logs:
            msg["logs"] = logs
        asyncio.create_task(broadcast(msg))
    except Exception:
        pass


@app.post("/api/runs")
async def run_queue(payload: dict, request: Request):
    """실행을 줄에 건다. 돌리는 것은 실행기가 한다."""
    cycle_id = str(payload.get("cycle_id") or "").strip()
    if not cycle_id:
        raise HTTPException(400, "cycle_id 가 필요합니다")
    cyc = await db.cycle_get(cycle_id)
    if cyc is None:
        raise HTTPException(404, "사이클을 찾을 수 없습니다")

    picked = [int(x) for x in (payload.get("pick") or []) if str(x).lstrip("-").isdigit()]
    items = cyc.get("items") if isinstance(cyc.get("items"), list) else []
    if not picked:
        picked = list(range(len(items)))
    if not picked:
        raise HTTPException(400, "돌릴 항목이 없습니다")

    # 같은 사이클을 둘이 동시에 돌리면 결과를 서로 덮는다. 한 번에 하나만.
    live = await db.run_active(cycle_id)
    if live:
        raise HTTPException(
            409,
            f"이 사이클은 이미 돌고 있습니다 ({live[0].get('started_by') or '누군가'}). "
            "끝나거나 멈춘 뒤에 다시 거세요.",
        )

    # 누가 걸었나. 화면에 「누가 돌리고 있나」 를 보여야 남이 멈추기 전에
    # 한 번 묻게 된다.
    who = ""
    try:
        who = _user_of(_token_from(request)) or ""
    except Exception:
        pass
    run_id = _uuid4().hex[:16]
    run = await db.run_create(
        run_id, cycle_id, str(cyc.get("name") or ""), picked,
        who or str(payload.get("who") or ""), len(picked),
    )
    await _run_push(run)
    return {"ok": True, "run": run}


@app.get("/api/runs")
async def run_list(
    cycle_id: str = "", active: int = 0, limit: int = 200,
    status: str = "", who: str = "", q: str = "",
):
    """실행 목록.

    `cycle_id` 가 있으면 그 사이클의 것, 없으면 **전부**. 「어제 밤에 뭐가
    돌았나」 는 사이클을 하나씩 열어서는 못 답한다 — Executions 화면이
    그 자리다.
    """
    if active:
        return {"runs": await db.run_active(cycle_id)}
    if cycle_id:
        return {"runs": await db.run_recent(cycle_id)}
    return {"runs": await db.run_all(limit, status, who, q), "people": await db.run_people()}


@app.get("/api/runs/{run_id}")
async def run_one(run_id: str, after: int = 0):
    """진행 + 지난 로그.

    브라우저를 닫았다 다시 열면 `after=0` 으로 통째로 받아 그대로 다시
    그린다. 이어서 볼 때는 마지막으로 본 seq 를 준다.
    """
    run = await db.run_get(run_id)
    if run is None:
        raise HTTPException(404, "실행을 찾을 수 없습니다")
    return {"run": run, "logs": await db.run_log_get(run_id, after)}


@app.post("/api/runs/{run_id}/stop")
async def run_stop(run_id: str):
    ok = await db.run_stop_ask(run_id)
    run = await db.run_get(run_id)
    if run:
        await _run_push(run)
    return {"ok": ok, "run": run}


# ── 여기부터는 실행기가 부르는 자리 ──────────────────────────────

@app.post("/api/runner/claim")
async def run_claim(payload: dict):
    _runner_guard(payload.get("key"))
    # 죽은 실행을 먼저 걷어낸다. 안 그러면 running 인 채로 영원히 남아
    # 화면은 계속 도는 줄 알고 기다린다.
    await db.run_sweep_dead()
    run = await db.run_claim(str(payload.get("worker") or "runner"))
    if run:
        await _run_push(run)
    return {"run": run}


@app.post("/api/runner/{run_id}/progress")
async def run_progress(run_id: str, payload: dict):
    _runner_guard(payload.get("key"))
    logs = payload.get("logs") if isinstance(payload.get("logs"), list) else []
    if logs:
        await db.run_log_add(run_id, logs)
    patch = {k: v for k, v in (payload.get("patch") or {}).items()}
    run = await db.run_progress(run_id, patch)
    if run is None:
        raise HTTPException(404, "실행을 찾을 수 없습니다")
    await _run_push(run, logs)
    # 멈춤을 부탁받았는지 실행기에게 알려준다 — 스텝 사이에서 스스로 내려온다
    return {"ok": True, "stop": bool(run.get("stop_asked"))}


@app.post("/api/runner/{run_id}/finish")
async def run_done(run_id: str, payload: dict):
    _runner_guard(payload.get("key"))
    logs = payload.get("logs") if isinstance(payload.get("logs"), list) else []
    if logs:
        await db.run_log_add(run_id, logs)
    run = await db.run_finish(
        run_id, str(payload.get("status") or "done"), str(payload.get("error") or "")
    )
    if run is None:
        raise HTTPException(404, "실행을 찾을 수 없습니다")
    await _run_push(run, logs)
    return {"ok": True}


@app.post("/api/runner/login")
async def runner_login(payload: dict):
    """실행기에게 보통 세션을 하나 내준다.

    실행기는 사이클·TC 를 읽고 장비 세션을 열고 결과를 저장한다 — 사람이
    하는 일과 똑같다. 그래서 인증 길을 따로 파지 않고 **평범한 토큰**을
    준다. 그러면 지금 있는 권한 검사가 그대로 적용된다.

    누가 무엇을 했는지는 남아야 하므로 실행기 자신의 이름으로 남긴다.
    """
    _runner_guard(payload.get("key"))
    token = _secrets.token_hex(16)
    SESSIONS[token] = {
        "username": "runner",
        "role": "관리자",
        "name": "실행 서버",
        "ts": datetime.now().timestamp(),
    }
    _save_one_session(token)
    return {"token": token}


# ══════════════════════════════════════════════════════════════════════
# Reports — 집계
#
# 「지금 이 장비 이 버전이 몇 % 왔나 · 어디가 깨졌나」 는 사이클을 하나씩
# 열어서는 못 답한다. 사이클 24건을 화면이 하나씩 읽게 하면 느리기도 하고,
# 무엇보다 같은 셈을 화면마다 다시 짜게 된다.
#
# 여기서 한 번에 세어 내려준다.
# ══════════════════════════════════════════════════════════════════════

@app.get("/api/report/summary")
async def report_summary():
    cycles = await db.cycle_list_meta()
    rows = []
    for c in cycles:
        cid = c.get("id")
        cname = c.get("name") or ""
        model = c.get("model") or ""
        version = c.get("version") or ""
        vg = c.get("version_group") or ""
        customer = str(c.get("customer") or "")
        cstatus = str(c.get("status") or "")
        for it in (c.get("items") or []):
            if not isinstance(it, dict):
                continue
            # 옛 자료(_verdict 가 없는 요약)는 여기서 다시 센다
            v = it.get("_verdict")
            if v is None:
                v = db.item_verdict(it, it.get("steps") or [])
            steps = it.get("steps") or []
            manual = [s for s in steps if s.get("manual") or s.get("action") == "수동"]
            rows.append({
                "cycle_id": cid,
                "cycle": cname or f"{model} {version}".strip(),
                "model": model,
                "version": version,
                "version_group": vg,
                "customer": customer,
                "cycle_status": cstatus,
                "tcid": it.get("tcid") or "",
                "name": it.get("name") or "",
                "req_id": it.get("req_id") or "",
                "severity": it.get("severity") or "",
                "assignee": it.get("assignee") or it.get("executed_by") or "",
                "executed_at": it.get("executed_at") or "",
                "auto": bool(it.get("executed_auto")),
                # 사람이 할 일인가 장비가 할 일인가
                "kind": ("manual" if manual and len(manual) == len(steps)
                         else "mixed" if manual else "auto" if steps else ""),
                "verdict": v or "",
                "steps": len(steps),
                "fail_steps": int(it.get("_steps_fail") or 0),
            })
    return {"rows": rows, "cycles": [
        {"id": c.get("id"), "name": c.get("name"), "model": c.get("model"),
         "version": c.get("version"), "version_group": c.get("version_group")}
        for c in cycles
    ]}


# ══════════════════════════════════════════════════════════════════════
# 결함 (defect) — 사이클에서 「이슈 생성」 으로 걸고, Defects 화면에서 Jira 로 민다
#
# 항목 하나에 결함 하나. 바로 Jira 로 올리지 않는다 — 64건 돌려 20건 깨지면
# 그중 열여덟은 같은 원인이거나 시험이 잘못된 것이다. UTOP 에 모아 사람이
# 추린 뒤에 민다.
# ══════════════════════════════════════════════════════════════════════

@app.get("/api/defects")
async def defect_list_api(status: str = "", cycle_id: str = ""):
    return {"defects": await db.defect_list(status, cycle_id)}


@app.get("/api/defects/for-item")
async def defect_for_item(cycle_id: str, tcid: str):
    """이 항목에 이미 건 결함이 있나 — 버튼이 「생성」/「봄」 을 가른다."""
    return {"defect": await db.defect_by_item(cycle_id, tcid)}


@app.post("/api/defects")
async def defect_create_api(payload: dict, request: Request):
    """사이클 항목에서 결함을 하나 만든다. 깨진 스텝 내용을 통째로 담는다."""
    cid = str(payload.get("cycle_id") or "").strip()
    tcid = str(payload.get("tcid") or "").strip()
    if not tcid:
        raise HTTPException(400, "tcid 가 필요합니다")
    # 같은 항목에 이미 있으면 그것을 돌려준다 — 항목 하나에 하나만.
    exist = await db.defect_by_item(cid, tcid)
    if exist:
        return {"defect": exist, "existed": True}
    who = ""
    try:
        who = _user_of(_token_from(request)) or ""
    except Exception:
        pass
    # ID 는 DEF-<프로젝트키>-<순번3>. 동시에 두 건이 같은 번호를 집으면
    # PK 가 겹치므로 그때만 번호를 다시 받아 온다.
    d = None
    did = ""
    for _ in range(3):
        did = await db.defect_next_id(str(payload.get("jira_project") or ""))
        try:
            d = await db.defect_create({
                "id": did,
                "title": str(payload.get("title") or payload.get("tc_name") or tcid),
                "severity": payload.get("severity"),
                "cycle_id": cid,
                "cycle_name": payload.get("cycle_name"),
                "tcid": tcid,
                "tc_name": payload.get("tc_name"),
                "model": payload.get("model"),
                "version": payload.get("version"),
                "steps": payload.get("steps") or [],
                "note": payload.get("note"),
                # 이슈 등록 칸 — 프로젝트 키·프로젝트명·이슈유형·우선순위·수정버전·구성요소·보고자
                "jira_project": payload.get("jira_project"),
                "project_name": payload.get("project_name"),
                "issue_type": payload.get("issue_type"),
                "priority": payload.get("priority"),
                "fix_version": payload.get("fix_version"),
                "component": payload.get("component"),
                "reporter": payload.get("reporter"),
                "panels": payload.get("panels") or {},
                "created_by": who,
            })
            break
        except Exception as e:
            if "duplicate key" not in str(e):
                raise
    if d is None:
        raise HTTPException(500, "결함 번호 발급이 계속 겹칩니다 — 다시 시도하세요")
    try: asyncio.create_task(broadcast({"type": "defect_updated", "id": did}))
    except Exception: pass
    return {"defect": d, "existed": False}


# 이슈 본문의 여섯 판 — 화면(Jira 프로젝트 패널 설정)과 같은 차례·같은 이름.
# 번호를 붙이는 것은 사람이 「3번 비었다」 고 말할 수 있게 하기 위해서다.
# 이슈 본문 판 — 웹의 WIKI_PANELS 와 **같은 차례·같은 열쇠**여야 한다.
# 여덟 판으로 바꿨다(지시). 이름이 바뀐 판도 열쇠는 그대로 두었다:
# 열쇠를 새로 지으면 이미 저장된 결함의 그 판이 빈 칸이 된다.
_DEFECT_PANELS = [
    ("symptom", "현상"),
    ("topo", "시험구성도"),
    ("steps", "시험절차"),
    ("detail", "시험내역"),
    ("config", "Configuration File (Config File)"),
    ("core", "Core File (Upload Core file)"),
    ("kernel", "Kernel Log & Syslog 조회"),
    ("attach", "첨부파일"),
]


def _defect_jira_body(d: dict) -> str:
    """결함을 Jira wiki 마크업 설명으로 편다.

    사람이 화면에서 채운 **여덟 판**이 있으면 그것으로 쓴다. 없으면 예전처럼
    깨진 스텝만 편다 — 옛 결함도 그대로 올라가야 한다.

    「시험절차」 는 비어 있으면 스텝으로 채운다. 그 판만큼은 화면이
    자동으로 만들어 주는 것이라, 사람이 손대지 않았다고 빼면 안 된다.
    """
    p = d.get("panels") or {}
    if any(str(p.get(k) or "").strip() for k, _ in _DEFECT_PANELS):
        L = []
        for i, (k, label) in enumerate(_DEFECT_PANELS, 1):
            v = str(p.get(k) or "").strip()
            if k == "steps" and not v:
                v = _defect_steps_body(d)
            if not v:
                continue
            L.append("{panel:title=%d. %s}" % (i, label))
            L.append(v)
            L.append("{panel}")
            L.append("")
        return "\n".join(L)
    return _defect_steps_body(d)


def _defect_steps_body(d: dict) -> str:
    """깨진 스텝을 Jira wiki 마크업으로 편다."""
    L = []
    L.append("h3. 시험 정보")
    L.append("|| 항목 || 내용 |")
    L.append(f"| 사이클 | {d.get('cycle_name') or d.get('cycle_id') or '-'} |")
    L.append(f"| 시험 | {d.get('tc_name') or ''} ({d.get('tcid') or ''}) |")
    L.append(f"| 모델 | {d.get('model') or '-'} |")
    L.append(f"| 버전 | {d.get('version') or '-'} |")
    steps = d.get("steps") or []
    if steps:
        L.append("")
        L.append("h3. 시험 절차 및 결과")
        for s in steps:
            no = s.get("no") or ""
            st = s.get("status") or ""
            desc = (s.get("desc") or s.get("cli") or "").strip()
            L.append(f"h4. #{no} {desc}  ({st})")
            if s.get("cli"):
                L.append("*명령*")
                L.append("{code}" + str(s["cli"]) + "{code}")
            if s.get("criteria"):
                L.append(f"*판정 기준*: {s['criteria']}")
            if s.get("reason"):
                L.append(f"*판정 근거*: {s['reason']}")
            out = str(s.get("output") or "").strip()
            if out:
                L.append("*출력*")
                L.append("{code}" + out[:3000] + "{code}")
    return "\n".join(L)


@app.post("/api/defects/{did}/push")
async def defect_push_jira(did: str, payload: dict = None):
    """결함을 Jira 이슈로 올린다. 프로젝트 키·이슈유형·우선순위·수정버전·구성요소·보고자를 함께 실어 보낸다."""
    payload = payload or {}
    d = await db.defect_get(did)
    if d is None:
        raise HTTPException(404, "결함을 찾을 수 없습니다")
    if d.get("jira_key"):
        return {"ok": True, "key": d["jira_key"], "existed": True, "defect": d}
    # 화면에서 고친 칸이 오면 그것으로 덮는다(먼저 저장하지 않았어도 밀 수 있게)
    proj = payload.get("jira_project") or d.get("jira_project")
    itype = payload.get("issue_type") or d.get("issue_type") or "Defect"
    if not proj:
        return {"ok": False, "error": "프로젝트 키가 필요합니다"}
    fields = {}
    fv = payload.get("fix_version") or d.get("fix_version")
    if fv:
        fields["fixVersions"] = [{"name": fv}]
    comp = payload.get("component") or d.get("component")
    if comp:
        fields["components"] = [{"name": comp}]
    rep = payload.get("reporter") or d.get("reporter")
    if rep:
        fields["reporter"] = {"name": rep}
    # 화면에서 방금 고친 판이 오면 그것으로 쓴다 — 「변경 저장」 을 먼저
    # 누르지 않고 바로 올리는 사람이 있다. 저장 안 했다고 옛 본문이 올라가면
    # 무엇이 올라갔는지 아무도 모른다.
    if isinstance(payload.get("panels"), dict):
        d = {**d, "panels": payload["panels"]}
    # 화면이 만든 본문이 오면 **그것을 그대로** 올린다. 미리보기와 같은
    # 함수에서 나온 글이라, 서버가 다시 만들면 화면에서 본 것과 Jira 에 남는
    # 것이 갈린다 — 그 어긋남은 이슈를 연 사람이 아니라 읽는 개발자가 겪는다.
    desc = str(payload.get("description") or "").strip() or _defect_jira_body(d)
    body = {
        "project": proj,
        "issuetype": itype,
        "summary": payload.get("title") or d.get("title") or d.get("tc_name") or did,
        "description": desc,
        "labels": [x for x in (payload.get("labels") or ["utop"]) if str(x).strip()],
    }
    prio = payload.get("priority") or d.get("priority")
    if prio:
        body["priority"] = prio
    # 화면이 그린 Jira 칸(createmeta)이 오면 함께 싣는다. 프로젝트마다 필수가
    # 다르고(사업자·이슈분류·시험시설…), 여기 없으면 Jira 가 그냥 물린다.
    # 화면이 준 값이 뒤에 온다 — 사람이 방금 고른 것이 이겨야 한다.
    if isinstance(payload.get("fields"), dict):
        fields.update(payload["fields"])
    if fields:
        body["fields"] = fields
    res = await jira_create_issue(body)
    if not res.get("ok"):
        return res
    key = res.get("key", "")
    upd = await db.defect_update(did, {
        "status": "pushed", "jira_key": key,
        "jira_project": proj, "issue_type": itype,
        "priority": prio, "fix_version": fv, "component": comp, "reporter": rep,
        **({"panels": payload["panels"]} if isinstance(payload.get("panels"), dict) else {}),
    })
    try: asyncio.create_task(broadcast({"type": "defect_updated", "id": did}))
    except Exception: pass
    return {"ok": True, "key": key, "url": res.get("url"), "defect": upd}


@app.patch("/api/defects/{did}")
async def defect_update_api(did: str, payload: dict):
    d = await db.defect_update(did, payload)
    if d is None:
        raise HTTPException(404, "결함을 찾을 수 없습니다")
    try: asyncio.create_task(broadcast({"type": "defect_updated", "id": did}))
    except Exception: pass
    return {"defect": d}


@app.delete("/api/defects/{did}")
async def defect_delete_api(did: str):
    ok = await db.defect_delete(did)
    try: asyncio.create_task(broadcast({"type": "defect_updated", "id": did}))
    except Exception: pass
    return {"ok": ok}


# ══════════════════════════════════════════════════════════════════════
# 자연어 시험 (AI Assistant) — 길 9개
#
# 다른 UTOP 서버에서 돌던 것을 옮겨 왔다. **여기가 맨 끝인 이유**: 그 모듈이
# 이 파일의 이름들(app · db · _ai_chat · run_cli …)을 받아 오므로, 위가 전부
# 자리를 잡은 뒤여야 한다.
#
# 못 붙어도 서버는 그대로 뜬다 — 자연어 시험만 죽고 나머지 화면은 산다.
# 시연을 앞두고 화면 한 칸 때문에 전체가 안 뜨는 일은 없어야 한다.
# ══════════════════════════════════════════════════════════════════════
try:
    import nl_test  # noqa: E402,F401
    print("[startup] 자연어 시험(nl_test) 붙음 — /api/ai/nl-* · /api/ai/examples", flush=True)
except Exception as _nl_e:  # pragma: no cover - 기동 로그로만 알린다
    print(f"[startup] 자연어 시험(nl_test) 못 붙임: {_nl_e}", flush=True)
