# backend/engine.py
# 실행 엔진 모듈: netmiko 기반 CLI 실행 · P/F 판정 · TC 개별 실행 · Cycle 자동 실행 · Baseline · PPTX
# 참조(reference) 백엔드에 APIRouter 로 접합한다. main.py 에서 broadcast 를 주입한다.
import asyncio
import json
import re
import textwrap
from datetime import datetime
from pathlib import Path
from typing import Any, Optional

from netmiko import ConnectHandler
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse

# ───────────────────────── 경로 ─────────────────────────
BASE_DIR = Path(__file__).parent.parent
DATA_DIR = BASE_DIR / "data"
TC_DIR = DATA_DIR / "tc"
CYCLE_DIR = DATA_DIR / "cycle"
DEVICES_FILE = DATA_DIR / "devices" / "devices.json"
PPTX_DIR = DATA_DIR / "pptx"
ARTIFACTS_DIR = DATA_DIR / "artifacts"
BASELINES_DIR = DATA_DIR / "baselines"
for _d in (PPTX_DIR, ARTIFACTS_DIR, BASELINES_DIR):
    try:
        _d.mkdir(parents=True, exist_ok=True)
    except Exception:
        pass


# ───────────────────────── 헬퍼 ─────────────────────────
# 판정 결과 상태 라벨 정규화 (내부 대문자 표기 → 저장/표시용 표기). 동작 불변.
_STATUS_LABELS = {"PASS": "Pass", "FAIL": "Fail", "N/A": "N/A"}


def load_json(path: Path) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def save_json(path: Path, data) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def safe_name(value: str) -> str:
    allowed = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789._-"
    return "".join(ch if ch in allowed else "_" for ch in str(value))[:180]


def tc_file_path(tc_id: str) -> Path:
    return TC_DIR / f"{tc_id}.json"


# broadcast 는 main.py 에서 주입한다 (engine.broadcast = broadcast)
async def broadcast(message: dict):  # noqa: D401 - placeholder, overridden by main
    return None


# Cycle 완료 후처리(RAG 자동 색인·AI 요약)도 main.py 에서 주입한다 (engine.on_cycle_complete = ...)
async def on_cycle_complete(cycle_id: str, cycle: dict):  # noqa: D401 - placeholder, overridden by main
    return None


# ───────────────────────── netmiko 실행 ─────────────────────────
def netmiko_device_type(device: Optional[dict] = None) -> str:
    device = device or {}
    explicit = str(device.get("netmiko_device_type") or device.get("device_type") or "").strip()
    if explicit:
        return explicit
    proto = str(device.get("protocol", "SSH")).upper()
    return "terminal_server" if proto == "SSH" else "generic_termserver"


def netmiko_params(device: dict) -> dict:
    return {
        "device_type": netmiko_device_type(device),
        "host": device["ip"],
        "port": int(device.get("port", 22)),
        "username": device.get("username", ""),
        "password": device.get("password", ""),
        "secret": device.get("secret", "") or device.get("enable_password", ""),
        "timeout": int(device.get("timeout", 10) or 10),
        "conn_timeout": int(device.get("conn_timeout", 10) or 10),
        "banner_timeout": int(device.get("banner_timeout", 15) or 15),
        "auth_timeout": int(device.get("auth_timeout", 10) or 10),
        "fast_cli": bool(device.get("fast_cli", False)),
    }


def split_cli_lines(value: Any) -> list:
    if not value:
        return []
    if isinstance(value, list):
        return [str(v).strip() for v in value if str(v).strip()]
    return [line.strip() for line in re.split(r"[\r\n;]+", str(value)) if line.strip()]




def enter_enable_mode(conn, device: dict, read_timeout: int) -> None:
    # netmiko 실행 경로(telnet 등)는 무조건 enable 시도 — User EXEC(>) 진입 방지 (사용자 요청)
    secret = str(device.get("secret") or device.get("enable_password") or device.get("password") or "").strip()
    enable_command = str(device.get("enable_command") or "enable").strip()
    if not enable_command:
        return
    # 이미 enable(#) 모드면 enable 단계 생략
    already = False
    try:
        already = conn.find_prompt().rstrip().endswith("#")
    except Exception:
        pass
    if not already:
        # enable 엔터 → 비밀번호 프롬프트가 뜨면 enable 비번(없으면 로그인 비번) 입력 (secret 없어도 enable은 시도)
        out = conn.send_command_timing(enable_command, read_timeout=read_timeout, strip_prompt=False, strip_command=False)
        if not out.rstrip().endswith("#") and re.search(r"password|passwd|암호|비밀번호|secret", out, re.I) and secret:
            conn.send_command_timing(secret, read_timeout=read_timeout, strip_prompt=False, strip_command=False)
    # enable(#) 후 페이징 끄기 — 유비쿼스/Ericsson-LG 등은 'terminal length 0'만 지원(terminal width 511 미지원)
    try:
        conn.send_command_timing("terminal length 0", read_timeout=read_timeout, strip_prompt=False, strip_command=False)
    except Exception:
        pass


def _clean_cli_output(raw: str, command: str) -> str:
    """netmiko 출력에서 에코된 명령(첫 줄)과 프롬프트(마지막 줄)를 제거해
    paramiko(직접 실행) 출력과 형식을 일치시킨다 → Baseline 비교 정확."""
    text = (raw or "").replace("\r\n", "\n").replace("\r", "\n")
    lines = text.split("\n")
    cmd = (command or "").strip()
    # 앞쪽: 명령 에코 줄 제거
    for i, ln in enumerate(lines):
        if not ln.strip():
            continue
        if cmd and ln.strip() == cmd:
            lines = lines[i + 1:]
        break
    # 뒤쪽: 프롬프트 줄(토큰이 #/>/$/% 로 끝) 및 빈 줄 제거
    while lines and (not lines[-1].strip() or re.match(r"^[\w\.\-\(\)/@:]+[#>\$%]+\s*$", lines[-1].strip())):
        lines.pop()
    return "\n".join(lines)


def ssh_exec(ip: str, port: int, username: str, password: str, command: str, device: Optional[dict] = None) -> str:
    # Baseline 캡처 '실행'(blRunCapture)과 동일하게 paramiko exec_command 사용.
    # → 같은 모드(장비 exec 채널)에서 실행되어 enable 동작·출력 형식이 일치하고,
    #   에코된 명령/프롬프트가 붙지 않아 Baseline 비교가 정확해진다.
    import paramiko
    client = paramiko.SSHClient()
    client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
    client.connect(
        ip, port=int(port or 22), username=username, password=password,
        timeout=int((device or {}).get("conn_timeout", 8) or 8),
        allow_agent=False, look_for_keys=False, banner_timeout=8, auth_timeout=8,
    )
    try:
        read_timeout = int((device or {}).get("read_timeout", 30) or 30)
        stdin, stdout, stderr = client.exec_command(command, timeout=read_timeout)
        out = stdout.read().decode("utf-8", errors="replace")
        err = stderr.read().decode("utf-8", errors="replace")
        return _clean_cli_output(out + err, command)
    finally:
        try:
            client.close()
        except Exception:
            pass


# ───────────────────────── 장비 매칭 ─────────────────────────
def netmiko_exec(device: dict, command: str) -> str:
    """버그수정(#15): telnet(및 명시 device_type) 장비를 netmiko로 실행. SSH는 ssh_exec(paramiko) 유지(Baseline 정합성)."""
    proto = str(device.get("protocol", "SSH")).upper()
    params = netmiko_params(device)
    dt = str(device.get("netmiko_device_type") or device.get("device_type") or "").strip()
    if not dt:
        dt = "cisco_ios_telnet" if proto == "TELNET" else "cisco_ios"
    params["device_type"] = dt
    params["port"] = int(device.get("port") or (23 if proto == "TELNET" else 22))
    read_timeout = int(device.get("read_timeout", 30) or 30)
    conn = ConnectHandler(**params)
    try:
        try:
            enter_enable_mode(conn, device, read_timeout)
        except Exception:
            pass
        lines = split_cli_lines(command)
        is_config = (len(lines) > 1) or any(w in str(command).lower() for w in ("configure terminal", "conf t", "config t", "vlan database"))
        if is_config:
            out = conn.send_config_set(lines)
        else:
            out = conn.send_command_timing(command, read_timeout=read_timeout)
        return _clean_cli_output(out, command)
    finally:
        try:
            conn.disconnect()
        except Exception:
            pass

DEVICE_CATALOG_FILE = DATA_DIR / "state" / "device_catalog.json"  # 레거시 참조용(파일 삭제됨, DB app_kv 로 이전)

def _devices() -> list:
    # Device Registration(카탈로그) 은 DB(app_kv 'device_catalog') 에 저장됨.
    # main._kv_load_sync 는 startup 훅이 채운 캐시를 반환 → sync 컨텍스트에서 안전.
    # 레거시 devices.json 은 파일 그대로 병합(id/ip 중복 제거).
    devs = []
    seen = set()
    # 1) DB 캐시에서 device_catalog 읽기
    try:
        from main import _kv_load_sync  # lazy: main ↔ engine 순환 회피
        catalog = _kv_load_sync("device_catalog", {"devices": []}) or {}
        for d in (catalog.get("devices", []) or []):
            key = str(d.get("id", "")) or str(d.get("ip", ""))
            if key and key in seen:
                continue
            if key:
                seen.add(key)
            devs.append(d)
    except Exception as e:
        print(f"[engine._devices] device_catalog(DB) 읽기 실패: {e}", flush=True)
    # 2) 레거시 devices.json 파일도 병합
    try:
        if DEVICES_FILE.exists():
            for d in (load_json(DEVICES_FILE).get("devices", []) or []):
                key = str(d.get("id", "")) or str(d.get("ip", ""))
                if key and key in seen:
                    continue
                if key:
                    seen.add(key)
                devs.append(d)
    except Exception as e:
        print(f"[engine._devices] {DEVICES_FILE.name} 읽기 실패: {e}", flush=True)
    return devs


def find_cycle_device(cycle: dict) -> Optional[dict]:
    devices = _devices()
    model = str(cycle.get("model", "")).strip()
    if model:
        # 모델/이름 일치하는 첫 번째 장비 반환 — protocol(SSH/Telnet) 무관
        device = next((d for d in devices if str(d.get("model", "")).strip() == model or str(d.get("name", "")).strip() == model), None)
        if device:
            return device
    # 폴백: 모델 미지정이거나 못 찾으면 등록된 첫 번째 장비
    return devices[0] if devices else None


def resolve_step_device(step: dict, default_device: dict) -> dict:
    """참조 TC 스텝의 per-step device(id/model/name) 가 있으면 해당 장비로, 없으면 기본 장비."""
    sd = str(step.get("device", "") or "").strip()
    if not sd:
        return default_device
    devices = _devices()
    dev = (next((d for d in devices if str(d.get("id", "")) == sd), None)
           or next((d for d in devices if str(d.get("model", "")).strip() == sd), None)
           or next((d for d in devices if str(d.get("name", "")).strip() == sd), None))
    return dev or default_device


# ───────────────────────── 판정 ─────────────────────────
def is_executable_cli(command: str) -> bool:
    command = (command or "").strip()
    if not command:
        return False
    non_cli_markers = ["접속", "로그인", "확인한다", "준비", "사전", "console", "ssh/console"]
    if any(marker.lower() in command.lower() for marker in non_cli_markers):
        return False
    return True


def criteria_to_text(criteria: Any) -> str:
    if not criteria:
        return ""
    if isinstance(criteria, str):
        return criteria.strip()
    if isinstance(criteria, list):
        rows = []
        for item in criteria:
            if isinstance(item, str):
                value = item.strip()
            elif isinstance(item, dict):
                key = str(item.get("type") or item.get("key") or item.get("rule") or "").strip()
                val = str(item.get("value") or item.get("text") or "").strip()
                value = f"{key}:{val}" if key and val else str(item.get("criteria") or "").strip()
            else:
                value = str(item).strip()
            if value:
                rows.append(value)
        return "\n".join(rows)
    return str(criteria).strip()


def judge_by_criteria(output: str, criteria: Any, device_id: str = "") -> Optional[tuple]:
    output = output or ""
    criteria = criteria_to_text(criteria)
    if not criteria:
        return None
    # OR 그룹: '||' 단독 줄로 구분 — 그룹 중 하나라도 통과하면 PASS (그룹 내부는 AND)
    _glines = [l.strip() for l in re.split(r"[\r\n]+", criteria) if l.strip()]
    if "||" in _glines:
        _groups, _cur = [], []
        for _l in _glines:
            if _l == "||":
                if _cur:
                    _groups.append(_cur)
                _cur = []
            else:
                _cur.append(_l)
        if _cur:
            _groups.append(_cur)
        if len(_groups) > 1:
            _fails = []
            for _g in _groups:
                _r = judge_by_criteria(output, "\n".join(_g), device_id)
                if _r and _r[0] == "PASS":
                    return "PASS", "여러 조건 중 통과: " + _r[1]
                if _r:
                    _fails.append(_r[1].replace("\n", " "))
            return ("FAIL", "모든 조건 그룹 실패:\n" + ("\n".join(_fails))[:300]) if _fails else None
    out = output.lower()
    rules = [line.strip() for line in re.split(r"[\r\n;]+", criteria) if line.strip()]
    passed = []
    fails = []
    for rule in rules:
        key, _, value = rule.partition(":")
        key = key.strip().lower().replace("-", "_")
        value = value.strip()
        if not value and key not in ("interface_connected", "interfaces_connected", "all_interfaces_connected"):
            value = key
            key = "contains"
        result = None
        if key in ("contains", "contain", "include", "includes"):
            dval = value.replace("\\n", "\n")
            if "\n" in dval.strip():
                def _nb(t):
                    return "\n".join(re.sub(r"[ \t]+", " ", ln.strip()) for ln in t.splitlines() if ln.strip())
                ne, no = _nb(dval), _nb(output)
                if not ne:
                    continue
                result = ("PASS", "출력에 해당 내용 포함됨") if ne.lower() in no.lower() else ("FAIL", "출력에 해당 내용 없음")
            else:
                tokens = [x.strip() for x in dval.split(",") if x.strip()]
                if not tokens:
                    continue
                matched = [x for x in tokens if x.lower() in out]
                result = ("PASS", f"출력에 '{', '.join(matched)}' 포함됨") if matched else ("FAIL", f"출력에 '{', '.join(tokens)}' 없음")
        elif key in ("line_include", "line_includes", "line_contain", "line_contains"):
            # 문장(구절) 기준: 한 줄이면 그 문구가 든 줄을 찾고, 여러 줄이면 각 줄이 모두 출력에 있으면 통과
            dval = value.replace("\\n", "\n")
            vlines = [ln.strip() for ln in dval.splitlines() if ln.strip()]
            if not vlines:
                continue
            out_lines = output.splitlines()
            if len(vlines) > 1:
                missing = [lv for lv in vlines if not any(lv.lower() in ol.lower() for ol in out_lines)]
                result = ("PASS", f"{len(vlines)}개 줄 모두 확인됨") if not missing else ("FAIL", f"없는 줄: {missing[0][:50]}" + (f" 외 {len(missing) - 1}" if len(missing) > 1 else ""))
            else:
                phrase = vlines[0]
                hit = next((ln.strip() for ln in out_lines if phrase.lower() in ln.lower()), None)
                result = ("PASS", f"'{phrase}' 문구가 든 줄 있음") if hit is not None else ("FAIL", f"'{phrase}' 문구가 든 줄 없음")
        elif key in ("line_exclude", "line_excludes", "line_not_contain", "line_not_contains"):
            # 문장(구절) 기준: 입력 문구가 어느 줄에도 없으면 통과
            phrase = value.strip()
            if not phrase:
                continue
            hit = next((ln.strip() for ln in output.splitlines() if phrase.lower() in ln.lower()), None)
            result = ("FAIL", f"'{phrase}' 문구가 줄에 있음(제외 대상)") if hit is not None else ("PASS", f"'{phrase}' 문구 없음")
        elif key in ("line_field", "line_value", "field"):
            # 특정 줄(앵커)에서 일부 값 확인: 앵커가 든 줄에 기대값이 있으면 통과
            parts = value.split("|", 1)
            anchor = parts[0].strip()
            expect = parts[1].strip() if len(parts) > 1 else ""
            if not anchor:
                continue
            cand = [ln.strip() for ln in output.splitlines() if anchor.lower() in ln.lower()]
            if not cand:
                result = ("FAIL", f"'{anchor}' 가 든 줄 없음")
            elif not expect:
                result = ("PASS", f"'{anchor}' 줄 존재")
            elif any(expect.lower() in ln.lower() for ln in cand):
                result = ("PASS", f"'{anchor}' 줄에 '{expect}' 있음")
            else:
                result = ("FAIL", f"'{anchor}' 줄에 '{expect}' 없음 (실제: {cand[0][:60]})")
        elif key in ("range", "between", "num_range"):
            # 숫자 범위: range:앵커|최소~최대 (앵커 든 줄의 숫자가 범위 내면 통과)
            parts = value.split("|", 1)
            if len(parts) == 2:
                anchor, rng = parts[0].strip(), parts[1].strip()
            else:
                anchor, rng = "", parts[0].strip()
            mm = re.match(r"\s*(-?\d+(?:\.\d+)?)\s*[~,]\s*(-?\d+(?:\.\d+)?)", rng)
            if not mm:
                continue
            lo, hi = float(mm.group(1)), float(mm.group(2))
            if lo > hi:
                lo, hi = hi, lo
            cand_lines = [ln for ln in output.splitlines() if (not anchor or anchor.lower() in ln.lower())]
            hit_val, seen = None, []
            for ln in cand_lines:
                for ns in re.findall(r"-?\d+(?:\.\d+)?", ln):
                    seen.append(ns)
                    if lo <= float(ns) <= hi:
                        hit_val = ns
                        break
                if hit_val is not None:
                    break
            label = anchor or "값"
            if hit_val is not None:
                result = ("PASS", f"'{label}' {lo}~{hi} 범위 내: {hit_val}")
            else:
                result = ("FAIL", f"'{label}' {lo}~{hi} 범위 값 없음" + (f" (실제: {', '.join(seen[:3])})" if seen else ""))
        elif key in ("equals", "block", "match", "full"):
            # 여러 줄 블록 일치: 저장 시 줄바꿈을 \n 으로 인코딩. 줄/공백 정규화 후 부분일치 비교.
            expected = value.replace("\\n", "\n")
            def _norm_block(t):
                return "\n".join(re.sub(r"[ \t]+", " ", ln.strip()) for ln in t.splitlines() if ln.strip())
            ne, no = _norm_block(expected), _norm_block(output)
            if not ne:
                continue
            result = ("PASS", "출력 블록 일치") if ne.lower() in no.lower() else ("FAIL", "출력 블록 불일치 (기대한 내용이 출력에 없음)")
        elif key in ("contains_all", "include_all", "includes_all"):
            # 줄바꿈 또는 콤마 구분: 모든 토큰이 출력에 있어야 PASS (선택 줄/문구 검증)
            tokens = [x.strip() for x in re.split(r"[,\n]", value) if x.strip()]
            if not tokens:
                continue
            missing = [x for x in tokens if x.lower() not in out]
            result = ("FAIL", f"필수 문자열 없음: {', '.join(missing)}") if missing else ("PASS", f"필수 문자열 모두 확인: {', '.join(tokens)}")
        elif key in ("not_contains", "exclude", "excludes"):
            tokens = [x.strip() for x in value.split(",") if x.strip()]
            if not tokens:
                continue
            found = [x for x in tokens if x.lower() in out]
            result = ("FAIL", f"금지 문자열 발견: {', '.join(found)}") if found else ("PASS", "금지 문자열 없음")
        elif key == "regex":
            if not value:
                continue
            try:
                result = ("PASS", f"정규식 일치: {value}") if re.search(value, output, re.I | re.M) else ("FAIL", f"정규식 미일치: {value}")
            except re.error as exc:
                return "FAIL", f"정규식 오류: {exc}"
        elif key in ("interface_connected", "interfaces_connected", "all_interfaces_connected"):
            bad_words = ["down", "disable", "disabled", "disconnect", "notconnect", "not-connect", "inactive", "fail", "fault"]
            lines = [line.strip() for line in output.splitlines() if line.strip()]
            data_lines = [line for line in lines if re.search(r"\b(eth|ge|xe|pon|port|interface|gi|te)\S*", line, re.I)]
            bad = [line for line in data_lines if any(word in line.lower() for word in bad_words)]
            if bad:
                result = ("FAIL", "비정상 interface 발견: " + " / ".join(bad[:3]))
            elif data_lines and any(re.search(r"\b(connect|connected|up)\b", line, re.I) for line in data_lines):
                result = ("PASS", "interface 상태 정상")
            else:
                result = ("FAIL", "interface 상태를 판정할 출력이 부족함")
        elif key in ("baseline", "bl"):
            if "/" in value:
                bl_dev, bl_key = value.split("/", 1)
                bl_dev, bl_key = bl_dev.strip(), bl_key.strip()
            else:
                bl_dev, bl_key = device_id, value.strip()
            if not bl_key:
                continue
            bl = load_baseline(bl_dev, bl_key)
            if not bl:
                result = ("FAIL", f"Baseline '{bl_dev}/{bl_key}' 없음")
            else:
                extra = bl.get("masks", [])
                masked_bl = normalize_for_baseline(apply_baseline_masks(bl.get("raw", ""), extra))
                masked_out = normalize_for_baseline(apply_baseline_masks(output, extra))
                if masked_bl == masked_out:
                    result = ("PASS", f"Baseline '{bl_key}' 일치")
                else:
                    bl_lines, out_lines = masked_bl.splitlines(), masked_out.splitlines()
                    diffs = []
                    if len(bl_lines) != len(out_lines):
                        diffs.append(f"라인 수 기대 {len(bl_lines)} 실제 {len(out_lines)}")
                    for i in range(max(len(bl_lines), len(out_lines))):
                        bl_l = bl_lines[i] if i < len(bl_lines) else ""
                        out_l = out_lines[i] if i < len(out_lines) else ""
                        if bl_l != out_l:
                            diffs.append(f"라인{i+1} 불일치: 기대[{bl_l[:70]}] 실제[{out_l[:70]}]")
                        if len(diffs) >= 5:
                            break
                    result = ("FAIL", f"Baseline '{bl_key}' 불일치 - {' / '.join(diffs) if diffs else '내용 불일치'}")
        if not result:
            continue
        if result[0] == "FAIL":
            fails.append(result[1])
        else:
            passed.append(result[1])
    if fails:
        return "FAIL", "\n".join(fails)
    if passed:
        return "PASS", "\n".join(passed)
    return None


def judge_cli_result(output: str, expected: str, criteria: Any = "", device_id: str = "") -> tuple:
    output = output or ""
    expected = (expected or "").strip()
    fail_patterns = ["[오류]", "% invalid input", "invalid input", "unknown command", "command not found", "syntax error", "permission denied", "authentication failed"]
    if any(p in output.lower() for p in fail_patterns):
        return "FAIL", "장비 출력에 오류 패턴이 포함됨"
    criteria_result = judge_by_criteria(output, criteria, device_id)
    if criteria_result:
        return criteria_result
    if expected:
        exp, out = expected.lower(), output.lower()
        if exp in out:
            return "PASS", "기대 문자열이 출력에 포함됨"
        natural_ok = ["정상", "출력", "표시", "확인"]
        if any(k in exp for k in natural_ok) and output.strip():
            if "version" in out and ("버전" in expected or "version" in exp):
                return "PASS", "버전 정보가 출력됨"
            return "PASS", "오류 없이 출력이 반환됨"
        return "FAIL", "기대 결과 문자열을 출력에서 찾지 못함"
    return ("PASS", "오류 없이 출력이 반환됨") if output.strip() else ("FAIL", "장비 출력이 비어 있음")


def step_command(step: dict) -> str:
    # 참조 TC 스텝은 cli, 내 형식은 input/command
    return (step.get("input") or step.get("cli") or step.get("command") or "").strip()


def step_criteria(step: dict) -> str:
    return criteria_to_text(
        step.get("pass_criteria_list") or step.get("pass_criteria") or step.get("criteria") or step.get("pass") or ""
    )


async def run_tc_step(device: dict, step: dict, command: Optional[str] = None) -> str:
    device = resolve_step_device(step, device)
    if not device:
        return "[오류] 실행할 장비 정보를 찾을 수 없습니다."
    command = step_command(step) if command is None else command
    if not command:
        return "[오류] 실행할 CLI 명령어가 없습니다."
    if not is_executable_cli(command):
        return f"[SKIP] 자동 CLI 실행 대상이 아닌 절차입니다: {command}"
    proto = str(device.get("protocol", "")).upper()
    ip = device.get("ip", "")
    if proto == "TELNET":
        port = int(device.get("port") or 23)
        header = f"[Telnet → {ip}:{port}]\n"
        out = await asyncio.to_thread(netmiko_exec, device, command)
        return header + out
    if proto and proto != "SSH":
        return f"[오류] TC 자동 실행은 SSH/TELNET만 지원합니다. 현재 프로토콜: {proto}"
    port = int(device.get("port") or 22)
    header = f"[SSH → {ip}:{port}]\n"
    out = await asyncio.to_thread(ssh_exec, ip, port, device.get("username", ""), device.get("password", ""), command, device)
    return header + out


# ───────────────────────── Baseline ─────────────────────────
DEFAULT_MASKS = [
    (r"\b\d+w\d+d\b", "**"),
    (r"\b\d+h\d+m\b", "**"),
    (r"\b\d{1,3}:\d{2}:\d{2}\b", "**"),
    (r"\b\d+ days?,\s*\d{1,2}:\d{2}:\d{2}\b", "**"),
    (r"\b\d[\d,]*\s+(packets|bytes|errors|dropped|frames|resets|collisions|overruns|ignored|watchdog)\b", r"** \1"),
    (r"\b\d+(?:\.\d+)?\s+(Kbps|Mbps|Gbps|bps|Kbit/s|Mbit/s|Gbit/s)\b", r"** \1"),
    (r"\b\d+ bits/sec\b", "** bits/sec"),
    (r"\b\d+ packets/sec\b", "** packets/sec"),
    (r"(Last input|Last output|Last clearing of)\s+\S+", r"\1 **"),
]


def apply_baseline_masks(text: str, extra_masks: list) -> str:
    for pattern, replace in DEFAULT_MASKS:
        try:
            text = re.sub(pattern, replace, text, flags=re.I | re.M)
        except re.error:
            pass
    for m in (extra_masks or []):
        try:
            text = re.sub(m["pattern"], m.get("replace", "**"), text, flags=re.I | re.M)
        except (re.error, KeyError):
            pass
    return text


def normalize_for_baseline(text: str) -> str:
    lines = [re.sub(r"[ \t]+", " ", line.strip()) for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


def baseline_file_path(device_id: str, key: str) -> Path:
    return BASELINES_DIR / safe_name(device_id) / f"{safe_name(key)}.json"


def load_baseline(device_id: str, key: str) -> Optional[dict]:
    path = baseline_file_path(device_id, key)
    if not path.exists():
        return None
    try:
        return load_json(path)
    except Exception:
        return None


# ───────────────────────── PPTX ─────────────────────────
def terminal_image(text: str, path: Path):
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as exc:
        raise RuntimeError("Pillow가 설치되어 있지 않습니다.") from exc
    font = ImageFont.load_default()
    lines = []
    for line in (text or "").splitlines() or [""]:
        wrapped = textwrap.wrap(line, width=118, replace_whitespace=False, drop_whitespace=False) or [""]
        lines.extend(wrapped)
    lines = lines[:80]
    char_w, line_h, width = 7, 16, 920
    height = max(180, 26 + line_h * len(lines))
    img = Image.new("RGB", (width, height), "#0a0c10")
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, width, 24), fill="#1e2130")
    draw.text((12, 7), "CLI Result", fill="#e8eaf0", font=font)
    y = 34
    for line in lines:
        draw.text((14, y), line[: int((width - 28) / char_w)], fill="#a8ff78", font=font)
        y += line_h
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path)


def create_cycle_pptx(cycle: dict) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("python-pptx가 설치되어 있지 않습니다.") from exc
    PPTX_DIR.mkdir(parents=True, exist_ok=True)
    ARTIFACTS_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]

    def add_textbox(slide, left, top, width, height, text, size=14, bold=False):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(text or "")
        run.font.size = Pt(size)
        run.font.bold = bold
        return box

    results = cycle.get("results", [])
    total = len(cycle.get("items", [])) or len([x for x in cycle.get("tc_ids", []) if x and x != "on"]) or len(results)
    passed = len([r for r in results if r.get("status") == "PASS"])
    failed = len([r for r in results if r.get("status") == "FAIL"])
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.45, 0.35, 9.0, 0.5, "TC Cycle Result", 24, True)
    add_textbox(slide, 0.45, 0.95, 9.0, 0.35, f"{cycle.get('model','-')} / {cycle.get('version','-')} / {cycle.get('date', cycle.get('created_at','-'))}", 14)
    add_textbox(slide, 0.45, 1.45, 9.0, 0.4, f"Total {total} / PASS {passed} / FAIL {failed}", 16, True)
    add_textbox(slide, 0.45, 2.05, 9.0, 3.8, cycle.get("note", ""), 12)

    for result in results:
        tcid = result.get("tcid", "")
        if not tcid or tcid == "on":
            continue
        tc_path = tc_file_path(tcid)
        tc = load_json(tc_path) if tc_path.exists() else {"tcid": tcid}
        slide = prs.slides.add_slide(blank)
        add_textbox(slide, 0.35, 0.25, 9.3, 0.35, f"{tcid} - {result.get('status','대기')}", 17, True)
        add_textbox(slide, 0.35, 0.65, 9.3, 0.3, tc.get("name", ""), 11)
        add_textbox(slide, 0.35, 1.0, 9.3, 0.65, tc.get("overview") or tc.get("object") or "", 10)
        y = 1.75
        for idx, step in enumerate((result.get("steps") or [])[:3], start=1):
            command = step.get("command", "")
            output = step.get("output", "")
            status = step.get("result", "")
            add_textbox(slide, 0.35, y, 9.2, 0.25, f"Step {idx} [{status}] {step.get('desc','')}", 11, True)
            add_textbox(slide, 0.35, y + 0.25, 9.2, 0.22, f"$ {command}", 9)
            img_path = ARTIFACTS_DIR / safe_name(cycle.get("id", "cycle")) / f"{safe_name(tcid)}_step{idx}.png"
            terminal_image(output, img_path)
            slide.shapes.add_picture(str(img_path), Inches(0.35), Inches(y + 0.55), width=Inches(9.0))
            y += 2.05
            if y > 5.2:
                break

    pptx_path = PPTX_DIR / f"{safe_name(cycle.get('id','cycle'))}.pptx"
    prs.save(pptx_path)
    return pptx_path


# ───────────────────────── TC 절차 실행 (공통) ─────────────────────────
def _resolve_step_model(step, model):
    """Cycle 모델에 맞는 모델 오버라이드(cli/criteria)를 적용. 없으면 기본값."""
    cli = step_command(step)
    crit = step_criteria(step)
    m = str(model or "").strip()
    if m:
        for ov in (step.get("overrides") or []):
            if str(ov.get("model", "")).strip() == m:
                if str(ov.get("cli", "")).strip():
                    cli = str(ov.get("cli")).strip()
                if str(ov.get("criteria", "")).strip():
                    crit = criteria_to_text(ov.get("criteria"))
                break
    return cli, crit


async def _execute_tc(tc: dict, tc_id: str, device: dict, ev_prefix: str, cycle_id: str = "", owner: str = "", version: str = "", model: str = ""):
    """TC 의 steps 를 실행·판정하고 (status, memo, step_results) 반환. TC 파일에 결과 기록."""
    steps = tc.get("steps") or []
    if not steps:
        memo = "등록된 시험 절차가 없습니다"
        tc["status"] = "N/A"
        save_json(tc_file_path(tc_id), tc)
        return "N/A", memo, []

    step_results = []
    for idx, step in enumerate(steps):
        started = datetime.now().isoformat()
        command, _crit = _resolve_step_model(step, model)
        await broadcast({"type": f"{ev_prefix}_step_start", "tcid": tc_id, "cycle_id": cycle_id, "seq": idx + 1, "desc": step.get("desc", ""), "command": command})
        try:
            output = await run_tc_step(device, step, command)
        except Exception as exc:
            output = f"[오류] {exc}"
        criteria = _crit
        if output.startswith("[SKIP]"):
            result, reason = "N/A", output.replace("[SKIP] ", "")
        else:
            result, reason = judge_cli_result(output, step.get("expected", ""), criteria, device.get("id", ""))
        # 프론트(cycleItemStatus)·DB 일관성을 위해 Pass/Fail/N/A 라벨로 정규화
        result = _STATUS_LABELS.get(result, result)
        step["output"] = output
        step["result"] = result
        step["reason"] = reason
        step["pass_criteria"] = criteria
        step["executed_at"] = started
        step["device_id"] = device.get("id", "")
        step["device_model"] = device.get("model", "")
        sr = {"seq": idx + 1, "desc": step.get("desc", ""), "command": command, "expected": step.get("expected", ""),
              "pass_criteria": criteria, "result": result, "reason": reason, "output": output}
        step_results.append(sr)
        await broadcast({"type": f"{ev_prefix}_step_done", "tcid": tc_id, "cycle_id": cycle_id, **sr})

    executable = [s for s in step_results if s["result"] != "N/A"]
    status = "FAIL" if any(s["result"] == "Fail" for s in executable) else "PASS" if executable else "N/A"
    reasons = [f"Step {s['seq']}: {s['reason']}" for s in step_results if s["result"] in ("Fail", "N/A")]
    memo = "; ".join(reasons[:3])
    tc["status"] = status
    tc.setdefault("history", []).append({
        "date": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "status": status, "owner": owner, "version": version,
        "memo": memo or f"자동 실행 ({device.get('model','')})",
    })
    save_json(tc_file_path(tc_id), tc)
    return status, memo, step_results


def _cycle_tcids(cycle: dict) -> list:
    """참조(items) / 레거시(tc_ids) 모두에서 실행 대상 tcid 목록 추출."""
    if cycle.get("items"):
        return [str(it.get("tcid", "")).strip() for it in cycle["items"] if it.get("tcid")]
    return [x for x in cycle.get("tc_ids", []) if x and x != "on"]


def _apply_item_result(cycle: dict, tcid: str, status: str, memo: str, step_results: list, exec_seconds: float = None):
    """참조 cycle items[] 에 실행 결과 반영 (overall + per-step result)."""
    label = _STATUS_LABELS.get(status, status)
    today = datetime.now().strftime("%Y-%m-%d")
    for it in cycle.get("items", []):
        if str(it.get("tcid", "")).strip() != tcid:
            continue
        it["status"] = status
        it["result"] = label
        it["memo"] = memo
        it["last_run"] = datetime.now().isoformat()
        if exec_seconds is not None:
            it["exec_seconds"] = exec_seconds
        isteps = it.get("steps") or []
        for i, sr in enumerate(step_results):
            r = _STATUS_LABELS.get(sr["result"], sr["result"])
            if i < len(isteps):
                isteps[i]["result"] = r
                isteps[i]["date"] = today
                isteps[i]["output"] = sr.get("output", "")
                isteps[i]["reason"] = sr.get("reason", "")
                isteps[i]["manual"] = False
            else:
                isteps.append({"criteria": sr.get("command", ""), "desc": sr.get("desc", ""), "result": r, "date": today, "output": sr.get("output", ""), "reason": sr.get("reason", ""), "manual": False})
        it["steps"] = isteps  # 버그수정: 잘라내지 않음(데이터 손실 방지)
        break


async def _execute_cycle_item(item, device, cycle_id, owner="", version=""):
    """버그수정: 사이클 item의 steps(action/cli/criteria/type)를 직접 실행·판정. item.steps in-place 갱신."""
    import re as _re2
    isteps = item.get("steps") or []
    tcid = str(item.get("tcid", "")).strip()
    if not isteps:
        return "N/A", "실행할 스텝이 없습니다"
    today = datetime.now().strftime("%Y-%m-%d")
    _ctmap = {"contains": "contains", "line": "line_include", "notcontains": "not_contains", "not_contains": "not_contains", "diff": "equals", "table": "contains"}
    results = []
    for idx, st in enumerate(isteps):
        action = str(st.get("action") or "CLI").strip()
        cli = step_command(st)
        crit = step_criteria(st)
        ctype = str(st.get("type", "")).strip().lower()
        if crit and ctype and (":" not in crit) and ("||" not in crit):
            crit = _ctmap.get(ctype, ctype) + ":" + crit
        await broadcast({"type": "cycle_step_start", "cycle_id": cycle_id, "tcid": tcid, "seq": idx + 1, "command": cli})
        if action in ("대기", "Wait", "wait"):
            _m = _re2.search(r"(\d+)", (cli or "") + " " + str(st.get("desc", "")))
            sec = min(int(_m.group(1)) if _m else 1, 30)
            await asyncio.sleep(sec)
            output, result, reason = (str(sec) + "초 대기 완료"), "N/A", "대기 스텝"
        elif action.startswith("Ping") or action.startswith("SNMP") or action in ("호출",):
            output, result, reason = "", "N/A", action + " 스텝은 서버 자동실행 미지원 (수동 확인)"
        elif (not cli) or (not is_executable_cli(cli)):
            output, result, reason = "", "N/A", "실행 대상 CLI 아님"
        else:
            try:
                output = await run_tc_step(device, st, cli)
            except Exception as exc:
                output = "[오류] " + str(exc)
            if output.startswith("[SKIP]"):
                result, reason = "N/A", output.replace("[SKIP] ", "")
            else:
                result, reason = judge_cli_result(output, st.get("expected", ""), crit, device.get("id", ""))
            result = _STATUS_LABELS.get(result, result)
        st["output"] = output
        st["result"] = result
        st["reason"] = reason
        st["date"] = today
        st["manual"] = False
        results.append((idx, result, reason))
        await broadcast({"type": "cycle_step_done", "cycle_id": cycle_id, "tcid": tcid, "seq": idx + 1, "result": result, "reason": reason, "output": output})
    executable = [r for r in results if r[1] != "N/A"]
    status = "FAIL" if any(r[1] == "Fail" for r in executable) else ("PASS" if executable else "N/A")
    fails = ["Step " + str(i + 1) + ": " + rs for (i, rv, rs) in results if rv == "Fail"]
    memo = ("; ".join(fails))[:200] or ("자동 실행 (" + (device.get("model", "") or device.get("name", "")) + ")")
    return status, memo

# ───────────────────────── 라우터 ─────────────────────────
router = APIRouter()


@router.post("/api/tc/{tc_id}/run")
async def run_single_tc(tc_id: str, body: Optional[dict] = None):
    """단일 TC 개별 실행. body: device_id(특정 장비), cycle_id(해당 Cycle 장비로 실행하고 결과 반영)."""
    body = body or {}
    path = tc_file_path(tc_id)
    if not path.exists():
        raise HTTPException(404, "TC를 찾을 수 없습니다")
    tc = load_json(path)

    cycle = None
    cycle_path = None
    cycle_id = str(body.get("cycle_id", "")).strip()
    if cycle_id:
        cycle_path = CYCLE_DIR / f"{cycle_id}.json"
        if cycle_path.exists():
            cycle = load_json(cycle_path)

    devices = _devices()
    device = None
    device_id = str(body.get("device_id", "")).strip()
    if device_id:
        device = next((d for d in devices if str(d.get("id", "")) == device_id), None)
        if not device:
            raise HTTPException(404, "선택한 장비를 찾을 수 없습니다")
    if not device and cycle:
        device = find_cycle_device(cycle)
    if not device:
        model = str(tc.get("model", "")).strip()
        if model:
            device = next((d for d in devices if str(d.get("model", "")).strip() == model or str(d.get("name", "")).strip() == model), None)
        if not device:
            device = devices[0] if devices else None
    if not device:
        raise HTTPException(404, "실행할 장비를 찾을 수 없습니다")

    await broadcast({"type": "tc_run_start", "tcid": tc_id, "cycle_id": cycle_id,
                     "device": f"{device.get('model','')} / {device.get('ip','')}", "total": len(tc.get("steps") or [])})
    _t0 = datetime.now()
    status, memo, step_results = await _execute_tc(tc, tc_id, device, "tc_run", cycle_id,
                                                   owner=(cycle or {}).get("owner", ""), version=(cycle or {}).get("version", ""), model=(cycle or {}).get("model", "") or tc.get("model", ""))
    _elapsed = round((datetime.now() - _t0).total_seconds(), 1)
    if cycle is not None and cycle_path is not None:
        _apply_item_result(cycle, tc_id, status, memo, step_results, _elapsed)
        results = cycle.get("results") or []
        entry = {"tcid": tc_id, "status": status, "memo": memo, "reason": memo, "steps": step_results}
        for i, rr in enumerate(results):
            if rr.get("tcid") == tc_id:
                results[i] = entry
                break
        else:
            results.append(entry)
        cycle["results"] = results
        save_json(cycle_path, cycle)
        await broadcast({"type": "cycle_tc_done", "cycle_id": cycle_id, "tcid": tc_id, "status": status, "memo": memo, "exec_seconds": _elapsed})

    await broadcast({"type": "tc_run_done", "tcid": tc_id, "cycle_id": cycle_id, "status": status, "memo": memo, "steps": step_results})
    return {"success": True, "tcid": tc_id, "cycle_id": cycle_id, "status": status, "memo": memo, "steps": step_results,
            "device": {"id": device.get("id", ""), "model": device.get("model", ""), "ip": device.get("ip", "")}}


@router.post("/api/cycle/{cycle_id}/run")
async def run_cycle(cycle_id: str):
    cycle_path = CYCLE_DIR / f"{cycle_id}.json"
    if not cycle_path.exists():
        raise HTTPException(404, "Cycle을 찾을 수 없습니다")
    cycle = load_json(cycle_path)
    device = find_cycle_device(cycle)
    if not device:
        raise HTTPException(404, "실행할 장비를 찾을 수 없습니다")

    tcids = _cycle_tcids(cycle)
    await broadcast({"type": "cycle_run_start", "cycle_id": cycle_id,
                     "device": f"{device.get('model','')} / {device.get('ip','')}", "total": len(tcids)})
    run_results = []
    for tcid in tcids:
        await broadcast({"type": "cycle_tc_start", "cycle_id": cycle_id, "tcid": tcid})
        item = next((it for it in cycle.get("items", []) if str(it.get("tcid", "")).strip() == tcid), None)
        if item and item.get("steps"):
            _t0 = datetime.now()
            status, memo = await _execute_cycle_item(item, device, cycle_id, cycle.get("owner", ""), cycle.get("version", ""))
            _elapsed = round((datetime.now() - _t0).total_seconds(), 1)
            item["status"] = status
            item["result"] = _STATUS_LABELS.get(status, status)
            item["memo"] = memo
            item["last_run"] = datetime.now().isoformat()
            item["exec_seconds"] = _elapsed
            run_results.append({"tcid": tcid, "status": status, "memo": memo, "reason": memo, "steps": item.get("steps", []), "exec_seconds": _elapsed})
            await broadcast({"type": "cycle_tc_done", "cycle_id": cycle_id, "tcid": tcid, "status": status, "memo": memo, "exec_seconds": _elapsed})
            continue
        path = tc_file_path(tcid)
        if not path.exists():
            memo = "TC 파일 없음"
            run_results.append({"tcid": tcid, "status": "FAIL", "memo": memo, "reason": memo, "steps": []})
            _apply_item_result(cycle, tcid, "FAIL", memo, [])
            await broadcast({"type": "cycle_tc_done", "cycle_id": cycle_id, "tcid": tcid, "status": "FAIL", "memo": memo})
            continue
        tc = load_json(path)
        _t0 = datetime.now()
        status, memo, step_results = await _execute_tc(tc, tcid, device, "cycle", cycle_id,
                                                       owner=cycle.get("owner", ""), version=cycle.get("version", ""), model=cycle.get("model", ""))
        _elapsed = round((datetime.now() - _t0).total_seconds(), 1)
        run_results.append({"tcid": tcid, "status": status, "memo": memo, "reason": memo, "steps": step_results, "exec_seconds": _elapsed})
        _apply_item_result(cycle, tcid, status, memo, step_results, _elapsed)
        await broadcast({"type": "cycle_tc_done", "cycle_id": cycle_id, "tcid": tcid, "status": status, "memo": memo, "exec_seconds": _elapsed})

    cycle["results"] = run_results
    cycle["executed_at"] = datetime.now().isoformat()
    cycle["device_id"] = device.get("id", "")
    cycle["device_model"] = device.get("model", "")
    if cycle.get("make_ppt"):
        try:
            pptx_path = create_cycle_pptx(cycle)
            cycle["pptx"] = str(pptx_path.relative_to(BASE_DIR))
        except Exception as exc:
            cycle["pptx_error"] = str(exc)
    save_json(cycle_path, cycle)
    await broadcast({"type": "cycle_run_done", "cycle_id": cycle_id, "cycle": cycle})
    # 완료 후처리(RAG 색인·AI 요약)는 백그라운드 — 응답을 지연시키지 않고, 실패해도 무시
    try:
        asyncio.create_task(on_cycle_complete(cycle_id, cycle))
    except Exception:
        pass
    return {"success": True, "cycle": cycle}


@router.get("/api/cycle/{cycle_id}/ppt")
async def download_cycle_ppt(cycle_id: str):
    cycle_path = CYCLE_DIR / f"{cycle_id}.json"
    if not cycle_path.exists():
        raise HTTPException(404, "Cycle을 찾을 수 없습니다")
    cycle = load_json(cycle_path)
    pptx_path = PPTX_DIR / f"{safe_name(cycle_id)}.pptx"
    if not pptx_path.exists():
        try:
            pptx_path = create_cycle_pptx(cycle)
            cycle["pptx"] = str(pptx_path.relative_to(BASE_DIR))
            save_json(cycle_path, cycle)
        except Exception as exc:
            raise HTTPException(500, str(exc))
    return FileResponse(str(pptx_path),
                        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        filename=f"{safe_name(cycle_id)}.pptx")


# ── Baseline API ──
@router.get("/api/baselines")
async def list_all_baselines():
    result = {}
    if BASELINES_DIR.exists():
        for device_dir in sorted(BASELINES_DIR.iterdir()):
            if not device_dir.is_dir():
                continue
            baselines = []
            for f in sorted(device_dir.glob("*.json")):
                try:
                    b = load_json(f)
                    baselines.append({"key": f.stem, "command": b.get("command", ""), "captured_at": b.get("captured_at", ""), "description": b.get("description", "")})
                except Exception:
                    pass
            if baselines:
                result[device_dir.name] = baselines
    return result


@router.get("/api/baselines/{device_id}")
async def list_device_baselines(device_id: str):
    device_dir = BASELINES_DIR / safe_name(device_id)
    baselines = []
    if device_dir.exists():
        for f in sorted(device_dir.glob("*.json")):
            try:
                b = load_json(f)
                baselines.append({"key": f.stem, "command": b.get("command", ""), "captured_at": b.get("captured_at", ""),
                                  "description": b.get("description", ""), "raw": b.get("raw", ""), "masks": b.get("masks", [])})
            except Exception:
                pass
    return {"device_id": device_id, "baselines": baselines}


@router.get("/api/baselines/{device_id}/{key}")
async def get_baseline_api(device_id: str, key: str):
    bl = load_baseline(device_id, key)
    if not bl:
        raise HTTPException(404, "Baseline not found")
    return bl


@router.post("/api/baselines/{device_id}/{key}")
async def save_baseline_api(device_id: str, key: str, data: dict):
    path = baseline_file_path(device_id, key)
    path.parent.mkdir(parents=True, exist_ok=True)
    saved = {"device_id": device_id, "key": key, "command": data.get("command", ""), "raw": data.get("raw", ""),
             "masks": data.get("masks", []), "captured_at": data.get("captured_at", datetime.now().isoformat()), "description": data.get("description", "")}
    save_json(path, saved)
    return saved


@router.delete("/api/baselines/{device_id}/{key}")
async def delete_baseline_api(device_id: str, key: str):
    path = baseline_file_path(device_id, key)
    if path.exists():
        path.unlink()
    return {"ok": True}


# 장비에서 Baseline 캡처 (현재 CLI 출력 저장)
@router.post("/api/baselines/{device_id}/{key}/capture")
async def capture_baseline(device_id: str, key: str, body: dict):
    devices = _devices()
    device = next((d for d in devices if str(d.get("id", "")) == device_id), None)
    if not device:
        raise HTTPException(404, "장비를 찾을 수 없습니다")
    command = body.get("command", "")
    if not command:
        raise HTTPException(400, "command 가 필요합니다")
    try:
        raw = await asyncio.to_thread(ssh_exec, device["ip"], int(device.get("port", 22)), device.get("username", ""), device.get("password", ""), command, device)
    except Exception as exc:
        raise HTTPException(500, f"명령 실행 실패: {exc}")
    path = baseline_file_path(device_id, key)
    saved = {"device_id": device_id, "key": key, "command": command, "raw": raw, "masks": body.get("masks", []),
             "captured_at": datetime.now().isoformat(), "description": body.get("description", "")}
    save_json(path, saved)
    return saved
